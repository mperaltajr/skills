"""
check_pptx_hygiene.py — Programmatic pre-pass for slide QC
==========================================================

Catches hygiene issues that don't need vision — they can be answered
deterministically by reading the PPTX XML. Runs BEFORE the visual QC
pass so the visual pass can focus on what only eyes can catch.

Detects:
  - Lorem ipsum / placeholder text residue
  - "[Insert X]" / "Subtitle goes here" / TODO / FIXME / XXX in text
  - Hidden slides leaking into the file
  - Comments left attached to slides
  - Speaker notes with non-substantive content
  - File name patterns that signal workspace dumps
    ("Final_final_v3.pptx", "Copy of...", "deck (3).pptx")

Footer / page-number presence is NOT detected here — it lives in the
visual QC pass instead, because most decks inherit those from the
slide master or layout and python-pptx doesn't see the inherited text
at the slide level. The visual pass reads the rendered PNG, where
inheritance is resolved, so it gives a reliable answer.

Outputs JSON to stdout — the SKILL invoker reads this and merges into
the unified Critical/Major/Advisory table. JSON shape:

  {
    "pptx": "<absolute path>",
    "slides": <int>,
    "violations": [
      {
        "slide": <int or null>,           # null for deck-level
        "severity": "Critical|Major|Advisory",
        "category": "<short tag>",
        "issue": "<one-line description>"
      },
      ...
    ]
  }

The script never modifies the PPTX.

Usage:
    py -3 check_pptx_hygiene.py "<path to deck.pptx>"
"""

from __future__ import annotations

import argparse
import json
import pathlib
import re
import sys

from pptx import Presentation


# ---------------------------------------------------------------------------
# Detection patterns
# ---------------------------------------------------------------------------

# Cross-skill placeholder contract.
#
# slide-builder/twins/helpers.py add_footer() emits these EXACT strings as
# INTENTIONAL presenter prompts when the caller passes footnote=None or
# source=None. Mirror of:
#   INTENTIONAL_FOOTNOTE_PLACEHOLDER = "[add footnote here or delete]"
#   INTENTIONAL_SOURCE_PLACEHOLDER   = "[add source here or delete]"
# in slide-builder/twins/helpers.py. Identity-matched, not regex-matched —
# do not generalize. If slide-builder changes the wording, this list must
# change in lockstep (and the cross-skill contract test should catch drift).
INTENTIONAL_PLACEHOLDER_STRINGS: tuple[str, ...] = (
    "[add footnote here or delete]",
    "[add source here or delete]",
)


LOREM_PATTERNS: list[tuple[str, str]] = [
    (r"\blorem\s+ipsum\b",         "Lorem ipsum"),
    (r"\bdolor\s+sit\s+amet\b",    "Lorem ipsum (dolor sit amet)"),
    (r"\[insert[^\]]{0,50}\]",     "[Insert ...] placeholder"),
    (r"\[your[^\]]{0,50}\]",       "[Your ...] placeholder"),
    (r"\bsubtitle\s+goes\s+here\b", "'Subtitle goes here' placeholder"),
    (r"\bplaceholder\s+text\b",    "'Placeholder text'"),
    (r"\bclick\s+to\s+add\b",      "'Click to add' template prompt"),
    (r"\bTODO\b",                  "TODO marker"),
    (r"\bFIXME\b",                 "FIXME marker"),
    (r"\bXXX\b",                   "XXX marker"),
    # Sweep gap from Pre-Workshop Scope Alignment build:
    # Hardline #2 placeholders (worker emits these when the brief doesn't
    # enumerate cell/list content, with a named source the user should
    # consult). These are NOT the intentional-presenter-prompt convention
    # (slide-builder/twins/helpers.py constants) — they mean unfinished
    # work and stay Critical. Visual disambiguation remains an open item.
    (r"\[[^\]]{1,60}\s+—\s+fill\s+from\b[^\]]{0,60}\]", "Hardline #2 '[X — fill from ...]' placeholder"),
    (r"\[[^\]]{1,30}\s+goes\s+here\]", "'[X goes here]' template prompt"),
    (r"\[\s*(TBD|TK|TKTK)\s*\]",        "[TBD] / [TK] marker"),
    (r"\b(TBD|TKTK)\b",                  "TBD / TKTK marker (bare)"),
]


def _is_intentional_placeholder(text: str) -> bool:
    """Return True if `text` contains a known slide-builder intentional
    presenter-prompt placeholder. Identity-match against the cross-skill
    contract — these strings are deliberately rendered for the presenter
    to fill or delete before showing the deck, not lorem-ipsum residue.
    """
    return any(ph in text for ph in INTENTIONAL_PLACEHOLDER_STRINGS)

FILENAME_BAD_PATTERNS: list[tuple[str, str]] = [
    (r"final[-_ ]?final",                                "'final_final' in filename"),
    (r"v\d+[-_ ]?final",                                 "'vN_final' pattern in filename"),
    (r"USE[-_ ]?THIS[-_ ]?ONE",                          "'USE_THIS_ONE' in filename"),
    (r"^copy\s+of\s+",                                   "starts with 'Copy of'"),
    (r"\s*\(\d+\)\.pptx$",                               "'(N).pptx' duplicate-file suffix"),
    # Standalone draft/wip/tmp/temp markers — surrounded by non-letters so
    # "template" / "templated" don't match "temp".
    (r"(?<![a-z])(draft|wip|tmp|temp)(?![a-z])",         "draft/wip/temp marker in filename"),
]

# Patterns indicating speaker notes are scratch, not substantive
NOTES_JUNK_PATTERNS = [
    r"\bTODO\b",
    r"\bFIXME\b",
    r"\bdraft\b",
    r"\bfix\s+this\b",
    r"\bscratch\b",
    r"\bWIP\b",
    r"\basdf+\b",
    r"\btest\s+test\b",
    r"\bxxxxx+",
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _gather_slide_text(slide) -> str:
    """Concatenate all run text in a slide for full-text scanning."""
    chunks: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text:
                    chunks.append(run.text)
    return " ".join(chunks)


def _slide_is_hidden(slide) -> bool:
    """Slides hidden in PowerPoint set `show="0"` on the sld element."""
    return slide._element.get("show") == "0"


def _slide_has_comments(slide) -> bool:
    """Check for comment relationships attached to the slide part."""
    try:
        for rel in slide.part.rels.values():
            if "comment" in (rel.reltype or "").lower():
                return True
    except Exception:
        pass
    return False


def _emu_to_in(emu) -> float | None:
    if emu is None:
        return None
    try:
        return int(emu) / 914400
    except (TypeError, ValueError):
        return None


# ---------------------------------------------------------------------------
# Checks
# ---------------------------------------------------------------------------

def check_filename(pptx_path: pathlib.Path) -> list[dict]:
    name = pptx_path.name
    name_lower = name.lower()
    for pat, label in FILENAME_BAD_PATTERNS:
        if re.search(pat, name_lower):
            return [{
                "slide": None,
                "severity": "Major",
                "category": "file-hygiene",
                "issue": f"Filename smells like a workspace dump ({label}): '{name}'. Rename before sharing.",
            }]
    return []


def check_placeholders_in_text(slide, slide_num: int) -> list[dict]:
    violations: list[dict] = []
    text = _gather_slide_text(slide)
    if not text:
        return violations

    # Intentional presenter prompts emitted by
    # slide-builder/twins/helpers.py add_footer() get an Advisory note,
    # not a Critical. They're the cross-skill convention: presenter is
    # expected to fill or delete in PowerPoint before showing the deck.
    if _is_intentional_placeholder(text):
        violations.append({
            "slide": slide_num,
            "severity": "Advisory",
            "category": "intentional-placeholder",
            "issue": (
                f"Slide {slide_num} renders a slide-builder intentional "
                f"presenter prompt (e.g., '[add footnote here or delete]'). "
                f"Fill or delete in PowerPoint before showing the deck."
            ),
        })
        # Strip the intentional placeholder strings before the Critical
        # regex pass so the existing LOREM_PATTERNS don't false-Critical
        # on a known-intentional substring.
        scrub = text
        for ph in INTENTIONAL_PLACEHOLDER_STRINGS:
            scrub = scrub.replace(ph, "")
        text = scrub

    for pat, label in LOREM_PATTERNS:
        if re.search(pat, text, re.IGNORECASE):
            violations.append({
                "slide": slide_num,
                "severity": "Critical",
                "category": "placeholder-residue",
                "issue": f"Placeholder text detected on Slide {slide_num}: {label}.",
            })
    return violations


def check_hidden_slides(prs) -> list[dict]:
    violations: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        if _slide_is_hidden(slide):
            violations.append({
                "slide": i,
                "severity": "Major",
                "category": "hidden-slide",
                "issue": f"Slide {i} is marked hidden — will not appear in presentation mode. Confirm intentional or unhide.",
            })
    return violations


def check_comments(prs) -> list[dict]:
    violations: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        if _slide_has_comments(slide):
            violations.append({
                "slide": i,
                "severity": "Major",
                "category": "comments-left",
                "issue": f"Slide {i} has reviewer comments still attached. Resolve or delete before sharing.",
            })
    return violations


def check_speaker_notes(slide, slide_num: int) -> list[dict]:
    if not slide.has_notes_slide:
        return []
    text = slide.notes_slide.notes_text_frame.text or ""
    if not text.strip():
        return []
    for pat in NOTES_JUNK_PATTERNS:
        if re.search(pat, text, re.IGNORECASE):
            return [{
                "slide": slide_num,
                "severity": "Major",
                "category": "speaker-notes",
                "issue": f"Slide {slide_num} speaker notes contain scratch content (matched '{pat}'). Clean up or delete before sharing.",
            }]
    return []


# Footer / page-number detection was attempted here as a python-pptx pass but
# produced too many false positives — most decks inherit the footer from the
# slide master or layout, and python-pptx does not see the inherited text at
# the slide level. The visual QC pass already catches missing footers reliably
# because it reads the rendered PNG, where master inheritance is resolved.
# So this check intentionally lives in the visual pass, not here.


# ---------------------------------------------------------------------------
# Entry
# ---------------------------------------------------------------------------

def run_all_checks(pptx_path: pathlib.Path) -> dict:
    prs = Presentation(str(pptx_path))
    violations: list[dict] = []

    violations.extend(check_filename(pptx_path))
    violations.extend(check_hidden_slides(prs))
    violations.extend(check_comments(prs))

    for i, slide in enumerate(prs.slides, start=1):
        violations.extend(check_placeholders_in_text(slide, i))
        violations.extend(check_speaker_notes(slide, i))

    # Sort: Critical first, then Major, then Advisory; within severity by slide number
    severity_order = {"Critical": 0, "Major": 1, "Advisory": 2}
    violations.sort(key=lambda v: (severity_order.get(v["severity"], 9),
                                    v["slide"] if v["slide"] is not None else 0))

    return {
        "pptx": str(pptx_path),
        "slides": len(prs.slides),
        "violations": violations,
    }


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Programmatic hygiene pre-pass for slide QC. Emits JSON to stdout."
    )
    parser.add_argument("pptx", help="Path to the PPTX file to inspect")
    args = parser.parse_args()

    pptx_path = pathlib.Path(args.pptx).resolve()
    if not pptx_path.exists():
        print(json.dumps({"error": f"File not found: {pptx_path}"}), file=sys.stderr)
        return 1
    if not pptx_path.is_file() or pptx_path.suffix.lower() != ".pptx":
        print(json.dumps({"error": f"Not a .pptx file: {pptx_path}"}), file=sys.stderr)
        return 1

    try:
        result = run_all_checks(pptx_path)
    except Exception as exc:
        print(json.dumps({"error": f"Failed to inspect {pptx_path.name}: {exc}"}), file=sys.stderr)
        return 1

    print(json.dumps(result, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
