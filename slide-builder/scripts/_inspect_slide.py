"""Quick shape inspector for one slide of a PPTX."""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if len(sys.argv) != 3:
    print("Usage: python _inspect_slide.py <pptx_path> <slide_number>")
    sys.exit(1)

pptx_path = sys.argv[1]
slide_num = int(sys.argv[2])

prs = Presentation(pptx_path)
slide = prs.slides[slide_num - 1]

print(f"Slide {slide_num} of {pptx_path}")
print(f"Slide width x height: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f}")
print(f"Layout name: {slide.slide_layout.name}")
print(f"Total shapes on slide: {len(slide.shapes)}")
print()
print(f"{'#':>3}  {'NAME':<30} {'TYPE':<25} {'POS (in)':<20} {'SIZE':<18} {'TEXT'}")
print("-" * 130)
for i, s in enumerate(slide.shapes):
    try:
        type_name = str(s.shape_type).replace("MSO_SHAPE_TYPE.", "")
    except Exception:
        type_name = "?"
    try:
        pos = f"({s.left.inches:.2f}, {s.top.inches:.2f})"
    except Exception:
        pos = "?"
    try:
        size = f"{s.width.inches:.2f} x {s.height.inches:.2f}"
    except Exception:
        size = "?"
    text = ""
    if s.has_text_frame:
        text = s.text_frame.text[:40].replace("\n", " | ")
    print(f"{i:>3}  {(s.name or '?'):<30} {type_name:<25} {pos:<20} {size:<18} {text}")

print()
print(f"Layout shapes ({len(slide.slide_layout.shapes)}):")
for i, s in enumerate(slide.slide_layout.shapes):
    try:
        type_name = str(s.shape_type).replace("MSO_SHAPE_TYPE.", "")
    except Exception:
        type_name = "?"
    text = ""
    if s.has_text_frame:
        text = s.text_frame.text[:40].replace("\n", " | ")
    print(f"  {i}: {s.name} ({type_name}) {text}")
