"""
build_slide.py — v3.3 prototype: HTML-mockup-as-spec build path with optional layout-aware mode.

The Builder reads the approved HTML mockup directly and emits a PPTX. The
mockup is the spec — no blueprint JSON, no body kinds, no content JSON.

Two modes (controlled per-slide by mockup attributes):
  - BLANK MODE (default): pick a "Blank" layout from the template, draw
    everything from the mockup on top.
  - LAYOUT-AWARE MODE: when the mockup root has data-layout-master and
    data-layout-index attributes, the builder uses that specific layout
    from the template and fills its placeholders by matching mockup
    elements that carry data-placeholder="name". Untagged elements are
    rendered as overlay shapes on top of the layout.

How it works:
  1. Render the mockup HTML in headless Chromium (playwright).
  2. For each picked slide-option, find its <div data-slide-index data-option> root.
  3. Walk every visible descendant (bbox + computed style + inline runs).
  4. Pick the slide layout: layout-aware if root has data-layout-*, else blank.
  5. Fill data-placeholder elements into matching layout placeholders.
  6. Render remaining elements as shapes/textboxes/pictures on top.
  7. Theme-color binding: mockup hex matching the template's theme slot
     emits as <a:schemeClr val="accentN"/> for brand inheritance.
  8. Theme-font binding: every text run gets <a:latin typeface=".."/> from
     the template's theme so fonts work even when not installed locally.
  9. Chart screenshot: data-chart="true" elements get screenshotted from the
     HTML and inserted as PNG. data-chart-data attribute populates a
     ThinkCell-compatible xlsx companion.

Modes:
    Build a deck:
        python build_slide.py --mockup mockups.html --picks 1A,2A \
            --target deck.pptx --client-template template.pptx

    Print theme (colors + fonts) as JSON for slide-helper to inject into mockup:
        python build_slide.py --print-theme template.pptx

    Catalog the template's layouts as JSON (for layout-aware mockup authoring):
        python build_slide.py --catalog-layouts template.pptx

Mockup attributes recognized:
    Slide root <div>:
        data-slide-index="3" data-option="A"           (required)
        data-layout-master="6" data-layout-index="11"  (optional, layout-aware mode)
        data-layout-name="..."                         (informational only)

    Any element:
        data-role="source|footnote|page-number"  -> filled into master placeholder
        data-placeholder="title|body|stat-1|..."  -> filled into layout placeholder
        data-chart="true"                          -> screenshot as picture
        data-chart-data='{...JSON...}'             -> appended to xlsx companion
        data-chart-title="..."                     -> sheet name in xlsx

Dependencies:
    pip install python-pptx lxml playwright openpyxl
    playwright install chromium
"""

import argparse
import json
import sys
import re
import zipfile

# Windows console defaults to cp1252; force UTF-8 so Unicode in print() doesn't crash.
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# --- Constants -----------------------------------------------------------

PX_TO_EMU = 9525  # at 96 DPI; standard MBB canvas is 1280x720 px
SLIDE_W_PX = 1280
SLIDE_H_PX = 720
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

def _emu(px_value):
    return Emu(int(px_value * PX_TO_EMU))


# --- Theme color extraction ---------------------------------------------

def extract_theme_fonts(template_path):
    """Read the client template's theme1.xml fontScheme. Return dict
    {"major": "FedEx Sans Bold", "minor": "FedEx Sans"} — the major/heading
    and minor/body font names defined by the template's theme.

    These are the fonts PowerPoint will use when a run references the theme
    font tokens (+mj-lt / +mn-lt). We also use them as the explicit
    font-family that gets stamped onto every run, so the PPTX displays
    correctly even when the rendering machine doesn't have the font
    available for the HTML mockup preview."""
    out = {"major": None, "minor": None}
    with zipfile.ZipFile(template_path, "r") as z:
        theme_xml = None
        for n in z.namelist():
            if n == "ppt/theme/theme1.xml":
                theme_xml = z.read(n); break
        if theme_xml is None:
            for n in z.namelist():
                if n.startswith("ppt/theme/") and n.endswith(".xml"):
                    theme_xml = z.read(n); break
    if theme_xml is None:
        return out
    root = etree.fromstring(theme_xml)
    fontScheme = root.find(f".//{{{NS_A}}}fontScheme")
    if fontScheme is None:
        return out
    major = fontScheme.find(f"{{{NS_A}}}majorFont")
    minor = fontScheme.find(f"{{{NS_A}}}minorFont")
    if major is not None:
        latin = major.find(f"{{{NS_A}}}latin")
        if latin is not None:
            out["major"] = latin.get("typeface")
    if minor is not None:
        latin = minor.find(f"{{{NS_A}}}latin")
        if latin is not None:
            out["minor"] = latin.get("typeface")
    return out


def extract_theme_hex(template_path):
    """Read the client template's theme1.xml clrScheme. Return dict
    {"accent1": "4D148C", ..., "dk1": "000000", "lt1": "FFFFFF"}."""
    out = {}
    with zipfile.ZipFile(template_path, "r") as z:
        theme_xml = None
        for n in z.namelist():
            if n == "ppt/theme/theme1.xml":
                theme_xml = z.read(n); break
        if theme_xml is None:
            for n in z.namelist():
                if n.startswith("ppt/theme/") and n.endswith(".xml"):
                    theme_xml = z.read(n); break
    if theme_xml is None:
        return out
    root = etree.fromstring(theme_xml)
    clrScheme = root.find(f".//{{{NS_A}}}clrScheme")
    if clrScheme is None:
        return out
    for child in clrScheme:
        local = etree.QName(child).localname
        for c in child:
            cl = etree.QName(c).localname
            if cl == "srgbClr":
                out[local] = c.get("val", "").upper()
            elif cl == "sysClr":
                out[local] = (c.get("lastClr") or "000000").upper()
    return out


def make_color_resolver(theme_hex):
    """Return a function (hex_str) -> {"kind": "schemeClr"|"srgbClr", "val": ...}.

    For shape fills and run colors, we emit literal srgbClr (hex) values to
    ensure rendering correctness across PowerPoint, LibreOffice, and other
    OOXML renderers. Theme binding via <a:schemeClr> is sound in spec but
    rendering engines vary in resolving non-accent slots (dk2, lt2, etc.) —
    LibreOffice in particular drops the fill silently in some cases.

    Tradeoff: brand colors won't auto-swap when the template changes
    (the mockup encodes literal hex, not theme tokens). To preserve
    brand-swap, run --print-theme on the new template, update the
    mockup CSS hex values to the new brand colors, and rebuild.
    """
    def resolve(h):
        if not h:
            return None
        h = h.upper()
        return {"kind": "srgbClr", "val": h}
    return resolve


# --- CSS color parsing --------------------------------------------------

def parse_css_color(value):
    """Parse 'rgb(r,g,b)', 'rgba(r,g,b,a)', '#hex' to (hex, alpha) where
    alpha is 0..1. Returns (None, 1.0) if not parseable or transparent."""
    if not value:
        return (None, 1.0)
    v = value.strip().lower()
    if v in ("transparent", "inherit", "initial", "unset", "none"):
        return (None, 0.0)
    m = re.match(r"^#([0-9a-f]{3})$", v)
    if m:
        h = m.group(1)
        return (f"{h[0]*2}{h[1]*2}{h[2]*2}".upper(), 1.0)
    m = re.match(r"^#([0-9a-f]{6})([0-9a-f]{2})?$", v)
    if m:
        h = m.group(1).upper()
        if m.group(2):
            a = int(m.group(2), 16) / 255.0
            return (h, a)
        return (h, 1.0)
    m = re.match(r"^rgba?\(\s*([\d.]+)\s*,\s*([\d.]+)\s*,\s*([\d.]+)\s*(?:,\s*([\d.]+)\s*)?\)$", v)
    if m:
        r, g, b = int(float(m.group(1))), int(float(m.group(2))), int(float(m.group(3)))
        a = float(m.group(4)) if m.group(4) is not None else 1.0
        return (f"{r:02X}{g:02X}{b:02X}", a)
    return (None, 1.0)


def is_visible_background(bg_value):
    """True if the bg color is non-transparent enough to render."""
    h, a = parse_css_color(bg_value)
    return h is not None and a > 0.01


def _extract_gradient_color(bg_image):
    """Return the first solid color stop from a CSS gradient string, or None.

    Gradients are stored in backgroundImage, not backgroundColor. Chrome
    returns rgba(0,0,0,0) for backgroundColor on gradient elements, so the
    builder falls back to this to avoid rendering a transparent rectangle.
    The gradient is approximated as its first stop color — lossy but correct.
    """
    if not bg_image or "gradient" not in bg_image.lower():
        return None
    m = re.search(r"(#[0-9a-fA-F]{3,8}|rgba?\([^)]+\))", bg_image)
    return m.group(1) if m else None


def _effective_bg(el):
    """Return (hex, alpha) for an element's background.

    Prefers solid backgroundColor; falls back to the first gradient stop
    when backgroundColor is transparent (the common gradient case).
    Returns (None, 0) when the element has no visible background.
    """
    hex_, alpha = parse_css_color(el.get("bgColor"))
    if hex_ and alpha > 0.01:
        return hex_, alpha
    grad_color = _extract_gradient_color(el.get("bgImage") or "")
    if grad_color:
        hex_, alpha = parse_css_color(grad_color)
        if hex_ and alpha > 0.01:
            return hex_, alpha
    return None, 0


def snap_pt(raw_pt, minimum_pt=10.0):
    """Snap a raw pt value to the nearest 0.5pt increment and enforce a floor.

    Body text floor: 10pt. Footer/source/chart-annotation floor: 8pt.
    Use _font_floor(el) to select the right floor before calling parse_px_to_pt.

    Snap table (common CSS px inputs, body floor):
      12px -> 9.0pt raw  -> 10.0pt (body floor)
      13px -> 9.75pt raw -> 10.0pt (body floor)
      14px -> 10.5pt raw -> 10.5pt
      16px -> 12.0pt raw -> 12.0pt
      18px -> 13.5pt raw -> 13.5pt
      20px -> 15.0pt raw -> 15.0pt
      24px -> 18.0pt raw -> 18.0pt
    """
    snapped = round(raw_pt * 2) / 2   # nearest 0.5pt
    return max(snapped, minimum_pt)


def _font_floor(el):
    """Return the minimum pt floor for this element.

    Footer text, source lines, and bottom-zone caption text are intentionally
    small (8-9pt in consulting decks). Everything else uses the 10pt body floor.
    """
    role = (el.get("role") or "").lower()
    if role in ("footer", "source", "caption"):
        return 8.0
    y = el.get("y") or 0
    if y > SLIDE_H_PX - 60:   # bottom 60px = footer zone
        return 8.0
    return 10.0


def parse_px_to_pt(value, minimum_pt=10.0):
    """CSS px font-size to pt: pt = px * 0.75 (96 DPI assumption).

    Snaps to nearest 0.5pt and enforces the given floor. Pass
    minimum_pt=_font_floor(el) at call sites for context-aware flooring.
    """
    if not value:
        return None
    m = re.match(r"^(-?\d+(?:\.\d+)?)px$", value.strip())
    if not m:
        return None
    raw_pt = float(m.group(1)) * 0.75
    snapped = snap_pt(raw_pt, minimum_pt=minimum_pt)
    if raw_pt < minimum_pt:
        print(f"  WARNING: font-size {value} -> {raw_pt:.2f}pt raw -> {snapped}pt (floored to {minimum_pt}pt). "
              f"Use 14px+ for body text, 11px+ for footer/caption text.")
    return snapped


# --- Playwright DOM walker ----------------------------------------------

# JS that runs in the page and returns a list of visible element descriptors
# for the slide root passed in as `slideEl`.
DOM_WALKER_JS = r"""
(slideEl) => {
  // Get computed style + bbox for each visible descendant.
  // Returns [{idx, parentIdx, tag, x, y, w, h, ...style..., directText, runs, role, isChart}]
  const slideRect = slideEl.getBoundingClientRect();
  const out = [];
  let counter = 0;

  function getDirectText(el) {
    // Concatenate text from this element's text nodes + INLINE-displayed
    // children (so a div with a <span> inline returns the full text).
    // Block-level children are skipped (they have their own bbox).
    // Whitespace is collapsed per CSS white-space:normal rules.
    let s = '';
    for (const node of el.childNodes) {
      if (node.nodeType === Node.TEXT_NODE) {
        s += node.nodeValue.replace(/\s+/g, ' ');
      } else if (node.nodeType === Node.ELEMENT_NODE) {
        if (node.tagName === 'BR') { s += '\n'; continue; }
        const cs = getComputedStyle(node);
        if (cs.display === 'inline' || cs.display === 'inline-block') {
          s += getDirectText(node);
        }
      }
    }
    // Trim whitespace at element boundaries (matches browser rendering for
    // block-level containers — leading/trailing ws is invisible).
    return s.replace(/^[ \t]+/, '').replace(/[ \t]+$/, '');
  }

  function getInlineRuns(el) {
    // Yield runs of (text, style) for the element's content. Each run is a
    // contiguous span of text under one inline styling context.
    // Whitespace is collapsed per CSS white-space:normal rules.
    const runs = [];
    for (const node of el.childNodes) {
      if (node.nodeType === Node.TEXT_NODE) {
        const t = node.nodeValue.replace(/\s+/g, ' ');
        if (t) runs.push({ text: t, color: null, fontWeight: null, fontStyle: null });
      } else if (node.nodeType === Node.ELEMENT_NODE) {
        if (node.tagName === 'BR') {
          runs.push({ text: '\n', color: null, fontWeight: null, fontStyle: null });
          continue;
        }
        const cs = getComputedStyle(node);
        if (cs.display !== 'inline' && cs.display !== 'inline-block') continue;
        const childRuns = getInlineRuns(node);
        for (const r of childRuns) {
          runs.push({
            text: r.text,
            color: r.color || cs.color,
            fontWeight: r.fontWeight || cs.fontWeight,
            fontStyle: r.fontStyle || cs.fontStyle,
            fontSize: r.fontSize || cs.fontSize,
          });
        }
      }
    }
    // Trim whitespace at the block boundary: strip leading ws from first
    // run, trailing ws from last run. Matches browser's whitespace handling
    // for block-level containers.
    if (runs.length > 0) {
      runs[0].text = runs[0].text.replace(/^[ \t]+/, '');
      runs[runs.length - 1].text = runs[runs.length - 1].text.replace(/[ \t]+$/, '');
      // Drop any runs that became empty after trimming
      for (let i = runs.length - 1; i >= 0; i--) {
        if (runs[i].text === '') runs.splice(i, 1);
      }
    }
    return runs;
  }

  function isBlockElement(el) {
    const tag = el.tagName;
    if (tag === 'SCRIPT' || tag === 'STYLE' || tag === 'BR') return false;
    const cs = getComputedStyle(el);
    if (cs.display === 'none' || cs.visibility === 'hidden') return false;
    return true;
  }

  function visit(el, parentIdx) {
    if (!isBlockElement(el)) return -1;
    const cs = getComputedStyle(el);
    // SVG elements have display:inline by default, but we want to capture
    // them as chart screenshots, so don't fold them into parent runs.
    if (cs.display === 'inline' && el.tagName.toLowerCase() !== 'svg') {
      // Inline elements are folded into parent's runs
      return -1;
    }
    const rect = el.getBoundingClientRect();
    if (rect.width === 0 || rect.height === 0) {
      // Still recurse — children might be visible
      // (some empty wrappers contain absolutely-positioned children)
    }
    const idx = counter++;
    const parentRect = el.parentElement ? el.parentElement.getBoundingClientRect() : null;
    const parentCS = el.parentElement ? getComputedStyle(el.parentElement) : null;

    // Table detection: capture table structure so the builder can emit a
    // real PPTX table instead of a bunch of stacked textboxes.
    let tableData = null;
    if (el.tagName.toLowerCase() === 'table') {
      tableData = { rows: [] };
      const trs = el.querySelectorAll('tr');
      trs.forEach(tr => {
        const trCS = getComputedStyle(tr);
        const row = { bg: trCS.backgroundColor, cells: [] };
        const cells = tr.querySelectorAll('th, td');
        cells.forEach(cell => {
          const ccs = getComputedStyle(cell);
          const cellRect = cell.getBoundingClientRect();
          row.cells.push({
            isHeader: cell.tagName.toLowerCase() === 'th',
            text: cell.innerText || cell.textContent || '',
            runs: getInlineRuns(cell),
            bg: ccs.backgroundColor,
            color: ccs.color,
            fontSize: ccs.fontSize,
            fontWeight: ccs.fontWeight,
            fontFamily: ccs.fontFamily,
            textAlign: ccs.textAlign,
            colspan: parseInt(cell.getAttribute('colspan') || '1'),
            rowspan: parseInt(cell.getAttribute('rowspan') || '1'),
            width: cellRect.width,
          });
        });
        tableData.rows.push(row);
      });
    }

    const item = {
      idx: idx,
      parentIdx: parentIdx,
      tag: el.tagName.toLowerCase(),
      x: rect.x - slideRect.x,
      y: rect.y - slideRect.y,
      w: rect.width,
      h: rect.height,
      bgColor: cs.backgroundColor,
      bgImage: cs.backgroundImage,
      color: cs.color,
      fontSize: cs.fontSize,
      fontWeight: cs.fontWeight,
      fontFamily: cs.fontFamily,
      fontStyle: cs.fontStyle,
      textAlign: cs.textAlign,
      textTransform: cs.textTransform,
      letterSpacing: cs.letterSpacing,
      lineHeight: cs.lineHeight,
      opacity: cs.opacity,
      borderLeftWidth: cs.borderLeftWidth,
      borderLeftColor: cs.borderLeftColor,
      borderLeftStyle: cs.borderLeftStyle,
      borderTopWidth: cs.borderTopWidth,
      borderTopColor: cs.borderTopColor,
      borderTopStyle: cs.borderTopStyle,
      directText: getDirectText(el),
      runs: getInlineRuns(el),
      role: el.getAttribute('data-role'),
      // Auto-detect SVG as chart: walker can't walk SVG children individually,
      // so screenshot the whole SVG element as a PNG. data-chart="true" still
      // works for non-SVG chart containers.
      //
      // EXCEPTION: full-slide SVG overlays (pointer-events:none AND covering
      // ≥90% of slide dimensions) are NOT screenshotted — doing so captures the
      // entire rendered slide as a flat raster placed on top of all PPTX shapes,
      // producing a double-layer effect. These overlays (flowchart connector SVGs)
      // are skipped entirely in Phase B; arrows must be added manually in PowerPoint.
      // See known-issues-and-improvements.md Issue #10 Option B for the long-term fix.
      isChart: (el.tagName.toLowerCase() === 'svg'
               && !(cs.pointerEvents === 'none'
                    && rect.width  >= slideRect.width  * 0.9
                    && rect.height >= slideRect.height * 0.9))
               || el.getAttribute('data-chart') === 'true'
               || el.getAttribute('data-chart') === '1',
      chartData: el.getAttribute('data-chart-data') || null,
      chartTitle: el.getAttribute('data-chart-title') || null,
      placeholder: el.getAttribute('data-placeholder') || null,
      whitespace: cs.whiteSpace || null,
      tableData: tableData,
      // Parent geometry — used to detect flex-centered children that need parent's width
      parentX: parentRect ? (parentRect.x - slideRect.x) : null,
      parentY: parentRect ? (parentRect.y - slideRect.y) : null,
      parentW: parentRect ? parentRect.width : null,
      parentH: parentRect ? parentRect.height : null,
      parentDisplay: parentCS ? parentCS.display : null,
      parentJustifyContent: parentCS ? parentCS.justifyContent : null,
      parentAlignItems: parentCS ? parentCS.alignItems : null,
    };
    out.push(item);
    if (item.isChart) {
      // Don't recurse into chart contents — the screenshot captures it.
      return idx;
    }
    if (tableData) {
      // Don't recurse into table cells — the table data captures them.
      return idx;
    }
    for (const child of el.children) {
      visit(child, idx);
    }
    return idx;
  }

  visit(slideEl, -1);
  return out;
}
"""


def render_and_walk(html_path, slide_index, option):
    """Render html_path in playwright, find the slide-option root, walk its
    descendants. Returns (elements_list, screenshot_fn) where screenshot_fn
    takes (element_idx, out_path) and saves a high-DPI screenshot of that
    element's bbox.

    Returns (None, None) if the option isn't found.
    """
    from playwright.sync_api import sync_playwright

    p = sync_playwright().start()
    browser = p.chromium.launch()
    page = browser.new_page(viewport={"width": SLIDE_W_PX + 40, "height": SLIDE_H_PX + 40},
                            device_scale_factor=2.0)
    page.goto(f"file://{Path(html_path).resolve()}")
    # Find the slide-option root
    selector = f'[data-slide-index="{slide_index}"][data-option="{option}"]'
    locator = page.locator(selector)
    if locator.count() == 0:
        browser.close(); p.stop()
        return (None, None, None, None, None)
    handle = locator.first.element_handle()
    elements = page.evaluate(DOM_WALKER_JS, handle)
    # Capture layout-aware attrs from the slide root
    layout_attrs = page.evaluate("""(el) => ({
        master: el.getAttribute('data-layout-master'),
        index:  el.getAttribute('data-layout-index'),
        name:   el.getAttribute('data-layout-name'),
    })""", handle)
    # Build a screenshot helper that captures a specific element's bounding box
    # We'll use the index assigned by the walker to find the element again.
    # IMPORTANT: this MUST use the same DOM-walking rules as DOM_WALKER_JS,
    # otherwise the counter goes out of sync and the wrong element gets shot.
    def screenshot_element(idx, out_path):
        # Find the element by walking the DOM the same way and counting
        target_handle = page.evaluate_handle(f"""(slideEl) => {{
            let counter = 0;
            let target = null;
            function isBlock(el) {{
                const tag = el.tagName;
                if (tag === 'SCRIPT' || tag === 'STYLE' || tag === 'BR') return false;
                const cs = getComputedStyle(el);
                if (cs.display === 'none' || cs.visibility === 'hidden') return false;
                // SVG has display:inline by default but we walk into it
                if (cs.display === 'inline' && el.tagName.toLowerCase() !== 'svg') return false;
                return true;
            }}
            function isChartLike(el) {{
                return el.tagName.toLowerCase() === 'svg'
                    || el.getAttribute('data-chart') === 'true'
                    || el.getAttribute('data-chart') === '1';
            }}
            function visit(el) {{
                if (!isBlock(el)) return;
                if (counter === {idx}) {{ target = el; return; }}
                counter++;
                if (target) return;
                // Don't recurse into chart-like elements (matches main walker)
                if (isChartLike(el)) return;
                // Don't recurse into <table> children (matches main walker)
                if (el.tagName.toLowerCase() === 'table') return;
                for (const ch of el.children) {{ visit(ch); if (target) return; }}
            }}
            visit(slideEl);
            return target;
        }}""", handle)
        if target_handle:
            target_handle.screenshot(path=out_path)
            return True
        return False

    return (elements, screenshot_element, browser, p, layout_attrs)


# --- Layout / blank slide selection -------------------------------------

def catalog_layouts(template_path):
    """Walk every master and every layout in the template; return a list
    suitable for JSON dump. Used by --catalog-layouts to give slide-helper
    the menu of corporate-approved layouts available in this template.

    Each entry:
        {
            "master_index": int,    # 1-based, for display
            "layout_index": int,    # 1-based within master
            "name": str,
            "n_shapes": int,
            "n_placeholders": int,
            "placeholders": [
                {"idx": int, "type": str, "name": str,
                 "x_emu": int, "y_emu": int, "w_emu": int, "h_emu": int,
                 "x_px": float, "y_px": float, "w_px": float, "h_px": float}
            ],
            "background": "light|dark|unknown"
        }
    """
    pres = Presentation(str(template_path))
    out = []
    for mi, master in enumerate(pres.slide_masters, start=1):
        for li, layout in enumerate(master.slide_layouts, start=1):
            placeholders = []
            for shape in layout.shapes:
                if not shape.is_placeholder:
                    continue
                ph = shape.placeholder_format
                # Convert EMU → px for mockup authoring (96 dpi assumption)
                try:
                    x_emu, y_emu = int(shape.left or 0), int(shape.top or 0)
                    w_emu, h_emu = int(shape.width or 0), int(shape.height or 0)
                except Exception:
                    x_emu = y_emu = w_emu = h_emu = 0
                placeholders.append({
                    "idx": ph.idx if ph.idx is not None else -1,
                    "type": str(ph.type) if ph.type else None,
                    "name": shape.name,
                    "x_emu": x_emu, "y_emu": y_emu, "w_emu": w_emu, "h_emu": h_emu,
                    "x_px": round(x_emu / PX_TO_EMU, 1),
                    "y_px": round(y_emu / PX_TO_EMU, 1),
                    "w_px": round(w_emu / PX_TO_EMU, 1),
                    "h_px": round(h_emu / PX_TO_EMU, 1),
                })
            bg = "light" if _layout_has_light_bg(layout) else "dark"
            out.append({
                "master_index": mi,
                "layout_index": li,
                "name": layout.name or "",
                "n_shapes": len(layout.shapes),
                "n_placeholders": len(placeholders),
                "placeholders": placeholders,
                "background": bg,
            })
    return out


def _layout_has_light_bg(layout):
    bg = layout.element.find(qn("p:cSld")).find(qn("p:bg"))
    if bg is None:
        return True
    bgPr = bg.find(qn("p:bgPr"))
    if bgPr is not None:
        sf = bgPr.find(qn("a:solidFill"))
        if sf is not None and len(sf):
            clr = sf[0]
            local = etree.QName(clr).localname
            if local == "schemeClr":
                return clr.get("val") in {"bg1", "bg2", "lt1", "lt2"}
            if local == "srgbClr":
                hex_val = clr.get("val", "")
                try:
                    r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
                    return (r + g + b) / 3 > 200
                except Exception:
                    return False
        return False
    bgRef = bg.find(qn("p:bgRef"))
    if bgRef is not None and len(bgRef):
        clr = bgRef[0]
        if etree.QName(clr).localname == "schemeClr":
            return clr.get("val") in {"bg1", "bg2", "lt1", "lt2"}
    return False


def _master_decoration_count(master):
    """Count non-placeholder decorative shapes on the master itself.
    High count means the master has logos/branding bars baked in that
    will show through any layout selected from this master."""
    return sum(1 for s in master.shapes if not s.is_placeholder)


def find_blank_layout(master_or_pres):
    """Find the cleanest blank-style layout. Searches across all masters
    when given a Presentation, or within one master when given a Master.
    Prefers (in order):
      1. A layout literally named "Blank" with light background
      2. A layout with "blank" in its name
      3. The layout with the fewest total decorations (master + layout)
    """
    candidates = []
    try:
        for m in master_or_pres.slide_masters:
            for L in m.slide_layouts:
                candidates.append((m, L))
    except AttributeError:
        for L in master_or_pres.slide_layouts:
            candidates.append((master_or_pres, L))

    light = [(m, L) for m, L in candidates if _layout_has_light_bg(L)]
    cands = light if light else candidates

    def cost(ml):
        m, L = ml
        return _master_decoration_count(m) + len([s for s in L.shapes if not s.is_placeholder])

    exact = [(m, L) for m, L in cands
             if (L.name or "").strip().lower() in ("blank", "blank slide")]
    if exact:
        return min(exact, key=cost)[1]

    sub = [(m, L) for m, L in cands if "blank" in (L.name or "").lower()]
    if sub:
        return min(sub, key=cost)[1]

    return min(cands, key=cost)[1]


def find_named_layout(pres, master_index, layout_index):
    """Look up a specific layout by 1-based master and layout indexes
    (matching the catalog_layouts output). Returns None if out of range."""
    masters = list(pres.slide_masters)
    if master_index < 1 or master_index > len(masters):
        return None
    layouts = list(masters[master_index - 1].slide_layouts)
    if layout_index < 1 or layout_index > len(layouts):
        return None
    return layouts[layout_index - 1]


def add_blank_slide(pres):
    return pres.slides.add_slide(find_blank_layout(pres))


def add_layout_slide(pres, master_index, layout_index):
    """Add a slide using a specific layout. Falls back to blank if the
    indexes are invalid."""
    layout = find_named_layout(pres, master_index, layout_index)
    if layout is None:
        return add_blank_slide(pres)
    return pres.slides.add_slide(layout)


def _pre_build_checks(html_path: Path) -> None:
    """Scan the mockup HTML for known build-breaking issues before running Playwright.

    Prints warnings for soft issues; exits with an error for hard blockers.
    """
    import re as _re
    html = html_path.read_text(encoding="utf-8", errors="replace")
    fatal = []
    warnings = []

    # --- Canvas size check ---------------------------------------------------
    # Detect 4:3 legacy canvas (1024×768). Any slide div wider than 0 but not
    # 1280px wide is suspicious; 1024 is the most common wrong value.
    if _re.search(r"width\s*:\s*1024\s*px", html, _re.IGNORECASE):
        fatal.append(
            "Canvas is 1024px wide (4:3 legacy format). "
            "Change every slide div to width:1280px; height:720px before building. "
            "The template is 16:9 (1280x720); a 4:3 canvas will cause content to "
            "bleed into the template's right-panel area."
        )
    if _re.search(r"height\s*:\s*768\s*px", html, _re.IGNORECASE):
        fatal.append(
            "Canvas is 768px tall (4:3 legacy format). "
            "Change every slide div to height:720px before building."
        )

    # --- data-placeholder misuse check ---------------------------------------
    # data-placeholder="title" must ONLY appear on the actual slide title element.
    # Common misuse: applied to quote blocks, subtitles, body content divs.
    ph_matches = _re.findall(
        r'<([^>]+?)\s+data-placeholder=["\']title["\'][^>]*>([^<]{0,80})',
        html, _re.IGNORECASE
    )
    for tag_attrs, content_preview in ph_matches:
        tag = tag_attrs.strip().split()[0].lower()
        # blockquote, p, span inside body, or very long content are suspicious
        if tag in ("blockquote", "p", "span", "figcaption") or len(content_preview.strip()) > 60:
            warnings.append(
                f'data-placeholder="title" on <{tag}> element '
                f'(preview: {content_preview.strip()[:50]!r}). '
                f'Only the actual slide title should carry this attribute — '
                f'all other text must use explicit position:absolute coordinates. '
                f'Misuse routes text into the layout title placeholder at an '
                f'unexpected position, which may render it invisible.'
            )

    # --- Unsupported CSS features --------------------------------------------
    if _re.search(r"border-radius\s*:\s*50%", html):
        warnings.append(
            "border-radius:50% detected. The builder renders divs as rectangles "
            "regardless of border-radius — circles must use SVG <circle> elements."
        )
    if _re.search(r"transform\s*:\s*translate", html, _re.IGNORECASE):
        warnings.append(
            "CSS transform:translate detected. The builder does not support CSS "
            "transforms — use explicit top:/left: coordinates instead."
        )

    # --- Native <table> elements ---------------------------------------------
    # Native <table>/<tr>/<td> return h=0 in Playwright inside positioned
    # containers. place_table() silently skips h<=0, producing a blank slide.
    if _re.search(r"<table[\s>]", html, _re.IGNORECASE):
        fatal.append(
            "Native <table> element detected. Playwright returns h=0 for <table> "
            "inside position:absolute containers — place_table() will silently skip "
            "it and the slide area will be blank with no error. "
            "Convert all tabular data to div-based rows (display:flex per row) "
            "before building. See known-issues-and-improvements.md §18."
        )

    # --- inset:0 shorthand ---------------------------------------------------
    if _re.search(r"\binset\s*:\s*0", html, _re.IGNORECASE):
        warnings.append(
            "inset:0 detected. The builder does not parse the CSS `inset` shorthand. "
            "Replace with explicit coordinates: top:0; left:0; right:0; bottom:0. "
            "Without this fix, background fill shapes will have zero/wrong dimensions "
            "and the layout template background will bleed through."
        )

    # --- Positioned div with left but no right/width -------------------------
    # A div with only `left` and no `right` or `width` will produce a very
    # narrow text shape, potentially narrower than the layout placeholder below it,
    # causing stray placeholder content (bullets, text) to bleed through.
    positioned_divs = _re.findall(
        r'<div[^>]+style\s*=\s*["\']([^"\']*)["\']',
        html, _re.IGNORECASE
    )
    for style_val in positioned_divs:
        style_lower = style_val.lower()
        has_left = _re.search(r'\bleft\s*:', style_lower)
        has_right = _re.search(r'\bright\s*:', style_lower)
        has_width = _re.search(r'\bwidth\s*:', style_lower)
        if has_left and not has_right and not has_width:
            # Extract a snippet for context
            preview = style_val[:80].strip()
            warnings.append(
                f"Positioned div has `left` but no `right` or `width` "
                f"(style: {preview!r}). The builder will create a narrow shape that "
                f"may not cover the layout placeholder zone, causing placeholder "
                f"content (bullets, stray text) to bleed through. Add an explicit "
                f"`width` or `right` value."
            )

    # --- Report --------------------------------------------------------------
    if warnings:
        print("\nPRE-BUILD WARNINGS:")
        for w in warnings:
            print(f"  ⚠  {w}")
    if fatal:
        print("\nPRE-BUILD ERRORS (build aborted):", file=sys.stderr)
        for e in fatal:
            print(f"  ✗  {e}", file=sys.stderr)
        sys.exit(1)
    if not warnings and not fatal:
        print("Pre-build checks: OK")


def purge_sections(pres):
    """Remove all PowerPoint sections inherited from the template.

    Sections are stored as a <p14:sectionLst> extension inside the
    presentation's <p:extLst>. When building a new deck we never want
    template section headers showing in the slide panel.
    """
    SECTION_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"
    prs_elm = pres.part.presentation._element
    ext_lst = prs_elm.find(qn("p:extLst"))
    if ext_lst is None:
        return
    for ext in list(ext_lst):
        if ext.get("uri") == SECTION_URI:
            ext_lst.remove(ext)
    # If extLst is now empty, remove it too (keeps XML tidy)
    if len(ext_lst) == 0:
        prs_elm.remove(ext_lst)


def purge_starter_slides(pres):
    sld_id_lst = pres.slides._sldIdLst
    pres_part = pres.part
    package = pres_part.package
    targets = []
    for s in list(sld_id_lst):
        rId = s.get(qn("r:id"))
        try:
            slide_part = pres_part.related_part(rId)
            targets.append((rId, slide_part.partname, s))
        except (KeyError, AttributeError):
            targets.append((rId, None, s))
    for rId, partname, sld_id_elem in targets:
        sld_id_lst.remove(sld_id_elem)
        try:
            pres_part.drop_rel(rId)
        except Exception:
            pass
        if partname is not None:
            try:
                package.drop_rel(partname)
            except Exception:
                pass


# --- Chart-data Excel companion (ThinkCell-compatible) ------------------

def collect_chart_data(elements, slide_label):
    """Pull out chart-data JSON blocks from chart elements, return list of
    (sheet_name, data_dict) for the xlsx writer. data_dict has 'categories',
    'series', optional 'notes'."""
    out = []
    chart_idx = 0
    for el in elements:
        if not el.get("isChart"):
            continue
        raw = el.get("chartData")
        if not raw:
            continue
        try:
            data = json.loads(raw)
        except Exception as e:
            print(f"WARN: slide {slide_label} chart {chart_idx}: invalid data-chart-data JSON ({e})", file=sys.stderr)
            chart_idx += 1
            continue
        title = el.get("chartTitle") or data.get("title")
        # Sheet name: short, file-safe
        chart_idx += 1
        suffix = f"_chart{chart_idx}" if chart_idx > 1 else ""
        # Excel sheet names are <=31 chars, no [ ] : * ? / \
        base = (title or f"slide{slide_label}").strip()
        base = re.sub(r"[\[\]:*?/\\]", "", base)[:25]
        sheet_name = f"{base}{suffix}"[:31] or f"slide{slide_label}{suffix}"
        out.append((sheet_name, data))
    return out


def write_chart_data_xlsx(all_chart_data, xlsx_path):
    """Write a ThinkCell-compatible xlsx.
    all_chart_data: list of (sheet_name, {"categories": [...], "series": [{"name":..., "values":[...]}], "notes": "..."})

    Layout per sheet (ThinkCell paste-friendly):
        A1: ""        B1: series1_name    C1: series2_name    ...
        A2: cat1      B2: val             C2: val             ...
        A3: cat2      B3: val             C3: val             ...
    Plus a final 'Notes' sheet aggregating any notes.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment
    except ImportError:
        print("WARN: openpyxl not installed — skipping chart-data xlsx (pip install openpyxl)", file=sys.stderr)
        return False
    if not all_chart_data:
        return False

    wb = Workbook()
    # Remove default sheet — we'll add our own
    default = wb.active
    wb.remove(default)

    notes_blocks = []
    used_names = set()
    for sheet_name, data in all_chart_data:
        # Disambiguate duplicates
        nm = sheet_name; n = 1
        while nm in used_names:
            n += 1
            nm = f"{sheet_name[:28]}_{n}"
        used_names.add(nm)
        ws = wb.create_sheet(title=nm)

        cats = data.get("categories") or []
        series = data.get("series") or []
        # Header row: A1 blank, then series names
        ws.cell(row=1, column=1, value="")
        for j, s in enumerate(series, start=2):
            c = ws.cell(row=1, column=j, value=s.get("name", f"Series {j-1}"))
            c.font = Font(bold=True)
        # Data rows
        for i, cat in enumerate(cats, start=2):
            ws.cell(row=i, column=1, value=cat).font = Font(bold=True)
            for j, s in enumerate(series, start=2):
                vals = s.get("values") or []
                v = vals[i-2] if i-2 < len(vals) else None
                ws.cell(row=i, column=j, value=v)
        # Auto-width columns (rough)
        for col_idx in range(1, len(series) + 2):
            max_len = 8
            for row_idx in range(1, len(cats) + 2):
                v = ws.cell(row=row_idx, column=col_idx).value
                if v is not None:
                    max_len = max(max_len, len(str(v)) + 2)
            from openpyxl.utils import get_column_letter
            ws.column_dimensions[get_column_letter(col_idx)].width = min(max_len, 30)

        if data.get("notes"):
            notes_blocks.append((nm, data["notes"]))

    if notes_blocks:
        ns = wb.create_sheet(title="Notes")
        ns.cell(row=1, column=1, value="Chart").font = Font(bold=True)
        ns.cell(row=1, column=2, value="Notes / methodology").font = Font(bold=True)
        for i, (nm, note) in enumerate(notes_blocks, start=2):
            ns.cell(row=i, column=1, value=nm)
            c = ns.cell(row=i, column=2, value=note)
            c.alignment = Alignment(wrap_text=True, vertical="top")
        ns.column_dimensions["A"].width = 24
        ns.column_dimensions["B"].width = 80

    xlsx_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(xlsx_path))
    return True


# --- XML emit helpers (theme-bound colors) -------------------------------

def _emit_solidFill(parent, color_spec, alpha=1.0, *, prepend=False):
    """Emit <a:solidFill><a:srgbClr|schemeClr val=.../></a:solidFill> into parent."""
    if color_spec is None:
        return None
    sf = etree.Element(qn("a:solidFill"))
    if color_spec["kind"] == "schemeClr":
        clr = etree.SubElement(sf, qn("a:schemeClr"))
    else:
        clr = etree.SubElement(sf, qn("a:srgbClr"))
    clr.set("val", color_spec["val"])
    if alpha < 0.99:
        a = etree.SubElement(clr, qn("a:alpha"))
        a.set("val", str(int(alpha * 100000)))
    if prepend:
        parent.insert(0, sf)
    else:
        parent.append(sf)
    return sf


def apply_shape_fill(shape, color_spec, alpha=1.0):
    spPr = shape.fill._xPr
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill", "a:blipFill"):
        for old in spPr.findall(qn(tag)):
            spPr.remove(old)
    _emit_solidFill(spPr, color_spec, alpha)


# --- Apply text + runs --------------------------------------------------

CSS_FONT_WEIGHT_BOLD = {"bold", "bolder", "600", "700", "800", "900"}


def apply_paragraph_runs(text_frame, el, color_resolver, theme_fonts=None):
    """Populate the text_frame from el's `runs` (preferred) or `directText`.
    Each run becomes its own <a:r> with its own font properties.
    Newlines split paragraphs.

    theme_fonts: optional {"major": "...", "minor": "..."} from
    extract_theme_fonts(). If provided, each run gets an explicit
    <a:latin typeface=".."/> stamped on it so PowerPoint uses the client's
    template font even when the HTML preview couldn't render it locally.
    Heading-ish runs (large or bold) use major; body uses minor.
    """
    runs = el.get("runs") or []
    text = el.get("directText", "") or ""
    # Pre-clear paragraphs
    p_el = text_frame.paragraphs[0]._p
    for r in p_el.findall(qn("a:r")):
        p_el.remove(r)
    # If runs list is empty but we have direct text, synthesize a single run
    if not runs and text.strip():
        runs = [{"text": text, "color": None, "fontWeight": None, "fontStyle": None}]

    # The element-level style is the default
    el_color = el.get("color")
    el_fw = (el.get("fontWeight") or "").strip().lower()
    el_fs_pt = parse_px_to_pt(el.get("fontSize"))
    el_fst = (el.get("fontStyle") or "").lower()
    el_align = (el.get("textAlign") or "left").lower()
    el_xform = (el.get("textTransform") or "").lower()
    el_ls = el.get("letterSpacing") or "normal"

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                 "right": PP_ALIGN.RIGHT, "start": PP_ALIGN.LEFT, "end": PP_ALIGN.RIGHT}
    text_frame.paragraphs[0].alignment = align_map.get(el_align, PP_ALIGN.LEFT)

    # Group runs by paragraph (split on \n in run text)
    paragraphs = [[]]
    for r in runs:
        rt = (r.get("text") or "").replace("\r", "")
        if "\n" in rt:
            parts = rt.split("\n")
            for i, part in enumerate(parts):
                if part:
                    paragraphs[-1].append({**r, "text": part})
                if i < len(parts) - 1:
                    paragraphs.append([])
        else:
            if rt:
                paragraphs[-1].append({**r, "text": rt})

    while len(paragraphs) > len(text_frame.paragraphs):
        text_frame.add_paragraph()

    for pi, runs_in_para in enumerate(paragraphs):
        para = text_frame.paragraphs[pi]
        para.alignment = align_map.get(el_align, PP_ALIGN.LEFT)
        # Clear runs of this paragraph
        para_el = para._p
        for old in para_el.findall(qn("a:r")):
            para_el.remove(old)
        for run_data in runs_in_para:
            run = para.add_run()
            run_text = run_data.get("text", "")
            if el_xform == "uppercase":
                run_text = run_text.upper()
            elif el_xform == "lowercase":
                run_text = run_text.lower()
            run.text = run_text

            # Run-level overrides fall back to element-level
            r_fs = parse_px_to_pt(run_data.get("fontSize")) or el_fs_pt
            r_fw = (run_data.get("fontWeight") or el_fw or "").strip().lower()
            r_fst = (run_data.get("fontStyle") or el_fst or "").lower()
            r_color = run_data.get("color") or el_color

            if r_fs:
                run.font.size = Pt(r_fs)
            if r_fw in CSS_FONT_WEIGHT_BOLD:
                run.font.bold = True
            if r_fst == "italic":
                run.font.italic = True

            color_hex, color_alpha = parse_css_color(r_color)
            if color_hex:
                cs = color_resolver(color_hex)
                rPr = run._r.get_or_add_rPr()
                for old in rPr.findall(qn("a:solidFill")):
                    rPr.remove(old)
                _emit_solidFill(rPr, cs, color_alpha, prepend=True)

            # letter-spacing in em
            if el_ls and el_ls != "normal" and el_fs_pt:
                m = re.match(r"^([\d.]+)(em|px)$", el_ls)
                if m:
                    val, unit = float(m.group(1)), m.group(2)
                    if unit == "em":
                        spc_val = int(val * el_fs_pt * 100)
                    else:
                        spc_val = int(val * 0.75 * 100)
                    if abs(spc_val) > 0:
                        rPr = run._r.get_or_add_rPr()
                        rPr.set("spc", str(spc_val))

            # Font family — explicitly stamp the template's typeface so the
            # PPTX renders correctly even when the HTML preview lacked the
            # font. Heading-ish runs (large or bold) use major; body uses minor.
            if theme_fonts:
                is_heading = (r_fs and r_fs >= 18) or (r_fw in CSS_FONT_WEIGHT_BOLD)
                typeface = theme_fonts.get("major" if is_heading else "minor")
                if typeface:
                    rPr = run._r.get_or_add_rPr()
                    # Remove any prior <a:latin> to avoid duplicates
                    for old in rPr.findall(qn("a:latin")):
                        rPr.remove(old)
                    latin = etree.SubElement(rPr, qn("a:latin"))
                    latin.set("typeface", typeface)


# --- Element placement on slide -----------------------------------------

INVARIANT_ROLES = {"source", "footnote", "page-number"}


def _has_text(el):
    return bool((el.get("directText") or "").strip())


def _has_visible_background(el):
    hex_, alpha = _effective_bg(el)
    return hex_ is not None


def _has_left_border(el):
    if (el.get("borderLeftWidth") or "0px") == "0px":
        return False
    if (el.get("borderLeftStyle") or "none") == "none":
        return False
    return True


def place_left_border(slide, el, color_resolver):
    """Render a border-left as a thin shape on the element's left edge."""
    width_str = el.get("borderLeftWidth", "0px")
    m = re.match(r"^([\d.]+)px$", width_str)
    if not m:
        return
    bw = float(m.group(1))
    color_hex, _ = parse_css_color(el.get("borderLeftColor"))
    if not color_hex:
        return
    x, y, w, h = el["x"], el["y"], el.get("w") or 0, el.get("h") or 0
    if h <= 0:
        return
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(x), _emu(y), _emu(bw), _emu(h))
    apply_shape_fill(bar, color_resolver(color_hex))
    bar.line.fill.background()


def place_element(slide, el, color_resolver, theme_fonts=None):
    """Translate a single element to a PPTX shape/textbox/picture."""
    if el.get("idx") == 0:
        return  # slide root — that's the canvas, don't draw it on top of the chart
    if el.get("isChart"):
        return  # handled by chart screenshot path
    if el.get("tableData"):
        return  # handled by place_table path
    if el.get("role") in INVARIANT_ROLES:
        return  # handled by invariant placement path
    if el.get("placeholder") and el.get("_placeholder_filled"):
        return  # successfully filled into a layout placeholder; don't double-render

    x = el.get("x", 0); y = el.get("y", 0)
    w = el.get("w") or 0; h = el.get("h") or 0
    if w <= 0 or h <= 0:
        return

    has_bg = _has_visible_background(el)
    has_text = _has_text(el)
    has_lb = _has_left_border(el)

    # Flex-centering correction: when this element is a child of a flex container
    # with justify-content centering, its bbox is content-tight. PowerPoint's
    # text metrics differ slightly from Chrome's, so a content-tight box can
    # cause unwanted wrapping. Expand the textbox to the parent's width and
    # rely on text-align:center to keep the visual centered.
    eff_x, eff_w = x, w
    eff_align_center = False
    # Respect CSS white-space:nowrap — always disable word wrap regardless of bbox heuristic.
    css_nowrap = (el.get("whitespace") or "").lower() in ("nowrap", "pre", "pre-line", "pre-wrap")
    is_single_line = css_nowrap
    if not has_bg and has_text:
        parent_disp = (el.get("parentDisplay") or "").lower()
        parent_just = (el.get("parentJustifyContent") or "").lower()
        if parent_disp in ("flex", "inline-flex") and parent_just in ("center", "space-around", "space-evenly"):
            if el.get("parentX") is not None and el.get("parentW"):
                eff_x = el["parentX"]
                eff_w = el["parentW"]
                eff_align_center = True
        else:
            # Horizontal slack: PowerPoint renders glyphs wider than Chrome,
            # so content-tight bboxes wrap unexpectedly. Scale with font size.
            # 10px font needs ~8px; 84px hero font needs much more.
            fs_pt = parse_px_to_pt(el.get("fontSize")) or 12
            slack = max(8, int(fs_pt * 0.6))
            eff_w = w + slack
        # Detect single-line text: bbox height is approximately one line.
        # In that case, disable word_wrap so PPT doesn't reflow the text
        # back onto a second line where it overlaps adjacent shapes.
        fs_pt2 = parse_px_to_pt(el.get("fontSize")) or 12
        # Roughly: if the bbox h is less than 1.5 * font-size-in-px, it's one line.
        fs_px = fs_pt2 / 0.75
        if h <= fs_px * 1.5:
            is_single_line = True

    if has_bg:
        bg_hex, bg_alpha = _effective_bg(el)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(x), _emu(y), _emu(w), _emu(h))
        apply_shape_fill(shape, color_resolver(bg_hex), bg_alpha)
        shape.line.fill.background()
        if has_text:
            tf = shape.text_frame
            tf.margin_left = tf.margin_right = Emu(45720)  # ~5px
            tf.margin_top = tf.margin_bottom = Emu(45720)
            tf.word_wrap = not css_nowrap  # honour white-space:nowrap for badge pills etc.
            # Vertical anchor: middle if flex with align-items:center
            if (el.get("parentAlignItems") or "").lower() == "center":
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                tf.vertical_anchor = MSO_ANCHOR.TOP
            apply_paragraph_runs(tf, el, color_resolver, theme_fonts)
    elif has_text:
        # Multi-line boxes get +8px width buffer: PowerPoint's font metrics are
        # slightly wider than Chrome's, causing trailing words or periods to wrap
        # to their own line and appear as stray bullets. Single-line boxes use
        # word_wrap=False instead so the buffer isn't needed there.
        box_w = eff_w if is_single_line else eff_w + 8
        box = slide.shapes.add_textbox(_emu(eff_x), _emu(y), _emu(box_w), _emu(h))
        tf = box.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = not is_single_line
        apply_paragraph_runs(tf, el, color_resolver, theme_fonts)
        if eff_align_center:
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER

    if has_lb:
        place_left_border(slide, el, color_resolver)


# --- Chart screenshotting -----------------------------------------------

def place_chart(slide, el, screenshot_fn, target_dir, slide_label=""):
    if not screenshot_fn:
        return
    # Include slide label in filename so multi-slide builds don't have
    # screenshot collisions when SVGs share idx across slides.
    suffix = f"_{slide_label}" if slide_label else ""
    out_png = target_dir / f"chart{suffix}_{el['idx']}.png"
    ok = screenshot_fn(el["idx"], str(out_png))
    if not ok or not out_png.exists():
        return
    x, y, w, h = el["x"], el["y"], el.get("w") or 0, el.get("h") or 0
    if w <= 0 or h <= 0:
        return
    slide.shapes.add_picture(str(out_png), _emu(x), _emu(y), _emu(w), _emu(h))


# --- Real PPTX table emission -------------------------------------------

def place_table(slide, el, color_resolver, theme_fonts=None):
    """Emit a real PPTX table from a captured <table> element.

    The table data was captured by the walker as a 2D structure of cells
    with per-cell text/styling. python-pptx's add_table() produces a
    real, editable PPTX table object."""
    td = el.get("tableData")
    if not td or not td.get("rows"):
        return
    rows = td["rows"]
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(r["cells"]) for r in rows)
    if n_cols == 0:
        return

    x = el["x"]; y = el["y"]
    w = el.get("w") or 0; h = el.get("h") or 0
    if w <= 0 or h <= 0:
        return

    table_shape = slide.shapes.add_table(n_rows, n_cols, _emu(x), _emu(y), _emu(w), _emu(h))
    table = table_shape.table

    # Set column widths from the first row's cell widths if available
    first_row_cells = rows[0]["cells"]
    if len(first_row_cells) == n_cols:
        for ci, cell_data in enumerate(first_row_cells):
            cw = cell_data.get("width") or (w / n_cols)
            try:
                table.columns[ci].width = _emu(cw)
            except Exception:
                pass

    for ri, row_data in enumerate(rows):
        cells_data = row_data.get("cells", [])
        for ci in range(n_cols):
            if ci >= len(cells_data):
                continue
            cell_data = cells_data[ci]
            cell = table.cell(ri, ci)
            # Background fill: prefer cell bg, fall back to row bg
            bg_color_str = cell_data.get("bg")
            if not is_visible_background(bg_color_str):
                bg_color_str = row_data.get("bg")
            if is_visible_background(bg_color_str):
                bg_hex, bg_alpha = parse_css_color(bg_color_str)
                if bg_hex:
                    cell.fill.solid()
                    # Use the color resolver so theme-bound hex still works
                    cs = color_resolver(bg_hex)
                    if cs and cs.get("kind") == "srgbClr":
                        from pptx.dml.color import RGBColor
                        cell.fill.fore_color.rgb = RGBColor.from_string(cs["val"])
            # Text content
            tf = cell.text_frame
            tf.margin_left = tf.margin_right = Emu(54864)  # ~6px
            tf.margin_top = tf.margin_bottom = Emu(36576)  # ~4px
            tf.word_wrap = True
            # Clear any default paragraph runs
            p_el = tf.paragraphs[0]._p
            for r in p_el.findall(qn("a:r")):
                p_el.remove(r)
            # Build a fake "element" to feed apply_paragraph_runs
            fake_el = {
                "directText": cell_data.get("text", ""),
                "runs": cell_data.get("runs", []),
                "color": cell_data.get("color"),
                "fontSize": cell_data.get("fontSize"),
                "fontWeight": cell_data.get("fontWeight"),
                "fontFamily": cell_data.get("fontFamily"),
                "textAlign": cell_data.get("textAlign", "left"),
                "textTransform": "",
                "letterSpacing": "normal",
            }
            apply_paragraph_runs(tf, fake_el, color_resolver, theme_fonts)


# --- Layout-aware placeholder filling -----------------------------------

# Common keyword aliases — what `data-placeholder="title"` etc. should match
# on a layout's placeholder name. Order matters: more specific aliases first.
PLACEHOLDER_ALIASES = {
    "title":       ["title"],
    "subtitle":    ["subtitle", "sub-title", "sub title"],
    "body":        ["body", "text placeholder", "content placeholder"],
    "intro":       ["intro"],
    "caption":     ["caption"],
    "date":        ["date"],
    "footer":      ["footer"],
    "slide-number": ["slide number", "page number"],
    # Numbered groups: stat-1, stat-2, stat-3 etc.
    # When the user writes data-placeholder="stat-N", we look for "Statistic Placeholder N"
    "stat":        ["statistic placeholder", "stat placeholder", "stat"],
    "sub":         ["subheading / body placeholder", "subheading", "sub-heading", "sub heading"],
    "callout":     ["call-out", "callout"],
    "header":      ["header"],
    "image":       ["picture placeholder", "image placeholder", "photo placeholder"],
    "icon":        ["icon placeholder"],
}


def _parse_placeholder_token(token):
    """Parse 'stat-2' -> ('stat', 2). Returns (base, number_or_None)."""
    m = re.match(r"^([a-zA-Z][a-zA-Z\-_]*?)[-_ ]?(\d+)?$", token.strip())
    if not m:
        return (token.strip().lower(), None)
    return (m.group(1).lower(), int(m.group(2)) if m.group(2) else None)


def _find_placeholder_by_name(slide, token):
    """Find a placeholder on the slide that matches `token` (e.g. 'title',
    'stat-2', 'caption', 'body').

    The matching is done against the LAYOUT's placeholder names (which retain
    semantic names like "Statistic Placeholder 2"), then the corresponding
    placeholder on the slide is resolved by placeholder index. Slides
    inherit position/font/style from layout placeholders but use generic
    names like "Text Placeholder 2" so matching the slide directly fails.

    Tries:
      1. Keyword + number combined (e.g. "Statistic Placeholder 1")
      2. Keyword match (no number requirement)
      3. Literal token base substring
    """
    base, num = _parse_placeholder_token(token)
    keywords = PLACEHOLDER_ALIASES.get(base, [base])

    layout = slide.slide_layout
    # Build map of {placeholder.idx: name} from the layout
    layout_phs = []
    for shape in layout.shapes:
        if not shape.is_placeholder:
            continue
        try:
            ph_idx = shape.placeholder_format.idx
        except Exception:
            continue
        layout_phs.append((ph_idx, (shape.name or "").lower()))

    # Determine which layout placeholder idx the token refers to
    matched_idx = None

    # Pass 1: keyword + number combined
    if num is not None:
        for ph_idx, sname in layout_phs:
            for kw in keywords:
                if kw and kw in sname and str(num) in sname:
                    matched_idx = ph_idx
                    break
            if matched_idx is not None:
                break
    # Pass 2: keyword match (skip if name's number contradicts)
    if matched_idx is None:
        for ph_idx, sname in layout_phs:
            matched = False
            for kw in keywords:
                if kw and kw in sname:
                    matched = True; break
            if not matched:
                continue
            if num is not None:
                other_nums = re.findall(r"\d+", sname)
                if other_nums and str(num) not in other_nums:
                    continue
            matched_idx = ph_idx
            break
    # Pass 3: literal token base substring against layout placeholder names
    if matched_idx is None:
        for ph_idx, sname in layout_phs:
            if base in sname:
                matched_idx = ph_idx
                break

    if matched_idx is None:
        return None

    # Now resolve to the slide-level placeholder with that idx
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        try:
            if shape.placeholder_format.idx == matched_idx:
                return shape
        except Exception:
            continue
    return None


def _find_placeholder_by_index(slide_or_layout, idx):
    """Find a placeholder by its placeholder index (for tags like
    data-placeholder="ph:5")."""
    for shape in slide_or_layout.shapes:
        if not shape.is_placeholder:
            continue
        try:
            if shape.placeholder_format.idx == idx:
                return shape
        except Exception:
            continue
    return None


def fill_layout_placeholders(slide, elements, color_resolver, theme_fonts):
    """For each element with data-placeholder, find a matching placeholder
    on the slide (which inherits from the layout) and pump the element's
    text/runs into it. Mark the element as _placeholder_filled so
    place_element skips it later.

    The layout's placeholder retains its position, font, and (typically)
    color from the master/layout — only the text content + inline run
    overrides (color/weight via <span>) come from the mockup.
    """
    for el in elements:
        ph_name = el.get("placeholder")
        if not ph_name:
            continue
        # Support "ph:5" syntax for explicit placeholder index
        target = None
        if ph_name.lower().startswith("ph:"):
            try:
                target = _find_placeholder_by_index(slide, int(ph_name[3:]))
            except ValueError:
                target = None
        else:
            target = _find_placeholder_by_name(slide, ph_name)
        if target is None or not target.has_text_frame:
            print(f"WARN: data-placeholder={ph_name!r} no matching placeholder on slide", file=sys.stderr)
            continue
        # Pump runs into the placeholder's text frame
        # We deliberately do NOT pass theme_fonts here — the layout's font
        # is already correct for the placeholder. We also skip color
        # unless the mockup explicitly set one via inline span.
        tf = target.text_frame
        # Clear existing text
        for p_el in list(tf.paragraphs[0]._p.getparent()):
            if p_el.tag.endswith("}p") and p_el is not tf.paragraphs[0]._p:
                p_el.getparent().remove(p_el)
        for r in tf.paragraphs[0]._p.findall(qn("a:r")):
            tf.paragraphs[0]._p.remove(r)
        apply_paragraph_runs(tf, el, color_resolver, theme_fonts=None)
        el["_placeholder_filled"] = True


# --- Invariant zone (source/footnote/page-number) -----------------------

def _find_layout_text_shape(layout, *keywords):
    for shape in layout.shapes:
        if not shape.has_text_frame:
            continue
        name = (shape.name or "").lower()
        if any(kw in name for kw in keywords):
            return shape
    return None


def _fill_at_layout_shape(slide, layout_shape, text):
    if not text:
        return
    try:
        x, y, w, h = (int(layout_shape.left), int(layout_shape.top),
                      int(layout_shape.width), int(layout_shape.height))
    except Exception:
        return
    fs = 9
    try:
        for para in layout_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    fs = run.font.size.pt; break
            if fs != 9:
                break
    except Exception:
        pass
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)


def _is_likely_invariant(el):
    """Heuristic for mockups that didn't tag invariants with data-role.
    Bottom-of-slide tiny-font text is treated as invariant."""
    fs_pt = parse_px_to_pt(el.get("fontSize")) or 12
    if fs_pt > 11:
        return None
    y = el.get("y", 0)
    if y < SLIDE_H_PX - 50:
        return None
    text = (el.get("directText") or "").strip().lower()
    if text.startswith("source"):
        return "source"
    if text.startswith("footnote"):
        return "footnote"
    # page number: a short token, often a single digit / Roman / "1 of 4"
    if len(text) <= 8:
        return "page-number"
    return None


def annotate_invariants(elements):
    """Mutate elements: assign role for any element heuristically detected
    as an invariant, so place_element will skip it."""
    for el in elements:
        if not el.get("role"):
            r = _is_likely_invariant(el)
            if r:
                el["role"] = r


def place_invariants(slide, elements):
    layout = slide.slide_layout
    for el in elements:
        role = el.get("role")
        if role not in INVARIANT_ROLES:
            continue
        text = (el.get("directText") or "").strip()
        if not text:
            continue
        if role == "source":
            ph = _find_layout_text_shape(layout, "source")
        elif role == "footnote":
            ph = _find_layout_text_shape(layout, "footnote", "disclaimer")
        elif role == "page-number":
            ph = _find_layout_text_shape(layout, "slide number", "page", "tracker")
        else:
            ph = None
        if ph is not None:
            _fill_at_layout_shape(slide, ph, text)
        else:
            # Fall back to placing as a textbox at element coords
            box = slide.shapes.add_textbox(_emu(el["x"]), _emu(el["y"]),
                                           _emu(el.get("w") or 200),
                                           _emu(el.get("h") or 20))
            tf = box.text_frame
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            fs_pt = parse_px_to_pt(el.get("fontSize"))
            if fs_pt:
                run.font.size = Pt(fs_pt)


# --- Draft badge ---------------------------------------------------------

def _add_draft_badge(slide, pres):
    """Stamp a grey DRAFT pill at the top-center of the slide.
    Placed last so it sits on top of all other shapes.
    Delete this shape manually when the deck is ready to share with a client."""
    sw = pres.slide_width   # EMU
    badge_w = Emu(1_371_600)   # ~108px at 96dpi (1.43 in)
    badge_h = Emu(342_900)     # ~27px (0.36 in)
    badge_x = (sw - badge_w) // 2
    badge_y = Emu(0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, badge_x, badge_y, badge_w, badge_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xE0, 0x00)  # yellow
    shape.line.fill.background()

    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Emu(91_440)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "DRAFT"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)  # red


# --- Picks parsing -------------------------------------------------------

def parse_picks(picks_str):
    out = []
    for tok in picks_str.split(","):
        tok = tok.strip()
        if not tok:
            continue
        m = re.match(r"^(\d+)\s*([A-Za-z]+)$", tok)
        if not m:
            print(f"WARN: can't parse pick '{tok}'", file=sys.stderr); continue
        out.append((m.group(1), m.group(2).upper()))
    return out


# --- Main ---------------------------------------------------------------

def build_one_slide(html_path, slide_index, option, pres, color_resolver, screenshot_dir, theme_fonts=None):
    elements, screenshot_fn, browser, pw, layout_attrs = render_and_walk(html_path, slide_index, option)
    try:
        if elements is None:
            print(f"WARN: slide {slide_index}{option} not found in mockup", file=sys.stderr)
            return False, []

        # Decide which layout to use: layout-aware (if mockup root has
        # data-layout-master + data-layout-index) or blank.
        layout_aware = False
        if layout_attrs and layout_attrs.get("master") and layout_attrs.get("index"):
            try:
                mi = int(layout_attrs["master"])
                li = int(layout_attrs["index"])
                slide = add_layout_slide(pres, mi, li)
                layout_aware = True
                lname = layout_attrs.get("name") or ""
                print(f"  slide {slide_index}{option}: using layout master {mi} #{li} {lname!r}")
            except (ValueError, TypeError):
                slide = add_blank_slide(pres)
                print(f"WARN: slide {slide_index}{option}: invalid data-layout attrs, using blank", file=sys.stderr)
        else:
            slide = add_blank_slide(pres)

        # Clear master-level FOOTER placeholders so they don't bleed through.
        # Templates that define a footer text placeholder at the master level
        # will push that text onto every slide unless it is explicitly cleared here.
        for _ph in slide.placeholders:
            if _ph.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                try:
                    _ph.text_frame.clear()
                except Exception:
                    pass

        # Tag invariants (source/footnote/page-number) before placement
        # so they aren't double-rendered.
        annotate_invariants(elements)

        # If layout-aware, fill placeholders FIRST so the placeholder-
        # tagged elements get marked _placeholder_filled and place_element
        # skips them.
        if layout_aware:
            fill_layout_placeholders(slide, elements, color_resolver, theme_fonts)

        # Place tables (real PPTX tables, not stacked textboxes)
        for el in elements:
            if el.get("tableData"):
                place_table(slide, el, color_resolver, theme_fonts)
        # Place shape/text elements (skips placeholders, charts, invariants)
        # IMPORTANT: this must happen BEFORE chart pictures, otherwise a wrapper
        # div with a background fill (e.g. #fafafa around an SVG) would paint
        # over the chart picture and hide it.
        for el in elements:
            place_element(slide, el, color_resolver, theme_fonts)
        # Place chart screenshots ON TOP of any background containers so the
        # chart is visible in the final render.
        slide_label = f"{slide_index}{option}"
        for el in elements:
            if el.get("isChart"):
                place_chart(slide, el, screenshot_fn, screenshot_dir, slide_label)
        # Place invariants last (master placeholder fills sit on top)
        place_invariants(slide, elements)
        # DRAFT badge — always on top, always last shape added
        _add_draft_badge(slide, pres)
        # Collect chart-data JSON for the xlsx companion
        chart_data = collect_chart_data(elements, f"{slide_index}{option}")
        return True, chart_data
    finally:
        if browser:
            try: browser.close()
            except Exception: pass
        if pw:
            try: pw.stop()
            except Exception: pass


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--mockup", required=True, help="Path to mockup HTML file")
    ap.add_argument("--picks", required=True,
                    help="Comma-separated picks like '1A,2B,3A,4B'")
    ap.add_argument("--target", required=True, help="Output PPTX path")
    ap.add_argument("--client-template", required=True, help="Client PPTX template")
    ap.add_argument("--slide", type=int, default=None,
                    help="Build only this 1-based slide index into a single-slide temp PPTX. "
                         "Used by parallel build workers — the pick for this slide must be "
                         "included in --picks (other picks are ignored).")
    args = ap.parse_args()

    html_path = Path(args.mockup).expanduser().resolve()
    target_path = Path(args.target).expanduser()
    template_path = Path(args.client_template).expanduser()

    if not html_path.exists():
        print(f"ERROR: mockup not found: {html_path}", file=sys.stderr); sys.exit(1)
    if not template_path.exists():
        print(f"ERROR: client template not found: {template_path}", file=sys.stderr); sys.exit(1)

    # Pre-build sanity checks on the mockup HTML
    _pre_build_checks(html_path)

    target_path.parent.mkdir(parents=True, exist_ok=True)
    # Always start from the client template — never open an existing target file.
    # Opening target_path would append slides onto a prior build instead of
    # starting fresh, producing a deck with duplicate or stale slides.
    pres = Presentation(str(template_path))
    purge_starter_slides(pres)
    purge_sections(pres)

    theme_hex = extract_theme_hex(str(template_path))
    color_resolver = make_color_resolver(theme_hex)
    theme_fonts = extract_theme_fonts(str(template_path))

    picks = parse_picks(args.picks)
    if not picks:
        print("ERROR: no valid picks parsed from --picks", file=sys.stderr); sys.exit(1)

    # --slide N: build only that one slide (parallel worker mode).
    # Filter picks to just the requested slide index.
    if args.slide is not None:
        picks = [(s, o) for s, o in picks if s == args.slide]
        if not picks:
            print(f"ERROR: --slide {args.slide} not found in --picks", file=sys.stderr)
            sys.exit(1)

    # Consolidate all chart screenshots into one shared _screenshots/ folder
    # inside the same directory as the mockup (typically _session/).
    # This prevents per-build _chart_screenshots_<name> folders from cluttering
    # the project directory.
    screenshot_dir = html_path.parent / "_screenshots"
    screenshot_dir.mkdir(parents=True, exist_ok=True)

    built = 0
    all_chart_data = []
    for s, o in picks:
        ok, chart_data = build_one_slide(html_path, s, o, pres, color_resolver, screenshot_dir, theme_fonts)
        if ok:
            built += 1
        all_chart_data.extend(chart_data)

    pres.save(str(target_path))
    print(f"OK: built {built} slide(s) from {html_path.name} into {target_path.name}")
    print(f"Theme colors detected: {theme_hex}")
    print(f"Theme fonts detected: major={theme_fonts.get('major')!r} minor={theme_fonts.get('minor')!r}")

    # --- Post-build geometry QA -------------------------------------------------
    # Run check_slide_geometry checks on every slide in the saved PPTX.
    # Import lazily so the rest of the build is unaffected if the module is absent.
    try:
        from check_slide_geometry import parse_slide as _parse_slide
        from check_slide_geometry import check_overlaps, check_boundaries, check_typography, check_empty_bottom
        import io

        # Standard widescreen EMU dimensions (13.33" x 7.5" at 914400 EMU/inch)
        SLIDE_W_EMU = 12_192_000
        SLIDE_H_EMU = 6_858_000

        qa_pass = True
        print("\n--- Geometry QA ---")
        with zipfile.ZipFile(str(target_path), "r") as zf:
            slide_names = sorted(
                n for n in zf.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                and "/_rels/" not in n
            )
            for sname in slide_names:
                label = sname.split("/")[-1]
                xml_bytes = zf.read(sname)
                import tempfile, os
                with tempfile.NamedTemporaryFile(suffix=".xml", delete=False) as tf:
                    tf.write(xml_bytes)
                    tf_path = Path(tf.name)
                try:
                    shapes, sizes = _parse_slide(tf_path)
                finally:
                    os.unlink(tf_path)

                overlap_errs  = check_overlaps(shapes)
                boundary_errs = check_boundaries(shapes, SLIDE_W_EMU, SLIDE_H_EMU)
                typo_errs     = check_typography(sizes)
                empty_warns   = check_empty_bottom(shapes, SLIDE_H_EMU)

                all_errs = overlap_errs + boundary_errs + typo_errs
                if all_errs:
                    qa_pass = False
                    print(f"  FAIL {label}:")
                    for e in all_errs:
                        print(f"    - {e}")
                elif empty_warns:
                    print(f"  WARN {label}: {empty_warns[0]}")
                else:
                    print(f"  PASS {label}")

        if qa_pass:
            print("Geometry QA: all slides clean.")
        else:
            print("\nGeometry QA: FAILURES detected above — fix the mockup and rebuild before delivering.")
    except ImportError:
        print("(check_slide_geometry not found — skipping geometry QA)")
    except Exception as _qa_err:
        print(f"(Geometry QA skipped — unexpected error: {_qa_err})")
    # ---------------------------------------------------------------------------

    # Write the ThinkCell-compatible chart-data xlsx companion
    if all_chart_data:
        xlsx_path = target_path.with_name(target_path.stem + "-chart-data.xlsx")
        if write_chart_data_xlsx(all_chart_data, xlsx_path):
            print(f"OK: wrote chart-data companion -> {xlsx_path.name}")
            print(f"     ({len(all_chart_data)} chart(s); paste any range into ThinkCell)")


def print_theme_main():
    """Helper invocation: `python build_slide.py --print-theme template.pptx`
    Prints the template's color and font scheme as JSON, for slide-helper
    to inject into the mockup CSS at the start of a session."""
    ap = argparse.ArgumentParser()
    ap.add_argument("--print-theme", required=True, help="Template PPTX")
    args = ap.parse_args()
    template_path = Path(args.print_theme).expanduser()
    if not template_path.exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr); sys.exit(1)
    theme_hex = extract_theme_hex(str(template_path))
    theme_fonts = extract_theme_fonts(str(template_path))
    print(json.dumps({"colors": theme_hex, "fonts": theme_fonts}, indent=2))


def catalog_layouts_main():
    """Helper invocation: `python build_slide.py --catalog-layouts template.pptx`
    Prints a JSON manifest of every master/layout in the template, with
    placeholder positions. Used by slide-helper to present corporate-approved
    layouts to the user during Phase A."""
    ap = argparse.ArgumentParser()
    ap.add_argument("--catalog-layouts", required=True, help="Template PPTX")
    ap.add_argument("--filter", default=None,
                    help="Optional case-insensitive substring filter on layout name")
    args = ap.parse_args()
    template_path = Path(args.catalog_layouts).expanduser()
    if not template_path.exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr); sys.exit(1)
    cat = catalog_layouts(str(template_path))
    if args.filter:
        f = args.filter.lower()
        cat = [L for L in cat if f in (L["name"] or "").lower()]
    print(json.dumps(cat, indent=2))


if __name__ == "__main__":
    if "--print-theme" in sys.argv:
        print_theme_main()
    elif "--catalog-layouts" in sys.argv:
        catalog_layouts_main()
    else:
        main()
