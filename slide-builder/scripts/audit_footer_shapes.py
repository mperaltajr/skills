"""
audit_footer_shapes.py — inspect footer shapes in all skeleton PPTXs.
Prints shape names, positions (inches), and text for any shape in the
bottom 1" of the slide, plus all shapes whose text contains {{FOOTNOTE}},
{{SOURCE}}, or {{PAGE_NUMBER}}.
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Emu

SKEL_DIR = pathlib.Path(__file__).resolve().parent.parent / "skeletons"
FOOTER_TOKENS = {"FOOTNOTE", "SOURCE", "PAGE_NUMBER"}
SLIDE_H = 6858000  # EMU: 7.5"


def emu_to_in(emu):
    return round(emu / 914400, 3)


for skel_dir in sorted(SKEL_DIR.iterdir()):
    pptx = skel_dir / f"{skel_dir.name}.pptx"
    if not pptx.exists():
        continue
    prs = Presentation(str(pptx))
    slide = prs.slides[0]
    print(f"\n{'='*60}")
    print(f"Skeleton: {skel_dir.name}")
    print(f"{'='*60}")
    found_tokens = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        # Check for footer tokens
        for tok in FOOTER_TOKENS:
            if f"{{{{{tok}}}}}" in text:
                found_tokens.add(tok)
        # Print shapes in footer zone (bottom 1")
        try:
            y = shape.top
            h = shape.height
            if y + h > SLIDE_H - Inches(1):
                x_in = emu_to_in(shape.left)
                y_in = emu_to_in(shape.top)
                w_in = emu_to_in(shape.width)
                h_in = emu_to_in(shape.height)
                print(f"  [{shape.name}]")
                print(f"    pos: ({x_in}\", {y_in}\")  size: {w_in}\" x {h_in}\"")
                print(f"    text: {text[:80]!r}")
                # Print font size if available
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                print(f"    font: {run.font.size.pt}pt")
                                break
                except Exception:
                    pass
        except Exception:
            pass
    missing = FOOTER_TOKENS - found_tokens
    print(f"  Tokens present: {found_tokens or 'none'}")
    print(f"  Tokens MISSING: {missing or 'none'}")
