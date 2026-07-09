"""Slide Lab deck compiler — combine user-picked themed slides into one final deck.

The sys.path setup points at slide-builder/'s shared infrastructure explicitly
since this script does not carry its own twins/ module.

Inputs:
  --out PATH    Orchestrator output dir (the one that has _meta.json, slide_NN/ dirs, etc.)
  --picks       Either a JSON file path OR a JSON string mapping slide_NN -> letter.
                Example string: {"slide_01":"A","slide_02":"C",...}
                If omitted, reads <out>/picks.json if present.
  --final PATH  Final deck path (default: <out>/final_deck.pptx)

What it does:
  1. Open the client template (path from <out>/_meta.json).
  2. _clear_existing_slides() — strip template stock slides + named sections.
  3. For each slide in numeric order, open <out>/slide_NN/option_<X>.pptx
     and copy its single slide's shapes into a new blank slide in the final deck
     (deepcopy(shape.element) + append to _spTree).
  4. Save to --final.
  5. Render every slide of the final deck to PNG via render_libre → output to
     <out>/final_pngs/.
  6. Write <out>/COMPILED.md summarizing picks, output path, slide count,
     render success/fail, opens-cleanly status.

The themed PPTX is already client-branded by finalize_deck.py — we do
NOT re-graft or re-theme here. Just combine. Fallback options
were already assembled into themed PPTX at finalize time and look identical
to native options from this script's perspective.
"""
from __future__ import annotations

import argparse
import json
import re
import sys
import traceback
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Optional

# Path setup: twins/ is local to this skill; render_slides lives in the
# sibling slide-qc skill.
SKILL_ROOT = Path(__file__).resolve().parents[1]              # slide-builder/
QC_SCRIPTS = SKILL_ROOT.parent / "slide-qc" / "scripts"       # render_slides
sys.path.insert(0, str(SKILL_ROOT))
sys.path.insert(0, str(QC_SCRIPTS))

import _paths as _p  # noqa: E402

from pptx import Presentation  # noqa: E402
from twins.composer import (  # noqa: E402
    _clear_existing_slides,
    _find_blank_layout,
    _find_named_layout,
    _strip_layout_placeholders,
    remove_empty_placeholders,
)


# ---------------------------------------------------------------------------
# Pick parsing
# ---------------------------------------------------------------------------
def parse_picks(arg: Optional[str], out_dir: Path) -> dict[str, str]:
    """Resolve --picks. Accepts:
      - None             -> read <out>/picks.json
      - existing file    -> json.load
      - JSON string      -> json.loads
    Normalizes keys to 'slide_NN' (zero-padded) and uppercases letters.
    """
    if arg is None:
        candidate = _p.picks_json(out_dir)
        if not candidate.exists():
            raise SystemExit(f"--picks not given and {candidate} does not exist")
        raw = candidate.read_text(encoding="utf-8")
    else:
        as_path = Path(arg)
        if as_path.exists():
            raw = as_path.read_text(encoding="utf-8")
        else:
            raw = arg

    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        raise SystemExit(f"--picks: invalid JSON ({e})")

    if not isinstance(data, dict):
        raise SystemExit(f"--picks: expected object, got {type(data).__name__}")

    normalized: dict[str, str] = {}
    for k, v in data.items():
        m = re.match(r"slide[_\-]?(\d+)", str(k), re.IGNORECASE)
        if not m:
            raise SystemExit(f"--picks: bad key {k!r}, expected like 'slide_01'")
        key = _p.slide_key(int(m.group(1)))
        letter = str(v).strip().upper()
        if not letter:
            raise SystemExit(f"--picks: empty letter for {key}")
        normalized[key] = letter
    return normalized


# ---------------------------------------------------------------------------
# Slide copying
# ---------------------------------------------------------------------------
def copy_picked_slide_into(dst_prs, src_pptx: Path,
                           layout_name: str = "",
                           layout_chrome=None,
                           keep_master_shapes: bool = False,
                           page_position: Optional[int] = None) -> int:
    """Open `src_pptx`, append its first slide's shapes onto a new slide in
    `dst_prs`.

    Body-canonical inheritance: when layout_chrome is a body-canonical
    chrome (with title_placeholder_idx set), the target layout's inherited
    placeholders + decorative shapes are KEPT (not stripped); the layout name
    drives the layout pick. Otherwise this falls back to the blank-layout
    behavior (blank layout, strip placeholders).

    `page_position` (1-based) re-stamps the slide-number placeholder to the
    slide's actual position in the compiled deck. finalize bakes the page number
    at build time from the slide's original number, so after a rebuild/insert
    renumber the baked value can be stale; compile order is the source of truth.
    """
    src_prs = Presentation(str(src_pptx))
    src_slide = src_prs.slides[0]

    _is_body_canonical = (
        layout_chrome is not None
        and getattr(layout_chrome, "layout_class", None) == "body-canonical"
        and getattr(layout_chrome, "title_placeholder_idx", None) is not None
    )

    if _is_body_canonical and layout_name:
        target = _find_named_layout(dst_prs, layout_name) or _find_blank_layout(dst_prs)
    else:
        target = _find_blank_layout(dst_prs)

    new_slide = dst_prs.slides.add_slide(target)
    if not _is_body_canonical:
        _strip_layout_placeholders(new_slide, keep_master_shapes=keep_master_shapes)

    sp_tree = new_slide.shapes._spTree
    count = 0
    for shape in src_slide.shapes:
        sp_tree.append(deepcopy(shape.element))
        count += 1

    # For body-canonical
    # destinations, the new_slide already has layout-inherited placeholders
    # (empty Title 1, Text Placeholder 2, etc.). The source slide carries
    # its OWN populated copies of the same placeholders (finalize_deck wrote
    # text into them on the source's instance of the layout). After deepcopy,
    # the destination ends up with TWO placeholders per (type, idx) key:
    # one empty inherited, one populated from source. LibreOffice suppresses
    # the empty one; PowerPoint renders BOTH — the empty one shows the
    # "Click to add title" prompt overlaying the real title.
    #
    # Dedupe: walk all placeholder shapes on the new slide, group by
    # (type, idx), and when two exist for the same key, drop the empty one.
    if _is_body_canonical:
        _dedupe_placeholder_duplicates(new_slide)
        # Drop any inherited placeholder left empty (e.g. the layout's content
        # placeholder — the body renders as free-floating shapes) so PowerPoint
        # doesn't show its 'Click to add text' prompt in edit mode. Runs AFTER
        # dedupe so the populated title/slide-number survive.
        remove_empty_placeholders(new_slide)
    if page_position is not None:
        _restamp_page_number(new_slide, page_position)
    return count


def _restamp_page_number(slide, position: int) -> None:
    """Set the slide-number placeholder's text to `position` (its 1-based place
    in the compiled deck). Only re-stamps a placeholder that currently holds a
    bare integer — the literal finalize bakes in — so auto-number fields and
    cover slides (no page-number placeholder) are left alone. Preserves the
    existing run's formatting by editing the run text in place."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    for shape in slide.shapes:
        try:
            if not shape.is_placeholder:
                continue
            if shape.placeholder_format.type != PP_PLACEHOLDER.SLIDE_NUMBER:
                continue
        except Exception:
            continue
        try:
            tf = shape.text_frame
            if (tf.text or "").strip().isdigit():
                paras = tf.paragraphs
                if paras and paras[0].runs:
                    paras[0].runs[0].text = str(position)
                    for extra in paras[0].runs[1:]:
                        extra.text = ""
                else:
                    tf.text = str(position)
        except Exception:
            pass
        return


def _dedupe_zip_entries(pptx_path: Path) -> int:
    """Rewrite a PPTX zip in place, keeping only the LAST entry for each
    name. Returns count of duplicates removed.

    LibreOffice rejects zips with duplicate entry names as corrupt; python-
    pptx occasionally produces them on multi-slide saves where shape
    relationships overlap. This dedupe pass is folded into compile_picks to
    keep the saved deck loadable.
    """
    import zipfile
    import io
    if not pptx_path.exists():
        return 0
    try:
        raw = pptx_path.read_bytes()
        with zipfile.ZipFile(io.BytesIO(raw), "r") as zin:
            names = zin.namelist()
            if len(names) == len(set(names)):
                return 0  # no duplicates; cheap exit
            # Walk infos so we keep LAST occurrence per name (overrides earlier)
            kept: dict[str, tuple] = {}
            for info in zin.infolist():
                kept[info.filename] = (info, zin.read(info.filename))
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for info, data in kept.values():
                zout.writestr(info, data)
        pptx_path.write_bytes(buf.getvalue())
        n_removed = len(names) - len(set(names))
        print(f"  dedupe: removed {n_removed} duplicate zip entries")
        return n_removed
    except Exception as exc:
        sys.stderr.write(
            f"  WARN: zip dedupe failed ({type(exc).__name__}: {exc}). "
            f"Deck may not open in LibreOffice if it had duplicate entries.\n"
        )
        return 0


def _dedupe_placeholder_duplicates(slide) -> int:
    """Remove duplicate placeholder shapes that share (type, idx). When two
    placeholders have the same key, keep the one with text content and drop
    the other. Returns number of shapes removed.

    Used after compile_picks's deepcopy path to clean up the layout-inherited
    + source-copied duplicate pair that otherwise renders as visible empty
    placeholder prompts in PowerPoint.
    """
    from collections import defaultdict
    groups: dict[tuple, list] = defaultdict(list)
    for shp in list(slide.shapes):
        try:
            pf = shp.placeholder_format
        except Exception:
            continue
        if pf is None:
            continue
        try:
            key = (int(pf.type), int(pf.idx))
        except Exception:
            continue
        groups[key].append(shp)

    sp_tree = slide.shapes._spTree
    removed = 0
    for key, shapes in groups.items():
        if len(shapes) <= 1:
            continue
        # Prefer the shape that has actual text content; drop empties.
        def _has_text(shp) -> bool:
            try:
                return bool((shp.text_frame.text or "").strip())
            except Exception:
                return False
        with_text = [s for s in shapes if _has_text(s)]
        empties  = [s for s in shapes if not _has_text(s)]
        # If at least one has text, drop ALL the empties.
        # If none have text, keep the first (arbitrary; both empty anyway).
        if with_text:
            to_drop = empties
        else:
            to_drop = shapes[1:]
        for shp in to_drop:
            try:
                sp_tree.remove(shp.element)
                removed += 1
            except (ValueError, AttributeError):
                pass
    return removed


# ---------------------------------------------------------------------------
# COMPILED.md
# ---------------------------------------------------------------------------
COMPILED_TEMPLATE = """# Slide Lab compiled deck

Generated: {ts}

Out dir   : `{out}`
Template  : `{template}`
Final deck: `{final}`

## Picks

| Slide | Pick | Source PPTX | Shapes copied | Status |
|-------|------|-------------|---------------|--------|
{rows}

## Result

- Final slide count: **{slide_count}**
- Opens cleanly (python-pptx reload): **{opens}**
- Renders attempted: **{render_total}**
- Renders succeeded: **{render_ok}**
- Renders failed   : **{render_fail}**

## Failures

{failures}
"""


# ---------------------------------------------------------------------------
# Splice-back (option 6b) — replace slide(s) in an EXTERNAL deck in place
# ---------------------------------------------------------------------------
def run_splice(out_dir: Path, meta: dict, picks: dict, template_path: Path,
               original: Path, final: Optional[Path]) -> int:
    """Splice picked, rebuilt slides back into `original` at their positions,
    preserving every other slide. Writes a NEW file; never touches the original.

    Each picks key `slide_NN` is the 1-based position IN THE ORIGINAL DECK
    (adopt_deck.py records it that way). For each, we append the rebuilt themed
    slide, move it to position NN in `<p:sldIdLst>`, and drop the old slide
    there. Slide count is unchanged.
    """
    from pptx.oxml.ns import qn  # noqa: E402
    if not original.exists():
        print(f"ERROR: --splice-into deck not found: {original}")
        return 2

    # Per-slide layout + chrome + keep-master, same resolution as a compile.
    from twins.client_theme import load_client_theme  # noqa: E402
    from _chrome_schema import load_chrome_yml  # noqa: E402
    try:
        _theme = load_client_theme(str(template_path))
        keep_master = not bool(getattr(_theme, "strip_master_backgrounds", True))
    except Exception:
        keep_master = True  # external deck may not be registered; keep master chrome
    chrome_spec = None
    try:
        chrome_spec = load_chrome_yml(_p.chrome_yml(template_path))
    except Exception:
        chrome_spec = None
    slide_layouts: dict[str, str] = {}
    for s in meta.get("slides", []):
        n = s.get("n")
        if isinstance(n, int):
            slide_layouts[_p.slide_key(n)] = (s.get("layout") or "").strip()

    out = final or original.with_name(f"{original.stem}_slidelab{original.suffix}")
    prs = Presentation(str(original))
    sldIdLst = prs.slides._sldIdLst
    n_slides = len(list(sldIdLst))

    print("=" * 72)
    print("Slide Lab splice-back (option 6b)")
    print(f"  original : {original}  ({n_slides} slides)")
    print(f"  out      : {out}")
    print("=" * 72)

    spliced, failures = [], []
    for key in sorted(picks.keys(), key=lambda k: int(k.split("_")[1])):
        raw = picks[key]
        letter = raw[0] if isinstance(raw, list) else raw
        pos = int(key.split("_")[1])  # 1-based position in the ORIGINAL deck
        src = out_dir / key / _p.option_pptx_name(letter)
        if not (1 <= pos <= n_slides):
            failures.append(f"{key}: position {pos} out of range (deck has {n_slides})")
            continue
        if not src.exists():
            failures.append(f"{key}: missing rebuilt slide {src}")
            continue
        layout_name = slide_layouts.get(key, "")
        layout_chrome = chrome_spec.layouts.get(layout_name) if (chrome_spec and layout_name) else None
        ids_before = list(sldIdLst)          # snapshot BEFORE the append
        copy_picked_slide_into(prs, src, layout_name=layout_name,
                               layout_chrome=layout_chrome,
                               keep_master_shapes=keep_master, page_position=pos)
        new_sldId = list(sldIdLst)[-1]        # the just-appended slide
        old_sldId = ids_before[pos - 1]       # the slide currently at position pos
        old_sldId.addprevious(new_sldId)      # move rebuilt slide into position pos
        rId = old_sldId.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(old_sldId)            # drop the old slide (count unchanged)
        spliced.append(pos)
        print(f"  spliced {key} (pick {letter}) into position {pos}  ok")

    if not spliced:
        print("ERROR: nothing spliced.")
        for f in failures:
            print(f"  - {f}")
        return 2
    try:
        prs.save(str(out))
    except (PermissionError, OSError) as exc:
        print(f"ERROR: could not write {out.name}: {type(exc).__name__}: {exc} "
              f"(is it open in PowerPoint?)")
        return 3
    _dedupe_zip_entries(out)

    # Verify the count is preserved (the whole point of a splice vs a compile).
    verify = Presentation(str(out))
    final_count = len(verify.slides)
    ok = final_count == n_slides
    print(f"\n[verify] slide count {final_count} (was {n_slides}) — "
          f"{'ok' if ok else 'MISMATCH'}")
    print(f"  spliced positions: {spliced}")
    if failures:
        print("  skipped:")
        for f in failures:
            print(f"    - {f}")
    print(f"\nSpliced deck: {out}")
    print("Run slide-qc on it before sending (a deck isn't done until QC has run).")
    return 0 if ok else 1


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main() -> int:
    ap = argparse.ArgumentParser(description="Compile picked themed slides into a final deck.")
    ap.add_argument("--out", required=True, type=Path, help="Orchestrator output dir")
    ap.add_argument("--picks", default=None, help="picks.json path OR JSON string")
    ap.add_argument("--final", default=None, type=Path,
                    help="Final deck path (default: <out>/final_deck.pptx)")
    ap.add_argument("--all-variations", action="store_true",
                    help="Instead of "
                         "selecting one pick per slide, iterate ALL options "
                         "present per slide and stack them into the final "
                         "deck. Useful for shipping a stakeholder-review "
                         "artifact where the audience picks among variants. "
                         "Mutually exclusive with --picks.")
    ap.add_argument("--splice-into", default=None, type=Path,
                    help="Option 6b (external-deck redesign): splice the picked, "
                         "rebuilt slide(s) back into THIS original .pptx at their "
                         "positions, keeping every other slide untouched. Writes a "
                         "new file (never overwrites the original). Use this for a "
                         "deck adopted via adopt_deck.py — a plain compile would "
                         "drop the un-rebuilt slides.")
    args = ap.parse_args()
    if args.all_variations and args.picks:
        print("ERROR: --all-variations is mutually exclusive with --picks. "
              "Choose one: pick-driven compile (--picks) OR all-variations "
              "stack (--all-variations).")
        return 2

    from _log import attach as _log_attach
    _log_attach(args.out, "compile_picks.py")

    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    out_dir: Path = args.out
    if not out_dir.exists():
        print(f"ERROR: out dir not found: {out_dir}")
        return 2

    meta_path = _p.meta_json(out_dir)
    if not meta_path.exists():
        print(f"ERROR: _meta.json not found at {meta_path}")
        return 2
    meta = json.loads(meta_path.read_text(encoding="utf-8"))
    from _meta_schema import validate_warn
    validate_warn(meta, source="compile_picks")
    template_path = Path(meta["template"])
    if not template_path.exists():
        print(f"ERROR: template (from _meta.json) not found: {template_path}")
        return 2

    # All-variations mode synthesizes a
    # picks-like structure that carries ALL options per slide instead of
    # one. Downstream rendering iterates the same loop; the only difference
    # is the picks dict value is a LIST of letters instead of a single str.
    if args.all_variations:
        picks: dict = {}
        for s in meta.get("slides", []):
            n = s.get("n")
            if isinstance(n, int):
                sdir = out_dir / _p.slide_key(n)
                # Iterate the options actually built on disk — the count is
                # now settings-driven (default 1) and a reviewer may have
                # requested extra options on specific slides, so a hardcoded
                # A/B/C would emit FAIL rows for options that never existed.
                letters = sorted(
                    p.stem.split("_")[1] for p in sdir.glob("option_?.pptx")
                )
                picks[_p.slide_key(n)] = letters or list(_p.option_letters())
        total = sum(len(v) for v in picks.values())
        print(f"  all-variations mode: {len(picks)} slides, "
              f"{total} options present = {total} planned outputs")
        final_path: Path = args.final or (out_dir / "final_deck_all_variations.pptx")
    else:
        picks = parse_picks(args.picks, out_dir)
        final_path = args.final or (out_dir / "final_deck.pptx")

    # Option 6b: splice rebuilt slides back into the external original in place.
    if args.splice_into:
        return run_splice(out_dir, meta, picks, template_path,
                          Path(args.splice_into), args.final)

    # Guard: an adopted external deck must NOT be compiled from scratch — the
    # clear-and-rebuild path would drop every slide the user didn't rebuild.
    # adopt_deck.py stamps `adopted_source`; force the splice path instead.
    if meta.get("adopted_source") and not args.all_variations:
        print("ERROR: this build was adopted from an external deck "
              f"({meta.get('adopted_source')}).")
        print("       A plain compile would keep ONLY the rebuilt slides and drop the rest.")
        print("       Splice the rebuilt slide(s) back into the original instead:")
        print(f"         py -3 compile_picks.py --out {out_dir} "
              f"--splice-into \"{meta.get('adopted_source')}\"")
        return 2

    final_path.parent.mkdir(parents=True, exist_ok=True)

    print("=" * 72)
    print("Slide Lab deck compiler")
    print(f"  out      : {out_dir}")
    print(f"  template : {template_path}")
    print(f"  picks    : {len(picks)} entries")
    print(f"  final    : {final_path}")
    print("=" * 72)

    print("\n[1] Open template + clear existing slides + sections")
    # Open the normalized build copy when registration made one; sidecar lookups
    # below stay keyed off the ORIGINAL template_path.
    from twins.client_theme import resolve_build_template  # noqa: E402
    build_template_path = resolve_build_template(template_path)
    if build_template_path != template_path:
        print(f"  building on normalized copy: {build_template_path}")
    dst_prs = Presentation(str(build_template_path))
    _clear_existing_slides(dst_prs)

    # Load brand.yml to honor strip_master_backgrounds. When false (the
    # common client default), the master IS the brand chrome — top/bottom
    # bands must survive into the compiled deck.
    # See feedback_sidecar_fallback_must_be_loud.
    from twins.client_theme import load_client_theme  # noqa: E402
    _theme = load_client_theme(str(template_path))
    _keep_master = not bool(getattr(_theme, "strip_master_backgrounds", True))

    # Load chrome.yml so we can thread the per-slide
    # layout_chrome into copy_picked_slide_into. Without this, body-canonical
    # slides graft onto the blank layout and lose their decorative chrome
    # (top chevron / footer bar inherited from the named layout). See
    # feedback_sidecar_fallback_must_be_loud.
    from _chrome_schema import load_chrome_yml  # noqa: E402
    _chrome_spec = None
    try:
        _chrome_spec = load_chrome_yml(_p.chrome_yml(template_path))
    except Exception as _exc:
        print(f"  WARN: chrome.yml unavailable; compile will fall back to blank "
              f"layout for every slide (loses body-canonical chrome): "
              f"{type(_exc).__name__}: {_exc}")

    # Per-slide layout lookup from _meta.json.
    _slide_layouts: dict[str, str] = {}
    try:
        for s in meta.get("slides", []):
            n = s.get("n")
            if isinstance(n, int):
                _slide_layouts[_p.slide_key(n)] = (s.get("layout") or "").strip()
    except Exception:
        _slide_layouts = {}

    print("\n[2] Copy picked themed slides")
    rows: list[str] = []
    failures: list[str] = []
    ordered_keys = sorted(picks.keys(), key=lambda k: int(k.split("_")[1]))
    copied_count = 0
    for key in ordered_keys:
        # Normalize: pick-mode value is a str letter; all-variations is a list.
        raw = picks[key]
        letters = raw if isinstance(raw, list) else [raw]
        slide_layout_name = _slide_layouts.get(key, "")
        slide_layout_chrome = None
        if _chrome_spec is not None and slide_layout_name:
            slide_layout_chrome = _chrome_spec.layouts.get(slide_layout_name)
        for letter in letters:
            src = out_dir / key / _p.option_pptx_name(letter)
            if not src.exists():
                msg = f"missing source: {src}"
                failures.append(f"- **{key} pick {letter}**: {msg}")
                rows.append(f"| {key} | {letter} | `{src.name}` | - | FAIL ({msg}) |")
                print(f"  {key} pick {letter}  FAIL ({msg})")
                continue
            try:
                n_shapes = copy_picked_slide_into(
                    dst_prs, src,
                    layout_name=slide_layout_name,
                    layout_chrome=slide_layout_chrome,
                    keep_master_shapes=_keep_master,
                    page_position=copied_count + 1,
                )
                copied_count += 1
                rows.append(f"| {key} | {letter} | `{src.name}` | {n_shapes} | ok |")
                print(f"  {key} pick {letter}  ok (shapes={n_shapes})")
            except Exception as e:
                tb = traceback.format_exc().strip().splitlines()[-1]
                msg = f"{type(e).__name__}: {e} | {tb}"
                failures.append(f"- **{key} pick {letter}**: {msg}")
                rows.append(f"| {key} | {letter} | `{src.name}` | - | FAIL ({msg[:60]}...) |")
                print(f"  {key} pick {letter}  FAIL ({msg[:80]})")

    print(f"\n[3] Save final deck -> {final_path}")
    # Backup existing final deck before overwrite so re-runs don't silently
    # destroy hand-edits the user may have made between compiles.
    if final_path.exists():
        try:
            from datetime import datetime as _dt
            ts = _dt.now().strftime("%Y%m%dT%H%M%S")
            backup = final_path.with_name(f"{final_path.stem}.{ts}{final_path.suffix}")
            final_path.replace(backup)
            print(f"  backed up prior deck -> {backup.name}")
        except OSError as exc:
            sys.stderr.write(
                f"  WARN: could not back up prior {final_path.name}: {exc}\n"
                f"  Proceeding will overwrite. If you have unsaved edits in the prior deck, abort now (Ctrl+C).\n"
            )
    # The save itself can fail with
    # PermissionError if the destination final_deck.pptx is open in
    # PowerPoint, or with OSError under antivirus/file-lock conditions.
    # The backup above protects the prior copy; this guard protects
    # the user from a raw traceback on save.
    try:
        dst_prs.save(str(final_path))
    except (PermissionError, OSError) as exc:
        sys.stderr.write(
            f"\nERROR: could not write {final_path.name}: {type(exc).__name__}: {exc}\n"
            f"       The deck is locked - most commonly because the prior "
            f"copy is open in PowerPoint.\n"
            f"       Close PowerPoint (or pause antivirus on the build dir) "
            f"and re-run compile_picks.py.\n"
        )
        sys.exit(3)

    # python-pptx can write
    # zips with duplicate entry names when many slides graft from
    # similarly-structured source decks. LibreOffice rejects such PPTX files
    # as corrupt. Run an unconditional zip rewrite that keeps the LAST
    # occurrence of each name. Cheap on size, unbreakable for downstream readers.
    _dedupe_zip_entries(final_path)
    print(f"  saved ({final_path.stat().st_size:,} bytes)")

    print("\n[4] Verify opens cleanly")
    opens = False
    slide_count = 0
    try:
        verify = Presentation(str(final_path))
        slide_count = len(verify.slides)
        opens = True
        print(f"  ok - slide count = {slide_count}")
    except Exception as e:
        print(f"  FAIL: {type(e).__name__}: {e}")
        failures.append(f"- **reload**: {type(e).__name__}: {e}")

    print("\n[5] Render every slide to PNG")
    pngs_dir = out_dir / "final_pngs"
    render_total = 0
    render_ok = 0
    render_fail = 0
    try:
        from render_slides import render_libre
        render_libre(final_path, pngs_dir, dpi=120)
        pngs = sorted(pngs_dir.glob("slide_*.png"))
        render_total = max(slide_count, len(pngs))
        render_ok = len(pngs)
        render_fail = max(0, render_total - render_ok)
        print(f"  rendered {render_ok} png(s) into {pngs_dir}")
    except Exception as e:
        tb = traceback.format_exc().strip().splitlines()[-1]
        msg = f"{type(e).__name__}: {e} | {tb}"
        print(f"  FAIL: {msg}")
        failures.append(f"- **render**: {msg}")
        render_total = slide_count
        render_fail = slide_count

    print("\n[6] Write COMPILED.md")
    content = COMPILED_TEMPLATE.format(
        ts=datetime.now().isoformat(timespec="seconds"),
        out=out_dir,
        template=template_path,
        final=final_path,
        rows="\n".join(rows) if rows else "| (no picks) |",
        slide_count=slide_count,
        opens="yes" if opens else "NO",
        render_total=render_total,
        render_ok=render_ok,
        render_fail=render_fail,
        failures="\n".join(failures) if failures else "(none)",
    )
    compiled_md = out_dir / "COMPILED.md"
    compiled_md.write_text(content, encoding="utf-8")
    print(f"  {compiled_md}")

    print("\n" + "=" * 72)
    print("DONE - compile complete.")
    print(f"  Copied  : {copied_count} / {len(picks)}")
    print(f"  Opens   : {'yes' if opens else 'NO'}")
    print(f"  Slides  : {slide_count}")
    print(f"  Renders : {render_ok} / {render_total}")
    print(f"  Report  : {compiled_md}")
    print(f"  Deck    : {final_path}")
    print("=" * 72)
    print("\nNOT DONE YET — the deck must pass slide-qc before it is delivered.")
    print("A self-rendered PDF is not QC. Run the reviewer on the compiled deck:")
    print(f"  invoke the slide-qc skill on:  {final_path}")
    print("Do not tell the user the deck is finished or 'QC'd' until slide-qc has run.")

    return 0 if (opens and render_fail == 0 and not failures) else 1


if __name__ == "__main__":
    sys.exit(main())
