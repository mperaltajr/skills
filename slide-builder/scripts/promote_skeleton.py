"""
promote_skeleton.py — Convert any PPTX slide into a skeleton with {{TOKEN}} placeholders.

Usage:
    python promote_skeleton.py <pptx_path> <slide_number> <skeleton_name>

Arguments:
    pptx_path       Path to the source PPTX file
    slide_number    1-indexed slide number to promote
    skeleton_name   Name for the skeleton family (e.g. "two-panel", "cover-dark")

Output:
    slide-builder/skeletons/<skeleton_name>/
        <skeleton_name>.pptx   — tokenized PPTX
        skeleton.yaml           — metadata (edit after to fill page-types and best-for)

The tool prints a report of every shape: what was tokenized, what was kept structural, and why.
After running, review the YAML and rename any tokens that need clearer names.
"""

import sys
import os
import re
import yaml
import pathlib

_PATCHES_DIR = pathlib.Path(__file__).resolve().parent.parent / "patches"
if str(_PATCHES_DIR) not in sys.path:
    sys.path.insert(0, str(_PATCHES_DIR))
from patches import get_blank_layout

# Force UTF-8 output on Windows so arrow/dash characters don't crash
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# Structural text patterns — shapes matching these are NOT tokenized
# They are part of the fixed layout (labels, badges, page numbers, dividers)
# ---------------------------------------------------------------------------
STRUCTURAL_PATTERNS = [
    # Note: bare digits are handled separately in name_token (position-aware) —
    # they may be page numbers (footer zone) or hero numerals (body zone).
    r"^DRAFT$",
    r"^CONFIDENTIAL$",
    r"^STRICTLY CONFIDENTIAL$",
    r"^INTERNAL ONLY$",
    r"^TAKEAWAY\s*/\s*INSIGHT$",
    r"^KEY\s+TAKEAWAY$",
    r"^INSIGHT$",
    r"^FINDING$",
    r"^RECOMMENDATION$",
    r"^SUMMARY$",
    r"^\[Chart area.*\]$",             # chart placeholder label
    r"^\[.*placeholder.*\]$",
]

FOOTER_ZONE_THRESHOLD = 0.85  # shapes below this fraction of slide height are footer

FOOTNOTE_PATTERNS = [
    r"^\d+\.\s+",                      # "1. Some footnote"
    r"^Source:",
    r"^Sources:",
    r"^Note:",
    r"^Notes:",
]


def is_structural(text):
    t = text.strip()
    if not t:
        return True, "empty"
    for pat in STRUCTURAL_PATTERNS:
        if re.match(pat, t, re.IGNORECASE):
            return True, f"matches structural pattern: {pat}"
    # All-caps short text (≤ 5 words) — likely a hardcoded label.
    # Require at least one alphabetic character so bare digits (page numbers,
    # hero numerals) don't get swept up here — they're handled in name_token.
    words = t.split()
    has_alpha = any(c.isalpha() for c in t)
    if has_alpha and t == t.upper() and len(words) <= 5 and len(t) <= 40:
        return True, "all-caps short label"
    return False, ""


def is_footnote(text):
    t = text.strip()
    for pat in FOOTNOTE_PATTERNS:
        if re.match(pat, t, re.IGNORECASE):
            return True
    return False


def get_column_zone(shape_left_inches, slide_width_inches, num_columns):
    """Return which column zone a shape sits in (0-indexed)."""
    col_width = slide_width_inches / num_columns
    col = int(shape_left_inches / col_width)
    return min(col, num_columns - 1)


def infer_visual_structure(text_shapes, num_columns, slide_width):
    """Classify the slide's visual structure for skeleton-matching disambiguation.

    Vocabulary:
      hero                   — 1-2 main shapes (divider, hero stat, single quote)
      single-column          — vertical stack of items in one column
      two-panel              — two large side-by-side panels
      N-column-pillars       — N tall vertical columns with rich content per column
      N-column-cards         — N wide cards in a row, moderate body
      N-column-rows          — N-column row table (cells short, multiple rows stacked)

    The user's slide 5 ('3 tall columns with multiple bullets each') matched the
    existing three-column skeleton on page_types alone but the visual was wrong —
    that skeleton is 'three-column-rows', the user wanted 'three-column-pillars'.
    This field lets the matcher disambiguate without rendering the PPTX first.
    """
    n = len(text_shapes)
    if n <= 2:
        return "hero"
    if num_columns == 1:
        return "single-column"
    if num_columns == 2:
        return "two-panel"
    # Per-column tallest text shape — discriminates pillars vs rows vs cards
    col_w = slide_width / num_columns
    col_max_h = {}
    for s in text_shapes:
        if s.left is None or s.height is None:
            continue
        try:
            col = min(int(s.left.inches / col_w), num_columns - 1)
            col_max_h[col] = max(col_max_h.get(col, 0), s.height.inches)
        except Exception:
            continue
    if not col_max_h:
        return f"{num_columns}-column-cards"
    max_pillar_h = max(col_max_h.values())
    if max_pillar_h >= 2.5:
        return f"{num_columns}-column-pillars"
    if max_pillar_h < 1.2:
        return f"{num_columns}-column-rows"
    return f"{num_columns}-column-cards"


def infer_num_columns(shapes_with_text, slide_width):
    """Guess how many columns the layout has based on x-position clustering.
    Tries 4, then 3, then 2, then 1. Each tier requires every zone to be populated.
    """
    xs = []
    for s in shapes_with_text:
        if s.left is not None:
            xs.append(s.left.inches)
    if not xs:
        return 1
    xs.sort()
    # 4-column test
    quarter = slide_width / 4
    zone_counts_4 = [0, 0, 0, 0]
    for x in xs:
        z = min(int(x / quarter), 3)
        zone_counts_4[z] += 1
    if all(c >= 1 for c in zone_counts_4):
        return 4
    # 3-column test
    third = slide_width / 3
    zone_counts_3 = [0, 0, 0]
    for x in xs:
        z = min(int(x / third), 2)
        zone_counts_3[z] += 1
    if all(c >= 1 for c in zone_counts_3):
        return 3
    # 2-column test
    mid = slide_width / 2
    left_count = sum(1 for x in xs if x < mid - 0.5)
    right_count = sum(1 for x in xs if x > mid + 0.5)
    if left_count >= 1 and right_count >= 1:
        return 2
    return 1


def name_token(shape, shape_index, column_zone, num_columns, col_counters,
               title_done, subtitle_done, slide_h, numeral_done):
    """
    Return a token name for this shape based on its role and position.
    col_counters: dict mapping column_zone -> list of roles seen so far
    slide_h: slide height in inches (used to detect footer-zone shapes)
    numeral_done: 1-element list used as a mutable flag — only one HERO_NUMERAL per slide
    """
    name = shape.name or ""
    text = shape.text_frame.text.strip() if shape.has_text_frame else ""
    shape_top = shape.top.inches if shape.top else 0

    # Bare digit handling — position-aware AND size-aware
    if re.match(r"^\d+$", text):
        if shape_top > slide_h * FOOTER_ZONE_THRESHOLD:
            return "PAGE_NUMBER"
        # Small bare-digit shapes are decoration (numbered circles on cards,
        # step indicators on chevrons). Only large bare-digit shapes are
        # hero numerals (section dividers, KPI hero numbers).
        try:
            shape_h = shape.height.inches if shape.height else 0
        except Exception:
            shape_h = 0
        if shape_h < 1.0:
            # Decoration — keep structural by signaling caller to skip tokenization
            return None
        if not numeral_done[0]:
            numeral_done[0] = True
            return "HERO_NUMERAL"
        # Second large bare digit in body zone — fall through to positional naming

    # Shape-name-based overrides (python-pptx placeholder types)
    if re.match(r"Title\s+\d+", name) or re.match(r"Title$", name):
        if not title_done[0]:
            title_done[0] = True
            return "ACTION_TITLE"
    if re.match(r"Subtitle\s+\d+", name) or re.match(r"Subtitle$", name):
        if not subtitle_done[0]:
            subtitle_done[0] = True
            return "SUB_HEADLINE"

    # Footnote / source detection by content
    if is_footnote(text):
        if text.strip().lower().startswith("source"):
            return "SOURCE"
        return "FOOTNOTE"

    # Position-based naming
    if num_columns == 1:
        col_list = col_counters.setdefault(0, [])
        if not col_list:
            col_list.append("heading")
            return "HEADING"
        elif len(col_list) == 1:
            col_list.append("body")
            return "BODY_TEXT"
        else:
            idx = len(col_list)
            col_list.append(f"bullet_{idx}")
            return f"BULLET_{idx}"

    elif num_columns == 2:
        labels = ["LEFT", "RIGHT"]
        label = labels[min(column_zone, 1)]
        col_list = col_counters.setdefault(column_zone, [])
        if not col_list:
            col_list.append("heading")
            return f"{label}_HEADING"
        elif len(col_list) == 1:
            col_list.append("body")
            return f"{label}_BODY"
        else:
            idx = len(col_list)
            col_list.append(f"item_{idx}")
            return f"{label}_ITEM_{idx}"

    elif num_columns == 3:
        labels = ["COL_A", "COL_B", "COL_C"]
        label = labels[min(column_zone, 2)]
        col_list = col_counters.setdefault(column_zone, [])
        if not col_list:
            col_list.append("heading")
            return f"{label}_HEADING"
        elif len(col_list) == 1:
            col_list.append("body")
            return f"{label}_BODY"
        else:
            idx = len(col_list)
            col_list.append(f"item_{idx}")
            return f"{label}_ITEM_{idx}"

    elif num_columns == 4:
        labels = ["COL_A", "COL_B", "COL_C", "COL_D"]
        label = labels[min(column_zone, 3)]
        col_list = col_counters.setdefault(column_zone, [])
        if not col_list:
            col_list.append("heading")
            return f"{label}_HEADING"
        elif len(col_list) == 1:
            col_list.append("body")
            return f"{label}_BODY"
        else:
            idx = len(col_list)
            col_list.append(f"item_{idx}")
            return f"{label}_ITEM_{idx}"

    return f"TEXT_{shape_index:02d}"


def replace_text_in_shape(shape, token):
    """Replace all text in a shape's text frame with {{TOKEN}}, preserving run formatting."""
    tf = shape.text_frame
    token_str = "{{" + token + "}}"
    # Clear all paragraphs after the first, set first para's first run to token
    first_para = tf.paragraphs[0]
    # Preserve formatting of first run
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = token_str
        # Clear remaining runs in first paragraph
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        # No runs — set via paragraph XML directly
        from pptx.oxml.ns import qn
        from lxml import etree
        # Remove all 'a:r' elements, add one with the token
        for r_elem in first_para._p.findall(qn("a:r")):
            first_para._p.remove(r_elem)
        r_elem = etree.SubElement(first_para._p, qn("a:r"))
        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = token_str

    # Remove all paragraphs after the first
    for para in tf.paragraphs[1:]:
        p = para._p
        p.getparent().remove(p)


def promote(pptx_path, slide_number, skeleton_name, output_root):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_number - 1]
    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches

    # Collect shapes with text
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]

    # Sort by top position then left position for consistent naming order
    text_shapes.sort(key=lambda s: (s.top.inches if s.top else 0,
                                    s.left.inches if s.left else 0))

    # Infer column count from non-structural shapes
    content_shapes = []
    for s in text_shapes:
        structural, _ = is_structural(s.text_frame.text)
        if not structural:
            content_shapes.append(s)

    num_columns = infer_num_columns(content_shapes, slide_w)
    visual_structure = infer_visual_structure(content_shapes, num_columns, slide_w)

    # Process each shape
    report = []
    tokens = {}
    col_counters = {}
    title_done = [False]
    subtitle_done = [False]
    numeral_done = [False]

    for i, shape in enumerate(text_shapes):
        text = shape.text_frame.text.strip()
        structural, reason = is_structural(text)

        if structural:
            report.append({
                "shape": shape.name,
                "text": text[:60],
                "action": "KEPT (structural)",
                "reason": reason,
            })
            continue

        col_zone = get_column_zone(
            shape.left.inches if shape.left else 0,
            slide_w, num_columns
        )
        token = name_token(shape, i, col_zone, num_columns, col_counters,
                           title_done, subtitle_done, slide_h, numeral_done)

        if token is None:
            report.append({
                "shape": shape.name,
                "text": text[:60],
                "action": "KEPT (decoration)",
                "reason": "small bare-digit shape (sequence marker)",
            })
            continue

        replace_text_in_shape(shape, token)
        tokens[token] = text[:80]  # store original text as hint in YAML

        report.append({
            "shape": shape.name,
            "text": text[:60],
            "action": f"TOKENIZED → {{{{{token}}}}}",
            "reason": f"col_zone={col_zone}, num_columns={num_columns}",
        })

    # Extract the single promoted slide into a new presentation
    # python-pptx has no built-in slide extraction, so we copy the slide XML
    # into a fresh blank deck of the same dimensions.
    from pptx.util import Emu
    from lxml import etree
    import copy

    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    # Copy the slide layout from the source (use blank layout as base)
    slide_layout = get_blank_layout(new_prs)
    new_slide = new_prs.slides.add_slide(slide_layout)

    # Clear default placeholder shapes from blank layout
    sp_tree = new_slide.shapes._spTree
    for sp in sp_tree.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}sp"):
        sp_tree.remove(sp)
    # Remove all child elements except nvGrpSpPr and grpSpPr
    for child in list(sp_tree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            sp_tree.remove(child)

    # Resolve placeholder positions from the source layout before copying.
    # Placeholder shapes (Title, Subtitle) store position in the slide layout,
    # not in the slide XML. When copied to a new blank layout the position is lost.
    # Fix: read resolved positions first (before mutating XML), then write explicitly.
    PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
    DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Step 1: collect resolved positions for all placeholder shapes
    ph_positions = {}
    for shape in slide.shapes:
        sp_elem = shape._element
        ph_elems = sp_elem.findall(".//{%s}ph" % PML)
        if not ph_elems:
            continue
        try:
            ph_positions[id(sp_elem)] = (
                int(shape.left), int(shape.top),
                int(shape.width), int(shape.height)
            )
        except Exception:
            pass  # skip if position unresolvable

    # Step 2: write resolved positions explicitly into the shape XML
    for shape in slide.shapes:
        sp_elem = shape._element
        if id(sp_elem) not in ph_positions:
            continue
        left, top, width, height = ph_positions[id(sp_elem)]
        spPr = sp_elem.find("{%s}spPr" % PML)
        if spPr is None:
            continue
        xfrm = spPr.find("{%s}xfrm" % DML)
        if xfrm is None:
            xfrm = etree.Element("{%s}xfrm" % DML)
            spPr.insert(0, xfrm)
        off = xfrm.find("{%s}off" % DML)
        ext = xfrm.find("{%s}ext" % DML)
        if off is None:
            off = etree.SubElement(xfrm, "{%s}off" % DML)
        if ext is None:
            ext = etree.SubElement(xfrm, "{%s}ext" % DML)
        off.set("x", str(left))
        off.set("y", str(top))
        ext.set("cx", str(width))
        ext.set("cy", str(height))

    # Copy all shape elements from the source slide
    src_sp_tree = slide.shapes._spTree
    for child in src_sp_tree:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            sp_tree.append(copy.deepcopy(child))

    # Copy slide background if present
    if slide._element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg") is not None:
        bg = slide._element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
        new_slide._element.insert(0, copy.deepcopy(bg))

    out_dir = os.path.join(output_root, skeleton_name)
    os.makedirs(out_dir, exist_ok=True)
    pptx_out = os.path.join(out_dir, f"{skeleton_name}.pptx")
    new_prs.save(pptx_out)

    # Record source metadata so smoke tests / builds can graft onto the right layout
    try:
        source_layout = slide.slide_layout.name
    except Exception:
        source_layout = ""

    # Write skeleton.yaml
    yaml_data = {
        "id": skeleton_name,
        "category": skeleton_name.split("-")[0] if "-" in skeleton_name else skeleton_name,
        "page_types": [],   # fill in after
        "best_for": "",     # fill in after
        "num_columns": num_columns,
        # Critical for skeleton matching — see feedback-2026-05-18.md Error 2.
        # Without this, the model matches on page_types/best_for alone and gets
        # the visual structure wrong (e.g. row-table vs vertical-pillars).
        "visual_structure": visual_structure,
        "source_pptx": os.path.basename(pptx_path),
        "source_slide_num": slide_number,
        "source_layout_name": source_layout,
        "tokens": {k: f"[{v}]" for k, v in tokens.items()},
    }
    yaml_out = os.path.join(out_dir, "skeleton.yaml")
    with open(yaml_out, "w", encoding="utf-8") as f:
        yaml.dump(yaml_data, f, allow_unicode=True, sort_keys=False, default_flow_style=False)

    # Print report
    print(f"\n{'='*60}")
    print(f"  promote_skeleton — {skeleton_name}")
    print(f"  Source: {os.path.basename(pptx_path)}, slide {slide_number}")
    print(f"  Columns detected: {num_columns}")
    print(f"  Visual structure: {visual_structure}")
    print(f"{'='*60}")
    print(f"\n{'ACTION':<30} {'SHAPE':<25} {'TEXT PREVIEW'}")
    print("-" * 80)
    for r in report:
        print(f"{r['action']:<30} {r['shape']:<25} {r['text']}")

    print(f"\nTokens created: {len(tokens)}")
    for k, v in tokens.items():
        print(f"  {{{{ {k} }}}}  ←  \"{v[:50]}\"")

    print(f"\nOutput:")
    print(f"  PPTX:  {pptx_out}")
    print(f"  YAML:  {yaml_out}")
    print(f"\nNext: edit skeleton.yaml to fill in page_types and best_for.")
    print(f"      Rename any token keys that need clearer names (update both YAML and PPTX text).")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Usage: python promote_skeleton.py <pptx_path> <slide_number> <skeleton_name>")
        print("Example: python promote_skeleton.py my_deck.pptx 3 two-panel-finding")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slide_num = int(sys.argv[2])
    skeleton_name = sys.argv[3]

    # Output root = slide-builder/skeletons/ relative to this script's location
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_root = os.path.join(script_dir, "..", "skeletons")
    output_root = os.path.normpath(output_root)

    promote(pptx_path, slide_num, skeleton_name, output_root)
