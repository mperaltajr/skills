#!/usr/bin/env python3
"""adopt_deck.py — bring an EXTERNAL .pptx into the pipeline so specific slides
can be rebuilt/redesigned at Slide-Lab quality (option 6b), then spliced back.

Slide Lab's per-slide rebuild (`build_deck.py --slide N`) needs the pipeline's
`_meta.json` + a recorded brief + a registered template. An external client
deck has none of these. `adopt_deck.py` synthesizes them WITHOUT rebuilding the
deck from scratch:

  1. Requires the deck to be registered **as its own template** first (so the
     rebuilt slide matches its neighbors' masters/layouts/brand). Registration
     is a standalone step — if the sidecars are missing this STOPS and tells you
     to run option 7, never registers inline.
  2. Extracts each slide's text into a brief (`mode: rebuild-slice`, so the
     narrative gate is bypassed).
  3. Reads each slide's real layout name and writes `_meta.json` with an
     `adopted_source` marker + a copy of the original at `adopted_source.pptx`.

After adopting: `build_deck.py --slide N` → worker → `finalize_deck.py --slide N`
→ pick in REVIEW → `compile_picks.py --splice-into <original>` (replaces slide N
in a copy, keeps every other slide). The `adopted_source` marker makes a plain
compile refuse (it would drop the un-rebuilt slides).

Run:  py -3 adopt_deck.py <external.pptx> --slides 3,7 --out <out_dir>
      (python3 on macOS/Linux)
"""
from __future__ import annotations

import argparse
import datetime as _dt
import json
import shutil
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import re as _re  # noqa: E402
import _paths as _p  # noqa: E402
from _meta_schema import META_SCHEMA_VERSION_CURRENT  # noqa: E402
from pptx import Presentation  # noqa: E402


def _slug(text: str) -> str:
    s = _re.sub(r"[^a-z0-9]+", "-", (text or "").lower()).strip("-")
    return s or "adopted"


def _registration_missing(deck: Path) -> list[str]:
    """Return the list of missing registration sidecars for `deck` (empty = ok)."""
    missing = []
    for label, path in (("brand.yml", _p.brand_yml(deck)),
                        ("chrome.yml", _p.chrome_yml(deck)),
                        ("theme.json", _p.theme_json(deck))):
        if not Path(path).exists():
            missing.append(label)
    return missing


def _extract_slide(slide) -> tuple[str, list[str], str]:
    """Return (title, body_lines, layout_name) from a slide's text shapes.

    Title = the title placeholder if present, else the first text shape.
    Body = every other non-empty text shape, one bullet per paragraph.
    """
    try:
        layout_name = slide.slide_layout.name or ""
    except Exception:
        layout_name = ""
    title = ""
    body: list[str] = []
    title_shape = None
    # Prefer a real title placeholder.
    for sh in slide.shapes:
        if sh.has_text_frame and getattr(sh, "is_placeholder", False):
            try:
                if sh.placeholder_format.idx == 0 or "title" in (sh.name or "").lower():
                    _t = sh.text_frame.text.strip()
                    # First line only — a multi-line title would orphan text
                    # above the brief's field labels (SLIDE_HEADER_RE is single-line).
                    title = _t.splitlines()[0].strip() if _t else ""
                    title_shape = sh
                    break
            except Exception:
                pass
    for sh in slide.shapes:
        if not sh.has_text_frame or sh is title_shape:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        if not title:  # no placeholder title — first text shape becomes title
            title = txt.splitlines()[0].strip()
            title_shape = sh
            rest = txt.splitlines()[1:]
            body.extend(l.strip() for l in rest if l.strip())
            continue
        body.extend(l.strip() for l in txt.splitlines() if l.strip())
    return title, body, layout_name


def _slide_section(n: int, title: str, body: list[str], layout: str) -> str:
    ev = "\n".join(f"- {b}" for b in body) if body else "- (no extractable text — enrich before rebuilding)"
    return (
        f"## Slide {n} — {title or f'Slide {n}'}\n"
        f"**Layout:** {layout}\n"
        f"**Slide type:** Adopted (rebuild)\n"
        f"**Governing thought:** {title or '(fill in)'}\n"
        f"**The takeaway:** {title or '(fill in)'}\n"
        f"**Evidence / content:**\n{ev}\n\n"
    )


def main(argv) -> int:
    ap = argparse.ArgumentParser(description="Adopt an external .pptx for per-slide rebuild (option 6b).")
    ap.add_argument("deck", help="The external .pptx to adopt.")
    ap.add_argument("--out", required=True, help="Output dir for the adopted build (the pipeline's <out>).")
    ap.add_argument("--slides", default="", help="Comma-separated 1-based slide numbers you intend to rebuild (informational; the brief carries all slides).")
    args = ap.parse_args(argv)

    deck = Path(args.deck).resolve()
    if not deck.exists():
        print(f"[error] deck not found: {deck}", file=sys.stderr)
        return 2

    # 1. Registration gate — standalone, never inline (absolute rule).
    missing = _registration_missing(deck)
    if missing:
        print(f"ERROR: this deck isn't registered as a template yet (missing: {', '.join(missing)}).")
        print("       Register it as its OWN template first (option 7), so the rebuilt")
        print("       slide matches the rest of the deck, then re-run adopt:")
        print(f"         py -3 register_template.py propose \"{deck}\"")
        print("       (confirm the colors + mock slide, then register_template.py confirm)")
        return 7

    out_dir = Path(args.out).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(deck))
    slides = list(prs.slides)
    n_slides = len(slides)

    # 2. Extract every slide into a rebuild-slice brief.
    layouts: list[str] = []
    titles: list[str] = []
    sections = []
    for i, slide in enumerate(slides, start=1):
        title, body, layout = _extract_slide(slide)
        titles.append(title)
        layouts.append(layout)
        sections.append(_slide_section(i, title, body, layout))
    default_layout = layouts[0] if layouts else ""
    brief_text = (
        f"---\n"
        f"client_template: {deck}\n"
        f"deck_type: Adopted external deck\n"
        f"default_layout: {default_layout}\n"
        f"mode: rebuild-slice\n"
        f"---\n\n"
        f"## Deck-level design notes\n\n"
        f"Adopted from an external deck ({deck.name}) for per-slide rebuild. "
        f"Extracted text is a starting point — enrich the target slide's Evidence "
        f"before rebuilding (charts/images/SmartArt don't extract).\n\n"
        + "".join(sections)
    )
    brief_path = out_dir / "adopted_brief.md"
    brief_path.write_text(brief_text, encoding="utf-8")

    # 3. Synthesize _meta.json + copy the original as the splice target.
    # Shape must satisfy _meta_schema.MetaJson — build_deck.py --slide N
    # hard-validates it (schema_version, out, client_slug, deck_meta required).
    meta = {
        "schema_version": META_SCHEMA_VERSION_CURRENT,
        "template": str(deck),
        "brief": str(brief_path),
        "out": str(out_dir),
        "client_slug": _slug(deck.stem),
        "adopted_source": str(deck),
        "slide_count": n_slides,
        "generated_at": _dt.datetime.now().isoformat(timespec="seconds"),
        "slides": [
            {"n": i, "title": titles[i - 1] or f"Slide {i}",
             "layout": layouts[i - 1], "page_type": ""}
            for i in range(1, n_slides + 1)
        ],
        "deck_meta": {
            "deck_type": "Adopted external deck",
            "governing_thought": "",
            "audience": "",
        },
    }
    (out_dir / "_meta.json").write_text(
        json.dumps(meta, indent=2, ensure_ascii=False), encoding="utf-8")
    shutil.copy2(deck, out_dir / "adopted_source.pptx")

    # Pre-create the slide dirs the user plans to rebuild (build_deck --slide
    # will populate them; creating them now makes the next step obvious).
    want = [s.strip() for s in args.slides.split(",") if s.strip()]
    targets = []
    for s in want:
        try:
            n = int(s)
        except ValueError:
            continue
        if 1 <= n <= n_slides:
            (out_dir / _p.slide_key(n)).mkdir(exist_ok=True)
            targets.append(n)

    print("=" * 72)
    print(f"Adopted {deck.name}  ({n_slides} slides)  ->  {out_dir}")
    print("=" * 72)
    print(f"  brief:          {brief_path}")
    print(f"  _meta.json:     {out_dir / '_meta.json'}  (adopted_source marked)")
    print(f"  splice target:  {out_dir / 'adopted_source.pptx'}")
    print("\nNext — for each slide you want to redesign (e.g. "
          + (", ".join(str(t) for t in targets) if targets else "N") + "):")
    print(f"  1. Enrich that slide's Evidence in {brief_path.name} if the extract is thin.")
    print(f"  2. py -3 build_deck.py --slide N --out {out_dir} --template \"{deck}\"")
    print(f"  3. dispatch the slide-builder-worker for slide N. If it writes option_A.html")
    print(f"     (sketch path, the default), also dispatch slide-builder-translator on the")
    print(f"     pick to make option_A_native.py. Then:")
    print(f"       py -3 finalize_deck.py --slide N --out {out_dir} --template \"{deck}\"")
    print(f"  4. pick in REVIEW.html, then splice back (keeps every other slide):")
    print(f"       py -3 compile_picks.py --out {out_dir} --splice-into \"{deck}\"")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
