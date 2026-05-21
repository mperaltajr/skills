"""
extract_icons.py — One-time icon extraction from Accenture icon library PPTX.

Uses POSITION-BASED lookup: each icon in icon-index.json has source_x_pct /
source_y_pct coordinates that point to the exact icon shape as visually identified
in the exported PNG slides. This is far more reliable than keyword-proximity
matching (which was returning wrong icons due to label ambiguity).

Extracted shapes are saved as raw XML to icons/<name>.xml.
Build-time insertion reads the .xml files directly — no PPTX opening at build time.

Usage:
    py -3 extract_icons.py                          # extract all icons
    py -3 extract_icons.py --icon gear              # extract one icon
    py -3 extract_icons.py --icon gear --verify     # extract + open PNG preview
    py -3 extract_icons.py --inspect 131 0.71 0.76  # list shapes near a coordinate
                                                    # (use to find coords for new icons)
"""

from __future__ import annotations

import argparse
import json
import math
import pathlib
import sys

from lxml import etree
from pptx import Presentation

_PATCHES_DIR = pathlib.Path(__file__).resolve().parent.parent / "patches"
if str(_PATCHES_DIR) not in sys.path:
    sys.path.insert(0, str(_PATCHES_DIR))
from patches import get_blank_layout

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

import os as _os

_SCRIPTS_DIR = pathlib.Path(__file__).resolve().parent
_SKILLS_DIR = _SCRIPTS_DIR.parent
_ICONS_DIR = _SKILLS_DIR / "icons"
_ICON_INDEX = _ICONS_DIR / "icon-index.json"

# Resolve the source icon-library PPTX in priority order:
#   1. SLIDE_LAB_ICON_LIBRARY env var (absolute path)
#   2. --icon-library CLI arg (set _ICON_PPTX before extraction calls)
#   3. default placeholder — extraction will error with a clear install hint
_DEFAULT_ICON_PPTX = _ICONS_DIR / "_source-library.pptx"
_ICON_PPTX = pathlib.Path(_os.environ.get("SLIDE_LAB_ICON_LIBRARY") or _DEFAULT_ICON_PPTX)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ---------------------------------------------------------------------------
# Shape enumeration helpers
# ---------------------------------------------------------------------------

def _get_slide_shapes(slide):
    """Return all direct children of spTree (sp, grpSp, pic, cxnSp)."""
    sp_tree = slide._element.find(f".//{{{_NS_P}}}spTree")
    if sp_tree is None:
        return []
    shapes = []
    for child in sp_tree:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("sp", "grpSp", "pic", "cxnSp"):
            shapes.append(child)
    return shapes


def _get_text(elem) -> str:
    """Extract all visible text from an element."""
    parts = [t.text.strip() for t in elem.iter(f"{{{_NS_A}}}t") if t.text and t.text.strip()]
    return "; ".join(parts)


def _get_bounding_box(elem):
    """
    Return (x, y, cx, cy) in EMU for a shape element, or None if not available.

    For grpSp: reads from grpSpPr/xfrm (outer slide-space bounding box).
    For sp/pic/cxnSp: reads from spPr/xfrm.
    Falls back to first xfrm found anywhere in the element.
    """
    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

    if tag == "grpSp":
        grp_pr = elem.find(f"{{{_NS_P}}}grpSpPr")
        xfrm = grp_pr.find(f"{{{_NS_A}}}xfrm") if grp_pr is not None else None
    else:
        sp_pr = elem.find(f"{{{_NS_P}}}spPr")
        xfrm = sp_pr.find(f"{{{_NS_A}}}xfrm") if sp_pr is not None else None

    # Fallback: search anywhere
    if xfrm is None:
        xfrm = elem.find(f".//{{{_NS_A}}}xfrm")

    if xfrm is None:
        return None

    off = xfrm.find(f"{{{_NS_A}}}off")
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if off is None or ext is None:
        return None

    try:
        x  = int(off.get("x", 0))
        y  = int(off.get("y", 0))
        cx = int(ext.get("cx", 0))
        cy = int(ext.get("cy", 0))
    except (ValueError, TypeError):
        return None

    return (x, y, cx, cy)


# ---------------------------------------------------------------------------
# Position-based shape finder  (primary approach)
# ---------------------------------------------------------------------------

def find_shape_at_position(slide, x_pct: float, y_pct: float, prs: Presentation):
    """
    Find the non-text shape that best matches the given fractional slide position.

    Strategy:
      1. Convert (x_pct, y_pct) -> EMU using actual slide dimensions.
      2. Collect all non-text shapes whose bounding box contains the target point.
      3. Return the one with the smallest bounding-box area (most specific hit).
      4. If no containment hit, fall back to the shape whose center is nearest.

    Args:
        slide:  python-pptx Slide object
        x_pct:  fractional x position (0.0 = left edge, 1.0 = right edge)
        y_pct:  fractional y position (0.0 = top edge, 1.0 = bottom edge)
        prs:    Presentation object (to read actual slide dimensions)

    Returns shape element or None.
    """
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    target_x = int(x_pct * slide_w)
    target_y = int(y_pct * slide_h)

    shapes = _get_slide_shapes(slide)

    # Pass 1: containment — smallest bounding box wins
    best_contain = None
    best_area = float("inf")

    for sh in shapes:
        # Skip pure text boxes
        tag = sh.tag.split("}")[-1] if "}" in sh.tag else sh.tag
        if tag == "sp" and _get_text(sh):
            continue

        bbox = _get_bounding_box(sh)
        if bbox is None:
            continue
        x, y, cx, cy = bbox
        if cx <= 0 or cy <= 0:
            continue

        if x <= target_x <= x + cx and y <= target_y <= y + cy:
            area = cx * cy
            if area < best_area:
                best_area = area
                best_contain = sh

    if best_contain is not None:
        return best_contain

    # Pass 2: proximity fallback — nearest center
    best_near = None
    best_dist = float("inf")

    for sh in shapes:
        tag = sh.tag.split("}")[-1] if "}" in sh.tag else sh.tag
        if tag == "sp" and _get_text(sh):
            continue
        bbox = _get_bounding_box(sh)
        if bbox is None:
            continue
        x, y, cx, cy = bbox
        if cx <= 0 or cy <= 0:
            continue
        cx_center = x + cx // 2
        cy_center = y + cy // 2
        dist = math.sqrt((cx_center - target_x) ** 2 + (cy_center - target_y) ** 2)
        if dist < best_dist:
            best_dist = dist
            best_near = sh

    return best_near


# ---------------------------------------------------------------------------
# Keyword-based shape finder  (legacy / fallback for icons without coordinates)
# ---------------------------------------------------------------------------

def _keyword_matches(shape_text: str, keyword: str) -> bool:
    if not keyword:
        return False
    sl = shape_text.lower()
    for token in [t.strip() for t in keyword.lower().split(";") if t.strip()][:3]:
        if len(token) > 2 and token in sl:
            return True
    return False


def _shape_center(elem):
    bbox = _get_bounding_box(elem)
    if bbox is None:
        return (0.0, 0.0)
    x, y, cx, cy = bbox
    return (x + cx / 2.0, y + cy / 2.0)


def find_icon_shape_by_keyword(slide, keyword: str):
    """
    Legacy keyword-proximity lookup. Kept as fallback when position coords
    are not available. Finds the icon shape nearest to keyword text boxes.
    """
    shapes = _get_slide_shapes(slide)
    text_shapes, icon_shapes = [], []
    for sh in shapes:
        tag = sh.tag.split("}")[-1] if "}" in sh.tag else sh.tag
        txt = _get_text(sh)
        if txt and tag == "sp":
            text_shapes.append((sh, txt))
        else:
            icon_shapes.append(sh)

    matches = [(sh, txt) for sh, txt in text_shapes if _keyword_matches(txt, keyword)]
    if not matches:
        return None

    best, best_d = None, float("inf")
    for text_sh, _ in matches:
        tc = _shape_center(text_sh)
        for icon_sh in icon_shapes:
            d = math.sqrt(sum((a - b) ** 2 for a, b in zip(tc, _shape_center(icon_sh))))
            if d < best_d:
                best_d = d
                best = icon_sh
    return best


# ---------------------------------------------------------------------------
# Color normalization: strip all fill colors to black (tinted at build time)
# ---------------------------------------------------------------------------

def normalize_colors(elem) -> None:
    for solid_fill in elem.iter(f"{{{_NS_A}}}solidFill"):
        srgb = solid_fill.find(f"{{{_NS_A}}}srgbClr")
        if srgb is not None:
            srgb.set("val", "000000")


# ---------------------------------------------------------------------------
# Extract one icon
# ---------------------------------------------------------------------------

def extract_icon(prs: Presentation, entry: dict, output_path: pathlib.Path) -> bool:
    """
    Extract a single icon from the source PPTX and save as XML.

    Prefers position-based lookup (source_x_pct / source_y_pct).
    Falls back to keyword-based lookup if coordinates are absent.
    """
    import copy

    name        = entry["name"]
    slide_num   = entry["source_slide"]
    x_pct       = entry.get("source_x_pct")
    y_pct       = entry.get("source_y_pct")
    keyword     = entry.get("concept_label", entry.get("keyword", ""))

    icon_elem = None
    method_used = ""

    # --- Position-based lookup ---
    if x_pct is not None and y_pct is not None:
        if 1 <= slide_num <= len(prs.slides):
            icon_elem = find_shape_at_position(
                prs.slides[slide_num - 1], x_pct, y_pct, prs
            )
            if icon_elem is not None:
                method_used = f"pos ({x_pct:.2f}, {y_pct:.2f}) on slide {slide_num}"

    # --- Keyword fallback ---
    if icon_elem is None and keyword:
        if 1 <= slide_num <= len(prs.slides):
            icon_elem = find_icon_shape_by_keyword(prs.slides[slide_num - 1], keyword)
            if icon_elem is not None:
                method_used = f"keyword '{keyword[:30]}' on slide {slide_num}"

    if icon_elem is None:
        print(f"  FAIL  '{name}' -- not found (slide {slide_num})")
        return False

    saved = copy.deepcopy(icon_elem)
    normalize_colors(saved)

    xml_bytes = etree.tostring(saved, pretty_print=True, xml_declaration=False)
    output_path.write_bytes(xml_bytes)

    size_kb = output_path.stat().st_size // 1024
    tag = saved.tag.split("}")[-1] if "}" in saved.tag else saved.tag
    bbox = _get_bounding_box(saved)
    bbox_str = f" bbox=({bbox[0]//914},{bbox[1]//914},{bbox[2]//914},{bbox[3]//914})pt" if bbox else ""
    print(f"  OK    '{name}' -> {output_path.name} ({size_kb}KB, {tag}){bbox_str}")
    print(f"        [{method_used}]")
    return True


# ---------------------------------------------------------------------------
# Inspect mode: list all shapes near a coordinate (for finding new icon coords)
# ---------------------------------------------------------------------------

def inspect_shapes(prs: Presentation, slide_num: int, x_pct: float, y_pct: float,
                   radius_pct: float = 0.10) -> None:
    """
    Print all non-text shapes within radius of (x_pct, y_pct) on a slide.
    Use this to find the correct coordinates for a new icon.

    Output includes: shape type, bounding box in points, distance, and any text.
    """
    if slide_num < 1 or slide_num > len(prs.slides):
        print(f"ERROR: slide {slide_num} out of range (1-{len(prs.slides)})")
        return

    slide   = prs.slides[slide_num - 1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    tx      = int(x_pct * slide_w)
    ty      = int(y_pct * slide_h)
    r_emu   = int(radius_pct * max(slide_w, slide_h))

    EMU_PER_PT = 12700

    print(f"\nSlide {slide_num} — shapes near ({x_pct:.3f}, {y_pct:.3f})")
    print(f"  Target EMU: ({tx}, {ty})  radius: {r_emu // EMU_PER_PT}pt")
    print(f"  {'Type':<8} {'x,y (pt)':<18} {'w,h (pt)':<16} {'dist(pt)':<10} Text")
    print("  " + "-" * 80)

    shapes = _get_slide_shapes(slide)
    hits = []

    for sh in shapes:
        bbox = _get_bounding_box(sh)
        if bbox is None:
            continue
        x, y, cx, cy = bbox
        cx_c = x + cx // 2
        cy_c = y + cy // 2
        dist = math.sqrt((cx_c - tx) ** 2 + (cy_c - ty) ** 2)
        if dist > r_emu:
            continue
        tag  = sh.tag.split("}")[-1] if "}" in sh.tag else sh.tag
        text = _get_text(sh)[:40] if _get_text(sh) else ""
        hits.append((dist, tag, x, y, cx, cy, text))

    hits.sort(key=lambda h: h[0])

    for dist, tag, x, y, cx, cy in [(h[0], h[1], h[2], h[3], h[4], h[5]) for h in hits]:
        text = next((h[6] for h in hits if h[1] == tag and h[2] == x), "")
        print(f"  {tag:<8} ({x//EMU_PER_PT:>5},{y//EMU_PER_PT:>5})   "
              f"({cx//EMU_PER_PT:>5},{cy//EMU_PER_PT:>5})   "
              f"{dist//EMU_PER_PT:>6}pt   {text}")

    if not hits:
        print(f"  (no shapes within {r_emu // EMU_PER_PT}pt — try a larger radius)")


# ---------------------------------------------------------------------------
# Verify: build a test PPTX and export PNG
# ---------------------------------------------------------------------------

def verify_icon(icon_name: str) -> None:
    """Quick visual verify: insert icon into blank slide, export PNG, open it."""
    sys.path.insert(0, str(_SCRIPTS_DIR))
    from icon_helper import insert_icon
    from pptx import Presentation as _Prs
    from pptx.util import Emu

    prs = _Prs()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    slide = prs.slides.add_slide(get_blank_layout(prs))

    W = H = Emu(914400)   # 1 inch
    L = Emu(9144000 // 2 - 457200)
    T = Emu(5143500 // 2 - 457200)

    ok  = insert_icon(icon_name, slide, int(L), int(T), int(W), int(H), "#A100FF")
    out = pathlib.Path(f"C:/Temp/verify_{icon_name}.pptx")
    prs.save(str(out))
    print(f"  Saved verify PPTX: {out}")

    export_script = _SCRIPTS_DIR.parent.parent / "slide-qc" / "scripts" / "export_slides.py"
    png_out       = pathlib.Path(f"C:/Temp/verify_{icon_name}_png")
    import subprocess
    result = subprocess.run(
        ["py", "-3", str(export_script), str(out), "--out", str(png_out)],
        capture_output=True, text=True
    )
    print(result.stdout.strip())
    if result.returncode == 0:
        png_file = png_out / "slide_01.png"
        if png_file.exists():
            import os
            os.startfile(str(png_file))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Extract icons from Accenture icon library PPTX"
    )
    parser.add_argument("--icon",    help="Extract only this one icon name (default: all)")
    parser.add_argument("--verify",  action="store_true",
                        help="After extraction, open a PNG preview (requires --icon)")
    parser.add_argument("--pptx",    help="Override icon PPTX path")
    parser.add_argument("--inspect", nargs=3, metavar=("SLIDE", "X_PCT", "Y_PCT"),
                        help="List all shapes near a coordinate — use to find coords for new icons. "
                             "Example: --inspect 131 0.71 0.76")
    args = parser.parse_args()

    pptx_path = pathlib.Path(args.pptx) if args.pptx else _ICON_PPTX
    if not pptx_path.exists():
        print(f"ERROR: Icon PPTX not found: {pptx_path}")
        sys.exit(1)

    print(f"Opening icon library: {pptx_path.name}")
    prs = Presentation(str(pptx_path))
    print(f"  {len(prs.slides)} slides loaded")

    # Inspect mode
    if args.inspect:
        slide_num, x_pct, y_pct = int(args.inspect[0]), float(args.inspect[1]), float(args.inspect[2])
        inspect_shapes(prs, slide_num, x_pct, y_pct)
        return

    if not _ICON_INDEX.exists():
        print(f"ERROR: icon-index.json not found: {_ICON_INDEX}")
        sys.exit(1)

    with open(_ICON_INDEX, encoding="utf-8") as f:
        index_data = json.load(f)
    all_entries = index_data["icons"]

    if args.icon:
        entries = [e for e in all_entries if e["name"] == args.icon]
        if not entries:
            print(f"ERROR: '{args.icon}' not found in icon-index.json")
            sys.exit(1)
    else:
        entries = all_entries

    print(f"  Extracting {len(entries)} icon(s) to {_ICONS_DIR}/\n")

    ok_count = 0
    for entry in entries:
        out_file = _ICONS_DIR / f"{entry['name']}.xml"
        if extract_icon(prs, entry, out_file):
            ok_count += 1

    print(f"\nDone: {ok_count}/{len(entries)} icons extracted to {_ICONS_DIR}/")

    if args.verify and args.icon:
        print(f"\nVerifying '{args.icon}'...")
        verify_icon(args.icon)


if __name__ == "__main__":
    main()
