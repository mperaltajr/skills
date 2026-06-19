"""Slide Lab deck orchestrator — Part B: execute agent .py, graft, theme remap, render.

Originally forked from the legacy chassis-vocabulary skill; current behaviors:
  - sys.path includes <skill>/ (for local twins/) + slide-qc/scripts/ (for render_slides).
  - build_pptx() classifies each option_X.py by line-1 token before executing:
      * `# FALLBACK_MERMAID: ...`     -> render .mmd to PNG, assemble fallback PPTX
      * `# SKELETON_REJECTED: ...`    -> skip; rejection surfaces in REVIEW.html
      * normal native python-pptx     -> subprocess/runpy execution
  - _build_mermaid_fallback() invokes render_mermaid.py and assembles a PPTX
    with chrome from twins.helpers + the rendered PNG embedded at body-zone
    coordinates.

Inputs:
  --out PATH       output dir from build_deck.py
  --template PATH  client PPTX template (same one used in Part A)
  --skip-build     skip executing option_X.py (use existing option_X.pptx)
  --skip-render    skip rendering PNGs

Pipeline (unchanged from v1 except step 2's per-option branching):
  1. Discover every <out>/slide_NN/option_X.py and option_X_native.py
     (the Pattern B translator output added in M5).
  2. For each option:
       2a. Classify by line-1 / header token.
       2b. NATIVE                 -> run via subprocess to produce
                                     option_X.pptx (raw, pre-theme).
       2c. PATTERN_B_TRANSLATED   -> run via subprocess; finalize merges
                                     __template_fields__ into placeholders.
       2d. REJECTED               -> skip; logged for the report.
  3. Stash raw output to <out>/slide_NN/_raw/option_X.pptx.
  4. Graft + theme-remap each raw pptx onto the client template; write themed
     to <out>/slide_NN/option_X.pptx.
  5. Render each themed .pptx to PNG (LibreOffice headless, parallel x4).
  6. Run a deterministic per-option QC self-check.
  7. Write <out>/RESULT.md with per-slide status table.

Path convention after finalize, per slide_NN/ folder (unchanged from v1):
  option_X.py             — source (input)
  option_X.mmd            — Mermaid spec (FALLBACK_MERMAID only)
  option_X-mermaid.png    — rendered Mermaid PNG (FALLBACK_MERMAID only)
  option_X.pptx           — THEMED deliverable
  option_X.png            — PNG of themed
  option_X.qc.json        — per-option QC self-check report
  _raw/option_X.pptx      — raw pre-theme output (hidden, for debugging)
  _render_tmp/            — LibreOffice render staging (cleaned up)
"""
from __future__ import annotations

import argparse
import ast
import json
import os
import re
import runpy
import shutil
import subprocess
import sys
import traceback
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from copy import deepcopy
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

# Path setup: twins/ is local to this skill (re-homed Phase 2 cleanup
# 2026-05-26); render_slides lives in the sibling slide-qc skill.
SKILL_ROOT = Path(__file__).resolve().parents[1]              # slide-builder/
QC_SCRIPTS = SKILL_ROOT.parent / "slide-qc" / "scripts"       # render_slides
sys.path.insert(0, str(SKILL_ROOT))
sys.path.insert(0, str(QC_SCRIPTS))

import _paths as _p  # noqa: E402

import hashlib  # noqa: E402
import io  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from twins.client_theme import load_client_theme, apply_theme_to_shape_xml  # noqa: E402
from twins.composer import (  # noqa: E402
    _clear_existing_slides,
    _find_blank_layout,
    _find_named_layout,
    _strip_layout_placeholders,
    TemplatePlaceholderEmptyError,
)
from _chrome_schema import (  # noqa: E402
    ChromeSpec, ChromeSidecarMissingError, ChromeLayoutMissingError,
    load_chrome_yml,
)
import twins.helpers as _twins_helpers  # noqa: E402


# ---------------------------------------------------------------------------
# v2 ADDITION — option-classification tokens
# ---------------------------------------------------------------------------
SKELETON_REJECTED_TOKEN = "# SKELETON_REJECTED:"
# M5 (Pattern B production wiring, 2026-06-17): the Pattern B translator agent
# emits scripts whose header begins with `# CONTEXT_READ:` followed by
# `# BRIEF_IS_AUTHORITATIVE: True` and `# PATTERN: B-translated`. We detect by
# looking for the PATTERN: B-translated line in the first ~10 lines (Spec 4 §4).
PATTERN_B_TRANSLATED_TOKEN = "# PATTERN: B-translated"


def _parse_template_fields(py_path: Path) -> dict[str, str]:
    """Extract the ``__template_fields__`` dict from a translator script header.

    Pattern B translator output (Spec 4 §4) carries the chrome text
    (title/subtitle/footer/page_number) in a structured comment block at the
    top of `option_X_native.py`:

        # __template_fields__ = {
        #     "title": "...",
        #     "subtitle": "...",
        #     ...
        # }

    Returns the parsed dict, or an empty dict if the block is missing or
    malformed. Malformed blocks are non-blocking — finalize falls back to the
    brief's title/subtitle so existing builds keep shipping; the translation
    report flags the parse failure for QC.
    """
    if not py_path.exists():
        return {}
    try:
        text = py_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return {}
    match = re.search(
        r"#\s*__template_fields__\s*=\s*(\{.*?\n#\s*\})",
        text,
        re.DOTALL,
    )
    if not match:
        return {}
    raw = match.group(1)
    # Strip the leading "# " on every line so the result parses as a Python literal.
    payload = "\n".join(line.lstrip("# ").rstrip() for line in raw.splitlines())
    try:
        result = ast.literal_eval(payload)
        return result if isinstance(result, dict) else {}
    except (ValueError, SyntaxError) as exc:
        sys.stderr.write(
            f"  WARN: __template_fields__ in {py_path.name} failed to parse "
            f"({type(exc).__name__}: {exc}); finalize will fall back to the "
            f"brief's title/subtitle.\n"
        )
        return {}


def _classify_option(py_path: Path) -> tuple[str, str]:
    """Read the header of option_X.py and classify the option.

    Returns (status, reason) where status is one of:
        'native'              -> normal python-pptx execution path (Pattern C / legacy)
        'pattern_b_translated' -> Pattern B translator output (M5, 2026-06-17);
                                   executes like native but the parent flow extracts
                                   __template_fields__ for placeholder population.
        'skeleton_rejected'   -> skip; surface in REVIEW.html

    Discriminator is the line-1 token for legacy classifications; Pattern B
    translator output is detected by a `# PATTERN: B-translated` line in the
    first ~10 lines (Spec 4 §4 prescribes the header format).

    Mermaid fallback was retired in M7 (Decision 6, 2026-06-17) — Pattern B
    HTML→PNG supersedes it. Stale scripts carrying the old
    `# FALLBACK_MERMAID:` marker now fall through to the native classifier
    and fail loudly at execution time so the operator re-builds.
    """
    if not py_path.exists():
        return "missing", f"option script not found at {py_path}"
    try:
        text = py_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError) as exc:
        return "missing", f"could not read {py_path}: {exc}"

    line1 = text.splitlines()[0].strip() if text else ""

    if line1.startswith(SKELETON_REJECTED_TOKEN):
        reason = line1[len(SKELETON_REJECTED_TOKEN):].strip()
        return "skeleton_rejected", reason or "SKELETON_REJECTED (no reason given)"
    # M5: detect Pattern B translator output via the PATTERN: B-translated token
    # in the script header (Spec 4 §4 mandates this format).
    head = "\n".join(text.splitlines()[:10])
    if PATTERN_B_TRANSLATED_TOKEN in head:
        return "pattern_b_translated", "Pattern B translator output"
    return "native", ""


# ---------------------------------------------------------------------------
# Per-option render-QC self-check — verbatim from v1
# ---------------------------------------------------------------------------
_SRGB_RE = re.compile(r'srgbClr\s+val="([0-9A-Fa-f]{6})"')

_PLACEHOLDER_PATTERNS = (
    "Click to add",
    "Lorem ipsum",
    "Proceed with Option B",
)
# F1-B (2026-05-26): import the canonical cross-skill placeholder
# constants from twins/helpers.py so finalize_deck doesn't drift from
# slide-qc or slide-builder on the intentional-presenter-prompt contract.
try:
    from twins.helpers import (
        INTENTIONAL_FOOTNOTE_PLACEHOLDER,
        INTENTIONAL_SOURCE_PLACEHOLDER,
    )
    _PLACEHOLDER_ALLOWED = (
        INTENTIONAL_FOOTNOTE_PLACEHOLDER,
        INTENTIONAL_SOURCE_PLACEHOLDER,
    )
except ImportError:
    # Fallback if twins import path isn't set yet during early bootstrap.
    # The strings here MUST match the canonical constants verbatim.
    _PLACEHOLDER_ALLOWED = (
        "[add footnote here or delete]",
        "[add source here or delete]",
    )
_FOOTNOTE_NAME_PREFIXES = ("footnote", "source", "page-number")
_FOOTER_NUM_RE = re.compile(r"^\d+$")
# Two-tier font floor (committee-decision 2026-05-26):
#   - SOFT floor (10.5pt) applies to BODY-ROLE shapes only — named with
#     body|bullet|paragraph|narrative tokens. Workers tag long-form prose
#     this way; that's the text the floor protects.
#   - HARD floor (8.0pt) applies to everything else NOT footnote-like:
#     eyebrows, kickers, category labels, chart legends, axis labels,
#     decision-tree edge labels, column headers, sparkline deltas, etc.
#     These are *legitimately* small in client decks (8-10pt) but anything
#     below 8pt is illegible at projection scale regardless of role.
# This replaces a flat "everything must be >= 10.5pt unless explicitly
# excluded" rule that hit 10/12 false-positive last build because it had
# no concept of small-by-design labels. See `_decisions/cleanup-plan-master`
# committee log + reference/anti-patterns.md § Aesthetics #3.
_BODY_ROLE_NAME_TOKENS = ("body", "bullet", "paragraph", "narrative")
_BODY_FONT_FLOOR_PT = 10.5
_FONT_HARD_FLOOR_PT = 8.0


def _expected_palette_for_theme(theme) -> set:
    palette: set = set()
    try:
        for v in theme.color_map().values():
            if isinstance(v, str) and len(v) == 6:
                palette.add(v.upper())
    except Exception:
        pass
    for attr in ("dk1", "lt1", "dk2", "lt2",
                 "accent1", "accent2", "accent3", "accent4", "accent5", "accent6"):
        v = getattr(theme, attr, None)
        if isinstance(v, str) and len(v) == 6:
            palette.add(v.upper())
    palette.add("FFFFFF")
    palette.add("000000")
    return palette


def _hex_codes_in_pptx(path: Path) -> set:
    hexes: set = set()
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            for name in zf.namelist():
                if name.startswith("ppt/slides/") and name.endswith(".xml"):
                    try:
                        data = zf.read(name).decode("utf-8", errors="ignore")
                    except Exception:
                        continue
                    for m in _SRGB_RE.finditer(data):
                        hexes.add(m.group(1).upper())
    except Exception:
        pass
    return hexes


def _is_footnote_like_name(name: str) -> bool:
    n = (name or "").strip().lower()
    return any(n.startswith(p) for p in _FOOTNOTE_NAME_PREFIXES)


def _is_body_role_name(name: str) -> bool:
    """True iff a shape name marks long-form running text (the only role the
    10.5pt soft floor applies to). All other roles fall through to the 8pt
    hard floor. See the committee decision logged at the _BODY_ROLE_NAME_TOKENS
    constant above for rationale."""
    n = (name or "").strip().lower()
    return any(tok in n for tok in _BODY_ROLE_NAME_TOKENS)


def detect_installed_fonts() -> set:
    """Enumerate installed fonts via PowerShell. Returns set of family names
    (empty set on any failure, but the failure is now logged loudly).

    Pre-2026-06-02 the function silently swallowed every exception and
    returned an empty set, which suppressed font-mismatch warnings for
    operators on non-Windows shells or systems without PowerShell. Now the
    failure mode is loud so the operator knows font warnings are absent
    because detection failed, not because everything is installed.
    """
    try:
        cmd = (
            "Add-Type -AssemblyName System.Drawing; "
            "[System.Drawing.Text.InstalledFontCollection]::new().Families "
            "| ForEach-Object { $_.Name }"
        )
        result = subprocess.run(
            ["powershell.exe", "-NoProfile", "-NonInteractive", "-Command", cmd],
            capture_output=True, text=True, timeout=10,
        )
        if result.returncode != 0:
            sys.stderr.write(
                f"  WARN: PowerShell font enumeration exited {result.returncode}: "
                f"{(result.stderr or '').strip()[:200]}\n"
                f"  Client-font mismatch warnings will not appear this run.\n"
            )
            return set()
        return {line.strip() for line in result.stdout.splitlines() if line.strip()}
    except FileNotFoundError:
        sys.stderr.write(
            "  WARN: powershell.exe not found on PATH. Font detection skipped; "
            "client-font mismatch warnings will not appear this run.\n"
        )
        return set()
    except subprocess.TimeoutExpired:
        sys.stderr.write(
            "  WARN: PowerShell font enumeration timed out after 10s. "
            "Client-font mismatch warnings will not appear this run.\n"
        )
        return set()
    except Exception as exc:
        sys.stderr.write(
            f"  WARN: font enumeration failed ({type(exc).__name__}: {exc}). "
            f"Client-font mismatch warnings will not appear this run.\n"
        )
        return set()


def detect_missing_client_fonts(theme) -> list:
    installed = detect_installed_fonts()
    if not installed:
        return []
    inst_lower = {f.lower() for f in installed}
    candidates = []
    for f in (theme.major_font, theme.minor_font):
        f = (f or "").strip()
        if not f:
            continue
        f_lower = f.lower()
        if any(f_lower in i or i in f_lower for i in inst_lower):
            continue
        candidates.append(f)
    seen = set()
    missing = []
    for f in candidates:
        if f not in seen:
            seen.add(f)
            missing.append(f)
    return missing


def run_option_qc(themed_pptx_path: Path, png_path: Path, expected_palette: set,
                  page_type: str = "") -> dict:
    """Run 7 deterministic checks on a themed option.

    page_type (from _meta.json::slides[].page_type) lets the check exempt
    surfaces from rules that don't apply to them — most importantly, cover
    slides legitimately lack a page-number per consulting-deck convention,
    so footer_present passes by-design on `page_type == "cover"`.
    """
    checks: list = []
    png_ok = False
    png_detail = ""
    if not png_path.exists():
        png_detail = "PNG not found"
    else:
        # Floor was 50KB in v0; covers + sparse hero slides routinely render
        # at 20-35KB even with correct content. P6.35 (2026-05-26) lowered
        # to 12KB which is roughly the size of a near-blank PNG and lets
        # legitimate sparse slides through while still catching the "blank
        # canvas / render failed silently" failure mode.
        try:
            sz = png_path.stat().st_size
            if sz <= 12 * 1024:
                png_detail = f"PNG too small ({sz} bytes; floor 12KB)"
            else:
                png_ok = True
                png_detail = f"{sz} bytes"
        except Exception as e:
            png_detail = f"stat failed: {e}"
    checks.append({"check": "png_render_ok", "pass": png_ok, "severity": "block", "detail": png_detail})

    hexes = _hex_codes_in_pptx(themed_pptx_path)
    orphans = sorted(h for h in hexes if h not in expected_palette)
    if not orphans:
        pal_detail = f"all {len(hexes)} hex codes in palette"
        pal_ok = True
    else:
        sample = ", ".join(orphans[:8])
        more = "" if len(orphans) <= 8 else f" (+{len(orphans) - 8} more)"
        pal_detail = f"{len(orphans)} off-palette: {sample}{more}"
        pal_ok = False
    checks.append({"check": "palette_compliance", "pass": pal_ok, "severity": "warn", "detail": pal_detail})

    title_ok = False
    title_detail = ""
    footer_ok = False
    footer_detail = ""
    body_ok = True
    body_offenders: list = []
    soft_floor_suppressed = 0  # runs in non-body roles between 8.0 and 10.5pt
    leak_ok = True
    leak_offenders: list = []
    shape_count = 0
    shape_count_ok = False
    shape_count_detail = ""
    open_err = ""
    try:
        prs = Presentation(str(themed_pptx_path))
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        shape_count = len(shapes)
        title_threshold = Pt(28)
        for shape in shapes:
            name = ""
            try:
                name = (shape.name or "").strip()
            except Exception:
                pass
            name_lower = name.lower()
            if not title_ok and name_lower.startswith("title"):
                title_ok = True
                title_detail = f"shape name '{name}' matches title*"
            if not footer_ok and name_lower == "page-number":
                footer_ok = True
                footer_detail = f"shape name '{name}' matches page-number"
            if not shape.has_text_frame:
                continue
            is_footnote_like = _is_footnote_like_name(name)
            tf = shape.text_frame
            try:
                shape_text = (tf.text or "").strip()
            except Exception:
                shape_text = ""
            if not footer_ok and shape_text and _FOOTER_NUM_RE.match(shape_text):
                footer_ok = True
                footer_detail = f"shape text '{shape_text}' is a page-number"
            for para in tf.paragraphs:
                for run in para.runs:
                    text = run.text or ""
                    sz = None
                    try:
                        sz = run.font.size
                    except Exception:
                        sz = None
                    if not title_ok and sz is not None and sz > title_threshold:
                        title_ok = True
                        title_detail = f"run > 28pt found (name='{name}', size={sz.pt:.1f}pt)"
                    if not is_footnote_like and sz is not None:
                        # Per-role floor: body-role shapes get the 10.5pt soft
                        # floor; everything else gets the 8.0pt hard floor.
                        if _is_body_role_name(name):
                            floor = _BODY_FONT_FLOOR_PT
                        else:
                            floor = _FONT_HARD_FLOOR_PT
                            if sz.pt < _BODY_FONT_FLOOR_PT:
                                # Not flagged, but counted so we don't lose
                                # the audit signal entirely.
                                soft_floor_suppressed += 1
                        if sz.pt < floor:
                            body_ok = False
                            body_offenders.append(
                                f"{name or '<unnamed>'} @ {sz.pt:.1f}pt"
                                f" (floor {floor:.1f}): {text[:40]!r}"
                            )
                    if text:
                        if not any(allowed in text for allowed in _PLACEHOLDER_ALLOWED):
                            for pat in _PLACEHOLDER_PATTERNS:
                                if pat in text:
                                    leak_ok = False
                                    leak_offenders.append(f"'{pat}' in {name or '<unnamed>'}: {text[:60]!r}")
                                    break
        if 5 <= shape_count <= 80:
            shape_count_ok = True
            shape_count_detail = f"{shape_count} shapes"
        else:
            shape_count_detail = f"{shape_count} shapes (expected 5..80)"
    except Exception as e:
        open_err = f"{type(e).__name__}: {e}"

    if open_err:
        title_detail = title_detail or f"could not open: {open_err}"
        footer_detail = footer_detail or f"could not open: {open_err}"
        shape_count_detail = shape_count_detail or f"could not open: {open_err}"
    if not title_detail:
        title_detail = "no title* shape and no run > 28pt"
    # Cover slides legitimately don't carry a page-number per consulting
    # convention (the title page is page 0 / unnumbered). Treat as pass.
    if (page_type or "").strip().lower() == "cover":
        footer_ok = True
        footer_detail = "cover slide — page-number not required by convention"
    if not footer_detail:
        footer_detail = "no shape named page-number and no bare-integer text"

    checks.append({"check": "title_present", "pass": title_ok, "severity": "warn", "detail": title_detail})
    checks.append({"check": "footer_present", "pass": footer_ok, "severity": "warn", "detail": footer_detail})
    if body_ok:
        body_detail = (
            f"body-role >= 10.5pt; other roles >= 8.0pt"
            f" (excluding footnote/source/page-number)"
        )
        if soft_floor_suppressed:
            body_detail += (
                f"; {soft_floor_suppressed} non-body-role runs in 8.0-10.5pt"
                f" band (eyebrows / labels / legends — by design)"
            )
    else:
        body_detail = (
            f"{len(body_offenders)} sub-floor runs: "
            + "; ".join(body_offenders[:4])
            + ("" if len(body_offenders) <= 4 else f" (+{len(body_offenders) - 4} more)")
        )
    checks.append({
        "check": "body_font_floor", "pass": body_ok, "severity": "warn",
        "detail": body_detail,
    })
    checks.append({
        "check": "placeholder_leak", "pass": leak_ok, "severity": "block",
        "detail": ("no placeholder strings found" if leak_ok
                   else f"{len(leak_offenders)} leak(s): " + "; ".join(leak_offenders[:3])
                        + ("" if len(leak_offenders) <= 3 else f" (+{len(leak_offenders) - 3} more)")),
    })
    checks.append({"check": "shape_count_sanity", "pass": shape_count_ok, "severity": "warn", "detail": shape_count_detail})

    # Defect 7 fix (CDIO QBR, 2026-06-15): worker geometry guards.
    # These three checks catch failure modes the original 7 checks missed,
    # surfaced by the CDIO QBR build: zero-fill card backgrounds (worker
    # forgot the fill_color arg), text-clip risk (text too long for box),
    # and overlap-with-chrome (worker placed content at canonical chrome y).
    zero_fill_ok = True
    zero_fill_offenders: list = []
    clip_risk_ok = True
    clip_risk_offenders: list = []
    chrome_overlap_ok = True
    chrome_overlap_offenders: list = []
    try:
        for shape in shapes:
            try:
                name = (shape.name or "").strip()
            except Exception:
                continue
            name_lower = name.lower()
            # (a) Zero-fill card / badge detection. Workers conventionally name
            # background rects card-*, badge-*, pill-*, bg-*. If those AUTO_SHAPE
            # shapes ship with `fill.background()` (no fill), the slide
            # renders with invisible cards. Catches the slide-111 bug from CDIO.
            if (name_lower.startswith(("card-", "badge-", "pill-", "bg-"))
                    or "card_bg" in name_lower or "badge_bg" in name_lower):
                try:
                    fill_type = shape.fill.type
                    # fill.type values: None (no fill), 1 (solid), etc.
                    # python-pptx returns None when fill.background() was called.
                    if fill_type is None or int(fill_type or 0) == 5:  # 5 = BACKGROUND
                        zero_fill_ok = False
                        zero_fill_offenders.append(name)
                except Exception:
                    pass
            # (b) Text-clip risk heuristic. For text shapes with a known font
            # size, check whether the literal text is too long for the box.
            # Rough: monospace char width ~= 0.6 × font_size_pt; PowerPoint
            # auto-wraps but a long single token (e.g., "04" rendering as "0"
            # at slide 51) clips at single-character boundary.
            if shape.has_text_frame and shape.width and shape.height:
                try:
                    text = (shape.text_frame.text or "").strip()
                    if text:
                        # Find largest font size in the shape.
                        max_pt = 0.0
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                try:
                                    if r.font.size:
                                        pt_val = r.font.size.pt
                                        if pt_val > max_pt:
                                            max_pt = pt_val
                                except Exception:
                                    pass
                        if max_pt > 0:
                            w_emu = int(shape.width)
                            w_px = w_emu / 9525.0  # EMU → px @96dpi
                            # Heuristic: estimated rendered width = chars × 0.6 × pt → px
                            est_width_px = len(text) * 0.6 * max_pt * (96.0 / 72.0)
                            # Flag when estimate exceeds 1.5× box width AND box
                            # is narrow (< 200 px) — single-char-clip pattern.
                            if est_width_px > 1.5 * w_px and w_px < 200:
                                clip_risk_ok = False
                                clip_risk_offenders.append(
                                    f"{name!r} ({len(text)} chars @ {max_pt:.0f}pt, box={w_px:.0f}px)"
                                )
                except Exception:
                    pass
            # (c) Chrome-zone overlap: worker shapes whose top-y sits inside
            # the top invariant zone (< 60 px) or bottom invariant zone
            # (> 660 px). The invariant zones are reserved for chrome
            # (sources, footnotes, page numbers); content there collides
            # with the chrome graft. Per feedback_invariant_zone_chrome.
            try:
                top_emu = int(shape.top or 0)
                top_px = top_emu / 9525.0
                if shape.has_text_frame and (shape.text_frame.text or "").strip():
                    if 0 < top_px < 40 and not name_lower.startswith(
                        ("page-number", "footnote", "source", "header", "chrome")
                    ):
                        chrome_overlap_ok = False
                        chrome_overlap_offenders.append(f"{name!r}@y={top_px:.0f}px")
            except Exception:
                pass
    except Exception:
        # Don't fail the QC on geometry-walk error; record so operator sees it.
        zero_fill_ok = False
        zero_fill_offenders.append("(geometry-walk failed)")

    checks.append({
        "check": "zero_fill_card_bg", "pass": zero_fill_ok, "severity": "warn",
        "detail": ("all card/badge backgrounds have fill" if zero_fill_ok
                   else f"{len(zero_fill_offenders)} no-fill bg shape(s): "
                        + ", ".join(zero_fill_offenders[:3])),
    })
    checks.append({
        "check": "text_clip_risk", "pass": clip_risk_ok, "severity": "warn",
        "detail": ("no obvious clip risk" if clip_risk_ok
                   else f"{len(clip_risk_offenders)} shape(s) likely clipped: "
                        + "; ".join(clip_risk_offenders[:3])),
    })
    checks.append({
        "check": "chrome_zone_overlap", "pass": chrome_overlap_ok, "severity": "warn",
        "detail": ("content respects top invariant zone" if chrome_overlap_ok
                   else f"{len(chrome_overlap_offenders)} shape(s) in chrome zone: "
                        + ", ".join(chrome_overlap_offenders[:3])),
    })

    summary = {"pass": 0, "warn": 0, "block": 0}
    for c in checks:
        if c["pass"]:
            summary["pass"] += 1
        elif c["severity"] == "block":
            summary["block"] += 1
        else:
            summary["warn"] += 1
    return {"checks": checks, "summary": summary}


# ---------------------------------------------------------------------------
# M6 (Pattern B QC, 2026-06-17) — R4.1 through R4.8 checks per Spec 6
# ---------------------------------------------------------------------------
def _check_editability_structural(pptx_path: Path) -> dict:
    """Verify the translator's saved .pptx contains structurally editable text.

    Opens the file with python-pptx and walks every shape on every slide:
      - editable_text_shapes: shape has a text_frame with non-empty text
        and non-zero bounding box
      - picture_shapes: bitmap; counted as a flatten candidate but not a fail
        on its own (legitimate when a slide carries a real image)
      - zero_size_text: shape has text but zero width/height (unselectable
        in PowerPoint — equivalent to a flatten)

    Passes when at least one editable text shape exists AND no zero-size
    text frames are present. Used by R4.7 instead of the translator's
    SSIM-based self-report, which conflates editability with pre-graft
    visual fidelity (chrome zones are empty pre-graft by design, so
    title-zone SSIM tanks and a structurally-fine slide gets blocked).
    """
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        return {
            "pass": False,
            "detail": f"could not open native .pptx for introspection: {exc}",
            "counts": {},
        }

    editable_text_shapes = 0
    picture_shapes = 0
    zero_size_text = 0

    for slide in prs.slides:
        for shp in slide.shapes:
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_shapes += 1
                    continue
                if not getattr(shp, "has_text_frame", False):
                    continue
                tf = shp.text_frame
                has_text = any(p.text.strip() for p in tf.paragraphs)
                if not has_text:
                    continue
                if (shp.width or 0) <= 0 or (shp.height or 0) <= 0:
                    zero_size_text += 1
                    continue
                editable_text_shapes += 1
            except Exception:
                continue

    counts = {
        "editable_text_shapes": editable_text_shapes,
        "picture_shapes": picture_shapes,
        "zero_size_text": zero_size_text,
    }
    passed = editable_text_shapes >= 1 and zero_size_text == 0
    if passed:
        detail = f"{editable_text_shapes} editable text shapes, no flatten signals"
    elif editable_text_shapes == 0:
        detail = (
            f"no editable text shapes found "
            f"({picture_shapes} pictures, {zero_size_text} zero-size text)"
        )
    else:
        detail = (
            f"{editable_text_shapes} editable text shapes but "
            f"{zero_size_text} zero-size text frame(s) — unselectable"
        )
    return {"pass": passed, "detail": detail, "counts": counts}


def _check_r4_rules_for_pattern_b(st: "OptionStatus") -> list[dict]:
    """Run R4.1-R4.8 QC checks on a Pattern B translator output.

    Per Spec 6 (_decisions/pattern-b/spec-6-qc-rules-R4.md), Pattern B
    introduces a parallel rule family R4 that operates on the
    translator's output and report:

      R4.1 Strikethrough/underline on body text       (Critical)
      R4.2 Opacity / transparency on text             (Major)
      R4.3 Font-weight mismatch                       (Major)
      R4.4 Gradient conversion to solid               (Major)
      R4.5 CSS filters dropped                        (Advisory)
      R4.6 Icons rendering cleanly                    (Major)
      R4.7 Editability verified at build time         (Critical)
      R4.8 Text decorations audit (catch-all)         (Critical)

    Each return entry has the same shape as run_option_qc's `checks`
    items: {check, severity, pass, detail}. Severity values use the
    legacy run_option_qc vocabulary: "block" = Critical hard-fail,
    "warn" = Major surfaceable in REVIEW.html, "info" = Advisory.

    Returns an empty list when:
      - st.classification != "pattern_b_translated" (Pattern C / legacy),
      - the translator's translation_report.json is missing or malformed.
    """
    if st.classification != "pattern_b_translated":
        return []

    # Translation report lives next to the .py file; sibling convention.
    report_path = st.py_path.parent / f"option_{st.letter}_translation_report.json"
    report: dict = {}
    if report_path.exists():
        try:
            report = json.loads(report_path.read_text(encoding="utf-8"))
        except Exception:
            report = {}

    # Defensive accessors — every field below is optional in the report.
    # A malformed translator report could carry unexpected types (e.g.,
    # `"warnings": "error string"`); type-guard so finalize never crashes
    # on a bad report. Missing or wrong-type fields degrade to defaults.
    kill_list = report.get("css_kill_list_applied", {})
    if not isinstance(kill_list, dict):
        kill_list = {}
    warnings = report.get("warnings", [])
    if not isinstance(warnings, list):
        warnings = []
    warning_codes = {
        (w.get("code") or "").upper()
        for w in warnings
        if isinstance(w, dict)
    }

    # Script-source scan for R4.1 (defense-in-depth alongside translator
    # report). Pattern B output should never set strikethrough/underline
    # on body runs unless the brief explicitly authorizes — and the brief
    # gating happens at the translator, so any survivor here is a bug.
    script_has_strikethrough = False
    script_has_underline = False
    try:
        script_text = st.py_path.read_text(encoding="utf-8")
        script_has_strikethrough = "font.strikethrough = True" in script_text
        script_has_underline = "font.underline = True" in script_text
    except Exception:
        pass

    checks: list[dict] = []

    # R4.1 Strikethrough / underline on body text (Critical)
    text_decorations_stripped = int(kill_list.get("text_decorations_stripped", 0) or 0)
    r41_fail = (
        script_has_strikethrough
        or script_has_underline
        or "R4.1" in warning_codes
    )
    checks.append({
        "check": "R4.1 strikethrough/underline on body text",
        "severity": "block",
        "pass": not r41_fail,
        "detail": (
            f"strikethrough={script_has_strikethrough}, underline={script_has_underline}, "
            f"stripped_at_translate={text_decorations_stripped}"
        ),
    })

    # R4.2 Opacity / transparency on text (Major)
    opacity_normalized = int(kill_list.get("opacity_on_text_normalized", 0) or 0)
    r42_fail = opacity_normalized > 0 or "R4.2" in warning_codes
    checks.append({
        "check": "R4.2 opacity on text",
        "severity": "warn",
        "pass": not r42_fail,
        "detail": f"opacity_normalized_count={opacity_normalized}",
    })

    # R4.3 Font-weight mismatch (Major) — translator warns by emitting a
    # warning whose code contains FONT_WEIGHT.
    r43_fail = any("FONT_WEIGHT" in c for c in warning_codes)
    checks.append({
        "check": "R4.3 font-weight mismatch",
        "severity": "warn",
        "pass": not r43_fail,
        "detail": "translator reported font-weight fallback" if r43_fail else "no font-weight warnings",
    })

    # R4.4 Gradient conversion to solid (Major)
    gradients = int(kill_list.get("gradients_flattened", 0) or 0)
    r44_fail = gradients > 0
    checks.append({
        "check": "R4.4 gradient conversion",
        "severity": "warn",
        "pass": not r44_fail,
        "detail": f"gradients_flattened={gradients}",
    })

    # R4.5 CSS filters / shadows dropped (Advisory)
    filters_dropped = int(kill_list.get("filters_dropped", 0) or 0)
    shadows_stripped = int(kill_list.get("shadows_stripped", 0) or 0)
    r45_dropped = filters_dropped + shadows_stripped
    checks.append({
        "check": "R4.5 CSS filters / shadows dropped",
        "severity": "info",
        "pass": r45_dropped == 0,
        "detail": f"filters={filters_dropped}, shadows={shadows_stripped}",
    })

    # R4.6 Icons rendering (Major) — v0 trusts the translator's warning array.
    r46_fail = any("ICON" in c for c in warning_codes)
    checks.append({
        "check": "R4.6 icons rendering",
        "severity": "warn",
        "pass": not r46_fail,
        "detail": "translator flagged icon issue" if r46_fail else "no icon warnings",
    })

    # R4.7 Editability verified (Critical, load-bearing).
    # Structural introspection of the translator's saved .pptx — replaces the
    # earlier SSIM-based translator self-report (2026-06-18 cutover validation
    # found that self-report conflated editability with pre-graft visual
    # fidelity, producing 11/15 false-positive BLOCKs on the OTC deck).
    # Passes when the native .pptx contains at least one editable text shape
    # and no zero-size text frames. The translator's `editability_self_check`
    # subfield is now ignored — its prose is advisory only.
    if st.pptx_path.exists():
        struct = _check_editability_structural(st.pptx_path)
        r47_pass = struct["pass"]
        r47_detail = struct["detail"]
    else:
        r47_pass = False
        r47_detail = f"native .pptx missing at {st.pptx_path.name}"
    checks.append({
        "check": "R4.7 editability verified at build time",
        "severity": "block",
        "pass": r47_pass,
        "detail": r47_detail,
    })

    # R4.8 Text decorations audit catch-all (Critical) — any warning whose
    # code mentions a decoration we didn't explicitly bucket above.
    decoration_keywords = ("TEXT_DECORATION", "UNDERLINE", "STRIKE", "OVERLINE")
    r48_fail = any(
        any(k in c for k in decoration_keywords)
        for c in warning_codes
    )
    checks.append({
        "check": "R4.8 text decorations audit (catch-all)",
        "severity": "block",
        "pass": not r48_fail,
        "detail": "catch-all decoration warning from translator" if r48_fail else "no decoration warnings",
    })

    return checks


# ---------------------------------------------------------------------------
# Data classes — v1 base + v2 status enum
# ---------------------------------------------------------------------------
@dataclass
class OptionStatus:
    slide_n: int
    letter: str
    py_path: Path
    pptx_path: Path
    raw_archive_path: Path
    themed_pptx_path: Path
    themed_png_path: Path
    # v2 additions
    # M7 (Mermaid retirement, 2026-06-17): `fallback_mermaid` is no longer
    # produced by the classifier. `mmd_path` / `mermaid_png_path` removed.
    classification: str = "native"           # one of: native, pattern_b_translated, skeleton_rejected, missing
    classification_reason: str = ""
    # v1 fields
    built: Optional[bool] = None
    themed: Optional[bool] = None
    rendered: Optional[bool] = None
    n_shapes: int = 0
    n_subs: int = 0
    error: str = ""
    # v0.3 dark-variant collision tracking — populated by graft_and_theme
    # when slide.variant == "dark" and content fill/text colors collide with
    # brand.dark_bg_hex. Empty list means clean; non-empty triggers hard fail
    # in main() after all options are processed.
    dark_collisions: list = field(default_factory=list)
    # M6 (Pattern B QC, 2026-06-17): R4.1-R4.8 check results from
    # _check_r4_rules_for_pattern_b. Empty for legacy / Pattern C builds;
    # populated for Pattern B picks after graft+theme so QC surfaces in
    # REVIEW.html via the .qc.json file.
    r4_checks: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# Discovery — v1 base + v2 classification at scan time
# ---------------------------------------------------------------------------
def discover_options(out_dir: Path, expected: Optional[set] = None) -> list:
    """Discover per-slide option scripts produced by the parallel worker fanout.

    expected: optional set of (slide_n, letter) tuples the parent session
        promised in _meta.json. Any pair in `expected` that has no
        corresponding option_X.py file on disk is synthesized as an
        OptionStatus with classification="missing" so finalize surfaces the
        gap in RESULT.md instead of silently skipping the slide. This is the
        committee P4 fix (2026-05-26): treat a missing worker output the
        same as a SKELETON_REJECTED, ride the existing rejected-branch rails,
        no retries, no hard-fail. The operator sees the gap and re-dispatches
        the worker if they want a fuller deck.
    """
    statuses: list = []
    found_pairs: set = set()
    # M5: prefer Pattern B translator output (`option_X_native.py`) when both it
    # and a sibling `option_X.py` are present for the same slide+letter.
    # Build a map of which (slide_n, letter) pairs have a translator output so
    # we can skip the worker's `.py` if the translator already converted that
    # option. In production, Pattern B slides emit only HTML (no `option_X.py`
    # ever exists), so this guard mostly only matters for mixed-mode test runs.
    translator_pairs: set = set()
    for npy in out_dir.glob("slide_*/option_*_native.py"):
        try:
            slide_n = int(npy.parent.name.split("_")[1])
            # filename shape: option_<letter>_native.py -> letter = stem.split('_')[1]
            letter = npy.stem.split("_")[1]
        except (IndexError, ValueError):
            continue
        translator_pairs.add((slide_n, letter))

    # First pass: legacy `option_X.py` worker output (Pattern C / native).
    for py in sorted(out_dir.glob("slide_*/option_*.py")):
        # Skip translator-output files; they get their own pass below.
        if py.stem.endswith("_native"):
            continue
        slide_n = int(py.parent.name.split("_")[1])
        letter = py.stem.split("_")[1]
        # Translator output for this slide+letter exists → the translator's
        # script supersedes the worker's. Skip the worker's `.py`; the
        # translator's `_native.py` is appended in the second pass.
        if (slide_n, letter) in translator_pairs:
            continue
        slide_dir = py.parent
        classification, reason = _classify_option(py)
        statuses.append(OptionStatus(
            slide_n=slide_n,
            letter=letter,
            py_path=py,
            pptx_path=py.with_suffix(".pptx"),
            raw_archive_path=slide_dir / _p.RAW_DIR / _p.option_pptx_name(letter),
            themed_pptx_path=slide_dir / _p.option_pptx_name(letter),
            themed_png_path=slide_dir / _p.option_png_name(letter),
            classification=classification,
            classification_reason=reason,
        ))
        found_pairs.add((slide_n, letter))

    # M5: second pass — Pattern B translator output. Each `option_X_native.py`
    # gets an OptionStatus whose `py_path` is the translator script. The
    # downstream build_pptx + graft_and_theme paths execute it via the normal
    # native code path; the `__template_fields__` header is parsed for
    # placeholder population in `_apply_body_canonical_finishing`.
    for npy in sorted(out_dir.glob("slide_*/option_*_native.py")):
        try:
            slide_n = int(npy.parent.name.split("_")[1])
            letter = npy.stem.split("_")[1]
        except (IndexError, ValueError):
            continue
        slide_dir = npy.parent
        classification, reason = _classify_option(npy)
        statuses.append(OptionStatus(
            slide_n=slide_n,
            letter=letter,
            py_path=npy,
            pptx_path=npy.with_suffix(".pptx"),
            raw_archive_path=slide_dir / _p.RAW_DIR / _p.option_pptx_name(letter),
            themed_pptx_path=slide_dir / _p.option_pptx_name(letter),
            themed_png_path=slide_dir / _p.option_png_name(letter),
            classification=classification,
            classification_reason=reason,
        ))
        found_pairs.add((slide_n, letter))

    # P4 — synthesize missing options for every (slide_n, letter) the parent
    # session promised but no file landed for. Inserted in slide+letter order.
    if expected:
        for slide_n, letter in sorted(expected - found_pairs):
            slide_dir = out_dir / f"slide_{slide_n:02d}"
            statuses.append(OptionStatus(
                slide_n=slide_n,
                letter=letter,
                py_path=slide_dir / _p.option_py_name(letter),
                pptx_path=slide_dir / _p.option_pptx_name(letter),
                raw_archive_path=slide_dir / _p.RAW_DIR / _p.option_pptx_name(letter),
                themed_pptx_path=slide_dir / _p.option_pptx_name(letter),
                themed_png_path=slide_dir / _p.option_png_name(letter),
                classification="missing",
                classification_reason="worker did not produce option script",
            ))
        statuses.sort(key=lambda s: (s.slide_n, s.letter))
    return statuses


def stash_raw(st: OptionStatus) -> None:
    """Move the worker-produced PPTX to the raw archive path.

    Loud-failure contract (T2.13, audit 2026-05-26): if the rename fails —
    most commonly a Windows file lock from PowerPoint still holding the
    worker output open, or antivirus mid-scan — we MUST surface that.
    Silent return previously left `st.pptx_path` pointing at the original
    file, so `graft_and_theme` would then mutate the un-archived original
    and the raw pre-theme copy would never exist. Now we set
    `st.error = "stash_raw failed: ..."` so the failure shows up in
    RESULT.md per-option and the operator gets a real signal instead of
    a phantom "themed" status with a missing archive.
    """
    if not st.pptx_path.exists():
        return
    if st.pptx_path.resolve() == st.raw_archive_path.resolve():
        return
    st.raw_archive_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        st.pptx_path.replace(st.raw_archive_path)
    except OSError as exc:
        msg = f"stash_raw failed: {type(exc).__name__}: {exc}"
        sys.stderr.write(
            f"  WARN slide_{st.slide_n:02d} option_{st.letter}: {msg}\n"
            f"       (likely Windows file lock — close PowerPoint and re-run, "
            f"or pause antivirus on the build dir)\n"
        )
        # Preserve any prior error context; otherwise stamp this one so it
        # surfaces in RESULT.md downstream readers.
        if not st.error:
            st.error = msg
        return
    st.pptx_path = st.raw_archive_path


# ---------------------------------------------------------------------------
# v0.2 P1.3 — chrome.yml loader + sha8 freshness check
# ---------------------------------------------------------------------------

def _sha256_of_file(path: Path) -> str:
    """Hex-encoded sha256 of `path` content."""
    h = hashlib.sha256()
    with open(str(path), "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _load_and_verify_chrome(template_path: Path) -> ChromeSpec:
    """Load chrome.yml from the template's sidecar subfolder
    (`<template-stem>/chrome.yml`, v0.4+ layout) and hard-fail on SHA mismatch.

    Raises:
      ChromeSidecarMissingError if chrome.yml is absent (operator must run
        register_template.py propose -> commit).
      RuntimeError if the stored source_template_sha8 doesn't match the
        current template bytes (template mutated since registration —
        re-register).
    """
    chrome_yml_path = _p.chrome_yml(template_path)
    spec = load_chrome_yml(chrome_yml_path)
    current_sha8 = _sha256_of_file(template_path)[:8]
    if spec.source_template_sha8 != current_sha8:
        raise RuntimeError(
            f"chrome.yml source_template_sha8 mismatch: "
            f"stored={spec.source_template_sha8!r} vs current={current_sha8!r}. "
            f"Template {template_path.name} changed since registration. "
            f"Re-register: register_template.py propose -> commit."
        )
    return spec


def _pick_default_layout_name(spec: ChromeSpec) -> str:
    """Return a sensible default layout name when the meta lacks one.

    Priority:
      1. First body-canonical layout (with light bg if any).
      2. First body-canonical layout (any bg).
      3. First layout in the spec (alphabetical fallback).

    P1.3 ships before P1.4 wires per-slide layout into _meta.json. Until P1.4
    lands, every slide gets the deck-wide default returned here. After P1.4,
    this function is only consulted for slides that lack a layout field AND
    have no default_layout in front-matter.
    """
    light_canonical = [n for n, lc in spec.layouts.items()
                       if lc.layout_class == "body-canonical"
                       and lc.background == "light"]
    if light_canonical:
        return sorted(light_canonical)[0]
    canonical = [n for n, lc in spec.layouts.items()
                 if lc.layout_class == "body-canonical"]
    if canonical:
        return sorted(canonical)[0]
    if spec.layouts:
        return sorted(spec.layouts)[0]
    return ""

# Body-zone embed coordinates — matches the conventions baked into prompt.md
# and reference/fallback.md. 1240x540 PNG centered horizontally at y=110.
# Slide canvas: 1280x720 EMUs; we work in pixels and convert.
PX_TO_EMU = 9525  # python-pptx convention: 1 px = 9525 EMU at 96 DPI



# ---------------------------------------------------------------------------
# Step 1: build pptx — v1 native path + Pattern B translator path (M5)
# ---------------------------------------------------------------------------
# M7 (Mermaid retirement, Decision 6, 2026-06-17): the Mermaid fallback path
# was removed here. `_resolve_mermaid_theme`, `_render_mermaid_png`, and
# `_assemble_fallback_pptx` were deleted. Pattern B HTML→PNG (per
# `_decisions/pattern-b/SPEC.md`) supersedes Mermaid for curved-container
# diagrams. Stale `option_X.py` scripts carrying `# FALLBACK_MERMAID:` on
# line 1 now fall through to the `native` classifier and fail loudly at
# execution time — the operator re-builds.
# ---------------------------------------------------------------------------
def _try_load_brief_title_from_prompt(st: OptionStatus) -> str:
    """Extract the slide title from <slide_dir>/_prompt.md so the fallback
    assembler can stamp the right title block. Best-effort; returns '' if not
    found."""
    prompt_path = st.py_path.parent / _p.PROMPT_MD
    if not prompt_path.exists():
        return ""
    try:
        text = prompt_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return ""
    m = re.search(r"\*\*Slide title:\*\*\s*(.+?)\s*$", text, re.MULTILINE)
    return m.group(1).strip() if m else ""


def _run_subprocess(st: OptionStatus, env: dict[str, str] | None = None):
    try:
        result = subprocess.run(
            [sys.executable, str(st.py_path.resolve())],
            capture_output=True, text=True, timeout=120,
            cwd=str(st.py_path.parent.resolve()),
            env=env,
        )
        if result.returncode != 0:
            return False, (result.stderr or result.stdout or "non-zero exit")[-400:]
        if not st.pptx_path.exists():
            return False, "script ran but no .pptx produced"
        return True, ""
    except subprocess.TimeoutExpired:
        return False, "timeout (120s)"
    except Exception as e:
        return False, f"subprocess: {type(e).__name__}: {e}"


def _run_runpy(st: OptionStatus):
    saved_cwd = os.getcwd()
    saved_argv = sys.argv[:]
    try:
        os.chdir(str(st.py_path.parent))
        sys.argv = [str(st.py_path)]
        runpy.run_path(str(st.py_path), run_name="__main__")
        if not st.pptx_path.exists():
            return False, "runpy: script ran but no .pptx produced"
        return True, ""
    except SystemExit as e:
        if e.code in (0, None) and st.pptx_path.exists():
            return True, ""
        return False, f"runpy SystemExit({e.code})"
    except Exception as e:
        tb_last = traceback.format_exc().strip().splitlines()[-1] if traceback else ""
        return False, f"runpy: {type(e).__name__}: {e} | {tb_last}"
    finally:
        os.chdir(saved_cwd)
        sys.argv = saved_argv


def _option_subprocess_env(chrome_yml_path: Path, layout_name: str) -> dict[str, str]:
    """Env for subprocess option_X.py runs — supplies chrome to twins.helpers."""
    env = os.environ.copy()
    env["SLIDE_LAB_CHROME_YML"] = str(chrome_yml_path)
    env["SLIDE_LAB_LAYOUT_NAME"] = layout_name
    return env


def build_pptx(st: OptionStatus,
               chrome_yml_path: Path | None = None,
               layout_chrome=None,
               layout_name: str = "",
               prefer_runpy: bool = False) -> None:
    """v2: branch on classification before invoking the native path.

    skeleton_rejected     -> mark not-built; rejection surfaces in REVIEW.html
    pattern_b_translated  -> Pattern B translator output (M5); runs via the
                              native path; finalize merges __template_fields__
                              into the placeholder population step.
    native                -> v1's subprocess/runpy execution

    M7 (Mermaid retirement, Decision 6, 2026-06-17): the `fallback_mermaid`
    branch was removed. `mermaid_theme` parameter was retired (it was only
    consumed by the Mermaid fallback's `_render_mermaid_png` call).
    """
    if st.pptx_path.exists():
        try:
            st.pptx_path.unlink()
        except Exception:
            pass

    # v2 skeleton-rejected branch — no build attempt
    if st.classification == "skeleton_rejected":
        st.built = False
        st.error = f"SKELETON_REJECTED: {st.classification_reason}"
        return

    # v2 missing branch — option script not found
    if st.classification == "missing":
        st.built = False
        st.error = st.classification_reason
        return

    # Native path — subprocess gets chrome via env vars; runpy gets it via
    # the in-process module attribute (option script imports twins.helpers
    # in this same process).
    subprocess_env = None
    if chrome_yml_path is not None and layout_name:
        subprocess_env = _option_subprocess_env(chrome_yml_path, layout_name)
    if not prefer_runpy:
        ok, err = _run_subprocess(st, env=subprocess_env)
        if ok:
            st.built = True
            return
        sandbox_signals = ("WinError 5", "Access is denied", "PermissionError")
        if not any(sig in err for sig in sandbox_signals):
            st.built = False
            st.error = err
            return

    _twins_helpers.set_active_chrome(layout_chrome)
    try:
        ok, err = _run_runpy(st)
    finally:
        _twins_helpers.set_active_chrome(None)
    if ok:
        st.built = True
    else:
        st.built = False
        st.error = err


# ---------------------------------------------------------------------------
# Step 2: graft + theme remap — verbatim from v1
# ---------------------------------------------------------------------------
def _apply_body_canonical_finishing(new_slide, prs, layout_chrome,
                                     src_slide, slide_n: int,
                                     fallback_title: str = "",
                                     fallback_subtitle: str = "",
                                     dark_variant: bool = False,
                                     dark_bg_hex: str = "",
                                     brand_ttf_path: str = "",
                                     template_fields_override: Optional[dict] = None) -> None:
    """Body-canonical chrome population + title/subtitle finishing.

    M5 (Pattern B production wiring, 2026-06-17): `template_fields_override`
    is the dict parsed from a Pattern B translator's `__template_fields__`
    header. When non-None and a `title` / `subtitle` key is present, that
    value takes priority over `fallback_title` / `fallback_subtitle` (which
    come from the brief). For Pattern C / legacy builds, this kwarg is
    always None and behavior is byte-identical to pre-M5.
    """
    """Body-canonical post-graft.

    Light variant (default): populate inherited title/footer/page-number
    placeholders. Layout's decorative chrome inherits as designed.

    Dark variant (per-slide ``**Variant:** dark`` in brief): insert a
    full-bleed overlay with the brand's ``dark_bg_hex``, draw the title
    fresh as a slide-level white text box at the layout's title-placeholder
    position. The inherited title placeholder is hidden under the overlay,
    so we draw on top instead of populating it.

    Title text priority: source slide's 'title' shape -> fallback_title
    (passed from _meta.json) -> empty. Page number gets the slide index.
    """
    from twins.composer import _populate_layout_placeholders
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    # Title text: prefer source-slide 'title' shape, then meta fallback.
    src_title = ""
    src_subtitle = ""
    for shape in src_slide.shapes:
        try:
            name = (shape.name or "").strip().lower()
        except Exception:
            name = ""
        if (not src_title) and (name == "title" or name.startswith("title")):
            try:
                src_title = (shape.text_frame.text or "").strip()
            except Exception:
                src_title = ""
            continue
        # v2.1 (2026-06-02, Bug #2): also harvest the source slide's
        # subtitle textbox so finalize can populate the layout's inherited
        # subtitle placeholder. Worker scripts name the shape "subtitle"
        # by convention when they author a free-floating one.
        if (not src_subtitle) and (name == "subtitle" or name.startswith("subtitle")):
            try:
                src_subtitle = (shape.text_frame.text or "").strip()
            except Exception:
                src_subtitle = ""
            continue
    # M5 (Pattern B production wiring, 2026-06-17): translator's
    # __template_fields__ takes priority over both source-shape harvest and
    # brief fallback. Pattern B HTML stores title/subtitle as DATA fields
    # (data-template-field attributes), extracted by the translator and
    # carried in the script header. They are NEVER positioned as freeform
    # shapes in the body, so the source-slide harvest above never finds
    # them — this branch is how chrome lands for Pattern B picks.
    if template_fields_override:
        if template_fields_override.get("title"):
            src_title = str(template_fields_override["title"]).strip()
        if template_fields_override.get("subtitle"):
            src_subtitle = str(template_fields_override["subtitle"]).strip()
    if not src_title:
        src_title = (fallback_title or "").strip()
    # v2.2 (2026-06-05, SLIDE_LAB_FEEDBACK_LOG #1): fall back to meta-supplied
    # subtitle when source slide has no "subtitle" shape. Workers built before
    # the body-canonical default fix don't name a subtitle shape, so without
    # this fallback the so-what from the brief never lands on the slide.
    if not src_subtitle:
        src_subtitle = (fallback_subtitle or "").strip()

    if dark_variant and dark_bg_hex:
        # Dark path: full-bleed brand-dark overlay, title drawn fresh on top.
        # NOTE: the strip happens upstream in graft_and_theme BEFORE the
        # graft loop (so we don't wipe grafted source content). Here we only
        # add the overlay + title — content is already on the slide.
        hex_s = (dark_bg_hex or "").lstrip("#")
        if len(hex_s) != 6:
            sys.stderr.write(
                f"  WARN: dark_bg_hex {dark_bg_hex!r} malformed on slide "
                f"{slide_n}; skipping dark overlay\n"
            )
            return
        try:
            r = int(hex_s[0:2], 16); g = int(hex_s[2:4], 16); b = int(hex_s[4:6], 16)
        except ValueError:
            sys.stderr.write(
                f"  WARN: dark_bg_hex {dark_bg_hex!r} non-hex on slide "
                f"{slide_n}\n"
            )
            return

        # Full-bleed overlay at z-order position 2 (above layout chrome,
        # below any subsequently appended slide content)
        overlay = new_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, prs.slide_width, prs.slide_height,
        )
        overlay.name = "dark-variant-overlay"
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(r, g, b)
        overlay.line.fill.background()
        sp_tree = new_slide.shapes._spTree
        elem = overlay.element
        sp_tree.remove(elem)
        sp_tree.insert(2, elem)

        # Title at the layout's title-placeholder position (from chrome.yml
        # body-canonical fields), drawn white. body_top_y_px is the title
        # placeholder bottom in our schema, so use a default 40-px-from-top
        # position. If the layout exposes a real title placeholder bound,
        # honor its top. Source slide's grafted shapes will sit above this
        # text box because they were appended last to spTree.
        title_x_emu = Emu(40 * 9525)
        title_y_emu = Emu(40 * 9525)
        title_w_emu = Emu(1200 * 9525)
        title_h_emu = Emu(44 * 9525)
        tb = new_slide.shapes.add_textbox(title_x_emu, title_y_emu,
                                            title_w_emu, title_h_emu)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = src_title or ""
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        return

    # v2.2 (2026-06-05, SLIDE_LAB_FEEDBACK_LOG #3): Mario's design rule —
    # "if title wraps to >2 lines, no subtitle." Use Pillow + brand TTF
    # measurement (NOT char count — proportional fonts misfire on char
    # thresholds). Only fires when the layout actually exposes a subtitle
    # slot (otherwise there's nothing to drop) and src_subtitle is set.
    # Gated to LIGHT path only — dark-variant titles are drawn freshly and
    # don't share the subtitle box.
    _subtitle_idx = getattr(layout_chrome, "subtitle_placeholder_idx", None)
    if src_title and src_subtitle and _subtitle_idx is not None:
        from _chrome_schema import (
            count_wrapped_lines,
            _find_brand_ttf,
            TitleMetricsUnavailableError,
        )
        # Prefer brand-yml-recorded TTF path (Gate A.3, persisted at register
        # time); fall back to layout-chrome override; final fallback is a
        # transitional disk scan with a WARN telling the operator to
        # re-register.
        _ttf = (
            getattr(layout_chrome, "title_font_ttf_path", None)
            or brand_ttf_path
            or _find_brand_ttf()
        )
        # Title box width in px: layout's title_box_width_px if available,
        # else canonical CANONICAL_TITLE_W. Title font size: layout's
        # title_font_pt if available, else 28pt (canonical for body-canonical
        # FedEx layouts).
        _title_w_px = getattr(layout_chrome, "title_box_width_px", None) or 1190
        _title_pt = getattr(layout_chrome, "title_font_pt", None) or 28
        try:
            _n_lines = count_wrapped_lines(src_title, _ttf, _title_pt, _title_w_px)
            if _n_lines >= 3:
                sys.stderr.write(
                    f"  INFO: slide {slide_n} title wraps to {_n_lines} lines "
                    f"at {_title_pt}pt in {_title_w_px}px box; dropping "
                    f"subtitle per Mario's >2-line rule.\n"
                )
                src_subtitle = ""
        except TitleMetricsUnavailableError as _exc:
            # Transitional fallback: brand.yml doesn't yet record the title
            # TTF (older registration). Char-count proxy at ≥110 chars is
            # imperfect for proportional fonts but better than skipping the
            # rule entirely. Re-register the template to fix permanently.
            sys.stderr.write(
                f"  WARN: slide {slide_n} title-wrap measurement unavailable "
                f"({_exc}); falling back to char-count proxy. Recovery: "
                f"re-run register_template to record the brand TTF path.\n"
            )
            if len(src_title) > 110:
                sys.stderr.write(
                    f"  INFO: slide {slide_n} title is {len(src_title)} chars "
                    f"(>110, char-count proxy); dropping subtitle per "
                    f"Mario's >2-line rule.\n"
                )
                src_subtitle = ""

    # Light path: keep existing behavior — overlay only if the layout's
    # chrome ships a body_overlay_hex (legacy per-layout dark path), then
    # populate inherited placeholders.
    from twins.composer import _insert_dark_overlay
    overlay_hex = getattr(layout_chrome, "body_overlay_hex", None)
    body_top_px = getattr(layout_chrome, "body_top_y_px", None)
    body_bot_px = getattr(layout_chrome, "body_bottom_y_px", None)
    if overlay_hex and body_top_px is not None and body_bot_px is not None:
        _insert_dark_overlay(
            new_slide, prs,
            int(body_top_px) * 9525,
            int(body_bot_px) * 9525,
            overlay_hex,
        )

    # v2.1 (2026-06-02, Bug #2): pass subtitle through so finalize populates
    # the inherited subtitle placeholder. Composer no-ops cleanly when
    # subtitle is None or the layout has no subtitle placeholder type.
    # v2.2 (2026-06-05, SLIDE_LAB_FEEDBACK_LOG #2 follow-up): pass the
    # chrome-registered placeholder ids so the composer can match by idx
    # FIRST. FedEx layout's "subtitle" slot is actually a BODY-type
    # placeholder at idx=10; strict type matching alone would silently
    # miss it.
    _title_idx = getattr(layout_chrome, "title_placeholder_idx", None)
    _subtitle_idx = getattr(layout_chrome, "subtitle_placeholder_idx", None)
    found = _populate_layout_placeholders(
        new_slide,
        title=src_title or None,
        subtitle=src_subtitle or None,
        footer=None,
        page_num=str(slide_n),
        title_idx=_title_idx,
        subtitle_idx=_subtitle_idx,
    )

    # v2.2 (2026-06-05, SLIDE_LAB_FEEDBACK_LOG #2): loud-fail when non-empty
    # title or subtitle text was supplied but the layout had no matching
    # placeholder type to populate. The FedEx OTC bug shipped because this
    # check did not exist — every slide's subtitle silently disappeared
    # because the bespoke layout had no SUBTITLE placeholder. Honor
    # feedback_sidecar_fallback_must_be_loud: invisible data loss is never
    # a recoverable runtime state.
    layout_name = getattr(layout_chrome, "name", "(unknown)") if layout_chrome else "(unknown)"
    if src_title and not found.get("title"):
        raise TitleDropError(
            f"slide {slide_n}: title text supplied ({len(src_title)} chars) "
            f"but layout {layout_name!r} has no TITLE placeholder to host it.\n"
            f"  This is the silent-drop bug class — every word of the title "
            f"would have been thrown away.\n"
            f"  Recovery: pick a body-canonical layout for this slide (edit "
            f"_meta.json slides[].layout), or re-register the template so "
            f"default_content_layout points to a body-canonical layout."
        )
    if src_subtitle and not found.get("subtitle"):
        raise SubtitleDropError(
            f"slide {slide_n}: subtitle text supplied ({len(src_subtitle)} chars) "
            f"but layout {layout_name!r} has no SUBTITLE placeholder to host it.\n"
            f"  This is the silent-drop bug class — the so-what for this slide "
            f"would have been invisibly dropped.\n"
            f"  Recovery: pick a body-canonical layout for this slide (edit "
            f"_meta.json slides[].layout to a layout with "
            f"subtitle_placeholder_idx set), or omit the subtitle field if "
            f"this slide intentionally has no so-what."
        )


class TitleDropError(RuntimeError):
    """Raised when title text was supplied but the slide's layout has no
    TITLE placeholder. Per feedback_sidecar_fallback_must_be_loud — silent
    data loss is never a recoverable runtime state. The FedEx OTC bug
    (2026-06-05) shipped a deck with every subtitle silently dropped
    because this check did not exist for subtitle; add for both to prevent
    recurrence in either direction.
    """


class SubtitleDropError(RuntimeError):
    """Raised when subtitle text was supplied but the slide's layout has
    no SUBTITLE placeholder. See TitleDropError for the design rationale.
    """


class TemplateLayoutMissingError(RuntimeError):
    """Raised when a per-slide layout name is supplied but no layout in the
    registered template matches that name.

    Killing Feedback Log Deferred #1 (2026-06-05): _find_named_layout returned
    None silently and the caller fell back to _find_blank_layout, hiding the
    fact that the slide is now on the wrong layout. The most common real-world
    cause: template was re-saved in PowerPoint and a layout got renamed (or
    deleted) after registration; per-slide _meta.json still references the old
    name. Silent fallback = subtitle drop + chrome misalignment with no signal.

    Per feedback_sidecar_fallback_must_be_loud — silent template-vs-meta drift
    is invisible data loss.
    """


class PrimaryAccentCollisionError(RuntimeError):
    """Raised when the loaded ClientTheme's primary_hex and accent_hex are
    too close in RGB space — a signal that the registration step picked the
    wrong color slots (e.g., swapped primary and accent, or both pointing at
    the same brand color).

    Catches a class of error the visual-confirmation gate at registration is
    designed to prevent, but acts as a backstop at build time in case the
    user accepted a bad auto-pick without confirming the preview.
    """


class DarkVariantCollisionError(RuntimeError):
    """Raised when a dark-variant slide has content whose color collides with
    the brand's ``dark_bg_hex``, making the content invisible against the
    full-bleed dark overlay.

    Hard fail (not a warning) because invisible content is worse than a build
    failure — at least the failure tells the user exactly what to fix.

    Per ``feedback_sidecar_fallback_must_be_loud``: silent fallback on
    rendering issues is the v1 bug class we're not repeating.
    """


def _hex_to_rgb_tuple(hex_str: str) -> tuple[int, int, int]:
    s = (hex_str or "").lstrip("#")
    if len(s) != 6:
        return (0, 0, 0)
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return (0, 0, 0)


def _rgb_distance(hex_a: str, hex_b: str) -> float:
    """Euclidean distance in RGB space between two hex strings.
    Max value ~441 (#000000 vs #FFFFFF). 0 = identical.
    """
    a = _hex_to_rgb_tuple(hex_a)
    b = _hex_to_rgb_tuple(hex_b)
    return ((a[0] - b[0]) ** 2 + (a[1] - b[1]) ** 2 + (a[2] - b[2]) ** 2) ** 0.5


# Distance under which a fill/text color is considered "collides with dark_bg"
# and would render effectively invisible. Tuned empirically: identical hex = 0,
# imperceptible difference < 30, low contrast < 60. We hard-fail at < 40 to
# catch exact-and-near-exact collisions without flagging legitimately distinct
# colors that just happen to share a hue family.
_DARK_VARIANT_COLLISION_THRESHOLD = 40.0


def _check_dark_variant_collisions(slide, dark_bg_hex: str,
                                    slide_n: int) -> list[str]:
    """Walk every shape on a slide; flag any fill or text-run color whose RGB
    distance to ``dark_bg_hex`` is below the collision threshold.

    Returns a list of human-readable collision descriptions (empty list = no
    collisions). Skips the dark-overlay shape itself (it's SUPPOSED to be the
    dark_bg color).
    """
    issues: list[str] = []
    if not dark_bg_hex:
        return issues
    dark_hex = (dark_bg_hex or "").lstrip("#").upper()

    def _check(label: str, hex_in: str) -> None:
        if not hex_in:
            return
        h = (hex_in or "").lstrip("#").upper()
        if len(h) != 6:
            return
        dist = _rgb_distance(dark_hex, h)
        if dist < _DARK_VARIANT_COLLISION_THRESHOLD:
            issues.append(
                f"slide {slide_n}: {label} (#{h}) collides with "
                f"dark_bg_hex (#{dark_hex}), distance={dist:.1f}"
            )

    for shape in slide.shapes:
        try:
            shape_name = (shape.name or "<unnamed>").strip()
        except Exception:
            shape_name = "<unnamed>"
        # Skip the overlay itself
        if shape_name == "dark-variant-overlay":
            continue

        # Shape solid fill: collision = shape disappears against the dark bg
        shape_has_solid_fill = False
        try:
            fill = getattr(shape, "fill", None)
            if fill is not None and fill.type == 1:  # MSO_FILL_TYPE.SOLID
                shape_has_solid_fill = True
                try:
                    rgb = fill.fore_color.rgb
                    if rgb is not None:
                        _check(f"shape {shape_name!r} fill", str(rgb))
                except Exception:
                    pass
        except Exception:
            pass

        # Text run colors — only flag if the containing shape has NO solid
        # fill. If the shape has a solid fill, the text sits on THAT fill
        # (not on the dark bg), so a "collision with dark_bg" here is a
        # false positive (e.g., dark text inside a white bar on a purple bg).
        if not shape_has_solid_fill:
            try:
                tf = getattr(shape, "text_frame", None)
                if tf is not None:
                    for p in tf.paragraphs:
                        for run in p.runs:
                            try:
                                c = run.font.color
                                if c is not None and c.type is not None:
                                    try:
                                        rgb = c.rgb
                                        if rgb is not None:
                                            _check(
                                                f"text in {shape_name!r} "
                                                f"(\"{(run.text or '')[:30]}\")",
                                                str(rgb),
                                            )
                                    except Exception:
                                        pass
                            except Exception:
                                pass
            except Exception:
                pass

    return issues


def graft_and_theme(st: OptionStatus, template_path: Path, theme, color_map,
                    layout_name: str = "",
                    layout_chrome=None,
                    slide_title: str = "",
                    slide_subtitle: str = "",
                    slide_variant: str = "") -> None:
    try:
        src_prs = Presentation(str(st.pptx_path))
        src_slide = src_prs.slides[0]
        st.n_shapes = len(src_slide.shapes)

        prs = Presentation(str(template_path))
        _clear_existing_slides(prs)
        # v0.2 P1.3: graft onto the layout named in chrome.yml / meta when one
        # is supplied; fall back to blank for body-canonical and for slides
        # that didn't carry a layout (P1.4 wires per-slide layout via meta).
        #
        # Gate B.2 (2026-06-08, SLIDE_LAB_FEEDBACK_LOG Deferred #1): when a
        # non-empty layout_name was requested but not found, raise loudly
        # instead of silently falling back to blank. The blank fallback is
        # ONLY for slides that asked for no specific layout (layout_name == "").
        target_layout = _find_named_layout(prs, layout_name)
        if target_layout is None:
            if (layout_name or "").strip():
                available = sorted({
                    (lyt.name or "").strip()
                    for m in prs.slide_masters
                    for lyt in m.slide_layouts
                    if (lyt.name or "").strip()
                })
                raise TemplateLayoutMissingError(
                    f"This slide asks to use layout {layout_name!r}, but no "
                    f"layout with that name exists in your template "
                    f"({template_path.name!r}).\n\n"
                    f"  Most common cause: someone renamed or deleted that "
                    f"layout in PowerPoint after the template was registered "
                    f"with Slide Lab.\n\n"
                    f"  Layouts currently in this template: {available}\n\n"
                    f"  What to do:\n"
                    f"    1. If you (or the client) recently re-saved the "
                    f"template, re-register it:\n"
                    f"         py -3 scripts/register_template.py propose "
                    f"\"{template_path}\"\n"
                    f"    2. Otherwise, open _meta.json in the build output "
                    f"folder and change this slide's 'layout' field to one "
                    f"of the names above."
                )
            target_layout = _find_blank_layout(prs)
        new_slide = prs.slides.add_slide(target_layout)

        # v0.3 (2026-05-28) body-canonical layout inheritance:
        # When the chosen layout is body-canonical AND chrome.yml carries the
        # v2 layout-inheritance fields (title_placeholder_idx set), KEEP the
        # layout's inherited placeholders + decorative shapes. Slide Lab's
        # source-slide content gets grafted on top; the inherited title +
        # footer placeholders get populated with text from the source slide
        # after the graft. Otherwise (bespoke / cover / v1 chrome) strip the
        # placeholders so Slide Lab redraws everything at canonical position.
        # Dark-variant slides also strip — they want a clean canvas under
        # the full-bleed overlay (no gradient bar / orange bar bleeding
        # around the dark fill), and they draw their own title white on the
        # overlay rather than populating the inherited title placeholder.
        # MUST happen BEFORE the graft loop, otherwise the strip wipes out
        # the just-grafted source content (audit-failed 2026-05-28).
        _is_body_canonical = (
            layout_chrome is not None
            and getattr(layout_chrome, "layout_class", None) == "body-canonical"
            and getattr(layout_chrome, "title_placeholder_idx", None) is not None
        )
        _is_dark_variant = (slide_variant or "").strip().lower() == "dark"
        # Dark variant: full strip (hide master too) — gives a clean canvas
        # under the full-bleed dark overlay so the master's purple bars /
        # header rules don't bleed around the dark fill.
        # Bespoke non-dark: respect brand.yml `strip_master_backgrounds`.
        # When false (default for FedEx-style brands), the master IS the
        # brand chrome — keep the top/bottom bands inherited from the master.
        # When true (legacy templates with photographic backgrounds), strip.
        # Body-canonical: no strip — inherits everything including placeholders
        # for the post-graft populate step.
        if _is_dark_variant:
            _strip_layout_placeholders(new_slide)
        elif not _is_body_canonical:
            _keep_master = not bool(getattr(theme, "strip_master_backgrounds", True))
            _strip_layout_placeholders(new_slide, keep_master_shapes=_keep_master)

        sp_tree = new_slide.shapes._spTree
        # Gate B.1 fix (2026-06-08, SLIDE_LAB_FEEDBACK_LOG #8 root cause):
        # When body-canonical, new_slide already has title/subtitle/footer
        # placeholders via layout inheritance — _populate_layout_placeholders
        # writes into them after the graft. If we also deepcopy the src_slide's
        # placeholder shapes here, we end up with N+1 copies of each placeholder
        # on every finalize re-run (1 inherited + 1 from src). Five re-runs in
        # one session → 6 stacked title placeholders, all populated identically,
        # rendering as garbled anti-aliased text. compile_picks.py had a
        # post-hoc dedupe; this is the real fix.
        _PH_TAG = qn("p:ph")
        for shape in src_slide.shapes:
            # Picture shapes carry a blipFill that references an image via
            # an rId in the source part's relationships. A naive deepcopy
            # of shape.element preserves the rId string but does NOT carry
            # the relationship onto the target slide's part, leaving an
            # orphan reference that LibreOffice and PowerPoint both render
            # as nothing. Re-embed the blob via add_picture so python-pptx
            # registers a fresh, valid rId on the target. See finalize_deck
            # committee 2026-05-26 root-cause report — Mermaid fallback live
            # test surfaced this.
            #
            # Custom <p:pic> children (srcRect crop, recolor, alpha, rot,
            # picture-effects) are not preserved here. Today's only Picture
            # consumer (_assemble_fallback_pptx) writes plain add_picture
            # calls with no transforms, so this is safe. Add transform
            # preservation here if a future fallback writes Picture shapes
            # with effects.
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    new_slide.shapes.add_picture(
                        io.BytesIO(blob),
                        left=shape.left, top=shape.top,
                        width=shape.width, height=shape.height,
                    )
                    continue
                except Exception:
                    # Fall through to naive copy; it will render blank but
                    # preserves the (broken) shape so the operator can
                    # diagnose. Better than swallowing the picture entirely.
                    pass
            # Body-canonical: skip placeholder shapes — new_slide already
            # has them via inheritance, populate step writes the content.
            if _is_body_canonical and shape.element.find(".//" + _PH_TAG) is not None:
                continue
            sp_tree.append(deepcopy(shape.element))

        # v0.3 body-canonical: insert dark-overlay rectangle (if any) so it
        # sits BEFORE the source-slide content in z-order, and write the source
        # slide's title + footer text into the layout's inherited placeholders.
        if _is_body_canonical:
            try:
                _is_dark = (slide_variant or "").strip().lower() == "dark"
                _dark_bg = getattr(theme, "dark_bg_hex", "") or ""
                # M5: for Pattern B translator output, parse the script's
                # __template_fields__ header. Result is {} for Pattern C /
                # legacy classifications (the parser exits early if the
                # marker isn't found), so this is safe for every code path.
                _tf_override = None
                if st.classification == "pattern_b_translated":
                    _tf_override = _parse_template_fields(st.py_path)
                _apply_body_canonical_finishing(
                    new_slide, prs, layout_chrome, src_slide, st.slide_n,
                    fallback_title=slide_title,
                    fallback_subtitle=slide_subtitle,
                    dark_variant=_is_dark,
                    dark_bg_hex=_dark_bg,
                    brand_ttf_path=getattr(theme, "title_font_ttf_path", "") or "",
                    template_fields_override=_tf_override,
                )
                # v0.3: hard-fail collision check on dark-variant slides.
                # If any shape fill or text color collides with dark_bg_hex,
                # the content would render invisible. Record on st so the
                # main loop can refuse the build after all options finish.
                if _is_dark and _dark_bg:
                    st.dark_collisions = _check_dark_variant_collisions(
                        new_slide, _dark_bg, st.slide_n,
                    )
            except (TitleDropError, SubtitleDropError,
                    TemplatePlaceholderEmptyError):
                # v2.2 (2026-06-05, SLIDE_LAB_FEEDBACK_LOG #2): silent-drop
                # errors are hard fails per feedback_sidecar_fallback_must_be_loud.
                # Re-raise so the graft fails loudly with the named exception
                # and stops the build — invisible data loss is worse than
                # a stopped build.
                # Gate B.3 (2026-06-08): TemplatePlaceholderEmptyError is the
                # same silent-drop bug class one layer deeper; same handling.
                raise
            except Exception as _exc:
                # Don't fail the graft on overlay/populate trouble — surface
                # in qc later. The grafted slide is still valid.
                sys.stderr.write(
                    f"  WARN: body-canonical finishing skipped on slide "
                    f"{st.slide_n}: {type(_exc).__name__}: {_exc}\n"
                )

        subs = 0
        for shape in new_slide.shapes:
            subs += apply_theme_to_shape_xml(
                shape.element, color_map,
                major_font=theme.major_font, minor_font=theme.minor_font,
            )
        st.n_subs = subs

        st.themed_pptx_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(st.themed_pptx_path))
        st.themed = True
    except (TitleDropError, SubtitleDropError,
            TemplateLayoutMissingError, TemplatePlaceholderEmptyError):
        # Gate B fixes (2026-06-08): named silent-drop / layout-drift errors
        # must propagate out of graft_and_theme so main() halts the whole
        # build (not just records one option as failed). Per
        # feedback_sidecar_fallback_must_be_loud — invisible data loss is
        # worse than a stopped build. Without this, the outer generic
        # `except Exception` below would catch them and the build would
        # continue with the bad option recorded only in RESULT.md.
        st.themed = False
        raise
    except Exception as e:
        st.themed = False
        st.error = f"graft: {type(e).__name__}: {e}"


# ---------------------------------------------------------------------------
# Step 3: render to PNG — verbatim from v1
# ---------------------------------------------------------------------------
def render_one(st: OptionStatus) -> OptionStatus:
    from render_slides import render_libre
    pptx = st.themed_pptx_path
    tmp = _p.render_tmp_dir(pptx)
    tmp.mkdir(parents=True, exist_ok=True)
    try:
        render_libre(pptx, tmp, dpi=120)
        src_png = tmp / "slide_01.png"
        if src_png.exists():
            src_png.replace(st.themed_png_path)
            st.rendered = True
        else:
            st.rendered = False
            st.error = "render: no png produced"
    except Exception as e:
        st.rendered = False
        st.error = f"render: {type(e).__name__}: {e}"
    return st


# ---------------------------------------------------------------------------
# Step 4: RESULT.md — v1 base + v2 classification column
# ---------------------------------------------------------------------------
RESULT_TEMPLATE = """# Slide Lab v2 deck — finalize result

Generated: {ts}

Out: `{out}`
Template: `{template}`

## Counts

- Total options: **{total}**
- Native      : {native_count}
- Pattern B   : {pattern_b_count}
- Rejected    : {rejected_count}
- Missing     : {missing_count}
- Built       : **{built_ok} / {total}**
- Themed      : **{themed_ok} / {total}**
- Rendered    : **{rendered_ok} / {total}**

## Per-option status

| Slide | Option | Class | Built | Themed | Rendered | Shapes | Subs | Error/Reason |
|-------|--------|-------|-------|--------|----------|--------|------|--------------|
{rows}

## Outputs

- **Themed PPTX**: `<out>/slide_NN/option_X.pptx`
- **PNG thumbnails**: `<out>/slide_NN/option_X.png`
- **QC self-check**: `<out>/slide_NN/option_X.qc.json`
- **Raw pre-theme PPTX**: `<out>/slide_NN/_raw/option_X.pptx`

## Failures + rejections

{failures}
"""


def _mark(ok):
    if ok is True:
        return "ok"
    if ok is False:
        return "FAIL"
    return "-"


def _class_short(c: str) -> str:
    return {
        "native": "native",
        "pattern_b_translated": "pattern_b",
        "skeleton_rejected": "rejected",
        "missing": "missing",
    }.get(c, c)


def write_result(out_dir: Path, template_path: Path, statuses: list) -> Path:
    rows = []
    for st in statuses:
        err_or_reason = st.error or (st.classification_reason if st.classification != "native" else "")
        rows.append(
            f"| {st.slide_n:>2} | {st.letter} | {_class_short(st.classification)} | "
            f"{_mark(st.built)} | {_mark(st.themed)} | {_mark(st.rendered)} | "
            f"{st.n_shapes} | {st.n_subs} | "
            f"{(err_or_reason[:80] + '...') if len(err_or_reason) > 80 else err_or_reason} |"
        )

    failures = []
    for st in statuses:
        if st.classification == "skeleton_rejected":
            failures.append(f"- **slide {st.slide_n:02d} option {st.letter}** [REJECTED]: {st.classification_reason}")
        elif st.classification == "missing":
            failures.append(f"- **slide {st.slide_n:02d} option {st.letter}** [MISSING]: {st.classification_reason} — re-dispatch the worker if a fuller deck is wanted")
        elif st.error:
            failures.append(f"- **slide {st.slide_n:02d} option {st.letter}**: {st.error}")

    total = len(statuses)
    content = RESULT_TEMPLATE.format(
        ts=datetime.now().isoformat(timespec="seconds"),
        out=out_dir,
        template=template_path,
        total=total,
        native_count=sum(1 for s in statuses if s.classification == "native"),
        pattern_b_count=sum(1 for s in statuses if s.classification == "pattern_b_translated"),
        rejected_count=sum(1 for s in statuses if s.classification == "skeleton_rejected"),
        missing_count=sum(1 for s in statuses if s.classification == "missing"),
        built_ok=sum(1 for s in statuses if s.built),
        themed_ok=sum(1 for s in statuses if s.themed),
        rendered_ok=sum(1 for s in statuses if s.rendered),
        rows="\n".join(rows) if rows else "| (no options found) |",
        failures="\n".join(failures) if failures else "(none)",
    )
    target = out_dir / "RESULT.md"
    target.write_text(content, encoding="utf-8")
    return target


# ---------------------------------------------------------------------------
# Main — v1 base + v2 Mermaid theme resolution + classification counters
# ---------------------------------------------------------------------------
def main() -> int:
    ap = argparse.ArgumentParser(description="Slide Lab v2 deck orchestrator — Part B (finalize)")
    ap.add_argument("--out", required=True, type=Path, help="Output dir from Part A")
    ap.add_argument("--template", required=True, type=Path, help="Client PPTX template")
    ap.add_argument("--skip-build", action="store_true", help="Skip executing option_X.py")
    ap.add_argument("--skip-render", action="store_true", help="Skip PNG rendering")
    ap.add_argument("--allow-missing", action="store_true",
                    help="Proceed even when some expected option_X.py files are absent "
                         "(e.g., an interrupted worker). Default: hard-fail with a re-dispatch "
                         "punch list so the deck doesn't ship with hidden gaps.")
    args = ap.parse_args()

    from _log import attach as _log_attach
    _log_attach(args.out, "finalize_deck.py")

    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

    if not args.out.exists():
        print(f"ERROR: out dir not found: {args.out}")
        return 2
    if not args.template.exists():
        print(f"ERROR: template not found: {args.template}")
        return 2

    # M7 (Mermaid retirement, 2026-06-17): the Mermaid theme resolution step
    # was removed here. `_resolve_mermaid_theme` and the `--theme` CLI arg
    # are retired; Pattern B HTML→PNG (per Decision 6) supersedes Mermaid for
    # curved-container diagrams.

    # M1 — Pattern B routing read (2026-06-16). Defensive only: log the
    # effective pattern from _meta.json so operator output makes routing
    # visible. Legacy mode (pattern_default absent or "legacy") proceeds
    # with the existing code path verbatim. Translator dispatch (Stage 3.5)
    # lands in M4; until then, even a non-legacy pattern_default takes the
    # current path. This keeps M1 byte-identical to pre-M1 for legacy builds.
    try:
        _meta_path = _p.meta_json(args.out)
        _meta_dict = json.loads(_meta_path.read_text(encoding="utf-8"))
        _pattern_default = _meta_dict.get("pattern_default")
        if _pattern_default and _pattern_default != "legacy":
            print(
                f"  [pattern routing] _meta.json::pattern_default = "
                f"{_pattern_default!r}; per-slide map = "
                f"{_meta_dict.get('pattern_per_slide', {})}"
            )
            print(
                f"  [pattern routing] NOTE: Stage 3.5 (translator dispatch) "
                f"lands in M4. Proceeding via legacy path for this build."
            )
    except (OSError, json.JSONDecodeError):
        # Defensive read; never block the build on a parse error here.
        pass

    # v0.2 P1.3: load + verify chrome.yml. Hard-fail on sha mismatch — refuse
    # to build slides against a chrome spec the template has drifted away from.
    try:
        chrome_spec = _load_and_verify_chrome(args.template)
    except ChromeSidecarMissingError as exc:
        print(
            f"ERROR: chrome sidecar missing.\n{exc}\n\n"
            f"Re-register the template (one-time):\n"
            f"  py -3 slide-builder/scripts/register_template.py propose "
            f"\"{args.template}\"\n"
            f"  # chat shows preview + picks; copy picks.json back\n"
            f"  py -3 slide-builder/scripts/register_template.py commit "
            f"\"{args.template}\" --picks <picks.json>",
            file=sys.stderr,
        )
        return 7
    except RuntimeError as exc:
        print(f"ERROR: chrome.yml is stale.\n{exc}", file=sys.stderr)
        return 7
    chrome_yml_path = _p.chrome_yml(args.template)
    default_layout_name = _pick_default_layout_name(chrome_spec)

    print("=" * 72)
    print("Slide Lab v2 deck orchestrator — Part B (finalize)")
    print(f"  out           : {args.out}")
    print(f"  template      : {args.template}")
    print(f"  chrome.yml    : {chrome_yml_path}")
    print(f"  layouts       : {len(chrome_spec.layouts)} "
          f"(default for slides w/o layout field: {default_layout_name!r})")
    print("=" * 72)

    print("\n[1] Discover option_X.py files + classify by line-1 token")
    # Compute the expected (slide_n, letter) set from _meta.json so
    # discover_options can surface missing worker outputs as classification
    # "missing" — committee P4 fix (2026-05-26). If _meta.json is missing
    # or unparseable, fall back to filesystem-only discovery (legacy
    # behavior).
    expected_pairs: set = set()
    try:
        meta_dict = json.loads(_p.meta_json(args.out).read_text(encoding="utf-8"))
        from _meta_schema import validate_warn
        validate_warn(meta_dict, source="finalize_deck:expected_pairs")
        for s in meta_dict.get("slides", []):
            n = s.get("n")
            if isinstance(n, int):
                for letter in _p.option_letters():
                    expected_pairs.add((n, letter))
    except Exception:
        expected_pairs = set()

    # v0.2 P1.3: slide_n -> layout name lookup. meta.slides[i].layout is
    # added in P1.4; until then every slide gets default_layout_name.
    # v0.3: slide_n -> variant lookup ("dark" or "" / "light"). Read from
    # meta.slides[i].variant; defaults to light.
    _slide_layout_names: dict[int, str] = {}
    _slide_titles: dict[int, str] = {}
    _slide_subtitles: dict[int, str] = {}
    _slide_variants: dict[int, str] = {}
    try:
        _meta_dict_layouts = json.loads(_p.meta_json(args.out).read_text(encoding="utf-8"))
        for s in _meta_dict_layouts.get("slides", []):
            n = s.get("n")
            if isinstance(n, int):
                _slide_layout_names[n] = (s.get("layout") or "").strip() or default_layout_name
                _slide_titles[n] = (s.get("title") or "").strip()
                _slide_subtitles[n] = (s.get("subtitle") or "").strip()
                _slide_variants[n] = (s.get("variant") or "").strip().lower()
    except Exception:
        _slide_layout_names = {}
        _slide_titles = {}
        _slide_subtitles = {}
        _slide_variants = {}

    def _layout_name_for(slide_n: int) -> str:
        return _slide_layout_names.get(slide_n, default_layout_name)

    def _slide_title_for(slide_n: int) -> str:
        return _slide_titles.get(slide_n, "")

    def _slide_subtitle_for(slide_n: int) -> str:
        return _slide_subtitles.get(slide_n, "")

    def _slide_variant_for(slide_n: int) -> str:
        return _slide_variants.get(slide_n, "")

    def _layout_chrome_for(slide_n: int):
        name = _layout_name_for(slide_n)
        lc = chrome_spec.layouts.get(name)
        if lc is None:
            # Loud-fail per feedback_sidecar_fallback_must_be_loud. Silent
            # fallback to a default layout was the OTC Sizing bug class —
            # slides got built against geometry that didn't match the
            # template's actual layout shapes.
            available = ", ".join(sorted(chrome_spec.layouts.keys())) or "(none)"
            raise ChromeLayoutMissingError(
                f"slide {slide_n} references layout {name!r}, which is not "
                f"in chrome.yml.\n"
                f"  available layouts: {available}\n"
                f"  Either fix the slide's 'layout' field in _meta.json, or "
                f"re-register the template:\n"
                f"    py -3 slide-builder/scripts/register_template.py propose "
                f"\"<template.pptx>\"\n"
                f"    py -3 slide-builder/scripts/register_template.py commit "
                f"\"<template.pptx>\" --picks <picks.json>"
            )
        return lc

    statuses = discover_options(args.out, expected=expected_pairs)
    n_native = sum(1 for s in statuses if s.classification == "native")
    n_pattern_b = sum(1 for s in statuses if s.classification == "pattern_b_translated")
    n_rejected = sum(1 for s in statuses if s.classification == "skeleton_rejected")
    n_missing = sum(1 for s in statuses if s.classification == "missing")
    if n_missing:
        # Pre-flight gate (RC-5 fix, 2026-06-02): if any expected option_X.py
        # is absent, halt BEFORE the build loop with a punch list of slides
        # that need re-dispatch. Otherwise finalize runs for several minutes,
        # then crashes mid-build when it tries to execute the missing script,
        # and the operator has to scroll through logs to find which slides
        # failed. Override with --allow-missing for the rare partial-rebuild
        # case (e.g., the operator only wants to re-finalize 5/6 slides).
        missing_groups: dict[int, list[str]] = {}
        for s in statuses:
            if s.classification == "missing":
                missing_groups.setdefault(s.slide_n, []).append(s.letter)
        if not args.allow_missing:
            sys.stderr.write(
                f"\nERROR: {n_missing} expected worker output(s) absent. "
                f"finalize would crash mid-build.\n\n"
                f"Slides needing re-dispatch:\n"
            )
            for slide_n in sorted(missing_groups):
                letters = ", ".join(sorted(missing_groups[slide_n]))
                slide_dir = _p.slide_dir(args.out, slide_n)
                prompt_path = _p.prompt_md(args.out, slide_n)
                sys.stderr.write(
                    f"  slide_{slide_n:02d}  (missing option(s): {letters})\n"
                    f"    prompt:  {prompt_path}\n"
                    f"    out dir: {slide_dir}\n"
                )
            sys.stderr.write(
                f"\nRe-dispatch the worker agent for each slide above using "
                f"the prompt file, then re-run finalize_deck.py.\n"
                f"Override with --allow-missing to proceed with gaps "
                f"(surfaced as classification=missing in RESULT.md).\n"
            )
            return 11
        print(f"  [missing] {n_missing} worker output(s) absent — proceeding "
              f"under --allow-missing (will be surfaced in RESULT.md / REVIEW.html)")
    print(f"  found {len(statuses)} option scripts across "
          f"{len({s.slide_n for s in statuses})} slides")
    print(f"  classification: native={n_native}  pattern_b={n_pattern_b}  rejected={n_rejected}")
    if not statuses:
        print("  nothing to do.")
        write_result(args.out, args.template, statuses)
        return 0

    print("\n[2] Execute / assemble / skip per option (serial)")
    for i, st in enumerate(statuses, 1):
        if args.skip_build and st.pptx_path.exists():
            st.built = True
            print(f"  [{i:>3}/{len(statuses)}] slide_{st.slide_n:02d}/option_{st.letter}  [{_class_short(st.classification)}]  skip-build (exists)")
            continue
        build_pptx(
            st,
            chrome_yml_path=chrome_yml_path,
            layout_chrome=_layout_chrome_for(st.slide_n),
            layout_name=_layout_name_for(st.slide_n),
        )
        if st.classification == "skeleton_rejected":
            flag = f"rejected ({st.classification_reason[:50]})"
        elif st.built:
            flag = "ok"
        else:
            flag = f"FAIL ({st.error[:50]})"
        print(f"  [{i:>3}/{len(statuses)}] slide_{st.slide_n:02d}/option_{st.letter}  [{_class_short(st.classification)}]  {flag}")

    built_statuses = [s for s in statuses if s.built]
    print(f"  built: {len(built_statuses)} / {len(statuses)}")

    print("\n[3] Load client theme")
    theme = load_client_theme(str(args.template))
    color_map = theme.color_map()
    expected_palette = _expected_palette_for_theme(theme)
    print(f"  raw dk2={theme.dk2}  lt2={theme.lt2}  font={theme.minor_font}")
    print(f"  brand_primary=#{theme.brand_primary} (source: {theme.brand_primary_source})")
    print(f"  brand_accent =#{theme.brand_accent} (source: {theme.brand_accent_source})")
    print(f"  color_map entries: {len(color_map)}")
    print(f"  expected palette  : {len(expected_palette)} hex codes")

    # Gate B.4 (2026-06-08): primary/accent RGB-distance backstop.
    # Catches the registration-best-guess-inverted class of failure where
    # _best_guess_primary_accent() picks the same slot for both (or two
    # slots whose colors are visually indistinguishable). The visual-
    # confirmation gate at registration is the primary defense; this is the
    # backstop that catches builds where the user accepted a bad auto-pick.
    # Threshold tuned conservatively — primary/accent on a real brand should
    # be at least 60 units apart in RGB-Euclidean (e.g., FedEx purple
    # #4D148C vs orange #FF6600 = ~310 units).
    try:
        _pa_dist = _rgb_distance(theme.brand_primary or "", theme.brand_accent or "")
        if _pa_dist < 30.0:
            raise PrimaryAccentCollisionError(
                f"Your brand's primary color (#{theme.brand_primary}) and "
                f"accent color (#{theme.brand_accent}) are nearly identical "
                f"(RGB distance {_pa_dist:.1f}; we need at least 30).\n\n"
                f"  Why this matters: if primary and accent look the same, "
                f"every accent stripe, callout, and emphasis color on every "
                f"slide will be invisible against the primary fill — the deck "
                f"will look flat and chromatically wrong.\n\n"
                f"  Most common cause: when you registered the template, the "
                f"primary and accent swatches in register.html were the same "
                f"color (or two near-identical swatches).\n\n"
                f"  What to do: re-register the template and pick visibly "
                f"distinct primary and accent colors:\n"
                f"    py -3 scripts/register_template.py propose <your-template.pptx>\n"
                f"  Then open register.html and pick TWO different swatches for "
                f"primary and accent."
            )
        if _pa_dist < 60.0:
            print(
                f"  WARN: primary/accent RGB distance is {_pa_dist:.1f} (soft "
                f"threshold=60). Colors may be visually indistinct on some "
                f"slides. Re-register if this surprises you."
            )
    except PrimaryAccentCollisionError as exc:
        # Surface the operator-facing message cleanly with non-zero exit,
        # not a raw stack trace.
        print(f"\nERROR: brand color collision.\n{exc}", file=sys.stderr)
        return 8

    print("\n[3.1] Detect client-template fonts on local machine")
    font_missing = detect_missing_client_fonts(theme)
    if font_missing:
        print(f"  WARNING: client fonts not installed locally: {', '.join(font_missing)}")
        print(f"  PNG thumbnails will show font substitution; the .pptx itself is correct.")
    else:
        print(f"  all client fonts available (or detection inconclusive)")

    finalize_meta = {
        "template": str(args.template),
        "major_font": theme.major_font,
        "minor_font": theme.minor_font,
        "font_missing": font_missing,
        "classification_counts": {"native": n_native, "pattern_b_translated": n_pattern_b, "skeleton_rejected": n_rejected},
    }
    try:
        _p.finalize_meta_json(args.out).write_text(json.dumps(finalize_meta, indent=2), encoding="utf-8")
    except Exception as e:
        print(f"  WARNING: could not write _finalize_meta.json: {e}")

    print("\n[3.5] Stash raw pre-theme PPTX")
    for st in built_statuses:
        stash_raw(st)

    print("\n[4] Graft + theme remap (serial — python-pptx not thread-safe)")
    for i, st in enumerate(built_statuses, 1):
        try:
            graft_and_theme(
                st, args.template, theme, color_map,
                layout_name=_layout_name_for(st.slide_n),
                layout_chrome=_layout_chrome_for(st.slide_n),
                slide_title=_slide_title_for(st.slide_n),
                slide_subtitle=_slide_subtitle_for(st.slide_n),
                slide_variant=_slide_variant_for(st.slide_n),
            )
        except (TitleDropError, SubtitleDropError,
                TemplateLayoutMissingError,
                TemplatePlaceholderEmptyError) as _exc:
            # Gate B fixes (2026-06-08): silent-drop / layout-drift errors
            # halt the build cleanly with a recovery hint instead of a raw
            # stack trace. The named exception messages were written for the
            # operator; relay them verbatim and exit non-zero.
            print(
                f"  [{i:>3}/{len(built_statuses)}] slide_"
                f"{st.slide_n:02d}/option_{st.letter}  FAIL ({type(_exc).__name__})"
            )
            print(f"\nERROR: build halted — {type(_exc).__name__} on "
                  f"slide {st.slide_n} option {st.letter}.\n{_exc}",
                  file=sys.stderr)
            return 8
        flag = f"ok (shapes={st.n_shapes} subs={st.n_subs})" if st.themed else f"FAIL ({st.error[:50]})"
        print(f"  [{i:>3}/{len(built_statuses)}] slide_{st.slide_n:02d}/option_{st.letter}  {flag}")

    themed_statuses = [s for s in built_statuses if s.themed]
    print(f"  themed: {len(themed_statuses)} / {len(built_statuses)}")

    # v0.3 dark-variant collision gate. Before any further work, refuse to
    # ship a build that contains shapes/text whose color would render
    # invisible against the brand's dark background. Loud, named failure.
    _all_dark_issues: list[str] = []
    for s in themed_statuses:
        _all_dark_issues.extend(s.dark_collisions)
    if _all_dark_issues:
        print("\n[4b] DARK-VARIANT COLLISION CHECK — BUILD REFUSED")
        print(f"  Found {len(_all_dark_issues)} content color(s) that collide "
              f"with brand dark_bg_hex.")
        print(f"  These shapes/text runs would render invisible on the dark "
              f"overlay. Build aborted.\n")
        for issue in _all_dark_issues:
            print(f"    - {issue}")
        print("")
        print("  To fix: open the affected option scripts and change the")
        print("  colliding fill / text color to something with contrast "
              "against dark_bg_hex.")
        print("  Or change brand.yml dark_bg_hex to a color that doesn't "
              "match your content.")
        raise DarkVariantCollisionError(
            f"{len(_all_dark_issues)} dark-variant collision(s); "
            f"build refused. See [4b] output above for the offending slides."
        )

    if not args.skip_render:
        print(f"\n[5] Render themed .pptx -> .png (parallel x4)")
        with ThreadPoolExecutor(max_workers=4) as pool:
            futures = {pool.submit(render_one, st): st for st in themed_statuses}
            done = 0
            for fut in as_completed(futures):
                st = fut.result()
                done += 1
                flag = "ok" if st.rendered else f"FAIL ({st.error[:50]})"
                print(f"  [{done:>3}/{len(themed_statuses)}] slide_{st.slide_n:02d}/option_{st.letter}  {flag}")

        for tmp in args.out.glob("slide_*/_render_tmp"):
            shutil.rmtree(tmp, ignore_errors=True)
    else:
        print("\n[5] Render skipped (--skip-render)")

    print("\n[5b] Per-option QC self-check")
    # Build a slide_n -> page_type lookup so the QC can exempt cover slides
    # from page-number requirement etc. _meta.json was already validated +
    # loaded upstream; re-read here for narrow scope.
    _slide_page_types: dict[int, str] = {}
    try:
        _meta_dict = json.loads(_p.meta_json(args.out).read_text(encoding="utf-8"))
        # No validate_warn here — finalize_deck:expected_pairs already
        # validated this meta upstream; second pass would double-warn.
        for s in _meta_dict.get("slides", []):
            n = s.get("n")
            if isinstance(n, int):
                _slide_page_types[n] = (s.get("page_type") or "").strip()
    except Exception:
        pass

    qc_counts = {"all_ok": 0, "warn_only": 0, "block": 0}
    for st in themed_statuses:
        try:
            pt = _slide_page_types.get(st.slide_n, "")
            result = run_option_qc(st.themed_pptx_path, st.themed_png_path, expected_palette,
                                   page_type=pt)
            # M6 (Pattern B, 2026-06-17): for Pattern B picks, run R4.1-R4.8
            # checks against the translator's report + script. Merge into the
            # same QC JSON so build_review.py renders them uniformly with
            # R1-R3 results. Pattern C / legacy returns [] and is a no-op.
            r4 = _check_r4_rules_for_pattern_b(st)
            if r4:
                st.r4_checks = r4
                result["checks"] = list(result.get("checks", [])) + r4
                summary = result.setdefault("summary", {"pass": 0, "block": 0, "warn": 0})
                for c in r4:
                    if c["pass"]:
                        summary["pass"] = summary.get("pass", 0) + 1
                    elif c["severity"] == "block":
                        summary["block"] = summary.get("block", 0) + 1
                    else:
                        summary["warn"] = summary.get("warn", 0) + 1
            qc_path = st.themed_pptx_path.parent / (st.themed_pptx_path.stem + ".qc.json")
            qc_path.write_text(json.dumps(result, indent=2), encoding="utf-8")
            summ = result["summary"]
            if summ["block"] > 0:
                qc_counts["block"] += 1
                flag = f"BLOCK ({summ['block']}b/{summ['warn']}w)"
            elif summ["warn"] > 0:
                qc_counts["warn_only"] += 1
                flag = f"warn ({summ['warn']}w)"
            else:
                qc_counts["all_ok"] += 1
                flag = "ok"
            print(f"  slide_{st.slide_n:02d}/option_{st.letter}  {flag}")
        except Exception as e:
            print(f"  slide_{st.slide_n:02d}/option_{st.letter}  qc-FAIL ({type(e).__name__}: {e})")
    print(f"  QC totals: all-ok={qc_counts['all_ok']}  warn-only={qc_counts['warn_only']}  block={qc_counts['block']}")

    print("\n[6] Write RESULT.md")
    result_path = write_result(args.out, args.template, statuses)
    print(f"  {result_path}")

    print("\n[7] Generate Gate 3 visual preview (GATE3-PREVIEW.html)")
    preview_script = Path(__file__).resolve().parent / "build_gate_preview.py"
    try:
        result = subprocess.run(
            [sys.executable, str(preview_script), "--out", str(args.out)],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            # Echo the script's output (it prints path + slide count)
            for line in (result.stdout or "").splitlines():
                if line.strip():
                    print(f"  {line}")
        else:
            print(f"  WARNING: build_gate_preview.py exited {result.returncode}")
            for line in (result.stderr or "").splitlines():
                if line.strip():
                    print(f"  {line}")
    except subprocess.TimeoutExpired:
        print(f"  WARNING: build_gate_preview.py timed out (60s)")
    except Exception as exc:
        print(f"  WARNING: build_gate_preview.py invocation failed: {type(exc).__name__}: {exc}")

    print("\n" + "=" * 72)
    print("DONE — Part B complete.")
    print(f"  Native     : {n_native}")
    print(f"  Pattern B  : {n_pattern_b}")
    print(f"  Rejected   : {n_rejected}")
    print(f"  Built    : {sum(1 for s in statuses if s.built)} / {len(statuses)}")
    print(f"  Themed   : {sum(1 for s in statuses if s.themed)} / {len(statuses)}")
    print(f"  Rendered : {sum(1 for s in statuses if s.rendered)} / {len(statuses)}")
    print(f"  Report   : {result_path}")
    print("=" * 72)
    return 0


if __name__ == "__main__":
    sys.exit(main())
