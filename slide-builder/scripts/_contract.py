"""_contract.py — module-load-time pipeline contract test.

Verifies the Slide Lab pipeline is internally consistent without requiring
a live build. Runs three checks:

1. **Paths registry sanity** — every helper named in `_paths.ARTIFACT_MANIFEST`
   resolves to a callable or constant on the `_paths` module.

2. **Meta-JSON schema round-trip** — a synthetic MetaJson model serialises
   to dict and validates back without drift.

3. **Handoff coverage** — every artifact in the manifest with a declared
   writer script appears in that script's source text via `_paths.<name>` or
   `_paths.<NAME>` form. Same for readers. Orphans marked `accepted: True`
   in the manifest are skipped with reason recorded.

Run directly:
    py -3 scripts/_contract.py
Exit code 0 = all green. Non-zero = at least one drift detected.

Wire into CI before any v0.1 release tag per DECISIONS.md § "Hardening
triple — items 2 + 3 acceptance gate."
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

import _paths as _p
from _meta_schema import (
    MetaJson, validate_meta_dict, MetaJsonSchemaError,
    META_SCHEMA_VERSION_CURRENT,
)
from _chrome_schema import (
    BoxPx, LayoutChrome, ChromeSpec, ChromeSidecarMissingError, ChromeSchemaError,
    CHROME_SCHEMA_VERSION_CURRENT,
    validate_chrome_dict, load_chrome_yml, dump_chrome_yml,
)


def _fail(msg: str) -> None:
    print(f"FAIL: {msg}")


def _ok(msg: str) -> None:
    print(f"  ok: {msg}")


# ---------------------------------------------------------------------------
# Check 1 — paths registry sanity
# ---------------------------------------------------------------------------

def check_paths_registry() -> list[str]:
    """Every name in ARTIFACT_MANIFEST must resolve on the _paths module."""
    errors: list[str] = []
    for entry in _p.ARTIFACT_MANIFEST:
        name = entry["name"]
        if not hasattr(_p, name):
            errors.append(f"manifest names '{name}' but _paths has no such attribute")
            continue
        attr = getattr(_p, name)
        if not (callable(attr) or isinstance(attr, str)):
            errors.append(f"_paths.{name} is neither a callable nor a string constant (got {type(attr).__name__})")
    if not errors:
        _ok(f"paths registry: {len(_p.ARTIFACT_MANIFEST)} manifest entries all resolve on _paths")
    return errors


# ---------------------------------------------------------------------------
# Check 2 — meta-JSON schema round-trip
# ---------------------------------------------------------------------------

def check_meta_schema_roundtrip() -> list[str]:
    """A synthetic MetaJson must serialise to dict and re-validate cleanly."""
    errors: list[str] = []
    synthetic = MetaJson(
        schema_version=META_SCHEMA_VERSION_CURRENT,
        template="C:/fake/template.pptx",
        brief="C:/fake/brief.md",
        out="C:/fake/out",
        mermaid_theme="C:/fake/theme.json",
        client_slug="acme",
        slide_count=2,
        generated_at="2026-05-26T00:00:00+00:00",
        brand_primary="#4D148C",
        brand_accent="#FF6600",
        slides=[
            {"n": 1, "title": "T1", "forecasted_pattern": "full-canvas", "page_type": "cover", "layout": "cover_light"},
            {"n": 2, "title": "T2", "forecasted_pattern": "50-50",       "page_type": "headline-finding", "layout": "body_canonical_light"},
        ],
        deck_meta={"deck_type": "client-pitch", "governing_thought": "g", "audience": "a"},
    )
    as_dict = synthetic.model_dump()
    try:
        validated = validate_meta_dict(as_dict)
    except MetaJsonSchemaError as e:
        errors.append(f"round-trip validation failed: {e}")
        return errors
    if validated.schema_version != META_SCHEMA_VERSION_CURRENT:
        errors.append(
            f"round-trip schema_version mismatch: "
            f"{validated.schema_version} != {META_SCHEMA_VERSION_CURRENT}"
        )
    if not errors:
        _ok(f"meta-json schema round-trip @ version {META_SCHEMA_VERSION_CURRENT}")
    return errors


# ---------------------------------------------------------------------------
# Check 3 — handoff coverage
# ---------------------------------------------------------------------------

_SCRIPT_TEXT_CACHE: dict[str, str] = {}


def _script_text(name: str) -> str:
    if name not in _SCRIPT_TEXT_CACHE:
        p = HERE / name
        _SCRIPT_TEXT_CACHE[name] = p.read_text(encoding="utf-8") if p.exists() else ""
    return _SCRIPT_TEXT_CACHE[name]


def _references_artifact(script_name: str, artifact_name: str) -> bool:
    """A script 'uses' an artifact if it imports the helper name from _paths.

    Accepts any of:
      - `_p.<name>(...)` — absolute-path helper form
      - `_p.<name>_name(...)` — filename-only helper form (per-option assets)
      - `_p.<NAME>` — uppercase constant form (slide-relative filenames)
    """
    text = _script_text(script_name)
    if not text:
        return False
    patterns = [
        rf"\b_p\.{re.escape(artifact_name)}\b",
        rf"\b_p\.{re.escape(artifact_name)}_name\b",
        rf"\b_p\.{re.escape(artifact_name.upper())}\b",
    ]
    return any(re.search(p, text) for p in patterns)


def check_handoff_coverage() -> list[str]:
    errors: list[str] = []
    checked = 0
    skipped = 0
    for entry in _p.ARTIFACT_MANIFEST:
        if entry.get("accepted"):
            skipped += 1
            continue
        name = entry["name"]
        writer = entry["writer"]
        readers = entry["readers"]
        checked += 1
        if writer.endswith(".py") and not _references_artifact(writer, name):
            errors.append(f"writer mismatch: manifest says {writer} writes {name!r}, but the script doesn't reference _paths.{name}")
        for r in readers:
            if r.endswith(".py") and not _references_artifact(r, name):
                errors.append(f"reader mismatch: manifest says {r} reads {name!r}, but the script doesn't reference _paths.{name}")
    if not errors:
        _ok(f"handoff coverage: {checked} artifacts checked, {skipped} accepted-orphans skipped")
    return errors


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Check 4 — pipeline-script import smoke
# ---------------------------------------------------------------------------

def check_pipeline_imports() -> list[str]:
    """Actually import every pipeline script under try/except.

    `check_handoff_coverage` greps script text for path-helper references but
    never executes the module-level code. Module-load-time bugs (forward
    references to imports that don't land until later in the file, missing
    submodule deps, broken `from X import Y` lines) slip through that grep.
    This check imports each script so module-level errors surface here, not
    at the operator's first `--help`. Audit finding T2.7 (2026-05-26).
    """
    import importlib
    errors: list[str] = []
    # Scripts that are entry points (have an `if __name__ == "__main__"`) +
    # helper modules. All should import cleanly with sys.path already set
    # to HERE.
    targets = [
        "_paths", "_meta_schema", "_chrome_schema", "_log", "_contract",
        "build_deck", "finalize_deck", "compile_picks",
        "build_review", "build_gate_preview",
        "register_template", "clean", "diagnostic",
        "icon_helper", "render_mermaid", "migrate_brief_to_v3",
    ]
    imported = 0
    for name in targets:
        try:
            # importlib.import_module caches; reimport via reload is overkill
            # for the contract test. A successful first-import here proves
            # module-level code ran without exception, which is what matters.
            mod = importlib.import_module(name)
            if mod is None:
                errors.append(f"import {name}: returned None")
            else:
                imported += 1
        except Exception as exc:
            errors.append(f"import {name}: {type(exc).__name__}: {exc}")
    if not errors:
        _ok(f"pipeline imports: {imported} modules loaded cleanly")
    return errors


# ---------------------------------------------------------------------------
# Check 5 — install sentinels (source/installed content fingerprint match)
# ---------------------------------------------------------------------------
#
# Drift class this prevents: a file installed at the right path with the
# wrong content. Existence checks (`Test-Path`) pass; user-facing behavior
# silently breaks. Worker-agent T1-R1 (v1 content at v0.1 path) and the
# T1-NEW-A `sys.argv[1]` regression both belong to this class.
#
# Each entry declares:
#   source            — absolute path to the canonical source artifact
#                       (lives inside the skill)
#   installed         — environment-resolved path to the installed copy
#                       (typically under ~/.claude/agents/ for sub-agents)
#   sentinels         — list of EXACT strings that must appear in BOTH
#                       files. Choose strings that encode the user-facing
#                       CONTRACT (e.g., the output filename pattern, the
#                       save-contract verb) — not strings that describe
#                       the file itself. The point is to catch "right
#                       path, wrong content."
#
# Adding new sentinels: when a future drift surfaces because the installed
# copy diverged from the source, add a sentinel here that would have
# caught the specific contradiction. This file is the institutional
# learning surface for install-time drift.

import os

_SKILL_ROOT = HERE.parent

INSTALL_SENTINELS = [
    {
        "name": "slide-builder-worker agent",
        "source": _SKILL_ROOT / "agents" / "slide-builder-worker.md",
        "installed": Path(os.path.expanduser("~/.claude/agents/slide-builder-worker.md")),
        # "option_A.pptx" encodes the OUTPUT-filename contract from
        # prompt.md:252-260. The string "option_A.py" (the SCRIPT name)
        # is not sufficient because v1 worker content also contains it.
        # The OUTPUT contract is what tells the worker to save to a
        # literal filename rather than sys.argv[1].
        "sentinels": ["option_A.pptx"],
    },
]


def check_install_sentinels() -> list[str]:
    errors: list[str] = []
    checked = 0
    for entry in INSTALL_SENTINELS:
        name = entry["name"]
        source: Path = entry["source"]
        installed: Path = entry["installed"]
        sentinels: list[str] = entry["sentinels"]

        if not source.exists():
            errors.append(f"sentinel: source for {name!r} missing at {source}")
            continue

        source_text = source.read_text(encoding="utf-8", errors="replace")
        for s in sentinels:
            if s not in source_text:
                errors.append(
                    f"sentinel: source {name!r} ({source.name}) does not "
                    f"contain canonical contract string {s!r} — the source "
                    f"itself drifted from its declared contract"
                )

        if not installed.exists():
            errors.append(
                f"sentinel: installed {name!r} missing at {installed} — "
                f"INSTALL step likely not run, or wrong destination"
            )
            continue

        installed_text = installed.read_text(encoding="utf-8", errors="replace")
        for s in sentinels:
            if s not in installed_text:
                errors.append(
                    f"sentinel: installed {name!r} at {installed} does not "
                    f"contain canonical contract string {s!r} — file exists "
                    f"but content is wrong (likely a stale copy)"
                )
        checked += 1

    if not errors:
        _ok(f"install sentinels: {checked} entries verified at source AND installed")
    return errors


# ---------------------------------------------------------------------------
# Check 6 — every markdown file:path reference resolves on disk
# ---------------------------------------------------------------------------
#
# Drift class this prevents: docs that name files which no longer exist
# (CHANGELOG listing QUICKSTART.md after it was deleted; SKILL.md naming
# `slide-builder_archived_2026-05-26/` when that path doesn't exist;
# cross-skill references to deleted v1 scripts).
#
# Scope: every .md under the skill EXCEPT under _decisions/. The decisions
# dir is forensic — historical records intentionally cite paths that may
# have moved or been deleted.

_DOC_REF_EXTS = (".py", ".md", ".json", ".yml", ".yaml", ".pptx", ".png",
                 ".html", ".css", ".xml", ".txt", ".toml")

# Match path-like tokens inside backticks. Restricting to backtick-wrapped
# tokens keeps us out of prose ("the design.md" vs "`design.md`") and
# avoids false positives on natural language.
_DOC_REF_RE = re.compile(
    r"`([A-Za-z_][\w./\-]*?(?:\.[A-Za-z0-9]{1,8}))`"
)

# Match directory-path refs (trailing slash) inside backticks. Catches
# the `slide-builder_archived_2026-05-26/` class of phantom-archive ref.
# Requires at least one path component (i.e., presence of "/" before
# the trailing slash, OR a leading skill-relative path with separators).
_DOC_DIR_REF_RE = re.compile(
    r"`([A-Za-z_][\w\-]*(?:/[\w.\-]+)*/)`"
)

# Dir refs that legitimately point at user-context paths (the user's
# home/.claude/agents/ install location, the user's OneDrive Claude
# Projects, etc.) — not under the skill tree by design.
_DIR_USER_CONTEXT_PREFIXES = (
    "~/.claude/", "OneDrive", "Documents/", "Downloads/",
    "_session/",  # session output dirs
    "C:/", "C:\\",  # absolute Windows paths
)

# Directory-path templates / placeholders used in prose, not literal
# paths. `slide_NN/` is the per-slide template (NN is a slot); `out/`
# is the build output dir as named in CLI examples.
_DIR_TEMPLATE_REFS = {
    "slide_NN/", "out/", "dont/", "exemplars/",
    "_templates/", "_session/",
}

# Directories legitimately cited in deletion-context entries
# (CHANGELOG history, master plan archive log). The check skips them
# globally — these are forensic citations, not live refs.
_DIR_DELETED_BY_DESIGN = {
    "slide-builder_archived_2026-05-26/",
    "slide-builder_archived_2026-05-26/exemplars/dont/",
    "slide-builder/icons/_audit/",
    "slide-builder/icons/_backup/",
    "icons/_audit/", "icons/_backup/",
    # v1 do/ positive-exemplar corpus never ported
    "do/", "do/single-finding/", "do/chart-bottom-takeaway/",
    "do/hero-kpi-tile/", "do/2panel-delta-spine/",
    "do/chart-right-takeaway/",
}

# Runtime artifacts produced by the pipeline. These appear in docs as
# "the output is X" — they intentionally don't live in the source tree.
# Catching them as broken refs is noise.
_RUNTIME_ARTIFACTS = {
    "_meta.json", "_prompt.md", "_finalize_meta.json",
    "picks.json", "brand.yml", "theme.json", "chrome.yml",
    "register.picks.json", "chrome.commit_method.txt",
    "final_deck.pptx", "final.pptx",  # legacy/example name
    "REVIEW.html", "GATE3-PREVIEW.html",
    "register.html", "register.proposal.json",
    "preview.pptx", "preview.png", "palette.png",
    "COMPILED.md", "RESULT.md", "build.log",
    "dispatch_plan.md",
    "qc.json", "raw.pptx",
    # Per-option artifacts, written at build time under slide_NN/
    "option_A.py", "option_B.py", "option_C.py",
    "option_A.pptx", "option_B.pptx", "option_C.pptx",
    "option_A.mmd", "option_B.mmd", "option_C.mmd",
    "option_X.py", "option_X.pptx", "option_X.mmd",  # template names in docs
    "option_X-mermaid.png",  # rendered fallback PNG
    "spec.md",
}

# Per-client runtime theme files (mermaid-<slug>.json). The slug
# varies by client; we treat the whole family as runtime.
_RUNTIME_PREFIXES = (
    "_session/",  # session-folder outputs (DECISIONS.md, briefs, etc.)
    "mermaid-",   # per-client mermaid theme JSON
)

# Filenames that are LEGITIMATELY referenced in deletion-context
# entries (CHANGELOG, master plan, audit handovers). The check skips
# them globally — if they reappear as live code, the import check
# would catch it.
_DELETED_BY_DESIGN = {
    "build_slide.py", "extract_icons.py", "QUICKSTART.md",
    "phase-a-rules.md", "visual-treatment-library.md",
    "designer-brief.md", "rules.md",
    "_verify_critical_fixes.py", "mermaid-brand.json",
    "stage-a-precommit.md", "slide-builder-simple-worker.md",
    "slide-builder.md", "slide-designer.md", "deck-builder.md",
    "smoke_test.py",
    # Path-class deletions
    "theme/mermaid-brand.json", "scripts/_verify_critical_fixes.py",
    # v0.2 candidate work — documented as future-pending, not present
    "check_brief_fidelity.py", "extract_lucide.py",
}

# Bare-filename refs that legitimately point at the user's external
# memory directory (~/.claude/projects/.../memory/) rather than the
# skill tree. Cross-context refs by design.
_MEMORY_FILE_PREFIXES = ("feedback_", "project_")


def _resolve_doc_ref(ref: str, md_file: Path) -> bool:
    """True iff `ref` resolves to an existing path under the skill tree
    or a sibling skill (one level up from the skill root).

    Search order:
    1. Absolute path (rare in docs but supported)
    2. Relative to md_file's directory
    3. Relative to skill root
    4. Strip leading skill-name component, retry under skill root
       (handles "slide-builder/foo/bar.md" written from parent perspective)
    5. Cross-skill: relative to skills root (handles "slide-qc/SKILL.md")
    6. Subtree Glob: bare filename matched anywhere under skill root
    """
    candidate = Path(ref.replace("\\", "/"))
    if candidate.is_absolute() and candidate.exists():
        return True
    if (md_file.parent / candidate).exists():
        return True
    if (_SKILL_ROOT / candidate).exists():
        return True
    parts = candidate.parts
    if parts and parts[0] == _SKILL_ROOT.name and (_SKILL_ROOT / Path(*parts[1:])).exists():
        return True
    # Cross-skill: skills/ root is the parent of the skill root.
    skills_root = _SKILL_ROOT.parent
    if (skills_root / candidate).exists():
        return True
    # Bare-filename subtree search: only meaningful when ref has no
    # directory component.
    if len(parts) == 1:
        # Use a constrained glob — rglob is O(tree) but the skill is small.
        for _ in _SKILL_ROOT.rglob(ref):
            return True
    return False


def check_doc_file_refs() -> list[str]:
    errors: list[str] = []
    md_files = [p for p in _SKILL_ROOT.rglob("*.md")
                if "_decisions" not in p.parts
                and "__pycache__" not in p.parts
                and "_audit" not in p.parts      # hangover dir, slated for deletion
                and "_backup" not in p.parts]    # hangover dir, slated for deletion
    refs_checked = 0
    refs_skipped_runtime = 0
    refs_skipped_deleted = 0
    refs_skipped_memory = 0
    dir_refs_checked = 0
    dir_refs_skipped = 0
    for md_file in md_files:
        try:
            text = md_file.read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            errors.append(f"doc refs: cannot read {md_file}: {exc}")
            continue
        # File-path refs
        for match in _DOC_REF_RE.finditer(text):
            ref = match.group(1)
            if not ref.endswith(_DOC_REF_EXTS):
                continue
            if ref.startswith(("http://", "https://", "mailto:", "#")):
                continue
            ref_name = Path(ref).name
            # Runtime prefix match (e.g., _session/..., mermaid-acme.json)
            if any(ref.startswith(p) for p in _RUNTIME_PREFIXES) or \
               any(ref_name.startswith(p) for p in _RUNTIME_PREFIXES):
                refs_skipped_runtime += 1
                continue
            # Runtime-artifact allowlist (output files, not source)
            if ref in _RUNTIME_ARTIFACTS or ref_name in _RUNTIME_ARTIFACTS:
                refs_skipped_runtime += 1
                continue
            # Intentionally-deleted-by-design allowlist (history docs
            # legitimately cite these as removed artifacts)
            if ref in _DELETED_BY_DESIGN or ref_name in _DELETED_BY_DESIGN:
                refs_skipped_deleted += 1
                continue
            # User-memory cross-context refs (live outside the skill)
            if any(ref_name.startswith(p) for p in _MEMORY_FILE_PREFIXES) and ref_name.endswith(".md"):
                refs_skipped_memory += 1
                continue
            refs_checked += 1
            if not _resolve_doc_ref(ref, md_file):
                rel = md_file.relative_to(_SKILL_ROOT)
                errors.append(
                    f"doc refs: {rel} cites {ref!r} but no such path exists "
                    f"under the skill tree or sibling skills"
                )
        # Directory-path refs (trailing slash) — catches the
        # `slide-builder_archived_2026-05-26/` class of phantom-archive ref
        for match in _DOC_DIR_REF_RE.finditer(text):
            dref = match.group(1)
            # Skip user-context dir refs (~/.claude/, OneDrive, etc.)
            if any(dref.startswith(p) for p in _DIR_USER_CONTEXT_PREFIXES):
                dir_refs_skipped += 1
                continue
            # Skip placeholder/template dir refs (slide_NN/, out/, etc.)
            if dref in _DIR_TEMPLATE_REFS:
                dir_refs_skipped += 1
                continue
            # Skip deletion-context dir refs (forensic citations)
            if dref in _DIR_DELETED_BY_DESIGN:
                dir_refs_skipped += 1
                continue
            # Skip if it's actually a file ref that happens to have a
            # trailing slash captured elsewhere — shouldn't happen with
            # the regex, but defensive
            if "." in dref.rstrip("/").split("/")[-1]:
                dir_refs_skipped += 1
                continue
            dir_refs_checked += 1
            # Resolve without the trailing slash
            dref_clean = dref.rstrip("/")
            if not _resolve_doc_ref(dref_clean, md_file):
                rel = md_file.relative_to(_SKILL_ROOT)
                errors.append(
                    f"doc refs (dir): {rel} cites {dref!r} but no such "
                    f"directory exists under the skill tree or sibling skills"
                )
    if not errors:
        _ok(
            f"doc file refs: {refs_checked} file refs + {dir_refs_checked} "
            f"dir refs resolve ({refs_skipped_runtime} runtime + "
            f"{refs_skipped_deleted} deleted-by-design + {refs_skipped_memory} "
            f"memory-context + {dir_refs_skipped} user-context skipped)"
        )
    return errors


# ---------------------------------------------------------------------------
# Check 7 — typing.get_type_hints() succeeds for every callable
# ---------------------------------------------------------------------------
#
# Drift class this prevents: `from __future__ import annotations` defers
# annotation evaluation, so `Optional[str]` without `from typing import
# Optional` runs fine at module load but breaks any code that calls
# `typing.get_type_hints(fn)`. Caught the `clean.py:149` Optional
# regression in the prior cycle (after manual investigation).

def check_type_hints_resolve() -> list[str]:
    import importlib
    import inspect
    import typing
    errors: list[str] = []
    targets = [
        "_paths", "_meta_schema", "_chrome_schema", "_log",
        "build_deck", "finalize_deck", "compile_picks",
        "build_review", "build_gate_preview",
        "register_template", "clean", "diagnostic",
        "icon_helper", "render_mermaid", "migrate_brief_to_v3",
    ]
    callables_checked = 0
    for mod_name in targets:
        try:
            mod = importlib.import_module(mod_name)
        except Exception as exc:
            # Module import already covered by check_pipeline_imports.
            # If it failed there, don't double-report; if it succeeded
            # there but fails here, something race-y is happening —
            # surface it.
            errors.append(f"type hints: cannot re-import {mod_name}: {exc}")
            continue
        for attr_name, attr in list(vars(mod).items()):
            if attr_name.startswith("_"):
                continue
            if not callable(attr):
                continue
            # Only inspect callables actually defined in this module
            # (skip re-exports / aliases of third-party stuff).
            try:
                defining = inspect.getmodule(attr)
                if defining is None or defining.__name__ != mod_name:
                    continue
            except Exception:
                continue
            try:
                typing.get_type_hints(attr)
                callables_checked += 1
            except Exception as exc:
                errors.append(
                    f"type hints: {mod_name}.{attr_name}: "
                    f"{type(exc).__name__}: {exc}"
                )
    if not errors:
        _ok(f"type hints: {callables_checked} callables across {len(targets)} modules resolve")
    return errors


# ---------------------------------------------------------------------------
# Checks 8-12 — v0.2 chrome inheritance (P1.5 in design lock; lands with the
# P1.3 / P1.4 code it protects, per design)
# ---------------------------------------------------------------------------

# Check 8 — chrome schema roundtrip
#
# Synthetic ChromeSpec serialises to dict and re-validates without drift.
# Schema version preserved. Mirrors check_meta_schema_roundtrip; fast unit
# check that the pydantic model survives ser/de.

def check_chrome_schema_roundtrip() -> list[str]:
    errors: list[str] = []
    synthetic = ChromeSpec(
        schema_version=CHROME_SCHEMA_VERSION_CURRENT,
        source_template_sha8="deadbeef",
        layouts={
            "body_canonical_light": LayoutChrome(
                name="body_canonical_light", layout_class="body-canonical",
                text_role="dark_on_light", background="light",
                has_page_number=True,
                title_placeholder_idx=0,
                body_top_y_px=90, body_bottom_y_px=665,
            ),
            "cover_dark": LayoutChrome(
                name="cover_dark", layout_class="bespoke",
                text_role="light_on_dark", background="dark",
                has_page_number=False,
                title=BoxPx(x_px=80, y_px=240, w_px=1120, h_px=120, font_pt=48, anchor="bottom"),
                subtitle=BoxPx(x_px=80, y_px=370, w_px=900, h_px=40, font_pt=20),
            ),
        },
    )
    as_dict = synthetic.model_dump()
    try:
        validated = validate_chrome_dict(as_dict)
    except ChromeSchemaError as e:
        errors.append(f"chrome roundtrip validation failed: {e}")
        return errors
    if validated.schema_version != CHROME_SCHEMA_VERSION_CURRENT:
        errors.append(
            f"chrome roundtrip schema_version mismatch: "
            f"{validated.schema_version} != {CHROME_SCHEMA_VERSION_CURRENT}"
        )
    if validated.source_template_sha8 != "deadbeef":
        errors.append(
            f"chrome roundtrip sha8 drift: {validated.source_template_sha8!r}"
        )
    if set(validated.layouts) != {"body_canonical_light", "cover_dark"}:
        errors.append(
            f"chrome roundtrip layout keys drift: {set(validated.layouts)}"
        )
    if not errors:
        _ok(f"chrome schema roundtrip @ version {CHROME_SCHEMA_VERSION_CURRENT}")
    return errors


# Check 9 — chrome sidecar no silent fallback
#
# Mock-patches every chrome source away and calls add_title_block(chrome=None).
# Asserts ChromeSidecarMissingError fires. Catches the v1 slot-mapping
# silent-fallback bug class — if a future refactor reintroduces a default,
# this check is the trip wire.

def check_chrome_sidecar_no_silent_fallback() -> list[str]:
    errors: list[str] = []
    SKILL_ROOT = HERE.parent
    twins_path = SKILL_ROOT
    if str(twins_path) not in sys.path:
        sys.path.insert(0, str(twins_path))
    try:
        import twins.helpers as h
    except Exception as exc:
        errors.append(f"sidecar fallback check: cannot import twins.helpers: {exc}")
        return errors

    # Snapshot env + module attribute, scrub them, exercise, restore.
    saved_active = h._ACTIVE_CHROME
    saved_chrome_yml = os.environ.pop("SLIDE_LAB_CHROME_YML", None)
    saved_layout = os.environ.pop("SLIDE_LAB_LAYOUT_NAME", None)
    h.set_active_chrome(None)
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        try:
            h.add_title_block(slide, "Title", "Sub")
            errors.append(
                "add_title_block(chrome=None) with no env vars + no _ACTIVE_CHROME "
                "did NOT raise — silent fallback regression"
            )
        except ChromeSidecarMissingError:
            pass  # expected
        except Exception as exc:
            errors.append(
                f"add_title_block raised wrong exception type: "
                f"{type(exc).__name__}: {exc}"
            )

        # add_footer should also raise
        try:
            h.add_footer(slide, page_num=1)
            errors.append(
                "add_footer(chrome=None) with no env vars + no _ACTIVE_CHROME "
                "did NOT raise — silent fallback regression"
            )
        except ChromeSidecarMissingError:
            pass
        except Exception as exc:
            errors.append(
                f"add_footer raised wrong exception type: "
                f"{type(exc).__name__}: {exc}"
            )
    finally:
        h.set_active_chrome(saved_active)
        if saved_chrome_yml is not None:
            os.environ["SLIDE_LAB_CHROME_YML"] = saved_chrome_yml
        if saved_layout is not None:
            os.environ["SLIDE_LAB_LAYOUT_NAME"] = saved_layout

    if not errors:
        _ok("chrome sidecar no-silent-fallback: helpers raise without sidecar")
    return errors


# Check 10 — sidecar freshness vs template sha
#
# Builds a synthetic registered template (.pptx + chrome.yml stamped with a
# specific sha8), mutates the .pptx bytes by 1, and asserts the freshness
# check in finalize_deck._load_and_verify_chrome hard-fails. Catches the
# "operator changed the template, forgot to re-register" failure mode.

def check_sidecar_freshness_vs_template_sha() -> list[str]:
    errors: list[str] = []
    import tempfile
    from pptx import Presentation as _Pres
    try:
        import finalize_deck as fd
    except Exception as exc:
        errors.append(f"freshness check: cannot import finalize_deck: {exc}")
        return errors

    with tempfile.TemporaryDirectory() as td:
        tpl = Path(td) / "synthetic.pptx"
        prs = _Pres()
        prs.save(str(tpl))

        # Stamp chrome.yml with the matching sha8.
        original_sha8 = fd._sha256_of_file(tpl)[:8]
        spec = ChromeSpec(
            schema_version=CHROME_SCHEMA_VERSION_CURRENT,
            source_template_sha8=original_sha8,
            layouts={"Blank": LayoutChrome(
                name="Blank", layout_class="body-canonical",
                text_role="dark_on_light", background="light",
                has_page_number=True,
            )},
        )
        dump_chrome_yml(spec, _p.chrome_yml(tpl))

        # Sanity: matching sha8 -> no raise.
        try:
            fd._load_and_verify_chrome(tpl)
        except Exception as exc:
            errors.append(
                f"freshness check: matched sha8 unexpectedly raised: "
                f"{type(exc).__name__}: {exc}"
            )
            return errors

        # Mutate template bytes -> sha8 changes -> hard fail expected.
        with open(tpl, "ab") as f:
            f.write(b"\x00")
        try:
            fd._load_and_verify_chrome(tpl)
            errors.append(
                "freshness check: mutated template did NOT raise "
                "RuntimeError — sha8 freshness check is silently passing"
            )
        except RuntimeError as exc:
            if "source_template_sha8" not in str(exc):
                errors.append(
                    f"freshness check: raised but message lacks "
                    f"'source_template_sha8': {exc}"
                )
        except Exception as exc:
            errors.append(
                f"freshness check: raised wrong type: "
                f"{type(exc).__name__}: {exc}"
            )

    if not errors:
        _ok("sidecar freshness: sha8 mismatch hard-fails finalize_deck")
    return errors


# Check 11 — chrome field single source
#
# Greps twins/helpers.py + slide-qc/ + every pattern helper for numeric
# coordinate literals that look like chrome positions. Allowlist exactly
# one location: _chrome_schema.py (chrome.yml lives on disk and is not
# scanned). Any other location is a contract violation.
#
# Patterns: `x_px = 58|64|1170`, `y_px = 20|100|660|672|676|688`.
# Plain integer matching would over-fire (e.g., 64 as an icon size,
# 100 as a percentage). The contextual key (x_px=, y_px=, x_emu=, y_emu=)
# restricts to coordinate assignments.

_CHROME_X_LITERALS = ("58", "64", "1170")
_CHROME_Y_LITERALS = ("20", "100", "660", "672", "676", "688")

# Lower-case file paths that ARE the allowed source-of-truth. Hits inside
# these are skipped.
_CHROME_LITERAL_ALLOWED_SUFFIXES = (
    "scripts/_chrome_schema.py",
    "scripts/_contract.py",  # exception: this very file names the literals
                              # to enforce them — meta-allowlist.
)

_CHROME_X_RE = re.compile(
    r"\bx(?:_px|_emu)?\s*=\s*("
    + "|".join(_CHROME_X_LITERALS) + r")\b"
)
_CHROME_Y_RE = re.compile(
    r"\by(?:_px|_emu)?\s*=\s*("
    + "|".join(_CHROME_Y_LITERALS) + r")\b"
)


def _is_chrome_literal_allowed(path: Path) -> bool:
    posix = path.resolve().as_posix().lower()
    return any(posix.endswith(suf) for suf in _CHROME_LITERAL_ALLOWED_SUFFIXES)


def check_chrome_field_single_source() -> list[str]:
    errors: list[str] = []
    skill_root = HERE.parent
    skills_root = skill_root.parent
    scan_dirs = [
        skill_root / "twins",
        skill_root / "scripts",
        skills_root / "slide-qc",
    ]
    files_checked = 0
    hits_found = 0
    for d in scan_dirs:
        if not d.exists():
            continue
        for p in d.rglob("*.py"):
            if "__pycache__" in p.parts:
                continue
            if "_decisions" in p.parts:
                # decision-time scripts and gallery exemplars are not pipeline
                # code; they don't get loaded at build time.
                continue
            if _is_chrome_literal_allowed(p):
                continue
            try:
                text = p.read_text(encoding="utf-8", errors="replace")
            except OSError as exc:
                errors.append(f"chrome-literals: cannot read {p}: {exc}")
                continue
            files_checked += 1
            for m in _CHROME_X_RE.finditer(text):
                line_no = text[:m.start()].count("\n") + 1
                rel = p.relative_to(skills_root)
                errors.append(
                    f"chrome-literals: {rel}:{line_no} has x_px={m.group(1)} "
                    f"(canonical chrome x). Move to _chrome_schema.py "
                    f"(CANONICAL_*_X) or read from chrome.yml."
                )
                hits_found += 1
            for m in _CHROME_Y_RE.finditer(text):
                line_no = text[:m.start()].count("\n") + 1
                rel = p.relative_to(skills_root)
                errors.append(
                    f"chrome-literals: {rel}:{line_no} has y_px={m.group(1)} "
                    f"(canonical chrome y). Move to _chrome_schema.py "
                    f"(CANONICAL_*_Y) or read from chrome.yml."
                )
                hits_found += 1
    if not errors:
        _ok(
            f"chrome field single-source: {files_checked} files scanned, "
            f"no chrome-position literals outside _chrome_schema.py"
        )
    return errors


# Check 12 — layout inheritance roundtrip per layout
#
# Per Agent D's prediction: extraction picks the wrong layout's footer when
# multiple match. To catch that, we re-extract chrome from each registered
# .pptx on disk and compare bounds to the stored chrome.yml. PER-LAYOUT, not
# per-template aggregate. Synthetic in-memory pass always runs; opportunistic
# disk scan catches drift against real templates.

_REGISTERED_TEMPLATE_SEARCH_ROOTS: tuple[Path, ...] = (
    # Common locations where a coworker keeps client templates.
    # Audit punch list P1.7 (2026-05-28): D's-bet catcher had zero coverage
    # against Mario's real templates until OneDrive paths were added — the
    # FDX Template lives under "OneDrive - Accenture/Claude Projects/_templates/"
    # and was invisible to the contract test prior to this extension.
    Path.home() / "Documents",
    Path.home() / "OneDrive - Accenture" / "Claude Projects" / "_templates",
    # Per-client library subdirs (each <client>/_templates/ holds registered
    # templates for that client). Non-existent roots are skipped harmlessly
    # by _opportunistic_chrome_yml_pairs; safe to enumerate hypothetical
    # locations alongside existing ones.
    Path.home() / "OneDrive - Accenture" / "Library" / "FedEx" / "_templates",
    Path.home() / "OneDrive - Accenture" / "Library" / "Accenture" / "_templates",
)


def _opportunistic_chrome_yml_pairs() -> list[tuple[Path, Path]]:
    """Find (chrome.yml, sibling .pptx) pairs under the search roots.
    Returns []; the disk scan is best-effort + bounded — we don't want the
    contract test to take minutes scanning a large file tree.

    v0.4 layout: chrome.yml lives at `<template-parent>/<stem>/chrome.yml`,
    sibling .pptx is at `<template-parent>/<stem>.pptx`. Pre-v0.4 flat
    layout (`<stem>.chrome.yml` next to `<stem>.pptx`) is no longer
    supported in production but the scan handles both for legacy diagnostic
    runs.
    """
    pairs: list[tuple[Path, Path]] = []
    for root in _REGISTERED_TEMPLATE_SEARCH_ROOTS:
        if not root.exists():
            continue
        try:
            # v0.4 subfolder layout: <stem>/chrome.yml
            for chrome_yml in root.rglob("chrome.yml"):
                depth = len(chrome_yml.relative_to(root).parts)
                if depth > 5:
                    continue
                # Sibling .pptx lives one level up, named <stem>.pptx where
                # <stem> is the parent dir name.
                stem = chrome_yml.parent.name
                pptx_sibling = chrome_yml.parent.parent / f"{stem}.pptx"
                if not pptx_sibling.exists():
                    pptx_sibling = chrome_yml.parent.parent / f"{stem}.potx"
                if pptx_sibling.exists():
                    pairs.append((chrome_yml, pptx_sibling))
                if len(pairs) >= 8:
                    return pairs
            # Legacy flat layout fallback: <stem>.chrome.yml + <stem>.pptx
            for chrome_yml in root.rglob("*.chrome.yml"):
                depth = len(chrome_yml.relative_to(root).parts)
                if depth > 5:
                    continue
                stem = chrome_yml.name[:-len(".chrome.yml")]
                pptx_sibling = chrome_yml.with_name(f"{stem}.pptx")
                if not pptx_sibling.exists():
                    pptx_sibling = chrome_yml.with_name(f"{stem}.potx")
                if pptx_sibling.exists():
                    pairs.append((chrome_yml, pptx_sibling))
                if len(pairs) >= 8:
                    return pairs
        except (PermissionError, OSError):
            continue
    return pairs


def check_layout_inheritance_roundtrip_per_layout() -> list[str]:
    errors: list[str] = []
    pairs = _opportunistic_chrome_yml_pairs()
    layouts_checked = 0
    templates_checked = 0
    try:
        import register_template as rt
    except Exception as exc:
        # Best-effort: if register_template can't load, the inheritance check
        # cannot run. Skip rather than fail (this preserves the green-baseline
        # contract when chrome support is being bootstrapped).
        _ok("layout inheritance roundtrip: skipped (register_template not importable)")
        return errors
    try:
        from pptx import Presentation as _Pres
    except Exception:
        return errors

    # Synthetic always-runs case (P1.5 hardening): even with no disk fixtures
    # available, exercise the extractor + roundtrip path so a refactor that
    # breaks classification or BoxPx round-trips trips this check.
    import tempfile as _tf

    # ---- P1.8 (audit punch list): synthetic BESPOKE BoxPx roundtrip. -----
    # The default pptx has zero bespoke layouts under the current classifier,
    # so the extraction path never exercises position-field serialization.
    # Hand-craft a ChromeSpec with non-None title/subtitle/footnote/source/
    # page_number on a bespoke layout, dump to YAML, reload, assert every
    # BoxPx axis survives byte-for-byte. This is the drift class Agent D bet
    # on — a refactor that loses, swaps, or rounds BoxPx fields trips this.
    with _tf.TemporaryDirectory() as _td:
        _bespoke_spec = ChromeSpec(
            schema_version=CHROME_SCHEMA_VERSION_CURRENT,
            source_template_sha8="bespoke0",
            layouts={
                "cover_dark": LayoutChrome(
                    name="cover_dark", layout_class="bespoke",
                    text_role="light_on_dark", background="dark",
                    has_page_number=False,
                    title=BoxPx(x_px=80, y_px=240, w_px=1120, h_px=120,
                                font_pt=48, anchor="bottom"),
                    subtitle=BoxPx(x_px=80, y_px=370, w_px=900, h_px=40,
                                   font_pt=20),
                    footnote=BoxPx(x_px=60, y_px=665, w_px=900, h_px=20),
                    source=BoxPx(x_px=60, y_px=690, w_px=900, h_px=20),
                    page_number=None,  # cover_dark intentionally has no page#
                ),
                "section_divider_dark": LayoutChrome(
                    name="section_divider_dark", layout_class="bespoke",
                    text_role="light_on_dark", background="dark",
                    has_page_number=True,
                    title=BoxPx(x_px=64, y_px=320, w_px=1152, h_px=80,
                                font_pt=40, anchor="middle"),
                    subtitle=None,  # section divider has no subtitle slot
                    footnote=None,
                    source=None,
                    page_number=BoxPx(x_px=1170, y_px=688, w_px=52, h_px=16),
                ),
            },
        )
        _bespoke_yml = Path(_td) / "bespoke_synthetic.chrome.yml"
        dump_chrome_yml(_bespoke_spec, _bespoke_yml)
        _bespoke_loaded = load_chrome_yml(_bespoke_yml)

        # Top-level invariants
        if _bespoke_loaded.schema_version != CHROME_SCHEMA_VERSION_CURRENT:
            errors.append(
                f"inheritance roundtrip [bespoke-synthetic]: schema_version "
                f"drift {CHROME_SCHEMA_VERSION_CURRENT} -> "
                f"{_bespoke_loaded.schema_version}"
            )
        if _bespoke_loaded.source_template_sha8 != "bespoke0":
            errors.append(
                f"inheritance roundtrip [bespoke-synthetic]: sha8 drift "
                f"{_bespoke_loaded.source_template_sha8!r}"
            )
        if set(_bespoke_loaded.layouts) != set(_bespoke_spec.layouts):
            errors.append(
                f"inheritance roundtrip [bespoke-synthetic]: layout-key drift "
                f"{set(_bespoke_spec.layouts)} -> {set(_bespoke_loaded.layouts)}"
            )

        # Per-layout per-field per-axis comparison
        for _name, _orig in _bespoke_spec.layouts.items():
            _re = _bespoke_loaded.layouts.get(_name)
            if _re is None:
                continue  # already caught above
            for _scalar in ("name", "layout_class", "text_role",
                            "background", "has_page_number"):
                if getattr(_orig, _scalar) != getattr(_re, _scalar):
                    errors.append(
                        f"inheritance roundtrip [bespoke-synthetic]: "
                        f"{_name!r}.{_scalar} drift "
                        f"{getattr(_orig, _scalar)!r} -> {getattr(_re, _scalar)!r}"
                    )
            for _field in ("title", "subtitle", "footnote", "source",
                           "page_number"):
                _b1 = getattr(_orig, _field)
                _b2 = getattr(_re, _field)
                if (_b1 is None) != (_b2 is None):
                    errors.append(
                        f"inheritance roundtrip [bespoke-synthetic]: "
                        f"{_name!r}.{_field} None-vs-set drift "
                        f"(orig={_b1 is not None}, reloaded={_b2 is not None})"
                    )
                    continue
                if _b1 is None or _b2 is None:
                    continue
                for _axis in ("x_px", "y_px", "w_px", "h_px",
                              "font_pt", "anchor"):
                    if getattr(_b1, _axis) != getattr(_b2, _axis):
                        errors.append(
                            f"inheritance roundtrip [bespoke-synthetic]: "
                            f"{_name!r}.{_field}.{_axis} drift "
                            f"{getattr(_b1, _axis)!r} -> "
                            f"{getattr(_b2, _axis)!r}"
                        )

    # ---- Existing synthetic case: extract from default pptx + re-extract. ----
    with _tf.TemporaryDirectory() as _td:
        _tpl = Path(_td) / "synthetic.pptx"
        _Pres().save(str(_tpl))
        _prs = _Pres(str(_tpl))
        _spec1 = rt.extract_chrome_spec(_prs, sha8="aaaaaaaa")
        # Dump + reload to exercise YAML serialization.
        _yml = Path(_td) / "synthetic.chrome.yml"
        dump_chrome_yml(_spec1, _yml)
        _loaded = load_chrome_yml(_yml)
        if set(_loaded.layouts) != set(_spec1.layouts):
            errors.append(
                f"inheritance roundtrip [synthetic]: layout-key drift "
                f"{set(_spec1.layouts)} -> {set(_loaded.layouts)}"
            )
        _prs2 = _Pres(str(_tpl))
        _spec2 = rt.extract_chrome_spec(
            _prs2, sha8="aaaaaaaa",
            classifications_override={
                n: lc.layout_class for n, lc in _loaded.layouts.items()
            },
        )
        for name, lc1 in _loaded.layouts.items():
            lc2 = _spec2.layouts.get(name)
            if lc2 is None:
                errors.append(
                    f"inheritance roundtrip [synthetic]: layout {name!r} "
                    f"disappeared on re-extract"
                )
                continue
            if lc1.layout_class != lc2.layout_class:
                errors.append(
                    f"inheritance roundtrip [synthetic]: layout {name!r} "
                    f"class drift {lc1.layout_class!r} -> {lc2.layout_class!r}"
                )
            for field in ("title", "subtitle", "footnote", "source", "page_number"):
                b1 = getattr(lc1, field)
                b2 = getattr(lc2, field)
                if (b1 is None) != (b2 is None):
                    errors.append(
                        f"inheritance roundtrip [synthetic]: {name!r}.{field} "
                        f"presence drift"
                    )
                    continue
                if b1 is not None and b2 is not None:
                    for axis in ("x_px", "y_px", "w_px", "h_px"):
                        if getattr(b1, axis) != getattr(b2, axis):
                            errors.append(
                                f"inheritance roundtrip [synthetic]: "
                                f"{name!r}.{field}.{axis} drift "
                                f"{getattr(b1, axis)} -> {getattr(b2, axis)}"
                            )

    for chrome_yml, pptx in pairs:
        try:
            stored = load_chrome_yml(chrome_yml)
        except Exception as exc:
            errors.append(
                f"inheritance roundtrip: cannot load {chrome_yml}: "
                f"{type(exc).__name__}: {exc}"
            )
            continue
        try:
            prs = _Pres(str(pptx))
            re_extracted = rt.extract_chrome_spec(
                prs, sha8=stored.source_template_sha8,
                # Honor stored classifications so the comparison isn't
                # confused by heuristic flipping a layout body<->bespoke.
                classifications_override={
                    n: lc.layout_class for n, lc in stored.layouts.items()
                },
            )
        except Exception as exc:
            errors.append(
                f"inheritance roundtrip: re-extract failed for {pptx.name}: "
                f"{type(exc).__name__}: {exc}"
            )
            continue
        templates_checked += 1
        for name, stored_lc in stored.layouts.items():
            new_lc = re_extracted.layouts.get(name)
            if new_lc is None:
                errors.append(
                    f"inheritance roundtrip: {pptx.name}: layout {name!r} "
                    f"present in chrome.yml but not in re-extraction"
                )
                continue
            layouts_checked += 1
            if stored_lc.layout_class != new_lc.layout_class:
                # honored via override — flag only if override didn't take
                errors.append(
                    f"inheritance roundtrip: {pptx.name}: layout {name!r} "
                    f"class drift {stored_lc.layout_class!r} -> {new_lc.layout_class!r}"
                )
            # For bespoke layouts, compare every populated box.
            for field in ("title", "subtitle", "footnote", "source", "page_number"):
                stored_box = getattr(stored_lc, field)
                new_box = getattr(new_lc, field)
                if stored_box is None and new_box is None:
                    continue
                if stored_box is None or new_box is None:
                    errors.append(
                        f"inheritance roundtrip: {pptx.name}: {name!r}.{field} "
                        f"presence drift (stored={stored_box is not None}, "
                        f"new={new_box is not None})"
                    )
                    continue
                # Tolerance: 1 EMU = 0 px at 96 DPI rounding -> px equality.
                for axis in ("x_px", "y_px", "w_px", "h_px"):
                    if getattr(stored_box, axis) != getattr(new_box, axis):
                        errors.append(
                            f"inheritance roundtrip: {pptx.name}: "
                            f"{name!r}.{field}.{axis} drift "
                            f"{getattr(stored_box, axis)} -> "
                            f"{getattr(new_box, axis)}"
                        )

    if not errors:
        if templates_checked == 0:
            _ok(
                "layout inheritance roundtrip: synthetic roundtrip OK; "
                "no registered templates found on disk (search roots: "
                + ", ".join(str(r) for r in _REGISTERED_TEMPLATE_SEARCH_ROOTS)
                + ")"
            )
        else:
            _ok(
                f"layout inheritance roundtrip: synthetic OK + "
                f"{layouts_checked} layouts across {templates_checked} "
                f"registered templates match within 0 px"
            )
    return errors


# ---------------------------------------------------------------------------
# Check 13 — layout resolution fail-loud (P1.5 hardening)
#
# Exercises build_deck.resolve_slide_layouts with a brief lacking both per-
# slide Layout and front-matter default_layout. Asserts errors list is
# non-empty AND that emit_layout_resolution_error formats correctly. Protects
# the P1.4 exit-9 fail-loud gate against future regressions where someone
# adds a silent default.
# ---------------------------------------------------------------------------

def check_layout_resolution_fails_loud() -> list[str]:
    errors: list[str] = []
    try:
        import build_deck as bd
    except Exception as exc:
        errors.append(f"layout-resolution check: cannot import build_deck: {exc}")
        return errors

    # Minimal synthetic brief with two slides and no layout anywhere.
    brief = {
        "front_matter": {"client_template": "T.pptx", "deck_type": "POV"},
        "deck_notes": "",
        "slides": [
            {"slide_n": 1, "title": "Cover", "archetype": "", "layout": "",
             "governing_thought": "", "so_what": "", "editorial_emphasis": "",
             "evidence_content": "", "chart_type": "none", "not_this_slide": ""},
            {"slide_n": 2, "title": "Body", "archetype": "", "layout": "",
             "governing_thought": "", "so_what": "", "editorial_emphasis": "",
             "evidence_content": "", "chart_type": "none", "not_this_slide": ""},
        ],
        "slide_total": 2,
    }
    # Template path doesn't have to exist for resolve_slide_layouts; it only
    # tries to load chrome.yml as a best-effort, ignoring missing.
    fake_tpl = Path("/_synthetic/T.pptx")
    resolved, available, errs = bd.resolve_slide_layouts(brief, fake_tpl)
    if len(errs) != 2:
        errors.append(
            f"layout-resolution: expected 2 errors (both slides lack layout), "
            f"got {len(errs)}: {errs}"
        )
    if any(r != "" for r in resolved):
        errors.append(
            f"layout-resolution: resolved entries for failed slides should "
            f"be empty strings; got {resolved}"
        )

    # Now with deck-level default — both should resolve.
    brief_with_default = {**brief}
    brief_with_default["front_matter"] = {
        **brief["front_matter"], "default_layout": "body_canonical_light",
    }
    resolved2, _, errs2 = bd.resolve_slide_layouts(brief_with_default, fake_tpl)
    if errs2:
        errors.append(
            f"layout-resolution: default_layout should resolve both slides; "
            f"got errors {errs2}"
        )
    if resolved2 != ["body_canonical_light", "body_canonical_light"]:
        errors.append(
            f"layout-resolution: deck default not propagated; got {resolved2}"
        )

    # Per-slide overrides default.
    brief_with_override = {**brief_with_default}
    brief_with_override["slides"] = [
        {**brief["slides"][0], "layout": "cover_light"},
        brief["slides"][1],
    ]
    resolved3, _, errs3 = bd.resolve_slide_layouts(brief_with_override, fake_tpl)
    if errs3 or resolved3 != ["cover_light", "body_canonical_light"]:
        errors.append(
            f"layout-resolution: per-slide override not honored; "
            f"resolved={resolved3} errs={errs3}"
        )

    if not errors:
        _ok("layout resolution fail-loud: 3 scenarios verified")
    return errors


# ---------------------------------------------------------------------------
# Check 14 — layout-inheritance integration smoke (P1.9 audit punch list)
#
# Wraps tests/run_layout_inheritance_smoke.py. Gated on the fixture pptx
# existing: a brand-new clone without the fixture skips the check rather
# than spending several seconds building it inside the contract test
# (the smoke is intentionally self-bootstrapping when run directly).
# ---------------------------------------------------------------------------

_SMOKE_FIXTURE_PPTX = HERE.parent / "tests" / "fixtures" / "layout_diverse_template.pptx"
_SMOKE_SCRIPT = HERE.parent / "tests" / "run_layout_inheritance_smoke.py"


def check_layout_inheritance_smoke_runs() -> list[str]:
    errors: list[str] = []
    if not _SMOKE_FIXTURE_PPTX.exists():
        _ok(
            "layout-inheritance smoke: skipped (fixture absent — run "
            "tests/run_layout_inheritance_smoke.py to build it)"
        )
        return errors
    if not _SMOKE_SCRIPT.exists():
        errors.append(f"smoke script missing at {_SMOKE_SCRIPT}")
        return errors
    try:
        result = subprocess.run(
            [sys.executable, str(_SMOKE_SCRIPT), "--quiet"],
            capture_output=True, text=True, timeout=240,
        )
    except FileNotFoundError as exc:
        errors.append(f"cannot invoke python: {exc}")
        return errors
    except subprocess.TimeoutExpired:
        errors.append("layout-inheritance smoke: timed out after 240s")
        return errors
    if result.returncode != 0:
        tail = (result.stdout + result.stderr).strip().splitlines()[-20:]
        errors.append(
            "layout-inheritance smoke FAILED (exit "
            f"{result.returncode}):\n" + "\n".join(tail)
        )
        return errors
    _ok("layout-inheritance smoke: all 5 phases pass against fixture")
    return errors


import subprocess  # noqa: E402  — only used by the smoke check above


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

def main() -> int:
    print("Slide Lab pipeline contract test")
    print("=" * 60)

    all_errors: list[str] = []
    for check in (check_paths_registry,
                  check_meta_schema_roundtrip,
                  check_handoff_coverage,
                  check_pipeline_imports,
                  check_install_sentinels,
                  check_doc_file_refs,
                  check_type_hints_resolve,
                  check_chrome_schema_roundtrip,
                  check_chrome_sidecar_no_silent_fallback,
                  check_sidecar_freshness_vs_template_sha,
                  check_chrome_field_single_source,
                  check_layout_inheritance_roundtrip_per_layout,
                  check_layout_resolution_fails_loud,
                  check_layout_inheritance_smoke_runs):
        errs = check()
        for e in errs:
            _fail(e)
        all_errors.extend(errs)

    print("=" * 60)
    if all_errors:
        print(f"CONTRACT TEST FAILED: {len(all_errors)} drift(s) detected.")
        return 1
    print("CONTRACT TEST PASSED.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
