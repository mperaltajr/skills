"""Register a client PPTX template — one-time setup per template.

Produces per-template sidecars in a subfolder NEXT TO the .pptx (v0.4+):
  <template-parent>/<stem>/
    brand.yml             : human-editable, ~15 lines
    theme.json            : machine-generated, SHA-stamped audit
    chrome.yml            : per-layout chrome geometry
    preview.pptx / .png   : 3-surface registration preview
    palette.png           : swatch grid
    register.html         : interactive picker UI
    register.proposal.json: Phase-1 output
    register.picks.json   : user's confirmed picks (when used)
    thumbnails/           : per-layout PNG thumbnails

Four CLI subcommands:

  propose      (chat-driven flow, default for coworker setup)
  -----------------------------------------------------------------------
  py -3 scripts/register_template.py propose <template.pptx>
      Extracts theme XML colors + fonts. Best-guess primary/accent. Builds
      preview PPTX (3 surfaces). Renders to PNG. Writes a proposal JSON +
      preview + palette + register.html into <stem>/. NO interactive
      prompts, NO TTY gate, NO writes to brand.yml/theme.json.
      The orchestrator (chat layer) reads the proposal JSON, shows the
      register.html to the user, takes picks in chat, then invokes
      `commit` (browser available) or `commit-cli` (chat-only).

  commit       (chat-driven flow, picks via JSON file)
  -----------------------------------------------------------------------
  py -3 scripts/register_template.py commit <template.pptx> --picks <picks.json>
      Consumes a picks JSON with the user's chosen primary/accent (or
      {"accept": true} to use the Phase-1 best-guess unchanged) and writes
      <stem>/brand.yml + <stem>/theme.json + <stem>/chrome.yml. NO TTY gate.
      Use this when register.html is openable in a real browser and the
      chat can receive the live-edited picks JSON back.

  commit-cli   (chat-driven flow, picks via CLI flags — RC-2 fallback)
  -----------------------------------------------------------------------
  py -3 scripts/register_template.py commit-cli <template.pptx> \\
      --primary-slot dk2 --primary-hex 4D148C \\
      --accent-slot lt2  --accent-hex FF6600
      Same write side as `commit`, but takes picks from CLI flags instead
      of a JSON file. Use when the chat cannot render register.html with
      live JS (e.g., chat-only preview panels). The orchestrator shows
      palette.png in chat, takes user picks conversationally, then runs
      this. Optional flags: --cover-bg-slot/--cover-bg-hex,
      --dark-bg-slot/--dark-bg-hex, --strip-master-backgrounds,
      repeatable --layout-class "NAME=body-canonical|bespoke", --accept.

  interactive  (legacy PowerShell flow, power users with a real TTY)
  -----------------------------------------------------------------------
  py -3 scripts/register_template.py interactive <template.pptx>
  py -3 scripts/register_template.py interactive <template.pptx> --auto-accept-phase1
      The original 4-phase interactive flow with input() prompts and the
      TTY confirmation gate. For batch registration / migration / power
      users running from a real terminal. NOT for coworkers — see `propose`
      / `commit` / `commit-cli` for the chat-driven paths.

Picks JSON shape (consumed by `commit`):

  {
    "accept": false,                 # if true, ignore the color fields below
    "primary_slot": "dk2",           # one of the theme slot names
    "primary_hex":  "4D148C",        # 6-char hex without leading #
    "accent_slot":  "lt2",
    "accent_hex":   "FF6600",
    "cover_bg_slot": "dk2",          # optional; defaults to primary
    "cover_bg_hex":  "4D148C",
    "strip_master_backgrounds": false, # optional; default false (keep master)
    "default_content_layout": "2_Title & Text 01",  # optional but recommended;
                                     # the layout the user picked as the
                                     # default for content slides. Validated
                                     # against chrome.yml at commit time.
                                     # Read by build_deck.py as the template-
                                     # level fallback before sole-body-
                                     # canonical auto-detection.
    "layout_classifications": {      # optional; per-layout class override
      "Cover": "bespoke",
      "1_Title & Text 01": "body-canonical"
    }
  }

Outputs (all subcommands): absolute paths printed to stdout for each
file written, plus exit code 0 on success.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import sys
import tempfile
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, Optional

SKILL_ROOT = Path(__file__).resolve().parents[1]
QC_SCRIPTS = SKILL_ROOT.parent / "slide-qc" / "scripts"
sys.path.insert(0, str(SKILL_ROOT))
sys.path.insert(0, str(QC_SCRIPTS))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

import _paths as _p  # noqa: E402
from _chrome_schema import (  # noqa: E402
    BoxPx, LayoutChrome, ChromeSpec,
    CHROME_SCHEMA_VERSION_CURRENT,
    CANONICAL_BODY_TOP_Y, CANONICAL_BODY_BOTTOM_Y,
    dump_chrome_yml, load_chrome_yml,
)

try:
    import yaml
except ImportError:
    print("ERROR: PyYAML not installed. pip install pyyaml")
    sys.exit(2)


# ---------------------------------------------------------------------------
# Theme1.xml parsing — minimal, no external deps
# ---------------------------------------------------------------------------
_COLOR_SLOTS = ("dk1", "lt1", "dk2", "lt2",
                "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                "hlink", "folHlink")


def _read_theme_xml(template_path: Path) -> str:
    """Extract the first theme1.xml from a .pptx/.potx zip."""
    with zipfile.ZipFile(str(template_path), "r") as zf:
        for name in zf.namelist():
            ln = name.lower()
            if "/theme/theme" in ln and ln.endswith(".xml"):
                return zf.read(name).decode("utf-8", errors="replace")
    raise ValueError(f"No theme XML found in {template_path}")


def _extract_color(theme_xml: str, slot: str) -> str:
    """Pull a 6-char uppercase hex for the given slot name, or '000000'."""
    m = re.search(rf"<a:{slot}>(.*?)</a:{slot}>", theme_xml, re.DOTALL)
    if not m:
        return "000000"
    inner = m.group(1)
    srgb = re.search(r'srgbClr\s+val="([0-9A-Fa-f]{6})"', inner)
    sysc = re.search(r'sysClr[^/>]*lastClr="([0-9A-Fa-f]{6})"', inner)
    if srgb:
        return srgb.group(1).upper()
    if sysc:
        return sysc.group(1).upper()
    return "000000"


_FONT_STYLE_SUFFIXES = (
    " Regular", " Bold", " Italic", " Bold Italic", " Light", " Medium",
    " SemiBold", " ExtraBold", " Black", " Thin", " ExtraLight", " Heavy",
    " Oblique", " Light Italic", " Medium Italic", " SemiBold Italic",
    " ExtraBold Italic", " Black Italic",
)


def _strip_font_style(name: str) -> str:
    n = name.strip()
    for s in _FONT_STYLE_SUFFIXES:
        if n.endswith(s):
            return n[:-len(s)].strip()
    return n


def _extract_font(theme_xml: str, kind: str) -> str:
    """kind = 'majorFont' or 'minorFont'."""
    m = re.search(rf"<a:{kind}>(.*?)</a:{kind}>", theme_xml, re.DOTALL)
    if not m:
        return ""
    latin = re.search(r'<a:latin\s+typeface="([^"]+)"', m.group(1))
    return _strip_font_style(latin.group(1)) if latin else ""


# ---------------------------------------------------------------------------
# Color math (HSV) — used to best-guess primary/accent
# ---------------------------------------------------------------------------
def _hex_to_hsv(hex_color: str) -> tuple:
    r, g, b = [int(hex_color[i:i + 2], 16) / 255.0 for i in (0, 2, 4)]
    cmax, cmin = max(r, g, b), min(r, g, b)
    delta = cmax - cmin
    V = cmax
    S = (delta / cmax) if cmax > 0 else 0.0
    if delta == 0:
        H = 0.0
    elif cmax == r:
        H = ((g - b) / delta) % 6
    elif cmax == g:
        H = (b - r) / delta + 2
    else:
        H = (r - g) / delta + 4
    return (H * 60, S, V)


def _is_branded(hex_color: str) -> bool:
    """Saturated + not too dark/light — a candidate brand color."""
    try:
        _, s, v = _hex_to_hsv(hex_color)
    except Exception:
        return False
    return s >= 0.30 and v >= 0.20


# ---------------------------------------------------------------------------
# Fix #1 — TTY read helper. Bypasses piped stdin to require a keystroke at
# the controlling terminal. Used to gate the disk-write commit.
# ---------------------------------------------------------------------------
def _tty_input(prompt: str) -> str:
    """Read one line from the controlling terminal regardless of stdin redir.

    Raises RuntimeError if no controlling terminal exists (CI / headless).
    """
    try:
        if sys.platform == "win32":
            tty = open("CONIN$", "r")
        else:
            tty = open("/dev/tty", "r")
    except OSError as e:
        raise RuntimeError(
            "no controlling terminal available — Y-gate cannot be confirmed "
            "(re-run from an interactive terminal)"
        ) from e
    try:
        sys.stdout.write(prompt); sys.stdout.flush()
        line = tty.readline()
        return line.rstrip("\r\n")
    finally:
        tty.close()


# ---------------------------------------------------------------------------
# Fix #2 — Visual-similarity detector. Flags pairs of candidates that share
# a hue family and saturation band, prompting a typed-hex confirm step in
# Phase 3.
# ---------------------------------------------------------------------------
def _is_visually_similar(hex_a: str, hex_b: str,
                         hue_thresh: float = 30.0,
                         sat_thresh: float = 0.20,
                         val_thresh: float = 0.25) -> bool:
    """True if two hexes are close enough in HSV that swatch labels alone may
    not visually disambiguate them.

    Defaults are tuned for the ACN case: #460073 (deep purple) vs #A100FF
    (bright purple) — same hue family (~270°), saturations 1.0 vs 1.0,
    values 0.45 vs 1.0. Hue_diff ~ 6° → similar. Value_diff 0.55 → still
    flagged since both purple. We trip on hue+sat ONLY (value can vary).
    """
    try:
        ha, sa, va = _hex_to_hsv(hex_a)
        hb, sb, vb = _hex_to_hsv(hex_b)
    except Exception:
        return False
    dh = abs(ha - hb) % 360
    if dh > 180:
        dh = 360 - dh
    return (dh < hue_thresh) and (abs(sa - sb) < sat_thresh)


def _ambiguous_neighbors(hex_picked: str,
                         entries: list,
                         exclude_slot: Optional[str] = None) -> list:
    """Return entries from `entries` that are visually similar to hex_picked.

    `entries` is a list of (slot, hex) tuples (from show_palette).
    `exclude_slot` skips a slot from the result (used to skip the picked entry
    itself when entries contains duplicate hexes under different slots).
    """
    similars = []
    for slot, hexv in entries:
        if hexv.upper() == hex_picked.upper():
            continue
        if exclude_slot and slot == exclude_slot:
            continue
        if _is_visually_similar(hex_picked, hexv):
            similars.append((slot, hexv.upper()))
    return similars


def _confirm_hex_if_ambiguous(label: str, slot: str, hexv: str,
                              entries: list) -> bool:
    """If the picked color has visually-similar siblings in the palette, force
    the user to retype the 6-char hex. Returns True if confirmed, False if
    mismatch (caller re-runs the picker).
    """
    similars = _ambiguous_neighbors(hexv, entries, exclude_slot=slot)
    if not similars:
        return True
    print()
    print(f"  {label}: slot={slot} hex=#{hexv.upper()}")
    print(f"  Visually similar in palette: "
          + ", ".join(f"#{h} ({s})" for s, h in similars))
    print(f"  Re-type the 6-char hex (no leading #) to confirm:")
    typed = input(f"    {label} hex: ").strip().upper().lstrip("#")
    if typed != hexv.upper():
        print(f"  Mismatch: typed #{typed!r} vs picked #{hexv.upper()} — "
              f"picker will re-run.")
        return False
    print(f"  Confirmed: {label} = #{hexv.upper()}")
    return True


def _count_color_uses_in_masters(template_path: Path, hex_codes: list) -> Dict[str, int]:
    """Count how many times each hex appears in slideMaster XMLs + theme1.xml.

    Used to bias best-guess: a color used heavily across the masters is more
    likely to be brand chrome.
    """
    counts: Dict[str, int] = {h.upper(): 0 for h in hex_codes}
    try:
        with zipfile.ZipFile(str(template_path), "r") as zf:
            for name in zf.namelist():
                ln = name.lower()
                if not (ln.startswith("ppt/slidemasters/") or "/theme/theme" in ln):
                    continue
                if not ln.endswith(".xml"):
                    continue
                try:
                    data = zf.read(name).decode("utf-8", errors="ignore").upper()
                except Exception:
                    continue
                for h in counts:
                    counts[h] += data.count(f'VAL="{h}"')
    except Exception:
        pass
    return counts


def _best_guess_primary_accent(colors: Dict[str, str], template_path: Path) -> tuple:
    """Pick (primary_hex, primary_slot, accent_hex, accent_slot) using:

      1. Filter to "branded" candidates (S>=0.30, V>=0.20).
      2. Score each by (saturation * 100) + (master-XML use-count).
      3. Primary = highest-scoring. Accent = next-highest with different hue.

    If nothing passes the branded gate, fall back to dk2/lt2 raw.
    """
    candidates = []
    for slot in ("dk2", "lt2", "accent1", "accent2", "accent3",
                 "accent4", "accent5", "accent6"):
        v = colors.get(slot, "")
        if not v or len(v) != 6:
            continue
        if _is_branded(v):
            candidates.append((slot, v))

    if not candidates:
        return (colors.get("dk2", "000000"), "dk2",
                colors.get("lt2", "FFFFFF"), "lt2")

    # Count uses across masters
    use_counts = _count_color_uses_in_masters(
        template_path, [v for _, v in candidates]
    )

    scored = []
    for slot, hexv in candidates:
        _, s, _ = _hex_to_hsv(hexv)
        use = use_counts.get(hexv, 0)
        # Saturation dominates; usage breaks ties (a single point per usage is
        # plenty since saturation is in [0,100] after the *100 scale).
        score = s * 100 + use * 0.5
        scored.append((score, slot, hexv))
    scored.sort(reverse=True)

    primary_slot, primary_hex = scored[0][1], scored[0][2]

    # Accent: next branded candidate with a different hue (>= 30deg apart) and
    # a different hex. If no hue-distinct candidate exists, take next-by-score.
    primary_h, _, _ = _hex_to_hsv(primary_hex)
    accent_slot, accent_hex = None, None
    for _, slot, hexv in scored[1:]:
        if hexv == primary_hex:
            continue
        h2, _, _ = _hex_to_hsv(hexv)
        dh = abs(primary_h - h2) % 360
        if dh > 180:
            dh = 360 - dh
        if dh >= 30:
            accent_slot, accent_hex = slot, hexv
            break
    if accent_hex is None and len(scored) > 1:
        # Fallback — second by score, even if same hue family
        accent_slot, accent_hex = scored[1][1], scored[1][2]
    if accent_hex is None:
        accent_slot, accent_hex = "lt2", colors.get("lt2", "FFFFFF")

    return primary_hex, primary_slot, accent_hex, accent_slot


# ---------------------------------------------------------------------------
# Preview slide builder
# ---------------------------------------------------------------------------
def _rgb(hexs: str) -> RGBColor:
    return RGBColor(int(hexs[0:2], 16), int(hexs[2:4], 16), int(hexs[4:6], 16))


def _clear_existing_slides(prs) -> None:
    """Remove all slides from the presentation (keep masters/layouts).

    python-pptx has no public `remove_slide`; we drop the relationship via
    `part.drop_rel(rId)` AND remove the entry from sldIdLst. Without dropping
    the rel, the old slide parts remain referenced and python-pptx writes
    them back into the saved zip — which collides with new slides that take
    the same /ppt/slides/slideN.xml names. LibreOffice then errors on the
    duplicate zip entries.
    """
    R_ID = qn("r:id")
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(R_ID)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def _pick_blank_layout(prs):
    """Pick the best body-style layout for preview rendering.

    The preview surfaces (color swatches, KPI tiles, bullet bars) need to
    land somewhere LEGIBLE. The cover / title layout typically ships with
    art-direction baked into the master (FedEx's huge X graphic, NFL
    stadium photos, Accenture branding) that swamps anything drawn on top.
    Picking a body-canonical-style layout puts the preview where the
    eventual deck body slides will live — which is also where master
    decoration is lightest.

    Priority:
      1. A layout literally named "Blank" (case-insensitive, after
         stripping numeric prefixes like "1_" or "12__"), searched across
         EVERY master — `prs.slide_layouts` only returns master[0]'s
         layouts, but FedEx-style templates put the Blank layout on a
         later master.
      2. Otherwise the layout with the FEWEST non-placeholder shapes —
         that's the closest analogue to a body-canonical layout when no
         explicit Blank exists.
      3. Last resort: the first layout in the presentation.
    """
    import re as _re
    # Pass 1: find "blank" by name across every master.
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            n = (layout.name or "").strip().lower()
            n = _re.sub(r"^[\d_]+", "", n)  # strip "1_", "12__"
            if n == "blank":
                return layout
    # Pass 2: prefer the layout with the fewest non-placeholder shapes.
    best_layout = None
    best_count = None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            try:
                count = sum(
                    1 for shp in layout.shapes
                    if not getattr(shp, "is_placeholder", False)
                )
            except Exception:
                continue
            if best_count is None or count < best_count:
                best_count = count
                best_layout = layout
    if best_layout is not None:
        return best_layout
    # Pass 3: legacy fallback.
    if prs.slide_masters and list(prs.slide_masters[0].slide_layouts):
        return prs.slide_masters[0].slide_layouts[0]
    layouts = list(prs.slide_layouts)
    return layouts[0] if layouts else None


def _add_text(slide, x_px, y_px, w_px, h_px, *, text, size_pt, bold=False,
              color=(0x1A, 0x1A, 0x2E), font_name=None, align=PP_ALIGN.LEFT,
              anchor=None):
    """Helper: add a text-bearing rectangle (no fill, no line)."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(x_px * 9525), Emu(y_px * 9525),
        Emu(w_px * 9525), Emu(h_px * 9525),
    )
    box.fill.background(); box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    tf.margin_right = Emu(0); tf.margin_bottom = Emu(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size_pt); r.font.bold = bold
    r.font.color.rgb = RGBColor(*color)
    if font_name:
        r.font.name = font_name
    return box


def _add_fill(slide, x_px, y_px, w_px, h_px, hex_color):
    """Helper: add a filled rectangle with no border."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(x_px * 9525), Emu(y_px * 9525),
        Emu(w_px * 9525), Emu(h_px * 9525),
    )
    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(hex_color)
    shp.line.fill.background()
    return shp


def _surface_title(slide, primary_hex, accent_hex, cover_bg_hex,
                   font_heading, font_body):
    """Surface 1 — title bar (primary) + body + accent stripe + cover-bg card
    + swatch row. NO white background — master decoration shows through."""
    _add_fill(slide, 40, 40, 1200, 96, primary_hex)
    _add_text(slide, 56, 48, 1170, 80,
              text="Brand preview — title using PRIMARY",
              size_pt=28, bold=True, color=(0xFF, 0xFF, 0xFF),
              font_name=font_heading, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, 40, 156, 760, 90,
              text=("Body paragraph in the body font. The accent stripe below "
                    "uses ACCENT. The card uses COVER-BG. Master decoration "
                    "(if any) is visible behind this slide."),
              size_pt=14, color=(0x1A, 0x1A, 0x2E), font_name=font_body)
    _add_fill(slide, 40, 270, 1200, 12, accent_hex)  # accent stripe
    _add_fill(slide, 40, 314, 560, 280, cover_bg_hex)  # cover-bg card
    _add_text(slide, 60, 360, 520, 60,
              text="COVER-BG card",
              size_pt=22, bold=True, color=(0xFF, 0xFF, 0xFF),
              font_name=font_heading)
    _add_text(slide, 60, 420, 520, 140,
              text=("This box uses cover_bg_hex with white text — what the "
                    "cover slide will look like."),
              size_pt=13, color=(0xFF, 0xFF, 0xFF), font_name=font_body)
    # Swatch row (right)
    for i, (label, hexv) in enumerate(
        [("PRIMARY", primary_hex), ("ACCENT", accent_hex),
         ("COVER-BG", cover_bg_hex)]
    ):
        x = 640 + i * 200
        _add_fill(slide, x, 314, 180, 140, hexv)
        _add_text(slide, x, 460, 180, 22, text=label, size_pt=11, bold=True,
                  color=(0x33, 0x33, 0x33), font_name=font_body,
                  align=PP_ALIGN.CENTER)
        _add_text(slide, x, 484, 180, 20, text=f"#{hexv.upper()}", size_pt=10,
                  color=(0x66, 0x66, 0x66), font_name=font_body,
                  align=PP_ALIGN.CENTER)
    _add_text(slide, 40, 660, 1200, 30,
              text=("preview surface 1/3 — title · master decoration visible "
                    "behind this slide if present"),
              size_pt=9, color=(0x88, 0x88, 0x88), font_name=font_body)


def _surface_content(slide, primary_hex, accent_hex, secondaries,
                     font_heading, font_body):
    """Surface 2 — content body. Sub-header bar (primary), bullets with one
    accent-emphasized word, sidebar KPI box, secondary-color swatch row to
    exercise accent3/4/5/6."""
    _add_fill(slide, 40, 40, 800, 56, primary_hex)
    _add_text(slide, 56, 48, 770, 40,
              text="Content surface — bullets + sidebar KPI",
              size_pt=18, bold=True, color=(0xFF, 0xFF, 0xFF),
              font_name=font_heading, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, 40, 120, 700, 30,
              text="Sub-headline in body font with one ACCENT word inline.",
              size_pt=14, color=(0x1A, 0x1A, 0x2E), font_name=font_body)
    # Bullet list
    bullet_y = 170
    bullets = [
        "First bullet — body font, dk1 text, default size.",
        "Second bullet — exercises body-font line wrap behavior.",
        "Third bullet — checks fallback if FONT is unavailable.",
        "Fourth bullet — long enough to test word wrap at the right edge.",
    ]
    for i, b in enumerate(bullets):
        # Bullet dot in accent
        _add_fill(slide, 50, bullet_y + i * 40 + 8, 8, 8, accent_hex)
        _add_text(slide, 70, bullet_y + i * 40, 700, 36,
                  text=b, size_pt=13, color=(0x1A, 0x1A, 0x2E),
                  font_name=font_body)
    # Sidebar KPI box
    _add_fill(slide, 880, 120, 360, 220, primary_hex)
    _add_text(slide, 900, 140, 320, 50,
              text="KPI", size_pt=14, bold=True,
              color=(0xFF, 0xFF, 0xFF), font_name=font_body)
    _add_text(slide, 900, 178, 320, 110,
              text="73%", size_pt=72, bold=True,
              color=(0xFF, 0xFF, 0xFF), font_name=font_heading)
    _add_text(slide, 900, 290, 320, 40,
              text="Pipeline coverage — last quarter",
              size_pt=11, color=(0xFF, 0xFF, 0xFF), font_name=font_body)
    # Secondary swatch row — exercises chart-series colors
    _add_text(slide, 40, 410, 800, 24,
              text="Secondary accents (chart series 3-6):",
              size_pt=11, bold=True, color=(0x55, 0x55, 0x55),
              font_name=font_body)
    for i, (slot, hexv) in enumerate(secondaries[:6]):
        x = 40 + i * 175
        _add_fill(slide, x, 440, 160, 80, hexv)
        _add_text(slide, x, 524, 160, 18, text=f"{slot}",
                  size_pt=10, bold=True, color=(0x33, 0x33, 0x33),
                  font_name=font_body, align=PP_ALIGN.CENTER)
        _add_text(slide, x, 544, 160, 18, text=f"#{hexv.upper()}",
                  size_pt=9, color=(0x66, 0x66, 0x66), font_name=font_body,
                  align=PP_ALIGN.CENTER)
    _add_text(slide, 40, 660, 1200, 30,
              text=("preview surface 2/3 — content · body font wrap + secondary "
                    "accents · master decoration visible behind"),
              size_pt=9, color=(0x88, 0x88, 0x88), font_name=font_body)


def _surface_dashboard(slide, primary_hex, accent_hex, secondaries,
                       font_heading, font_body):
    """Surface 3 — dashboard. 2×2 KPI tile grid using primary + 3 secondary
    accents. Largest single use of brand chrome — best surface for spotting
    photographic bleed-through under strip_master_backgrounds=false."""
    _add_text(slide, 40, 30, 1200, 50,
              text="Dashboard surface — 2×2 KPI tiles",
              size_pt=22, bold=True, color=(0x1A, 0x1A, 0x2E),
              font_name=font_heading)
    # Choose 4 tile colors: primary, then up to 3 secondaries (fall back to
    # accent_hex if not enough secondaries)
    fills = [primary_hex]
    for slot, hexv in secondaries[:3]:
        fills.append(hexv)
    while len(fills) < 4:
        fills.append(accent_hex)
    tile_w, tile_h, gap = 580, 240, 20
    coords = [(40, 100), (40 + tile_w + gap, 100),
              (40, 100 + tile_h + gap), (40 + tile_w + gap, 100 + tile_h + gap)]
    labels = [("Revenue", "$4.2B"), ("Margin", "18.7%"),
              ("Pipeline", "$1.1B"), ("Win rate", "62%")]
    for (lx, ly), color, (lbl, val) in zip(coords, fills, labels):
        _add_fill(slide, lx, ly, tile_w, tile_h, color)
        _add_text(slide, lx + 20, ly + 16, tile_w - 40, 28,
                  text=lbl.upper(), size_pt=12, bold=True,
                  color=(0xFF, 0xFF, 0xFF), font_name=font_body)
        _add_text(slide, lx + 20, ly + 50, tile_w - 40, 140,
                  text=val, size_pt=64, bold=True,
                  color=(0xFF, 0xFF, 0xFF), font_name=font_heading)
    _add_text(slide, 40, 660, 1200, 30,
              text=("preview surface 3/3 — dashboard · widest brand-chrome "
                    "coverage; bleed-through (if any) most visible here"),
              size_pt=9, color=(0x88, 0x88, 0x88), font_name=font_body)


def build_preview_pptx(template_path: Path, primary_hex: str, accent_hex: str,
                     cover_bg_hex: str, font_heading: str, font_body: str,
                     all_colors: dict, out_path: Path) -> None:
    """Build a 3-surface preview deck OPENED AGAINST the client template.

    Surfaces:
      1. Title — brand chrome + swatch row
      2. Content — bullets + sidebar KPI + secondary-accent swatches
      3. Dashboard — 2×2 KPI grid (widest brand coverage)

    The deck is built on top of the client template's master/layouts, so
    photographic decoration baked into the master (NFL stadium photos, etc.)
    bleeds through behind every preview slide. This is what makes
    strip_master_backgrounds=true visibly necessary at registration time.

    `all_colors` is the full theme palette dict (dk1..folHlink) — surface 2
    uses accent3-6 to exercise chart-series secondaries.
    """
    try:
        prs = Presentation(str(template_path))
    except Exception as e:
        print(f"  WARNING: could not open template ({e}); falling back to blank.")
        prs = Presentation()
        prs.slide_width = Emu(1280 * 9525)
        prs.slide_height = Emu(720 * 9525)

    _clear_existing_slides(prs)
    blank = _pick_blank_layout(prs)
    if blank is None:
        raise RuntimeError("template has no usable slide layout")

    # Collect "secondary" colors for chart-series exercise: accent3..accent6
    # (skip accent1/accent2 since those often duplicate primary/accent variants)
    secondaries = []
    for slot in ("accent3", "accent4", "accent5", "accent6"):
        hexv = all_colors.get(slot, "")
        if hexv and len(hexv) == 6:
            secondaries.append((slot, hexv))

    _surface_title(prs.slides.add_slide(blank), primary_hex, accent_hex,
                   cover_bg_hex, font_heading, font_body)
    _surface_content(prs.slides.add_slide(blank), primary_hex, accent_hex,
                     secondaries, font_heading, font_body)
    _surface_dashboard(prs.slides.add_slide(blank), primary_hex, accent_hex,
                       secondaries, font_heading, font_body)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def render_pptx_to_png(pptx_path: Path, png_path: Path) -> bool:
    """Render the multi-surface preview PPTX and composite into one tall PNG.

    Reuses render_libre from slide-qc. If the PPTX has >1 slide, the surfaces
    are stacked vertically (gap of 20px) into a single PNG so Mario can review
    title + content + dashboard in one image.
    """
    try:
        from render_slides import render_libre
    except ImportError as e:
        print(f"WARNING: could not import render_libre: {e}")
        return False
    try:
        from PIL import Image
    except ImportError as e:
        print(f"WARNING: could not import PIL: {e}")
        return False
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        try:
            render_libre(pptx_path, tmp, dpi=120)
        except Exception as e:
            print(f"WARNING: render_libre failed: {e}")
            return False
        pngs = sorted(tmp.glob("slide_*.png"))
        if not pngs:
            print(f"WARNING: render produced no slide_*.png in {tmp}")
            return False
        images = [Image.open(p).convert("RGB") for p in pngs]
        gap = 20
        w = max(img.width for img in images)
        h = sum(img.height for img in images) + gap * (len(images) - 1)
        composite = Image.new("RGB", (w, h), (240, 240, 240))
        y = 0
        for img in images:
            # Center horizontally if any surface is narrower
            x = (w - img.width) // 2
            composite.paste(img, (x, y))
            y += img.height + gap
        png_path.parent.mkdir(parents=True, exist_ok=True)
        composite.save(png_path)
        # Close PIL Image handles so the temp dir can be cleaned up
        for img in images:
            img.close()
        return True


def render_layout_thumbnails(template_path: Path, out_dir: Path,
                              dpi: int = 96) -> dict:
    """Render one PNG thumbnail per layout in the client template.

    Opens a copy of the template, clears existing sample slides, adds one
    slide per layout (instantiating the layout untouched so its native
    placeholders + decorative chrome appear), saves a temp PPTX, renders
    via LibreOffice, and returns {layout_name: png_path}.

    Used by the registration HTML so the user can SEE each layout when
    deciding which to pick. Empty placeholders show their hint text
    ("Click to add title"), which is informative — Mario sees where the
    title sits, where body content lands, what the chrome looks like.
    """
    try:
        from render_slides import render_libre
    except ImportError as e:
        print(f"WARNING: could not import render_libre: {e}")
        return {}
    from pptx import Presentation as _Pres
    from copy import deepcopy as _deepcopy

    out_dir.mkdir(parents=True, exist_ok=True)

    prs = _Pres(str(template_path))
    _clear_existing_slides(prs)

    # Dedupe by name so the order/count matches propose_layouts_payload
    # (which also dedupes via a `seen` set). Without this, the dict mapping
    # name -> png is last-write-wins and points at the wrong rendered slide.
    layout_order = []
    _seen = set()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name = (layout.name or "").strip()
            if not name or name in _seen:
                continue
            _seen.add(name)
            layout_order.append(name)
            prs.slides.add_slide(layout)

    if not layout_order:
        return {}

    tmp_pptx = out_dir / "_layout_thumbnails.pptx"
    prs.save(str(tmp_pptx))

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        try:
            render_libre(tmp_pptx, tmp, dpi=dpi)
        except Exception as e:
            print(f"WARNING: render_libre failed for layout thumbnails: {e}")
            return {}
        pngs = sorted(tmp.glob("slide_*.png"))
        if len(pngs) < len(layout_order):
            print(f"WARNING: rendered {len(pngs)} PNGs but expected "
                  f"{len(layout_order)} layouts; some thumbnails missing")

        mapping = {}
        for i, name in enumerate(layout_order):
            if i >= len(pngs):
                break
            safe = re.sub(r"[^A-Za-z0-9_-]+", "_", name)[:80]
            dst = out_dir / f"thumb_{i:02d}_{safe}.png"
            try:
                shutil.copy(str(pngs[i]), str(dst))
                mapping[name] = dst
            except Exception as e:
                print(f"WARNING: could not copy thumbnail for {name!r}: {e}")

    try:
        tmp_pptx.unlink()
    except Exception:
        pass

    return mapping


# ---------------------------------------------------------------------------
# Palette swatch PNG — visual aid for chat-driven picks
# ---------------------------------------------------------------------------

def render_palette_swatches(palette: list, png_path: Path) -> bool:
    """Render a numbered swatch chart for the propose flow.

    Each palette entry gets a colored cell with index, slot name, and hex
    overlaid. The orchestrator shows this PNG in chat next to the preview
    PNG so the user picks by visible color, not by hex string.

    palette: list of {"index": int, "slot": str, "hex": str} dicts
    png_path: output PNG path (typically <template>.palette.png)
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError as e:
        print(f"  WARNING: PIL unavailable, skipping palette swatch PNG: {e}")
        return False

    if not palette:
        return False

    cols = 4
    rows = (len(palette) + cols - 1) // cols
    cell_w, cell_h = 240, 120
    pad = 12
    label_h = 44
    W = cols * cell_w + (cols + 1) * pad
    H = rows * cell_h + (rows + 1) * pad

    img = Image.new("RGB", (W, H), (245, 245, 245))
    draw = ImageDraw.Draw(img)

    # Best-effort font load. PIL's default font is small; try a few system fonts.
    font_idx = None
    font_lbl = None
    for fname in ("arialbd.ttf", "Arial-Bold.ttf", "DejaVuSans-Bold.ttf"):
        try:
            font_idx = ImageFont.truetype(fname, 28)
            font_lbl = ImageFont.truetype(fname.replace("bd", "").replace("-Bold", ""), 14)
            break
        except (OSError, IOError):
            continue
    if font_idx is None:
        font_idx = ImageFont.load_default()
        font_lbl = ImageFont.load_default()

    for i, entry in enumerate(palette):
        col = i % cols
        row = i // cols
        x0 = pad + col * (cell_w + pad)
        y0 = pad + row * (cell_h + pad)
        x1 = x0 + cell_w
        y1 = y0 + cell_h

        try:
            r = int(entry["hex"][0:2], 16)
            g = int(entry["hex"][2:4], 16)
            b = int(entry["hex"][4:6], 16)
        except (KeyError, ValueError):
            r, g, b = 200, 200, 200

        # Color block (top portion)
        draw.rectangle([x0, y0, x1, y1 - label_h], fill=(r, g, b), outline=(80, 80, 80), width=1)
        # Label block (bottom strip, light gray)
        draw.rectangle([x0, y1 - label_h, x1, y1], fill=(255, 255, 255), outline=(80, 80, 80), width=1)

        # Auto-pick text color on the swatch by luminance (so dark swatches use white text, light use black)
        lum = (0.299 * r + 0.587 * g + 0.114 * b)
        ink = (255, 255, 255) if lum < 140 else (30, 30, 30)

        # Index, big and prominent, top-left of the color block
        draw.text((x0 + 12, y0 + 8), str(entry["index"]), fill=ink, font=font_idx)
        # Slot name, top-right of the color block
        slot_text = entry["slot"]
        try:
            tw = draw.textlength(slot_text, font=font_lbl)
        except AttributeError:
            tw = font_lbl.getsize(slot_text)[0] if hasattr(font_lbl, "getsize") else 60
        draw.text((x1 - 12 - tw, y0 + 14), slot_text, fill=ink, font=font_lbl)

        # Hex on the white label strip
        hex_text = f"#{entry['hex']}"
        draw.text((x0 + 12, y1 - label_h + 12), hex_text, fill=(30, 30, 30), font=font_lbl)

    png_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(png_path)
    return True


# ---------------------------------------------------------------------------
# Interactive register.html — one-page registration UX
# ---------------------------------------------------------------------------

def render_register_html(proposal: dict, html_path: Path) -> bool:
    """Write a self-contained interactive registration page.

    The page embeds the preview PNG + a clickable palette + the strip-bg
    toggle. The user picks primary / accent / cover-bg by clicking a
    swatch (no hex typing), toggles strip_master_backgrounds, and the
    page shows a live picks.json payload they can copy to clipboard
    and paste back to the chat orchestrator.

    proposal: the dict returned by _extract_theme_and_preview()
    html_path: output HTML path (typically <template>.register.html)
    """
    import html as _html

    tpl_name = Path(proposal["template"]).name
    preview_rel = Path(proposal["preview_png"]).name if proposal.get("preview_png") else None
    palette = proposal.get("palette", [])
    bg = proposal.get("best_guess", {})
    bg_primary = bg.get("primary_slot", "")
    bg_accent  = bg.get("accent_slot", "")
    bg_cover   = bg.get("cover_bg_slot", "")
    sha8 = proposal.get("sha8", "")
    fonts_head = proposal.get("font_heading", "")
    fonts_body = proposal.get("font_body", "")
    n_master = proposal.get("n_master", 0)
    n_layout = proposal.get("n_layout", 0)

    # Build palette JSON for the JS (the page is interactive, JS holds state)
    palette_json = json.dumps(palette)
    layouts_json = json.dumps(proposal.get("layouts", []))

    preview_block = (
        f'<img src="{_html.escape(preview_rel)}" alt="Preview composite" class="preview-img">'
        if preview_rel else
        '<div class="preview-missing">The preview image was not rendered (LibreOffice unavailable). '
        f'Inspect the preview PPTX directly: <code>{_html.escape(Path(proposal["preview_pptx"]).name)}</code></div>'
    )

    html_doc = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>Register template — TPL_NAME</title>
<style>
  :root {
    --bg:        #FAFAFA;
    --panel:     #FFFFFF;
    --rule:      #E1E1E6;
    --text-dk:   #1A1A1A;
    --text-mid:  #555;
    --text-dim:  #888;
    --accent:    #4D148C;
    --ok:        #10B981;
    --warn:      #F59E0B;
  }
  * { box-sizing: border-box; margin: 0; padding: 0; }
  body {
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", system-ui, sans-serif;
    background: var(--bg);
    color: var(--text-dk);
    padding: 24px 32px 64px;
    line-height: 1.5;
    max-width: 1280px;
    margin: 0 auto;
  }
  header {
    border-bottom: 2px solid var(--accent);
    padding-bottom: 16px;
    margin-bottom: 24px;
  }
  header h1 {
    font-size: 24px;
    font-weight: 700;
    color: var(--accent);
    margin-bottom: 6px;
  }
  header .sub {
    font-size: 13px;
    color: var(--text-mid);
  }
  header code { font-family: ui-monospace, Consolas, monospace; font-size: 12px; }
  section {
    background: var(--panel);
    border: 1px solid var(--rule);
    border-radius: 8px;
    padding: 20px 24px;
    margin-bottom: 20px;
  }
  section h2 {
    font-size: 16px;
    font-weight: 700;
    color: var(--text-dk);
    margin-bottom: 12px;
  }
  section .help {
    font-size: 13px;
    color: var(--text-mid);
    margin-bottom: 16px;
  }
  .preview-img {
    display: block;
    max-width: 100%;
    border: 1px solid var(--rule);
    border-radius: 4px;
  }
  .preview-img-caption {
    font-size: 11px;
    color: var(--text-dim);
    text-align: center;
    margin-top: 6px;
    font-style: italic;
  }
  .preview-missing {
    padding: 16px;
    background: #FEF3C7;
    border: 1px solid var(--warn);
    border-radius: 4px;
    font-size: 13px;
  }
  .picker { margin-bottom: 28px; }
  .picker:last-child { margin-bottom: 0; }
  .picker-label {
    font-size: 14px;
    font-weight: 700;
    color: var(--text-dk);
    margin-bottom: 4px;
    display: flex;
    align-items: center;
    gap: 8px;
  }
  .picker-label .current {
    font-family: ui-monospace, Consolas, monospace;
    font-size: 12px;
    color: var(--text-mid);
    font-weight: 400;
  }
  .role-explainer {
    font-size: 12.5px;
    color: var(--text-mid);
    margin-bottom: 10px;
    padding: 8px 12px;
    background: #F9FAFB;
    border-left: 3px solid var(--accent);
    border-radius: 0 4px 4px 0;
    line-height: 1.55;
  }
  .role-explainer strong { color: var(--text-dk); }
  .role-explainer em { color: var(--text-dk); font-style: normal; font-weight: 600; }
  .swatches {
    display: grid;
    grid-template-columns: repeat(12, 1fr);
    gap: 6px;
  }
  .swatch {
    aspect-ratio: 1 / 1;
    border: 2px solid var(--rule);
    border-radius: 4px;
    cursor: pointer;
    position: relative;
    transition: transform 80ms ease, border-color 80ms ease;
  }
  .swatch:hover { transform: scale(1.08); border-color: var(--text-mid); }
  .swatch.selected {
    border-color: var(--accent);
    border-width: 3px;
    box-shadow: 0 0 0 2px rgba(77, 20, 140, 0.18);
  }
  .swatch .idx {
    position: absolute;
    top: 4px;
    left: 4px;
    font-size: 11px;
    font-weight: 700;
    text-shadow: 0 0 2px rgba(0,0,0,0.6), 0 0 2px rgba(0,0,0,0.6);
    color: #fff;
    pointer-events: none;
  }
  .strip-toggle {
    display: flex;
    align-items: flex-start;
    gap: 12px;
    padding: 16px;
    background: #F9FAFB;
    border-radius: 6px;
  }
  .strip-toggle input { margin-top: 3px; transform: scale(1.3); cursor: pointer; }
  .strip-toggle .label-block { flex: 1; }
  .strip-toggle .label-block strong { display: block; font-size: 14px; margin-bottom: 4px; }
  .strip-toggle .label-block .note { font-size: 12px; color: var(--text-mid); }
  .picks-block {
    background: #1A1A1A;
    color: #E5E7EB;
    border-radius: 6px;
    padding: 16px;
    font-family: ui-monospace, Consolas, monospace;
    font-size: 13px;
    white-space: pre-wrap;
    word-break: break-all;
    margin-bottom: 12px;
  }
  .picks-block .json-key   { color: #93C5FD; }
  .picks-block .json-str   { color: #FBBF24; }
  .picks-block .json-bool  { color: #F472B6; }
  .copy-btn {
    background: var(--accent);
    color: #fff;
    border: none;
    border-radius: 6px;
    padding: 10px 18px;
    font-size: 13px;
    font-weight: 600;
    cursor: pointer;
    font-family: inherit;
  }
  .copy-btn:hover { opacity: 0.9; }
  .copy-btn.ok { background: var(--ok); }
  .hint {
    font-size: 12px;
    color: var(--text-mid);
    margin-top: 8px;
  }
  .layout-row {
    display: flex;
    align-items: center;
    gap: 16px;
    padding: 10px 12px;
    border: 1px solid var(--rule);
    border-radius: 4px;
    margin-bottom: 6px;
    background: #FFFFFF;
  }
  .layout-meta { flex: 1; min-width: 0; }
  .layout-name {
    font-weight: 600;
    font-family: ui-monospace, Consolas, monospace;
    font-size: 13px;
    color: var(--text-dk);
    overflow: hidden;
    text-overflow: ellipsis;
    white-space: nowrap;
  }
  .layout-sub {
    font-size: 11.5px;
    color: var(--text-mid);
    margin-top: 2px;
  }
  .layout-sub code {
    font-family: ui-monospace, Consolas, monospace;
    font-size: 11px;
    color: var(--accent);
  }
  .layout-toggle {
    display: flex;
    gap: 12px;
    font-size: 12.5px;
    color: var(--text-dk);
    white-space: nowrap;
  }
  .layout-toggle label { cursor: pointer; }
  .layout-toggle input[type="radio"] { margin-right: 4px; cursor: pointer; }
</style>
</head>
<body>

<header>
  <h1>Register template — TPL_NAME</h1>
  <div class="sub">
    SHA8 <code>SHA8</code> · MASTERS masters / LAYOUTS layouts ·
    fonts heading <code>FONT_HEAD</code> / body <code>FONT_BODY</code>
  </div>
</header>

<section id="preview-section">
  <h2>1. Does this preview look like the client's brand?</h2>
  <div class="help">
    Three sample slides built with the auto-best-guess colors so you can eyeball them before
    committing. Top is a cover-style slide; middle is a typical content layout with a KPI
    callout; bottom is a 2&times;2 dashboard. If the dominant color looks like the client's
    brand, you're done &mdash; jump to step 4 and copy the JSON. If it looks wrong (inverted,
    washed out, off-brand), fix in step 2.
  </div>
  PREVIEW_BLOCK
</section>

<section id="picks-section">
  <h2>2. Pick brand colors</h2>
  <div class="help">
    Click a swatch to assign it. Each row picks one color from the template's palette.
    The auto-best-guess is pre-selected (purple ring). Auto-pick has historically inverted
    on most client templates &mdash; if the preview above looked wrong, this is where you
    correct it.
  </div>

  <div class="picker" data-role="primary_slot">
    <div class="picker-label">
      Primary brand color
      <span class="current" id="cur-primary"></span>
    </div>
    <div class="role-explainer">
      <strong>The dominant brand color.</strong> Used for title bars, headers, dark-canvas
      fills, anything that should read as "this is a &lt;client&gt; deck." For most clients
      this is the color in their logo and on their website headers
      (e.g. FedEx purple, Accenture purple, IBM blue).
    </div>
    <div class="swatches" id="sw-primary"></div>
  </div>

  <div class="picker" data-role="accent_slot">
    <div class="picker-label">
      Accent / highlight color
      <span class="current" id="cur-accent"></span>
    </div>
    <div class="role-explainer">
      <strong>The "look at this" color.</strong> Used sparingly &mdash; one accent moment
      per slide (a single hero KPI, a highlight stripe, a recommended row). Never for body
      text or background. Typically the brand's secondary
      color (e.g. FedEx orange, Accenture violet/magenta). Should
      contrast strongly with primary.
    </div>
    <div class="swatches" id="sw-accent"></div>
  </div>

  <div class="picker" data-role="cover_bg_slot">
    <div class="picker-label">
      Cover-slide background color
      <span class="current" id="cur-cover"></span>
    </div>
    <div class="role-explainer">
      <strong>The fill behind the cover slide only.</strong> Drives the dramatic full-bleed
      background on slide 1 (the deck title page). Most often the same as primary, but some
      brands prefer a lighter or darker variant for the cover. <em>Default to primary unless
      you have a specific cover style in mind.</em>
    </div>
    <div class="swatches" id="sw-cover"></div>
  </div>

  <div class="picker" data-role="dark_bg_slot">
    <div class="picker-label">
      Dark-variant background color
      <span class="current" id="cur-dark"></span>
    </div>
    <div class="role-explainer">
      <strong>The fill behind any slide marked as a dark variant.</strong> When a brief
      sets <code>**Variant:** dark</code> on a slide, Slide Lab strips the layout&rsquo;s bright
      chrome and paints this color full-bleed instead. <em>Default is your primary brand
      color &mdash; change only if your brand has a specific dark-variant shade.</em>
    </div>
    <div class="swatches" id="sw-dark"></div>
  </div>
</section>

<section id="layouts-section">
  <h2>3. Pick the slide design you want</h2>
  <div id="layouts-pick-cluster"></div>
</section>

<section id="strip-section">
  <h2>4. Strip baked-in master backgrounds?</h2>
  <div class="help">
    For most client templates the master <strong>is</strong> the brand chrome — FedEx purple bars,
    Accenture branding, header rules, footer marks. Keeping it is the right default. Tick this box
    <strong>only</strong> if the master ships with unwanted photographic decoration
    (stadium photos, building shots, watermarks you don't want behind every slide).
    Stripping is destructive — leave unchecked if in doubt.
    <strong>Default: keep.</strong>
  </div>
  <label class="strip-toggle">
    <input type="checkbox" id="strip-bg">
    <div class="label-block">
      <strong>Strip master-slide background photos / decoration</strong>
      <div class="note">
        Opt in only when the master has photographic decoration you do NOT want behind
        every built slide. Stripping the master removes the brand chrome along with the
        decoration — for FedEx / Accenture / most consulting templates that's the wrong move.
      </div>
    </div>
  </label>
</section>

<section id="commit-section">
  <h2>5. Copy picks JSON and paste back to the chat</h2>
  <div class="help">
    Click the button to copy the picks JSON to your clipboard. Paste it back to the chat
    where you opened this page. The orchestrator will write it to disk and run
    <code>register_template.py commit</code>.
  </div>
  <div class="picks-block" id="picks-json"></div>
  <button class="copy-btn" id="copy-btn">Copy picks JSON to clipboard</button>
  <div class="hint">
    Or save the JSON to a file and run:
    <code>py -3 register_template.py commit "TPL_PATH" --picks &lt;picks.json&gt;</code>
  </div>
</section>

<script>
const palette = PALETTE_JSON;
const layouts = LAYOUTS_JSON;
const bestGuess = {
  primary_slot: "BG_PRIMARY",
  accent_slot:  "BG_ACCENT",
  cover_bg_slot:"BG_COVER",
};
const layoutClassifications = {};
layouts.forEach(l => { layoutClassifications[l.name] = l.auto_class; });

// Build swatches for one picker row.
function buildSwatches(rowId, role) {
  const row = document.getElementById(rowId);
  row.innerHTML = "";
  palette.forEach(entry => {
    const sw = document.createElement("div");
    sw.className = "swatch";
    sw.style.backgroundColor = "#" + entry.hex;
    sw.title = `${entry.index}. slot ${entry.slot} - #${entry.hex}`;
    sw.dataset.slot = entry.slot;
    sw.dataset.hex  = entry.hex;
    sw.dataset.idx  = entry.index;
    const idx = document.createElement("span");
    idx.className = "idx";
    idx.textContent = entry.index;
    sw.appendChild(idx);
    sw.addEventListener("click", () => selectSwatch(role, entry));
    row.appendChild(sw);
  });
}

// Active selections
const state = {
  primary_slot: bestGuess.primary_slot,
  primary_hex:  paletteHex(bestGuess.primary_slot),
  accent_slot:  bestGuess.accent_slot,
  accent_hex:   paletteHex(bestGuess.accent_slot),
  cover_bg_slot:bestGuess.cover_bg_slot,
  cover_bg_hex: paletteHex(bestGuess.cover_bg_slot),
  dark_bg_slot: bestGuess.primary_slot,   // default to primary
  dark_bg_hex:  paletteHex(bestGuess.primary_slot),
  strip_master_backgrounds: false,
};

function paletteHex(slot) {
  const e = palette.find(p => p.slot === slot);
  return e ? e.hex : "";
}

function selectSwatch(role, entry) {
  if (role === "primary_slot") {
    state.primary_slot = entry.slot;
    state.primary_hex  = entry.hex;
  } else if (role === "accent_slot") {
    state.accent_slot = entry.slot;
    state.accent_hex  = entry.hex;
  } else if (role === "cover_bg_slot") {
    state.cover_bg_slot = entry.slot;
    state.cover_bg_hex  = entry.hex;
  } else if (role === "dark_bg_slot") {
    state.dark_bg_slot = entry.slot;
    state.dark_bg_hex  = entry.hex;
  }
  refreshUI();
}

function refreshUI() {
  // Highlight selected in each row
  document.querySelectorAll("#sw-primary .swatch").forEach(s => {
    s.classList.toggle("selected", s.dataset.slot === state.primary_slot);
  });
  document.querySelectorAll("#sw-accent .swatch").forEach(s => {
    s.classList.toggle("selected", s.dataset.slot === state.accent_slot);
  });
  document.querySelectorAll("#sw-cover .swatch").forEach(s => {
    s.classList.toggle("selected", s.dataset.slot === state.cover_bg_slot);
  });
  document.querySelectorAll("#sw-dark .swatch").forEach(s => {
    s.classList.toggle("selected", s.dataset.slot === state.dark_bg_slot);
  });

  // Update "current" labels
  document.getElementById("cur-primary").textContent = `selected: #${state.primary_hex} (slot ${state.primary_slot})`;
  document.getElementById("cur-accent").textContent  = `selected: #${state.accent_hex} (slot ${state.accent_slot})`;
  document.getElementById("cur-cover").textContent   = `selected: #${state.cover_bg_hex} (slot ${state.cover_bg_slot})`;
  document.getElementById("cur-dark").textContent    = `selected: #${state.dark_bg_hex} (slot ${state.dark_bg_slot})`;

  // Update strip toggle
  state.strip_master_backgrounds = document.getElementById("strip-bg").checked;

  // Update picks JSON preview. v0.4.1 (2026-06-02): no longer emits
  // layout_classifications — that field is architectural metadata
  // (does Slide Lab strip placeholders or inherit them?) that the user
  // shouldn't have to think about. register_template.py commit recomputes
  // it from the template's actual placeholder structure. The picks JSON
  // now carries only editorial intent: brand colors + the layout the user
  // picked as the default for content slides.
  const payload = {
    accept: false,
    primary_slot:  state.primary_slot,
    primary_hex:   state.primary_hex,
    accent_slot:   state.accent_slot,
    accent_hex:    state.accent_hex,
    cover_bg_slot: state.cover_bg_slot,
    cover_bg_hex:  state.cover_bg_hex,
    dark_bg_slot:  state.dark_bg_slot,
    dark_bg_hex:   state.dark_bg_hex,
    strip_master_backgrounds: state.strip_master_backgrounds,
    default_content_layout: layoutPicks.body || "",
  };
  document.getElementById("picks-json").textContent = JSON.stringify(payload, null, 2);

  // Reset copy button if it was showing "Copied!"
  const btn = document.getElementById("copy-btn");
  btn.classList.remove("ok");
  btn.textContent = "Copy picks JSON to clipboard";
}

// Pick state: which named layout is the body. Covers stay on Slide Lab's
// draw-everything path (no layout pick — Slide Lab builds them from scratch
// on a blank layout the same way it did before v0.3).
const layoutPicks = { body: null };

function autoBodyGuess() {
  // First preference: a non-dark layout whose name contains "default slide
  // template" — that's the common convention for the workhorse body layout.
  let g = layouts.find(l => l.background !== "dark"
                        && (l.name || "").toLowerCase().includes("default slide template"));
  if (g) return g.name;
  // Fallback: first non-dark, non-cover layout with >= 1 placeholder.
  g = layouts.find(l => l.background !== "dark"
                    && !(l.name || "").toLowerCase().includes("cover")
                    && (l.n_placeholders || 0) >= 1);
  return g ? g.name : (layouts[0] && layouts[0].name) || null;
}

function buildLayoutGrid() {
  const container = document.getElementById("layouts-pick-cluster");
  if (!container) return;
  container.innerHTML = "";
  if (!layouts.length) {
    container.innerHTML = '<div class="help">No slide designs found.</div>';
    return;
  }
  if (!layoutPicks.body) layoutPicks.body = autoBodyGuess();

  const help = document.createElement("div");
  help.style.fontSize = "13px";
  help.style.color = "var(--text-mid)";
  help.style.marginBottom = "16px";
  help.innerHTML =
    "Slide Lab has rendered every layout in your template below. Click " +
    "<strong>Use for body</strong> on the layout you want for ~95% of pages. " +
    "<br><br>" +
    "<strong>Cover and section-divider slides are handled separately</strong> " +
    "— Slide Lab draws those from scratch on a blank canvas the way it always " +
    "has. You don't pick a cover layout here.";
  container.appendChild(help);

  const grid = document.createElement("div");
  grid.style.display = "grid";
  grid.style.gridTemplateColumns = "repeat(auto-fill, minmax(320px, 1fr))";
  grid.style.gap = "16px";

  layouts.forEach(l => {
    const card = document.createElement("div");
    card.className = "layout-row";
    card.style.flexDirection = "column";
    card.style.alignItems = "stretch";
    card.style.padding = "12px";
    const isBody = layoutPicks.body === l.name;
    if (isBody) {
      card.style.borderColor = "var(--accent)";
      card.style.borderWidth = "2px";
    }

    // Thumbnail
    const thumbWrap = document.createElement("div");
    thumbWrap.style.cssText = "background: #F0F0F0; border: 1px solid var(--rule);"
      + "border-radius: 4px; overflow: hidden; aspect-ratio: 16 / 9;"
      + "display: flex; align-items: center; justify-content: center;"
      + "margin-bottom: 10px;";
    if (l.thumbnail) {
      const img = document.createElement("img");
      img.src = l.thumbnail;
      img.alt = l.name;
      img.style.cssText = "width: 100%; height: 100%; object-fit: contain; display: block;";
      thumbWrap.appendChild(img);
    } else {
      thumbWrap.innerHTML = "<span style='color: var(--text-dim); font-size: 12px;'>(no thumbnail)</span>";
    }
    card.appendChild(thumbWrap);

    // BODY badge if picked
    const badges = document.createElement("div");
    badges.style.cssText = "display: flex; gap: 6px; margin-bottom: 8px; min-height: 18px;";
    if (isBody) {
      const b = document.createElement("span");
      b.style.cssText = "background: var(--accent); color: white; font-size: 11px;"
        + "font-weight: 600; padding: 2px 8px; border-radius: 3px;";
      b.textContent = "BODY";
      badges.appendChild(b);
    }
    card.appendChild(badges);

    // Layout name (raw)
    const nameEl = document.createElement("div");
    nameEl.style.cssText = "font-family: ui-monospace, Consolas, monospace;"
      + "font-size: 12px; color: var(--text-mid); margin-bottom: 10px;"
      + "word-break: break-all;";
    nameEl.textContent = l.name;
    card.appendChild(nameEl);

    // Single button: Use for body
    const bodyBtn = document.createElement("button");
    bodyBtn.className = "copy-btn";
    bodyBtn.style.cssText = "width: 100%; padding: 8px 10px; font-size: 13px;";
    bodyBtn.textContent = isBody ? "✓ Picked for body" : "Use for body";
    bodyBtn.addEventListener("click", () => pickBody(l.name));
    card.appendChild(bodyBtn);

    grid.appendChild(card);
  });

  container.appendChild(grid);
}

function pickBody(name) {
  // Single-pick semantics (2026-06-02): clicking a layout in the grid
  // selects it as THE default content layout. Reset every other layout
  // to its auto-detected classification first so prior clicks don't
  // accumulate (the additive behavior was the 2026-06-02 register-html
  // bug — every click added another "body-canonical" without unsetting
  // the previous pick, ending up with 12+ body-canonicals).
  layoutPicks.body = name;
  layouts.forEach(l => {
    if (l.name === layoutPicks.body) {
      // Picked layout: classify as body-canonical so chrome resolution
      // populates the inherited title/footer placeholders post-graft.
      layoutClassifications[l.name] = "body-canonical";
    } else {
      // Everything else: revert to Slide Lab's auto guess. The user can
      // still flip individual layouts via the advanced "Show all"
      // disclosure below; this just removes accidental additive picks.
      layoutClassifications[l.name] = l.auto_class;
    }
  });
  buildLayoutGrid();
  refreshUI();
}

// Build per-layout classification rows (collapsed under "Show all" details)
function buildLayoutRows() {
  const container = document.getElementById("layouts-list");
  if (!container) return;
  container.innerHTML = "";
  if (!layouts.length) {
    container.innerHTML = '<div class="help">No layouts detected in this template.</div>';
    return;
  }
  layouts.forEach(l => {
    const row = document.createElement("div");
    row.className = "layout-row";
    const auto = l.auto_class;
    const bgLabel = l.background === "dark" ? "dark bg" : "light bg";
    const pnLabel = l.has_page_number ? "page#" : "no-page#";
    row.innerHTML = `
      <div class="layout-meta">
        <div class="layout-name">${l.name}</div>
        <div class="layout-sub">
          auto: <code>${auto}</code> &middot;
          ${l.n_placeholders} placeholders, ${l.n_shapes} shapes &middot;
          ${bgLabel} &middot; ${pnLabel}
        </div>
      </div>
      <div class="layout-toggle">
        <label><input type="radio" name="lc-${l.name}" value="body-canonical"
               ${auto === 'body-canonical' ? 'checked' : ''}> body-canonical</label>
        <label><input type="radio" name="lc-${l.name}" value="bespoke"
               ${auto === 'bespoke' ? 'checked' : ''}> bespoke</label>
      </div>
    `;
    container.appendChild(row);
    row.querySelectorAll(`input[name="lc-${l.name}"]`).forEach(input => {
      input.addEventListener("change", (e) => {
        layoutClassifications[l.name] = e.target.value;
        refreshUI();
      });
    });
  });
}

// Wire up
buildSwatches("sw-primary", "primary_slot");
buildSwatches("sw-accent",  "accent_slot");
buildSwatches("sw-cover",   "cover_bg_slot");
buildSwatches("sw-dark",    "dark_bg_slot");
buildLayoutGrid();
document.getElementById("strip-bg").addEventListener("change", refreshUI);
document.getElementById("copy-btn").addEventListener("click", () => {
  const txt = document.getElementById("picks-json").textContent;
  navigator.clipboard.writeText(txt).then(() => {
    const btn = document.getElementById("copy-btn");
    btn.textContent = "Copied! Paste back to the chat.";
    btn.classList.add("ok");
  }).catch(err => {
    alert("Clipboard copy failed: " + err + "\\nSelect the JSON above and copy manually.");
  });
});
refreshUI();
</script>

</body>
</html>
"""
    # Substitute tokens (using plain replace, not f-string, so the embedded JS/CSS
    # with its own {} braces stays intact)
    html_doc = (html_doc
        .replace("TPL_NAME",    _html.escape(tpl_name))
        .replace("TPL_PATH",    _html.escape(proposal["template"]).replace("\\", "\\\\"))
        .replace("SHA8",        _html.escape(sha8))
        .replace("MASTERS",     str(n_master))
        .replace("FONT_HEAD",   _html.escape(fonts_head))
        .replace("FONT_BODY",   _html.escape(fonts_body))
        .replace("PREVIEW_BLOCK", preview_block)
        .replace("PALETTE_JSON", palette_json)
        .replace("LAYOUTS_JSON", layouts_json)
        .replace("LAYOUTS",     str(n_layout))
        .replace("BG_PRIMARY",  _html.escape(bg_primary))
        .replace("BG_ACCENT",   _html.escape(bg_accent))
        .replace("BG_COVER",    _html.escape(bg_cover))
    )

    html_path.write_text(html_doc, encoding="utf-8")
    return True


# ---------------------------------------------------------------------------
# Template SHA
# ---------------------------------------------------------------------------
def sha256_of_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(str(path), "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


# ---------------------------------------------------------------------------
# brand.yml / theme.json writers
# ---------------------------------------------------------------------------
BRAND_YML_TEMPLATE = """# {stem}.brand.yml
# Registered {date} - Template SHA: {sha8}
# Source: theme1.xml + register_template.py

# Ground-truth hexes - these drive every build.
# Edit and re-run register_template.py to re-validate.
primary_hex:     "#{primary_hex}"
accent_hex:      "#{accent_hex}"
cover_bg_hex:    "#{cover_bg_hex}"
dark_bg_hex:     "#{dark_bg_hex}"

# Informational only - slot names from PowerPoint (not used by builds,
# kept for human reference)
primary_slot:    {primary_slot}
accent_slot:     {accent_slot}
cover_bg_slot:   {cover_bg_slot}
dark_bg_slot:    {dark_bg_slot}

# Fonts
font_heading:    "{font_heading}"
font_body:       "{font_body}"

# Background handling
# Default false — KEEP the master decoration. For most client templates
# (FedEx, Accenture, etc.) the master IS the brand chrome and stripping
# makes every slide look off-spec. Set true ONLY when the master ships
# unwanted photographic decoration (stadium photos, watermarks) that you
# do NOT want behind every slide.
strip_master_backgrounds: {strip_master_bg}
"""


def write_brand_yml(path: Path, *, primary_hex: str, accent_hex: str,
                    cover_bg_hex: str, primary_slot: str, accent_slot: str,
                    cover_bg_slot: str,
                    dark_bg_hex: str, dark_bg_slot: str,
                    font_heading: str, font_body: str,
                    strip_master_backgrounds: bool, sha8: str) -> None:
    content = BRAND_YML_TEMPLATE.format(
        stem=path.stem.replace(".brand", ""),
        date=datetime.now().strftime("%Y-%m-%d"),
        sha8=sha8,
        primary_hex=primary_hex.upper(),
        accent_hex=accent_hex.upper(),
        cover_bg_hex=cover_bg_hex.upper(),
        dark_bg_hex=dark_bg_hex.upper(),
        primary_slot=primary_slot,
        accent_slot=accent_slot,
        cover_bg_slot=cover_bg_slot,
        dark_bg_slot=dark_bg_slot,
        font_heading=font_heading or "",
        font_body=font_body or "",
        strip_master_bg="true" if strip_master_backgrounds else "false",
    )
    path.write_text(content, encoding="utf-8")


def write_theme_json(path: Path, *, template_path: Path, sha: str,
                     brand_dict: dict, all_colors: dict, all_fonts: dict,
                     master_count: int, layout_count: int,
                     brand_yml_path: Path,
                     default_content_layout: str = "") -> None:
    """Schema v2 (2026-06-02): adds `default_content_layout` — the layout
    the user picked as the default for content slides in briefs that don't
    override per-slide. Read by build_deck.py:resolve_slide_layouts as the
    template-level fallback. Empty string = not set (orchestrator should
    ask, or build_deck falls back to sole-body-canonical heuristic)."""
    obj = {
        "schema_version": 2,
        "template_path": str(template_path),
        "template_sha": sha,
        "registered_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "registered_by": os.environ.get("USERNAME") or os.environ.get("USER") or "",
        "default_content_layout": default_content_layout or "",
        "brand": brand_dict,
        "extracted": {
            "all_colors": all_colors,
            "all_fonts": all_fonts,
            "master_count": master_count,
            "layout_count": layout_count,
        },
        "brand_yml_path": str(brand_yml_path),
    }
    path.write_text(json.dumps(obj, indent=2), encoding="utf-8")


# ---------------------------------------------------------------------------
# Master/layout counts (helps audit)
# ---------------------------------------------------------------------------
def count_masters_and_layouts(template_path: Path) -> tuple:
    n_master, n_layout = 0, 0
    try:
        with zipfile.ZipFile(str(template_path), "r") as zf:
            for name in zf.namelist():
                ln = name.lower()
                if ln.startswith("ppt/slidemasters/") and ln.endswith(".xml") \
                        and "/_rels/" not in ln:
                    n_master += 1
                elif ln.startswith("ppt/slidelayouts/") and ln.endswith(".xml") \
                        and "/_rels/" not in ln:
                    n_layout += 1
    except Exception:
        pass
    return n_master, n_layout


# ---------------------------------------------------------------------------
# Interactive swatch picker
# ---------------------------------------------------------------------------
def _ansi_swatch(hex_color: str) -> str:
    """Return an ANSI escape that prints a colored block. Empty string if
    terminal probably doesn't support it (we don't actually check)."""
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
    except Exception:
        return ""
    # Foreground 38;2;r;g;b for a colored block (back to default after)
    return f"\033[38;2;{r};{g};{b}m████\033[0m"


def show_palette(colors: Dict[str, str]) -> list:
    """Print a numbered list of theme colors. Returns the list in order."""
    entries = []
    for slot in _COLOR_SLOTS:
        hexv = colors.get(slot, "")
        if not hexv:
            continue
        entries.append((slot, hexv))
    print()
    for i, (slot, hexv) in enumerate(entries, 1):
        block = _ansi_swatch(hexv)
        print(f"  {i:>2}. {slot:<10} #{hexv}  {block}")
    print()
    return entries


def pick_index(prompt: str, n_max: int) -> int:
    while True:
        s = input(prompt).strip()
        if not s:
            continue
        try:
            i = int(s)
        except ValueError:
            print(f"  not a number; pick 1..{n_max}")
            continue
        if 1 <= i <= n_max:
            return i - 1
        print(f"  out of range; pick 1..{n_max}")


# ---------------------------------------------------------------------------
# chrome.yml writer — used by every commit path (interactive, propose/commit,
# auto-accept). Re-opens the template to extract; stamps sha8; writes to
# _p.chrome_yml(tpl). Returns 0 on success, 2 on failure.
# ---------------------------------------------------------------------------

def _write_chrome_yml_for(tpl: Path, *, sha8: str,
                          layout_overrides: dict | None,
                          commit_method: str | None = None) -> int:
    """Extract and dump chrome.yml next to the template (subfolder layout).

    sha8 freshness-stamps the spec so finalize_deck.py can hard-fail when the
    .pptx mutates after registration. layout_overrides lets picks.json
    reclassify any layout body-canonical <-> bespoke.

    commit_method (optional): an audit-trail string appended to the YAML
    document as a top-level field. Used to mark sidecars written by the
    implementer's default-pick path (commit_method="implementer_default")
    vs. real user input via the registration page.
    """
    chrome_yml_path = _p.chrome_yml(tpl)
    chrome_yml_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs = Presentation(str(tpl))
        spec = extract_chrome_spec(prs, sha8=sha8,
                                   classifications_override=layout_overrides)
        dump_chrome_yml(spec, chrome_yml_path)
    except Exception as exc:
        print(f"  ERROR: chrome.yml extraction failed: "
              f"{type(exc).__name__}: {exc}")
        return 2
    # Keep commit_method outside the validated ChromeSpec by writing it to a
    # sibling marker file. Cleaner than weakening pydantic strict validation.
    if commit_method:
        sib = _p.chrome_commit_method_txt(tpl)
        sib.write_text(str(commit_method).strip() + "\n", encoding="utf-8")
        print(f"  commit_method marker: {sib} -> {commit_method}")
    n_layouts = len(spec.layouts)
    n_bespoke = sum(1 for v in spec.layouts.values()
                    if v.layout_class == "bespoke")
    n_canonical = n_layouts - n_bespoke
    print(f"  chrome.yml:  {chrome_yml_path}")
    print(f"               {n_layouts} layouts "
          f"({n_canonical} body-canonical, {n_bespoke} bespoke); "
          f"sha8={sha8}")
    return 0


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def _main_interactive(args) -> int:
    sys.stdout.reconfigure(encoding="utf-8")

    tpl = args.template.resolve()
    if not tpl.exists():
        print(f"ERROR: template not found: {tpl}")
        return 2

    # ---- Phase 1: auto-extract ----
    print("=" * 72)
    print(f"register_template.py - {tpl.name}")
    print("=" * 72)
    print("\n[Phase 1] Auto-extract theme1.xml")

    theme_xml = _read_theme_xml(tpl)
    colors = {slot: _extract_color(theme_xml, slot) for slot in _COLOR_SLOTS}
    font_heading = _extract_font(theme_xml, "majorFont")
    font_body = _extract_font(theme_xml, "minorFont")

    n_master, n_layout = count_masters_and_layouts(tpl)
    primary_hex, primary_slot, accent_hex, accent_slot = \
        _best_guess_primary_accent(colors, tpl)
    cover_bg_hex, cover_bg_slot = primary_hex, primary_slot

    print(f"  colors extracted: {len([v for v in colors.values() if v])}")
    print(f"  fonts: heading={font_heading!r}  body={font_body!r}")
    print(f"  masters: {n_master}  layouts: {n_layout}")
    print(f"  best-guess primary:  #{primary_hex} (slot: {primary_slot})")
    print(f"  best-guess accent:   #{accent_hex} (slot: {accent_slot})")
    print(f"  best-guess cover-bg: #{cover_bg_hex} (slot: {cover_bg_slot})")

    # Default: KEEP master backgrounds. Stripping is destructive — for many
    # client templates (FedEx purple chrome, Accenture branding, etc.) the
    # master IS the brand identity and stripping makes every built slide look
    # off-spec. Opt in to strip only when the master carries photographic
    # decoration the user explicitly does not want behind every slide.
    strip_master_backgrounds = False

    # ---- Phase 2: preview build (always) ----
    _p.template_sidecar_dir(tpl).mkdir(parents=True, exist_ok=True)
    preview_path = _p.preview_pptx(tpl)
    preview_png = _p.preview_png(tpl)
    print("\n[Phase 2] Preview build (3 surfaces: title + content + dashboard)")
    build_preview_pptx(
        template_path=tpl,
        primary_hex=primary_hex, accent_hex=accent_hex,
        cover_bg_hex=cover_bg_hex,
        font_heading=font_heading, font_body=font_body,
        all_colors=colors,
        out_path=preview_path,
    )
    print(f"  preview pptx: {preview_path}")

    rendered = render_pptx_to_png(preview_path, preview_png)
    if rendered:
        print(f"  preview png:  {preview_png}  (composite of 3 surfaces)")
    else:
        print(f"  preview png: (render failed - inspect {preview_path} directly)")

    sha = sha256_of_file(tpl)
    sha8 = sha[:8]

    if args.auto_accept_phase1:
        print("\n[Phase 4] Auto-accept (--auto-accept-phase1)")
        _write_outputs(
            tpl, sha, sha8, primary_hex, primary_slot, accent_hex, accent_slot,
            cover_bg_hex, cover_bg_slot,
            dark_bg_hex=primary_hex, dark_bg_slot=primary_slot,
            font_heading=font_heading, font_body=font_body,
            strip_master_backgrounds=strip_master_backgrounds,
            colors=colors, n_master=n_master, n_layout=n_layout,
        )
        rc = _write_chrome_yml_for(tpl, sha8=sha8, layout_overrides=None)
        if rc != 0:
            return rc
        print(f"\n  Preview PNG: {preview_png}")
        print(f"  Review and re-run register_template.py (without --auto-accept-phase1)")
        print(f"  to confirm or override.")
        return 0

    # ---- Phase 2 confirmation ----
    while True:
        print()
        print(f"  Preview PNG: {preview_png}")
        print(f"  Open it. Does this look right?")
        ans = input("  [Y]es / [N]o pick swatches / [E]dit brand.yml: ").strip().lower()
        if ans in ("y", "yes"):
            # Fix #1: TTY gate on the write commit. Piped 'y' is not enough —
            # the user must press a key at the controlling terminal.
            try:
                tty_ans = _tty_input(
                    "\n  TTY CONFIRM REQUIRED — type 'y' at the keyboard "
                    "to write brand.yml + theme.json (n = re-pick): "
                ).strip().lower()
            except RuntimeError as e:
                print(f"  ABORT: {e}")
                return 2
            if tty_ans in ("y", "yes"):
                print("  TTY confirmation received.")
                break
            print("  TTY confirmation not received; re-running picker.")
            ans = "n"  # fall through to Phase 3
        if ans in ("e", "edit"):
            # Write defaults then open in editor
            _write_outputs(
                tpl, sha, sha8, primary_hex, primary_slot,
                accent_hex, accent_slot, cover_bg_hex, cover_bg_slot,
                dark_bg_hex=primary_hex, dark_bg_slot=primary_slot,
                font_heading=font_heading, font_body=font_body,
                strip_master_backgrounds=strip_master_backgrounds,
                colors=colors, n_master=n_master, n_layout=n_layout,
            )
            brand_yml = _p.brand_yml(tpl)
            try:
                os.startfile(str(brand_yml))
            except Exception as e:
                print(f"  could not open editor: {e}")
            print(f"  Brand.yml: {brand_yml}")
            print(f"  Edit it and re-run register_template.py to re-validate.")
            return 0
        if ans in ("n", "no"):
            # Phase 3: swatch picker
            print("\n[Phase 3] Swatch picker")
            entries = show_palette(colors)

            # Fix #2: typed-hex confirmation if picked color has visually-similar
            # siblings in the palette. Loops until user lands a confirmed pick.
            while True:
                i = pick_index("  Pick PRIMARY by number: ", len(entries))
                primary_slot, primary_hex = entries[i]
                if _confirm_hex_if_ambiguous(
                    "PRIMARY", primary_slot, primary_hex, entries
                ):
                    break
            while True:
                i = pick_index("  Pick ACCENT by number: ", len(entries))
                accent_slot, accent_hex = entries[i]
                if _confirm_hex_if_ambiguous(
                    "ACCENT", accent_slot, accent_hex, entries
                ):
                    break

            ans2 = input("  Pick COVER BG by number (or Enter = same as PRIMARY): ").strip()
            if not ans2:
                cover_bg_slot, cover_bg_hex = primary_slot, primary_hex
            else:
                try:
                    i = int(ans2) - 1
                    if 0 <= i < len(entries):
                        cover_bg_slot, cover_bg_hex = entries[i]
                        if not _confirm_hex_if_ambiguous(
                            "COVER_BG", cover_bg_slot, cover_bg_hex, entries
                        ):
                            print("  defaulting COVER_BG to PRIMARY.")
                            cover_bg_slot, cover_bg_hex = primary_slot, primary_hex
                    else:
                        cover_bg_slot, cover_bg_hex = primary_slot, primary_hex
                except ValueError:
                    cover_bg_slot, cover_bg_hex = primary_slot, primary_hex
            ans3 = input(
                "  Strip baked-in master backgrounds on content slides?\n"
                "    Default NO — keeps the client's brand chrome (FedEx purple,\n"
                "    Accenture branding, etc.) behind every slide. Answer yes ONLY\n"
                "    if the master ships unwanted photographic decoration.\n"
                "  [y/N]: "
            ).strip().lower()
            strip_master_backgrounds = (ans3 == "y" or ans3 == "yes")
            # Rebuild preview (3 surfaces, against client template)
            print(f"\n  Rebuilding preview with primary={primary_hex} accent={accent_hex} cover_bg={cover_bg_hex}")
            build_preview_pptx(
                template_path=tpl,
                primary_hex=primary_hex, accent_hex=accent_hex,
                cover_bg_hex=cover_bg_hex,
                font_heading=font_heading, font_body=font_body,
                all_colors=colors,
                out_path=preview_path,
            )
            render_pptx_to_png(preview_path, preview_png)
            # loop back to confirmation
            continue
        print("  please answer Y, N, or E")

    # ---- Phase 4: write canonical files ----
    print("\n[Phase 4] Write brand.yml + theme.json + chrome.yml")
    _write_outputs(
        tpl, sha, sha8, primary_hex, primary_slot,
        accent_hex, accent_slot, cover_bg_hex, cover_bg_slot,
        # Interactive Phase 1 has no dark-bg picker; default to primary.
        dark_bg_hex=primary_hex, dark_bg_slot=primary_slot,
        font_heading=font_heading, font_body=font_body,
        strip_master_backgrounds=strip_master_backgrounds,
        colors=colors, n_master=n_master, n_layout=n_layout,
    )
    rc = _write_chrome_yml_for(tpl, sha8=sha8, layout_overrides=None)
    if rc != 0:
        return rc
    return 0


def _write_outputs(tpl: Path, sha: str, sha8: str,
                   primary_hex: str, primary_slot: str,
                   accent_hex: str, accent_slot: str,
                   cover_bg_hex: str, cover_bg_slot: str,
                   *,
                   dark_bg_hex: str, dark_bg_slot: str,
                   font_heading: str, font_body: str,
                   strip_master_backgrounds: bool,
                   colors: dict, n_master: int, n_layout: int,
                   default_content_layout: str = "") -> None:
    _p.template_sidecar_dir(tpl).mkdir(parents=True, exist_ok=True)
    brand_yml = _p.brand_yml(tpl)
    theme_json = _p.theme_json(tpl)

    write_brand_yml(
        brand_yml,
        primary_hex=primary_hex, accent_hex=accent_hex,
        cover_bg_hex=cover_bg_hex,
        primary_slot=primary_slot, accent_slot=accent_slot,
        cover_bg_slot=cover_bg_slot,
        dark_bg_hex=dark_bg_hex, dark_bg_slot=dark_bg_slot,
        font_heading=font_heading, font_body=font_body,
        strip_master_backgrounds=strip_master_backgrounds,
        sha8=sha8,
    )

    brand_dict = {
        "primary_hex": f"#{primary_hex.upper()}",
        "accent_hex": f"#{accent_hex.upper()}",
        "cover_bg_hex": f"#{cover_bg_hex.upper()}",
        "dark_bg_hex": f"#{dark_bg_hex.upper()}",
        "primary_slot": primary_slot,
        "accent_slot": accent_slot,
        "cover_bg_slot": cover_bg_slot,
        "dark_bg_slot": dark_bg_slot,
        "font_heading": font_heading,
        "font_body": font_body,
        "strip_master_backgrounds": strip_master_backgrounds,
    }
    write_theme_json(
        theme_json,
        template_path=tpl, sha=sha, brand_dict=brand_dict,
        all_colors=colors, all_fonts={"heading": font_heading, "body": font_body},
        master_count=n_master, layout_count=n_layout,
        brand_yml_path=brand_yml,
        default_content_layout=default_content_layout,
    )

    print(f"  Registered.")
    print(f"  brand.yml:  {brand_yml}")
    print(f"  theme.json: {theme_json}")


# ---------------------------------------------------------------------------
# Chrome extraction (v0.2 P1.2) — body-canonical vs bespoke classifier
# ---------------------------------------------------------------------------
#
# For each slide_layout in the template:
#   1. Classify as body-canonical (chrome frame only) or bespoke (art-directed).
#      Heuristic: count non-placeholder shapes. <=2 -> body-canonical. The user
#      can override per-layout in register.html before commit.
#   2. For bespoke layouts: walk placeholders, extract bounds in EMU, convert
#      to px (96 DPI), populate the title/subtitle/footnote/source/page_number
#      fields. Missing placeholder -> field stays None (no silent canonical
#      fallback — that's the v1 slot-mapping bug class).
#   3. For body-canonical layouts: leave position fields None. helpers.py
#      reads canonical_*_box() constants from _chrome_schema.py at build time.
#   4. Detect text_role from master+layout background luminance.
#
# Result lands in <template-stem>.chrome.yml next to brand.yml / theme.json.

_EMU_PER_PX: int = 9525  # 96 DPI: 1 px = 9525 EMU. python-pptx convention.


def _emu_to_px(emu: int) -> int:
    """Convert EMU to integer px at 96 DPI. None-safe (returns 0 for None)."""
    if emu is None:
        return 0
    return int(round(int(emu) / _EMU_PER_PX))


def _luminance_from_hex(hex_str: str) -> float:
    """Rec. 709 luminance in [0, 255] from a 6-char hex string. NaN on bad input."""
    s = (hex_str or "").lstrip("#")
    if len(s) != 6:
        return float("nan")
    try:
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
    except ValueError:
        return float("nan")
    return 0.299 * r + 0.587 * g + 0.114 * b


def _walk_for_first_srgb(element) -> str:
    """Return the first <a:srgbClr val="..."> hex found under `element`, or ''."""
    for el in element.iter():
        tag = el.tag.split("}", 1)[-1] if "}" in el.tag else el.tag
        if tag == "srgbClr":
            val = el.get("val")
            if val and len(val) == 6:
                return val.upper()
    return ""


def _bg_is_dark(element) -> bool | None:
    """Inspect the element's <p:cSld><p:bg> and return True/False/None.

    True  -> luminance < 128 (dark background)
    False -> luminance >= 128 (light background)
    None  -> no resolvable srgb fill (caller falls back to caller's default)
    """
    bg_el = None
    csld = element.find(qn("p:cSld"))
    if csld is not None:
        bg_el = csld.find(qn("p:bg"))
    if bg_el is None:
        return None
    hex_val = _walk_for_first_srgb(bg_el)
    if not hex_val:
        return None
    lum = _luminance_from_hex(hex_val)
    if lum != lum:  # NaN
        return None
    return lum < 128.0


def _detect_text_role_for_layout(layout) -> tuple[str, str]:
    """Return (text_role, background) for one layout.

    Inspects the LAYOUT's own background first; falls back to its slide
    MASTER's background; defaults to light_on_dark only when there's an
    explicit dark fill signal. Defaults to dark_on_light otherwise.
    """
    layout_dark = _bg_is_dark(layout.element)
    master_dark = _bg_is_dark(layout.slide_master.element)
    is_dark = layout_dark if layout_dark is not None else master_dark
    if is_dark is True:
        return "light_on_dark", "dark"
    return "dark_on_light", "light"


def _classify_layout(layout) -> str:
    """Heuristic classifier: body-canonical vs bespoke.

    v0.4.1 (2026-06-02) — wide-net rule. A layout is body-canonical if it
    has BOTH a top-anchored title placeholder AND at least one content
    placeholder (body / object / chart / table / picture). Decorative
    shapes (FedEx chevron + footer band, Accenture brand rules, etc.)
    don't trigger bespoke anymore — those are inherited chrome, not
    content geometry, and they don't interfere with Slide Lab's
    title-plus-body content flow.

    The old rule counted non-placeholder shapes (≤2 → body-canonical).
    That misclassified every FedEx/OTC workhorse layout as bespoke
    because the chevron + footer band each count as one non-placeholder
    shape. The resulting bespoke chrome.yml entries had null footnote /
    source / page-number positions (the layout has no such placeholders
    to extract from), and finalize_deck failed loud when trying to render
    those zones.

    Bespoke is reserved for true art-directed layouts: covers (centered
    title only, no body), section dividers (large title or quote, no
    body), thank-you slides, and similar — layouts where the layout
    itself dictates content position rather than serving as a workhorse
    title+body chassis.
    """
    has_top_title = False
    has_body_content = False
    has_only_center_title = False
    title_count = 0

    BODY_CONTENT_TYPES = {
        PP_PLACEHOLDER.BODY,
        PP_PLACEHOLDER.OBJECT,
        PP_PLACEHOLDER.CHART,
        PP_PLACEHOLDER.TABLE,
        PP_PLACEHOLDER.PICTURE,
        PP_PLACEHOLDER.MEDIA_CLIP,
    }

    for ph in layout.placeholders:
        try:
            t = ph.placeholder_format.type
        except Exception:
            continue
        if t == PP_PLACEHOLDER.TITLE:
            has_top_title = True
            title_count += 1
        elif t == PP_PLACEHOLDER.CENTER_TITLE:
            # Center-title is cover/divider hallmark, NOT a workhorse title.
            has_only_center_title = True
            title_count += 1
        elif t in BODY_CONTENT_TYPES:
            has_body_content = True

    # Body-canonical: top-anchored title + at least one content placeholder.
    if has_top_title and has_body_content:
        return "body-canonical"

    # Cover / divider / quote / thank-you style: center title only, or
    # no title, or no body — these are art-directed by the layout itself.
    return "bespoke"


def _placeholder_role(shape) -> str | None:
    """Map a placeholder shape to a chrome role.

    Returns one of {'title', 'subtitle', 'footnote', 'source', 'page_number'}
    or None when the placeholder doesn't correspond to a chrome slot.
    """
    try:
        pf = shape.placeholder_format
    except Exception:
        return None
    if pf is None:
        return None
    t = pf.type
    name_lower = (shape.name or "").lower()
    if t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
        return "title"
    if t == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if t == PP_PLACEHOLDER.SLIDE_NUMBER:
        return "page_number"
    if t == PP_PLACEHOLDER.FOOTER:
        # 'footer' in client templates is usually the source line; some
        # templates split source vs footnote into separate text placeholders.
        # Disambiguate by name token; otherwise call it 'source'.
        if "footnote" in name_lower:
            return "footnote"
        return "source"
    # Some templates expose dated/footer via BODY placeholder with named
    # idx. We don't map those to chrome roles — only the standard types.
    return None


def _box_px_from_shape(shape) -> BoxPx:
    """Extract a BoxPx from a python-pptx shape (EMU -> px)."""
    return BoxPx(
        x_px=_emu_to_px(shape.left),
        y_px=_emu_to_px(shape.top),
        w_px=_emu_to_px(shape.width),
        h_px=_emu_to_px(shape.height),
        font_pt=None,
        anchor="top",
    )


def _extract_bespoke_boxes(layout) -> dict[str, BoxPx | None]:
    """Walk a bespoke layout's placeholders, return {role: BoxPx or None}.

    Missing roles stay None. No silent fallback — helpers raise at build
    time if a bespoke layout has None for an expected field.
    """
    out: dict[str, BoxPx | None] = {
        "title": None, "subtitle": None,
        "footnote": None, "source": None, "page_number": None,
    }
    for shape in layout.placeholders:
        role = _placeholder_role(shape)
        if role and out.get(role) is None:
            out[role] = _box_px_from_shape(shape)
    return out


def _extract_body_zone_for_canonical(layout) -> dict:
    """For a body-canonical layout, return v2 inheritance fields:
       title_placeholder_idx, subtitle_placeholder_idx, body_top_y_px,
       body_bottom_y_px.

    title_placeholder_idx: the .placeholder_format.idx of the FIRST title-
       type inherited placeholder (None if no title placeholder).
    subtitle_placeholder_idx: the .placeholder_format.idx of the FIRST
       subtitle-type inherited placeholder (None if no subtitle placeholder).
       Added 2026-06-02 for Bug #2: subtitle was rendering at hardcoded
       canonical position because finalize_deck had no idx to write into.
    body_top_y_px: bottom of the title placeholder in px (or 110 fallback).
    body_bottom_y_px: top of the footer placeholder in px (or 660 fallback).
    """
    title_idx = None
    subtitle_idx = None
    title_bottom_px = None
    footer_top_px = None
    for ph in layout.placeholders:
        try:
            pf = ph.placeholder_format
            t = int(pf.type)
            idx = int(pf.idx)
        except Exception:
            continue
        # type values: 1=TITLE, 13=CENTER_TITLE, 4=SUBTITLE, 15=FOOTER
        if t in (1, 13) and title_idx is None:
            title_idx = idx
            try:
                title_bottom_px = _emu_to_px(int(ph.top) + int(ph.height))
            except Exception:
                pass
        if t == 4 and subtitle_idx is None:
            subtitle_idx = idx
        if t == 15 and footer_top_px is None:
            try:
                footer_top_px = _emu_to_px(int(ph.top))
            except Exception:
                pass
    if title_bottom_px is None:
        title_bottom_px = CANONICAL_BODY_TOP_Y
    if footer_top_px is None:
        footer_top_px = CANONICAL_BODY_BOTTOM_Y
    return {
        "title_placeholder_idx": title_idx,
        "subtitle_placeholder_idx": subtitle_idx,
        "body_top_y_px": int(title_bottom_px),
        "body_bottom_y_px": int(footer_top_px),
    }


def _propose_layout_chromes(prs, classifications_override: dict[str, str] | None = None
                            ) -> dict[str, LayoutChrome]:
    """Build LayoutChrome objects for every named slide_layout in `prs`.

    classifications_override: optional {layout_name: 'body-canonical'|'bespoke'}
    to apply user-confirmed reclassifications from picks.json. Layout names
    not in the dict use the auto-detected classification.
    """
    out: dict[str, LayoutChrome] = {}
    seen_names: set[str] = set()
    # Pass 1: build all layouts with their own classification.
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name = (layout.name or "").strip()
            if not name:
                continue
            # De-dup by name across masters (FedEx-style templates put each
            # logical layout on one master only; if two masters share a layout
            # name, keep the first).
            if name in seen_names:
                continue
            seen_names.add(name)

            auto_class = _classify_layout(layout)
            final_class = (classifications_override or {}).get(name, auto_class)
            text_role, background = _detect_text_role_for_layout(layout)

            # Detect page-number presence: any placeholder with role
            # 'page_number' OR layout shape named like a page-number slot.
            has_pn = False
            for shape in layout.placeholders:
                if _placeholder_role(shape) == "page_number":
                    has_pn = True
                    break

            position_fields: dict[str, BoxPx | None] = {
                "title": None, "subtitle": None,
                "footnote": None, "source": None, "page_number": None,
            }
            if final_class == "bespoke":
                position_fields = _extract_bespoke_boxes(layout)

            # v2 (2026-05-28) + v2.1 (2026-06-02) body-canonical inheritance fields
            inherit_fields = {
                "title_placeholder_idx": None,
                "subtitle_placeholder_idx": None,
                "body_top_y_px": None,
                "body_bottom_y_px": None,
                "body_overlay_hex": None,
            }
            if final_class == "body-canonical":
                inherit_fields.update(_extract_body_zone_for_canonical(layout))

            out[name] = LayoutChrome(
                name=name,
                layout_class=final_class,  # type: ignore[arg-type]
                text_role=text_role,  # type: ignore[arg-type]
                background=background,  # type: ignore[arg-type]
                has_page_number=has_pn,
                title=position_fields["title"],
                subtitle=position_fields["subtitle"],
                footnote=position_fields["footnote"],
                source=position_fields["source"],
                page_number=position_fields["page_number"],
                title_placeholder_idx=inherit_fields["title_placeholder_idx"],
                subtitle_placeholder_idx=inherit_fields["subtitle_placeholder_idx"],
                body_top_y_px=inherit_fields["body_top_y_px"],
                body_bottom_y_px=inherit_fields["body_bottom_y_px"],
                body_overlay_hex=inherit_fields["body_overlay_hex"],
            )
    return out


def extract_chrome_spec(prs, sha8: str,
                        classifications_override: dict[str, str] | None = None
                        ) -> ChromeSpec:
    """Extract a full ChromeSpec from an opened Presentation.

    `sha8` is the first 8 hex chars of sha256(template.pptx) — stamped into
    the spec so finalize_deck can hard-fail on template-mutation drift.
    """
    return ChromeSpec(
        schema_version=CHROME_SCHEMA_VERSION_CURRENT,
        source_template_sha8=sha8,
        layouts=_propose_layout_chromes(prs, classifications_override),
    )


def propose_layouts_payload(prs) -> list[dict]:
    """Build a JSON-serializable list of per-layout proposals for register.html.

    Each entry: {name, auto_class, text_role, background, n_shapes,
                  n_placeholders, has_page_number}.
    The HTML page renders this list with a per-layout class toggle; the user's
    choices land in picks.json under `layout_classifications`.
    """
    items: list[dict] = []
    seen: set[str] = set()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            name = (layout.name or "").strip()
            if not name or name in seen:
                continue
            seen.add(name)
            auto_class = _classify_layout(layout)
            text_role, background = _detect_text_role_for_layout(layout)
            n_shapes = sum(1 for _ in layout.shapes)
            n_placeholders = sum(1 for _ in layout.placeholders)
            has_pn = any(_placeholder_role(s) == "page_number"
                         for s in layout.placeholders)
            items.append({
                "name": name,
                "auto_class": auto_class,
                "text_role": text_role,
                "background": background,
                "n_shapes": n_shapes,
                "n_placeholders": n_placeholders,
                "has_page_number": has_pn,
            })
    return items


# ---------------------------------------------------------------------------
# Chat-driven subcommands (propose / commit) — non-interactive, no TTY gate
# ---------------------------------------------------------------------------

def _extract_theme_and_preview(tpl: Path) -> dict:
    """Phase 1 + Phase 2 shared core.

    Extracts theme XML, best-guess primary/accent, builds preview PPTX, renders
    to PNG. Returns a dict carrying every value the chat orchestrator (and
    the commit subcommand) need. Does NOT prompt or write brand.yml/theme.json.
    """
    theme_xml = _read_theme_xml(tpl)
    colors = {slot: _extract_color(theme_xml, slot) for slot in _COLOR_SLOTS}
    font_heading = _extract_font(theme_xml, "majorFont")
    font_body    = _extract_font(theme_xml, "minorFont")
    n_master, n_layout = count_masters_and_layouts(tpl)
    primary_hex, primary_slot, accent_hex, accent_slot = \
        _best_guess_primary_accent(colors, tpl)
    cover_bg_hex, cover_bg_slot = primary_hex, primary_slot

    _p.template_sidecar_dir(tpl).mkdir(parents=True, exist_ok=True)
    preview_pptx     = _p.preview_pptx(tpl)
    preview_png      = _p.preview_png(tpl)
    palette_png      = _p.palette_png(tpl)
    register_html    = _p.register_html(tpl)
    build_preview_pptx(
        template_path=tpl,
        primary_hex=primary_hex, accent_hex=accent_hex,
        cover_bg_hex=cover_bg_hex,
        font_heading=font_heading, font_body=font_body,
        all_colors=colors,
        out_path=preview_pptx,
    )
    rendered = render_pptx_to_png(preview_pptx, preview_png)

    sha  = sha256_of_file(tpl)
    sha8 = sha[:8]

    palette = [
        {"index": i + 1, "slot": slot, "hex": hexv}
        for i, (slot, hexv) in enumerate(
            (s, colors[s]) for s in _COLOR_SLOTS if colors[s] and colors[s] != "000000"
        )
    ]
    palette_rendered = render_palette_swatches(palette, palette_png)

    # v0.3 (2026-05-28): render one PNG per layout so the registration HTML
    # can show the user what each layout actually looks like before they
    # pick. v0.4: stored in <sidecar-dir>/thumbnails/ (sibling of register.html).
    thumbnails_dir = _p.thumbnails_dir(tpl)
    try:
        layout_thumbnails = render_layout_thumbnails(tpl, thumbnails_dir, dpi=96)
        print(f"  layout thumbnails: {len(layout_thumbnails)} rendered "
              f"to {thumbnails_dir}")
    except Exception as _exc:
        print(f"  WARNING: layout thumbnails render failed: "
              f"{type(_exc).__name__}: {_exc}")
        layout_thumbnails = {}

    # v0.2 P1.2: per-layout chrome proposals for register.html UI. Auto-detected
    # classification + text_role per slide_layout. User reclassifies in the
    # picks JSON before commit.
    try:
        _proposal_prs = Presentation(str(tpl))
        layout_proposals = propose_layouts_payload(_proposal_prs)
    except Exception as _exc:
        print(f"  WARNING: layout proposal scan failed: {type(_exc).__name__}: {_exc}")
        layout_proposals = []
    # Attach thumbnail paths (relative to the register.html location) so the
    # <img src=...> loads under file:// without security blocks.
    # Thumbnails live at <sidecar-dir>/thumbnails/<name>.png; register.html
    # lives at <sidecar-dir>/register.html. So relative paths from the HTML
    # are just `thumbnails/<name>.png`.
    _sidecar_dir = _p.template_sidecar_dir(tpl)
    for lp in layout_proposals:
        name = lp.get("name")
        if name in layout_thumbnails:
            abs_path = layout_thumbnails[name]
            try:
                rel = abs_path.relative_to(_sidecar_dir)
                lp["thumbnail"] = str(rel).replace("\\", "/")
            except ValueError:
                lp["thumbnail"] = str(abs_path).replace("\\", "/")
        else:
            lp["thumbnail"] = None

    proposal_dict = {
        "template":        str(tpl),
        "sha":             sha,
        "sha8":            sha8,
        "n_master":        n_master,
        "n_layout":        n_layout,
        "colors":          colors,
        "font_heading":    font_heading,
        "font_body":       font_body,
        "best_guess": {
            "primary_slot":  primary_slot, "primary_hex":  primary_hex,
            "accent_slot":   accent_slot,  "accent_hex":   accent_hex,
            "cover_bg_slot": cover_bg_slot, "cover_bg_hex": cover_bg_hex,
        },
        "preview_pptx":      str(preview_pptx),
        "preview_png":       str(preview_png) if rendered else None,
        "preview_rendered":  rendered,
        "palette_png":     str(palette_png) if palette_rendered else None,
        "palette":         palette,
        "layouts":         layout_proposals,
    }

    # Render the one-page interactive register.html using the proposal dict
    html_rendered = False
    try:
        html_rendered = render_register_html(proposal_dict, register_html)
    except Exception as exc:
        print(f"  WARNING: register.html render failed: {type(exc).__name__}: {exc}")
    proposal_dict["register_html"] = str(register_html) if html_rendered else None

    return proposal_dict


def _main_propose(args) -> int:
    """Phase 1 + Phase 2: extract + preview build. Write proposal JSON. No prompts."""
    sys.stdout.reconfigure(encoding="utf-8")

    tpl = args.template.resolve()
    if not tpl.exists():
        print(f"ERROR: template not found: {tpl}")
        return 2

    print(f"register_template propose - {tpl.name}")
    print(f"  extracting theme XML + colors + fonts...")

    # Clean up legacy *.smoke.* artifacts from earlier register runs that
    # used the v0 "smoke" naming. Avoids two parallel naming conventions
    # on disk (.smoke.* + .preview.*). Skipped silently if nothing to clean.
    for legacy in (tpl.with_name(f"{tpl.stem}.smoke.pptx"),
                   tpl.with_name(f"{tpl.stem}.smoke.png")):
        try:
            if legacy.exists():
                legacy.unlink()
                print(f"  cleaned legacy artifact: {legacy.name}")
        except OSError:
            pass

    proposal = _extract_theme_and_preview(tpl)

    proposal_path = _p.register_proposal_json(tpl)
    proposal_path.parent.mkdir(parents=True, exist_ok=True)
    proposal_path.write_text(json.dumps(proposal, indent=2), encoding="utf-8")

    print(f"  best-guess primary: #{proposal['best_guess']['primary_hex']} "
          f"(slot {proposal['best_guess']['primary_slot']})")
    print(f"  best-guess accent:  #{proposal['best_guess']['accent_hex']} "
          f"(slot {proposal['best_guess']['accent_slot']})")
    print(f"  preview pptx:    {proposal['preview_pptx']}")
    if proposal["preview_rendered"]:
        print(f"  preview png:     {proposal['preview_png']}")
    else:
        print(f"  preview png:     (render failed - inspect preview pptx directly)")
    if proposal.get("palette_png"):
        print(f"  palette png:   {proposal['palette_png']}")
    if proposal.get("register_html"):
        print(f"  register page: {proposal['register_html']}")
    print(f"  proposal:      {proposal_path}")
    print()
    print(f"  Next step: open register.html in the chat preview panel. The user clicks")
    print(f"  swatches to pick primary/accent/cover-bg, toggles strip-master-backgrounds,")
    print(f"  and copies the live-updated picks JSON. The orchestrator saves it and runs:")
    print(f"    py -3 register_template.py commit {tpl} --picks <picks.json>")
    return 0


def _warn_on_picks_drift(tpl: Path) -> None:
    """Detect multiple picks-like files in the template's sidecar dir.

    The canonical filename is `register.picks.json` (written by the
    registration scripts via _p.register_picks_json). But the `commit`
    subcommand accepts an arbitrary path via `--picks`, so an orchestrator
    can drift into writing `picks.json` (or any other name) instead. When
    that happens we end up with multiple competing files describing what
    SHOULD be a single source of truth. Surface that loudly so the
    operator can resolve it before the next commit silently picks one.

    The 2026-06-02 OTC rebuild surfaced this: `picks.json` (today) and
    `register.picks.json` (May 29) co-existed with different layout
    classifications, and the operator had to manually diff them to figure
    out which one was authoritative.
    """
    sidecar_dir = _p.template_sidecar_dir(tpl)
    if not sidecar_dir.exists():
        return
    candidates = sorted(p for p in sidecar_dir.glob("*picks*.json"))
    if len(candidates) <= 1:
        return
    canonical = _p.register_picks_json(tpl)
    print(f"  WARNING: multiple picks files found in {sidecar_dir}:")
    for c in candidates:
        marker = " (canonical)" if c == canonical else " (stray)"
        print(f"    {c.name}{marker}")
    print(f"  Only {canonical.name} should exist. Consolidate manually "
          f"before the next commit, or the orchestrator's choice of "
          f"--picks path silently determines which one wins.")


def _commit_from_picks_dict(tpl: Path, picks: dict, *, source: str) -> int:
    """Shared commit core: takes a picks dict (from any source — JSON file
    or CLI flags), runs the extract pass, writes brand.yml + theme.json +
    chrome.yml. Returns 0 on success, non-zero on error.

    `source` is a one-line human-readable label printed at the top of the
    commit output for audit ("picks JSON", "CLI flags", "Phase-1 best-guess").
    """
    # Recompute extraction so we have the fonts/colors/sha — even on "accept",
    # we don't trust the picks file to carry every field.
    proposal = _extract_theme_and_preview(tpl)
    bg = proposal["best_guess"]

    if picks.get("accept", False):
        primary_slot = bg["primary_slot"]; primary_hex = bg["primary_hex"]
        accent_slot  = bg["accent_slot"];  accent_hex  = bg["accent_hex"]
        cover_bg_slot = bg["cover_bg_slot"]; cover_bg_hex = bg["cover_bg_hex"]
        dark_bg_slot = primary_slot
        dark_bg_hex  = primary_hex
        strip_master_backgrounds = False
        source = f"{source} (accept=true → Phase-1 best-guess)"
    else:
        # Required fields when not accepting:
        for k in ("primary_slot", "primary_hex", "accent_slot", "accent_hex"):
            if k not in picks:
                print(f"ERROR: picks missing required field '{k}' "
                      f"(when accept is not true). Got keys: {list(picks.keys())}")
                return 2
        primary_slot = picks["primary_slot"]
        primary_hex  = picks["primary_hex"].upper().lstrip("#")
        accent_slot  = picks["accent_slot"]
        accent_hex   = picks["accent_hex"].upper().lstrip("#")
        cover_bg_slot = picks.get("cover_bg_slot", primary_slot)
        cover_bg_hex  = picks.get("cover_bg_hex",  primary_hex).upper().lstrip("#")
        dark_bg_slot = picks.get("dark_bg_slot", primary_slot)
        dark_bg_hex  = picks.get("dark_bg_hex",  primary_hex).upper().lstrip("#")
        strip_master_backgrounds = picks.get("strip_master_backgrounds", False)

    print(f"register_template commit - {tpl.name}")
    print(f"  source:        {source}")
    print(f"  primary:       #{primary_hex} (slot {primary_slot})")
    print(f"  accent:        #{accent_hex} (slot {accent_slot})")
    print(f"  cover_bg:      #{cover_bg_hex} (slot {cover_bg_slot})")
    print(f"  dark_bg:       #{dark_bg_hex} (slot {dark_bg_slot})")
    print(f"  strip masters: {strip_master_backgrounds}")

    # default_content_layout (2026-06-02): the layout the user wants for
    # content slides in briefs that don't override. Stored in theme.json
    # so build_deck.py can use it as the template-level fallback before
    # the sole-body-canonical auto-detect heuristic fires.
    # Validated against the template's actual layout names below (after
    # chrome.yml lands) so we never persist a phantom name.
    default_content_layout = (picks.get("default_content_layout") or "").strip()

    _write_outputs(
        tpl, proposal["sha"], proposal["sha8"],
        primary_hex, primary_slot, accent_hex, accent_slot,
        cover_bg_hex, cover_bg_slot,
        dark_bg_hex=dark_bg_hex, dark_bg_slot=dark_bg_slot,
        font_heading=proposal["font_heading"], font_body=proposal["font_body"],
        strip_master_backgrounds=strip_master_backgrounds,
        colors=proposal["colors"],
        n_master=proposal["n_master"], n_layout=proposal["n_layout"],
        default_content_layout=default_content_layout,
    )

    layout_overrides = picks.get("layout_classifications") or {}
    if not isinstance(layout_overrides, dict):
        print(f"  WARNING: picks.layout_classifications is not a dict; ignored.")
        layout_overrides = {}
    rc = _write_chrome_yml_for(tpl, sha8=proposal["sha8"],
                               layout_overrides=layout_overrides,
                               commit_method=picks.get("commit_method"))
    if rc != 0:
        return rc

    # Validate default_content_layout against the newly-written chrome.yml.
    # Loud-fail if the user picked a layout name that doesn't exist in the
    # template — better to halt registration than to ship a theme.json
    # pointing at a phantom layout. See feedback_sidecar_fallback_must_be_loud.
    if default_content_layout:
        try:
            chrome_spec = load_chrome_yml(_p.chrome_yml(tpl))
            if default_content_layout not in chrome_spec.layouts:
                print(
                    f"  ERROR: default_content_layout {default_content_layout!r} "
                    f"is not a layout in {tpl.name}.\n"
                    f"  available layouts: "
                    f"{', '.join(sorted(chrome_spec.layouts.keys())) or '(none)'}"
                )
                return 2
            klass = getattr(chrome_spec.layouts[default_content_layout],
                            "layout_class", None)
            print(f"  default_content_layout: {default_content_layout!r} "
                  f"(class: {klass})")
        except Exception as _exc:
            print(f"  WARNING: could not validate default_content_layout "
                  f"against chrome.yml: {type(_exc).__name__}: {_exc}")

    brand_yml  = _p.brand_yml(tpl)
    theme_json = _p.theme_json(tpl)
    print(f"  brand.yml:   {brand_yml}")
    print(f"  theme.json:  {theme_json}")
    return 0


def _main_commit(args) -> int:
    """Phase 4: read picks JSON, write brand.yml + theme.json. No TTY gate."""
    sys.stdout.reconfigure(encoding="utf-8")

    tpl = args.template.resolve()
    if not tpl.exists():
        print(f"ERROR: template not found: {tpl}")
        return 2

    picks_path = args.picks.resolve()
    if not picks_path.exists():
        print(f"ERROR: picks JSON not found: {picks_path}")
        return 2

    try:
        picks = json.loads(picks_path.read_text(encoding="utf-8"))
    except Exception as e:
        print(f"ERROR: could not parse picks JSON: {e}")
        return 2

    _warn_on_picks_drift(tpl)

    rc = _commit_from_picks_dict(tpl, picks, source="picks JSON")
    if rc != 0:
        return rc

    # Normalize: ensure register.picks.json (the canonical filename) reflects
    # this commit's picks, regardless of which path --picks pointed at.
    # Prevents the drift class that produced two competing picks files on
    # the OTC template (see _warn_on_picks_drift docstring).
    canonical = _p.register_picks_json(tpl)
    canonical.parent.mkdir(parents=True, exist_ok=True)
    canonical.write_text(json.dumps(picks, indent=2), encoding="utf-8")
    if picks_path.resolve() != canonical.resolve():
        print(f"  normalized: copied picks -> {canonical}")
    return 0


def _main_commit_cli(args) -> int:
    """Phase 4 variant: build picks from CLI flags. No picks JSON, no
    register.html required — for chat orchestrators that can't render local
    HTML with live JS. RC-2 fix (2026-06-02).

    The orchestrator reads `palette.png` in chat, decides colors with the
    user, then invokes this subcommand with the slot + hex flags. Bypasses
    the register.html / picks.json round-trip entirely.
    """
    sys.stdout.reconfigure(encoding="utf-8")

    tpl = args.template.resolve()
    if not tpl.exists():
        print(f"ERROR: template not found: {tpl}")
        return 2

    # Build the same dict shape _commit_from_picks_dict expects.
    picks: dict = {}
    if args.accept:
        picks["accept"] = True
    else:
        # When not accepting, primary/accent are required.
        missing = []
        if not args.primary_slot or not args.primary_hex:
            missing.append("--primary-slot + --primary-hex")
        if not args.accent_slot or not args.accent_hex:
            missing.append("--accent-slot + --accent-hex")
        if missing:
            print(f"ERROR: --accept not set; missing required: {', '.join(missing)}")
            return 2
        picks["primary_slot"]  = args.primary_slot
        picks["primary_hex"]   = args.primary_hex
        picks["accent_slot"]   = args.accent_slot
        picks["accent_hex"]    = args.accent_hex
        if args.cover_bg_slot: picks["cover_bg_slot"] = args.cover_bg_slot
        if args.cover_bg_hex:  picks["cover_bg_hex"]  = args.cover_bg_hex
        if args.dark_bg_slot:  picks["dark_bg_slot"]  = args.dark_bg_slot
        if args.dark_bg_hex:   picks["dark_bg_hex"]   = args.dark_bg_hex
        picks["strip_master_backgrounds"] = args.strip_master_backgrounds

    # Parse --layout-class flags into the layout_classifications dict.
    # Format: --layout-class "Layout Name=body-canonical"
    layout_classifications: dict = {}
    for entry in (args.layout_class or []):
        if "=" not in entry:
            print(f"ERROR: --layout-class entry must be NAME=CLASS, got {entry!r}")
            return 2
        name, _, klass = entry.partition("=")
        klass = klass.strip().lower()
        if klass not in ("body-canonical", "bespoke"):
            print(f"ERROR: --layout-class CLASS must be 'body-canonical' or "
                  f"'bespoke', got {klass!r} for layout {name.strip()!r}")
            return 2
        layout_classifications[name.strip()] = klass
    if layout_classifications:
        picks["layout_classifications"] = layout_classifications

    if args.default_content_layout:
        picks["default_content_layout"] = args.default_content_layout

    picks["commit_method"] = "cli_flags"

    _warn_on_picks_drift(tpl)

    rc = _commit_from_picks_dict(tpl, picks, source="CLI flags")
    if rc != 0:
        return rc

    # Write the synthesized picks to the canonical filename so a later
    # `commit --picks <register.picks.json>` re-run replays this state.
    canonical = _p.register_picks_json(tpl)
    canonical.parent.mkdir(parents=True, exist_ok=True)
    canonical.write_text(json.dumps(picks, indent=2), encoding="utf-8")
    print(f"  wrote picks audit: {canonical}")
    return 0


# ---------------------------------------------------------------------------
# Top-level dispatcher
# ---------------------------------------------------------------------------

def main() -> int:
    ap = argparse.ArgumentParser(
        description="Register a client PPTX template (chat-driven or interactive).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    sub = ap.add_subparsers(dest="cmd", required=True, metavar="<cmd>")

    p_prop = sub.add_parser("propose",
        help="Phase 1 + Phase 2: extract + preview. Write proposal JSON. No prompts.")
    p_prop.add_argument("template", type=Path, help="Path to .pptx/.potx")

    p_comm = sub.add_parser("commit",
        help="Phase 4: read picks JSON, write brand.yml + theme.json. No TTY.")
    p_comm.add_argument("template", type=Path, help="Path to .pptx/.potx")
    p_comm.add_argument("--picks", type=Path, required=True,
        help="Path to picks JSON (see module docstring for shape).")

    # RC-2 fix (2026-06-02): CLI-flag variant of commit for chat orchestrators
    # that cannot render register.html. The orchestrator inspects palette.png
    # in chat, decides colors with the user, then invokes this subcommand
    # with the slot + hex flags directly. No picks.json round-trip needed.
    p_cli = sub.add_parser("commit-cli",
        help="Phase 4 variant: take picks from CLI flags (no picks.json, no register.html). "
             "Use when the chat cannot open register.html with live JS.")
    p_cli.add_argument("template", type=Path, help="Path to .pptx/.potx")
    p_cli.add_argument("--accept", action="store_true",
        help="Accept Phase-1 best-guess colors verbatim. When set, every other "
             "color flag is ignored.")
    p_cli.add_argument("--primary-slot", type=str, default="",
        help="Theme slot name for the primary brand color (e.g. dk2, lt2, accent1).")
    p_cli.add_argument("--primary-hex", type=str, default="",
        help="Primary hex (6-char, with or without #).")
    p_cli.add_argument("--accent-slot", type=str, default="",
        help="Theme slot name for the accent color.")
    p_cli.add_argument("--accent-hex", type=str, default="",
        help="Accent hex (6-char, with or without #).")
    p_cli.add_argument("--cover-bg-slot", type=str, default="",
        help="Optional. Cover slide background slot. Defaults to primary-slot.")
    p_cli.add_argument("--cover-bg-hex", type=str, default="",
        help="Optional. Cover slide background hex. Defaults to primary-hex.")
    p_cli.add_argument("--dark-bg-slot", type=str, default="",
        help="Optional. Dark-variant background slot. Defaults to primary-slot.")
    p_cli.add_argument("--dark-bg-hex", type=str, default="",
        help="Optional. Dark-variant background hex. Defaults to primary-hex.")
    p_cli.add_argument("--strip-master-backgrounds", action="store_true",
        help="Strip master decoration from every built slide. Default: false "
             "(keep the master — it is usually the brand chrome).")
    p_cli.add_argument("--layout-class", action="append", default=[],
        metavar="NAME=CLASS",
        help="Override auto-detected layout classification. CLASS must be "
             "'body-canonical' or 'bespoke'. Repeatable. "
             "Example: --layout-class \"Title and Content=body-canonical\"")
    p_cli.add_argument("--default-content-layout", type=str, default="",
        metavar="NAME",
        help="The layout the user picked as the default for content slides. "
             "Stored in theme.json and used by build_deck.py as the "
             "template-level layout fallback when a brief doesn't specify "
             "per-slide Layout: or front-matter default_layout:. Validated "
             "against the template's registered layouts at commit time. "
             "Example: --default-content-layout \"2_Title & Text 01\"")

    p_int = sub.add_parser("interactive",
        help="Legacy 4-phase interactive flow with TTY confirmation gate.")
    p_int.add_argument("template", type=Path, help="Path to .pptx/.potx")
    p_int.add_argument("--auto-accept-phase1", action="store_true",
        help="Skip interactive confirmation; write brand.yml using Phase-1 best-guesses.")

    args = ap.parse_args()
    if args.cmd == "propose":
        return _main_propose(args)
    if args.cmd == "commit":
        return _main_commit(args)
    if args.cmd == "commit-cli":
        return _main_commit_cli(args)
    if args.cmd == "interactive":
        return _main_interactive(args)
    ap.print_help()
    return 2


if __name__ == "__main__":
    sys.exit(main())
