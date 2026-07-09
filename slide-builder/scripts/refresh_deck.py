#!/usr/bin/env python3
"""refresh_deck.py — content-only refresh of a recurring/PMO deck (option 6c).

The recurring case: a client owns a status/PMO deck on a FIXED template that
gets re-issued each cycle with new numbers, same design. Slide Lab did not
build it, so the pipeline's rebuild path (which needs `_meta.json`) doesn't
apply, and a full build would REDESIGN it (wrong — the format must stay
identical). This does neither: it edits text runs in place, design frozen.

Two subcommands:

  spec  <deck.pptx> [--out spec.json]
        Dump every text shape (per slide) into a refresh spec the user/Claude
        edits: fill `new_text` for the fields that change this cycle, leave the
        rest null. Non-destructive — only reads the deck.

  apply <deck.pptx> <spec.json> [--out <dated-copy.pptx>]
        Apply the spec to a COPY of the deck (never the original). For each
        entry with a non-null `new_text`, locate the shape by (slide, shape
        index), verify it still matches the recorded `current_text` (drift
        guard), and replace the text while preserving the first run's
        formatting. Prints a changed / skipped / mismatch report.

Design is frozen: shapes, positions, layouts, and every non-listed run are
untouched. Run slide-qc on the output copy before sending (a deck isn't done
until QC has run).

Run:  py -3 refresh_deck.py spec  deck.pptx            (python3 on macOS/Linux)
      py -3 refresh_deck.py apply deck.pptx spec.json
"""
from __future__ import annotations

import argparse
import datetime as _dt
import json
import shutil
import sys
from pathlib import Path

from pptx import Presentation  # noqa: E402


def _iter_text_shapes(prs):
    """Yield (slide_no_1based, shape_idx, shape) for every shape with text.

    shape_idx is the position within `slide.shapes` — stable as long as the
    deck's shape layout doesn't change between cycles (the drift guard in
    apply() catches it if it does).
    """
    for s_i, slide in enumerate(prs.slides, start=1):
        for sh_i, shape in enumerate(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                yield s_i, sh_i, shape


def _shape_name(shape) -> str:
    try:
        return shape.name or ""
    except Exception:
        return ""


def cmd_spec(args) -> int:
    deck = Path(args.deck)
    if not deck.exists():
        print(f"[error] deck not found: {deck}", file=sys.stderr)
        return 2
    prs = Presentation(str(deck))
    fields = []
    for slide_no, shape_idx, shape in _iter_text_shapes(prs):
        fields.append({
            "slide": slide_no,
            "shape_idx": shape_idx,
            "shape_name": _shape_name(shape),
            "current_text": shape.text_frame.text,
            "new_text": None,  # fill this for the fields that change this cycle
        })
    spec = {
        "source_deck": str(deck.resolve()),
        "generated_at": _dt.datetime.now().isoformat(timespec="seconds"),
        "_how_to": ("Fill `new_text` for the fields that change this cycle; "
                    "leave the rest null. Then: refresh_deck.py apply <deck> <this file>."),
        "fields": fields,
    }
    out = Path(args.out) if args.out else deck.with_name(deck.stem + "_refresh_spec.json")
    out.write_text(json.dumps(spec, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"[ok] wrote refresh spec: {out}")
    print(f"     {len(fields)} text fields found across {len(prs.slides._sldIdLst)} slides.")
    print("     Edit `new_text` for the fields to change, then run `apply`.")
    return 0


def _set_text_preserving_format(shape, new_text: str) -> None:
    """Replace a shape's visible text with `new_text`, keeping the first run's
    formatting (font, size, color, bold). Extra runs/paragraphs are removed so
    the result is a single clean run — right for a number/label swap.
    """
    tf = shape.text_frame
    # Keep the first paragraph; drop the others.
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    para = tf.paragraphs[0]
    runs = list(para.runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        # No run to inherit formatting from — add one.
        para.add_run().text = new_text


def cmd_apply(args) -> int:
    deck = Path(args.deck)
    spec_path = Path(args.spec)
    if not deck.exists():
        print(f"[error] deck not found: {deck}", file=sys.stderr)
        return 2
    if not spec_path.exists():
        print(f"[error] spec not found: {spec_path}", file=sys.stderr)
        return 2
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    changes = [f for f in spec.get("fields", []) if f.get("new_text") is not None]
    if not changes:
        print("[error] no fields have `new_text` set — nothing to change. "
              "Edit the spec first.", file=sys.stderr)
        return 2

    # Always work on a dated copy; never touch the original.
    if args.out:
        out = Path(args.out)
    else:
        stamp = _dt.datetime.now().strftime("%Y%m%d")
        out = deck.with_name(f"{deck.stem}_refreshed_{stamp}{deck.suffix}")
    out.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(deck, out)

    prs = Presentation(str(out))
    slides = list(prs.slides)
    applied, mismatched, out_of_range = 0, [], []
    for f in changes:
        s_no, sh_idx = f.get("slide"), f.get("shape_idx")
        if not isinstance(s_no, int) or not (1 <= s_no <= len(slides)):
            out_of_range.append(f); continue
        shapes = list(slides[s_no - 1].shapes)
        if not isinstance(sh_idx, int) or not (0 <= sh_idx < len(shapes)):
            out_of_range.append(f); continue
        shape = shapes[sh_idx]
        if not shape.has_text_frame:
            out_of_range.append(f); continue
        # Drift guard: only overwrite if the shape still holds the text the spec
        # was generated against. Protects against a template that moved shapes.
        if shape.text_frame.text != f.get("current_text"):
            mismatched.append(f); continue
        _set_text_preserving_format(shape, f["new_text"])
        applied += 1

    prs.save(str(out))
    print(f"[ok] wrote refreshed copy: {out}")
    print(f"     applied {applied} / {len(changes)} field(s); "
          f"{len(mismatched)} skipped (text drifted), "
          f"{len(out_of_range)} skipped (slide/shape not found).")
    if mismatched:
        print("     Drifted fields (left unchanged — the deck's shapes moved since the spec was made):")
        for f in mismatched[:10]:
            print(f"       slide {f.get('slide')} shape {f.get('shape_idx')} "
                  f"({f.get('shape_name') or '?'}): expected {f.get('current_text')!r}")
    if out_of_range:
        print("     Not-found fields (slide/shape index out of range or not text):")
        for f in out_of_range[:10]:
            print(f"       slide {f.get('slide')} shape {f.get('shape_idx')}")
    print("     Design is unchanged — run slide-qc on the copy before sending.")
    # Non-zero only if NOTHING applied; partial drift is a warning, not a hard fail.
    return 0 if applied else 1


def main(argv) -> int:
    ap = argparse.ArgumentParser(description="Content-only refresh of a recurring/PMO deck (option 6c).")
    sub = ap.add_subparsers(dest="cmd", required=True)
    sp = sub.add_parser("spec", help="Dump text shapes into an editable refresh spec.")
    sp.add_argument("deck")
    sp.add_argument("--out", default=None)
    ap_ = sub.add_parser("apply", help="Apply a filled spec to a dated copy of the deck.")
    ap_.add_argument("deck")
    ap_.add_argument("spec")
    ap_.add_argument("--out", default=None)
    args = ap.parse_args(argv)
    if args.cmd == "spec":
        return cmd_spec(args)
    if args.cmd == "apply":
        return cmd_apply(args)
    ap.print_help()
    return 2


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
