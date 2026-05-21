"""
add_footer_tokens.py — add missing FOOTNOTE, SOURCE, PAGE_NUMBER token shapes
to skeleton PPTXs that are currently missing them.

Adds shapes directly to each skeleton PPTX file.

Run once: python scripts/add_footer_tokens.py
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

SKILL_DIR = pathlib.Path(__file__).resolve().parent.parent
SKEL_DIR  = SKILL_DIR / "skeletons"

# Footer zone geometry (inches) — matches the two-panel/three-column/process-chevron standard
FOOTNOTE_X, FOOTNOTE_Y  = 0.607, 7.004
FOOTNOTE_W, FOOTNOTE_H  = 12.12, 0.135

SOURCE_X, SOURCE_Y      = 0.607, 7.139
SOURCE_W, SOURCE_H      = 7.522, 0.135

PAGENUM_X, PAGENUM_Y    = 9.333, 7.139
PAGENUM_W, PAGENUM_H    = 3.393, 0.135

FOOTER_PT = 8.0  # font size for all footer tokens


def _add_textbox(slide, x_in, y_in, w_in, h_in, text, pt=FOOTER_PT,
                 align=PP_ALIGN.LEFT, shape_name=None):
    """Add a textbox with a single run to slide. Returns the shape."""
    from pptx.util import Emu
    box = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    if shape_name:
        box.name = shape_name
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(pt)
    return box


def _existing_tokens(slide):
    """Return set of token names already in slide shapes."""
    tokens = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for tok in ("FOOTNOTE", "SOURCE", "PAGE_NUMBER"):
            if "{{" + tok + "}}" in shape.text_frame.text:
                tokens.add(tok)
    return tokens


def fix_skeleton(skel_name):
    pptx_path = SKEL_DIR / skel_name / f"{skel_name}.pptx"
    if not pptx_path.exists():
        print(f"  SKIP {skel_name}: pptx not found")
        return
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    existing = _existing_tokens(slide)
    added = []

    if "FOOTNOTE" not in existing:
        _add_textbox(slide, FOOTNOTE_X, FOOTNOTE_Y, FOOTNOTE_W, FOOTNOTE_H,
                     "{{FOOTNOTE}}", shape_name="footer_footnote")
        added.append("FOOTNOTE")

    if "SOURCE" not in existing:
        _add_textbox(slide, SOURCE_X, SOURCE_Y, SOURCE_W, SOURCE_H,
                     "{{SOURCE}}", shape_name="footer_source")
        added.append("SOURCE")

    if "PAGE_NUMBER" not in existing:
        _add_textbox(slide, PAGENUM_X, PAGENUM_Y, PAGENUM_W, PAGENUM_H,
                     "{{PAGE_NUMBER}}", align=PP_ALIGN.RIGHT, shape_name="footer_page_number")
        added.append("PAGE_NUMBER")

    if added:
        prs.save(str(pptx_path))
        print(f"  OK {skel_name}: added {added}")
    else:
        print(f"  OK {skel_name}: nothing to add (already complete)")


SKELETONS = [
    "cover",
    "single-finding",
    "two-panel",
    "three-column",
    "recommendation",
    "process-chevron",
    "pull-quote",
]

print("Adding missing footer tokens to skeleton PPTXs...\n")
for name in SKELETONS:
    fix_skeleton(name)
print("\nDone.")
