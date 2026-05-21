"""
extract_icons_bulk.py — Extract ALL icons from the Accenture icon library PPTX.

Systematically discovers every icon section on every slide using text label
positions, extracts the gray outline variant of each icon, and names it from
the first keyword in the label. Handles naming conflicts and skips non-icon
slides (index/cover slides).

Output:
    icons/<canonical_name>.xml   — one file per unique icon concept
    icons/icon-index-full.json   — complete index of all extracted icons

Usage:
    py -3 extract_icons_bulk.py
    py -3 extract_icons_bulk.py --slides 13 14 15    # specific slides only
    py -3 extract_icons_bulk.py --category BUSINESS  # slides matching category
    py -3 extract_icons_bulk.py --dry-run            # list what would be extracted

The 15 hand-curated icons in icon-index.json stay unchanged.
Bulk-extracted icons go into icon-index-full.json and are named
<keyword>-<slide_num> to avoid conflicts (e.g., bar-chart-13, calendar-152).
"""

from __future__ import annotations

import argparse
import copy
import json
import math
import pathlib
import re
import sys

from lxml import etree
from pptx import Presentation

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

import os as _os

_SCRIPTS_DIR = pathlib.Path(__file__).resolve().parent
_SKILLS_DIR  = _SCRIPTS_DIR.parent
_ICONS_DIR   = _SKILLS_DIR / "icons"

# Source icon-library PPTX: see SLIDE_LAB_ICON_LIBRARY env var.
_DEFAULT_ICON_PPTX = _ICONS_DIR / "_source-library.pptx"
_ICON_PPTX = pathlib.Path(_os.environ.get("SLIDE_LAB_ICON_LIBRARY") or _DEFAULT_ICON_PPTX)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Slides that are index/cover/nav/template pages — skip these
# Slides 1-2: index/cover
# Slides 3-11: "how to use" template demos — contain text boxes and layout shapes, not icons
_SKIP_SLIDES = {1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11}

# Minimum icon size — shapes smaller than this are decorative lines/borders
_MIN_ICON_PT = 15  # pt


# ---------------------------------------------------------------------------
# Shape helpers  (duplicated here so this script is standalone)
# ---------------------------------------------------------------------------

def _get_slide_shapes(slide):
    sp_tree = slide._element.find(f".//{{{_NS_P}}}spTree")
    if sp_tree is None:
        return []
    return [c for c in sp_tree
            if c.tag.split("}")[-1] in ("sp", "grpSp", "pic", "cxnSp")]


def _get_text(elem) -> str:
    parts = [t.text.strip() for t in elem.iter(f"{{{_NS_A}}}t")
             if t.text and t.text.strip()]
    return "; ".join(parts)


def _get_bbox(elem):
    """Return (x, y, cx, cy) in EMU for a shape, or None."""
    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

    if tag == "grpSp":
        grp_pr = elem.find(f"{{{_NS_P}}}grpSpPr")
        xfrm   = grp_pr.find(f"{{{_NS_A}}}xfrm") if grp_pr is not None else None
    else:
        sp_pr = elem.find(f"{{{_NS_P}}}spPr")
        xfrm  = sp_pr.find(f"{{{_NS_A}}}xfrm") if sp_pr is not None else None

    if xfrm is None:
        xfrm = elem.find(f".//{{{_NS_A}}}xfrm")
    if xfrm is None:
        return None

    off = xfrm.find(f"{{{_NS_A}}}off")
    ext = xfrm.find(f"{{{_NS_A}}}ext")
    if off is None or ext is None:
        return None

    try:
        return (int(off.get("x", 0)), int(off.get("y", 0)),
                int(ext.get("cx", 0)), int(ext.get("cy", 0)))
    except (ValueError, TypeError):
        return None


# ---------------------------------------------------------------------------
# Naming helpers
# ---------------------------------------------------------------------------

def _canonical_name(label_text: str, slide_num: int) -> str:
    """
    Build a canonical icon name from the label text + slide number.

    Strategy:
    - Take the first keyword (before first semicolon)
    - Lowercase, strip punctuation, replace spaces with hyphens
    - Append slide number to ensure uniqueness (same keyword can appear on different slides)
    - Result: 'bar-chart-13', 'calendar-152', 'gear-131'
    """
    first_kw = label_text.split(";")[0].strip()
    # Remove characters that aren't alphanumeric, space, or hyphen
    clean    = re.sub(r"[^a-zA-Z0-9\s\-]", "", first_kw)
    # Collapse whitespace -> hyphens
    slug     = re.sub(r"[\s\-]+", "-", clean.strip()).lower().strip("-")
    if not slug:
        slug = "icon"
    return f"{slug}-{slide_num}"


# ---------------------------------------------------------------------------
# Per-slide icon discovery
# ---------------------------------------------------------------------------

def discover_icons_on_slide(slide, slide_num: int, prs: Presentation):
    """
    Discover all unique icon concepts on a slide and return their gray icons.

    Algorithm:
    1. Collect all text label shapes (sp with text, short text = icon label)
    2. For each label, find non-text shapes in the same column ABOVE the label
    3. Pick the leftmost top shape (gray outline variant)

    Returns list of dicts: {label, canonical_name, shape_elem}
    """
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    EMU_PT  = 12700

    all_shapes = _get_slide_shapes(slide)

    # Separate labels from icon shapes
    # Labels: short text (<200 chars), not too tall
    labels = []
    icons  = []

    for sh in all_shapes:
        bbox = _get_bbox(sh)
        if bbox is None:
            continue
        x, y, cx, cy = bbox

        text = _get_text(sh)
        tag  = sh.tag.split("}")[-1] if "}" in sh.tag else sh.tag

        if tag == "sp" and text:
            # Filter: skip slide title / category header (very wide, very top)
            if cy < EMU_PT * _MIN_ICON_PT:  # too thin — decorative
                continue
            if len(text) > 200:             # too long — paragraph text
                continue
            if y < slide_h * 0.04:          # in header zone
                continue
            labels.append((sh, text, bbox))
        elif cx >= EMU_PT * _MIN_ICON_PT and cy >= EMU_PT * _MIN_ICON_PT:
            icons.append((sh, bbox))

    results = []
    seen_shapes = set()  # avoid duplicating the same shape

    for _, label_text, label_bbox in labels:
        lx, ly, lcx, lcy = label_bbox
        label_cx = lx + lcx / 2.0  # horizontal center of label

        # Icons for this section are in a column x-band and above the label
        # Section width ≈ label text box width (labels are centered in their section)
        half_section = max(lcx * 0.8, slide_w * 0.06)
        x_left  = label_cx - half_section
        x_right = label_cx + half_section

        # Icons sit in the row above the label — search from label top up ~40% of slide
        y_search_bottom = ly             # must be above the label
        y_search_top    = ly - slide_h * 0.40  # don't go too far up

        candidates = []
        for icon_sh, icon_bbox in icons:
            ix, iy, icx, icy = icon_bbox
            icon_center_x = ix + icx / 2.0
            icon_center_y = iy + icy / 2.0

            if not (x_left <= icon_center_x <= x_right):
                continue
            if not (y_search_top <= iy < y_search_bottom):
                continue

            candidates.append((icon_sh, icon_bbox))

        if not candidates:
            continue

        # Among candidates, pick the leftmost + topmost (= gray outline variant)
        # Sort by x first (left icon), then by y (top icon)
        candidates.sort(key=lambda c: (c[1][0], c[1][1]))
        gray_icon = candidates[0][0]

        shape_id = id(gray_icon)
        if shape_id in seen_shapes:
            continue
        seen_shapes.add(shape_id)

        cname = _canonical_name(label_text, slide_num)
        results.append({
            "canonical_name": cname,
            "label_text":     label_text,
            "slide":          slide_num,
            "shape_elem":     gray_icon,
            "bbox_pt":        tuple(v // EMU_PT for v in _get_bbox(gray_icon))
        })

    return results


# ---------------------------------------------------------------------------
# Color normalization
# ---------------------------------------------------------------------------

def normalize_colors(elem) -> None:
    for sf in elem.iter(f"{{{_NS_A}}}solidFill"):
        srgb = sf.find(f"{{{_NS_A}}}srgbClr")
        if srgb is not None:
            srgb.set("val", "000000")


# ---------------------------------------------------------------------------
# Slide category extraction (reads the header text)
# ---------------------------------------------------------------------------

def get_slide_category(slide) -> str:
    """Read the category name from the slide header (large bold text, top-left)."""
    shapes = _get_slide_shapes(slide)
    NS_A   = _NS_A
    for sh in shapes:
        bbox = _get_bbox(sh)
        if bbox is None:
            continue
        x, y, cx, cy = bbox
        # Category text is near the top-left of the slide
        if y > 50 * 12700:  # below 50pt from top
            continue
        text = _get_text(sh)
        if text and len(text) < 60:  # short — not a label list
            return text.upper()
    return "UNKNOWN"


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def _unique_name(base: str, used: set) -> str:
    """Return base if unused, otherwise base-2, base-3, etc."""
    if base not in used:
        return base
    counter = 2
    while f"{base}-{counter}" in used:
        counter += 1
    return f"{base}-{counter}"


def main():
    parser = argparse.ArgumentParser(
        description="Bulk-extract all icons from Accenture icon library PPTX"
    )
    parser.add_argument("--pptx",      help="Override icon PPTX path")
    parser.add_argument("--slides",    nargs="+", type=int,
                        help="Extract only from these slide numbers")
    parser.add_argument("--category",  help="Filter to slides with this category name (partial match)")
    parser.add_argument("--dry-run",   action="store_true",
                        help="List what would be extracted without saving files")
    parser.add_argument("--overwrite", action="store_true",
                        help="Overwrite existing XML files (default: skip existing)")
    args = parser.parse_args()

    pptx_path = pathlib.Path(args.pptx) if args.pptx else _ICON_PPTX
    if not pptx_path.exists():
        print(f"ERROR: Icon PPTX not found: {pptx_path}")
        sys.exit(1)

    print(f"Opening: {pptx_path.name}")
    prs = Presentation(str(pptx_path))
    print(f"  {len(prs.slides)} slides")
    print(f"  Slide size: {int(prs.slide_width)//12700} x {int(prs.slide_height)//12700} pt\n")

    _ICONS_DIR.mkdir(exist_ok=True)

    # Determine which slides to process
    slide_indices = args.slides or list(range(1, len(prs.slides) + 1))
    slide_indices = [s for s in slide_indices if s not in _SKIP_SLIDES]

    all_entries   = []
    used_names    = set()   # tracks every name assigned this run — no collisions
    saved         = 0
    skipped_noicons = 0
    total_found   = 0

    for slide_num in slide_indices:
        if slide_num < 1 or slide_num > len(prs.slides):
            continue

        slide    = prs.slides[slide_num - 1]
        category = get_slide_category(slide)

        if args.category and args.category.upper() not in category:
            continue

        icons_on_slide = discover_icons_on_slide(slide, slide_num, prs)

        if not icons_on_slide:
            skipped_noicons += 1
            continue

        total_found += len(icons_on_slide)
        print(f"Slide {slide_num:3d} [{category}] — {len(icons_on_slide)} icons")

        for entry in icons_on_slide:
            base_name = entry["canonical_name"]
            cname     = _unique_name(base_name, used_names)
            used_names.add(cname)

            label    = entry["label_text"]
            shape_el = entry["shape_elem"]
            bbox_pt  = entry["bbox_pt"]
            out_path = _ICONS_DIR / f"{cname}.xml"

            entry_record = {
                "name":           cname,
                "concept":        label.split(";")[0].strip(),
                "source_slide":   slide_num,
                "label_text":     label,
                "slide_category": category,
                "bbox_pt":        list(bbox_pt),
                "source_x_pct":   round((bbox_pt[0] + bbox_pt[2] / 2) / (int(prs.slide_width)  // 12700), 3),
                "source_y_pct":   round((bbox_pt[1] + bbox_pt[3] / 2) / (int(prs.slide_height) // 12700), 3),
            }
            all_entries.append(entry_record)

            if args.dry_run:
                print(f"    [DRY] {cname}.xml  ({bbox_pt[2]}x{bbox_pt[3]}pt)  '{label[:55]}'")
                continue

            saved_elem = copy.deepcopy(shape_el)
            normalize_colors(saved_elem)
            xml_bytes  = etree.tostring(saved_elem, pretty_print=True, xml_declaration=False)
            out_path.write_bytes(xml_bytes)
            size_kb = out_path.stat().st_size // 1024
            print(f"    SAVE {cname}.xml ({size_kb}KB)")
            saved += 1

    print(f"\n{'[DRY RUN] ' if args.dry_run else ''}Results:")
    print(f"  Icons found / saved: {total_found}")
    print(f"  Empty slides skipped: {skipped_noicons}")

    # Write full index
    full_index_path = _ICONS_DIR / "icon-index-full.json"
    if not args.dry_run and all_entries:
        with open(full_index_path, "w", encoding="utf-8") as f:
            json.dump({
                "_source_pptx": str(_ICON_PPTX),
                "_note": "Auto-extracted bulk icon index. Use icon-index.json for the 15 curated vocabulary icons.",
                "_count": len(all_entries),
                "icons": all_entries
            }, f, indent=2)
        print(f"\n  Full index written: {full_index_path.name} ({len(all_entries)} icons)")
        print(f"  Unique names: {len(used_names)}")


if __name__ == "__main__":
    main()
