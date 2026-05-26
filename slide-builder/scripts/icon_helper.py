"""
icon_helper.py — Slide Lab icon insertion (build-time, zero PPTX overhead)

Reads pre-extracted icon XML files from icons/<name>.xml and injects the vector
shape directly into the target slide's spTree. No PPTX opening, no COM, no
file format conversion — just XML read + inject.

v0.1 ships icons pre-extracted under slide-builder/icons/; there is no live
extraction step. Missing names fall through to a labeled dashed-border
placeholder. See slide-builder/icons/README.md for the vocabulary + fallback.

Usage:
    from icon_helper import insert_icon

    insert_icon(
        icon_name="gear",
        target_slide=slide,       # python-pptx Slide object
        left_emu=914400,
        top_emu=1143000,
        width_emu=457200,
        height_emu=457200,
        accent_color="#A100FF"    # template accent1 hex
    )

If the .xml file is missing, inserts a labeled dashed-border placeholder
and prints a warning — does not raise an exception.
"""

from __future__ import annotations

import copy
import pathlib
import sys
from typing import Optional

from lxml import etree

_PATCHES_DIR = pathlib.Path(__file__).resolve().parent.parent / "patches"
if str(_PATCHES_DIR) not in sys.path:
    sys.path.insert(0, str(_PATCHES_DIR))
from patches import get_blank_layout

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_SKILLS_DIR = pathlib.Path(__file__).resolve().parent.parent
_ICONS_DIR = _SKILLS_DIR / "icons"

# PPTX namespaces
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Module-level XML cache: parse each icon file once per process.
# insert_icon() deepcopies from this cache so it is never mutated.
_ICON_CACHE: dict[str, "etree._Element"] = {}


# ---------------------------------------------------------------------------
# Color tinting
# ---------------------------------------------------------------------------

def _apply_accent_color(elem, hex_color: str) -> None:
    """Replace all srgbClr fills in the shape tree with hex_color.
    Skips schemeClr (theme colors) — those are intentional."""
    color = hex_color.lstrip("#").upper()
    for solid_fill in elem.iter(f"{{{_NS_A}}}solidFill"):
        srgb = solid_fill.find(f"{{{_NS_A}}}srgbClr")
        if srgb is not None:
            srgb.set("val", color)


# ---------------------------------------------------------------------------
# Position and scale
# ---------------------------------------------------------------------------

def _reposition(elem, left_emu: int, top_emu: int, width_emu: int, height_emu: int) -> bool:
    """
    Update the element's xfrm to the requested bounding box.

    PPTX grpSp coordinate system:
      off/ext    = group position/size in slide EMU space
      chOff/chExt = local coordinate origin/scale for child shapes
      Children use coords in chOff..chOff+chExt range — NOT slide EMUs.

    Correct approach for grpSp:
      1. Update off  → new slide position
      2. Update ext  → new slide size
      3. Scale chExt proportionally (so children render at new size)
      4. Leave chOff and all child xfrm elements untouched — the rendering engine
         maps children from local space to the new physical ext automatically.

    Nested grpSp inside a grpSp: their off/ext are in the PARENT's local coordinate
    space and do NOT need to be changed — they scale automatically with the outer
    group's chExt. Only the outermost group's xfrm needs updating.

    Returns True if repositioned, False if the shape had no usable xfrm.
    """
    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

    if tag == "grpSp":
        # For a group, position data is in p:grpSpPr/a:xfrm
        grp_pr = elem.find(f"{{{_NS_P}}}grpSpPr")
        if grp_pr is None:
            # Some icons serialise grpSpPr using the DML namespace — try that fallback
            grp_pr = elem.find(f"{{{_NS_A}}}grpSpPr")
        if grp_pr is None:
            print(f"  [icon] WARNING: grpSp has no grpSpPr — cannot reposition")
            return False

        xfrm = grp_pr.find(f"{{{_NS_A}}}xfrm")
        if xfrm is None:
            print(f"  [icon] WARNING: grpSpPr has no a:xfrm — cannot reposition")
            return False

        old_ext   = xfrm.find(f"{{{_NS_A}}}ext")
        old_chExt = xfrm.find(f"{{{_NS_A}}}chExt")

        old_cx = int(old_ext.get("cx", 1))  if old_ext  is not None else 1
        old_cy = int(old_ext.get("cy", 1))  if old_ext  is not None else 1

        off = xfrm.find(f"{{{_NS_A}}}off")
        if off is not None:
            off.set("x", str(left_emu))
            off.set("y", str(top_emu))
        if old_ext is not None:
            old_ext.set("cx", str(width_emu))
            old_ext.set("cy", str(height_emu))
        # chExt is deliberately NOT rescaled — keeping the original local-coord
        # extent means the renderer's ext/chExt ratio changes when ext changes,
        # which makes children scale proportionally into the new physical box.
        # (Previous version scaled chExt by new_ext/old_ext which kept the
        # ratio invariant — children rendered at original authored size
        # regardless of target box. That was the icon-too-small bug.)
        # chOff and child xfrm elements deliberately left untouched.
        return True

    else:
        # Simple shape or pic — xfrm is in p:spPr, coords are slide EMUs directly
        sp_pr = elem.find(f"{{{_NS_P}}}spPr")
        if sp_pr is None:
            # Fallback: scan for any xfrm in the element
            xfrm = elem.find(f".//{{{_NS_A}}}xfrm")
        else:
            xfrm = sp_pr.find(f"{{{_NS_A}}}xfrm")
            if xfrm is None:
                xfrm = elem.find(f".//{{{_NS_A}}}xfrm")

        if xfrm is None:
            print(f"  [icon] WARNING: shape <{tag}> has no a:xfrm — cannot reposition")
            return False

        off = xfrm.find(f"{{{_NS_A}}}off")
        ext = xfrm.find(f"{{{_NS_A}}}ext")
        if off is not None:
            off.set("x", str(left_emu))
            off.set("y", str(top_emu))
        if ext is not None:
            ext.set("cx", str(width_emu))
            ext.set("cy", str(height_emu))
        return True


# ---------------------------------------------------------------------------
# Placeholder fallback
# ---------------------------------------------------------------------------

def _insert_placeholder(target_slide, left_emu, top_emu, width_emu, height_emu,
                         icon_name: str) -> None:
    from pptx.util import Pt, Emu as _Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    shape = target_slide.shapes.add_shape(1, left_emu, top_emu, width_emu, height_emu)
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    shape.line.width = _Emu(12700)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[{icon_name}]"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def insert_icon(
    icon_name: str,
    target_slide,
    left_emu: int,
    top_emu: int,
    width_emu: int,
    height_emu: int,
    accent_color: str = "#A100FF",
) -> bool:
    """
    Insert a pre-extracted icon into target_slide at the given bounding box.

    Reads icons/<icon_name>.xml, tints with accent_color, repositions to
    (left_emu, top_emu, width_emu, height_emu), and injects into the slide's
    spTree. No PPTX opening, no COM — just XML read + inject.

    Returns True on success, False if placeholder was inserted instead.

    If the XML file is missing, the icon was not pre-extracted into the
    skill's icons/ directory. v0.1 ships pre-extracted icons under
    slide-builder/icons/*.xml; missing names fall through to a labeled
    dashed-border placeholder (no extraction step in the v0.1 pipeline).
    """
    xml_path = _ICONS_DIR / f"{icon_name}.xml"

    if not xml_path.exists():
        print(f"  [icon] WARNING: {xml_path.name} not found in slide-builder/icons/ - inserting labeled placeholder.")
        _insert_placeholder(target_slide, left_emu, top_emu, width_emu, height_emu, icon_name)
        return False

    # Load from cache (parse once per process) and deepcopy for this insertion.
    # Deepcopy is mandatory: _apply_accent_color and _reposition mutate the element
    # in-place, and lxml elements can only have one parent — appending to a second
    # slide's spTree would silently detach the element from the first.
    cache_key = str(xml_path)
    if cache_key not in _ICON_CACHE:
        try:
            _ICON_CACHE[cache_key] = etree.parse(cache_key).getroot()
        except Exception as e:
            print(f"  [icon] ERROR reading {xml_path.name}: {e} — inserting placeholder.")
            _insert_placeholder(target_slide, left_emu, top_emu, width_emu, height_emu, icon_name)
            return False

    # Work on a deep copy of the cached tree — never mutate the cache
    elem = copy.deepcopy(_ICON_CACHE[cache_key])

    # Apply accent color tint
    _apply_accent_color(elem, accent_color)

    # Reposition to target bounding box
    repositioned = _reposition(elem, left_emu, top_emu, width_emu, height_emu)
    if not repositioned:
        # Shape had no usable xfrm — insert placeholder so the slot is visible
        print(f"  [icon] WARNING: '{icon_name}' could not be repositioned — inserting placeholder.")
        _insert_placeholder(target_slide, left_emu, top_emu, width_emu, height_emu, icon_name)
        return False

    # Inject into slide spTree
    sp_tree = target_slide.shapes._spTree
    if sp_tree is None:
        print(f"  [icon] ERROR: spTree not found in target slide — inserting placeholder.")
        _insert_placeholder(target_slide, left_emu, top_emu, width_emu, height_emu, icon_name)
        return False

    sp_tree.append(elem)
    return True


# ---------------------------------------------------------------------------
# CLI: verify a single icon
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import argparse
    from pptx import Presentation
    from pptx.util import Emu

    parser = argparse.ArgumentParser(description="Verify an icon by building a test PPTX")
    parser.add_argument("icon_name", help="Icon name (e.g. gear, lightbulb)")
    parser.add_argument("output", help="Output PPTX path")
    parser.add_argument("--color", default="#A100FF", help="Accent color hex")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    slide = prs.slides.add_slide(get_blank_layout(prs))

    W = H = Emu(914400)
    L = Emu(9144000 // 2 - 457200)
    T = Emu(5143500 // 2 - 457200)

    ok = insert_icon(args.icon_name, slide, int(L), int(T), int(W), int(H), args.color)
    prs.save(args.output)
    print(f"Saved: {args.output} [{'OK' if ok else 'PLACEHOLDER'}]")
