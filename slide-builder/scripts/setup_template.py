"""
setup_template.py — One-time client template onboarding tool
=============================================================
Extracts EVERYTHING the slide builder needs to know about a client's
PowerPoint template — recursively, with full property capture.

Why "EVERYTHING": earlier versions of this script captured only shape
position + text preview. That missed:
  - shapes nested inside <p:grpSp> groups (silently dropped)
  - every fill colour (srgb AND scheme-resolved)
  - every line / stroke
  - text run properties (size, font, bold, italic, colour)
  - master text-style defaults (titleStyle/bodyStyle/otherStyle)
  - layout-level placeholder style overrides
  - z-order of layered shapes

That gap produced wrong specs in real work. This version captures all of
the above, with scheme-colour references resolved to hex via the theme
clrScheme, and preserves XML document order for z-order.

Run this ONCE per client when you receive their template PPTX.

Usage:
    py -3 setup_template.py <template.pptx> <client-name>
    py -3 setup_template.py "Contoso Template.pptx" contoso
    py -3 setup_template.py "Contoso Template.pptx" contoso --master 0

Output (default: skill's clients/ folder; use --output to write to a Claude Project folder):
    <output-dir>/template.json   — full structured spec
    <output-dir>/template.md     — human-readable summary

Recommended usage — write directly into the client's Claude Project folder:
    py -3 setup_template.py "FedEx-Template.pptx" fedex \
        --output "<path>/Claude Projects/FedEx"

The client name becomes the folder name — use lowercase-hyphenated.
--master defaults to 0 (first master). Pass --master N for a different one.
"""

from __future__ import annotations

import argparse
import json
import pathlib
import re
import sys
import zipfile
from datetime import date

from lxml import etree
from pptx import Presentation

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_SCRIPTS_DIR = pathlib.Path(__file__).resolve().parent
_SKILLS_DIR  = _SCRIPTS_DIR.parent
_CLIENTS_DIR = _SKILLS_DIR / "clients"

EMU_PER_IN = 914400

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NSMAP = {"a": NS_A, "p": NS_P, "r": NS_R}

# ---------------------------------------------------------------------------
# Small XML helpers
# ---------------------------------------------------------------------------

def _emu_to_in(emu) -> float | None:
    if emu is None:
        return None
    try:
        return round(int(emu) / EMU_PER_IN, 3)
    except (TypeError, ValueError):
        return None


def _emu_attr(elem, attr) -> int | None:
    """Read an EMU attribute, returning int or None."""
    if elem is None:
        return None
    v = elem.get(attr)
    if v is None:
        return None
    try:
        return int(v)
    except ValueError:
        return None


def _find(elem, path):
    """elem.find with NSMAP."""
    if elem is None:
        return None
    return elem.find(path, NSMAP)


def _findall(elem, path):
    if elem is None:
        return []
    return elem.findall(path, NSMAP)


# ---------------------------------------------------------------------------
# Theme — colour scheme + fonts
# ---------------------------------------------------------------------------

_COLOR_ROLES = [
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
]


def _extract_theme(pptx_path: pathlib.Path) -> dict:
    """Read every theme1.xml field we care about — colors, fonts, font fallbacks."""
    with zipfile.ZipFile(pptx_path) as z:
        theme_files = sorted(n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml", n))
        if not theme_files:
            return {"colors": {}, "fonts": {}}
        theme_xml = z.read(theme_files[0])

    root = etree.fromstring(theme_xml)

    # Colors — each role maps to either srgbClr or sysClr (sysClr has lastClr fallback)
    colors: dict[str, str] = {}
    clr_scheme = _find(root, ".//a:clrScheme")
    if clr_scheme is not None:
        for child in clr_scheme:
            role = etree.QName(child).localname
            if role not in _COLOR_ROLES:
                continue
            srgb = _find(child, "a:srgbClr")
            sysclr = _find(child, "a:sysClr")
            if srgb is not None:
                colors[role] = (srgb.get("val") or "").upper()
            elif sysclr is not None:
                # sysClr has both 'val' (system colour name) and 'lastClr' (cached hex)
                colors[role] = (sysclr.get("lastClr") or sysclr.get("val") or "").upper()

    # Fonts — major (heading) and minor (body), each with latin/ea/cs typeface
    fonts: dict[str, dict[str, str]] = {"major": {}, "minor": {}}
    font_scheme = _find(root, ".//a:fontScheme")
    if font_scheme is not None:
        for role_xml, role_key in [("a:majorFont", "major"), ("a:minorFont", "minor")]:
            mf = _find(font_scheme, role_xml)
            if mf is None:
                continue
            for script, attr in [("a:latin", "latin"), ("a:ea", "east_asian"), ("a:cs", "complex_script")]:
                el = _find(mf, script)
                if el is not None and el.get("typeface"):
                    fonts[role_key][attr] = el.get("typeface")

    return {"colors": colors, "fonts": fonts}


# ---------------------------------------------------------------------------
# Scheme colour resolver
# ---------------------------------------------------------------------------

# Scheme tokens that show up in <a:schemeClr val="X"/> — map to theme role.
# tx1/bg1/tx2/bg2 are aliases for dk1/lt1/dk2/lt2 ("text 1", "background 1", ...).
_SCHEME_ALIASES = {
    "tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2",
    "phClr": None,  # placeholder colour — resolved at use site, not here
}


def _resolve_scheme(scheme_val: str, theme_colors: dict) -> str | None:
    """Given a schemeClr val like 'accent1' or 'tx1', return its hex from theme.colors."""
    if scheme_val is None:
        return None
    role = _SCHEME_ALIASES.get(scheme_val, scheme_val)
    if role is None:
        return None
    hex_val = theme_colors.get(role)
    return hex_val.upper() if hex_val else None


# ---------------------------------------------------------------------------
# Fill / line / text-run extractors
# ---------------------------------------------------------------------------

def _extract_color_ref(parent, theme_colors: dict) -> dict | None:
    """Extract a colour reference from a <a:solidFill> or <a:fillRef> parent.
    Returns dict with 'hex', 'scheme', and 'lumMod'/'lumOff'/'shade'/'tint' modifiers.
    Returns None if no colour is set on this element.
    """
    if parent is None:
        return None

    srgb = _find(parent, "a:srgbClr")
    if srgb is not None:
        return {
            "kind": "srgb",
            "hex": (srgb.get("val") or "").upper(),
            "scheme": None,
            "modifiers": _color_modifiers(srgb),
        }

    scheme = _find(parent, "a:schemeClr")
    if scheme is not None:
        scheme_val = scheme.get("val")
        return {
            "kind": "scheme",
            "scheme": scheme_val,
            "hex": _resolve_scheme(scheme_val, theme_colors),
            "modifiers": _color_modifiers(scheme),
        }

    sysclr = _find(parent, "a:sysClr")
    if sysclr is not None:
        return {
            "kind": "sysClr",
            "scheme": None,
            "hex": (sysclr.get("lastClr") or "").upper() or None,
            "modifiers": _color_modifiers(sysclr),
        }

    return None


def _color_modifiers(clr_elem) -> dict:
    """Capture <a:lumMod>, <a:lumOff>, <a:shade>, <a:tint>, <a:alpha> on a colour element."""
    mods = {}
    for tag in ("lumMod", "lumOff", "shade", "tint", "alpha"):
        el = _find(clr_elem, f"a:{tag}")
        if el is not None:
            mods[tag] = el.get("val")
    return mods


def _extract_fill(sppr, theme_colors: dict) -> dict:
    """Extract the fill from a <p:spPr> or <a:lnPr> element."""
    if sppr is None:
        return {"kind": "inherited"}
    # No-fill
    if _find(sppr, "a:noFill") is not None:
        return {"kind": "noFill"}
    # Solid fill
    solid = _find(sppr, "a:solidFill")
    if solid is not None:
        return {"kind": "solid", "color": _extract_color_ref(solid, theme_colors)}
    # Gradient fill — capture stop list
    grad = _find(sppr, "a:gradFill")
    if grad is not None:
        stops = []
        for gs in _findall(grad, ".//a:gs"):
            stops.append({
                "pos": gs.get("pos"),
                "color": _extract_color_ref(gs, theme_colors),
            })
        return {"kind": "gradient", "stops": stops}
    # Blip (image) fill
    blip = _find(sppr, "a:blipFill")
    if blip is not None:
        embed = _find(blip, "a:blip")
        return {
            "kind": "blip",
            "embed_rId": embed.get(f"{{{NS_R}}}embed") if embed is not None else None,
        }
    # Pattern fill
    patt = _find(sppr, "a:pattFill")
    if patt is not None:
        return {"kind": "pattern", "preset": patt.get("prst")}
    # No explicit fill specified — inherits from parent (master/layout/theme)
    return {"kind": "inherited"}


def _extract_line(sppr, theme_colors: dict) -> dict:
    """Extract <a:ln> properties from a <p:spPr>."""
    if sppr is None:
        return {"kind": "inherited"}
    ln = _find(sppr, "a:ln")
    if ln is None:
        return {"kind": "inherited"}
    if _find(ln, "a:noFill") is not None:
        return {"kind": "none"}
    fill = _extract_fill(ln, theme_colors)
    width_emu = ln.get("w")
    dash = _find(ln, "a:prstDash")
    return {
        "kind": "line",
        "width_emu": int(width_emu) if width_emu else None,
        "width_pt": round(int(width_emu) / 12700, 2) if width_emu else None,
        "dash": dash.get("val") if dash is not None else None,
        "fill": fill,
    }


def _extract_run(run_elem, theme_colors: dict) -> dict:
    """Extract a single <a:r> (text run) — text + rPr."""
    rPr = _find(run_elem, "a:rPr")
    t = _find(run_elem, "a:t")
    info = {
        "text": (t.text or "") if t is not None else "",
    }
    if rPr is None:
        return info

    sz = rPr.get("sz")
    if sz:
        info["size_pt"] = int(sz) / 100  # sz is in hundredths of a point
    if rPr.get("b") is not None:
        info["bold"] = rPr.get("b") in ("1", "true")
    if rPr.get("i") is not None:
        info["italic"] = rPr.get("i") in ("1", "true")
    if rPr.get("u"):
        info["underline"] = rPr.get("u")

    # Font typeface — latin runs the show in our context
    for script, key in [("a:latin", "font"), ("a:ea", "font_east_asian"), ("a:cs", "font_complex")]:
        el = _find(rPr, script)
        if el is not None and el.get("typeface"):
            info[key] = el.get("typeface")

    # Colour
    solid = _find(rPr, "a:solidFill")
    if solid is not None:
        info["color"] = _extract_color_ref(solid, theme_colors)

    return info


def _extract_paragraph(p_elem, theme_colors: dict) -> dict:
    """Extract a single <a:p> — pPr + runs."""
    pPr = _find(p_elem, "a:pPr")
    info = {"runs": []}
    if pPr is not None:
        info["alignment"] = pPr.get("algn")
        info["level"] = pPr.get("lvl")
        info["margin_left_emu"] = pPr.get("marL")
        info["indent_emu"] = pPr.get("indent")
        # paragraph-level default rPr (used for empty paragraphs to seed default size)
        defRPr = _find(pPr, "a:defRPr")
        if defRPr is not None:
            # Use the run extractor on a synthetic <a:r> sibling structure:
            # easier: just read the attrs directly.
            d = {}
            if defRPr.get("sz"): d["size_pt"] = int(defRPr.get("sz")) / 100
            if defRPr.get("b") is not None: d["bold"] = defRPr.get("b") in ("1", "true")
            if defRPr.get("i") is not None: d["italic"] = defRPr.get("i") in ("1", "true")
            latin = _find(defRPr, "a:latin")
            if latin is not None and latin.get("typeface"):
                d["font"] = latin.get("typeface")
            solid = _find(defRPr, "a:solidFill")
            if solid is not None:
                d["color"] = _extract_color_ref(solid, theme_colors)
            if d:
                info["default_run_props"] = d

    for r in _findall(p_elem, "a:r"):
        info["runs"].append(_extract_run(r, theme_colors))
    return info


def _extract_text_frame(sp_elem, theme_colors: dict) -> dict | None:
    """Extract <p:txBody> — bodyPr + lstStyle + paragraphs."""
    txBody = _find(sp_elem, "p:txBody")
    if txBody is None:
        return None

    info = {"paragraphs": []}
    bodyPr = _find(txBody, "a:bodyPr")
    if bodyPr is not None:
        info["anchor"] = bodyPr.get("anchor")
        info["wrap"] = bodyPr.get("wrap")
        info["left_inset_emu"] = bodyPr.get("lIns")
        info["right_inset_emu"] = bodyPr.get("rIns")
        info["top_inset_emu"] = bodyPr.get("tIns")
        info["bottom_inset_emu"] = bodyPr.get("bIns")

    # lstStyle — per-level run-property defaults specific to this shape
    lstStyle = _find(txBody, "a:lstStyle")
    if lstStyle is not None:
        lvl_defaults = {}
        for child in lstStyle:
            local = etree.QName(child).localname
            if not (local.startswith("lvl") or local == "defPPr"):
                continue
            defRPr = _find(child, "a:defRPr")
            if defRPr is None:
                continue
            d = {}
            if defRPr.get("sz"): d["size_pt"] = int(defRPr.get("sz")) / 100
            if defRPr.get("b") is not None: d["bold"] = defRPr.get("b") in ("1", "true")
            latin = _find(defRPr, "a:latin")
            if latin is not None and latin.get("typeface"):
                d["font"] = latin.get("typeface")
            solid = _find(defRPr, "a:solidFill")
            if solid is not None:
                d["color"] = _extract_color_ref(solid, theme_colors)
            if d:
                lvl_defaults[local] = d
        if lvl_defaults:
            info["lstStyle_levels"] = lvl_defaults

    for p in _findall(txBody, "a:p"):
        info["paragraphs"].append(_extract_paragraph(p, theme_colors))
    return info


# ---------------------------------------------------------------------------
# Shape recursion — handles <p:sp>, <p:grpSp>, <p:pic>, <p:cxnSp>, <p:graphicFrame>
# ---------------------------------------------------------------------------

def _extract_xfrm(sp_elem) -> dict | None:
    """Pull position/size from a <p:spPr>/<a:xfrm> (or grpSpPr/xfrm)."""
    xfrm = _find(sp_elem, "p:spPr/a:xfrm")
    if xfrm is None:
        xfrm = _find(sp_elem, "p:grpSpPr/a:xfrm")
    if xfrm is None:
        return None
    off = _find(xfrm, "a:off")
    ext = _find(xfrm, "a:ext")
    chOff = _find(xfrm, "a:chOff")
    chExt = _find(xfrm, "a:chExt")
    info = {}
    if off is not None and ext is not None:
        info["left_emu"]   = _emu_attr(off, "x")
        info["top_emu"]    = _emu_attr(off, "y")
        info["width_emu"]  = _emu_attr(ext, "cx")
        info["height_emu"] = _emu_attr(ext, "cy")
        info["left_in"]    = _emu_to_in(info.get("left_emu"))
        info["top_in"]     = _emu_to_in(info.get("top_emu"))
        info["width_in"]   = _emu_to_in(info.get("width_emu"))
        info["height_in"]  = _emu_to_in(info.get("height_emu"))
    if chOff is not None and chExt is not None:
        info["chOff_emu"] = (_emu_attr(chOff, "x"), _emu_attr(chOff, "y"))
        info["chExt_emu"] = (_emu_attr(chExt, "cx"), _emu_attr(chExt, "cy"))
    rot = xfrm.get("rot")
    if rot is not None:
        info["rotation_60000ths_deg"] = int(rot)
    flipH = xfrm.get("flipH")
    if flipH is not None:
        info["flipH"] = flipH in ("1", "true")
    flipV = xfrm.get("flipV")
    if flipV is not None:
        info["flipV"] = flipV in ("1", "true")
    return info or None


def _extract_geom(sp_elem) -> dict | None:
    """Capture <a:prstGeom> or <a:custGeom> from <p:spPr>."""
    prst = _find(sp_elem, "p:spPr/a:prstGeom")
    if prst is not None:
        return {"kind": "preset", "preset": prst.get("prst")}
    cust = _find(sp_elem, "p:spPr/a:custGeom")
    if cust is not None:
        return {"kind": "custom"}
    return None


def _shape_record(sp_elem, theme_colors: dict, z_order: int) -> dict:
    """Build a full record for a single <p:sp> / <p:pic> / <p:cxnSp> / <p:graphicFrame>."""
    local = etree.QName(sp_elem).localname

    # Identification — walk candidates explicitly so lxml doesn't warn about
    # truth-testing Element objects (which is being deprecated).
    cNvPr = None
    for _path in (
        "p:nvSpPr/p:cNvPr",
        "p:nvPicPr/p:cNvPr",
        "p:nvCxnSpPr/p:cNvPr",
        "p:nvGraphicFramePr/p:cNvPr",
        "p:nvGrpSpPr/p:cNvPr",
    ):
        cNvPr = _find(sp_elem, _path)
        if cNvPr is not None:
            break

    name = cNvPr.get("name") if cNvPr is not None else None
    sp_id = cNvPr.get("id") if cNvPr is not None else None

    # Placeholder info
    ph = _find(sp_elem, "p:nvSpPr/p:nvPr/p:ph")
    ph_info = None
    if ph is not None:
        ph_info = {
            "type": ph.get("type"),     # e.g. 'title', 'subTitle', 'body'
            "idx":  ph.get("idx"),
            "sz":   ph.get("sz"),       # 'full', 'half', 'quarter'
        }

    sppr = _find(sp_elem, "p:spPr")

    record = {
        "tag":        local,        # 'sp', 'pic', 'cxnSp', 'graphicFrame', 'grpSp'
        "z_order":    z_order,
        "id":         sp_id,
        "name":       name,
        "placeholder": ph_info,
        "xfrm":       _extract_xfrm(sp_elem),
        "geometry":   _extract_geom(sp_elem),
        "fill":       _extract_fill(sppr, theme_colors),
        "line":       _extract_line(sppr, theme_colors),
        "text_frame": _extract_text_frame(sp_elem, theme_colors),
    }

    # If this is a group, recurse into children with their own z-order
    if local == "grpSp":
        children = []
        z = 0
        for child in sp_elem:
            cl = etree.QName(child).localname
            if cl in ("sp", "pic", "cxnSp", "graphicFrame", "grpSp"):
                children.append(_shape_record(child, theme_colors, z))
                z += 1
        record["children"] = children

    return record


def _walk_shapes(spTree_elem, theme_colors: dict) -> list:
    """Walk a <p:spTree> root, returning every top-level shape as a record.
    Group children are nested under their parent. Z-order is XML document order.
    """
    out = []
    z = 0
    for child in spTree_elem:
        local = etree.QName(child).localname
        if local in ("sp", "pic", "cxnSp", "graphicFrame", "grpSp"):
            out.append(_shape_record(child, theme_colors, z))
            z += 1
    return out


# ---------------------------------------------------------------------------
# Master text styles (titleStyle / bodyStyle / otherStyle)
# ---------------------------------------------------------------------------

def _extract_text_styles(master_elem, theme_colors: dict) -> dict:
    """Read <p:txStyles>/<p:titleStyle>, <p:bodyStyle>, <p:otherStyle>.
    Returns {style_name: {level_name: {size_pt, bold, font, color}}}.
    """
    out = {}
    txStyles = _find(master_elem, "p:txStyles")
    if txStyles is None:
        return out
    for style in txStyles:
        style_name = etree.QName(style).localname  # 'titleStyle' / 'bodyStyle' / 'otherStyle'
        levels = {}
        for child in style:
            local = etree.QName(child).localname
            if not (local.startswith("lvl") or local in ("defPPr",)):
                continue
            defRPr = _find(child, "a:defRPr")
            if defRPr is None:
                levels[local] = {}
                continue
            d = {}
            if defRPr.get("sz"): d["size_pt"] = int(defRPr.get("sz")) / 100
            if defRPr.get("b") is not None: d["bold"] = defRPr.get("b") in ("1", "true")
            if defRPr.get("i") is not None: d["italic"] = defRPr.get("i") in ("1", "true")
            latin = _find(defRPr, "a:latin")
            if latin is not None and latin.get("typeface"):
                d["font"] = latin.get("typeface")
            solid = _find(defRPr, "a:solidFill")
            if solid is not None:
                d["color"] = _extract_color_ref(solid, theme_colors)
            algn = child.get("algn")
            if algn:
                d["alignment"] = algn
            levels[local] = d
        out[style_name] = levels
    return out


# ---------------------------------------------------------------------------
# Master + layout extraction
# ---------------------------------------------------------------------------

def _layout_record(layout_part, theme_colors: dict, index: int) -> dict:
    """Build a full record for a slide layout."""
    layout_xml = layout_part.element  # <p:sldLayout>
    cSld = _find(layout_xml, "p:cSld")
    spTree = _find(cSld, "p:spTree") if cSld is not None else None
    shapes = _walk_shapes(spTree, theme_colors) if spTree is not None else []

    return {
        "index":  index,
        "name":   layout_part.name,
        "shapes": shapes,
    }


def _master_record(master, theme_colors: dict) -> dict:
    """Build a full record for the master slide (chrome + text styles)."""
    master_xml = master.element
    cSld = _find(master_xml, "p:cSld")
    spTree = _find(cSld, "p:spTree") if cSld is not None else None
    shapes = _walk_shapes(spTree, theme_colors) if spTree is not None else []

    return {
        "shapes":      shapes,
        "text_styles": _extract_text_styles(master_xml, theme_colors),
    }


# ---------------------------------------------------------------------------
# Effective zones (for legacy callers — title/subtitle position)
# ---------------------------------------------------------------------------

def _effective_zones_from_layouts(layouts: list, master_shapes: list) -> dict:
    """Find a representative title + subtitle placeholder position by scanning
    the master first, then layouts. Uses the first match.
    """
    zones = {}

    def _scan(shape_list):
        for sh in shape_list:
            ph = sh.get("placeholder")
            if ph is None:
                continue
            ph_type = (ph.get("type") or "").lower()
            xfrm = sh.get("xfrm") or {}
            if not xfrm.get("left_in") and xfrm.get("left_in") != 0:
                continue
            if ph_type in ("ctrtitle", "title") and "title" not in zones:
                zones["title"] = {
                    "left_in": xfrm.get("left_in"),
                    "top_in": xfrm.get("top_in"),
                    "width_in": xfrm.get("width_in"),
                    "height_in": xfrm.get("height_in"),
                }
            elif ph_type in ("subtitle", "body", "obj") and "subtitle" not in zones:
                zones["subtitle"] = {
                    "left_in": xfrm.get("left_in"),
                    "top_in": xfrm.get("top_in"),
                    "width_in": xfrm.get("width_in"),
                    "height_in": xfrm.get("height_in"),
                }
            # Recurse into groups
            for child in sh.get("children", []) or []:
                _scan([child])

    _scan(master_shapes)
    if "title" not in zones or "subtitle" not in zones:
        for layout in layouts:
            _scan(layout.get("shapes", []))
            if "title" in zones and "subtitle" in zones:
                break
    return zones


# ---------------------------------------------------------------------------
# Markdown summary writer (human-readable)
# ---------------------------------------------------------------------------

def _summarise_shape_for_md(shape, indent=0) -> list:
    """Return one or more markdown table rows for a shape (and its children)."""
    pad = "  " * indent
    name = shape.get("name") or "(unnamed)"
    tag = shape.get("tag")
    xfrm = shape.get("xfrm") or {}
    pos = f"({xfrm.get('left_in')}, {xfrm.get('top_in')}) {xfrm.get('width_in')}×{xfrm.get('height_in')}\"" \
        if xfrm.get("left_in") is not None else "(inherited)"
    fill = shape.get("fill") or {}
    if fill.get("kind") == "solid":
        c = fill.get("color") or {}
        fill_str = f"#{c.get('hex')}" if c.get("hex") else "scheme " + (c.get("scheme") or "?")
    elif fill.get("kind") == "noFill":
        fill_str = "noFill"
    else:
        fill_str = fill.get("kind") or "inherited"
    ph = shape.get("placeholder")
    ph_str = f"ph={ph['type']}/idx={ph['idx']}" if ph else ""
    rows = [f"{pad}- **{tag}** `{name}` {ph_str} — pos {pos} — fill {fill_str}"]
    # Text frame summary
    tf = shape.get("text_frame")
    if tf:
        for para in tf.get("paragraphs", []):
            for run in para.get("runs", []):
                txt = (run.get("text") or "")[:50]
                sz = run.get("size_pt", "inh")
                b = "B" if run.get("bold") else ""
                col = run.get("color") or {}
                col_str = f"#{col.get('hex')}" if col.get("hex") else (col.get("scheme") or "")
                rows.append(f"{pad}    · sz={sz}{b} font={run.get('font') or 'inh'} color={col_str} — {txt!r}")
    for child in shape.get("children", []) or []:
        rows.extend(_summarise_shape_for_md(child, indent + 1))
    return rows


def _write_md(data: dict, out_path: pathlib.Path) -> None:
    lines = [
        f"# Client Template — {data['client']}",
        "",
        f"Extracted: {data['extracted']}  |  Source: `{data['source_pptx']}`",
        "",
        f"Slide size: **{data['slide']['width_in']}\" × {data['slide']['height_in']}\"**",
        "",
        "## Theme — Colours",
        "",
        "| Role | Hex |",
        "|------|-----|",
    ]
    for role, hex_val in data["theme"]["colors"].items():
        lines.append(f"| `{role}` | `#{hex_val}` |")

    lines += ["", "## Theme — Fonts", "", "| Role | Latin | East Asian | Complex |",
              "|------|-------|-----------|---------|"]
    for role, fonts in data["theme"]["fonts"].items():
        lines.append(
            f"| `{role}` | `{fonts.get('latin','')}` "
            f"| `{fonts.get('east_asian','')}` | `{fonts.get('complex_script','')}` |"
        )

    lines += ["", "## Master — Text styles (titleStyle / bodyStyle / otherStyle)", ""]
    for style_name, levels in data["master"].get("text_styles", {}).items():
        lines.append(f"### `{style_name}`")
        lines.append("")
        lines.append("| Level | Size (pt) | Bold | Font | Color |")
        lines.append("|-------|-----------|------|------|-------|")
        for lvl_name, props in levels.items():
            col = props.get("color") or {}
            col_str = f"#{col.get('hex')}" if col.get("hex") else (col.get("scheme") or "")
            lines.append(
                f"| `{lvl_name}` | {props.get('size_pt','')} "
                f"| {'✓' if props.get('bold') else ''} "
                f"| {props.get('font','')} | {col_str} |"
            )
        lines.append("")

    lines += ["", "## Master — Chrome shapes (recursive, with z-order)", ""]
    for shape in data["master"]["shapes"]:
        lines.extend(_summarise_shape_for_md(shape))

    lines += ["", "## Effective zones (title / subtitle anchors)", ""]
    for zone, info in (data.get("effective_zones") or {}).items():
        if info:
            bottom = round((info.get("top_in") or 0) + (info.get("height_in") or 0), 3)
            lines.append(
                f"- **{zone}**: left={info.get('left_in')}\" top={info.get('top_in')}\" "
                f"w={info.get('width_in')}\" h={info.get('height_in')}\" → bottom at {bottom}\""
            )

    lines += ["", f"## Layouts ({len(data['layouts'])})", ""]
    for layout in data["layouts"]:
        lines.append(f"### Layout {layout['index']}: `{layout['name']}`")
        lines.append("")
        for shape in layout["shapes"]:
            lines.extend(_summarise_shape_for_md(shape))
        lines.append("")

    out_path.write_text("\n".join(lines), encoding="utf-8")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Extract complete client template spec from a PPTX. Output goes to the "
                    "project folder (--output) so template.json lives next to the brief, not "
                    "inside the skill directory."
    )
    parser.add_argument("pptx",   help="Path to the client template PPTX")
    parser.add_argument("client", help="Client name (lowercase-hyphenated, e.g. contoso)")
    parser.add_argument("--master", type=int, default=0,
                        help="Which slide master to catalog (default: 0)")
    parser.add_argument("--output", default=None,
                        help="Directory to write template.json and template.md. "
                             "Defaults to the project's clients/ folder inside the skill. "
                             "For new projects use the Claude Projects folder, e.g.: "
                             r'--output "C:\Users\...\Claude Projects\FedEx"')
    args = parser.parse_args()

    pptx_path = pathlib.Path(args.pptx)
    if not pptx_path.exists():
        print(f"ERROR: File not found: {pptx_path}")
        sys.exit(1)

    if args.output:
        client_dir = pathlib.Path(args.output)
    else:
        client_dir = _CLIENTS_DIR / args.client
    client_dir.mkdir(parents=True, exist_ok=True)

    print(f"Reading template: {pptx_path.name}")
    prs = Presentation(str(pptx_path))

    if args.master >= len(prs.slide_masters):
        print(f"ERROR: Template has {len(prs.slide_masters)} master(s); --master {args.master} is out of range.")
        sys.exit(1)

    master = prs.slide_masters[args.master]
    print(f"  Masters: {len(prs.slide_masters)}   "
          f"Layouts in master {args.master}: {len(master.slide_layouts)}")

    # 1. Theme (colours, fonts)
    theme = _extract_theme(pptx_path)
    print(f"  Theme colors: {len(theme['colors'])}   "
          f"Fonts: major={theme['fonts'].get('major', {}).get('latin', '?')}, "
          f"minor={theme['fonts'].get('minor', {}).get('latin', '?')}")
    theme_colors = theme["colors"]

    # 2. Master shapes + text styles (recursive)
    master_data = _master_record(master, theme_colors)
    n_master_shapes = sum(_count_shapes_recursively([s]) for s in master_data["shapes"]) \
        if master_data["shapes"] else 0
    print(f"  Master shapes (incl. group children): {n_master_shapes}")

    # 3. Layout catalog with full shape records
    layouts = [_layout_record(layout, theme_colors, i) for i, layout in enumerate(master.slide_layouts)]
    n_layout_shapes = sum(sum(_count_shapes_recursively([s]) for s in l["shapes"]) for l in layouts)
    print(f"  Layouts: {len(layouts)} (total shapes incl. groups: {n_layout_shapes})")

    # 4. Effective zones (legacy title/subtitle anchors)
    zones = _effective_zones_from_layouts(layouts, master_data["shapes"])

    # Assemble final data
    data = {
        "client":       args.client,
        "extracted":    str(date.today()),
        "source_pptx":  str(pptx_path.resolve()),
        "master_index": args.master,
        "schema_version": 2,   # v2 = full recursive capture (v1 = original partial capture)
        "slide": {
            "width_emu":   int(prs.slide_width),
            "height_emu":  int(prs.slide_height),
            "width_in":    _emu_to_in(int(prs.slide_width)),
            "height_in":   _emu_to_in(int(prs.slide_height)),
        },
        "theme":            theme,
        "master":           master_data,
        "layouts":          layouts,
        "effective_zones":  zones,
    }

    stem      = pptx_path.stem          # e.g. "fedex-template" from "fedex-template.pptx"
    json_path = client_dir / f"{stem}.json"
    md_path   = client_dir / f"{stem}.md"

    json_path.write_text(json.dumps(data, indent=2), encoding="utf-8")
    _write_md(data, md_path)

    print(f"\nSaved:")
    print(f"  {json_path}   ({json_path.stat().st_size:,} bytes)")
    print(f"  {md_path}")
    print(f"\nSchema version: 2 (full recursive shape capture; resolved scheme colours; "
          f"master text styles; layout placeholder fills)")
    print(f"\nTo use in a session: the skill looks for {json_path.name} next to the template")
    print(f"  or in the Claude Project folder. Found = reuse. Not found = re-run this script.")


def _count_shapes_recursively(shape_list):
    n = 0
    for s in shape_list:
        n += 1
        if s.get("children"):
            n += _count_shapes_recursively(s["children"])
    return n


if __name__ == "__main__":
    main()
