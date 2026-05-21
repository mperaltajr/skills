"""
merge_slides.py — combine single-slide temp PPTXs from parallel workers into deck.pptx

Usage:
    py -3 skills/slide-builder/scripts/merge_slides.py \
        --session _session \
        --client-template template.pptx \
        --target _session/deck.pptx

Expects temp files named _session/slide-N-temp.pptx (1-based, any count).
Collects chart-data xlsx sheets from _session/slide-N-temp-chart-data.xlsx if present.
"""

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
import copy
import lxml.etree as etree

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False


def _slide_number_from_name(p: Path) -> int:
    """Extract N from slide-N-temp.pptx."""
    m = re.search(r"slide-(\d+)-temp", p.stem)
    return int(m.group(1)) if m else 999999


def merge(session_dir: Path, template_path: Path, target_path: Path) -> None:
    temp_files = sorted(
        session_dir.glob("slide-*-temp.pptx"),
        key=_slide_number_from_name,
    )
    if not temp_files:
        print("ERROR: no slide-N-temp.pptx files found in session dir", file=sys.stderr)
        sys.exit(1)

    # Start from a clean copy of the template — never build on a prior deck.pptx.
    prs_out = Presentation(str(template_path))
    # Remove all starter slides from the output presentation.
    xml_slides = prs_out.slides._sldIdLst
    for _ in range(len(prs_out.slides)):
        xml_slides.remove(xml_slides[0])

    all_chart_data_paths: list[Path] = []

    for temp_path in temp_files:
        prs_src = Presentation(str(temp_path))
        if not prs_src.slides:
            print(f"WARN: {temp_path.name} has no slides — skipping", file=sys.stderr)
            continue

        # Copy each slide from the temp PPTX into the output presentation.
        # We use the XML deep-copy approach: clone the slide XML and re-register
        # its relationships in the output package.
        for slide_src in prs_src.slides:
            # Add a blank slide to the output using the first layout as a base.
            layout = prs_out.slide_layouts[0]
            slide_out = prs_out.slides.add_slide(layout)

            # Replace the blank slide's spTree with the source slide's spTree.
            sp_tree_src = slide_src.shapes._spTree
            sp_tree_out = slide_out.shapes._spTree

            # Clear all auto-added placeholder shapes from the blank slide.
            for child in list(sp_tree_out):
                tag = etree.QName(child.tag).localname
                if tag != "nvGrpSpPr" and tag != "grpSpPr":
                    sp_tree_out.remove(child)

            # Deep-copy every shape from source slide into output slide.
            for shape_elem in list(sp_tree_src):
                tag = etree.QName(shape_elem.tag).localname
                if tag in ("nvGrpSpPr", "grpSpPr"):
                    continue
                sp_tree_out.append(copy.deepcopy(shape_elem))

            # Copy slide background if present.
            src_bg = slide_src._element.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
            )
            if src_bg is not None:
                out_el = slide_out._element
                existing_bg = out_el.find(
                    ".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
                )
                if existing_bg is not None:
                    out_el.remove(existing_bg)
                out_el.insert(2, copy.deepcopy(src_bg))

        # Collect chart-data xlsx companion if present.
        chart_xlsx = temp_path.with_name(temp_path.stem + "-chart-data.xlsx")
        if chart_xlsx.exists():
            all_chart_data_paths.append(chart_xlsx)

        print(f"  merged: {temp_path.name}")

    target_path.parent.mkdir(parents=True, exist_ok=True)
    prs_out.save(str(target_path))
    print(f"OK: merged {len(temp_files)} slide(s) -> {target_path}")

    # Merge chart-data xlsx files into a single workbook.
    if all_chart_data_paths and _HAS_OPENPYXL:
        merged_xlsx = target_path.with_name(target_path.stem + "-chart-data.xlsx")
        wb_out = openpyxl.Workbook()
        wb_out.remove(wb_out.active)  # remove default Sheet
        for xlsx_path in all_chart_data_paths:
            wb_src = openpyxl.load_workbook(xlsx_path)
            for sheet_name in wb_src.sheetnames:
                ws_src = wb_src[sheet_name]
                # Deduplicate sheet names across slides.
                name = sheet_name
                suffix = 2
                while name in wb_out.sheetnames:
                    name = f"{sheet_name}_{suffix}"
                    suffix += 1
                ws_out = wb_out.create_sheet(title=name)
                for row in ws_src.iter_rows(values_only=True):
                    ws_out.append(list(row))
        wb_out.save(str(merged_xlsx))
        print(f"OK: merged chart-data -> {merged_xlsx}")
    elif all_chart_data_paths:
        print("WARN: openpyxl not installed — chart-data xlsx files not merged. "
              "Install with: pip install openpyxl")


def main() -> None:
    ap = argparse.ArgumentParser(
        description="Merge single-slide temp PPTXs from parallel workers into deck.pptx"
    )
    ap.add_argument("--session", required=True,
                    help="Session directory containing slide-N-temp.pptx files")
    ap.add_argument("--client-template", required=True,
                    help="Client PPTX template (used as the base presentation)")
    ap.add_argument("--target", required=True,
                    help="Output PPTX path")
    args = ap.parse_args()

    session_dir = Path(args.session).expanduser().resolve()
    template_path = Path(args.client_template).expanduser()
    target_path = Path(args.target).expanduser()

    if not session_dir.exists():
        print(f"ERROR: session dir not found: {session_dir}", file=sys.stderr)
        sys.exit(1)
    if not template_path.exists():
        print(f"ERROR: template not found: {template_path}", file=sys.stderr)
        sys.exit(1)

    merge(session_dir, template_path, target_path)


if __name__ == "__main__":
    main()
