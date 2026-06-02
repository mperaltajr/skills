"""
Shared helpers for hand-built PPTX twins.

Every pattern's twin builder imports these. Universal-invariant shapes (chrome,
title block, footer, convergence) get drawn here so each pattern builder only
has to add its body-specific shapes.

All position/size parameters are in CSS pixels (1280x720 canvas). Conversion to
EMU happens inside this module.

Chrome geometry contract (v0.2)
-------------------------------
Title / subtitle / footnote / source / page-number positions come from the
client template's chrome.yml (LayoutChrome model in scripts/_chrome_schema.py),
never from literals in this module. Helpers that draw chrome accept an
explicit `chrome=` kwarg; when None, _resolve_active_chrome() loads from
twins.helpers._ACTIVE_CHROME (set in-process by finalize_deck for the
Mermaid-fallback path) or from SLIDE_LAB_CHROME_YML + SLIDE_LAB_LAYOUT_NAME
env vars (set by finalize_deck before subprocess-invoking option_X.py). If
neither is set, ChromeSidecarMissingError is raised — no silent fallback to
hardcoded constants (that's the v1 slot-mapping bug class; see
feedback_sidecar_fallback_must_be_loud).
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# scripts/ holds _chrome_schema.py — single source of truth for chrome geometry
# and the LayoutChrome / ChromeSidecarMissingError types.
_SCRIPTS_DIR = Path(__file__).resolve().parent.parent / "scripts"
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from _chrome_schema import (  # noqa: E402
    LayoutChrome, ChromeSidecarMissingError,
    canonical_title_box, canonical_subtitle_box,
    canonical_footnote_box, canonical_source_box, canonical_page_number_box,
    CANONICAL_TITLE_X, CANONICAL_TITLE_FONT_PT, CANONICAL_SUBTITLE_FONT_PT,
    CANONICAL_FOOTNOTE_FONT_PX, CANONICAL_SOURCE_FONT_PX,
    CANONICAL_PAGE_NUMBER_FONT_PX,
    CANONICAL_BODY_TOP_Y, CANONICAL_BODY_BOTTOM_Y,
    load_chrome_yml,
)

# --- Brand palette (Accenture Graphik template — ACN Graphik Template.md) ---
# lt2  = #460073  Deep Purple     — dark backgrounds, primary dark fills
BRAND_PRIMARY     = RGBColor(0x46, 0x00, 0x73)
# accent1 = #7500C0  Medium Purple — mid-weight fills, sidebar panels
BRAND_PRIMARY_MID = RGBColor(0x75, 0x00, 0xC0)
# dk2  = #A100FF  Electric Purple — accent moments (one per slide)
BRAND_ACCENT      = RGBColor(0xA1, 0x00, 0xFF)
# accent2 = #C2A3FF  Light Purple  — soft accents, borders
BRAND_ACCENT_SOFT = RGBColor(0xC2, 0xA3, 0xFF)
# dk1  = #000000  Black            — primary body text
TEXT_DARK         = RGBColor(0x00, 0x00, 0x00)
# mid-gray for secondary labels and captions
TEXT_MID          = RGBColor(0x59, 0x59, 0x59)
# light-gray for tertiary labels, rules, footnotes (on LIGHT backgrounds)
TEXT_FAINT        = RGBColor(0x88, 0x88, 0x88)
TEXT_FAINT_ON_LIGHT = TEXT_FAINT
# Faint text variant for DARK backgrounds — needs to be light enough to read
# against light_on_dark layouts (covers, dark hero sections).
TEXT_FAINT_ON_DARK  = RGBColor(0xC8, 0xC8, 0xC8)
SLIDE_BG          = RGBColor(0xFF, 0xFF, 0xFF)
# accent3 = #E6DCFF  Very Light Purple — card / tile fills
CARD_BG           = RGBColor(0xE6, 0xDC, 0xFF)
# accent2 = #C2A3FF  Light Purple — card borders
CARD_BORDER       = RGBColor(0xC2, 0xA3, 0xFF)
DRAFT_BG          = RGBColor(0xFF, 0xD6, 0x00)
DRAFT_TEXT        = RGBColor(0xB7, 0x1C, 0x1C)
WHITE             = RGBColor(0xFF, 0xFF, 0xFF)
# Additional Accenture Graphik accent colors (available but not required on every slide)
ACCENT_PINK       = RGBColor(0xFF, 0x50, 0xA0)   # accent4 #FF50A0
ACCENT_BLUE       = RGBColor(0x22, 0x4B, 0xFF)   # accent5 #224BFF
ACCENT_TEAL       = RGBColor(0x05, 0xF2, 0xDB)   # accent6 #05F2DB

# --- Canvas constants ---
SLIDE_W_PX = 1280
SLIDE_H_PX = 720
SLIDE_W_EMU = Emu(SLIDE_W_PX * 9525)
SLIDE_H_EMU = Emu(SLIDE_H_PX * 9525)

# CSS px → EMU (96dpi assumed for screen pixels; PowerPoint uses 914400 EMU/in)
def px_to_emu(px):
    return Emu(int(px * 9525))

# CSS px → PowerPoint pt (1px @ 96dpi = 0.75pt @ 72dpi)
def px_to_pt(px):
    return Pt(px * 0.75)


def new_slide(prs=None):
    """Create a fresh 1280x720 slide on a blank layout. Returns (prs, slide)."""
    if prs is None:
        prs = Presentation()
        prs.slide_width = SLIDE_W_EMU
        prs.slide_height = SLIDE_H_EMU
    blank_layout = next(
        (l for l in prs.slide_layouts if l.name.strip().lower() == "blank"),
        prs.slide_layouts[6]
    )
    slide = prs.slides.add_slide(blank_layout)
    # Set white background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = SLIDE_BG
    return prs, slide


import re

_INLINE_TAG_RE = re.compile(r"<(strong|em|b|i)>(.*?)</\1>", re.IGNORECASE | re.DOTALL)


def _split_runs(text, *, base_bold=False, base_italic=False,
                base_color=None, emphasis_color=None):
    """Parse `<strong>X</strong>` / `<em>X</em>` (and `<b>` / `<i>`) tags into a
    list of (segment, bold, italic, color) tuples. Other inline content gets
    base styling. Returns [] if input is empty/None.

    `emphasis_color` (optional) tints strong/b segments — e.g., brand-primary
    highlights on a title.
    """
    if not text:
        return []
    runs = []
    pos = 0
    for m in _INLINE_TAG_RE.finditer(text):
        if m.start() > pos:
            runs.append((text[pos:m.start()], base_bold, base_italic, base_color))
        tag = m.group(1).lower()
        seg = m.group(2)
        is_strong = tag in ("strong", "b")
        is_em = tag in ("em", "i")
        runs.append((
            seg,
            base_bold or is_strong,
            base_italic or is_em,
            emphasis_color if (is_strong and emphasis_color) else base_color,
        ))
        pos = m.end()
    if pos < len(text):
        runs.append((text[pos:], base_bold, base_italic, base_color))
    return runs


def add_text(slide, shape_id, text, x_px, y_px, w_px, h_px, *,
             font_size_px=14, font_size_pt=None, color=TEXT_DARK, bold=False, italic=False,
             font_name="Inter", align="left", anchor="top",
             letter_spacing_px=0, uppercase=False, bg_fill=None,
             padding_px=(0, 0, 0, 0), emphasis_color=None):
    """
    Add a text shape with given content. `shape_id` becomes the shape name
    (used by the deck composer to find shapes for text/color substitution).

    If `text` contains `<strong>...</strong>` or `<em>...</em>` tags, the
    text is split into multiple runs. Strong runs get `bold=True` and (if
    provided) `emphasis_color` (typically brand-primary). Em runs get italic.
    Other tags pass through unchanged.
    """
    tb = slide.shapes.add_textbox(
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(w_px), px_to_emu(h_px),
    )
    tb.name = shape_id
    if bg_fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg_fill
        tb.line.fill.background()
    else:
        # Default: no fill, no line
        tb.fill.background()
        tb.line.fill.background()

    tf = tb.text_frame
    tf.word_wrap = True
    pt, pr, pb, pl = padding_px
    tf.margin_top = px_to_emu(pt)
    tf.margin_right = px_to_emu(pr)
    tf.margin_bottom = px_to_emu(pb)
    tf.margin_left = px_to_emu(pl)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    runs = _split_runs(
        text,
        base_bold=bold,
        base_italic=italic,
        base_color=color,
        emphasis_color=emphasis_color,
    )
    if not runs:
        runs = [(text or "", bold, italic, color)]

    for seg, seg_bold, seg_italic, seg_color in runs:
        run = p.add_run()
        run.text = seg.upper() if uppercase else seg
        f = run.font
        f.name = font_name
        f.size = Pt(font_size_pt) if font_size_pt is not None else px_to_pt(font_size_px)
        f.bold = seg_bold
        f.italic = seg_italic
        if seg_color is not None:
            f.color.rgb = seg_color

    return tb


def add_icon(slide, shape_id, x_px, y_px, size_px, glyph, *,
             color=None, font_name="Segoe UI Symbol"):
    """Add a glyph-based icon (Unicode character rendered in a sized text box).

    Examples of usable glyphs (Inter/Segoe-friendly):
      ☰  trigram (menu/lines)        →  arrow-right
      ✦  starburst                    ⊕  circled-plus
      ⊞  squared-plus                 ◇  diamond
      ▲  triangle                     ●  circle
      ✓  check                        ★  star
      ⓘ  info                         ⚙  gear (may not render in all fonts)
      ⚒  hammer-pick

    The icon is centered inside its bounding box. Use Inter or system fallback
    for most glyphs; Segoe UI Symbol gives the broadest coverage on Windows.
    """
    if color is None:
        color = BRAND_ACCENT
    tb = slide.shapes.add_textbox(
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(size_px), px_to_emu(size_px),
    )
    tb.name = shape_id
    tb.fill.background()
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = glyph
    f = run.font
    f.name = font_name
    f.size = px_to_pt(int(size_px * 0.85))
    f.color.rgb = color
    return tb


def add_icon_from_library(slide, shape_id, x_px, y_px, size_px, name, *,
                          color=None):
    """Insert a real vector icon from the Slide Lab icon library.

    PREFER THIS over `add_icon(...glyph=...)` whenever the slide calls for an icon
    (process / data / people / risk / decision / etc.). The library ships 1,143
    pre-extracted icons. Generic Unicode glyphs (☰ ✦ →) are a last resort only.

    The 15 standard icons (use these first; full catalog in icons/icon-index.json):
      gear            — process / operations / workflow
      wrench          — work in progress / tools
      people          — team / workforce / org
      chart-bar       — data / analytics / reporting
      compass         — strategy / direction / vision
      calendar        — timeline / schedule
      coins           — cost / budget / value
      shield-warning  — risk / controls / escalation
      diamond         — decision / approval / governance
      lightbulb       — insight / finding / idea
      globe           — external / market / scale
      clipboard-check — compliance / audit / sign-off
      chip            — technology / systems / AI
      speech          — communication / engagement / change
      package         — delivery / output / shipping

    The icon is tinted with `color` (defaults to BRAND_ACCENT). After the theme
    graft pipeline runs, BRAND_ACCENT gets remapped to the client's accent — so
    icons inherit client brand automatically.

    Example:
        add_icon_from_library(slide, "think-icon",
                              x_px=200, y_px=300, size_px=64,
                              name="lightbulb")

    Falls back to a labeled dashed-border placeholder if the named icon's XML
    doesn't exist in slide-builder/icons/ — the build never errors. v0.1
    ships pre-extracted icons; there is no live extraction step.
    """
    import sys
    from pathlib import Path as _Path
    _SCRIPTS = _Path(__file__).resolve().parent.parent / "scripts"
    if str(_SCRIPTS) not in sys.path:
        sys.path.insert(0, str(_SCRIPTS))
    from icon_helper import insert_icon

    if color is None:
        color = BRAND_ACCENT
    # RGBColor -> hex string for icon_helper
    if hasattr(color, "_RGBColor__valstr"):
        hex_color = "#" + color._RGBColor__valstr.upper()
    else:
        try:
            hex_color = "#{:02X}{:02X}{:02X}".format(*color)
        except Exception:
            hex_color = "#A100FF"
    insert_icon(
        icon_name=name,
        target_slide=slide,
        left_emu=int(px_to_emu(x_px)),
        top_emu=int(px_to_emu(y_px)),
        width_emu=int(px_to_emu(size_px)),
        height_emu=int(px_to_emu(size_px)),
        accent_color=hex_color,
    )
    # The icon is appended directly to spTree as the latest shape; expose it via
    # the standard shapes collection for downstream naming (set on the last shape).
    if len(slide.shapes) > 0:
        try:
            slide.shapes[-1].name = shape_id
        except Exception:
            pass
    return None  # icon_helper returns bool; the wrapper returns None to keep API simple


def add_circle(slide, shape_id, x_px, y_px, diameter_px, fill_color, *,
               no_line=True):
    """Add a filled circle (oval with equal width/height). The standard
    container for an icon when the brief calls for an "icon chip" treatment.

    Circles read as editorial; squares read as web-app tiles. Use circles.

    Example:
        cx, cy, d = 240, 280, 80
        add_circle(slide, "think-bg", cx - d//2, cy - d//2, d, BRAND_PRIMARY)
        add_icon_from_library(slide, "think-icon",
                              cx - 24, cy - 24, 48,
                              name="lightbulb", color=WHITE)
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(diameter_px), px_to_emu(diameter_px),
    )
    shape.name = shape_id
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    return shape


def add_rect(slide, shape_id, x_px, y_px, w_px, h_px, fill_color, *, no_line=True):
    """Add a filled rectangle. Used for rules, accent strips, color bars."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(w_px), px_to_emu(h_px),
    )
    shape.name = shape_id
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    return shape


# ============================================================================
# Active-chrome resolution
# ============================================================================
#
# Helpers that draw chrome (add_title_block, add_footer, add_source,
# add_footnote, add_convergence) read positions and text colors from a
# LayoutChrome instance. Three resolution paths, tried in order:
#
#   1. Explicit `chrome=` kwarg — tests and bespoke callers pass this directly.
#   2. Module attr `_ACTIVE_CHROME` — set by finalize_deck before in-process
#      Mermaid-fallback assembly.
#   3. Env vars `SLIDE_LAB_CHROME_YML` + `SLIDE_LAB_LAYOUT_NAME` — set by
#      finalize_deck before subprocess-launching option_X.py.
#
# If none resolve, ChromeSidecarMissingError is raised. No silent fallback —
# the v1 slot-mapping bug class lived inside that fallback path.

_ACTIVE_CHROME: LayoutChrome | None = None


def set_active_chrome(chrome: LayoutChrome | None) -> None:
    """Inject a LayoutChrome for in-process helper calls. Set by finalize_deck
    before running the Mermaid-fallback assembler; reset to None between
    slides to keep adjacent runs hermetic."""
    global _ACTIVE_CHROME
    _ACTIVE_CHROME = chrome


def _resolve_active_chrome() -> LayoutChrome:
    """Return the LayoutChrome the helpers should draw against.

    Raises ChromeSidecarMissingError when no source resolves — explicit
    fail-loud behavior. The contract test `check_chrome_sidecar_no_silent_fallback`
    asserts this path triggers on a missing sidecar.
    """
    if _ACTIVE_CHROME is not None:
        return _ACTIVE_CHROME
    chrome_yml_env = os.environ.get("SLIDE_LAB_CHROME_YML", "").strip()
    layout_name_env = os.environ.get("SLIDE_LAB_LAYOUT_NAME", "").strip()
    if chrome_yml_env and layout_name_env:
        try:
            spec = load_chrome_yml(Path(chrome_yml_env))
        except ChromeSidecarMissingError:
            raise
        lc = spec.layouts.get(layout_name_env)
        if lc is None:
            raise ChromeSidecarMissingError(
                f"chrome.yml at {chrome_yml_env} does not contain layout "
                f"{layout_name_env!r}. Available: {sorted(spec.layouts)}"
            )
        return lc
    raise ChromeSidecarMissingError(
        "No chrome resolved. Pass chrome= explicitly, call set_active_chrome(), "
        "or set SLIDE_LAB_CHROME_YML + SLIDE_LAB_LAYOUT_NAME env vars."
    )


def _text_colors_for(text_role: str) -> tuple[RGBColor, RGBColor]:
    """Return (primary_text, faint_text) for the given text_role.

    dark_on_light  -> (TEXT_DARK, TEXT_FAINT_ON_LIGHT)
    light_on_dark  -> (WHITE,     TEXT_FAINT_ON_DARK)
    """
    if text_role == "light_on_dark":
        return WHITE, TEXT_FAINT_ON_DARK
    return TEXT_DARK, TEXT_FAINT_ON_LIGHT


def _ensure_bespoke_box(chrome: LayoutChrome, attr_name: str):
    """Return chrome.<attr_name> for a bespoke layout; raise if None.

    Bespoke layouts MUST carry explicit positions for every chrome role they
    use. Missing field means the template wasn't re-registered after a
    layout change — raise loudly per the v0.2 design lock § 2.3.
    """
    box = getattr(chrome, attr_name, None)
    if box is None:
        raise ChromeSidecarMissingError(
            f"bespoke layout {chrome.name!r} has no {attr_name} box in "
            f"chrome.yml. Re-register the template "
            f"(register_template.py propose -> commit)."
        )
    return box


# ============================================================================
# Universal invariants — chrome + title block + footer + convergence
# ============================================================================

def add_chrome(slide):
    """No-op. The top chrome zone is reserved invariant space — only sources/
    footnotes/page numbers belong on the slide, all in the bottom footer.
    Kept as a function so existing builders don't need to be edited.
    """
    pass


# ---------------------------------------------------------------------------
# Cross-skill placeholder contract (F1-B, 2026-05-26)
# ---------------------------------------------------------------------------
# These two strings are INTENTIONAL presenter prompts that render in the
# footer when add_footer() is called with footnote=None or source=None.
# slide-qc imports them by name to allowlist the convention so its hygiene
# pre-pass and vision pass don't flag them as Critical lorem-ipsum residue.
#
# DO NOT change the wording without coordinating with slide-qc — the
# contract is identity-matched, not regex-guessed.
INTENTIONAL_FOOTNOTE_PLACEHOLDER = "[add footnote here or delete]"
INTENTIONAL_SOURCE_PLACEHOLDER = "[add source here or delete]"


def body_zone_for_chrome(chrome: LayoutChrome) -> tuple[int, int]:
    """Return (body_top_y_px, body_bottom_y_px) for the layout's body zone.

    body-canonical (v2 chrome): reads chrome.body_top_y_px / body_bottom_y_px.
      Raises ChromeSidecarMissingError if either is None — the layout was
      registered under v1 and must be re-registered to populate the zone.
    bespoke: returns conservative defaults that don't constrain (40, 660)
      px — covers/dividers compose their own art-directed body.
    """
    if chrome.layout_class == "body-canonical":
        top = getattr(chrome, "body_top_y_px", None)
        bot = getattr(chrome, "body_bottom_y_px", None)
        if top is None or bot is None:
            raise ChromeSidecarMissingError(
                f"body-canonical layout {chrome.name!r} has no body_top_y_px / "
                f"body_bottom_y_px in chrome.yml. Re-register the template "
                f"(register_template.py propose -> commit) so the v2 layout-"
                f"inheritance fields are populated."
            )
        return int(top), int(bot)
    # Bespoke: conservative defaults sourced from _chrome_schema constants
    return 40, CANONICAL_BODY_BOTTOM_Y


def _chrome_box_for(chrome: LayoutChrome, role: str):
    """Resolve a chrome BoxPx for the given role under the layout's class.

    body-canonical -> canonical_*_box() from _chrome_schema
    bespoke        -> chrome.<role> (raises if None)
    """
    if chrome.layout_class == "body-canonical":
        return {
            "title":       canonical_title_box(),
            "subtitle":    canonical_subtitle_box(),
            "footnote":    canonical_footnote_box(),
            "source":      canonical_source_box(),
            "page_number": canonical_page_number_box(),
        }[role]
    return _ensure_bespoke_box(chrome, role)


def add_footer(slide, page_num, source=None, footnote=None, *,
               chrome: LayoutChrome | None = None):
    """Footnote line + source line + page number.

    Positions come from `chrome` (LayoutChrome). When chrome is None,
    _resolve_active_chrome() runs — raises ChromeSidecarMissingError if no
    chrome source is available. No silent fallback to hardcoded literals.

    Invariant zone rule: ONLY sources, footnotes, and the page number may appear
    in the bottom invariant zone. No 'Slide Lab · YEAR · N' branding, no
    'CONFIDENTIAL' tag, no copyright. NO footer-rule divider line either —
    the user flagged that as visual noise.

    BOTH source and footnote lines are ALWAYS drawn — pass real text for each,
    or omit for placeholders. The user is expected to fill or delete in
    PowerPoint; the builder never guesses whether the slide needs them.

    The page-number is suppressed when chrome.has_page_number is False
    (cover-class layouts; replaces slide-qc's page_type == "cover" exemption).
    """
    if chrome is None:
        chrome = _resolve_active_chrome()
    _, faint_color = _text_colors_for(chrome.text_role)

    # v0.3 body-canonical inheritance: layout's own footer/page-number
    # placeholders are populated downstream by finalize_deck. Slide Lab's
    # footnote/source additions still draw at the canonical positions.
    _inherit_chrome = (chrome.layout_class == "body-canonical"
                       and getattr(chrome, "title_placeholder_idx", None) is not None)

    fn_box  = _chrome_box_for(chrome, "footnote")
    src_box = _chrome_box_for(chrome, "source")

    footnote_text = footnote if footnote else INTENTIONAL_FOOTNOTE_PLACEHOLDER
    add_text(
        slide, "footnote-1", f"1. {footnote_text}",
        x_px=fn_box.x_px, y_px=fn_box.y_px,
        w_px=fn_box.w_px, h_px=fn_box.h_px,
        font_size_px=CANONICAL_FOOTNOTE_FONT_PX, color=faint_color,
    )
    source_text = source if source else INTENTIONAL_SOURCE_PLACEHOLDER
    add_text(
        slide, "source", f"Source: {source_text}",
        x_px=src_box.x_px, y_px=src_box.y_px,
        w_px=src_box.w_px, h_px=src_box.h_px,
        font_size_px=CANONICAL_SOURCE_FONT_PX, color=faint_color, italic=True,
    )
    if chrome.has_page_number and not _inherit_chrome:
        pn_box = _chrome_box_for(chrome, "page_number")
        add_text(
            slide, "page-number", f"{page_num}",
            x_px=pn_box.x_px, y_px=pn_box.y_px,
            w_px=pn_box.w_px, h_px=pn_box.h_px,
            font_size_px=CANONICAL_PAGE_NUMBER_FONT_PX, color=faint_color,
            align="right",
        )


def add_title_block(slide, title, subtitle="", *,
                    chrome: LayoutChrome | None = None):
    """Title + optional subtitle. Positions and text colors come from chrome.

    Body-canonical layouts: when the active chrome's layout_class is
    body-canonical AND its title_placeholder_idx is set (v2 chrome), the
    inherited title placeholder will be populated downstream by finalize_deck;
    this helper draws nothing for the title in that case. The subtitle still
    draws as a free-floating textbox at the canonical subtitle position so
    patterns that pass a subtitle continue to render it. (v1 chrome with
    layout_class=body-canonical falls back to drawing the title textbox at
    canonical position, the pre-v0.3 behavior.)

    Bespoke layouts: title drawn at the chrome.title BoxPx position.

    Title is BOTTOM-anchored regardless of chrome's bespoke `anchor` field —
    feedback_title_bottom_anchor invariant: 2-line titles grow upward, never
    displace the subtitle. The chrome placeholder defines WHERE the title
    region lives; Slide Lab decides HOW text fills that region.

    Title supports inline `<strong>X</strong>` for brand-primary emphasis
    (matches every approved pattern's HTML convention).
    """
    if chrome is None:
        chrome = _resolve_active_chrome()
    text_color, _ = _text_colors_for(chrome.text_role)

    # v0.3 body-canonical inheritance: when the layout exposes an inherited
    # title placeholder via chrome (v2 schema), skip drawing the title
    # textbox. finalize_deck will write the title text into the inherited
    # placeholder after grafting.
    if (chrome.layout_class == "body-canonical"
            and getattr(chrome, "title_placeholder_idx", None) is not None):
        # v2.1 (2026-06-02, Bug #2): when the layout ALSO has a subtitle
        # placeholder index, skip drawing the free-floating subtitle textbox
        # too — finalize_deck will populate the subtitle placeholder
        # post-graft, so the subtitle ends up at the layout's intended
        # position instead of the hardcoded canonical box. When the layout
        # has a title placeholder but NO subtitle placeholder, fall back to
        # the legacy free-floating subtitle box (pre-v2.1 behavior).
        if subtitle and getattr(chrome, "subtitle_placeholder_idx", None) is None:
            subtitle_box = _chrome_box_for(chrome, "subtitle")
            add_text(
                slide, "subtitle", subtitle,
                x_px=subtitle_box.x_px, y_px=subtitle_box.y_px,
                w_px=subtitle_box.w_px, h_px=subtitle_box.h_px,
                font_size_pt=subtitle_box.font_pt or CANONICAL_SUBTITLE_FONT_PT,
                color=text_color, italic=True,
            )
        return None

    title_box = _chrome_box_for(chrome, "title")
    add_text(
        slide, "title", title,
        x_px=title_box.x_px, y_px=title_box.y_px,
        w_px=title_box.w_px, h_px=title_box.h_px,
        font_size_pt=title_box.font_pt or CANONICAL_TITLE_FONT_PT,
        color=text_color, bold=True,
        emphasis_color=BRAND_PRIMARY,
        anchor="bottom",
    )

    if subtitle:
        # For bespoke layouts without a subtitle slot, fall back to the
        # canonical subtitle box (positioned relative to canonical title).
        # body-canonical always has subtitle via canonical_subtitle_box.
        if chrome.layout_class == "bespoke" and chrome.subtitle is None:
            return
        subtitle_box = _chrome_box_for(chrome, "subtitle")
        add_text(
            slide, "subtitle", subtitle,
            x_px=subtitle_box.x_px, y_px=subtitle_box.y_px,
            w_px=subtitle_box.w_px, h_px=subtitle_box.h_px,
            font_size_pt=subtitle_box.font_pt or CANONICAL_SUBTITLE_FONT_PT,
            color=text_color, italic=True,
        )


def add_source(slide, source_text, *, chrome: LayoutChrome | None = None):
    """Source citation in the bottom invariant zone. Italic, faint text.
    Builders call this when their content requires citing — NOT auto-emitted
    by add_footer.
    """
    if chrome is None:
        chrome = _resolve_active_chrome()
    _, faint_color = _text_colors_for(chrome.text_role)
    src_box = _chrome_box_for(chrome, "source")
    add_text(
        slide, "source", f"Source: {source_text}",
        x_px=src_box.x_px, y_px=src_box.y_px,
        w_px=src_box.w_px, h_px=src_box.h_px,
        font_size_px=CANONICAL_SOURCE_FONT_PX, color=faint_color, italic=True,
    )


def add_footnote(slide, n, text, *, chrome: LayoutChrome | None = None):
    """Footnote (numbered) in the bottom invariant zone, just above source.
    Builders call this when their content references footnotes — NOT
    auto-emitted by add_footer.
    """
    if chrome is None:
        chrome = _resolve_active_chrome()
    _, faint_color = _text_colors_for(chrome.text_role)
    fn_box = _chrome_box_for(chrome, "footnote")
    add_text(
        slide, f"footnote-{n}", f"{n}. {text}",
        x_px=fn_box.x_px, y_px=fn_box.y_px,
        w_px=fn_box.w_px, h_px=fn_box.h_px,
        font_size_px=CANONICAL_FOOTNOTE_FONT_PX, color=faint_color,
    )


def add_convergence(slide, text, *, bottom_px=78, height_px=42):
    """Brand-primary band at bottom of body. White italic text.

    Body width is derived from CANONICAL_TITLE_X — the band sits inside the
    same horizontal margins as the title. Lives in _chrome_schema so it
    counts as a single-source coordinate reference.
    """
    y = SLIDE_H_PX - bottom_px - height_px
    body_w = SLIDE_W_PX - (2 * CANONICAL_TITLE_X)
    add_rect(
        slide, "convergence-bg",
        x_px=CANONICAL_TITLE_X, y_px=y, w_px=body_w, h_px=height_px,
        fill_color=BRAND_PRIMARY,
    )
    add_text(
        slide, "convergence", text,
        x_px=CANONICAL_TITLE_X, y_px=y, w_px=body_w, h_px=height_px,
        font_size_px=14, color=WHITE, italic=True, bold=False,
        anchor="middle", padding_px=(0, 22, 0, 22),
    )


def add_table(slide, shape_id, x_px, y_px, w_px, h_px, *,
              headers, rows,
              banded=True,
              header_fill=BRAND_PRIMARY,
              header_text=WHITE,
              body_text=TEXT_DARK,
              font_size_px=14,
              font_name="Inter"):
    """Add a native PowerPoint table with mandatory column headers.

    `headers` is REQUIRED — no default. Tables without column labels are
    unreadable; the API enforces what the design rule expects. Pass a list
    of header strings; the column count is derived from that list.

    `rows` is a list of row lists; each row MUST have the same length as
    `headers`. Cells contain plain text only.

    For cell-treatment patterns (harvey balls, RYG pills, arrows, mini-bars,
    sparklines), overlay shapes on top of the table after calling this helper.
    A future version may support cell-type dicts directly — see roadmap.

    Returns the python-pptx `GraphicFrame` containing the table.
    """
    n_cols = len(headers)
    if n_cols == 0:
        raise ValueError("add_table requires at least one header column")
    n_rows = len(rows) + 1  # +1 for the header row

    for i, row in enumerate(rows):
        if len(row) != n_cols:
            raise ValueError(
                f"row {i} has {len(row)} cells, expected {n_cols} to match headers"
            )

    # 14px body-floor enforcement (per final-findings empirical validation —
    # 3 of 20 agents cheated below 14px to fit dense tables; raising here
    # catches the violation at build time regardless of which agent slipped).
    if font_size_px < 14:
        raise ValueError(
            f"add_table font_size_px={font_size_px} is below the 14px body floor. "
            f"If the table doesn't fit at 14px, the content is too dense — shrink "
            f"the dataset, split across two slides, or restructure to a different "
            f"family. Do not cheat the body floor."
        )

    frame = slide.shapes.add_table(
        n_rows, n_cols,
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(w_px), px_to_emu(h_px),
    )
    frame.name = shape_id
    table = frame.table

    def _style_cell(cell, *, text, font_color, bold=False, fill=None):
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        cell.text = str(text)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = px_to_pt(font_size_px)
                run.font.bold = bold
                run.font.color.rgb = font_color

    # Header row
    for c, header in enumerate(headers):
        _style_cell(table.cell(0, c),
                    text=header, font_color=header_text,
                    bold=True, fill=header_fill)

    # Body rows
    for r, row_data in enumerate(rows, start=1):
        for c, val in enumerate(row_data):
            band_fill = CARD_BG if (banded and r % 2 == 0) else WHITE
            _style_cell(table.cell(r, c),
                        text=val, font_color=body_text,
                        bold=False, fill=band_fill)

    return frame
