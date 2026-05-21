"""
Builder for pattern 58: Decision frame — banner + status table + two-option ask.

Title (with right-aligned progress legend) + event banner (brand-primary) +
4-column status table (Initiative · Owner · Status · Comments) with Harvey-ball
progress markers + bottom binary decision ask strip.

Source HTML: _pattern-library/58_decision-frame-banner-table-ask.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from twins.helpers import px_to_emu

ROW_ALT = RGBColor(0xFA, 0xF7, 0xFD)
RULE = CARD_BORDER


def _harvey(slide, name, x, y, size, pct):
    """Draw a Harvey ball at (x, y) of given size. pct: 0, 25, 50, 75, 100."""
    # Empty circle base
    base = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    base.name = f"{name}-base"
    base.fill.solid()
    base.fill.fore_color.rgb = WHITE
    base.line.color.rgb = BRAND_PRIMARY_MID
    base.line.width = 19050

    if pct == 0:
        return
    if pct >= 100:
        base.fill.fore_color.rgb = BRAND_PRIMARY_MID
        return
    # For partial fills, approximate with a smaller circle overlay or pie shapes
    # PIE shape is more accurate but tricky; use overlay rects for quarter steps
    inner = size - 4
    ix = x + 2
    iy = y + 2
    if pct == 25:
        # Right half top quadrant (one quarter of the circle, top-right)
        # Use a PIE shape via PIE
        pie = slide.shapes.add_shape(
            MSO_SHAPE.PIE,
            px_to_emu(ix), px_to_emu(iy),
            px_to_emu(inner), px_to_emu(inner),
        )
        pie.name = f"{name}-fill"
        pie.fill.solid()
        pie.fill.fore_color.rgb = BRAND_PRIMARY_MID
        pie.line.fill.background()
        # default PIE goes from 0° to 90°; rotation alignment is approximate
    elif pct == 50:
        pie = slide.shapes.add_shape(
            MSO_SHAPE.PIE,
            px_to_emu(ix), px_to_emu(iy),
            px_to_emu(inner), px_to_emu(inner),
        )
        pie.name = f"{name}-fill"
        pie.fill.solid()
        pie.fill.fore_color.rgb = BRAND_PRIMARY_MID
        pie.line.fill.background()
    elif pct == 75:
        pie = slide.shapes.add_shape(
            MSO_SHAPE.PIE,
            px_to_emu(ix), px_to_emu(iy),
            px_to_emu(inner), px_to_emu(inner),
        )
        pie.name = f"{name}-fill"
        pie.fill.solid()
        pie.fill.fore_color.rgb = BRAND_PRIMARY_MID
        pie.line.fill.background()


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    # --- Title block ---
    add_text(
        slide, "eyebrow", "DECISION REQUIRED · HTR INITIATIVES",
        x_px=48, y_px=56, w_px=900, h_px=14,
        font_size_px=10, color=BRAND_ACCENT, bold=True,
        letter_spacing_px=2, uppercase=True,
    )
    add_text(
        slide, "title",
        "Accenture rolls off May 31 — <strong>Finance readiness on HTR initiatives</strong> determines coverage risk.",
        x_px=48, y_px=76, w_px=1040, h_px=64,
        font_size_px=22, color=TEXT_DARK, bold=True,
        emphasis_color=BRAND_PRIMARY,
    )
    add_text(
        slide, "subtitle",
        "Event up top, status in the middle, the binary ask at the bottom. Owner column makes the handoff explicit; Harvey balls show current progress.",
        x_px=48, y_px=144, w_px=1040, h_px=36,
        font_size_px=11, color=TEXT_MID,
    )
    add_rect(slide, "brand-rule", x_px=48, y_px=186, w_px=56, h_px=3, fill_color=BRAND_ACCENT)

    # Progress legend — BELOW subheadline (HARD RULE: top-y >= 230, right-aligned to x=1240)
    leg_w = 310
    leg_h = 38
    leg_x = 1240 - leg_w  # right edge at 1240
    leg_y = 230
    leg = add_rect(slide, "legend-bg", leg_x, leg_y, leg_w, leg_h, WHITE)
    leg.line.color.rgb = CARD_BORDER
    leg.line.width = 9525
    add_text(
        slide, "legend-title", "Progress",
        x_px=leg_x + 10, y_px=leg_y, w_px=58, h_px=leg_h,
        font_size_px=8, color=TEXT_MID, bold=True,
        letter_spacing_px=1.2, uppercase=True, anchor="middle",
    )
    add_rect(slide, "legend-divider", leg_x + 68, leg_y + 8, 1, leg_h - 16, RULE)
    legend_items = [(0, "0%"), (25, "~25%"), (50, "~50%"), (75, "~75%"), (100, "Done")]
    # All balls share a single horizontal midline: ball_y is identical for every item.
    ball_size = 14
    ball_y = leg_y + (leg_h - ball_size) // 2  # vertical center of legend row
    # Each item: [ball][label] with label vertically centered on same midline as ball.
    item_x = leg_x + 78
    item_step = 46  # ball(14) + gap(4) + label(~24) + gap(4)
    for i, (pct, lbl) in enumerate(legend_items):
        n = i + 1
        _harvey(slide, f"legend-{n}-ball", item_x, ball_y, ball_size, pct)
        add_text(
            slide, f"legend-{n}-lbl", lbl,
            x_px=item_x + ball_size + 2, y_px=ball_y, w_px=item_step - ball_size - 4, h_px=ball_size,
            font_size_px=8, color=TEXT_DARK, bold=True, anchor="middle",
        )
        item_x += item_step

    # --- Event banner --- (shifted down to clear legend bottom ~268)
    banner_y = 278
    banner_h = 50
    add_rect(slide, "event-banner-bg", 48, banner_y, 1280 - 96, banner_h, BRAND_PRIMARY)
    add_rect(slide, "event-banner-accent", 48, banner_y, 4, banner_h, BRAND_ACCENT)
    add_text(
        slide, "event-flag", "⚑ Project rolloff",
        x_px=70, y_px=banner_y, w_px=200, h_px=banner_h,
        font_size_px=12, color=BRAND_ACCENT, bold=True, anchor="middle",
        letter_spacing_px=1.5, uppercase=True,
    )
    add_text(
        slide, "event-date", "May 31, 2026",
        x_px=280, y_px=banner_y, w_px=220, h_px=banner_h,
        font_size_px=18, color=WHITE, bold=True, anchor="middle",
    )
    add_text(
        slide, "event-note",
        "Finance committed to own remaining HTR initiatives post-rolloff",
        x_px=720, y_px=banner_y, w_px=460, h_px=banner_h,
        font_size_px=11, color=BRAND_ACCENT_SOFT, anchor="middle", align="right",
    )

    # --- Status table --- (shifted to clear banner at y=278, h=50, bottom=328)
    table_top = 338
    table_left = 48
    table_w = 1280 - 96
    col_pct = [0.36, 0.18, 0.12, 0.34]
    col_widths = [int(table_w * p) for p in col_pct]
    col_x = [table_left]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)

    # Header
    header_h = 32
    add_rect(slide, "table-header-bg", table_left, table_top, table_w, header_h, BRAND_PRIMARY)
    headers = ["Initiative", "Owner", "Status", "Comments"]
    for i, h in enumerate(headers):
        n = i + 1
        align = "center" if i == 2 else "left"
        add_text(
            slide, f"table-header-{n}", h,
            x_px=col_x[i] + 12, y_px=table_top, w_px=col_widths[i] - 24, h_px=header_h,
            font_size_px=10, color=WHITE, bold=True, anchor="middle",
            letter_spacing_px=0.6, uppercase=False, align=align,
        )

    rows_data = [
        ("Background Checks Cost Reduction", "Client Finance", 25, "~25%", "Vendor RFP issued; awaiting bids"),
        ("Cost Per Hire Optimization", "Client Finance", 0, "Not started", "Owner unconfirmed — needs assignment"),
        ("Sourcing channel rationalization", "Client Finance", 50, "~50%", "Pilot underway in 2 regions"),
        ("Onboarding automation", "Client Finance", 0, "Not started", "Awaiting tool selection"),
        ("Vendor consolidation (Tier 2)", "Client Finance", 25, "~25%", "Inventory complete; negotiation pending"),
    ]
    row_h = 40
    body_top = table_top + header_h
    for i, (init, owner, pct, plabel, comment) in enumerate(rows_data):
        n = i + 1
        ry = body_top + i * row_h
        row_bg = ROW_ALT if i % 2 == 1 else WHITE
        add_rect(slide, f"table-row-{n}-bg", table_left, ry, table_w, row_h, row_bg)
        add_rect(slide, f"table-row-{n}-rule", table_left, ry + row_h, table_w, 1, RULE)
        # Initiative
        add_text(
            slide, f"table-row-{n}-init", init,
            x_px=col_x[0] + 12, y_px=ry, w_px=col_widths[0] - 24, h_px=row_h,
            font_size_px=12, color=TEXT_DARK, anchor="middle",
        )
        # Owner
        add_text(
            slide, f"table-row-{n}-owner", owner,
            x_px=col_x[1] + 12, y_px=ry, w_px=col_widths[1] - 24, h_px=row_h,
            font_size_px=11, color=TEXT_MID, anchor="middle",
        )
        # Status: harvey ball stacked above its label, all rows share the SAME ball top-y
        # so balls line up on a single horizontal midline across rows.
        row_ball_size = 16
        ball_x = col_x[2] + (col_widths[2] - row_ball_size) // 2
        ball_y = ry + 8  # constant offset from row top -> balls share horizontal midline
        _harvey(slide, f"table-row-{n}-ball", ball_x, ball_y, row_ball_size, pct)
        add_text(
            slide, f"table-row-{n}-status-label", plabel,
            x_px=col_x[2], y_px=ball_y + row_ball_size + 2, w_px=col_widths[2], h_px=10,
            font_size_px=7, color=TEXT_MID, bold=True, align="center",
            letter_spacing_px=0.6, uppercase=True,
        )
        # Comments
        add_text(
            slide, f"table-row-{n}-comment", comment,
            x_px=col_x[3] + 12, y_px=ry, w_px=col_widths[3] - 24, h_px=row_h,
            font_size_px=10, color=TEXT_MID, italic=True, anchor="middle",
        )

    # --- Decision ask strip --- (table bottom = 338 + 32 + 5*40 = 570, leave a small gap)
    ask_y = 580
    ask_h = 56
    ask_x = 48
    ask_w = 1280 - 96
    ask_box = add_rect(slide, "ask-strip-bg", ask_x, ask_y, ask_w, ask_h, CARD_BG)
    ask_box.line.color.rgb = BRAND_PRIMARY
    ask_box.line.width = 19050

    add_text(
        slide, "ask-label", "Direction needed",
        x_px=ask_x + 18, y_px=ask_y, w_px=120, h_px=ask_h,
        font_size_px=11, color=BRAND_PRIMARY, bold=True,
        letter_spacing_px=1.2, uppercase=True, anchor="middle",
    )

    # Two option boxes flanking an "or" divider
    opt_x = ask_x + 150
    or_w = 40
    opt_w = (ask_w - 150 - or_w - 32) // 2
    # Option 1 (primary)
    opt1 = add_rect(slide, "ask-option-1-bg", opt_x, ask_y + 10, opt_w, ask_h - 20, WHITE)
    opt1.line.color.rgb = BRAND_PRIMARY
    opt1.line.width = 19050
    add_text(
        slide, "ask-option-1-title",
        "Continue Accenture support past May 31",
        x_px=opt_x + 8, y_px=ask_y + 14, w_px=opt_w - 16, h_px=20,
        font_size_px=11, color=BRAND_PRIMARY, bold=True, align="center",
    )
    add_text(
        slide, "ask-option-1-sub",
        "Define scope · duration · revised SOW",
        x_px=opt_x + 8, y_px=ask_y + 32, w_px=opt_w - 16, h_px=14,
        font_size_px=10, color=TEXT_MID, align="center",
    )
    # OR divider
    add_text(
        slide, "or-divider", "or",
        x_px=opt_x + opt_w, y_px=ask_y, w_px=or_w, h_px=ask_h,
        font_size_px=10, color=TEXT_MID, bold=True, align="center", anchor="middle",
        letter_spacing_px=1.5, uppercase=True,
    )
    # Option 2 (alt)
    opt2_x = opt_x + opt_w + or_w
    opt2 = add_rect(slide, "ask-option-2-bg", opt2_x, ask_y + 10, opt_w, ask_h - 20, ROW_ALT)
    opt2.line.color.rgb = CARD_BORDER
    opt2.line.width = 19050
    add_text(
        slide, "ask-option-2-title",
        "Transition to Finance as planned",
        x_px=opt2_x + 8, y_px=ask_y + 14, w_px=opt_w - 16, h_px=20,
        font_size_px=11, color=TEXT_MID, bold=True, align="center",
    )
    add_text(
        slide, "ask-option-2-sub",
        "Confirm Finance readiness by May 31",
        x_px=opt2_x + 8, y_px=ask_y + 32, w_px=opt_w - 16, h_px=14,
        font_size_px=10, color=TEXT_MID, align="center",
    )

    add_footer(slide, page_num=58)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "58_decision-frame-banner-table-ask.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
