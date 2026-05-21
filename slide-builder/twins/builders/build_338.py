"""
Builder for pattern 338: 8-bucket octagon radial.

Left panel: 8 numbered circular nodes positioned at octagon vertices with a CORE
ovular center, spokes radiating out from CORE. python-pptx cannot render an
8-sided polygon outline cleanly, so we draw the 8 nodes + spokes only (no
octagon edge), approximating the visual.

Right panel: 8 numbered annotation cards in a vertical "workstream index" list.

Source HTML: _pattern-library/338_8bucket-octagon-radial.html
"""
from pathlib import Path
import sys
import math

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
    px_to_emu,
)
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


def _circle(slide, name, x, y, size, fill, *, text=None, text_color=WHITE, font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    shape.name = f"{name}-bg"
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text is not None:
        add_text(slide, f"{name}-text", text,
                 x_px=x, y_px=y, w_px=size, h_px=size,
                 font_size_px=font_size, color=text_color, bold=True,
                 align="center", anchor="middle")


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_title_block(
        slide,
        title="Eight strategic workstreams, <strong>one core ambition</strong>",
        subtitle="Each bucket connects to the centre — progress on any node strengthens the whole system.",
        title_x=48, title_y=44, title_w=1184, title_h=44,
        subtitle_h=20, brand_rule_w=64,
    )

    body_top = 200
    body_left = 48
    body_w = 1184
    body_h = 720 - 32 - 14 - body_top

    # Left panel: octagonal radial
    left_w = 540
    cx = body_left + left_w // 2
    cy = body_top + body_h // 2
    radius = 160
    node_size = 48
    core_size = 84

    # 8 vertices, starting at top going clockwise
    nodes = []
    for i in range(8):
        angle = -math.pi / 2 + i * (2 * math.pi / 8)
        nx = cx + radius * math.cos(angle) - node_size / 2
        ny = cy + radius * math.sin(angle) - node_size / 2
        nodes.append((nx, ny))

    # Spokes
    for i, (nx, ny) in enumerate(nodes):
        n = i + 1
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(nx + node_size / 2), px_to_emu(ny + node_size / 2),
            px_to_emu(cx), px_to_emu(cy),
        )
        connector.name = f"spoke-{n}"
        connector.line.color.rgb = BRAND_ACCENT_SOFT
        connector.line.width = 12700

    # Octagon edges (connect adjacent nodes)
    for i in range(8):
        n = i + 1
        x1, y1 = nodes[i]
        x2, y2 = nodes[(i + 1) % 8]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1 + node_size / 2), px_to_emu(y1 + node_size / 2),
            px_to_emu(x2 + node_size / 2), px_to_emu(y2 + node_size / 2),
        )
        conn.name = f"octagon-edge-{n}"
        conn.line.color.rgb = CARD_BORDER
        conn.line.width = 9525

    # Node labels (short)
    node_labels = [
        "Vision", "Strategy", "Operations", "Technology",
        "Talent", "Culture", "Governance", "Customers",
    ]
    for i, ((nx, ny), label) in enumerate(zip(nodes, node_labels)):
        n = i + 1
        _circle(slide, f"node-{n}",
                nx, ny, node_size, BRAND_PRIMARY_MID,
                text=str(n), text_color=WHITE, font_size=16)
        # Label below or above node depending on position
        # top node (i==0): label above; bottom (i==4): label below; sides: outside ring
        # Default place label below node
        label_w = 110
        lx = nx + node_size / 2 - label_w / 2
        # Use vertical position relative to center to decide
        if ny + node_size / 2 < cy - 20:
            ly = ny - 22
        else:
            ly = ny + node_size + 4
        add_text(slide, f"node-{n}-label", label,
                 x_px=int(lx), y_px=int(ly), w_px=label_w, h_px=18,
                 font_size_px=11, color=BRAND_PRIMARY, bold=True, align="center")

    # Core
    _circle(slide, "core",
            cx - core_size // 2, cy - core_size // 2, core_size,
            BRAND_ACCENT, text="CORE", text_color=WHITE, font_size=12)
    add_text(slide, "core-sub", "ambition",
             x_px=cx - core_size // 2, y_px=cy + 6, w_px=core_size, h_px=18,
             font_size_px=9, color=WHITE, align="center", italic=True)

    # Right panel: 8 annotation rows ("workstream index")
    right_x = body_left + left_w + 24
    right_w = body_w - left_w - 24

    add_text(slide, "index-header", "WORKSTREAM INDEX",
             x_px=right_x, y_px=body_top + 4, w_px=right_w, h_px=16,
             font_size_px=10, color=BRAND_PRIMARY, bold=True, letter_spacing_px=1.6)
    add_rect(slide, "index-rule", right_x, body_top + 22, right_w, 2, BRAND_ACCENT)

    items = [
        ("Vision", "Define and socialise the north star that anchors all decisions."),
        ("Strategy", "Set clear choices on where to play and how to win."),
        ("Operations", "Streamline processes that deliver day-to-day value."),
        ("Technology", "Build the platforms and data backbone for scale."),
        ("Talent", "Attract, grow, and retain the capabilities needed."),
        ("Culture", "Shape the behaviours and norms that make change stick."),
        ("Governance", "Establish accountabilities and decision-making authority."),
        ("Customers", "Keep the end-user experience at the centre of every choice."),
    ]
    rows_top = body_top + 32
    rows_h = body_h - 32
    row_h = rows_h // 8
    for i, (title, desc) in enumerate(items):
        n = i + 1
        ry = rows_top + i * row_h
        # Number circle
        ns = 22
        _circle(slide, f"ann-{n}-num",
                right_x, ry + (row_h - ns) // 2, ns, BRAND_PRIMARY_MID,
                text=str(n), text_color=WHITE, font_size=10)
        # Title
        add_text(slide, f"ann-{n}-title", title,
                 x_px=right_x + 32, y_px=ry + 4, w_px=120, h_px=row_h - 8,
                 font_size_px=11, color=BRAND_PRIMARY, bold=True, anchor="middle")
        # Desc
        add_text(slide, f"ann-{n}-desc", desc,
                 x_px=right_x + 160, y_px=ry + 4, w_px=right_w - 168, h_px=row_h - 8,
                 font_size_px=10, color=TEXT_MID, anchor="middle")

    add_footer(slide, page_num=338)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "338_8bucket-octagon-radial.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
