"""
Builder for pattern 79: ROI calculator.

Three columns with arrow connectors: INPUTS · CALCULATION · RESULT.
- Inputs: tinted card with assumption rows + investment summary row
- Calculation: stacked formula cards with caption labels
- Result: dark brand-primary card with massive ROI number, payback row, breakdown list
Convergence band at bottom.

Source HTML: _pattern-library/79_roi-calculator.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_convergence,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from twins.helpers import px_to_emu

WHITE_70 = RGBColor(0xD8, 0xCD, 0xE5)
WHITE_55 = RGBColor(0xBF, 0xB2, 0xCE)


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    # --- Title block ---
    add_text(
        slide, "title",
        "Slide Lab pilot ROI — <strong>24× return, payback in two weeks.</strong>",
        x_px=64, y_px=50, w_px=1152, h_px=44,
        font_size_px=24, color=TEXT_DARK, bold=True,
        emphasis_color=BRAND_PRIMARY,
    )
    add_text(
        slide, "subtitle",
        "Conservative assumptions: 30 senior consultants, 50 decks each per year, 9 hours saved per deck at $150/hr loaded rate.",
        x_px=64, y_px=98, w_px=1152, h_px=22,
        font_size_px=12, color=TEXT_MID, italic=True,
    )
    add_rect(slide, "brand-rule", x_px=64, y_px=132, w_px=56, h_px=3, fill_color=BRAND_ACCENT)

    # --- Body grid: inputs | arrow | calc | arrow | result ---
    body_top = 196
    body_bottom = 600
    body_h = body_bottom - body_top
    body_left = 64
    body_right = 1280 - 64
    body_w = body_right - body_left
    # Proportions: 27 | 4 | 38 | 4 | 35 = 108
    unit = body_w / 108
    inputs_w = int(unit * 27)
    arrow_w = int(unit * 4)
    calc_w = int(unit * 38)
    result_w = body_w - inputs_w - arrow_w - calc_w - arrow_w

    inputs_x = body_left
    arrow1_x = inputs_x + inputs_w
    calc_x = arrow1_x + arrow_w
    arrow2_x = calc_x + calc_w
    result_x = arrow2_x + arrow_w

    # --- Zone labels ---
    zone_label_h = 20
    zone_y = body_top
    for sid, label, zx, zw in [
        ("inputs-label", "Inputs · assumptions", inputs_x, inputs_w),
        ("calc-label", "Calculation · the math", calc_x, calc_w),
        ("result-label", "Result · the return", result_x, result_w),
    ]:
        add_text(
            slide, sid, label,
            x_px=zx, y_px=zone_y, w_px=zw, h_px=zone_label_h,
            font_size_px=9, color=TEXT_FAINT, bold=True,
            letter_spacing_px=2, uppercase=True,
        )

    content_top = zone_y + zone_label_h + 8
    content_h = body_h - (content_top - zone_y)

    # === INPUTS card ===
    inputs_card = add_rect(slide, "inputs-bg", inputs_x, content_top, inputs_w - 14, content_h, CARD_BG)
    inputs_card.line.color.rgb = CARD_BORDER
    inputs_card.line.width = 9525
    add_rect(slide, "inputs-accent", inputs_x, content_top, 3, content_h, BRAND_ACCENT_SOFT)

    input_items = [
        ("Senior consultants in scope", "30", False),
        ("Decks per consultant per year", "50", False),
        ("Hours saved per deck", "9", False),
        ("Loaded hourly rate", "$150", False),
        ("Pilot investment (Y1)", "$80K", True),
    ]
    pad_l = 18
    pad_r = 18
    inner_x = inputs_x + pad_l
    inner_w = inputs_w - 14 - pad_l - pad_r
    item_h = (content_h - 32) // len(input_items)
    for i, (label, value, is_invest) in enumerate(input_items):
        n = i + 1
        iy = content_top + 16 + i * item_h
        if is_invest:
            # Top rule (solid)
            add_rect(slide, f"input-{n}-rule", inner_x, iy - 4, inner_w, 1, CARD_BORDER)
            lbl_color = BRAND_PRIMARY
            lbl_bold = True
            val_color = BRAND_ACCENT
        else:
            lbl_color = TEXT_MID
            lbl_bold = False
            val_color = BRAND_PRIMARY
            # Bottom dashed-ish rule (use solid faint)
            if i < len(input_items) - 1:
                add_rect(slide, f"input-{n}-rule", inner_x, iy + item_h - 4, inner_w, 1, CARD_BORDER)
        add_text(
            slide, f"input-{n}-label", label,
            x_px=inner_x, y_px=iy, w_px=int(inner_w * 0.62), h_px=item_h - 8,
            font_size_px=11, color=lbl_color, bold=lbl_bold, anchor="middle",
        )
        add_text(
            slide, f"input-{n}-value", value,
            x_px=inner_x + int(inner_w * 0.62), y_px=iy, w_px=int(inner_w * 0.38), h_px=item_h - 8,
            font_size_px=14, color=val_color, bold=True, align="right", anchor="middle",
        )

    # === Arrow 1: inputs -> calc === (proper RIGHT_ARROW, no rotation)
    arrow_y = content_top + content_h // 2
    # Horizontal arrow body spans the full gap between cards; centerline = arrow_y
    head_h = 22
    arr1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        px_to_emu(arrow1_x + 2), px_to_emu(arrow_y - head_h // 2),
        px_to_emu(arrow_w - 4), px_to_emu(head_h),
    )
    arr1.name = "arrow-1-head"
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = BRAND_ACCENT
    arr1.line.fill.background()
    arr1.rotation = 0

    # === CALCULATION column ===
    calc_steps = [
        ("30 × 50 = 1,500", "Decks produced per year", False),
        ("1,500 × 9 = 13,500 hr", "Consultant hours saved per year", False),
        ("13,500 × $150 = $2.025M", "Gross value created per year", False),
        ("$2.025M − $80K = $1.945M", "Net value, Year 1", True),
    ]
    step_h = (content_h - 16) // len(calc_steps)
    step_inner_h = step_h - 8
    for i, (formula, caption, is_net) in enumerate(calc_steps):
        n = i + 1
        sy = content_top + i * step_h
        step_card = add_rect(slide, f"calc-{n}-bg", calc_x + 6, sy, calc_w - 12, step_inner_h, WHITE)
        step_card.line.color.rgb = CARD_BORDER
        step_card.line.width = 9525
        result_color = BRAND_ACCENT if is_net else BRAND_PRIMARY
        add_text(
            slide, f"calc-{n}-formula", formula,
            x_px=calc_x + 18, y_px=sy + 6, w_px=calc_w - 36, h_px=22,
            font_size_px=13, color=BRAND_PRIMARY, bold=True,
            emphasis_color=result_color,
        )
        add_text(
            slide, f"calc-{n}-caption", caption,
            x_px=calc_x + 18, y_px=sy + 30, w_px=calc_w - 36, h_px=14,
            font_size_px=9, color=TEXT_FAINT, bold=True,
            letter_spacing_px=0.8, uppercase=True,
        )
        # Down-arrow between steps — proper DOWN_ARROW, no rotation
        if i < len(calc_steps) - 1:
            d_w = 14
            d_h = step_h - step_inner_h - 2  # fill the gap between step cards
            if d_h < 8:
                d_h = 8
            darrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                px_to_emu(calc_x + calc_w // 2 - d_w // 2), px_to_emu(sy + step_inner_h + 1),
                px_to_emu(d_w), px_to_emu(d_h),
            )
            darrow.name = f"calc-{n}-arrow"
            darrow.fill.solid()
            darrow.fill.fore_color.rgb = BRAND_ACCENT_SOFT
            darrow.line.fill.background()
            darrow.rotation = 0

    # === Arrow 2: calc -> result === (proper RIGHT_ARROW, no rotation)
    arr2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        px_to_emu(arrow2_x + 2), px_to_emu(arrow_y - head_h // 2),
        px_to_emu(arrow_w - 4), px_to_emu(head_h),
    )
    arr2.name = "arrow-2-head"
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = BRAND_ACCENT
    arr2.line.fill.background()
    arr2.rotation = 0

    # === RESULT card (dark) ===
    rx = result_x + 14
    rw = result_w - 14
    rcard = add_rect(slide, "result-bg", rx, content_top, rw, content_h, BRAND_PRIMARY)
    add_rect(slide, "result-accent", rx + rw - 4, content_top, 4, content_h, BRAND_ACCENT)

    add_text(
        slide, "result-eyebrow", "Year 1 ROI",
        x_px=rx + 22, y_px=content_top + 18, w_px=rw - 44, h_px=14,
        font_size_px=9, color=BRAND_ACCENT_SOFT, bold=True,
        letter_spacing_px=2, uppercase=True,
    )
    add_text(
        slide, "result-roi-number", "24×",
        x_px=rx + 22, y_px=content_top + 38, w_px=rw - 44, h_px=86,
        font_size_px=64, color=WHITE, bold=True,
    )
    add_text(
        slide, "result-roi-label",
        "net return on $80K invested",
        x_px=rx + 22, y_px=content_top + 128, w_px=rw - 44, h_px=18,
        font_size_px=11, color=WHITE_70,
    )
    add_rect(slide, "result-divider", rx + 22, content_top + 156, rw - 44, 1, WHITE_55)

    # Payback row
    add_text(
        slide, "payback-value", "~2 wks",
        x_px=rx + 22, y_px=content_top + 168, w_px=80, h_px=22,
        font_size_px=18, color=BRAND_ACCENT_SOFT, bold=True,
    )
    add_text(
        slide, "payback-label",
        "payback period on cumulative savings",
        x_px=rx + 108, y_px=content_top + 172, w_px=rw - 130, h_px=22,
        font_size_px=10, color=WHITE_70,
    )

    # Breakdown
    bd_top = content_top + content_h - 88
    add_text(
        slide, "breakdown-title", "Breakdown",
        x_px=rx + 22, y_px=bd_top, w_px=rw - 44, h_px=14,
        font_size_px=8, color=WHITE_55, bold=True,
        letter_spacing_px=1.6, uppercase=True,
    )
    breakdown_items = [
        ("Gross value (Y1)", "$2.025M"),
        ("Pilot investment", "$80K"),
        ("Net value (Y1)", "$1.945M"),
    ]
    for i, (lbl, val) in enumerate(breakdown_items):
        n = i + 1
        by = bd_top + 20 + i * 18
        add_text(
            slide, f"breakdown-{n}-lbl", lbl,
            x_px=rx + 22, y_px=by, w_px=int(rw * 0.55), h_px=16,
            font_size_px=10, color=WHITE_70,
        )
        add_text(
            slide, f"breakdown-{n}-val", val,
            x_px=rx + rw - 100, y_px=by, w_px=80, h_px=16,
            font_size_px=10, color=WHITE, bold=True, align="right",
        )

    # --- Convergence band (custom: lighter purple band) ---
    conv_y = 614
    conv_h = 34
    add_rect(slide, "convergence-bg", 64, conv_y, 1280 - 128, conv_h, BRAND_PRIMARY)
    add_text(
        slide, "convergence",
        "The pilot pays for itself inside the first month — everything after that is <strong>net upside</strong>, and the model compounds as adoption widens.",
        x_px=80, y_px=conv_y, w_px=1280 - 128 - 32, h_px=conv_h,
        font_size_px=12, color=WHITE, italic=True, anchor="middle",
        emphasis_color=BRAND_ACCENT_SOFT,
    )

    add_footer(slide, page_num=79)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "79_roi-calculator.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
