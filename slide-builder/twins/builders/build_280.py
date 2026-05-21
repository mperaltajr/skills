"""
Builder for pattern 280: Customer Onboarding Journey (5 stages).

Source HTML: _pattern-library/280_customer-onboarding-journey.html

4 rows: stage headers / touchpoints / emotion line / opportunities.
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor

OPP_BG = RGBColor(0xF0, 0xFD, 0xF4)
OPP_BORDER = RGBColor(0xA7, 0xE5, 0xBF)
OPP_GREEN = RGBColor(0x16, 0xA3, 0x4A)
PAIN_BG = RGBColor(0xFE, 0xF2, 0xF2)
PAIN_BORDER = RGBColor(0xFE, 0xCA, 0xCA)
PAIN_RED = RGBColor(0xDC, 0x26, 0x26)
NEU_BG = RGBColor(0xF8, 0xF4, 0xFC)


def build():
    prs, slide = new_slide()
    add_chrome(slide)
    add_title_block(
        slide,
        title="Customer <strong>Onboarding Journey</strong> — Five Stages to Advocacy",
        subtitle="Touchpoints, emotional arc, and improvement opportunities across the full onboarding lifecycle.",
    )

    # Body area
    body_x = 56
    body_y = 234
    body_w = 1280 - 112
    body_h = 720 - body_y - 44

    stages = ["Sign-Up", "Welcome & Setup", "First Use", "Habit Formation", "Advocacy"]
    n_stages = 5
    gap = 8
    col_w = (body_w - gap * (n_stages - 1)) // n_stages

    # ROW 1: Stage headers
    hdr_y = body_y
    hdr_h = 38
    for i, name in enumerate(stages):
        sx = body_x + i * (col_w + gap)
        add_rect(slide, f"stage-{i+1}-header", sx, hdr_y, col_w, hdr_h, BRAND_PRIMARY)
        add_text(
            slide, f"stage-{i+1}-num", f"0{i+1}",
            x_px=sx + 10, y_px=hdr_y, w_px=30, h_px=hdr_h,
            font_size_px=10, color=BRAND_ACCENT_SOFT, bold=True, anchor="middle",
            letter_spacing_px=1.2,
        )
        add_text(
            slide, f"stage-{i+1}-name", name,
            x_px=sx + 42, y_px=hdr_y, w_px=col_w - 50, h_px=hdr_h,
            font_size_px=11, color=WHITE, bold=True, anchor="middle",
            uppercase=True, letter_spacing_px=0.6,
        )

    # ROW 2: Touchpoints (3 per stage)
    tp_y = hdr_y + hdr_h + 10
    tp_h = 130
    touchpoints = [
        ["Signup page", "Email confirmation", "Welcome email"],
        ["Setup wizard", "Configuration call", "Resource portal"],
        ["Tutorial", "Support chat", "Check-in email"],
        ["Weekly digest", "Usage tips", "Success team call"],
        ["NPS survey", "Case study ask", "Referral program"],
    ]
    for i, items in enumerate(touchpoints):
        sx = body_x + i * (col_w + gap)
        sub_h = (tp_h - 2 * 6) // 3
        for ji, item in enumerate(items):
            iy = tp_y + ji * (sub_h + 6)
            box = add_rect(slide, f"tp-{i+1}-{ji+1}-bg", sx, iy, col_w, sub_h, CARD_BG)
            box.line.color.rgb = CARD_BORDER
            box.line.width = 9525
            # dot
            add_rect(slide, f"tp-{i+1}-{ji+1}-dot",
                     sx + 10, iy + sub_h // 2 - 3, 6, 6, BRAND_ACCENT_SOFT)
            add_text(
                slide, f"tp-{i+1}-{ji+1}-text", item,
                x_px=sx + 22, y_px=iy, w_px=col_w - 28, h_px=sub_h,
                font_size_px=11, color=TEXT_DARK, anchor="middle",
            )

    # ROW 3: Emotion line — simulated by colored circles at heights, with connecting rule
    emo_y = tp_y + tp_h + 12
    emo_h = 80
    # Baseline track
    add_rect(slide, "emotion-baseline", body_x, emo_y + emo_h - 12, body_w, 1, CARD_BORDER)

    # Emotion scores per stage (positive = higher y, sad = lower)
    # score -> position from top: higher score = higher (smaller y)
    scores = [40, 55, 35, 65, 80]
    faces = ["😐", "😊", "😤", "😌", "🤩"]
    is_pain = [False, False, True, False, False]

    # Compute dot positions
    centers_x = [body_x + i * (col_w + gap) + col_w // 2 for i in range(n_stages)]
    # Map score 20..85 to y emo_y+10..emo_y+emo_h-20 (inverted)
    def y_for(score):
        return int(emo_y + 10 + (1 - (score - 20) / 65) * (emo_h - 30))

    points_y = [y_for(s) for s in scores]

    # Draw connecting line as a series of thin rects between centers
    from pptx.enum.shapes import MSO_SHAPE
    for i in range(n_stages - 1):
        # Simple straight rect approximation between two points
        x1, y1 = centers_x[i], points_y[i]
        x2, y2 = centers_x[i + 1], points_y[i + 1]
        # Use a connector line
        conn = slide.shapes.add_connector(1, x1 * 9525, y1 * 9525, x2 * 9525, y2 * 9525)
        conn.name = f"emotion-line-{i+1}"
        conn.line.color.rgb = BRAND_ACCENT
        conn.line.width = 25400  # ~2.5pt

    # Dots + faces + score labels
    for i, (cx, py, face, pain) in enumerate(zip(centers_x, points_y, faces, is_pain)):
        dot_color = PAIN_RED if pain else BRAND_ACCENT
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            (cx - 6) * 9525, (py - 6) * 9525, 12 * 9525, 12 * 9525,
        )
        dot.name = f"emotion-dot-{i+1}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.color.rgb = WHITE
        dot.line.width = 12700
        # Face above
        add_text(
            slide, f"emotion-face-{i+1}", face,
            x_px=cx - 16, y_px=py - 32, w_px=32, h_px=22,
            font_size_px=16, align="center",
        )
        # Score below
        add_text(
            slide, f"emotion-score-{i+1}", str(scores[i]),
            x_px=cx - 16, y_px=py + 8, w_px=32, h_px=14,
            font_size_px=9, color=TEXT_FAINT, align="center",
        )

    # Pain callout (above stage 3)
    pain_text_y = emo_y + emo_h + 6
    pain_w = 360
    pain_x = centers_x[2] - pain_w // 2
    pain = add_rect(slide, "pain-callout-bg", pain_x, pain_text_y, pain_w, 22, PAIN_BG)
    pain.line.color.rgb = PAIN_BORDER
    pain.line.width = 9525
    add_text(
        slide, "pain-callout-text",
        "⚠ 42% of users drop off at first complex configuration step",
        x_px=pain_x, y_px=pain_text_y, w_px=pain_w, h_px=22,
        font_size_px=10, color=PAIN_RED, bold=True, align="center", anchor="middle",
    )

    # ROW 4: Opportunities (2 per stage)
    opp_y = pain_text_y + 30
    opp_h = body_y + body_h - opp_y
    opportunities = [
        ["Simplify form fields & reduce friction at registration",
         "A/B test confirm email subject lines"],
        ["Personalise wizard steps by use case",
         "Offer self-serve config as alternative to call"],
        ["Break complex config into guided micro-steps",
         "Trigger proactive chat at drop-off point"],
        ["Tailor digest to usage patterns per segment",
         "Gamify milestones to reinforce habit loops"],
        ["Automate referral rewards at NPS ≥9",
         "Co-create case studies with champions"],
    ]
    for i, items in enumerate(opportunities):
        sx = body_x + i * (col_w + gap)
        is_critical = (i == 2)
        bg = PAIN_BG if is_critical else OPP_BG
        border = PAIN_BORDER if is_critical else OPP_BORDER
        check_color = PAIN_RED if is_critical else OPP_GREEN
        cell = add_rect(slide, f"opp-{i+1}-bg", sx, opp_y, col_w, opp_h, bg)
        cell.line.color.rgb = border
        cell.line.width = 9525
        for ji, item in enumerate(items):
            iy = opp_y + 8 + ji * 32
            add_text(
                slide, f"opp-{i+1}-{ji+1}-check", "✓",
                x_px=sx + 8, y_px=iy, w_px=14, h_px=14,
                font_size_px=11, color=check_color, bold=True,
            )
            add_text(
                slide, f"opp-{i+1}-{ji+1}-text", item,
                x_px=sx + 24, y_px=iy, w_px=col_w - 32, h_px=26,
                font_size_px=10, color=TEXT_MID, italic=True,
            )

    add_footer(slide, page_num=280)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "280_customer-onboarding-journey.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
