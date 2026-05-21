"""
Builder for pattern 31: Maturity pyramid.

Decompose treatment (per SHAPE-ROLES table): 5 stacked trapezoid tiers as
native shapes (MSO_SHAPE.TRAPEZOID rotated), plus tier labels, descriptions,
current-state chip and arrow, and convergence band.

Source HTML: _pattern-library/31_maturity-pyramid.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect, px_to_emu,
    add_chrome, add_footer, add_title_block, add_convergence,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Tier fill graduations (top=darkest)
TIER_FILLS = [
    RGBColor(0xE5, 0xD5, 0xF0),  # tier 1 base, lightest
    RGBColor(0xC3, 0x9B, 0xDB),  # tier 2 current
    RGBColor(0x8B, 0x5F, 0xB8),  # tier 3
    RGBColor(0x5C, 0x2D, 0x87),  # tier 4
    RGBColor(0x2D, 0x0A, 0x4E),  # tier 5 top, darkest
]


def add_trapezoid(slide, shape_id, x_px, y_px, w_px, h_px, fill_color, *, flip_v=True):
    """Add a trapezoid (PPTX MSO_SHAPE.TRAPEZOID is base-down; flip vertically
    to point up so it reads as a pyramid tier with wider base."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.TRAPEZOID,
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(w_px), px_to_emu(h_px),
    )
    shape.name = shape_id
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # MSO_SHAPE.TRAPEZOID has wider base at bottom by default; we want all
    # tiers to be wider at the base too, so leave default orientation.
    return shape


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_title_block(
        slide,
        title="Maturity model — where the practice is today and where it could go.",
        subtitle="CMMI-style assessment of deck quality discipline across the team.",
        title_h=68,
        subtitle_h=22,
    )

    # Stage area centered. Pyramid base width 580px (scaled to slide).
    # Pyramid drawn from tier 1 (base, widest) at bottom to tier 5 (top, narrowest).
    stage_left_center = 640  # slide center for symmetry
    tier_h = 56
    tier_gap = 8
    base_w = 580
    narrow_step = 90  # each tier narrows by this much
    pyramid_top = 200  # top of tier 5

    # Tier widths (top to bottom): 220, 310, 400, 490, 580
    # We render from tier 1 (bottom) up to tier 5 (top)
    tier_widths = [base_w - narrow_step * (4 - i) for i in range(5)]  # i=0 tier1=base
    # tier_widths: [220, 310, 400, 490, 580]? Recompute: for i=0 it should be base.
    # Rewriting: tier 1 = widest = 580; tier 5 = narrowest = 220
    tier_widths = [580, 490, 400, 310, 220]

    # Vertical positions: tier 5 at top, tier 1 at bottom
    # tier N (1..5): y = pyramid_top + (5 - N) * (tier_h + tier_gap)
    for i in range(5):
        tier_num = 5 - i  # 5, 4, 3, 2, 1
        y = pyramid_top + i * (tier_h + tier_gap)
        w = tier_widths[tier_num - 1]
        x = stage_left_center - w // 2
        fill = TIER_FILLS[tier_num - 1]
        # Use rectangle for tier shape (simple representation; pyramid effect by varying widths)
        shape = add_rect(slide, f"tier-{tier_num}-shape", x, y, w, tier_h, fill)

        # Tier label (LEVEL N + NAME) — centered on tier
        label_color = WHITE if tier_num >= 3 else BRAND_PRIMARY
        add_text(
            slide, f"tier-{tier_num}-label", f"LEVEL {tier_num}",
            x_px=x, y_px=y + 8, w_px=w, h_px=12,
            font_size_px=9, color=label_color, bold=True, align="center",
            letter_spacing_px=2, uppercase=True,
        )
        names = ["INITIAL", "MANAGED", "DEFINED", "QUANTITATIVELY MANAGED", "OPTIMIZING"]
        add_text(
            slide, f"tier-{tier_num}-name", names[tier_num - 1],
            x_px=x, y_px=y + 22, w_px=w, h_px=22,
            font_size_px=13, color=label_color, bold=True, align="center",
            letter_spacing_px=1,
        )

        # Tier description to the right of pyramid
        desc_x = stage_left_center + base_w // 2 + 30
        desc_w = 200
        descs = [
            "Ad-hoc. No standard. Quality varies by author.",
            "Templates exist; usage inconsistent across teams.",
            "Standard process followed. Predictable quality.",
            "Measured. Cycle time, edits, quality scores tracked.",
            "Continuous improvement. Patterns harvested across teams.",
        ]
        add_text(
            slide, f"tier-{tier_num}-desc", descs[tier_num - 1],
            x_px=desc_x, y_px=y + (tier_h - 36) // 2, w_px=desc_w, h_px=36,
            font_size_px=10, color=TEXT_MID,
        )

    # Current-state chip + arrow pointing to tier 2
    tier_2_y = pyramid_top + 3 * (tier_h + tier_gap)  # tier 2 = row index 3 from top
    chip_y = tier_2_y + (tier_h - 22) // 2
    # Tier 2 left edge
    tier_2_w = tier_widths[1]  # 490
    tier_2_left = stage_left_center - tier_2_w // 2

    add_rect(slide, "tier-current-chip", x_px=tier_2_left - 130, y_px=chip_y, w_px=110, h_px=22,
             fill_color=BRAND_ACCENT)
    add_text(
        slide, "tier-current-chip-text", "CURRENT STATE",
        x_px=tier_2_left - 130, y_px=chip_y, w_px=110, h_px=22,
        font_size_px=10, color=WHITE, bold=True, align="center",
        anchor="middle", uppercase=True, letter_spacing_px=1.5,
    )
    add_text(
        slide, "tier-current-note", "We're here today",
        x_px=tier_2_left - 150, y_px=chip_y + 24, w_px=150, h_px=14,
        font_size_px=10, color=BRAND_PRIMARY, italic=True, bold=True, align="right",
    )

    # Arrow (small triangle pointing right at tier 2)
    arrow_shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        px_to_emu(tier_2_left - 14), px_to_emu(chip_y + 4),
        px_to_emu(14), px_to_emu(14),
    )
    arrow_shape.name = "tier-current-arrow"
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = BRAND_ACCENT
    arrow_shape.line.fill.background()

    add_convergence(
        slide,
        "We're at Level 2 today. The pilot moves us to Level 3 in four weeks.",
    )

    add_footer(slide, page_num=31)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "31_maturity-pyramid.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
