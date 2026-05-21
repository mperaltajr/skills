"""
Builder for pattern 318: 3-bucket icon tiles (circle / diamond / triangle).

Three tile cards with brand-accent geometric icons, headline, rule, and body text.

Source HTML: _pattern-library/318_3bucket-icon-tiles.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect, add_icon,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.enum.shapes import MSO_SHAPE
from twins.helpers import px_to_emu


def _icon_circle(slide, name, x, y, size):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_ACCENT
    shape.line.fill.background()


def _icon_diamond(slide, name, x, y, size):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_ACCENT
    shape.line.fill.background()


def _icon_triangle(slide, name, x, y, size):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_ACCENT
    shape.line.fill.background()


def _tile(slide, n, x, y, w, h, *, icon_fn, headline, body):
    tile = add_rect(slide, f"tile-{n}-bg", x, y, w, h, CARD_BG)
    tile.line.color.rgb = CARD_BORDER
    tile.line.width = 9525
    # Icon centered
    icon_size = 48
    icon_x = x + (w - icon_size) // 2
    icon_y = y + 36
    icon_fn(slide, f"tile-{n}-icon", icon_x, icon_y, icon_size)
    # Headline
    add_text(
        slide, f"tile-{n}-headline", headline,
        x_px=x + 16, y_px=icon_y + icon_size + 18, w_px=w - 32, h_px=24,
        font_size_px=16, color=BRAND_PRIMARY, bold=True, align="center",
    )
    # Rule
    rule_w = 40
    add_rect(slide, f"tile-{n}-rule",
             x + (w - rule_w) // 2, icon_y + icon_size + 46, rule_w, 2, BRAND_ACCENT)
    # Body
    add_text(
        slide, f"tile-{n}-body", body,
        x_px=x + 22, y_px=icon_y + icon_size + 62, w_px=w - 44, h_px=180,
        font_size_px=12, color=TEXT_MID, align="center",
    )
    # Link
    add_text(
        slide, f"tile-{n}-link", "Learn more →",
        x_px=x + 16, y_px=y + h - 30, w_px=w - 32, h_px=18,
        font_size_px=10, color=BRAND_ACCENT, bold=True, align="center",
    )


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_title_block(
        slide,
        title="Three ways we create <strong>lasting impact</strong>",
        subtitle="A structured lens on where value is generated across the engagement",
        title_x=48, title_y=44, title_w=1184, title_h=44,
        subtitle_h=20, brand_rule_w=64,
    )

    body_top = 200
    body_left = 48
    body_w = 1184
    body_h = 720 - 32 - 14 - body_top
    gap = 18
    tile_w = (body_w - 2 * gap) // 3

    tiles = [
        (_icon_circle, "Operate at Scale",
         "Harmonize delivery across regions by standardizing core processes, tooling, and governance to reduce friction and unlock throughput at every tier."),
        (_icon_diamond, "Build for Resilience",
         "Architect platforms and operating models that absorb disruption, shorten recovery cycles, and maintain continuity even under sustained competitive pressure."),
        (_icon_triangle, "Accelerate Growth",
         "Redirect reinvested capacity toward high-value opportunities, compressing time-to-market and enabling the organization to capture demand faster than rivals."),
    ]
    for i, (icon_fn, head, body) in enumerate(tiles):
        n = i + 1
        cx = body_left + i * (tile_w + gap)
        _tile(slide, n, cx, body_top, tile_w, body_h,
              icon_fn=icon_fn, headline=head, body=body)

    add_footer(slide, page_num=318)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "318_3bucket-icon-tiles.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
