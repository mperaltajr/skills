"""
Builder for pattern 14: Gantt workplan timeline (12 weeks, 3 phases, 8 tasks).

Source HTML: _pattern-library/14_gantt-workplan-timeline.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect, px_to_emu,
    add_chrome, add_footer, add_title_block, add_convergence,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.enum.shapes import MSO_SHAPE


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_title_block(
        slide,
        title="Twelve-week pilot — what gets done in each phase.",
        subtitle="Setup is plumbing. Pilot is where the value shows. Scale is the decision we earn the right to make.",
        title_h=68,
        subtitle_h=22,
        brand_rule_w=56,
    )

    # Gantt grid geometry
    gantt_left = 56
    gantt_top = 230
    label_w = 200
    week_total_w = 1280 - 112 - label_w  # = 968
    week_w = week_total_w // 12  # = 80
    phase_h = 22
    week_head_h = 26
    task_h = 28
    n_tasks = 8

    # Phase bands (row 1)
    add_rect(slide, "gantt-phase-spacer", gantt_left, gantt_top, label_w, phase_h, WHITE)
    p1_x = gantt_left + label_w
    add_rect(slide, "phase-band-1-bg", p1_x, gantt_top, week_w * 4, phase_h, BRAND_PRIMARY_MID)
    add_text(
        slide, "phase-band-1-label", "PHASE 1 · SETUP",
        x_px=p1_x, y_px=gantt_top, w_px=week_w * 4, h_px=phase_h,
        font_size_px=10, color=WHITE, bold=True, uppercase=True,
        align="center", anchor="middle",
    )
    p2_x = p1_x + week_w * 4
    add_rect(slide, "phase-band-2-bg", p2_x, gantt_top, week_w * 4, phase_h, BRAND_ACCENT)
    add_text(
        slide, "phase-band-2-label", "PHASE 2 · PILOT",
        x_px=p2_x, y_px=gantt_top, w_px=week_w * 4, h_px=phase_h,
        font_size_px=10, color=WHITE, bold=True, uppercase=True,
        align="center", anchor="middle",
    )
    p3_x = p2_x + week_w * 4
    add_rect(slide, "phase-band-3-bg", p3_x, gantt_top, week_w * 4, phase_h, BRAND_PRIMARY_MID)
    add_text(
        slide, "phase-band-3-label", "PHASE 3 · SCALE",
        x_px=p3_x, y_px=gantt_top, w_px=week_w * 4, h_px=phase_h,
        font_size_px=10, color=WHITE, bold=True, uppercase=True,
        align="center", anchor="middle",
    )

    # Week headers (row 2)
    wh_y = gantt_top + phase_h
    add_rect(slide, "timeline-row-header-bg", gantt_left, wh_y, label_w, week_head_h, CARD_BG)
    add_text(
        slide, "timeline-row-header", "WORKSTREAM",
        x_px=gantt_left + 12, y_px=wh_y, w_px=label_w - 20, h_px=week_head_h,
        font_size_px=10, color=TEXT_FAINT, bold=True, uppercase=True, anchor="middle",
    )
    for i in range(12):
        wx = gantt_left + label_w + i * week_w
        add_rect(slide, f"timeline-tick-{i+1}-bg", wx, wh_y, week_w, week_head_h, CARD_BG)
        add_text(
            slide, f"timeline-tick-{i+1}-label", f"W{i+1}",
            x_px=wx, y_px=wh_y, w_px=week_w, h_px=week_head_h,
            font_size_px=11, color=BRAND_PRIMARY_MID, bold=True,
            align="center", anchor="middle",
        )

    # Task rows
    tasks = [
        ("Stakeholder alignment",     "completed",   1, 2, "W1–W2", None),
        ("Tool access provisioning",  "completed",   1, 2, "W1–W2", None),
        ("Onboarding sessions",       "completed",   3, 2, "W3–W4", None),
        ("First pilot decks",         "in-progress", 5, 3, "W5–W7 · in flight", None),
        ("Mid-pilot review",          "milestone",   6, 1, "",      1),  # diamond at W6
        ("Refinement",                "in-progress", 7, 3, "W7–W9", None),
        ("Scaling readiness check",   "milestone",   10, 1, "",     2),  # diamond at W10
        ("Q3 rollout decision",       "milestone",   12, 1, "",     3),  # diamond at W12
    ]
    tasks_top = wh_y + week_head_h
    for i, (label, status, start_week, span, bar_text, ms_n) in enumerate(tasks):
        n = i + 1
        ty = tasks_top + i * task_h
        # Label cell
        lbl_bg = add_rect(slide, f"task-{n}-label-bg", gantt_left, ty, label_w, task_h, WHITE)
        lbl_bg.line.color.rgb = CARD_BORDER
        lbl_bg.line.width = 9525
        add_text(
            slide, f"task-{n}-label", label,
            x_px=gantt_left + 12, y_px=ty, w_px=label_w - 20, h_px=task_h,
            font_size_px=11, color=TEXT_DARK, bold=False, anchor="middle",
        )
        if status == "milestone":
            # Diamond at start_week (centered in that week column)
            mx = gantt_left + label_w + (start_week - 1) * week_w + week_w // 2 - 7
            my = ty + task_h // 2 - 7
            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                px_to_emu(mx), px_to_emu(my),
                px_to_emu(14), px_to_emu(14),
            )
            diamond.name = f"milestone-{ms_n}"
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = BRAND_ACCENT
            diamond.line.color.rgb = WHITE
            diamond.line.width = 19050
        else:
            # Draw bar
            bar_x = gantt_left + label_w + (start_week - 1) * week_w + 4
            bar_w = span * week_w - 8
            bar_color = TEXT_FAINT if status == "completed" else BRAND_ACCENT
            add_rect(slide, f"task-{n}-bar", bar_x, ty + 6, bar_w, task_h - 12, bar_color)
            add_text(
                slide, f"task-{n}-bar-text", bar_text,
                x_px=bar_x, y_px=ty + 6, w_px=bar_w, h_px=task_h - 12,
                font_size_px=9, color=WHITE, bold=True,
                anchor="middle", padding_px=(0, 8, 0, 8),
            )

    # Vertical gridlines at each week-column boundary
    for i in range(1, 12):  # 11 internal boundaries
        vl_x = gantt_left + label_w + i * week_w
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            px_to_emu(vl_x), px_to_emu(tasks_top),
            px_to_emu(1), px_to_emu(n_tasks * task_h),
        )
        line.name = f"gantt-vgrid-{i}"
        from pptx.dml.color import RGBColor as _RGB
        line.fill.solid()
        line.fill.fore_color.rgb = _RGB(0xF1, 0xE8, 0xF8)
        line.line.fill.background()

    # Legend — top-right, below subheadline (y≈195), right-aligned to x≈1240
    legend_y = 195
    # Only rect swatches; diamond entry removed (milestone markers exist on body)
    legends = [
        ("Completed", TEXT_FAINT, 70),
        ("In progress", BRAND_ACCENT, 80),
        ("Milestone", BRAND_ACCENT, 60),
    ]
    item_gap = 18
    sw_w = 16
    # Calculate total width to right-align to x=1240
    total_w = sum(sw_w + 8 + lw + item_gap for _, _, lw in legends) - item_gap
    cur_x = 1240 - total_w
    for i, (label, color, label_w) in enumerate(legends):
        n = i + 1
        add_rect(slide, f"legend-{n}-swatch", cur_x, legend_y + 3, sw_w, 10, color)
        add_text(
            slide, f"legend-{n}-label", label,
            x_px=cur_x + sw_w + 6, y_px=legend_y, w_px=label_w, h_px=18,
            font_size_px=10, color=TEXT_MID, anchor="middle",
        )
        cur_x += sw_w + 8 + label_w + item_gap

    add_convergence(
        slide,
        "Milestones are decision points, not check-ins — each one earns the right to continue.",
    )

    add_footer(slide, page_num=14)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "14_gantt-workplan-timeline.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
