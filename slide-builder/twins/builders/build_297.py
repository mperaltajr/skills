"""
Builder for pattern 297: RACI Matrix (8 activities × 6 roles).

Source HTML: _pattern-library/297_raci-matrix.html

CRITICAL: Legend MUST sit BELOW the subheadline / brand-rule (top-y ≥ 230,
right-aligned to x ≈ 1240). Body table is shifted down to clear the legend
zone. Recipe from build_44.py + build_63.py.
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor

PILL_C_BG = RGBColor(0xF5, 0x9E, 0x0B)   # amber for Consulted
PILL_C_FG = RGBColor(0x1A, 0x1A, 0x00)


def build():
    prs, slide = new_slide()
    add_chrome(slide)
    add_title_block(
        slide,
        title="Project <strong>RACI Matrix</strong> — Roles & Accountability",
        subtitle="Who is Responsible, Accountable, Consulted, and Informed across key activities",
    )

    # === LEGEND — BELOW subheadline / brand-rule (top-y = 234, right-aligned) ===
    leg_h = 36
    leg_y = 234
    leg_w = 720
    leg_x = 1240 - leg_w   # right edge at 1240
    leg = add_rect(slide, "legend-bg", leg_x, leg_y, leg_w, leg_h, CARD_BG)
    leg.line.color.rgb = CARD_BORDER
    leg.line.width = 9525

    add_text(
        slide, "legend-label", "LEGEND",
        x_px=leg_x + 12, y_px=leg_y, w_px=60, h_px=leg_h,
        font_size_px=9, color=TEXT_FAINT, bold=True, anchor="middle",
        uppercase=True, letter_spacing_px=1.2,
    )

    raci_items = [
        ("R", "Responsible", "does the work", BRAND_ACCENT, WHITE),
        ("A", "Accountable", "owns the outcome", BRAND_PRIMARY, WHITE),
        ("C", "Consulted", "provides input", PILL_C_BG, PILL_C_FG),
        ("I", "Informed", "kept up to date", CARD_BG, TEXT_MID),
    ]
    item_x = leg_x + 78
    item_step = 158
    for i, (letter, label, desc, bg, fg) in enumerate(raci_items):
        n = i + 1
        sx = item_x + i * item_step
        # circular pill (use oval shape)
        from pptx.enum.shapes import MSO_SHAPE
        pill = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            sx * 9525, (leg_y + 8) * 9525, 20 * 9525, 20 * 9525,
        )
        pill.name = f"legend-{n}-swatch"
        pill.fill.solid()
        pill.fill.fore_color.rgb = bg
        if letter == "I":
            pill.line.color.rgb = CARD_BORDER
            pill.line.width = 12700
        else:
            pill.line.fill.background()
        # Letter text inside pill
        add_text(
            slide, f"legend-{n}-letter", letter,
            x_px=sx, y_px=leg_y + 8, w_px=20, h_px=20,
            font_size_px=10, color=fg, bold=True, align="center", anchor="middle",
        )
        # Label + description
        add_text(
            slide, f"legend-{n}-label", label,
            x_px=sx + 26, y_px=leg_y + 4, w_px=120, h_px=14,
            font_size_px=10, color=TEXT_DARK, bold=True,
        )
        add_text(
            slide, f"legend-{n}-desc", desc,
            x_px=sx + 26, y_px=leg_y + 18, w_px=120, h_px=14,
            font_size_px=9, color=TEXT_MID,
        )

    # === TABLE — pushed down to clear legend (legend bottom = 270) ===
    t_left = 48
    t_right = 1280 - 48
    t_top = 286
    t_bottom = 670  # stay clear of invariant zone (≥672)
    t_w = t_right - t_left
    t_h = t_bottom - t_top

    # Columns: activity (240px) + 6 role cols
    act_w = 240
    role_w = (t_w - act_w) // 6

    # Header row
    head_h = 44
    add_rect(slide, "table-head-bg", t_left, t_top, t_w, head_h, BRAND_PRIMARY)
    add_text(
        slide, "table-col-1-header", "ACTIVITY / DECISION",
        x_px=t_left + 14, y_px=t_top, w_px=act_w - 14, h_px=head_h,
        font_size_px=10, color=BRAND_ACCENT_SOFT, bold=True, anchor="middle",
        uppercase=True, letter_spacing_px=1,
    )
    roles = [
        ("Project", "Sponsor"),
        ("Program", "Manager"),
        ("Workstream", "Lead"),
        ("Finance", ""),
        ("IT", ""),
        ("Comms", ""),
    ]
    for i, (line1, line2) in enumerate(roles):
        n = i + 2
        cx = t_left + act_w + i * role_w
        add_text(
            slide, f"table-col-{n}-header", line1,
            x_px=cx, y_px=t_top + 6, w_px=role_w, h_px=18,
            font_size_px=10, color=WHITE, bold=True, align="center",
            uppercase=True, letter_spacing_px=0.6,
        )
        if line2:
            add_text(
                slide, f"table-col-{n}-header-line2", line2,
                x_px=cx, y_px=t_top + 24, w_px=role_w, h_px=14,
                font_size_px=10, color=WHITE, bold=True, align="center",
                uppercase=True, letter_spacing_px=0.6,
            )

    # Body rows
    activities = [
        ("Define scope",          ["A", "R", "C", "I", "I", "-"]),
        ("Approve budget",        ["A", "C", "-", "R", "-", "-"]),
        ("Execute delivery",      ["I", "A", "R", "I", "C", "-"]),
        ("Manage vendors",        ["-", "A", "R", "C", "C", "-"]),
        ("Report progress",       ["I", "R", "C", "I", "I", "I"]),
        ("Resolve issues",        ["C", "A", "R", "C", "R", "-"]),
        ("Sign-off",              ["A", "R", "C", "C", "-", "-"]),
        ("Communicate externally", ["A", "C", "I", "-", "-", "R"]),
    ]

    pill_color_map = {
        "R": (BRAND_ACCENT, WHITE),
        "A": (BRAND_PRIMARY, WHITE),
        "C": (PILL_C_BG, PILL_C_FG),
        "I": (CARD_BG, TEXT_MID),
    }

    row_h = (t_h - head_h) // len(activities)
    from pptx.enum.shapes import MSO_SHAPE

    for ri, (act_name, cells) in enumerate(activities):
        n = ri + 1
        ry = t_top + head_h + ri * row_h
        # Row background (alternating)
        bg = CARD_BG if ri % 2 == 1 else WHITE
        add_rect(slide, f"row-{n}-bg", t_left, ry, t_w, row_h, bg)
        # Activity label cell
        add_text(
            slide, f"table-row-{n}-label", act_name,
            x_px=t_left + 14, y_px=ry, w_px=act_w - 20, h_px=row_h,
            font_size_px=12, color=TEXT_DARK, anchor="middle",
        )
        # Right border of activity column
        add_rect(slide, f"row-{n}-act-border", t_left + act_w - 1, ry, 1, row_h, CARD_BORDER)

        # Role cells
        for ci, cell_val in enumerate(cells):
            cn = ci + 2
            cx = t_left + act_w + ci * role_w
            cell_id = f"table-row-{n}-cell-{cn}"
            if cell_val == "-":
                # tiny dot
                add_rect(
                    slide, cell_id,
                    cx + role_w // 2 - 3, ry + row_h // 2 - 3, 6, 6, CARD_BORDER,
                )
            else:
                bg_color, fg_color = pill_color_map[cell_val]
                # circular pill
                pill_size = 26
                px = cx + (role_w - pill_size) // 2
                py = ry + (row_h - pill_size) // 2
                pill = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    px * 9525, py * 9525, pill_size * 9525, pill_size * 9525,
                )
                pill.name = cell_id + "-pill"
                pill.fill.solid()
                pill.fill.fore_color.rgb = bg_color
                if cell_val == "I":
                    pill.line.color.rgb = CARD_BORDER
                    pill.line.width = 12700
                else:
                    pill.line.fill.background()
                add_text(
                    slide, cell_id, cell_val,
                    x_px=px, y_px=py, w_px=pill_size, h_px=pill_size,
                    font_size_px=11, color=fg_color, bold=True,
                    align="center", anchor="middle",
                )

    add_footer(slide, page_num=297)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "297_raci-matrix.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
