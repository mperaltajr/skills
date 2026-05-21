"""
Builder for pattern 268: Market share by segment — 4 donuts + summary table.

Source HTML: _pattern-library/268_market-share-segment.html

Layout: title + 4 donut cards (top) + 5-column summary table (bottom).

LEGEND PLACEMENT: Right-aligned below subheadline (top-y ≥ 230, right edge ≈ 1240).
Legend: Ours (brand-accent) / Competitor (card-border).
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect, px_to_emu,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

UP_COLOR = RGBColor(0x16, 0xA3, 0x4A)
DOWN_COLOR = RGBColor(0xDC, 0x26, 0x26)
FLAT_COLOR = RGBColor(0x64, 0x74, 0x8B)


def build():
    prs, slide = new_slide()
    add_chrome(slide)
    add_title_block(
        slide,
        title="Market share by segment — <strong>Enterprise and Government lead, SMB is the white space.</strong>",
        subtitle="Share vs. competitor across four addressable segments · FY2026 actuals · bubble size = total market",
    )

    # === LEGEND — below subheadline, right-aligned ===
    leg_y = 230
    leg_w = 240
    leg_x = 1240 - leg_w
    add_rect(slide, "legend-ours-swatch", leg_x, leg_y + 7, 10, 10, BRAND_ACCENT)
    add_text(slide, "legend-ours-label", "Ours",
             x_px=leg_x + 16, y_px=leg_y + 2, w_px=60, h_px=18,
             font_size_px=11, color=TEXT_MID, bold=True)
    add_rect(slide, "legend-comp-swatch", leg_x + 90, leg_y + 7, 10, 10, CARD_BORDER)
    add_text(slide, "legend-comp-label", "Competitor",
             x_px=leg_x + 106, y_px=leg_y + 2, w_px=120, h_px=18,
             font_size_px=11, color=TEXT_MID, bold=True)

    # Body
    body_top = 268
    body_bottom = 632
    body_h = body_bottom - body_top
    left_x = 48
    right_x = 1280 - 48
    body_w = right_x - left_x

    # Top — 4 donut cards
    donut_section_h = 180
    card_gap = 16
    card_w = (body_w - 3 * card_gap) // 4

    segments = [
        ("Enterprise", 34, "$4.2B total market", "up"),
        ("Mid-Market", 28, "$2.8B total market", "up"),
        ("SMB", 12, "$8.4B total market", "flat"),
        ("Government", 41, "$1.6B total market", "down"),
    ]
    for i, (name, pct, size, trend) in enumerate(segments):
        n = i + 1
        cx = left_x + i * (card_w + card_gap)
        cy = body_top
        # card bg
        card = add_rect(slide, f"donut-card-{n}", cx, cy, card_w, donut_section_h, WHITE)
        card.line.color.rgb = CARD_BORDER
        card.line.width = 9525
        # donut
        donut_size = 86
        dx = cx + (card_w - donut_size) // 2
        dy = cy + 10
        # outer track
        track = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            px_to_emu(dx), px_to_emu(dy), px_to_emu(donut_size), px_to_emu(donut_size))
        track.name = f"donut-track-{n}"
        track.fill.solid()
        track.fill.fore_color.rgb = CARD_BORDER
        track.line.fill.background()
        # arc approximation: full pie filled with brand-accent at pct (use full overlay since we cannot rotate)
        arc = slide.shapes.add_shape(MSO_SHAPE.PIE,
            px_to_emu(dx), px_to_emu(dy), px_to_emu(donut_size), px_to_emu(donut_size))
        arc.name = f"donut-our-{n}"
        arc.fill.solid()
        arc.fill.fore_color.rgb = BRAND_ACCENT
        arc.line.fill.background()
        # hole
        hole_size = 56
        hole = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            px_to_emu(dx + (donut_size - hole_size) // 2),
            px_to_emu(dy + (donut_size - hole_size) // 2),
            px_to_emu(hole_size), px_to_emu(hole_size))
        hole.name = f"donut-hole-{n}"
        hole.fill.solid()
        hole.fill.fore_color.rgb = WHITE
        hole.line.fill.background()
        # center pct
        add_text(slide, f"center-pct-{n}", f"{pct}%",
                 x_px=dx, y_px=dy + donut_size // 2 - 14, w_px=donut_size, h_px=24,
                 font_size_px=18, color=BRAND_PRIMARY, bold=True, align="center")
        # OUR SHARE label
        add_text(slide, f"donut-segment-label-{n}", "OUR SHARE",
                 x_px=cx, y_px=cy + 108, w_px=card_w, h_px=14,
                 font_size_px=8, color=TEXT_FAINT, bold=True,
                 letter_spacing_px=1.4, uppercase=True, align="center")
        # segment name + trend
        arrow_map = {"up": ("↑", UP_COLOR), "down": ("↓", DOWN_COLOR), "flat": ("→", FLAT_COLOR)}
        arrow, ac = arrow_map[trend]
        add_text(slide, f"seg-name-{n}", f"{name} {arrow}",
                 x_px=cx, y_px=cy + 126, w_px=card_w, h_px=20,
                 font_size_px=14, color=BRAND_PRIMARY, bold=True, align="center",
                 emphasis_color=ac)
        add_text(slide, f"seg-size-{n}", size,
                 x_px=cx, y_px=cy + 150, w_px=card_w, h_px=14,
                 font_size_px=10, color=TEXT_MID, align="center")

    # BOTTOM — summary table
    tbl_top = body_top + donut_section_h + 14
    tbl_h = body_bottom - tbl_top
    # 5 columns
    col_widths_pct = [22, 18, 16, 22, 22]
    col_widths = [int(body_w * p / 100) for p in col_widths_pct]
    col_x = [left_x]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)
    headers = ["Segment", "Market Size", "Our Share", "YoY Change", "Strategic Priority"]
    hdr_h = 30
    add_rect(slide, "tbl-header-bg", left_x, tbl_top, body_w, hdr_h, BRAND_PRIMARY)
    for ci, h in enumerate(headers):
        a = "left" if ci == 0 else "center"
        add_text(slide, f"th-{ci+1}", h.upper(),
                 x_px=col_x[ci] + 14, y_px=tbl_top, w_px=col_widths[ci] - 24, h_px=hdr_h,
                 font_size_px=10, color=WHITE, bold=True, anchor="middle", align=a,
                 letter_spacing_px=1.2, uppercase=True)

    rows = [
        ("Enterprise", "$4.2B", "34%", ("↑ Growing", UP_COLOR), ("Grow", "grow")),
        ("Mid-Market", "$2.8B", "28%", ("↑ Growing", UP_COLOR), ("Invest", "invest")),
        ("SMB", "$8.4B", "12%", ("→ Flat", FLAT_COLOR), ("Hold", "hold")),
        ("Government", "$1.6B", "41%", ("↓ Declining", DOWN_COLOR), ("Defend", "defend")),
    ]
    pill_styles = {
        "grow": (RGBColor(0xF3, 0xE8, 0xFF), BRAND_ACCENT),
        "invest": (RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x16, 0x65, 0x34)),
        "hold": (RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E)),
        "defend": (RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1E, 0x40, 0xAF)),
    }
    row_h = (tbl_h - hdr_h) // len(rows)
    for i, (seg, mkt, share, (yoy_t, yoy_c), (pill_t, pill_k)) in enumerate(rows):
        n = i + 1
        ry = tbl_top + hdr_h + i * row_h
        if i % 2 == 1:
            add_rect(slide, f"tr-{n}-bg", left_x, ry, body_w, row_h, CARD_BG)
        add_rect(slide, f"tr-{n}-rule", left_x, ry, body_w, 1, CARD_BORDER)
        # segment
        add_text(slide, f"td-seg-{n}", seg,
                 x_px=col_x[0] + 14, y_px=ry, w_px=col_widths[0] - 24, h_px=row_h,
                 font_size_px=13, color=BRAND_PRIMARY, bold=True, anchor="middle")
        # market size
        add_text(slide, f"td-market-{n}", mkt,
                 x_px=col_x[1], y_px=ry, w_px=col_widths[1], h_px=row_h,
                 font_size_px=13, color=TEXT_DARK, bold=True, align="center", anchor="middle")
        # share
        add_text(slide, f"td-share-{n}", share,
                 x_px=col_x[2], y_px=ry, w_px=col_widths[2], h_px=row_h,
                 font_size_px=16, color=BRAND_ACCENT, bold=True, align="center", anchor="middle")
        # yoy
        add_text(slide, f"td-yoy-{n}", yoy_t,
                 x_px=col_x[3], y_px=ry, w_px=col_widths[3], h_px=row_h,
                 font_size_px=13, color=yoy_c, bold=True, align="center", anchor="middle")
        # priority pill
        bg, fg = pill_styles[pill_k]
        pill_w = 90
        pill_h = 22
        px_p = col_x[4] + (col_widths[4] - pill_w) // 2
        py_p = ry + (row_h - pill_h) // 2
        add_rect(slide, f"td-pill-{n}", px_p, py_p, pill_w, pill_h, bg)
        add_text(slide, f"td-pill-{n}-text", pill_t.upper(),
                 x_px=px_p, y_px=py_p, w_px=pill_w, h_px=pill_h,
                 font_size_px=9, color=fg, bold=True, align="center", anchor="middle",
                 letter_spacing_px=0.8)

    add_footer(slide, page_num=268)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "268_market-share-segment.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
