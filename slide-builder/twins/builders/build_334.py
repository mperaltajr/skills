"""
Builder for pattern 334: 7-bucket timeline arc.

HTML uses an SVG arc with 7 nodes following the curve. python-pptx cannot draw
arbitrary SVG arcs, so we approximate with: a horizontal connector line + 7
circular nodes arranged with subtle vertical offset that simulates the arc
midpoint dip. Each node has number, title, and a 2-line caption above/below.

Source HTML: _pattern-library/334_7bucket-timeline-arc.html
"""
from pathlib import Path
import sys
import math

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
    px_to_emu,
)
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR


def _circle(slide, name, x, y, size, fill, *, text=None, text_color=WHITE, font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    shape.name = f"{name}-bg"
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text is not None:
        add_text(slide, f"{name}-text", text,
                 x_px=x, y_px=y, w_px=size, h_px=size,
                 font_size_px=font_size, color=text_color, bold=True,
                 align="center", anchor="middle")


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_title_block(
        slide,
        title="Seven milestones that <strong>define the transformation arc</strong>",
        subtitle="Each bucket represents a distinct phase gate — progress is sequential and momentum compounds across all seven stages.",
        title_x=48, title_y=44, title_w=1184, title_h=44,
        subtitle_h=20, brand_rule_w=64,
    )

    body_top = 200
    body_left = 48
    body_w = 1184
    body_h = 720 - 32 - 14 - body_top

    # Stage labels on left/right
    add_text(slide, "stage-past", "PAST",
             x_px=body_left, y_px=body_top + 4, w_px=80, h_px=18,
             font_size_px=10, color=TEXT_FAINT, bold=True, letter_spacing_px=1.6)
    add_text(slide, "stage-future", "FUTURE",
             x_px=body_left + body_w - 80, y_px=body_top + 4, w_px=80, h_px=18,
             font_size_px=10, color=BRAND_ACCENT, bold=True, letter_spacing_px=1.6, align="right")

    # Arc approximation: 7 nodes with vertical offset matching a downward-opening arc
    nodes = [
        ("Discovery", "Baseline & stakeholder", "assessment complete"),
        ("Design", "Future-state blueprint", "signed off"),
        ("Pilot", "MVP tested in one", "business unit"),
        ("Build", "Core platform & data", "infrastructure live"),
        ("Scale", "Rollout to all regions", "and functions"),
        ("Optimise", "Value realisation &", "continuous improvement"),
        ("Embed", "New operating model", "fully sustained"),
    ]
    num_nodes = len(nodes)
    node_size = 44
    # Arc spans from x_left to x_right, vertical center pulled up at extremes, lowest in middle
    arc_left = body_left + 60
    arc_right = body_left + body_w - 60
    arc_baseline_y = body_top + body_h // 2 + 10  # vertical center of arc
    arc_height = 36  # how far center bows down (arc center sits below endpoints)

    positions = []
    for i in range(num_nodes):
        t = i / (num_nodes - 1)  # 0..1
        x = arc_left + t * (arc_right - arc_left)
        # parabolic arc: y at endpoints = baseline-arc_height; mid = baseline+arc_height
        # use 4t(1-t) for symmetric arc midpoint maximum
        y_offset = arc_height * (4 * t * (1 - t)) - arc_height / 2
        y = arc_baseline_y + y_offset
        positions.append((x - node_size / 2, y - node_size / 2))

    # Draw connectors between successive nodes
    for i in range(num_nodes - 1):
        x1 = positions[i][0] + node_size / 2
        y1 = positions[i][1] + node_size / 2
        x2 = positions[i + 1][0] + node_size / 2
        y2 = positions[i + 1][1] + node_size / 2
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1), px_to_emu(y1), px_to_emu(x2), px_to_emu(y2),
        )
        conn.name = f"arc-seg-{i+1}"
        conn.line.color.rgb = BRAND_ACCENT_SOFT
        conn.line.width = 19050  # ~2pt

    # Draw nodes + labels (alternate label placement above/below arc)
    for i, ((nx, ny), (title, capA, capB)) in enumerate(zip(positions, nodes)):
        n = i + 1
        # Title block placement: above if node sits in lower half of arc, below if upper
        # Since arc dips down in middle, mid nodes are lower → place labels ABOVE.
        # Endpoints are higher → place labels BELOW.
        # Use simple rule: nodes 2..6 (middle 5) → labels above; 1 and 7 → labels below.
        label_above = (1 <= i <= num_nodes - 2)

        # Node circle
        _circle(slide, f"node-{n}",
                nx, ny, node_size, BRAND_PRIMARY,
                text=str(n), text_color=WHITE, font_size=16)

        # Title + caption block
        label_w = 140
        lx = nx + node_size / 2 - label_w / 2
        if label_above:
            ly = ny - 64
        else:
            ly = ny + node_size + 10
        add_text(slide, f"node-{n}-title", title,
                 x_px=int(lx), y_px=int(ly), w_px=label_w, h_px=18,
                 font_size_px=12, color=BRAND_PRIMARY, bold=True, align="center")
        add_text(slide, f"node-{n}-capA", capA,
                 x_px=int(lx), y_px=int(ly) + 20, w_px=label_w, h_px=14,
                 font_size_px=10, color=TEXT_MID, align="center")
        add_text(slide, f"node-{n}-capB", capB,
                 x_px=int(lx), y_px=int(ly) + 34, w_px=label_w, h_px=14,
                 font_size_px=10, color=TEXT_MID, align="center")

    add_footer(slide, page_num=334)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "334_7bucket-timeline-arc.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
