"""
Builder for pattern 68: Donut breakdown — donut chart + side breakdown table.

SVG-decompose (per SHAPE-ROLES table): build donut as a stacked-arc approximation
using primitive shapes. Since python-pptx pie/donut native charts are heavy,
we use a simpler oval+masking-like approach: render the donut as a single
chart-canvas oval placeholder (the segment IDs still address slice IDs by data).

Actually: since slice text/percentage labels live OUTSIDE the SVG (in the
breakdown panel), and the slices themselves are non-textual, we use an
oval placeholder for the visual but stamp donut-slice-N IDs on small swatches
in the breakdown panel that map 1:1 to the slices.

Source HTML: _pattern-library/68_donut-breakdown.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from twins.helpers import px_to_emu


DONUT_BG = RGBColor(0xF8, 0xF4, 0xFC)


def _add_oval(slide, shape_id, x, y, w, h, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px_to_emu(x), px_to_emu(y),
        px_to_emu(w), px_to_emu(h),
    )
    shape.name = shape_id
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_text(
        slide, "title",
        "Where the time goes — 30% to rework alone.",
        x_px=64, y_px=50, w_px=1100, h_px=40,
        font_size_px=28, color=TEXT_DARK, bold=True,
    )
    add_text(
        slide, "subtitle",
        "Per-deck baseline across the practice. The rework slice is the kill — Slide Lab targets it directly.",
        x_px=64, y_px=98, w_px=1100, h_px=22,
        font_size_px=14, color=TEXT_MID, italic=True,
    )
    add_rect(slide, "brand-rule", 64, 144, 56, 3, BRAND_ACCENT)

    # Body grid 45/55
    g_top = 200
    g_left = 64
    g_right = 1280 - 64
    g_bottom = 720 - 110  # leave room for convergence + footer
    g_w = g_right - g_left
    g_h = g_bottom - g_top
    donut_panel_w = int(g_w * 0.45)
    bd_x = g_left + donut_panel_w + 24
    bd_w = g_w - donut_panel_w - 24

    # Donut: outer oval BRAND_PRIMARY tint, inner oval white to make ring
    donut_d = 360
    dx = g_left + (donut_panel_w - donut_d) // 2
    dy = g_top + (g_h - donut_d) // 2

    # Use 5 ovals: track + 4 segments approximated by colored ring slices is hard
    # in PPTX. We render an outer oval as the focal accent ring, a primary-tint
    # main donut, and a hole. The 4 donut-slice-N IDs are anchored to small swatch
    # rectangles inside the donut wrap for addressability.
    _add_oval(slide, "donut-focal-ring", dx - 6, dy - 6, donut_d + 12, donut_d + 12, BRAND_ACCENT_SOFT)
    _add_oval(slide, "donut-track", dx, dy, donut_d, donut_d, DONUT_BG)

    # Slice placeholders (visual: 4 small filled quadrant-arc approximations
    # using off-center colored ovals masked by inner hole). For simplicity, we
    # render 4 invisible-id swatches around the donut perimeter as anchor shapes.
    slice_colors = [BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT_SOFT, BRAND_ACCENT]
    # Place tiny swatches at top, right, bottom, left of donut for slice IDs
    swatch_size = 24
    cx_center = dx + donut_d // 2
    cy_center = dy + donut_d // 2
    swatch_positions = [
        (cx_center - swatch_size // 2, dy - 10),  # top
        (dx + donut_d - 10, cy_center - swatch_size // 2),  # right
        (cx_center - swatch_size // 2, dy + donut_d - swatch_size + 10),  # bottom
        (dx - 14, cy_center - swatch_size // 2),  # left
    ]
    for si, (sx, sy) in enumerate(swatch_positions):
        sn = si + 1
        _add_oval(slide, f"donut-slice-{sn}", sx, sy, swatch_size, swatch_size, slice_colors[si])

    # Inner hole
    inner_d = donut_d - 110
    _add_oval(slide, "donut-hole", dx + 55, dy + 55, inner_d, inner_d, WHITE)

    # Donut center text
    add_text(
        slide, "donut-center-label", "PER DECK",
        x_px=dx, y_px=dy + 130, w_px=donut_d, h_px=14,
        font_size_px=10, color=TEXT_FAINT, bold=True, align="center", uppercase=True,
    )
    add_text(
        slide, "donut-center-value", "60 hr",
        x_px=dx, y_px=dy + 148, w_px=donut_d, h_px=60,
        font_size_px=56, color=BRAND_PRIMARY, bold=True, align="center",
    )
    add_text(
        slide, "donut-center-sub", "baseline effort, end to end",
        x_px=dx + 60, y_px=dy + 218, w_px=donut_d - 120, h_px=22,
        font_size_px=11, color=TEXT_MID, italic=True, align="center",
    )

    # Breakdown panel (right)
    # Header row
    bd_top = g_top + 14
    head_h = 26
    add_text(
        slide, "breakdown-col-name-header", "ACTIVITY",
        x_px=bd_x + 36, y_px=bd_top, w_px=bd_w - 200, h_px=14,
        font_size_px=10, color=TEXT_FAINT, bold=True, uppercase=True,
    )
    add_text(
        slide, "breakdown-col-num-header", "HOURS",
        x_px=bd_x + bd_w - 130, y_px=bd_top, w_px=60, h_px=14,
        font_size_px=10, color=TEXT_FAINT, bold=True, align="right", uppercase=True,
    )
    add_text(
        slide, "breakdown-col-pct-header", "SHARE",
        x_px=bd_x + bd_w - 60, y_px=bd_top, w_px=50, h_px=14,
        font_size_px=10, color=TEXT_FAINT, bold=True, align="right", uppercase=True,
    )
    add_rect(slide, "breakdown-head-rule", bd_x, bd_top + 20, bd_w, 1, CARD_BORDER)

    # 4 rows
    rows = [
        (BRAND_PRIMARY, "Storyline coaching", "framing, governing thought, structure", "10 hr", "17%", False),
        (BRAND_PRIMARY_MID, "Building slides", "layout, charts, visuals, formatting", "20 hr", "33%", False),
        (BRAND_ACCENT_SOFT, "Reviewing", "partner reads, comments, sign-off", "12 hr", "20%", False),
        (BRAND_ACCENT, "Reworking after reviews", "redo cycles, late-stage edits, version churn", "18 hr", "30%", True),
    ]
    row_top = bd_top + 32
    row_h = (g_h - 70) // 4

    # Focal tag above row 4
    focal_y = row_top + 3 * row_h - 22
    add_text(
        slide, "breakdown-focal-tag", "THE KILL · SLIDE LAB TARGETS THIS",
        x_px=bd_x + 5, y_px=focal_y, w_px=240, h_px=14,
        font_size_px=9, color=WHITE, bold=True, align="center", uppercase=True,
        bg_fill=BRAND_ACCENT, padding_px=(2, 8, 2, 8),
    )

    for ri, (swatch_color, name, sub, hours, pct, focal) in enumerate(rows):
        n = ri + 1
        ry = row_top + ri * row_h

        if focal:
            add_rect(slide, f"breakdown-{n}-bg", bd_x, ry, bd_w, row_h, CARD_BG)
            add_rect(slide, f"breakdown-{n}-stripe", bd_x, ry, 3, row_h, BRAND_ACCENT)
            sx_off = 8
        else:
            sx_off = 0

        # Swatch
        add_rect(slide, f"breakdown-{n}-swatch",
                 bd_x + 5 + sx_off, ry + (row_h - 14) // 2, 14, 14, swatch_color)

        # Name + sub
        add_text(
            slide, f"breakdown-{n}-name", name,
            x_px=bd_x + 30 + sx_off, y_px=ry + 6, w_px=bd_w - 220, h_px=20,
            font_size_px=13, color=BRAND_PRIMARY if focal else TEXT_DARK, bold=True,
        )
        add_text(
            slide, f"breakdown-{n}-name-sub", sub,
            x_px=bd_x + 30 + sx_off, y_px=ry + 26, w_px=bd_w - 220, h_px=18,
            font_size_px=10, color=TEXT_MID, italic=True,
        )

        # Hours
        add_text(
            slide, f"breakdown-{n}-num", hours,
            x_px=bd_x + bd_w - 130, y_px=ry + 12, w_px=60, h_px=22,
            font_size_px=16, color=BRAND_ACCENT if focal else BRAND_PRIMARY, bold=True, align="right",
        )
        # Pct
        add_text(
            slide, f"breakdown-{n}-pct", pct,
            x_px=bd_x + bd_w - 60, y_px=ry + 14, w_px=50, h_px=18,
            font_size_px=12, color=BRAND_PRIMARY_MID if focal else TEXT_MID,
            bold=True, align="right",
        )

    # Convergence (centered italic)
    add_text(
        slide, "convergence",
        "30% of deck time goes to rework. Slide Lab kills this — structure-first means edits land once.",
        x_px=64, y_px=720 - 60 - 24, w_px=1280 - 128, h_px=22,
        font_size_px=13, color=BRAND_PRIMARY, italic=True, align="center",
    )

    add_footer(slide, page_num=68)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "68_donut-breakdown.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
