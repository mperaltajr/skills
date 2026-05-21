"""
Builder for pattern 270: Disruption radar.

Source HTML: _pattern-library/270_disruption-radar.html

Layout: title + radar chart (6 axes, 2 series: Current vs 3-Year Forecast) on left,
top disruption risks panel on right.

LEGEND PLACEMENT: Right-aligned below subheadline (top-y ≥ 230, right edge ≈ 1240).
Series 1 = Current Exposure (solid brand-accent). Series 2 = 3-Year Forecast
(softer outline).
"""
from pathlib import Path
import sys
import math

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect, px_to_emu,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

RED_SEV = RGBColor(0xE5, 0x39, 0x35)
AMB_SEV = RGBColor(0xF5, 0x9E, 0x0B)


def _connector(slide, x1, y1, x2, y2, color, width_emu, dashed=False):
    line = slide.shapes.add_connector(1, 0, 0, 0, 0)
    line.begin_x = px_to_emu(x1)
    line.begin_y = px_to_emu(y1)
    line.end_x = px_to_emu(x2)
    line.end_y = px_to_emu(y2)
    line.line.color.rgb = color
    line.line.width = width_emu
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def build():
    prs, slide = new_slide()
    add_chrome(slide)
    add_title_block(
        slide,
        title="Disruption radar — <strong>where exposure is rising fastest over the next three years.</strong>",
        subtitle="Six disruption vectors scored 0–100; competitive and technological axes carry the highest current and forecast exposure.",
    )

    # === LEGEND — below subheadline, right-aligned ===
    leg_y = 230
    leg_w = 380
    leg_x = 1240 - leg_w
    add_rect(slide, "legend-current-swatch", leg_x, leg_y + 7, 10, 10, BRAND_ACCENT)
    add_text(slide, "legend-current", "Current Exposure",
             x_px=leg_x + 16, y_px=leg_y + 2, w_px=140, h_px=18,
             font_size_px=11, color=TEXT_MID, bold=True)
    add_rect(slide, "legend-forecast-swatch", leg_x + 160, leg_y + 7, 10, 10, BRAND_ACCENT_SOFT)
    add_text(slide, "legend-forecast", "3-Year Forecast",
             x_px=leg_x + 176, y_px=leg_y + 2, w_px=140, h_px=18,
             font_size_px=11, color=TEXT_MID, bold=True)

    # Body
    body_top = 268
    body_bottom = 635
    body_h = body_bottom - body_top
    left_x = 48
    right_x = 1280 - 48
    body_w = right_x - left_x

    radar_w = int(body_w * 0.58)
    gap = 18
    panel_w = body_w - radar_w - gap

    # Radar area
    radar_x = left_x
    radar_y = body_top
    radar_bg = add_rect(slide, "chart-canvas", radar_x, radar_y, radar_w, body_h, CARD_BG)
    radar_bg.line.color.rgb = CARD_BORDER
    radar_bg.line.width = 9525

    # Center of radar
    cx = radar_x + radar_w // 2
    cy = radar_y + body_h // 2 - 10
    R = 130

    # 6 axes at -90 + i*60 degrees
    axes = [
        ("Technological", 65, 82),
        ("Regulatory", 45, 68),
        ("Competitive", 72, 85),
        ("Customer", 58, 71),
        ("Operational", 41, 55),
        ("Economic", 38, 62),
    ]
    angles = []
    for i in range(6):
        deg = -90 + i * 60
        rad = math.radians(deg)
        angles.append((math.cos(rad), math.sin(rad)))

    # Concentric scale rings (4)
    for ring_idx, frac in enumerate([0.25, 0.5, 0.75, 1.0]):
        rr = int(R * frac)
        # draw hexagon
        pts = []
        for i in range(6):
            ux, uy = angles[i]
            pts.append((cx + ux * rr, cy + uy * rr))
        # connect points
        for i in range(6):
            x1, y1 = pts[i]
            x2, y2 = pts[(i + 1) % 6]
            ln = _connector(slide, x1, y1, x2, y2, CARD_BORDER, 6350)

    # Axis lines (center to outer vertex)
    for i, (ux, uy) in enumerate(angles):
        _connector(slide, cx, cy, cx + ux * R, cy + uy * R, CARD_BORDER, 6350)

    # Forecast polygon (draw first - underneath) - light fill via lines only
    f_pts = []
    for i, (ux, uy) in enumerate(angles):
        v = axes[i][2]
        f_pts.append((cx + ux * (v / 100) * R, cy + uy * (v / 100) * R))
    for i in range(6):
        x1, y1 = f_pts[i]
        x2, y2 = f_pts[(i + 1) % 6]
        _connector(slide, x1, y1, x2, y2, BRAND_ACCENT_SOFT, 19050)
    for i, (px, py) in enumerate(f_pts):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            px_to_emu(px - 4), px_to_emu(py - 4), px_to_emu(8), px_to_emu(8))
        d.name = f"dot-forecast-{i}"
        d.fill.solid()
        d.fill.fore_color.rgb = BRAND_ACCENT_SOFT
        d.line.fill.background()

    # Current polygon (on top)
    c_pts = []
    for i, (ux, uy) in enumerate(angles):
        v = axes[i][1]
        c_pts.append((cx + ux * (v / 100) * R, cy + uy * (v / 100) * R))
    for i in range(6):
        x1, y1 = c_pts[i]
        x2, y2 = c_pts[(i + 1) % 6]
        _connector(slide, x1, y1, x2, y2, BRAND_ACCENT, 25400)
    for i, (px, py) in enumerate(c_pts):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            px_to_emu(px - 5), px_to_emu(py - 5), px_to_emu(10), px_to_emu(10))
        d.name = f"dot-current-{i}"
        d.fill.solid()
        d.fill.fore_color.rgb = BRAND_ACCENT
        d.line.fill.background()

    # Axis labels (outside outer ring)
    label_positions = [
        (0, -1, "center"), (0.866, -0.5, "left"), (0.866, 0.5, "left"),
        (0, 1, "center"), (-0.866, 0.5, "right"), (-0.866, -0.5, "right"),
    ]
    for i, (ux, uy) in enumerate(angles):
        name = axes[i][0]
        lx = cx + ux * (R + 18) - 50
        ly = cy + uy * (R + 18) - 8
        a = label_positions[i][2]
        add_text(slide, f"label-{i}", name,
                 x_px=lx, y_px=ly, w_px=100, h_px=16,
                 font_size_px=10, color=BRAND_PRIMARY, bold=True, align=a)

    # Right risk panel
    pn_x = radar_x + radar_w + gap
    pn_y = body_top
    pn_bg = add_rect(slide, "risk-panel", pn_x, pn_y, panel_w, body_h, WHITE)
    pn_bg.line.color.rgb = CARD_BORDER
    pn_bg.line.width = 9525

    add_text(slide, "panel-header", "TOP DISRUPTION RISKS",
             x_px=pn_x + 18, y_px=pn_y + 16, w_px=panel_w - 36, h_px=16,
             font_size_px=11, color=BRAND_PRIMARY_MID, bold=True,
             letter_spacing_px=1.6, uppercase=True)

    risks = [
        ("Competitive", "72 → 85", "red",
         "AI-native entrants gaining SMB share at 3× our speed"),
        ("Technological", "65 → 82", "amber",
         "Generative AI reshapes product delivery economics"),
        ("Regulatory", "45 → 68", "amber",
         "EU AI Act compliance requirements from Q4 2026"),
    ]
    card_top = pn_y + 46
    card_area_h = body_h - 60
    card_gap = 10
    card_h = (card_area_h - 2 * card_gap) // 3

    for i, (axis, score, sev, desc) in enumerate(risks):
        n = i + 1
        cy_c = card_top + i * (card_h + card_gap)
        c = add_rect(slide, f"risk-card-{n}", pn_x + 14, cy_c, panel_w - 28, card_h, CARD_BG)
        c.line.color.rgb = CARD_BORDER
        c.line.width = 9525
        # severity stripe
        sev_color = RED_SEV if sev == "red" else AMB_SEV
        add_rect(slide, f"risk-{n}-stripe", pn_x + 14, cy_c, 4, card_h, sev_color)
        add_text(slide, f"risk-{n}-axis", axis.upper(),
                 x_px=pn_x + 24, y_px=cy_c + 10, w_px=140, h_px=14,
                 font_size_px=10, color=BRAND_PRIMARY, bold=True,
                 letter_spacing_px=1)
        add_text(slide, f"risk-{n}-score", score,
                 x_px=pn_x + panel_w - 100, y_px=cy_c + 10, w_px=80, h_px=14,
                 font_size_px=11, color=BRAND_ACCENT, bold=True, align="right")
        add_text(slide, f"risk-{n}-desc", desc,
                 x_px=pn_x + 24, y_px=cy_c + 30, w_px=panel_w - 50, h_px=card_h - 36,
                 font_size_px=11, color=TEXT_DARK)

    add_footer(slide, page_num=270)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "270_disruption-radar.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
