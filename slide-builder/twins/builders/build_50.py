"""
Builder for pattern 50: Two parallel paths converging.

Dark-gradient slide with two horizontal lanes (Path A behavioral, Path B knowledge),
each with 3 nodes connected by arrows, plus a right-side convergence panel.

Source HTML: _pattern-library/50_two-parallel-paths-converging.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    TEXT_FAINT, WHITE, DRAFT_BG, DRAFT_TEXT,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from twins.helpers import px_to_emu

# Dark theme tints
DARK_BG = RGBColor(0x2D, 0x0A, 0x4E)
DARK_BG_DEEP = RGBColor(0x1A, 0x05, 0x30)
PATH_A = BRAND_ACCENT
PATH_A_SOFT = BRAND_ACCENT_SOFT
PATH_B = BRAND_PRIMARY_MID
PATH_B_SOFT = RGBColor(0xB7, 0x9D, 0xD6)
NODE_A_INPUT_BG = RGBColor(0x4C, 0x18, 0x82)   # tinted accent over dark bg
NODE_A_OUT_BG = RGBColor(0x6A, 0x1A, 0xB0)
NODE_A_MID_BG = RGBColor(0x3A, 0x1A, 0x5E)
NODE_B_INPUT_BG = RGBColor(0x4B, 0x26, 0x70)
NODE_B_OUT_BG = RGBColor(0x57, 0x2C, 0x82)
NODE_B_MID_BG = RGBColor(0x3A, 0x1A, 0x5E)
WHITE_50 = RGBColor(0xC0, 0xB6, 0xCF)
WHITE_70 = RGBColor(0xDF, 0xD6, 0xEC)
CONV_PANEL_BG = RGBColor(0x14, 0x04, 0x28)


def build():
    prs, slide = new_slide()

    # Override background to dark. The HTML uses a 135deg gradient
    # 1A0530 -> 2D0A4E -> 3E1370; we approximate with a solid mid-tone (DARK_BG)
    # since python-pptx slide-background gradient API is unreliable across viewers.
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG

    # --- Top chrome (dark variant) ---

    # --- Title block (light text on dark) ---
    add_text(
        slide, "eyebrow", "The combined outcome",
        x_px=56, y_px=50, w_px=900, h_px=14,
        font_size_px=10, color=BRAND_ACCENT_SOFT, bold=True,
        letter_spacing_px=2, uppercase=True,
    )
    add_text(
        slide, "title",
        "Two parallel paths — <strong>each creates value alone; together they compound</strong> into something durable.",
        x_px=56, y_px=70, w_px=1168, h_px=70,
        font_size_px=24, color=WHITE, bold=True,
        emphasis_color=BRAND_ACCENT_SOFT,
    )
    add_rect(slide, "brand-rule", x_px=56, y_px=144, w_px=56, h_px=3, fill_color=BRAND_ACCENT)

    # --- Lane labels ---
    add_text(
        slide, "lane-a-tag", "Path 1",
        x_px=48, y_px=180, w_px=124, h_px=14,
        font_size_px=9, color=PATH_A_SOFT, bold=True, letter_spacing_px=1.4, uppercase=True,
    )
    add_text(
        slide, "lane-a-label", "Behavioral track",
        x_px=48, y_px=196, w_px=124, h_px=18,
        font_size_px=10, color=PATH_A_SOFT, bold=True, letter_spacing_px=1.6, uppercase=True,
    )
    add_text(
        slide, "lane-b-tag", "Path 2",
        x_px=48, y_px=408, w_px=124, h_px=14,
        font_size_px=9, color=PATH_B_SOFT, bold=True, letter_spacing_px=1.4, uppercase=True,
    )
    add_text(
        slide, "lane-b-label", "Knowledge track",
        x_px=48, y_px=424, w_px=124, h_px=18,
        font_size_px=10, color=PATH_B_SOFT, bold=True, letter_spacing_px=1.6, uppercase=True,
    )

    # --- Path A nodes (y=200, height=108) ---
    a_nodes = [
        (184, "a-input", "Behavioral input", "", "How we think, engage, and prioritize across the account", "input"),
        (414, "a-mid", "Creates", "Trusted relationships", "Client treats us as advisor, not vendor", "mid"),
        (656, "a-out", "Leads to", "INFLUENCE", "We shape decisions before they are made", "out"),
    ]
    for i, (x, sid, eyebrow, mid_title, desc, kind) in enumerate(a_nodes):
        n = i + 1
        w = 184 if i != 1 else 196
        node_bg = {"input": NODE_A_INPUT_BG, "mid": NODE_A_MID_BG, "out": NODE_A_OUT_BG}[kind]
        nrect = add_rect(slide, f"path-a-node-{n}", x, 200, w, 108, node_bg)
        nrect.line.color.rgb = PATH_A
        nrect.line.width = 19050
        add_text(
            slide, f"path-a-node-{n}-eyebrow", eyebrow,
            x_px=x + 14, y_px=208, w_px=w - 28, h_px=14,
            font_size_px=9, color=PATH_A_SOFT, bold=True,
            letter_spacing_px=1.4, uppercase=True,
        )
        if kind == "out":
            # Big outcome word
            add_text(
                slide, f"path-a-node-{n}-outcome", desc.split(" ")[0] if False else "INFLUENCE",
                x_px=x + 14, y_px=224, w_px=w - 28, h_px=22,
                font_size_px=17, color=PATH_A_SOFT, bold=True, letter_spacing_px=1,
            )
            # Bare description fallback already set
            add_text(
                slide, f"path-a-node-{n}-desc", "We shape decisions before they are made",
                x_px=x + 14, y_px=252, w_px=w - 28, h_px=50,
                font_size_px=10, color=WHITE_70,
            )
        elif kind == "mid":
            add_text(
                slide, f"path-a-node-{n}-title", mid_title,
                x_px=x + 14, y_px=224, w_px=w - 28, h_px=18,
                font_size_px=12, color=WHITE, bold=True,
            )
            add_text(
                slide, f"path-a-node-{n}-desc", desc,
                x_px=x + 14, y_px=246, w_px=w - 28, h_px=54,
                font_size_px=10, color=WHITE_70,
            )
        else:
            add_text(
                slide, f"path-a-node-{n}-desc", desc,
                x_px=x + 14, y_px=232, w_px=w - 28, h_px=70,
                font_size_px=10, color=WHITE_70,
            )

    # --- Path B nodes (y=428) ---
    b_nodes = [
        (184, "b-input", "Knowledge input", "", "What we know and deliver about the client's business", "input"),
        (414, "b-mid", "Generates", "Operational insights", "We surface value gaps the client didn't see coming", "mid"),
        (656, "b-out", "Builds", "CREDIBILITY", "Client listens — grounded in operations, not pitch", "out"),
    ]
    for i, (x, sid, eyebrow, mid_title, desc, kind) in enumerate(b_nodes):
        n = i + 1
        w = 184 if i != 1 else 196
        node_bg = {"input": NODE_B_INPUT_BG, "mid": NODE_B_MID_BG, "out": NODE_B_OUT_BG}[kind]
        nrect = add_rect(slide, f"path-b-node-{n}", x, 428, w, 108, node_bg)
        nrect.line.color.rgb = PATH_B_SOFT
        nrect.line.width = 19050
        add_text(
            slide, f"path-b-node-{n}-eyebrow", eyebrow,
            x_px=x + 14, y_px=436, w_px=w - 28, h_px=14,
            font_size_px=9, color=PATH_B_SOFT, bold=True,
            letter_spacing_px=1.4, uppercase=True,
        )
        if kind == "out":
            add_text(
                slide, f"path-b-node-{n}-outcome", mid_title or "CREDIBILITY",
                x_px=x + 14, y_px=452, w_px=w - 28, h_px=22,
                font_size_px=17, color=PATH_B_SOFT, bold=True, letter_spacing_px=1,
            )
            add_text(
                slide, f"path-b-node-{n}-desc", desc,
                x_px=x + 14, y_px=480, w_px=w - 28, h_px=50,
                font_size_px=10, color=WHITE_70,
            )
        elif kind == "mid":
            add_text(
                slide, f"path-b-node-{n}-title", mid_title,
                x_px=x + 14, y_px=452, w_px=w - 28, h_px=18,
                font_size_px=12, color=WHITE, bold=True,
            )
            add_text(
                slide, f"path-b-node-{n}-desc", desc,
                x_px=x + 14, y_px=474, w_px=w - 28, h_px=54,
                font_size_px=10, color=WHITE_70,
            )
        else:
            add_text(
                slide, f"path-b-node-{n}-desc", desc,
                x_px=x + 14, y_px=460, w_px=w - 28, h_px=70,
                font_size_px=10, color=WHITE_70,
            )

    # --- Arrows (small triangles between nodes) ---
    arrow_positions_a = [(372, 254), (614, 254), (844, 254)]
    arrow_positions_b = [(372, 482), (614, 482), (844, 482)]
    for i, (ax, ay) in enumerate(arrow_positions_a):
        n = i + 1
        # Stem
        add_rect(slide, f"arrow-a-{n}-stem", ax, ay - 1, 38, 2, PATH_A)
        # Triangle arrowhead
        tri = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            px_to_emu(ax + 36), px_to_emu(ay - 6),
            px_to_emu(12), px_to_emu(12),
        )
        tri.name = f"arrow-a-{n}-head"
        tri.fill.solid()
        tri.fill.fore_color.rgb = PATH_A
        tri.line.fill.background()
    for i, (ax, ay) in enumerate(arrow_positions_b):
        n = i + 1
        add_rect(slide, f"arrow-b-{n}-stem", ax, ay - 1, 38, 2, PATH_B_SOFT)
        tri = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            px_to_emu(ax + 36), px_to_emu(ay - 6),
            px_to_emu(12), px_to_emu(12),
        )
        tri.name = f"arrow-b-{n}-head"
        tri.fill.solid()
        tri.fill.fore_color.rgb = PATH_B_SOFT
        tri.line.fill.background()

    # --- Convergence panel (right side, vertical) ---
    cp_x = 896
    cp_y = 184
    cp_w = 1280 - 56 - cp_x
    cp_h = 720 - 96 - cp_y

    cp = add_rect(slide, "conv-panel-bg", cp_x, cp_y, cp_w, cp_h, CONV_PANEL_BG)
    cp.line.color.rgb = WHITE_50
    cp.line.width = 9525

    add_text(
        slide, "conv-eyebrow", "Combined effect",
        x_px=cp_x + 22, y_px=cp_y + 18, w_px=cp_w - 44, h_px=14,
        font_size_px=10, color=WHITE_50, bold=True,
        letter_spacing_px=1.8, uppercase=True,
    )
    # Row A
    add_rect(slide, "conv-pill-a", cp_x + 22, cp_y + 44, 78, 18, PATH_A)
    add_text(
        slide, "conv-pill-a-text", "INFLUENCE",
        x_px=cp_x + 22, y_px=cp_y + 44, w_px=78, h_px=18,
        font_size_px=10, color=WHITE, bold=True, align="center", anchor="middle",
        letter_spacing_px=1,
    )
    add_text(
        slide, "conv-row-a-text", "We shape the agenda before the client asks",
        x_px=cp_x + 110, y_px=cp_y + 44, w_px=cp_w - 130, h_px=36,
        font_size_px=11, color=WHITE_70,
    )
    # Row B
    add_rect(slide, "conv-pill-b", cp_x + 22, cp_y + 90, 78, 18, PATH_B)
    add_text(
        slide, "conv-pill-b-text", "CREDIBILITY",
        x_px=cp_x + 22, y_px=cp_y + 90, w_px=78, h_px=18,
        font_size_px=10, color=WHITE, bold=True, align="center", anchor="middle",
        letter_spacing_px=1,
    )
    add_text(
        slide, "conv-row-b-text", "Client listens because we understand their business",
        x_px=cp_x + 110, y_px=cp_y + 90, w_px=cp_w - 130, h_px=36,
        font_size_px=11, color=WHITE_70,
    )
    # Divider
    add_rect(slide, "conv-divider", cp_x + 22, cp_y + 140, cp_w - 44, 1, WHITE_50)
    # Equals statement
    add_text(
        slide, "conv-equals",
        "= New work. Expanded scope. <em>Durable success.</em>",
        x_px=cp_x + 22, y_px=cp_y + 152, w_px=cp_w - 44, h_px=60,
        font_size_px=17, color=WHITE, bold=True,
        emphasis_color=BRAND_ACCENT_SOFT,
    )

    # Callout
    co_y = cp_y + cp_h - 90
    add_rect(slide, "conv-callout-bg", cp_x + 22, co_y, cp_w - 44, 76, NODE_A_INPUT_BG)
    add_rect(slide, "conv-callout-accent", cp_x + 22, co_y, 3, 76, BRAND_ACCENT)
    add_text(
        slide, "conv-callout-title",
        "Without one path — value is left on the table.",
        x_px=cp_x + 34, y_px=co_y + 10, w_px=cp_w - 62, h_px=16,
        font_size_px=11, color=BRAND_ACCENT_SOFT, bold=True,
    )
    add_text(
        slide, "conv-callout-body",
        "Operating with only one track today. Adding the second completes the equation and compounds returns.",
        x_px=cp_x + 34, y_px=co_y + 30, w_px=cp_w - 62, h_px=44,
        font_size_px=10, color=WHITE_70,
    )

    # --- Footer (dark variant) ---
    add_text(slide, "footnote-1", "1. [add footnote here or delete]",
             x_px=58, y_px=672, w_px=1164, h_px=16,
             font_size_px=10, color=TEXT_FAINT)
    add_text(slide, "source", "Source: [add source here or delete]",
             x_px=58, y_px=688, w_px=1100, h_px=16,
             font_size_px=10, color=TEXT_FAINT, italic=True)
    add_text(slide, "page-number", "50",
             x_px=1170, y_px=688, w_px=52, h_px=16,
             font_size_px=11, color=WHITE, align="right")

    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "50_two-parallel-paths-converging.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
