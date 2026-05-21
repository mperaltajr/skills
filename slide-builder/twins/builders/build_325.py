"""
Builder for pattern 325: 5-bucket radial pentagon (SVG-based) + annotation cards on right.

Left panel: pentagon with 5 numbered nodes and center "CORE" — rendered as
chart-canvas placeholder (the SVG is too complex for native shapes; treated as
a picture-asset role per Slide Lab).
Right panel: 5 numbered annotation cards explaining each node.

Source HTML: _pattern-library/325_5bucket-radial-pentagon.html
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.enum.shapes import MSO_SHAPE
from twins.helpers import px_to_emu
import math


def _circle(slide, name, x, y, size, fill, text=None, text_color=WHITE,
            font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, px_to_emu(x), px_to_emu(y),
        px_to_emu(size), px_to_emu(size),
    )
    shape.name = f"{name}-bg"
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text is not None:
        add_text(
            slide, f"{name}-text", text,
            x_px=x, y_px=y, w_px=size, h_px=size,
            font_size_px=font_size, color=text_color, bold=True,
            align="center", anchor="middle",
        )


def build():
    prs, slide = new_slide()
    add_chrome(slide)

    add_title_block(
        slide,
        title="Five <strong>strategic priorities</strong> converge on a single core",
        subtitle="Each bucket feeds the centre — progress in any node accelerates the whole system.",
        title_x=48, title_y=44, title_w=1184, title_h=44,
        subtitle_h=20, brand_rule_w=64,
    )

    body_top = 200
    body_left = 48
    body_w = 1184
    body_h = 720 - 32 - 14 - body_top

    # Left panel: pentagon (drawn natively with 5 ovals + central CORE)
    left_w = 520
    cx = body_left + left_w // 2
    cy = body_top + body_h // 2
    radius = 150
    node_size = 56
    core_size = 68

    # Pentagon vertices (start at top, go clockwise)
    nodes = []
    for i in range(5):
        angle = -math.pi / 2 + i * (2 * math.pi / 5)
        nx = cx + radius * math.cos(angle) - node_size / 2
        ny = cy + radius * math.sin(angle) - node_size / 2
        nodes.append((nx, ny))

    # Spokes (lines from each node center to core center) — use thin rects
    core_cx = cx
    core_cy = cy
    for i, (nx, ny) in enumerate(nodes):
        n = i + 1
        # We approximate with a thin diagonal rect — PowerPoint rect cannot rotate
        # easily without complex XML, so use a connector-like rectangle at each
        # node→center direction. Use thin rect aligned along the spoke.
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_CONNECTOR
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(nx + node_size / 2), px_to_emu(ny + node_size / 2),
            px_to_emu(core_cx), px_to_emu(core_cy),
        )
        connector.name = f"spoke-{n}"
        connector.line.color.rgb = BRAND_ACCENT_SOFT
        connector.line.width = 12700  # 1pt

    # Node circles (number badges)
    for i, (nx, ny) in enumerate(nodes):
        n = i + 1
        _circle(slide, f"node-{n}", nx, ny, node_size, BRAND_PRIMARY_MID,
                text=str(n), font_size=18)
    # Core
    _circle(slide, "core",
            core_cx - core_size // 2, core_cy - core_size // 2, core_size,
            BRAND_ACCENT, text="CORE", font_size=11)

    # ── Right panel: 5 annotation cards ──
    right_x = body_left + left_w + 24
    right_w = body_w - left_w - 24
    cards = [
        ("Strategic Vision",
         "Define the north-star objective and align leadership around a single measurable outcome."),
        ("Data & Analytics",
         "Instrument every touchpoint; feed real-time signals back into the core decision layer."),
        ("Talent & Capability",
         "Upskill teams on priority tools and embed new ways of working into daily operating rhythm."),
        ("Technology Platform",
         "Modernise the stack to unlock scale; retire legacy systems that create compounding drag."),
        ("Change & Adoption",
         "Drive sustained adoption through structured comms, governance, and value-realisation tracking."),
    ]
    card_gap = 8
    card_h = (body_h - 4 * card_gap) // 5
    for i, (title, desc) in enumerate(cards):
        n = i + 1
        cy_ = body_top + i * (card_h + card_gap)
        card = add_rect(slide, f"ann-{n}-card", right_x, cy_, right_w, card_h, CARD_BG)
        card.line.color.rgb = CARD_BORDER
        card.line.width = 9525
        # Number circle on left
        ns = 28
        _circle(slide, f"ann-{n}-num", right_x + 12, cy_ + (card_h - ns) // 2,
                ns, BRAND_PRIMARY_MID, text=str(n), font_size=12)
        add_text(
            slide, f"ann-{n}-title", title,
            x_px=right_x + 52, y_px=cy_ + 14, w_px=right_w - 64, h_px=22,
            font_size_px=13, color=BRAND_PRIMARY, bold=True,
        )
        add_text(
            slide, f"ann-{n}-desc", desc,
            x_px=right_x + 52, y_px=cy_ + 36, w_px=right_w - 64, h_px=card_h - 44,
            font_size_px=11, color=TEXT_MID,
        )

    add_footer(slide, page_num=325)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "325_5bucket-radial-pentagon.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
