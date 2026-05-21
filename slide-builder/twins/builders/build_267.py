"""
Builder for pattern 267: Win/Loss analysis table.

Source HTML: _pattern-library/267_win-loss-analysis-table.html

Layout: title + LEFT summary panel (64% win rate big number, donut, 34/53,
$2.4M avg deal, 94d cycle) + RIGHT table (7 factor rows with win/loss bars
and insight column).

No top-level legend — column headers are the implicit legend (Win % green,
Loss % red).
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor

WIN_FILL = RGBColor(0x22, 0xC5, 0x5E)
WIN_TEXT = RGBColor(0x16, 0xA3, 0x4A)
LOSS_FILL = RGBColor(0xEF, 0x44, 0x44)
LOSS_TEXT = RGBColor(0xDC, 0x26, 0x26)
BAR_TRACK = RGBColor(0xEE, 0xE8, 0xF5)


def build():
    prs, slide = new_slide()
    add_chrome(slide)
    add_title_block(
        slide,
        title="Win/Loss Analysis — <strong>What Drives Our Results</strong>",
        subtitle="FY26 competitive pursuit review · 53 deals scored · source: CRM + debrief interviews",
    )

    # Body
    body_top = 230
    body_bottom = 632
    body_h = body_bottom - body_top
    left_x = 48
    right_x = 1280 - 48
    body_w = right_x - left_x

    # Split: left summary 28%, table 72%
    sum_w = int(body_w * 0.28)
    gap = 18
    tbl_w = body_w - sum_w - gap

    # LEFT — summary
    sum_x = left_x
    sum_y = body_top
    sum_bg = add_rect(slide, "left-panel", sum_x, sum_y, sum_w, body_h, CARD_BG)
    sum_bg.line.color.rgb = CARD_BORDER
    sum_bg.line.width = 9525
    add_rect(slide, "left-accent", sum_x, sum_y, 3, body_h, BRAND_ACCENT)

    add_text(slide, "win-rate-label", "WIN RATE",
             x_px=sum_x + 18, y_px=sum_y + 18, w_px=sum_w - 36, h_px=14,
             font_size_px=10, color=BRAND_PRIMARY, bold=True,
             letter_spacing_px=1.4, uppercase=True, align="center")
    add_text(slide, "win-rate-number", "64%",
             x_px=sum_x + 18, y_px=sum_y + 36, w_px=sum_w - 36, h_px=70,
             font_size_px=56, color=BRAND_ACCENT, bold=True, align="center")

    # Donut placeholder - drawn as outer circle minus inner using OVAL shapes
    from pptx.enum.shapes import MSO_SHAPE
    from twins.helpers import px_to_emu
    donut_size = 120
    donut_x = sum_x + (sum_w - donut_size) // 2
    donut_y = sum_y + 112
    # outer ring
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        px_to_emu(donut_x), px_to_emu(donut_y),
        px_to_emu(donut_size), px_to_emu(donut_size))
    outer.name = "donut-outer"
    outer.fill.solid()
    outer.fill.fore_color.rgb = CARD_BORDER
    outer.line.fill.background()
    # 64% arc approximation: overlay a brand-accent partial
    arc = slide.shapes.add_shape(MSO_SHAPE.PIE,
        px_to_emu(donut_x), px_to_emu(donut_y),
        px_to_emu(donut_size), px_to_emu(donut_size))
    arc.name = "donut-arc"
    arc.fill.solid()
    arc.fill.fore_color.rgb = BRAND_ACCENT
    arc.line.fill.background()
    # python-pptx doesn't let us set the pie sweep angle directly without XML hacking; leave as full pie
    # inner white circle to make donut hole
    inner_size = 80
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        px_to_emu(donut_x + (donut_size - inner_size) // 2),
        px_to_emu(donut_y + (donut_size - inner_size) // 2),
        px_to_emu(inner_size), px_to_emu(inner_size))
    inner.name = "donut-hole"
    inner.fill.solid()
    inner.fill.fore_color.rgb = CARD_BG
    inner.line.fill.background()
    # center label
    add_text(slide, "donut-center", "64%",
             x_px=donut_x, y_px=donut_y + donut_size // 2 - 16, w_px=donut_size, h_px=20,
             font_size_px=14, color=BRAND_ACCENT, bold=True, align="center")
    add_text(slide, "donut-sub", "WIN",
             x_px=donut_x, y_px=donut_y + donut_size // 2 + 4, w_px=donut_size, h_px=14,
             font_size_px=8, color=TEXT_MID, bold=True, align="center",
             letter_spacing_px=1.2)

    add_text(slide, "pursuits-line", "34 wins / 53 pursuits · FY26",
             x_px=sum_x + 18, y_px=sum_y + 244, w_px=sum_w - 36, h_px=16,
             font_size_px=11, color=TEXT_DARK, align="center")

    add_rect(slide, "stats-divider", sum_x + 28, sum_y + 268, sum_w - 56, 1, CARD_BORDER)

    # Stats
    stats = [("Avg Deal Size", "$2.4M", "per pursuit"),
             ("Avg Cycle Length", "94", "days")]
    for i, (lbl, val, unit) in enumerate(stats):
        sy = sum_y + 280 + i * 50
        add_text(slide, f"stat-{i+1}-label", lbl,
                 x_px=sum_x + 18, y_px=sy, w_px=sum_w - 36, h_px=14,
                 font_size_px=9, color=TEXT_MID, bold=True,
                 letter_spacing_px=1.2, uppercase=True)
        add_text(slide, f"stat-{i+1}-value", f"{val} {unit}",
                 x_px=sum_x + 18, y_px=sy + 16, w_px=sum_w - 36, h_px=20,
                 font_size_px=14, color=BRAND_PRIMARY, bold=True)

    # RIGHT — factors table
    tbl_x = sum_x + sum_w + gap
    tbl_y = body_top
    # Column widths
    name_w = 150
    win_w = (tbl_w - name_w - 230) // 2
    loss_w = win_w
    ins_w = 230
    col_x = [tbl_x, tbl_x + name_w, tbl_x + name_w + win_w, tbl_x + name_w + win_w + loss_w]
    col_widths = [name_w, win_w, loss_w, ins_w]

    # header
    hdr_h = 30
    add_rect(slide, "table-header-bg", tbl_x, tbl_y, tbl_w, hdr_h, BRAND_PRIMARY)
    add_text(slide, "th-title", "WHY WE WIN / LOSE",
             x_px=col_x[0] + 14, y_px=tbl_y, w_px=name_w - 24, h_px=hdr_h,
             font_size_px=10, color=WHITE, bold=True, anchor="middle",
             letter_spacing_px=1.2, uppercase=True)
    add_text(slide, "th-win", "WIN CONTRIBUTION",
             x_px=col_x[1], y_px=tbl_y, w_px=win_w, h_px=hdr_h,
             font_size_px=9, color=RGBColor(0xC8, 0xE6, 0xC9), bold=True, align="center", anchor="middle",
             letter_spacing_px=1.2, uppercase=True)
    add_text(slide, "th-loss", "LOSS REASON",
             x_px=col_x[2], y_px=tbl_y, w_px=loss_w, h_px=hdr_h,
             font_size_px=9, color=RGBColor(0xFF, 0xCD, 0xD2), bold=True, align="center", anchor="middle",
             letter_spacing_px=1.2, uppercase=True)
    add_text(slide, "th-insight", "KEY INSIGHT",
             x_px=col_x[3] + 10, y_px=tbl_y, w_px=ins_w - 20, h_px=hdr_h,
             font_size_px=9, color=BRAND_ACCENT_SOFT, bold=True, anchor="middle",
             letter_spacing_px=1.2, uppercase=True)

    factors = [
        ("Solution Fit", 78, 22, "Core differentiator in 8/10 wins"),
        ("Relationships", 71, 29, "C-suite access critical"),
        ("Price", 44, 56, "Lost 18 deals on price alone"),
        ("Proposal Quality", 67, 33, "Late submissions cost 6 deals"),
        ("Speed to Respond", 55, 45, "Average 4-day response time"),
        ("Technical Depth", 82, 18, "Our strongest attribute"),
        ("References", 69, 31, "3 relevant references needed"),
    ]
    row_h = (body_h - hdr_h) // len(factors)
    for i, (name, w_pct, l_pct, insight) in enumerate(factors):
        n = i + 1
        ry = tbl_y + hdr_h + i * row_h
        # row bg alternating
        if i % 2 == 1:
            add_rect(slide, f"row-{n}-bg", tbl_x, ry, tbl_w, row_h, CARD_BG)
        add_rect(slide, f"row-{n}-rule", tbl_x, ry, tbl_w, 1, CARD_BORDER)
        # factor name
        add_text(slide, f"row-{n}-name", name,
                 x_px=col_x[0] + 14, y_px=ry, w_px=name_w - 24, h_px=row_h,
                 font_size_px=11, color=BRAND_PRIMARY, bold=True, anchor="middle")
        # win bar
        bar_h = 12
        track_w = win_w - 30
        track_x = col_x[1] + 15
        track_y = ry + row_h // 2 - bar_h // 2 - 5
        add_rect(slide, f"row-{n}-win-track", track_x, track_y, track_w, bar_h, BAR_TRACK)
        add_rect(slide, f"row-{n}-win-fill", track_x, track_y,
                 int(track_w * w_pct / 100), bar_h, WIN_FILL)
        add_text(slide, f"row-{n}-win-pct", f"{w_pct}%",
                 x_px=track_x, y_px=track_y + bar_h + 2, w_px=track_w, h_px=14,
                 font_size_px=9, color=WIN_TEXT, bold=True, align="center")
        # loss bar
        track_x2 = col_x[2] + 15
        track_w2 = loss_w - 30
        add_rect(slide, f"row-{n}-loss-track", track_x2, track_y, track_w2, bar_h, BAR_TRACK)
        add_rect(slide, f"row-{n}-loss-fill", track_x2, track_y,
                 int(track_w2 * l_pct / 100), bar_h, LOSS_FILL)
        add_text(slide, f"row-{n}-loss-pct", f"{l_pct}%",
                 x_px=track_x2, y_px=track_y + bar_h + 2, w_px=track_w2, h_px=14,
                 font_size_px=9, color=LOSS_TEXT, bold=True, align="center")
        # insight
        add_text(slide, f"row-{n}-insight", insight,
                 x_px=col_x[3] + 10, y_px=ry, w_px=ins_w - 20, h_px=row_h,
                 font_size_px=10, color=TEXT_MID, italic=True, anchor="middle")

    add_footer(slide, page_num=267)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "267_win-loss-analysis-table.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
