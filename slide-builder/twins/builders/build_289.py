"""
Builder for pattern 289: Interview Synthesis Grid (5 themes × 5 roles + key themes panel).

Source HTML: _pattern-library/289_interview-synthesis-grid.html

Matrix on left, sentiment icons (pos/neu/neg) in each cell. Right panel: key themes.
"""
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from twins.helpers import (
    new_slide, add_text, add_rect,
    add_chrome, add_footer, add_title_block,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    CARD_BG, CARD_BORDER, TEXT_DARK, TEXT_MID, TEXT_FAINT, WHITE,
)
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

POS_BG = RGBColor(0xDC, 0xFC, 0xE7)
POS_FG = RGBColor(0x16, 0xA3, 0x4A)
NEU_BG = RGBColor(0xFE, 0xF9, 0xC3)
NEU_FG = RGBColor(0xCA, 0x8A, 0x04)
NEG_BG = RGBColor(0xFE, 0xE2, 0xE2)
NEG_FG = RGBColor(0xDC, 0x26, 0x26)


def sentiment(slide, sid, cx, cy, size, kind):
    if kind == "pos":
        bg, fg, glyph = POS_BG, POS_FG, "✓"
    elif kind == "neg":
        bg, fg, glyph = NEG_BG, NEG_FG, "✕"
    else:
        bg, fg, glyph = NEU_BG, NEU_FG, "–"
    pill = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx * 9525, cy * 9525, size * 9525, size * 9525,
    )
    pill.name = f"{sid}-bg"
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg
    pill.line.fill.background()
    add_text(
        slide, sid, glyph,
        x_px=cx, y_px=cy, w_px=size, h_px=size,
        font_size_px=13, color=fg, bold=True, align="center", anchor="middle",
    )


def build():
    prs, slide = new_slide()
    add_chrome(slide)
    add_title_block(
        slide,
        title="Interview Synthesis — <strong>Stakeholder Sentiment Grid</strong>",
        subtitle="Aggregated signals across 5 roles × 4 strategic themes · May 2026",
    )

    # Layout
    body_x = 48
    body_y = 232
    body_w = 1280 - 96
    body_h = 720 - body_y - 44

    # Right panel: Key Themes
    kt_w = 220
    kt_gap = 12
    matrix_w = body_w - kt_w - kt_gap

    # Matrix
    # Column headers: blank corner + 5 role avatars
    label_col_w = 160
    role_col_w = (matrix_w - label_col_w) // 5

    # Col header row
    ch_y = body_y
    ch_h = 60
    roles = [("CE", "CEO"), ("CF", "CFO"), ("VO", "VP Ops"),
             ("DI", "Dir IT"), ("MG", "Manager")]
    for i, (initials, role) in enumerate(roles):
        n = i + 2
        cx = body_x + label_col_w + i * role_col_w + role_col_w // 2
        # avatar
        av_size = 32
        ax = cx - av_size // 2
        ay = ch_y + 4
        av = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    ax * 9525, ay * 9525, av_size * 9525, av_size * 9525)
        av.name = f"col-{n}-avatar"
        av.fill.solid()
        av.fill.fore_color.rgb = BRAND_PRIMARY_MID
        av.line.fill.background()
        add_text(
            slide, f"col-{n}-avatar-text", initials,
            x_px=ax, y_px=ay, w_px=av_size, h_px=av_size,
            font_size_px=10, color=WHITE, bold=True, align="center", anchor="middle",
        )
        add_text(
            slide, f"col-{n}-role", role,
            x_px=cx - 50, y_px=ay + av_size + 4, w_px=100, h_px=16,
            font_size_px=10, color=BRAND_PRIMARY, bold=True, align="center",
        )

    # Data rows
    rows_top = ch_y + ch_h + 6
    rows_h = body_h - ch_h - 6
    themes = [
        ("Culture", "Org mindset & change readiness",
         ["pos", "neu", "pos", "neu", "neg"], False),
        ("Process", "Workflow efficiency & standards",
         ["neg", "neg", "neu", "pos", "neg"], True),
        ("Technology", "Tool adoption & data maturity",
         ["pos", "pos", "neu", "pos", "neu"], False),
        ("Leadership", "Vision alignment & sponsorship",
         ["pos", "neu", "neg", "neu", "pos"], True),
        ("Consensus", "Majority view per theme",
         ["pos", "neu", "neu", "pos", "neg"], False),  # consensus row special
    ]
    row_h = rows_h // len(themes)

    for ri, (theme, desc, sentiments, alt) in enumerate(themes):
        ry = rows_top + ri * row_h
        is_consensus = (theme == "Consensus")
        if is_consensus:
            row_bg = BRAND_PRIMARY
        elif alt:
            row_bg = CARD_BG
        else:
            row_bg = WHITE
        add_rect(slide, f"row-{ri+1}-bg", body_x, ry, matrix_w, row_h, row_bg)

        # Row header
        theme_color = WHITE if is_consensus else BRAND_PRIMARY
        desc_color = RGBColor(0xC7, 0xB0, 0xE0) if is_consensus else TEXT_MID
        add_text(
            slide, f"row-{ri+1}-theme", theme,
            x_px=body_x + 12, y_px=ry + 12, w_px=label_col_w - 24, h_px=16,
            font_size_px=12, color=theme_color, bold=True,
            uppercase=True, letter_spacing_px=0.4,
        )
        add_text(
            slide, f"row-{ri+1}-desc", desc,
            x_px=body_x + 12, y_px=ry + 30, w_px=label_col_w - 24, h_px=row_h - 36,
            font_size_px=10, color=desc_color,
        )

        # Sentiment cells
        for ci, kind in enumerate(sentiments):
            cn = ci + 2
            cx_center = body_x + label_col_w + ci * role_col_w + role_col_w // 2
            sz = 28
            sentiment(
                slide, f"row-{ri+1}-cell-{cn}",
                cx_center - sz // 2, ry + (row_h - sz) // 2, sz, kind,
            )

    # Key themes panel
    kt_x = body_x + matrix_w + kt_gap
    # Header
    add_text(
        slide, "kt-header", "KEY THEMES",
        x_px=kt_x, y_px=body_y, w_px=kt_w, h_px=18,
        font_size_px=10, color=BRAND_ACCENT, bold=True,
        uppercase=True, letter_spacing_px=1.4,
    )
    add_rect(slide, "kt-rule", kt_x, body_y + 22, kt_w, 2, BRAND_ACCENT)

    kt_items = [
        ("Culture", "Leadership and Ops are aligned; frontline staff remain hesitant to change.", False),
        ("Process", "Widespread friction in workflows; only IT sees structured improvement efforts.", False),
        ("Technology", "Strong top-down confidence; middle management adoption lags behind.", False),
        ("Leadership", "C-suite vision is set but VP Ops sponsorship gap creates execution risk.", False),
        ("Consensus", "Tech & Culture are strengths; Process is the primary transformation lever.", True),
    ]
    kt_top = body_y + 30
    kt_h_each = (body_h - 30 - (len(kt_items) - 1) * 4) // len(kt_items)
    for i, (label, text, accent) in enumerate(kt_items):
        ky = kt_top + i * (kt_h_each + 4)
        if accent:
            box = add_rect(slide, f"kt-item-{i+1}-bg", kt_x, ky, kt_w, kt_h_each, BRAND_PRIMARY)
            label_color = RGBColor(0xC7, 0xB0, 0xE0)
            text_color = WHITE
            bar_color = BRAND_ACCENT
        else:
            box = add_rect(slide, f"kt-item-{i+1}-bg", kt_x, ky, kt_w, kt_h_each, CARD_BG)
            box.line.color.rgb = CARD_BORDER
            box.line.width = 9525
            label_color = BRAND_PRIMARY_MID
            text_color = TEXT_DARK
            bar_color = BRAND_ACCENT_SOFT
        # left bar
        add_rect(slide, f"kt-item-{i+1}-bar", kt_x, ky, 3, kt_h_each, bar_color)
        add_text(
            slide, f"kt-item-{i+1}-label", label,
            x_px=kt_x + 10, y_px=ky + 6, w_px=kt_w - 16, h_px=12,
            font_size_px=9, color=label_color, bold=True,
            uppercase=True, letter_spacing_px=1,
        )
        add_text(
            slide, f"kt-item-{i+1}-text", text,
            x_px=kt_x + 10, y_px=ky + 20, w_px=kt_w - 16, h_px=kt_h_each - 24,
            font_size_px=10, color=text_color,
        )

    add_footer(slide, page_num=289)
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parents[2] / "_renders" / "twins" / "289_interview-synthesis-grid.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
