"""
Translator — applies small adjustments to a PPTX twin.

When you change something in an HTML pattern (move a shape, recolor a card,
resize a brand-rule, change a font), the translator applies the equivalent
change to the matching PPTX twin. It edits the twin in place by default; pass
`out_path` to write to a new file.

This is NOT a CSS parser. The caller supplies an explicit adjustments dict —
which shape, which property, what new value. No layout inference, no flex/grid
handling. Pure "find shape by name, set this property" operations.

Supported adjustments per shape:

    {
        "shape-id": {
            "text":        "new text content",
            "fill":        "#RRGGBB",        # background fill
            "left":        64,               # x in CSS px
            "top":         290,              # y in CSS px
            "width":       368,              # w in CSS px
            "height":      200,              # h in CSS px
            "font_size":   16,               # CSS px (auto converted to pt)
            "font_color":  "#RRGGBB",
            "bold":        True,
            "italic":      False,
            "align":       "left" | "center" | "right",
        }
    }

All keys optional. Apply only what's provided. Shapes not found are skipped
silently and logged.

Typical workflow:
    1. Designer edits an HTML pattern (e.g. shrinks the title from 32px to 28px,
       moves the brand-rule from x=64 to x=80)
    2. Designer or Claude builds an adjustments dict describing those changes
    3. translator.apply(twin_path, adjustments) updates the PPTX twin in place

Smoke test (run from slide-builder/):
    python -m twins.translator
"""
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


TWINS_DIR = Path(__file__).resolve().parent.parent / "_renders" / "twins"

_PX_TO_EMU = 9525        # 1 CSS px @ 96dpi → EMU
_PX_TO_PT = 0.75         # 1 CSS px @ 96dpi → pt @ 72dpi


def _hex_to_rgb(hex_str: str) -> RGBColor:
    s = hex_str.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _find_shapes_by_name(slide, name: str) -> List[Any]:
    """Return all shapes on the slide whose .name == name (usually 0 or 1)."""
    return [shape for shape in slide.shapes if shape.name == name]


def _apply_text(shape, new_text: str) -> None:
    """Replace text content; preserve first-run font properties."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    src_p = tf.paragraphs[0]
    src_run = src_p.runs[0] if src_p.runs else None
    src_font = None
    if src_run:
        src_font = {
            "name": src_run.font.name,
            "size": src_run.font.size,
            "bold": src_run.font.bold,
            "italic": src_run.font.italic,
            "color_rgb": (src_run.font.color.rgb if src_run.font.color and src_run.font.color.type else None),
        }
    src_alignment = src_p.alignment
    tf.clear()
    p = tf.paragraphs[0]
    if src_alignment is not None:
        p.alignment = src_alignment
    run = p.add_run()
    run.text = new_text
    if src_font:
        f = run.font
        if src_font["name"]:
            f.name = src_font["name"]
        if src_font["size"]:
            f.size = src_font["size"]
        if src_font["bold"] is not None:
            f.bold = src_font["bold"]
        if src_font["italic"] is not None:
            f.italic = src_font["italic"]
        if src_font["color_rgb"] is not None:
            f.color.rgb = src_font["color_rgb"]


def _apply_fill(shape, hex_str: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(hex_str)


def _apply_position(shape, *, left=None, top=None, width=None, height=None) -> None:
    if left is not None:
        shape.left = Emu(int(left * _PX_TO_EMU))
    if top is not None:
        shape.top = Emu(int(top * _PX_TO_EMU))
    if width is not None:
        shape.width = Emu(int(width * _PX_TO_EMU))
    if height is not None:
        shape.height = Emu(int(height * _PX_TO_EMU))


def _apply_font(shape, *, font_size=None, font_color=None, bold=None,
                italic=None, align=None) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(align, p.alignment)
        for run in p.runs:
            f = run.font
            if font_size is not None:
                f.size = Pt(font_size * _PX_TO_PT)
            if font_color is not None:
                f.color.rgb = _hex_to_rgb(font_color)
            if bold is not None:
                f.bold = bold
            if italic is not None:
                f.italic = italic


def _apply_one(shape, adjustment: Dict[str, Any]) -> None:
    """Dispatch all property changes for a single shape."""
    if "text" in adjustment:
        _apply_text(shape, adjustment["text"])
    if "fill" in adjustment:
        _apply_fill(shape, adjustment["fill"])
    pos_keys = {k: adjustment[k] for k in ("left", "top", "width", "height") if k in adjustment}
    if pos_keys:
        _apply_position(shape, **pos_keys)
    font_keys = {k: adjustment[k] for k in ("font_size", "font_color", "bold", "italic", "align") if k in adjustment}
    if font_keys:
        _apply_font(shape, **font_keys)


def apply(twin_path: str, adjustments: Dict[str, Dict[str, Any]], *,
          out_path: Optional[str] = None, verbose: bool = True) -> Path:
    """Apply adjustments to a PPTX twin. Returns the output path.

    By default, writes back to twin_path in place. Pass out_path to write
    elsewhere.
    """
    twin_path = Path(twin_path)
    if not twin_path.exists():
        raise FileNotFoundError(f"Twin not found: {twin_path}")
    out_path = Path(out_path) if out_path else twin_path
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(twin_path))
    if not prs.slides:
        raise ValueError(f"Twin has no slides: {twin_path}")
    slide = prs.slides[0]

    skipped: List[str] = []
    applied: List[str] = []
    for shape_id, adj in adjustments.items():
        shapes = _find_shapes_by_name(slide, shape_id)
        if not shapes:
            skipped.append(shape_id)
            continue
        for shape in shapes:
            _apply_one(shape, adj)
        applied.append(shape_id)

    prs.save(str(out_path))
    if verbose:
        print(f"Applied {len(applied)} adjustments to {out_path}")
        if skipped:
            print(f"  Skipped {len(skipped)} unknown shape ids: {skipped}")
    return out_path


def apply_from_yaml(adjustments_yaml: str, twin_path: Optional[str] = None,
                    out_path: Optional[str] = None, verbose: bool = True) -> Path:
    """Load adjustments from a YAML file. The YAML can either be a flat dict
    (used with an explicit twin_path) or include a top-level `twin:` key
    pointing at the target PPTX.

    Example YAML:
        twin: ../_renders/twins/01_anchor-with-cards-icons.pptx
        out: ../_renders/twins/01_anchor_v2.pptx
        adjustments:
          title:
            font_size: 28
            font_color: "#1A1A2E"
          brand-rule:
            left: 80
            width: 80
          card-1-heading:
            text: "New card 1 heading"
            font_color: "#2D0A4E"
    """
    import yaml
    spec_path = Path(adjustments_yaml)
    if not spec_path.exists():
        raise FileNotFoundError(f"Adjustments YAML not found: {spec_path}")
    data = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    if not isinstance(data, dict):
        raise ValueError("Adjustments YAML must be a mapping")
    if "adjustments" in data:
        target = twin_path or data.get("twin")
        target_out = out_path or data.get("out") or target
        adjustments = data["adjustments"] or {}
    else:
        # Flat: the whole file is the adjustments dict; caller must pass twin_path
        if twin_path is None:
            raise ValueError("Flat adjustments YAML requires twin_path argument")
        target = twin_path
        target_out = out_path or twin_path
        adjustments = data
    if not target:
        raise ValueError("Could not resolve target twin path")
    # Resolve relative to YAML
    target = Path(target)
    if not target.is_absolute():
        target = (spec_path.parent / target).resolve()
    target_out = Path(target_out)
    if not target_out.is_absolute():
        target_out = (spec_path.parent / target_out).resolve()
    return apply(str(target), adjustments, out_path=str(target_out), verbose=verbose)


if __name__ == "__main__":
    # Smoke test: shrink the brand-rule on pattern 01's twin and recolor a card.
    out = TWINS_DIR / "_translator-test-01.pptx"
    import shutil
    shutil.copy(TWINS_DIR / "01_anchor-with-cards-icons.pptx", out)
    apply(
        str(out),
        {
            "brand-rule": {"left": 80, "width": 120, "fill": "#16A34A"},
            "card-1-heading": {"text": "Latency", "font_color": "#DC2626"},
            "title": {"font_size": 28},
        },
    )
