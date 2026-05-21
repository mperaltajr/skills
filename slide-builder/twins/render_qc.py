"""
Render-time QC for composed twin PPTX slides.

Runs after a slide is composed (overrides applied, master placeholders
stripped) and BEFORE LibreOffice renders the PNG. Inspects the resulting
shapes geometrically — no vision pass needed because the composer already
knows where everything is.

Three verdicts:
  - clean    : no issues
  - warning  : at least one warning issue; option still selectable
  - critical : at least one critical issue; option blocked from selection

Checks are intentionally cheap (just iterate slide.shapes and look at
left/top/width/height + text_frame.text). They catch:
  C2  shape extends below y=672 (invariant zone reserved for source/
      footnote/page-num)
  C4  body content outside the safe vertical zone (y=220 to y=630)
  C7  forbidden placeholder text leaked into a content shape
  C8  empty content shape (orphan: brief didn't provide override AND
      _blank_unmatched_content_shapes blanked it)

A heavier vision-based pass (line wrap detection, contrast checks, visible
overlap z-fight) is intentionally out of scope here — it requires an LLM
or OCR step that the orchestrator can opt into separately.

Usage:
    from twins.render_qc import check_composed_pptx
    verdict = check_composed_pptx(pptx_path)
    # → {"verdict": "warning", "issues": [{"severity": "warning", "msg": "..."}]}
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, List

# Slide dimensions in pixels (twin canvas is 1280x720 @ 1px = 9525 EMU).
_SLIDE_W = 1280
_SLIDE_H = 720

# Invariant bottom zone — only source / footnote-N / page-number shapes
# are allowed here.
_INVARIANT_Y = 672
_INVARIANT_ALLOW_RE = re.compile(r"^(source|footnote-\d+|page-number|page-num|footer.*)$", re.IGNORECASE)

# Body safe zone — content text should fit in [220, 630].
_BODY_TOP_MIN = 220
_BODY_BOTTOM_MAX = 630

# Forbidden placeholder text patterns (subset of brief_qc list, plus the
# specific PowerPoint master placeholder strings).
_FORBIDDEN_TEXT_RE = re.compile(
    r"\b(TBD|Lorem|\[Client Name\]|xxxx+|Click to edit|placeholder|"
    r"\[insert.*?\]|Master title style|Master subtitle style|"
    r"FedEx Proprietary|Customize with Department)\b",
    re.IGNORECASE,
)

# Shape-id patterns that bear content (mirrors the allowlist in composer.py).
_CONTENT_SHAPE_RE = re.compile(
    r"^(title|subtitle|eyebrow|headline|hero-statement|hero-context|"
    r"hero-attribution|key-question|anchor-statement|tagline|"
    r"cover-(deck-title|title|wordmark|tagline|subtitle|eyebrow|pre-label|"
    r"presenter|presented-name|presented-label|client-name|brand-name|"
    r"date|meta(-\d+(-label|-value)?)?)|"
    r"(card|panel|column|pillar|step|option|bucket|col)-\d+-"
    r"(heading|name|body|label|title|description|eyebrow)|"
    r"(before|after)-panel-(heading|body|label)|"
    r"(before|after)-(heading|body|label)|"
    r"metric-\d+-(label|value|delta)|"
    r"sub-ask-\d+-(label|body)|primary-ask-text|convergence|takeaway)$"
)

_EMU_PER_PX = 9525


def _emu_to_px(emu: int) -> int:
    return int(round((emu or 0) / _EMU_PER_PX))


def _shape_bounds(shape) -> Dict[str, int]:
    try:
        return {
            "left": _emu_to_px(shape.left),
            "top": _emu_to_px(shape.top),
            "width": _emu_to_px(shape.width),
            "height": _emu_to_px(shape.height),
        }
    except Exception:
        return {"left": 0, "top": 0, "width": 0, "height": 0}


def _is_content_shape(name: str) -> bool:
    return bool(_CONTENT_SHAPE_RE.match(name or ""))


def check_composed_slide(slide) -> Dict[str, Any]:
    """Inspect a composed slide. Returns {verdict, issues}."""
    issues: List[Dict[str, str]] = []

    for shape in slide.shapes:
        try:
            name = shape.name or ""
        except Exception:
            continue
        b = _shape_bounds(shape)
        bottom = b["top"] + b["height"]
        right = b["left"] + b["width"]

        text = ""
        if getattr(shape, "has_text_frame", False):
            try:
                text = shape.text_frame.text or ""
            except Exception:
                text = ""

        # C2 — invariant bottom zone reserved for source/footnote/page-num
        if bottom > _INVARIANT_Y and not _INVARIANT_ALLOW_RE.match(name):
            # Don't flag chrome backgrounds/rects without text
            if text.strip() or _is_content_shape(name):
                issues.append({
                    "severity": "warning",
                    "msg": f"shape '{name}' extends to y={bottom} (>672, invariant zone) — C2",
                })

        # C4 — body content outside the safe vertical zone
        if _is_content_shape(name) and text.strip():
            # Title-family shapes legitimately sit ABOVE the body zone
            is_title_family = re.match(
                r"^(title|subtitle|eyebrow|headline|cover-)", name
            )
            if not is_title_family:
                if b["top"] < _BODY_TOP_MIN - 20:
                    issues.append({
                        "severity": "warning",
                        "msg": f"shape '{name}' top y={b['top']} above safe zone (<{_BODY_TOP_MIN}) — C4",
                    })
                if bottom > _BODY_BOTTOM_MAX + 20:
                    issues.append({
                        "severity": "warning",
                        "msg": f"shape '{name}' bottom y={bottom} below safe zone (>{_BODY_BOTTOM_MAX}) — C4",
                    })

        # C7 — forbidden placeholder text leaked
        if text and _FORBIDDEN_TEXT_RE.search(text):
            issues.append({
                "severity": "critical",
                "msg": f"shape '{name}' contains placeholder text: {text[:80]!r} — C7",
            })

        # C8 — empty content shape (orphan: brief didn't provide override)
        if _is_content_shape(name) and getattr(shape, "has_text_frame", False):
            if text == "":
                issues.append({
                    "severity": "warning",
                    "msg": f"content shape '{name}' is empty — brief did not provide content — C8",
                })

        # Off-canvas
        if b["left"] < -8 or b["top"] < -8 or right > _SLIDE_W + 8 or bottom > _SLIDE_H + 8:
            issues.append({
                "severity": "warning",
                "msg": f"shape '{name}' off-canvas: left={b['left']} top={b['top']} right={right} bottom={bottom}",
            })

    verdict = "clean"
    if any(i["severity"] == "critical" for i in issues):
        verdict = "critical"
    elif any(i["severity"] == "warning" for i in issues):
        verdict = "warning"
    return {"verdict": verdict, "issues": issues}


def check_composed_pptx(pptx_path: str) -> Dict[str, Any]:
    """Open the composed PPTX and run check_composed_slide on its first slide."""
    from pptx import Presentation
    p = Path(pptx_path)
    if not p.exists():
        return {"verdict": "critical", "issues": [{"severity": "critical", "msg": f"PPTX missing: {p}"}]}
    try:
        prs = Presentation(str(p))
    except Exception as e:
        return {"verdict": "critical", "issues": [{"severity": "critical", "msg": f"open failed: {e}"}]}
    if not prs.slides:
        return {"verdict": "critical", "issues": [{"severity": "critical", "msg": "no slides in PPTX"}]}
    return check_composed_slide(prs.slides[0])


if __name__ == "__main__":
    # Smoke test on a known twin
    test = Path(__file__).resolve().parent.parent / "_renders" / "twins" / "01_anchor-with-cards-icons.pptx"
    if test.exists():
        import json
        print(json.dumps(check_composed_pptx(str(test)), indent=2))
    else:
        print(f"Test pptx not found: {test}")
