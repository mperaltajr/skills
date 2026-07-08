"""Slide-graft chrome helpers.

Three small functions called by `scripts/finalize_deck.py::graft_and_theme`
and `scripts/compile_picks.py` to prepare a client-template `Presentation`
for slide grafting:

  - `_clear_existing_slides(prs)` — drop any sample slides shipped with the
    template + strip stale section groupings from the slide-panel sidebar.
  - `_find_blank_layout(prs)` — locate a 'Blank' slide layout across all
    masters (client templates often put the blank layout on a non-zero
    master; FedEx puts it on master 11).
  - `_strip_layout_placeholders(slide)` — remove inherited layout
    placeholders AND set `showMasterSp="0"` so master shapes don't bleed
    through onto the grafted content.

Historical note (Path D, 2026-05-26)
------------------------------------
This file used to host the legacy chassis-vocabulary skill's full
`compose_deck()` pipeline plus `_role_map_for`, `_blank_unmatched_via_role_map`,
`_PATTERN_LAYOUT_HINTS`, `_apply_overrides`, and ~25 supporting helpers — all
dependent on `twins.selector.load_catalog` and `twins.overrides_resolver`,
modules that no longer exist. v0.1 never reaches those code paths; v0.1 audit
(2026-05-26) flagged them as dead-broken imports (T1.4 in
`_decisions/v0.1-audit-handover-2026-05-26.md`). Deleted in this pass to
prevent future changes reaching them from exploding with ImportError.
The legacy chassis-vocabulary composer was archived and removed from disk
during Path D consolidation; if anyone needs to reference the full v1
composer they will need to recover it from version control history rather
than expect a live archive path.
"""
from __future__ import annotations

import re

from pptx.oxml.ns import qn
from pptx.util import Pt


class TemplatePlaceholderEmptyError(RuntimeError):
    """Raised when a target placeholder has no paragraph element to write
    into — i.e., the layout's placeholder XML is malformed or stripped.

    Gate B.3 (2026-06-08, SLIDE_LAB_FEEDBACK_LOG Deferred #3): the silent
    no-op in _write was the same bug class as TitleDropError but one layer
    deeper. The placeholder exists, the populate step found it, but the
    text write silently returns because tf.paragraphs is empty. Net effect:
    placeholder is found, write is "successful," but no text lands.

    Per feedback_sidecar_fallback_must_be_loud — silent placeholder no-op
    is invisible data loss.
    """


def _find_named_layout(prs, layout_name: str):
    """Find a slide layout by EXACT name across all masters.

    Used by finalize_deck.py to graft each option onto the layout the brief
    declared (chrome inheritance), instead of always falling back to blank.
    Returns None when the name doesn't match — caller falls back to
    _find_blank_layout.
    """
    if not layout_name:
        return None
    target = layout_name.strip()
    if not target:
        return None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            n = (layout.name or "").strip()
            if n == target:
                return layout
    return None


def _find_blank_layout(prs):
    """Find a blank slide layout across all masters in `prs`.

    Looks for a layout whose name (case-insensitive, after stripping
    whitespace and leading numeric prefixes like '1_') is exactly 'blank'.
    Scans every slide master because client templates often put their
    blank layout on a non-zero master (FedEx puts it on master 11).
    Falls back to the first master's first layout if no 'blank' is found.
    """
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            n = layout.name.strip().lower()
            n = re.sub(r"^[\d_]+", "", n)  # strip "1_", "12__", etc.
            if n == "blank":
                return layout
    return prs.slide_masters[0].slide_layouts[0]


def _strip_layout_placeholders(slide, *, keep_master_shapes: bool = False) -> int:
    """Remove layout-level placeholders that were copied onto the slide
    when `add_slide(layout)` was called (Click to add title, sample bullets).

    By default ALSO sets ``<p:sld showMasterSp="0">`` to hide master-level
    shapes (Click to edit Master title, footer placeholders) — the v1
    behavior for templates whose masters carried only placeholder leakage.

    When ``keep_master_shapes=True``: only the slide-level shape strip runs;
    master shapes (FedEx purple bars, header rules, OTC top/bottom bands)
    remain visible. Use this for bespoke layouts on templates where the
    master IS the brand chrome — i.e., brand.yml has
    ``strip_master_backgrounds: false``.

    Background, theme colors, and theme fonts inherit regardless of
    `showMasterSp` — those are not "shapes" in OOXML's sense.

    Returns the number of slide-level shapes removed.
    """
    spTree = slide.shapes._spTree
    removed = 0
    for shp in list(slide.shapes):
        try:
            spTree.remove(shp.element)
            removed += 1
        except (ValueError, AttributeError):
            pass

    if not keep_master_shapes:
        try:
            slide.element.set("showMasterSp", "0")
        except Exception:
            pass

    return removed


def _populate_layout_placeholders(slide, *, title=None, subtitle=None,
                                   footer=None, page_num=None,
                                   title_idx=None, subtitle_idx=None,
                                   title_font_pt=None):
    """Write text into the slide's inherited layout placeholders.

    For body-canonical grafts: instead of stripping the layout's title/footer/
    page-number placeholders and redrawing them as free-floating textboxes,
    keep them in place and write text into the FIRST placeholder of each
    matching type. Returns a dict {role: bool} flagging which roles were
    found-and-populated.

    title_idx / subtitle_idx (v2.2, SLIDE_LAB_FEEDBACK_LOG #2 follow-up):
    when chrome.yml registers a specific placeholder idx for title or
    subtitle (e.g., FedEx layouts where the "subtitle" slot is actually a
    BODY-type placeholder at idx=10, not a SUBTITLE-type placeholder),
    pass that idx in. The function will then match by idx FIRST and fall
    back to type-based matching only if idx-based lookup fails. This
    closes the silent-drop bug where BODY-type slots that serve as
    subtitles were skipped by strict type matching.

    Caller is responsible for choosing whether to invoke this (body-canonical)
    vs. _strip_layout_placeholders (bespoke / cover).
    """
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    # python-pptx PP_PLACEHOLDER values: TITLE=1, CENTER_TITLE=3, SUBTITLE=4,
    # FOOTER=15, SLIDE_NUMBER=13. (Earlier code used 13 for CENTER_TITLE and 12
    # for SLIDE_NUMBER — both wrong: 13 is the slide number, so page numbers
    # never populated and center-title covers never matched by type.)
    TITLE_TYPES = {1, 3}        # TITLE, CENTER_TITLE
    SUBTITLE_TYPES = {4}        # SUBTITLE
    FOOTER_TYPES = {15}         # FOOTER
    PAGE_NUM_TYPES = {13}       # SLIDE_NUMBER

    found = {"title": False, "subtitle": False,
             "footer": False, "page_number": False}

    def _write(ph, text, *, anchor_bottom: bool = False, font_pt=None):
        tf = ph.text_frame
        # Replace existing paragraph text in the first paragraph; clear extras.
        if not tf.paragraphs:
            # Gate B.3 (2026-06-08, SLIDE_LAB_FEEDBACK_LOG Deferred #3):
            # silent no-op was the v1 bug class. If the placeholder is
            # well-formed enough that we found it, it must have at least
            # one paragraph element to receive text. Empty text_frame means
            # malformed template — fail loudly with recovery instructions.
            try:
                _idx = ph.placeholder_format.idx
            except Exception:
                _idx = "?"
            try:
                _name = ph.name
            except Exception:
                _name = "<unnamed>"
            raise TemplatePlaceholderEmptyError(
                f"A placeholder in your template (named {_name!r}) is empty "
                f"in a way that prevents Slide Lab from writing text into it. "
                f"We tried to place: {(text or '')[:60]!r}\n\n"
                f"  Why this happens: every placeholder needs at least one "
                f"empty text line for content to land on. Some templates lose "
                f"that line after being edited and saved (often when a "
                f"placeholder is deleted and re-added in PowerPoint).\n\n"
                f"  What to do (~30 seconds in PowerPoint):\n"
                f"    1. Open the template in PowerPoint.\n"
                f"    2. View → Slide Master, find the layout this slide uses.\n"
                f"    3. Click inside the placeholder named {_name!r}, type any "
                f"character, then delete it. (This forces PowerPoint to "
                f"restore the empty text line.)\n"
                f"    4. Save the template (Ctrl+S), close the master view.\n"
                f"    5. Re-register the template:\n"
                f"         py -3 scripts/register_template.py propose <your-template.pptx>"
            )
        first = tf.paragraphs[0]
        # Clear existing runs
        for r in list(first.runs):
            r._r.getparent().remove(r._r)
        # Clear any auto-field (e.g. a slide-number <a:fld>) — we write literal
        # text so it renders reliably in LibreOffice + PowerPoint (auto-number
        # fields don't render under LibreOffice headless).
        for fld in first._p.findall(qn("a:fld")):
            first._p.remove(fld)
        run = first.add_run()
        run.text = str(text) if text is not None else ""
        # Controlled size (e.g. a consistent 28pt title) when the caller supplies
        # one; otherwise leave the run size unset so the template's placeholder
        # cascade governs.
        if font_pt:
            try:
                run.font.size = Pt(font_pt)
            except Exception:
                pass
        # Defect 4 fix (2026-06-15 CDIO QBR feedback): enforce
        # feedback_title_bottom_anchor on placeholder-populated titles too.
        # Without this, the inherited title placeholder uses the layout's
        # default vertical anchor (usually TOP), so 2-line titles grow DOWN
        # into the body content zone. With anchor_bottom=True, the title's
        # bottom-y stays fixed and longer titles grow UPWARD into the chrome
        # zone (which is reserved space — never displaces subtitle or body).
        if anchor_bottom:
            try:
                tf.vertical_anchor = MSO_ANCHOR.BOTTOM
            except Exception:
                pass

    # First pass: idx-based matching for any role that supplied an idx.
    # This honors chrome.yml's registered placeholder ids (e.g., FedEx
    # template's subtitle is at idx=10 as a BODY-type placeholder — strict
    # type matching would silently miss it).
    if title is not None and title_idx is not None:
        for ph in list(slide.placeholders):
            try:
                if ph.placeholder_format.idx == title_idx:
                    _write(ph, title, anchor_bottom=True, font_pt=title_font_pt)
                    found["title"] = True
                    break
            except Exception:
                continue
    if subtitle is not None and subtitle_idx is not None:
        for ph in list(slide.placeholders):
            try:
                if ph.placeholder_format.idx == subtitle_idx:
                    _write(ph, subtitle)
                    found["subtitle"] = True
                    break
            except Exception:
                continue

    # Second pass: type-based matching for roles not yet placed.
    for ph in list(slide.placeholders):
        try:
            t = int(ph.placeholder_format.type)
        except Exception:
            continue
        if title is not None and not found["title"] and t in TITLE_TYPES:
            # Title gets bottom-anchor invariant per feedback_title_bottom_anchor.
            _write(ph, title, anchor_bottom=True, font_pt=title_font_pt)
            found["title"] = True
            continue
        if subtitle is not None and not found["subtitle"] and t in SUBTITLE_TYPES:
            _write(ph, subtitle)
            found["subtitle"] = True
            continue
        if footer is not None and not found["footer"] and t in FOOTER_TYPES:
            _write(ph, footer)
            found["footer"] = True
            continue
        if page_num is not None and not found["page_number"] and t in PAGE_NUM_TYPES:
            _write(ph, page_num)
            found["page_number"] = True
            continue
    return found


def _compute_body_zone(layout):
    """Return (body_top_emu, body_bottom_emu) for the layout's drawable zone.

    Top edge is the bottom of the layout's title placeholder (if present),
    falling back to a conservative canonical top. Bottom edge is the top of
    the footer placeholder (if present), falling back to a conservative
    canonical bottom.

    Both edges are in EMU. The caller decides whether to use them
    (body-canonical grafts).
    """
    # 1280x720 px at 96 DPI = 12192000 EMU x 6858000 EMU
    title_bottom = None
    footer_top = None
    for ph in layout.placeholders:
        try:
            t = int(ph.placeholder_format.type)
        except Exception:
            continue
        if t in (1, 3) and title_bottom is None:  # TITLE, CENTER_TITLE
            try:
                title_bottom = int(ph.top) + int(ph.height)
            except Exception:
                pass
        if t == 15 and footer_top is None:
            try:
                footer_top = int(ph.top)
            except Exception:
                pass
    # Canonical fallbacks (in EMU). Constants live in _chrome_schema so the
    # single-source-for-geometry contract holds.
    import sys as _sys
    from pathlib import Path as _Path
    _scripts = _Path(__file__).resolve().parent.parent / "scripts"
    if str(_scripts) not in _sys.path:
        _sys.path.insert(0, str(_scripts))
    from _chrome_schema import CANONICAL_BODY_TOP_Y, CANONICAL_BODY_BOTTOM_Y
    if title_bottom is None:
        title_bottom = CANONICAL_BODY_TOP_Y * 9525
    if footer_top is None:
        footer_top = CANONICAL_BODY_BOTTOM_Y * 9525
    return title_bottom, footer_top


def _insert_dark_overlay(slide, prs, body_top_emu, body_bottom_emu, hex_color):
    """Insert a solid-fill rectangle covering the body zone as the FIRST
    slide-level shape (z-order: above layout chrome, below all subsequent
    slide content).

    hex_color is a 6-char hex string ('0A1A2E', no leading #). Raises
    ValueError if malformed.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    s = (hex_color or "").lstrip("#")
    if len(s) != 6:
        raise ValueError(f"body_overlay_hex must be 6-char hex; got {hex_color!r}")
    try:
        r = int(s[0:2], 16); g = int(s[2:4], 16); b = int(s[4:6], 16)
    except ValueError:
        raise ValueError(f"body_overlay_hex has non-hex chars: {hex_color!r}")

    width = int(prs.slide_width)
    height = int(body_bottom_emu) - int(body_top_emu)
    if height <= 0:
        return None
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, int(body_top_emu), width, height,
    )
    shape.name = "body-overlay"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()

    # Move overlay to position 0 in spTree so subsequent content draws on top.
    # add_shape appends at the end; bring it to the front of slide-level shapes
    # (after the required <p:nvGrpSpPr> + <p:grpSpPr> heads). spTree children
    # order: nvGrpSpPr, grpSpPr, then shape elements; insert overlay just
    # after grpSpPr (= index 2).
    sp_tree = slide.shapes._spTree
    elem = shape.element
    sp_tree.remove(elem)
    # Find first shape-child index (skip nvGrpSpPr + grpSpPr at positions 0,1)
    sp_tree.insert(2, elem)
    return shape


def clone_missing_chrome_placeholders(slide, layout):
    """Deepcopy the SLIDE_NUMBER placeholder from `layout` onto `slide` when it
    isn't already present.

    python-pptx's ``add_slide`` clones only the title/body placeholders — it
    skips SLIDE_NUMBER / FOOTER / DATE (treated as inherited chrome). But those
    aren't reliably rendered from the master by LibreOffice headless, so the
    page number silently vanished on every body-canonical build. Cloning the
    real placeholder here lets _populate_layout_placeholders write a literal
    page number into it (and compile_picks re-stamp it to the compile position).
    Returns the number of placeholders cloned.
    """
    from copy import deepcopy
    have = set()
    for ph in slide.placeholders:
        try:
            have.add(int(ph.placeholder_format.idx))
        except Exception:
            continue
    sp_tree = slide.shapes._spTree
    cloned = 0
    for ph in layout.placeholders:
        try:
            t = int(ph.placeholder_format.type)
            idx = int(ph.placeholder_format.idx)
        except Exception:
            continue
        if t == 13 and idx not in have:  # SLIDE_NUMBER
            sp_tree.append(deepcopy(ph._element))
            cloned += 1
    return cloned


def remove_empty_placeholders(slide) -> int:
    """Remove inherited placeholders left empty after population.

    Body-canonical grafts draw the slide's body as free-floating shapes, so the
    layout's content placeholder stays empty — and PowerPoint shows its
    'Click to add text' prompt in edit mode (harmless in slideshow/export, but it
    looks unfinished when the user opens the deck). Populated placeholders
    (title, subtitle, page number) are non-empty and kept. Returns the count
    removed.
    """
    removed = 0
    for ph in list(slide.placeholders):
        try:
            txt = (ph.text_frame.text or "").strip() if ph.has_text_frame else ""
        except Exception:
            txt = ""
        if not txt:
            parent = ph._element.getparent()
            if parent is not None:
                parent.remove(ph._element)
                removed += 1
    return removed


def _clear_existing_slides(prs):
    """Remove any slides already present in the presentation (client
    templates often ship with sample slides we don't want) AND remove
    section groupings from the slide-panel sidebar.
    """
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)

    # Strip <p14:sectionLst> from the presentation extLst so the slide
    # panel doesn't show stale section groupings ('Brand Elements', 'Begin
    # Building', etc.) from the original template. Sections live under a
    # <p:ext> whose child is <p14:sectionLst>, where p14 is the 2010
    # PowerPoint namespace. Earlier versions of this code used the 2012
    # namespace by mistake and the strip never fired — orphaned section
    # groups bled through grafted decks. Check BOTH known namespaces to
    # survive future schema variants.
    pres_el = prs.part._element
    _SECTION_LST_TAGS = (
        "{http://schemas.microsoft.com/office/powerpoint/2010/main}sectionLst",
        "{http://schemas.microsoft.com/office/powerpoint/2012/main}sectionLst",
    )
    for ext_lst in pres_el.findall(qn("p:extLst")):
        for ext in list(ext_lst):
            for tag in _SECTION_LST_TAGS:
                if ext.find(".//" + tag) is not None:
                    ext_lst.remove(ext)
                    break
