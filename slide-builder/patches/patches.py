"""
patches.py — Named, reusable python-pptx operations for Slide Lab.

The skill calls these by name rather than writing ad-hoc python-pptx code.
Every patch is tested and predictable. Patches are applied AFTER token fill.

Public API
----------
Template / slide construction (use these in every build script):
  get_blank_layout(prs)                                    -> layout   [scans ALL masters]
  get_named_layout(prs, layout_name)                       -> layout   [scans ALL masters]
  catalog_layouts(prs)                                     -> list of dicts
  clear_slides(prs)                                        -> None     [+ strips section breaks]
  new_slide_from_template(template_path, layout_name)      -> (prs, slide)
  skeleton_on_template(skel_name, template_path, out_path, layout_name) -> Path
  skeleton_from_pptx(skel_pptx_path, template_path, out_path, layout_name) -> Path

QC guards:
  check_no_bg_masking(prs)                                 -> list of (slide_idx, shape_name) issues

Shape lookup:
  find_shape_by_token(slide, token_name)          -> shape or None
  find_shape_by_name(slide, shape_name)           -> shape or None

Content:
  fill_tokens(slide, token_map)                   -> list of filled token names
  convert_to_bullets(shape, items, bullet_char)   -> None
  reposition_zone(shape, x_in, y_in)             -> None
  resize_zone(shape, width_in, height_in)         -> None
  add_hero_stat(slide, stat, label, x, y, w, h, color) -> shape
  set_shape_fill(shape, hex_color)               -> None
  set_rag_status(shape, status)                   -> None
  add_table_row(table_shape, row_data)            -> None
  add_comparison_table(slide, left_heading, right_heading, rows, ...) -> shape
  insert_icon(slide, icon_name, x_in, y_in, size_in, accent_hex) -> bool

Trigger catalogue: see CATALOGUE.md
"""

from __future__ import annotations

import sys
import os
import pathlib
from typing import Optional

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_PATCHES_DIR = pathlib.Path(__file__).resolve().parent
_SKILL_DIR   = _PATCHES_DIR.parent
_SCRIPTS_DIR = _SKILL_DIR / "scripts"

# Add scripts to path so we can import icon_helper
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

# RAG status colors
RAG_COLORS = {
    "green":  RGBColor(0x00, 0x70, 0x00),
    "g":      RGBColor(0x00, 0x70, 0x00),
    "amber":  RGBColor(0xFF, 0xA5, 0x00),
    "yellow": RGBColor(0xFF, 0xA5, 0x00),
    "a":      RGBColor(0xFF, 0xA5, 0x00),
    "red":    RGBColor(0xC0, 0x00, 0x00),
    "r":      RGBColor(0xC0, 0x00, 0x00),
}

DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
PML = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------
#
# CRITICAL: python-pptx's `prs.slide_layouts` only exposes layouts on slide
# master 0. Corporate templates routinely have multiple masters — the FedEx
# "Moving Forward" template has 12 masters and 178 layouts; its actual "1_Blank"
# layout is at prs.slide_masters[11].slide_layouts[0], invisible to the shortcut.
#
# Every helper below iterates ALL masters. Never use `prs.slide_layouts` for
# selection. See feedback-2026-05-18-template-layout.md for the failure mode
# this prevents (cover-photo bleed-through on every content slide).


def _iter_all_layouts(prs):
    """Yield (master_idx, layout_idx, layout) across every slide master."""
    for mi, master in enumerate(prs.slide_masters):
        for li, layout in enumerate(master.slide_layouts):
            yield (mi, li, layout)


def _has_picture(shapes_iter):
    """True if any shape in the iterable is a PICTURE or a picture placeholder."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes_iter:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
        except Exception:
            pass
        # Picture placeholder (MSO_PLACEHOLDER.PICTURE == 18)
        try:
            ph = shape.placeholder_format
            if ph is not None and ph.type == 18:
                return True
        except Exception:
            pass
    return False


def _layout_or_master_has_picture(layout):
    """True if a layout OR its master carries decorative picture content.

    FedEx puts the cover photo on the MASTER, not on individual layouts —
    so checking layout.shapes alone misses it. Walk the master too.
    """
    if _has_picture(layout.shapes):
        return True
    try:
        return _has_picture(layout.slide_master.shapes)
    except Exception:
        return False


def catalog_layouts(prs) -> list:
    """Return a list of dicts describing every layout across every master.

    Each entry: {master_idx, layout_idx, name, shape_count, has_picture}.
    Use this when picking a layout to assess minimality and picture presence.
    """
    out = []
    for mi, li, layout in _iter_all_layouts(prs):
        out.append({
            "master_idx": mi,
            "layout_idx": li,
            "name": layout.name,
            "shape_count": len(layout.shapes),
            "has_picture": _layout_or_master_has_picture(layout),
        })
    return out


def get_blank_layout(prs):
    """Find the cleanest blank layout across ALL slide masters.

    Selection order:
      1. Layout named exactly 'Blank' (case-insensitive) with no picture
      2. Layout whose name contains 'blank' with no picture
      3. Layout with zero shapes AND no master/layout picture
      4. Layout with fewest shapes AND no master/layout picture
      5. Fallback: prs.slide_layouts[0] (with loud warning — likely wrong)
    """
    candidates = list(_iter_all_layouts(prs))
    # Pass 1: exact "Blank" name, no picture
    for mi, li, layout in candidates:
        if layout.name.strip().lower() == "blank" and not _layout_or_master_has_picture(layout):
            return layout
    # Pass 2: substring "blank", no picture
    for mi, li, layout in candidates:
        if "blank" in layout.name.strip().lower() and not _layout_or_master_has_picture(layout):
            return layout
    # Pass 3: zero shapes AND no picture anywhere upstream
    for mi, li, layout in candidates:
        if len(layout.shapes) == 0 and not _layout_or_master_has_picture(layout):
            return layout
    # Pass 4: fewest shapes AND no picture anywhere upstream
    no_pic = [(len(l.shapes), mi, li, l) for mi, li, l in candidates if not _layout_or_master_has_picture(l)]
    if no_pic:
        no_pic.sort()
        return no_pic[0][3]
    print("  [patch] WARNING: no clean blank layout found across any master — "
          "falling back to slide_layouts[0]. Content slides will likely inherit "
          "template imagery.")
    return prs.slide_layouts[0]


def get_named_layout(prs, layout_name: str):
    """Return a slide layout by name across ALL slide masters.

    Tries (in order):
      1. Case-insensitive exact match on layout name
      2. Case-insensitive substring match
      3. Falls back to get_blank_layout() if not found

    Iterates every master so corporate templates with multi-master designs
    (FedEx 'Moving Forward', many BCG/MBB decks) are handled correctly.
    """
    name_lower = layout_name.strip().lower()
    candidates = list(_iter_all_layouts(prs))
    # Pass 1: exact name match
    for mi, li, layout in candidates:
        if layout.name.strip().lower() == name_lower:
            return layout
    # Pass 2: substring match
    for mi, li, layout in candidates:
        if name_lower in layout.name.strip().lower():
            return layout
    print(f"  [patch] WARNING: layout '{layout_name}' not found across {len(prs.slide_masters)} masters — falling back to blank.")
    return get_blank_layout(prs)


_R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def _strip_section_breaks(prs) -> int:
    """Remove inherited <p:sectionLst> entries from presentation.xml.

    When you clear a template's slides, the section breaks named in the
    template ("Title Slides", "Agenda", "Content", "Closing") remain in the
    Sections navigation pane of the final deck, grouping new slides under
    labels that no longer make sense. Strip them so the new deck starts clean.

    Returns the number of section breaks removed.
    """
    P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    prs_elem = prs.part._element  # <p:presentation>
    ext_lst = prs_elem.find(f"{{{PML}}}extLst")
    if ext_lst is None:
        return 0
    removed = 0
    for ext in list(ext_lst):
        section_lst = ext.find(f"{{{P14_NS}}}sectionLst")
        if section_lst is not None:
            ext.remove(section_lst)
            removed += 1
            # If the <p:ext> is now empty, drop it too
            if len(ext) == 0:
                ext_lst.remove(ext)
    return removed


def clear_slides(prs) -> None:
    """Remove every existing slide AND stale section breaks from a presentation.

    Section breaks are inherited from the source template's presentation.xml
    and survive slide deletion unless explicitly stripped — see
    feedback-2026-05-18-template-layout.md.
    """
    xml_slides = prs.slides._sldIdLst
    prs_part = prs.part
    for sld_id in list(xml_slides):
        rId = sld_id.get(_R_ID)
        slide_part = prs_part.related_part(rId)
        prs_part.drop_rel(rId)
        xml_slides.remove(sld_id)
        try:
            del prs_part.package._parts[slide_part.partname]
        except (KeyError, AttributeError):
            pass
    _strip_section_breaks(prs)


def check_no_bg_masking(prs) -> list:
    """Return a list of layout-failure issues — full-slide white rectangles
    stamped on top of (or beneath) the slide to mask unwanted template imagery.

    This pattern is ALWAYS a layout-selection bug. The correct fix is to pick
    a different layout (likely on a different master), not to paper over the
    template artifact with a rectangle. See feedback-2026-05-18-template-layout.md.

    Each issue is (slide_index, shape_name, reason).
    Slide_index is 1-based to match how users count slides.
    Empty list means clean.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    issues = []
    slide_w = prs.slide_width.emu
    slide_h = prs.slide_height.emu
    SLIDE_AREA = slide_w * slide_h
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            try:
                name = (shape.name or "").lower()
            except Exception:
                continue
            # Explicit named-mask shapes
            if "bg_white" in name or "background_mask" in name or "white_bg_rect" in name:
                issues.append((idx, shape.name, "named like a template-masking rectangle"))
                continue
            # Full-slide opaque rectangles (regardless of name) are suspicious
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                    continue
                # area within 5% of full slide?
                area = (shape.width.emu) * (shape.height.emu)
                if area > 0.95 * SLIDE_AREA:
                    issues.append((idx, shape.name, "full-slide-sized AUTO_SHAPE — likely a template artifact mask"))
            except Exception:
                continue
    return issues


def new_slide_from_template(template_path, layout_name: str = None) -> tuple:
    """Open a client template, clear its slides, return a fresh slide.

    Returns (prs, slide). The slide uses the blank layout by default so the
    master (footer, fonts, branding) carries through. Placeholder shapes from
    the layout are removed — the slide is ready for direct python-pptx content.

    Pass layout_name to select a named corporate layout (e.g. 'Cover', 'Title Slide')
    instead of blank. This is required for cover slides and any slide type that
    inherits significant branded decorations from a specific master layout.

    Use this for every direct-build slide (charts, tables, custom layouts).
    """
    from pptx import Presentation
    prs = Presentation(str(template_path))
    clear_slides(prs)
    layout = get_named_layout(prs, layout_name) if layout_name else get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            sp_tree.remove(child)
    return prs, slide


def skeleton_on_template(skeleton_name: str, template_path, out_path,
                         layout_name: str = None) -> pathlib.Path:
    """Graft skeleton shapes onto a slide from the client template.

    Opens the skeleton PPTX from the skill's skeletons library, deep-copies
    its shapes into a fresh slide built from template_path, saves to out_path.

    The result is a proper client-template PPTX — same master, fonts, and
    package relationships as the original template. It drops into any larger
    deck without conflicts.

    layout_name: optional named layout to use instead of blank (e.g. 'Cover',
    'Title Slide'). Pass this for cover slides and any skeleton that should
    inherit a specific master layout's decorations.

    Use this for every skeleton-based slide. Never copy the skeleton PPTX
    directly — that produces a slide on the skeleton's own master, not the
    client's.
    """
    import copy as _copy
    from pptx import Presentation
    skel_path = _SKILL_DIR / "skeletons" / skeleton_name / f"{skeleton_name}.pptx"
    skel_prs = Presentation(str(skel_path))
    skel_slide = skel_prs.slides[0]
    prs, slide = new_slide_from_template(template_path, layout_name=layout_name)
    sp_tree = slide.shapes._spTree
    for shape in skel_slide.shapes:
        sp_tree.append(_copy.deepcopy(shape._element))
    out_path = pathlib.Path(out_path)
    prs.save(str(out_path))
    return out_path


def skeleton_from_pptx(skel_pptx_path, template_path, out_path,
                       layout_name: str = None) -> pathlib.Path:
    """Graft shapes from an explicit skeleton PPTX path onto a client template slide.

    Same as skeleton_on_template but takes a full path to the skeleton PPTX
    instead of looking it up by name. Use for skeleton families stored flat
    (e.g. chart-with-takeaway variants: skeletons/chart-with-takeaway/<name>.pptx).

    layout_name: optional named layout (same as skeleton_on_template).
    """
    import copy as _copy
    from pptx import Presentation
    skel_prs = Presentation(str(skel_pptx_path))
    skel_slide = skel_prs.slides[0]
    prs, slide = new_slide_from_template(template_path, layout_name=layout_name)
    sp_tree = slide.shapes._spTree
    for shape in skel_slide.shapes:
        sp_tree.append(_copy.deepcopy(shape._element))
    out_path = pathlib.Path(out_path)
    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# Shape lookup
# ---------------------------------------------------------------------------

def find_shape_by_token(slide, token_name: str):
    """
    Return the shape whose text contains {{token_name}}.
    Use before fill to locate a shape; after fill the token is gone.
    Case-insensitive on the token name.
    """
    target = "{{" + token_name.upper() + "}}"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if target in shape.text_frame.text.upper():
            return shape
    return None


def find_shape_by_name(slide, shape_name: str):
    """Return the shape with the given name (exact match)."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    return None


# ---------------------------------------------------------------------------
# Token fill
# ---------------------------------------------------------------------------

def fill_tokens(slide, token_map: dict) -> list:
    """
    Replace {{KEY}} with value for every entry in token_map across all shapes.
    Returns list of token names that were successfully replaced.

    token_map keys are bare names (no braces): {"ACTION_TITLE": "My title", ...}

    Empty-value rule: if a token is replaced with "" and the resulting text has no
    word characters, the paragraph is removed. If a shape's ENTIRE content becomes
    empty (e.g. a single-paragraph bullet textbox like "• {{BULLET_4}}"), the whole
    shape is removed from the slide. This prevents orphaned "•" bullets and "1. "
    footnote stubs when optional slots are left blank.
    """
    import re as _re
    from pptx.oxml.ns import qn
    filled = []
    shapes_to_remove = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        shape_had_empty_token = False
        paras_to_remove = []

        for para in tf.paragraphs:
            para_had_empty_token = False
            for run in para.runs:
                for key, val in token_map.items():
                    token = "{{" + key + "}}"
                    if token in run.text:
                        if str(val) == "":
                            run.text = ""  # clear whole run so boilerplate prefix ("1. ", "• ") disappears
                            para_had_empty_token = True
                            shape_had_empty_token = True
                        else:
                            run.text = run.text.replace(token, str(val))
                        if key not in filled:
                            filled.append(key)
            if para_had_empty_token:
                combined = "".join(r.text for r in para.runs)
                if not _re.search(r'\w', combined):
                    paras_to_remove.append(para._p)

        txBody = tf._txBody
        for p_el in paras_to_remove:
            if len(txBody.findall(qn("a:p"))) > 1:
                txBody.remove(p_el)

        # If whole shape became empty after fills, remove it entirely
        if shape_had_empty_token and not _re.search(r'\w', tf.text):
            shapes_to_remove.append(shape._element)

    for sp_el in shapes_to_remove:
        parent = sp_el.getparent()
        if parent is not None:
            parent.remove(sp_el)

    return filled


# ---------------------------------------------------------------------------
# Patch: convert_to_bullets
# ---------------------------------------------------------------------------

def convert_to_bullets(shape, items: list, bullet_char: str = "•") -> None:
    """
    Replace all text in shape with a bulleted list of items.
    Preserves the font size and color of the first run in the shape.
    bullet_char defaults to "•" (U+2022).

    Trigger: brief has 3+ discrete items instead of continuous prose.
    """
    if not shape.has_text_frame or not items:
        return

    tf = shape.text_frame
    tf.word_wrap = True

    # Capture formatting from first existing run
    first_run = None
    for para in tf.paragraphs:
        if para.runs:
            first_run = para.runs[0]
            break

    font_size = first_run.font.size if first_run else Pt(11)
    font_bold = first_run.font.bold if first_run else False
    try:
        font_color = first_run.font.color.rgb if first_run else None
    except Exception:
        font_color = None

    # Clear all paragraphs from the XML directly
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    # Add one paragraph per bullet item
    for i, item in enumerate(items):
        p_elem = etree.SubElement(txBody, qn("a:p"))
        r_elem = etree.SubElement(p_elem, qn("a:r"))

        # Apply run properties if we have them
        rPr = etree.SubElement(r_elem, qn("a:rPr"))
        rPr.set("lang", "en-US")
        rPr.set("dirty", "0")
        if font_size:
            rPr.set("sz", str(int(font_size.pt * 100)))
        if font_bold:
            rPr.set("b", "1")
        if font_color:
            solidFill = etree.SubElement(rPr, qn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgbClr.set("val", str(font_color))

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = f"{bullet_char}  {item}"

        # Move rPr before t (PPTX requires rPr before t in a:r)
        r_elem.remove(rPr)
        r_elem.insert(0, rPr)


# ---------------------------------------------------------------------------
# Patch: reposition_zone
# ---------------------------------------------------------------------------

def reposition_zone(shape, x_in: float, y_in: float) -> None:
    """
    Move shape to (x_in, y_in) in inches. Width and height unchanged.

    Trigger: reading path check fails for default position (Rule 4A).
    Example: legend needs to move from full-width strip to top-right corner.
    """
    shape.left = Inches(x_in)
    shape.top  = Inches(y_in)


# ---------------------------------------------------------------------------
# Patch: resize_zone
# ---------------------------------------------------------------------------

def resize_zone(shape, width_in: float = None, height_in: float = None) -> None:
    """
    Resize shape. Pass only the dimension(s) to change.

    Trigger: content is longer or shorter than skeleton default.
    Example: 6-bullet list needs taller body zone than default 4-bullet height.
    """
    if width_in is not None:
        shape.width  = Inches(width_in)
    if height_in is not None:
        shape.height = Inches(height_in)


# ---------------------------------------------------------------------------
# Patch: add_hero_stat
# ---------------------------------------------------------------------------

def add_hero_stat(
    slide,
    stat_value: str,
    stat_label: str,
    x_in: float,
    y_in: float,
    w_in: float,
    h_in: float,
    color: RGBColor = RGBColor(0x00, 0x00, 0x00),
    stat_pt: float = 60,
    label_pt: float = 14,
) -> object:
    """
    Add a large-text stat with a smaller label below it.
    Returns the stat shape.

    Trigger: brief has a key quantitative finding to lead with.
    Example: "14%" cost increase as the hero number on a findings slide.
    """
    from pptx.util import Pt

    stat_h = Inches(h_in * 0.60)
    label_h = Inches(h_in * 0.35)
    gap = Inches(h_in * 0.05)

    # Stat value shape
    stat_shape = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), stat_h
    )
    tf = stat_shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = stat_value
    run.font.size = Pt(stat_pt)
    run.font.bold = True
    run.font.color.rgb = color

    # Label shape below
    label_shape = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in) + stat_h + gap, Inches(w_in), label_h
    )
    tf2 = label_shape.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = stat_label
    run2.font.size = Pt(label_pt)
    run2.font.color.rgb = color

    return stat_shape


# ---------------------------------------------------------------------------
# Patch: set_shape_fill
# ---------------------------------------------------------------------------

def set_shape_fill(shape, hex_color: str) -> None:
    """
    Set the fill color of a shape to hex_color (e.g. "#A100FF").

    Trigger: template-specific accent color needs to be applied to a zone.
    """
    color = hex_color.lstrip("#")
    r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Patch: set_rag_status
# ---------------------------------------------------------------------------

def set_rag_status(shape, status: str) -> None:
    """
    Set shape fill to the RAG color for status.
    status: "green"/"g", "amber"/"yellow"/"a", "red"/"r"

    Trigger: status field in brief maps to a RAG indicator shape.
    """
    color = RAG_COLORS.get(status.lower())
    if color is None:
        print(f"  [patch] WARNING: unknown RAG status '{status}'. Use green/amber/red.")
        return
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# Patch: add_table_row
# ---------------------------------------------------------------------------

def add_table_row(table_shape, row_data: list) -> None:
    """
    Append a row to a native PowerPoint table shape.
    row_data: list of strings, one per column. Extra values ignored,
              missing values filled with empty string.

    Trigger: table has a variable number of rows depending on brief data.
    """
    from pptx.oxml.ns import qn
    from copy import deepcopy

    tbl = table_shape.table
    num_cols = len(tbl.columns)

    # Clone the last row's XML as a template (preserves formatting)
    last_row_elem = tbl.rows[-1]._tr
    new_row_elem = deepcopy(last_row_elem)

    # Clear text in each cell of the new row
    for i, tc in enumerate(new_row_elem.findall(qn("a:tc"))):
        # Clear existing paragraphs
        txBody = tc.find(qn("a:txBody"))
        if txBody is not None:
            for p in txBody.findall(qn("a:p")):
                txBody.remove(p)
            # Add single paragraph with the new text
            p_elem = etree.SubElement(txBody, qn("a:p"))
            r_elem = etree.SubElement(p_elem, qn("a:r"))
            t_elem = etree.SubElement(r_elem, qn("a:t"))
            t_elem.text = row_data[i] if i < len(row_data) else ""

    # Insert the new row after the last row
    tbl_elem = tbl._tbl
    tbl_elem.append(new_row_elem)


# ---------------------------------------------------------------------------
# Patch: add_comparison_table
# ---------------------------------------------------------------------------

def add_comparison_table(
    slide,
    left_heading: str,
    right_heading: str,
    rows: list[tuple[str, str, str]],
    x_in: float = 0.607,
    y_in: float = 1.5,
    w_in: float = 12.12,
    h_row_in: float = 0.42,
    heading_pt: float = 12,
    body_pt: float = 11,
    criterion_pt: float = 10,
    left_col_w_pct: float = 0.22,
    header_fill: str = "#404040",
    alt_fill: str = "#F5F5F5",
    criterion_fill: str = "#E8E8E8",
) -> object:
    """
    Add a native two-column comparison table to the slide.

    Each row in `rows` is a 3-tuple: (criterion_label, left_value, right_value).
    The table renders as:
      | [criterion]  | [left_heading]     | [right_heading]    |
      | row 1 label  | row 1 left value   | row 1 right value  |
      | ...          | ...                | ...                |

    Replaces the 32-token text-box grid in the two-panel skeleton for cases
    where row count is variable. The seed table in the skeleton can be used for
    fixed 5-row comparisons; call this function for dynamic row counts.

    Returns the table shape.

    Trigger: two-panel brief has != 5 criteria rows, or brief data is tabular.
    """
    n_rows = len(rows) + 1  # +1 for header row
    n_cols = 3

    total_h = Inches(h_row_in * n_rows)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x_in), Inches(y_in),
        Inches(w_in), total_h,
    )
    tbl = table_shape.table

    # Column widths: criterion | left | right
    crit_w  = Inches(w_in * left_col_w_pct)
    half_w  = (Inches(w_in) - crit_w) // 2
    tbl.columns[0].width = int(crit_w)
    tbl.columns[1].width = int(half_w)
    tbl.columns[2].width = int(Inches(w_in) - crit_w - half_w)

    from pptx.oxml.ns import qn
    from lxml import etree

    def _fill_cell(cell, text, pt, bold=False, fill_hex=None, color_hex=None, align=PP_ALIGN.LEFT):
        if fill_hex:
            fh = fill_hex.lstrip("#")
            r, g, b = int(fh[0:2], 16), int(fh[2:4], 16), int(fh[4:6], 16)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(r, g, b)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(54864)
        tf.margin_top = tf.margin_bottom = Emu(36576)
        p_el = tf.paragraphs[0]._p
        for r in p_el.findall(qn("a:r")):
            p_el.remove(r)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.bold = bold
        if color_hex:
            ch = color_hex.lstrip("#")
            run.font.color.rgb = RGBColor(int(ch[0:2], 16), int(ch[2:4], 16), int(ch[4:6], 16))

    # Header row
    _fill_cell(tbl.cell(0, 0), "", heading_pt, bold=True,
               fill_hex=header_fill, color_hex="FFFFFF")
    _fill_cell(tbl.cell(0, 1), left_heading, heading_pt, bold=True,
               fill_hex=header_fill, color_hex="FFFFFF", align=PP_ALIGN.CENTER)
    _fill_cell(tbl.cell(0, 2), right_heading, heading_pt, bold=True,
               fill_hex=header_fill, color_hex="FFFFFF", align=PP_ALIGN.CENTER)

    # Data rows
    for i, (criterion, left_val, right_val) in enumerate(rows):
        row_fill = alt_fill if i % 2 == 1 else None
        _fill_cell(tbl.cell(i + 1, 0), criterion, criterion_pt, bold=True,
                   fill_hex=criterion_fill)
        _fill_cell(tbl.cell(i + 1, 1), left_val, body_pt, fill_hex=row_fill)
        _fill_cell(tbl.cell(i + 1, 2), right_val, body_pt, fill_hex=row_fill)

    return table_shape


# ---------------------------------------------------------------------------
# Patch: insert_icon (wrapper around icon_helper)
# ---------------------------------------------------------------------------

def insert_icon(
    slide,
    icon_name: str,
    x_in: float,
    y_in: float,
    size_in: float,
    accent_hex: str = "#000000",
) -> bool:
    """
    Insert a pre-extracted icon at (x_in, y_in) with square bounding box size_in.
    Color-swaps the icon fill to accent_hex.
    Returns True on success, False if placeholder was inserted.

    Trigger: brief mentions a concept that maps to an icon in the icon vocabulary.
    Available icons: any filename (without extension) under <skill>/icons/*.xml.
    """
    try:
        from icon_helper import insert_icon as _insert_icon
    except ImportError:
        print("  [patch] ERROR: icon_helper not found. Check scripts/ is in sys.path.")
        return False

    return _insert_icon(
        icon_name=icon_name,
        target_slide=slide,
        left_emu=int(Inches(x_in)),
        top_emu=int(Inches(y_in)),
        width_emu=int(Inches(size_in)),
        height_emu=int(Inches(size_in)),
        accent_color=accent_hex,
    )
