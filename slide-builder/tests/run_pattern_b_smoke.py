#!/usr/bin/env python3
"""End-to-end Pattern B Python-side smoke (M8 deliverable, 2026-06-17/18).

What this smoke covers
----------------------
The Pattern B pipeline has two execution surfaces:

  1. Chat-orchestrated steps  — agent dispatch (worker, translator).
  2. Python-side steps         — build_deck.py prompt rendering,
                                  render_html.py HTML→PNG, finalize_deck.py
                                  classifier + template-fields parsing +
                                  R4 QC checks + placeholder population.

Agent dispatch can't be exercised from a Python test (it's a chat-level
operation). This smoke covers the *Python* surface end-to-end so any
regression in the M1–M7 plumbing fails LOUDLY before it lands on Mario.

Specifically, the smoke:

  P1  Verifies build_placeholders() emits PATTERN: B|C based on the slide
      classifier output (M1 + M5).
  P2  Verifies scripts/render_html.py produces a 1280×720 PNG from a stub
      Pattern B HTML file (M3).
  P3  Verifies finalize_deck._classify_option() correctly recognizes a
      translator-style script header as ``pattern_b_translated`` (M5).
  P4  Verifies finalize_deck._parse_template_fields() extracts the
      __template_fields__ dict from a translator script header (M5).
  P5  Verifies finalize_deck._check_r4_rules_for_pattern_b() emits the
      eight R4 rule entries from Spec 6 with the locked severities
      (3 Critical / 4 Major / 1 Advisory) (M6).
  P6  Verifies the synthetic native script executes cleanly to produce a
      .pptx whose body shapes are editable (R4.7) (M5).
  P7  Verifies _meta_schema.MetaJson accepts both legacy metas (with
      mermaid_theme populated) and new metas (without), confirming the
      M7 retirement preserves backward compatibility.

Run
---
    py -3 slide-builder/tests/run_pattern_b_smoke.py

Exit
----
    0  -> all phases passed
    1  -> at least one assertion failed
"""
from __future__ import annotations

import io
import json
import os
import subprocess
import sys
import tempfile
import textwrap
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
SLIDE_BUILDER = REPO_ROOT / "slide-builder"
sys.path.insert(0, str(SLIDE_BUILDER / "scripts"))
sys.path.insert(0, str(SLIDE_BUILDER))

FAILURES: list[str] = []


def _check(label: str, ok: bool, detail: str = "") -> None:
    if ok:
        print(f"  ok: {label}")
    else:
        FAILURES.append(f"{label}: {detail}")
        print(f"  FAIL: {label}{(' — ' + detail) if detail else ''}")


# ---------------------------------------------------------------------------
# P1 — build_placeholders emits PATTERN per slide_pattern arg (M1 + M5)
# ---------------------------------------------------------------------------
def phase_1_build_placeholders() -> None:
    print("\n[P1] build_placeholders threads PATTERN field per M5")
    import build_deck  # noqa: E402  (placed after sys.path injection above)

    slide = {"slide_n": 7, "title": "T", "governing_thought": "g", "so_what": "s"}
    seeds = {
        "content_hash": "h", "pattern_pick_seed": "p",
        "variant_seed_a": "a", "variant_seed_b": "b", "variant_seed_c": "c",
    }
    ph_b = build_deck.build_placeholders(
        slide=slide, slide_total=8, deck_notes="",
        client_template_path=Path("."), output_dir=Path("."),
        seeds=seeds, likely_prior_patterns="", slide_pattern="B",
    )
    _check("PATTERN=B when slide_pattern='B'", ph_b.get("PATTERN") == "B",
           f"got {ph_b.get('PATTERN')!r}")
    ph_default = build_deck.build_placeholders(
        slide=slide, slide_total=8, deck_notes="",
        client_template_path=Path("."), output_dir=Path("."),
        seeds=seeds, likely_prior_patterns="",
    )
    _check("PATTERN=C when slide_pattern omitted (legacy default)",
           ph_default.get("PATTERN") == "C",
           f"got {ph_default.get('PATTERN')!r}")


# ---------------------------------------------------------------------------
# P2 — render_html.py produces a 1280×720 PNG (M3)
# ---------------------------------------------------------------------------
def phase_2_render_html(tmp: Path) -> Path:
    print("\n[P2] scripts/render_html.py turns HTML into a 1280×720 PNG")
    html_path = tmp / "stub.html"
    png_path = tmp / "stub.png"
    html_path.write_text(textwrap.dedent("""\
        <!DOCTYPE html>
        <html><head><meta charset="utf-8"><style>
          html,body{margin:0;padding:0;width:1280px;height:720px;background:#fff;
                    font-family:Segoe UI,sans-serif;}
          .t{position:absolute;top:30px;left:60px;font-size:28px;color:#4D148C;}
          .box{position:absolute;top:200px;left:80px;width:1120px;height:300px;
               background:#F4EEFB;border-radius:18px;}
        </style></head><body>
          <div class="t" data-template-field="title">Smoke title</div>
          <div class="box" data-shape-id="card">Body card</div>
        </body></html>"""), encoding="utf-8")
    r = subprocess.run(
        ["py", "-3", str(SLIDE_BUILDER / "scripts" / "render_html.py"),
         str(html_path), str(png_path)],
        capture_output=True, text=True, timeout=90,
    )
    _check("render_html.py exits 0", r.returncode == 0,
           f"stderr: {(r.stderr or '').strip()[:200]}")
    _check("PNG exists at the requested path", png_path.exists())
    if png_path.exists():
        from PIL import Image
        im = Image.open(png_path)
        _check("PNG dimensions are 1280×720", im.size == (1280, 720),
               f"got {im.size}")
    return png_path


# ---------------------------------------------------------------------------
# P3 + P4 — classifier + template-fields parser (M5)
# ---------------------------------------------------------------------------
TRANSLATOR_SCRIPT = textwrap.dedent("""\
    # CONTEXT_READ: smoke test fixture
    # BRIEF_IS_AUTHORITATIVE: True
    # PATTERN: B-translated
    # HTML_SOURCE: option_A.html (sha256: deadbeefdeadbeef)
    # PNG_TARGET: option_A.png (sha256: cafebabecafebabe)
    # __template_fields__ = {
    #     "title": "Smoke title",
    #     "subtitle": "Smoke subtitle",
    #     "footer": "Smoke footer",
    #     "page_number": "1 / 1",
    # }

    from pathlib import Path
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    def build():
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Single editable body shape; satisfies R4.7.
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Emu(96 * 9525), Emu(200 * 9525),
                                    Emu(1088 * 9525), Emu(300 * 9525))
        sh.adjustments[0] = 0.1
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF4, 0xEE, 0xFB)
        sh.line.color.rgb = RGBColor(0x4D, 0x14, 0x8C)
        sh.text_frame.text = "Body card"
        prs.save(str(Path(__file__).resolve().parent / "option_A_native.pptx"))

    if __name__ == "__main__":
        build()
""")

TRANSLATION_REPORT = {
    "slide_n": 1, "option_letter": "A",
    "html_source": "option_A.html", "png_target": "option_A.png",
    "html_sha256": "deadbeefdeadbeef", "png_sha256": "cafebabecafebabe",
    "template_fields": {
        "title": "Smoke title", "subtitle": "Smoke subtitle",
        "footer": "Smoke footer", "page_number": "1 / 1",
    },
    "shape_count_html": 1, "shape_count_native": 1, "shape_count_match": True,
    "translator_self_check": {
        "ssim_per_zone": {"title": 0.95, "subtitle": 0.93, "body": 0.91, "footer": 0.94},
        "lowest_zone_score": 0.91, "pass": True,
    },
    "css_kill_list_applied": {
        "gradients_flattened": 0, "shadows_stripped": 0, "filters_dropped": 0,
        "text_decorations_stripped": 0, "opacity_on_text_normalized": 0,
    },
    "editability_self_check": {
        "add_picture_with_text": False, "zero_size_text_frames": False,
        "chrome_zone_freeform_shapes": False, "shapes_outside_canvas": False,
        "pass": True,
    },
    "warnings": [],
}


def phase_3_classify_and_parse(tmp: Path) -> tuple[Path, Path]:
    print("\n[P3+P4] finalize_deck classifier + __template_fields__ parser (M5)")
    import finalize_deck  # noqa: E402

    slide_dir = tmp / "slide_01"
    slide_dir.mkdir()
    py = slide_dir / "option_A_native.py"
    py.write_text(TRANSLATOR_SCRIPT, encoding="utf-8")
    report = slide_dir / "option_A_translation_report.json"
    report.write_text(json.dumps(TRANSLATION_REPORT), encoding="utf-8")

    status, reason = finalize_deck._classify_option(py)
    _check("classifier returns 'pattern_b_translated' on translator output",
           status == "pattern_b_translated", f"got {status!r} ({reason})")

    fields = finalize_deck._parse_template_fields(py)
    _check("__template_fields__ has all 4 keys",
           set(fields) == {"title", "subtitle", "footer", "page_number"},
           f"got {list(fields)}")
    _check("title parsed verbatim", fields.get("title") == "Smoke title",
           f"got {fields.get('title')!r}")
    return py, report


# ---------------------------------------------------------------------------
# P5 — R4.1-R4.8 check helper emits all 8 entries with locked severities (M6)
# ---------------------------------------------------------------------------
def phase_5_r4_checks(tmp: Path, py_path: Path) -> None:
    print("\n[P5] R4.1-R4.8 QC helper emits 8 entries with Spec 6 severities")
    import finalize_deck  # noqa: E402
    from finalize_deck import OptionStatus

    st = OptionStatus(
        slide_n=1, letter="A", py_path=py_path,
        pptx_path=py_path.with_suffix(".pptx"),
        raw_archive_path=tmp / "_raw" / "option_A.pptx",
        themed_pptx_path=py_path.with_suffix(".pptx"),
        themed_png_path=py_path.with_suffix(".png"),
        classification="pattern_b_translated", classification_reason="",
    )
    checks = finalize_deck._check_r4_rules_for_pattern_b(st)
    _check("emits exactly 8 R4 checks", len(checks) == 8, f"got {len(checks)}")

    # Severity counts per Spec 6: 3 Critical / 4 Major / 1 Advisory
    block_n = sum(1 for c in checks if c.get("severity") == "block")
    warn_n  = sum(1 for c in checks if c.get("severity") == "warn")
    info_n  = sum(1 for c in checks if c.get("severity") == "info")
    _check("3 Critical (severity=block)", block_n == 3, f"got {block_n}")
    _check("4 Major (severity=warn)",     warn_n  == 4, f"got {warn_n}")
    _check("1 Advisory (severity=info)",  info_n  == 1, f"got {info_n}")

    # Every rule passes on the clean smoke fixture
    fails = [c for c in checks if not c.get("pass")]
    _check("all 8 checks pass on a clean fixture",
           not fails, f"failing checks: {[c.get('check') for c in fails]}")

    # Pattern C / legacy classification returns empty
    st.classification = "native"
    legacy = finalize_deck._check_r4_rules_for_pattern_b(st)
    _check("Pattern C / legacy classification returns []",
           legacy == [], f"got {legacy!r}")


# ---------------------------------------------------------------------------
# P6 — generated native script executes + produces editable PPTX (R4.7)
# ---------------------------------------------------------------------------
def phase_6_execute_and_assert_editable(tmp: Path, py_path: Path) -> None:
    print("\n[P6] native script executes; resulting PPTX is editable (R4.7)")
    r = subprocess.run(["py", "-3", str(py_path)],
                       capture_output=True, text=True, timeout=60,
                       cwd=str(py_path.parent))
    _check("native script exits 0", r.returncode == 0,
           f"stderr: {(r.stderr or '').strip()[:200]}")
    pptx = py_path.parent / "option_A_native.pptx"
    _check("PPTX produced", pptx.exists())
    if pptx.exists():
        from pptx import Presentation
        prs = Presentation(str(pptx))
        slide = prs.slides[0]
        n_text_frames = sum(
            1 for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text
        )
        n_pictures = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        _check("PPTX has at least one editable text frame",
               n_text_frames >= 1, f"got {n_text_frames}")
        _check("PPTX has zero picture-of-text shapes (R4.7)",
               n_pictures == 0, f"got {n_pictures}")


# ---------------------------------------------------------------------------
# P7 — schema compat for the M7 mermaid_theme optional default
# ---------------------------------------------------------------------------
def phase_7_schema_compat() -> None:
    print("\n[P7] MetaJson accepts legacy v3 (with mermaid_theme) AND new v3 (without)")
    from _meta_schema import validate_meta_dict

    base = {
        "schema_version": 3,
        "template": "/tmp/t.pptx", "brief": "/tmp/b.md", "out": "/tmp/o",
        "client_slug": "smoke", "slide_count": 1,
        "generated_at": "2026-06-17T00:00:00+00:00",
        "slides": [{"n": 1, "title": "T", "layout": "Use as default slide template"}],
        "deck_meta": {"deck_type": "test", "governing_thought": "g", "audience": "a"},
    }
    legacy = dict(base, mermaid_theme="/tmp/mermaid.json")
    legacy_meta = validate_meta_dict(legacy)
    _check("legacy v3 meta with mermaid_theme validates",
           legacy_meta.mermaid_theme == "/tmp/mermaid.json")

    new = dict(base)  # no mermaid_theme key
    new_meta = validate_meta_dict(new)
    _check("new v3 meta WITHOUT mermaid_theme validates (M7 optional default)",
           new_meta.mermaid_theme == "")


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------
def main() -> int:
    print(f"Pattern B end-to-end Python smoke (M8) — repo at {REPO_ROOT}")
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        phase_1_build_placeholders()
        phase_2_render_html(tmp)
        py, _report = phase_3_classify_and_parse(tmp)
        phase_5_r4_checks(tmp, py)
        phase_6_execute_and_assert_editable(tmp, py)
        phase_7_schema_compat()

    print()
    if FAILURES:
        print("=" * 72)
        print(f"SMOKE FAILED: {len(FAILURES)} assertion(s) failed.")
        for f in FAILURES:
            print(f"  - {f}")
        return 1
    print("All phases passed.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
