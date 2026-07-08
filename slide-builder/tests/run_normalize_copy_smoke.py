#!/usr/bin/env python3
"""Smoke test for the registration normalize step (Stage 3B).

Registration writes a normalized build copy and applies safe fixes to it while
never touching the user's original. This proves:
  - the ORIGINAL template is byte-identical before/after registration;
  - the build copy has its sample slides stripped;
  - a title placeholder that lost its empty text line is REPAIRED in the copy,
    while the original keeps the defect (the fix lands on the copy, not the file).

Run:  py -3 slide-builder/tests/run_normalize_copy_smoke.py
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import hashlib
import shutil
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import run_layout_inheritance_smoke as fixture
import _paths as _p
from pptx import Presentation
from pptx.oxml.ns import qn

_TITLE_TYPES = {1, 13}  # CENTER_TITLE, TITLE (python-pptx PP_PLACEHOLDER ints)


def _first_title_ph(layout):
    for ph in layout.placeholders:
        t = ph.placeholder_format.type
        if t is not None and int(t) in _TITLE_TYPES:
            return ph
    return None


def _para_count(ph) -> int:
    tb = ph._element.find(qn("p:txBody"))
    return 0 if tb is None else len(tb.findall(qn("a:p")))


def main() -> int:
    tmp = Path(tempfile.mkdtemp(prefix="normalize_copy_smoke_"))
    try:
        tpl = tmp / fixture.FIXTURE_PPTX.name
        shutil.copy2(fixture.FIXTURE_PPTX, tpl)

        # Craft a defective + non-empty original: strip <a:p> from one layout's
        # title placeholder, and add a sample slide.
        prs = Presentation(str(tpl))
        mi = li = None
        for _mi, master in enumerate(prs.slide_masters):
            for _li, layout in enumerate(master.slide_layouts):
                ph = _first_title_ph(layout)
                if ph is not None:
                    tb = ph._element.find(qn("p:txBody"))
                    if tb is not None and tb.find(qn("a:p")) is not None:
                        for p in tb.findall(qn("a:p")):
                            tb.remove(p)
                        mi, li = _mi, _li
                        break
            if mi is not None:
                break
        assert mi is not None, "could not find a title placeholder to break"
        prs.slides.add_slide(prs.slide_layouts[0])  # a sample slide to strip
        prs.save(str(tpl))

        prs0 = Presentation(str(tpl))
        assert _para_count(_first_title_ph(
            prs0.slide_masters[mi].slide_layouts[li])) == 0, "setup: title not broken"
        assert len(prs0.slides._sldIdLst) >= 1, "setup: no sample slide present"
        sha_before = hashlib.sha256(tpl.read_bytes()).hexdigest()

        print("[1] Register -> normalize the build copy")
        fixture.register_fixture(tpl)
        copy = _p.build_template_pptx(tpl)
        assert copy.exists(), f"build copy not created: {copy}"

        print("[2] Original is byte-identical (never modified)")
        assert hashlib.sha256(tpl.read_bytes()).hexdigest() == sha_before, "original modified!"
        prsO = Presentation(str(tpl))
        assert _para_count(_first_title_ph(
            prsO.slide_masters[mi].slide_layouts[li])) == 0, "original defect should remain"
        print("    ok: original unchanged, defect + sample slide retained")

        print("[3] Copy repaired the empty title placeholder")
        prsC = Presentation(str(copy))
        assert _para_count(_first_title_ph(
            prsC.slide_masters[mi].slide_layouts[li])) >= 1, "copy title not repaired"
        print("    ok: copy title placeholder has a text line")

        print("[4] Copy has no sample slides")
        assert len(prsC.slides._sldIdLst) == 0, "copy still has sample slides"
        print("    ok: copy stripped to 0 slides")

        print("\nSMOKE PASSED.")
        return 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
