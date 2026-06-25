#!/usr/bin/env python3
"""Smoke test for the font-size lock (PowerPoint default grid, floor 8pt).

Covers the three layers that keep every shipped text size on PowerPoint's
default grid:
  - twins.helpers.snap_font_pt: snaps a value to the nearest default, floor 8,
    ties rounding down,
  - finalize_deck._snap_font_sizes: locks every run in a slide (text boxes,
    table cells, grouped shapes),
  - the run_option_qc "font_size_locked" guard: flags off-grid sizes and passes
    once they're snapped.

Run:  py -3 slide-builder/tests/run_font_lock_smoke.py
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SKILL = HERE.parent
sys.path.insert(0, str(SKILL))
sys.path.insert(0, str(SKILL / "scripts"))

from pptx import Presentation
from pptx.util import Emu, Pt

from twins.helpers import snap_font_pt, ALLOWED_FONT_SIZES_PT
import finalize_deck as F


def main() -> int:
    print("[1] snap_font_pt")
    cases = {
        7.3: 8, 8.0: 8, 8.2: 8, 8.6: 9, 9.0: 9, 10.3: 10.5, 10.5: 10.5,
        11.0: 11, 12.75: 12, 13.0: 12, 13.5: 14, 6.0: 8, 100.0: 96, 96.0: 96,
    }
    for inp, exp in cases.items():
        got = snap_font_pt(inp)
        assert got == exp, f"snap_font_pt({inp}) = {got}, expected {exp}"
    assert snap_font_pt(None) is None
    assert all(snap_font_pt(s) == float(s) for s in ALLOWED_FONT_SIZES_PT), \
        "every grid value must snap to itself"
    print(f"    ok: {len(cases)} cases + grid fixed-points")

    print("[2] _snap_font_sizes (text box + table cell)")
    d = Path(tempfile.mkdtemp(prefix="font_lock_smoke_"))
    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    tf = sl.shapes.add_textbox(Emu(914400), Emu(914400), Emu(4000000), Emu(3000000)).text_frame
    for i, sz in enumerate([7.3, 8.2, 13.0, 11.0, 6.0]):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        r = p.add_run(); r.text = f"run{i}"; r.font.size = Pt(sz)
    gf = sl.shapes.add_table(1, 1, Emu(914400), Emu(4200000), Emu(2000000), Emu(500000))
    cr = gf.table.cell(0, 0).text_frame.paragraphs[0].add_run()
    cr.text = "cell"; cr.font.size = Pt(9.4)
    F._snap_font_sizes(sl.shapes)
    got = [round(r.font.size.pt, 2) for p in tf.paragraphs for r in p.runs if r.font.size]
    cell_pt = round(gf.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.size.pt, 2)
    assert got == [8, 8, 12, 11, 8], got
    assert cell_pt == 9, cell_pt
    print(f"    ok: text box {got}, table cell {cell_pt}pt")

    print("[3] run_option_qc font_size_locked guard")
    off = d / "off.pptx"
    prs2 = Presentation()
    s2 = prs2.slides.add_slide(prs2.slide_layouts[6])
    t2 = s2.shapes.add_textbox(Emu(914400), Emu(914400), Emu(4000000), Emu(2000000)).text_frame
    for i, sz in enumerate([8.2, 7.3, 13.0]):
        p = t2.add_paragraph() if i else t2.paragraphs[0]
        r = p.add_run(); r.text = f"body {i}"; r.font.size = Pt(sz)
    prs2.save(str(off))
    chk = {c["check"]: c for c in F.run_option_qc(off, d / "x.png", set(), page_type="body")["checks"]}
    assert chk["font_size_locked"]["pass"] is False, "off-grid should be flagged"
    prs3 = Presentation(str(off)); F._snap_font_sizes(prs3.slides[0].shapes); prs3.save(str(off))
    chk2 = {c["check"]: c for c in F.run_option_qc(off, d / "x.png", set(), page_type="body")["checks"]}
    assert chk2["font_size_locked"]["pass"] is True, "should pass after snap"
    print("    ok: flags off-grid, passes after snap")

    import shutil
    shutil.rmtree(d, ignore_errors=True)
    print("\nSMOKE PASSED.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
