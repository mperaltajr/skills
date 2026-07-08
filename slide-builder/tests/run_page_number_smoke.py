#!/usr/bin/env python3
"""Smoke test for page-number rendering on grafted slides.

python-pptx's add_slide clones only title/body placeholders — it skips the
SLIDE_NUMBER placeholder (treated as inherited chrome). The inherited auto-number
field also doesn't render under LibreOffice headless. So page numbers silently
vanished on every body-canonical build. The fix:
  - composer.clone_missing_chrome_placeholders() clones the SLIDE_NUMBER
    placeholder onto the slide, and
  - _populate_layout_placeholders() (PAGE_NUM_TYPES = {13} = SLIDE_NUMBER) writes
    a LITERAL page number into it, clearing any auto-field so it renders.

This test proves both, using python-pptx's default template (its Blank layout
carries a slide-number placeholder).

Run:  py -3 slide-builder/tests/run_page_number_smoke.py
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
SKILL_ROOT = HERE.parent
sys.path.insert(0, str(SKILL_ROOT))            # for `import twins`
sys.path.insert(0, str(SKILL_ROOT / "scripts"))  # for scripts-local imports

from pptx import Presentation
from pptx.oxml.ns import qn
from twins.composer import clone_missing_chrome_placeholders, _populate_layout_placeholders

SLIDE_NUMBER = 13


def _has_slide_number_ph(obj) -> bool:
    return any(int(ph.placeholder_format.type) == SLIDE_NUMBER
              for ph in obj.placeholders
              if ph.placeholder_format.type is not None)


def main() -> int:
    prs = Presentation()
    # Find a layout that carries a slide-number placeholder.
    layout = next((l for l in prs.slide_layouts if _has_slide_number_ph(l)), None)
    assert layout is not None, "default template has no layout with a slide-number placeholder"

    print("[1] add_slide does NOT clone the slide-number placeholder")
    slide = prs.slides.add_slide(layout)
    assert not _has_slide_number_ph(slide), "unexpected: add_slide cloned the slide-number ph"
    print("    ok: slide has no slide-number placeholder after add_slide")

    print("[2] clone_missing_chrome_placeholders adds it")
    n = clone_missing_chrome_placeholders(slide, layout)
    assert n >= 1, f"expected to clone >=1 placeholder, cloned {n}"
    assert _has_slide_number_ph(slide), "slide-number placeholder not present after clone"
    print(f"    ok: cloned {n}; slide now has a slide-number placeholder")

    print("[3] _populate writes a literal page number (no auto-field)")
    found = _populate_layout_placeholders(slide, page_num="7")
    assert found["page_number"], "page number did not populate"
    ph = next(p for p in slide.placeholders if int(p.placeholder_format.type) == SLIDE_NUMBER)
    assert ph.text_frame.text == "7", f"page number text wrong: {ph.text_frame.text!r}"
    # no residual auto-field (<a:fld>) that LibreOffice wouldn't render
    flds = ph._element.findall(".//" + qn("a:fld"))
    assert not flds, f"auto-field left in slide-number placeholder ({len(flds)})"
    print("    ok: slide number = '7', no auto-field")

    print("\nSMOKE PASSED.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
