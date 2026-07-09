#!/usr/bin/env python3
"""Smoke test for adopt_deck.py (option 6b ingest/scaffold).

Asserts:
  - the registration GATE fires on an unregistered deck (exit 7, tells you to
    register — never registers inline);
  - once the deck's sidecars exist, adopt writes: a rebuild-slice brief with a
    per-slide section carrying the extracted title, a `_meta.json` marked
    `adopted_source` with per-slide layouts, a copy of the original at
    adopted_source.pptx, and the requested slide dir.

The downstream build_deck --slide / finalize / splice reuse the existing,
already-smoke-tested rebuild path (splice is covered by run_splice_smoke.py).

Run:  py -3 slide-builder/tests/run_adopt_smoke.py   (python3 on macOS/Linux)
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import _paths as _p
from pptx import Presentation
from pptx.util import Inches


def _make_deck(path: Path):
    prs = Presentation()
    for t, b in [("Kickoff Overview", "Scope and goals"),
                ("Status: GREEN", "12 of 15 KPIs on track"),
                ("Next Steps", "Approve budget")]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tb.text_frame.text = t
        bb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        bb.text_frame.text = b
    prs.save(str(path))


def _fake_register(deck: Path):
    """Create minimal registration sidecars so the gate passes (existence-only)."""
    for path in (_p.brand_yml(deck), _p.chrome_yml(deck), _p.theme_json(deck)):
        Path(path).parent.mkdir(parents=True, exist_ok=True)
        Path(path).write_text("{}" if str(path).endswith(".json") else "brand: {}\n", encoding="utf-8")


def main() -> int:
    tmp = Path(tempfile.mkdtemp(prefix="adopt_smoke_"))
    try:
        deck = tmp / "client_pmo.pptx"
        _make_deck(deck)
        out_dir = tmp / "adopt_out"

        print("[1] registration gate fires on an unregistered deck")
        r = subprocess.run([sys.executable, str(SCRIPTS / "adopt_deck.py"), str(deck),
                            "--out", str(out_dir), "--slides", "2"],
                           capture_output=True, text=True)
        assert r.returncode == 7, f"expected exit 7, got {r.returncode}\n{r.stdout}"
        assert "register" in r.stdout.lower(), r.stdout
        print("    ok: unregistered deck refused with the register instruction")

        print("[2] adopt a registered deck")
        _fake_register(deck)
        r = subprocess.run([sys.executable, str(SCRIPTS / "adopt_deck.py"), str(deck),
                            "--out", str(out_dir), "--slides", "2"],
                           capture_output=True, text=True)
        assert r.returncode == 0, f"adopt failed:\n{r.stdout}\n{r.stderr}"

        print("[3] verify scaffolding")
        meta = json.loads((out_dir / "_meta.json").read_text(encoding="utf-8"))
        assert Path(meta.get("adopted_source", "")).resolve() == deck.resolve(), meta.get("adopted_source")
        assert meta.get("slide_count") == 3, meta
        assert len(meta["slides"]) == 3 and all("layout" in s for s in meta["slides"]), meta["slides"]
        assert meta["slides"][1]["title"] == "Status: GREEN", meta["slides"][1]
        brief = (out_dir / "adopted_brief.md").read_text(encoding="utf-8")
        assert "mode: rebuild-slice" in brief, "brief must bypass the narrative gate"
        assert "## Slide 2 — Status: GREEN" in brief, brief[:400]
        assert "12 of 15 KPIs on track" in brief, "body text not extracted"
        assert (out_dir / "adopted_source.pptx").exists(), "splice target copy missing"
        assert (out_dir / "slide_02").is_dir(), "target slide dir not created"
        # original untouched
        assert len(list(Presentation(str(deck)).slides)) == 3
        print("    ok: brief (rebuild-slice) + _meta (adopted_source, layouts) + copy + slide dir")

        print("\nSMOKE PASSED.")
        return 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
