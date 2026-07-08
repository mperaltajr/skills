#!/usr/bin/env python3
"""Smoke test for inserting a slide into an existing build (build_deck.py --insert N).

Builds a 3-slide deck, writes picks, then inserts a new slide at position 2 and
asserts the insert is surgical and correctly renumbered:
  - the deck grows to 4 slides, order 1..4, titles [old1, NEW, old2, old3];
  - only the new slide's dir is freshly prepped; the shifted slides keep their
    built output under their new numbers;
  - picks.json keys shift (>= N move up one; the new slide has no pick yet);
  - error paths (out-of-range N, brief not exactly one slide longer) exit 2.

Run:  py -3 slide-builder/tests/run_insert_slice_smoke.py
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import run_layout_inheritance_smoke as fixture  # reuse fixture template + registration


def _brief(tpl: Path, titles: list[str]) -> str:
    head = (
        f"---\nclient_template: {tpl}\ndeck_type: Smoke test\n"
        f"default_layout: body_canonical_light\nmode: template-fill\n---\n\n"
        f"## Deck-level design notes\n\nInsert smoke fixture.\n\n"
    )
    body = ""
    for i, title in enumerate(titles, 1):
        body += (
            f"## Slide {i} — {title}\n"
            f"**Layout:** body_canonical_light\n"
            f"**Archetype:** Synthesis / Findings\n"
            f"**Governing thought:** {title} governing thought.\n"
            f"**So-what:** {title} so-what.\n"
            f"**Evidence:** N/A.\n\n"
        )
    return head + body


def _build(*args: str) -> subprocess.CompletedProcess:
    return subprocess.run(
        [sys.executable, str(SCRIPTS / "build_deck.py"), *args],
        capture_output=True, text=True,
    )


def _titles(meta: dict) -> list[str]:
    return [s["title"] for s in sorted(meta["slides"], key=lambda e: e["n"])]


def main() -> int:
    tmp = Path(tempfile.mkdtemp(prefix="insert_slice_smoke_"))
    try:
        tpl = tmp / fixture.FIXTURE_PPTX.name
        shutil.copy2(fixture.FIXTURE_PPTX, tpl)
        fixture.register_fixture(tpl)
        out = tmp / "build"

        print("[1] Full 3-slide build")
        brief3 = tmp / "b3.md"
        brief3.write_text(_brief(tpl, ["Alpha", "Beta", "Gamma"]), encoding="utf-8")
        r = _build("--brief", str(brief3), "--template", str(tpl), "--out", str(out), "--confirm-template")
        assert r.returncode == 0, f"full build failed ({r.returncode}):\n{r.stderr[-800:]}"
        meta0 = json.loads((out / "_meta.json").read_text(encoding="utf-8"))
        assert meta0["slide_count"] == 3 and _titles(meta0) == ["Alpha", "Beta", "Gamma"], _titles(meta0)
        # seed picks so we can prove the shift
        (out / "picks.json").write_text(
            json.dumps({"slide_01": "A", "slide_02": "B", "slide_03": "C"}), encoding="utf-8")
        print("    ok: 3 slides [Alpha, Beta, Gamma] + picks A/B/C")

        print("[2] Insert a new slide at position 2")
        brief4 = tmp / "b4.md"
        brief4.write_text(_brief(tpl, ["Alpha", "INSERTED", "Beta", "Gamma"]), encoding="utf-8")
        r = _build("--brief", str(brief4), "--template", str(tpl), "--out", str(out),
                   "--insert", "2", "--confirm-template")
        assert r.returncode == 0, f"insert failed ({r.returncode}):\n{r.stderr[-800:]}"
        meta1 = json.loads((out / "_meta.json").read_text(encoding="utf-8"))
        assert meta1["slide_count"] == 4, meta1["slide_count"]
        assert [s["n"] for s in sorted(meta1["slides"], key=lambda e: e["n"])] == [1, 2, 3, 4]
        assert _titles(meta1) == ["Alpha", "INSERTED", "Beta", "Gamma"], _titles(meta1)
        assert (out / "slide_02" / "_prompt.md").exists(), "new slide 2 not prepped"
        assert (out / "slide_04").is_dir(), "shifted slide 4 dir missing"
        picks = json.loads((out / "picks.json").read_text(encoding="utf-8"))
        assert picks == {"slide_01": "A", "slide_03": "B", "slide_04": "C"}, picks
        # the human-readable dispatch plan renumbers to match the shifted deck
        plan = (out / "dispatch_plan.md").read_text(encoding="utf-8")
        for n, title in ((1, "Alpha"), (2, "INSERTED"), (3, "Beta"), (4, "Gamma")):
            assert f"| {n} | {title} |" in plan, f"dispatch_plan missing row {n} {title}:\n{plan}"
        assert "| 3 | Gamma |" not in plan, f"dispatch_plan kept stale numbering:\n{plan}"
        print("    ok: 4 slides [Alpha, INSERTED, Beta, Gamma]; picks shifted A/-/B/C; plan renumbered")

        print("[3] Error paths")
        r = _build("--brief", str(brief4), "--template", str(tpl), "--out", str(out),
                   "--insert", "9", "--confirm-template")
        assert r.returncode == 2, f"out-of-range insert: expected 2, got {r.returncode}"
        # brief not exactly one longer than the current 4-slide build
        r = _build("--brief", str(brief3), "--template", str(tpl), "--out", str(out),
                   "--insert", "2", "--confirm-template")
        assert r.returncode == 2, f"count-mismatch insert: expected 2, got {r.returncode}"
        print("    ok: out-of-range=2, count-mismatch=2")

        print("[4] Transactional: a mis-numbered brief must NOT shift the build")
        before = sorted(p.name for p in out.glob("slide_*"))
        misnum = tmp / "misnum.md"  # 5 slides but numbered 1,2,3,4,6 (gap at 5)
        head = (
            f"---\nclient_template: {tpl}\ndeck_type: Smoke test\n"
            f"default_layout: body_canonical_light\nmode: template-fill\n---\n\n"
            f"## Deck-level design notes\n\nMis-numbered.\n\n"
        )
        misbody = ""
        for n in (1, 2, 3, 4, 6):
            misbody += (
                f"## Slide {n} — S{n}\n**Layout:** body_canonical_light\n"
                f"**Archetype:** Synthesis / Findings\n**Governing thought:** g.\n"
                f"**So-what:** s.\n**Evidence:** N/A.\n\n"
            )
        misnum.write_text(head + misbody, encoding="utf-8")
        r = _build("--brief", str(misnum), "--template", str(tpl), "--out", str(out),
                   "--insert", "2", "--confirm-template")
        assert r.returncode == 2, f"mis-numbered insert: expected 2, got {r.returncode}"
        after = sorted(p.name for p in out.glob("slide_*"))
        assert after == before, f"build was shifted despite a rejected insert: {before} -> {after}"
        print("    ok: rejected (exit 2) and the build dirs are unchanged")

        print("[5] Page numbers re-stamp from compile position, not the baked value")
        import compile_picks
        from copy import deepcopy
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER

        def _slide_number_ph(slide):
            for ph in slide.placeholders:
                if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    return ph
            return None

        prs = Presentation()
        lay = prs.slide_layouts[6]  # Blank; its layout carries a slide-number placeholder
        sl = prs.slides.add_slide(lay)
        # The real pipeline deepcopies a finalized source slide (which carries a
        # populated slide-number placeholder) into the compiled deck; reproduce that
        # by copying the layout's slide-number placeholder onto the slide.
        lay_pn = next(ph for ph in lay.placeholders
                      if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER)
        sl.shapes._spTree.append(deepcopy(lay_pn._element))
        _slide_number_ph(sl).text_frame.text = "3"  # simulate finalize baking the original number
        compile_picks._restamp_page_number(sl, 7)
        got = _slide_number_ph(sl).text_frame.text
        assert got == "7", f"page number not re-stamped: {got!r}"
        # a non-numeric placeholder (empty / auto field) must be left alone
        _slide_number_ph(sl).text_frame.text = ""
        compile_picks._restamp_page_number(sl, 9)
        got2 = _slide_number_ph(sl).text_frame.text
        assert got2 == "", f"empty page-number placeholder should be left alone, got {got2!r}"
        print("    ok: baked '3' -> '7' by position; empty/field left alone")

        print("\nSMOKE PASSED.")
        return 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
