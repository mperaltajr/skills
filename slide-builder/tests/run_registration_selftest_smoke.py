#!/usr/bin/env python3
"""Smoke test for the registration mock-page self-test (Stage 3C).

Proves the render-based self-test in register_template._render_mock_page_selftest:
  - happy path: a well-formed template's default content layout builds a mock
    page where title/subtitle/footer land, and it renders to a PNG;
  - marquee path: a template whose default-layout title placeholder lost its
    empty text line is AUTO-FIXED in the build copy (Stage 3B), so the self-test
    then confirms the title lands — end-to-end proof that registration makes a
    broken template good.

Uses the layout-diverse fixture + LibreOffice (already a Slide Lab dependency).

Run:  py -3 slide-builder/tests/run_registration_selftest_smoke.py
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import json
import shutil
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

import run_layout_inheritance_smoke as fixture
import _paths as _p
import register_template as rt
from pptx import Presentation
from pptx.oxml.ns import qn

LAYOUT = "body_canonical_light"
_TITLE_TYPES = {1, 13}


def _set_default_layout(tpl: Path, name: str) -> None:
    tj = _p.theme_json(tpl)
    data = json.loads(tj.read_text(encoding="utf-8"))
    data["default_content_layout"] = name
    tj.write_text(json.dumps(data, indent=2), encoding="utf-8")


def _break_title_paragraphs(tpl: Path, layout_name: str) -> None:
    """Strip every <a:p> from the named layout's title placeholder in `tpl`."""
    prs = Presentation(str(tpl))
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if (layout.name or "").strip() != layout_name:
                continue
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if t is not None and int(t) in _TITLE_TYPES:
                    tb = ph._element.find(qn("p:txBody"))
                    if tb is not None:
                        for p in tb.findall(qn("a:p")):
                            tb.remove(p)
    prs.save(str(tpl))


def _register(tmp: Path) -> Path:
    tpl = tmp / fixture.FIXTURE_PPTX.name
    shutil.copy2(fixture.FIXTURE_PPTX, tpl)
    return tpl


def main() -> int:
    # --- Phase 1: happy path ---
    tmp1 = Path(tempfile.mkdtemp(prefix="selftest_smoke_ok_"))
    try:
        print("[1] Happy path: title/subtitle/footer land + render")
        tpl = _register(tmp1)
        fixture.register_fixture(tpl)          # builds + normalizes the build copy
        _set_default_layout(tpl, LAYOUT)
        fails, infos = rt._render_mock_page_selftest(tpl)
        assert fails == [], f"unexpected failures on a clean template: {fails}"
        mock = _p.selftest_pptx(tpl)
        assert mock.exists(), f"mock .pptx not saved for review: {mock}"
        mprs = Presentation(str(mock))
        assert len(mprs.slides._sldIdLst) == 1, "mock .pptx should have exactly one slide"
        print("    ok: landed; mock .pptx saved + openable; infos: " + " | ".join(infos))
    finally:
        shutil.rmtree(tmp1, ignore_errors=True)

    # --- Phase 2: broken template is auto-fixed, then self-test passes ---
    tmp2 = Path(tempfile.mkdtemp(prefix="selftest_smoke_fix_"))
    try:
        print("[2] Broken title placeholder -> auto-fixed in copy -> lands")
        tpl = _register(tmp2)
        _break_title_paragraphs(tpl, LAYOUT)   # defect in the ORIGINAL
        fixture.register_fixture(tpl)          # 3B repairs the COPY
        _set_default_layout(tpl, LAYOUT)
        # sanity: the original is still broken
        prsO = Presentation(str(tpl))
        for m in prsO.slide_masters:
            for lay in m.slide_layouts:
                if (lay.name or "").strip() == LAYOUT:
                    for ph in lay.placeholders:
                        t = ph.placeholder_format.type
                        if t is not None and int(t) in _TITLE_TYPES:
                            tb = ph._element.find(qn("p:txBody"))
                            assert tb is None or tb.find(qn("a:p")) is None, \
                                "setup: original title should still be broken"
        fails, infos = rt._render_mock_page_selftest(tpl)
        assert fails == [], f"auto-fix did not make the title land: {fails}"
        assert _p.selftest_pptx(tpl).exists(), "mock .pptx not saved for review"
        print("    ok: broken template auto-fixed; title lands on the copy")
    finally:
        shutil.rmtree(tmp2, ignore_errors=True)

    print("\nSMOKE PASSED.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
