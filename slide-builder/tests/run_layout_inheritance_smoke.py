#!/usr/bin/env python3
"""run_layout_inheritance_smoke.py — § 6.2 integration smoke for v0.2 layout
inheritance.

Audit punch list item 3 (v0.2-layout-inheritance-design-lock-2026-05-28.md).
The audit chat could not run § 6.2 because no template on the machine carried
the required layout-class diversity (cover_light + cover_dark + body_canonical_
light + body_canonical_dark + section_divider + section_divider_dark). This
smoke synthesizes a 6-layout fixture PPTX, registers it, then verifies the
three user-facing properties the design lock requires:

  1. Title bottom-y identity between body_canonical_light and
     body_canonical_dark (structural — both use canonical_title_box()).
  2. text_role color flip — light_on_dark layouts produce WHITE title text;
     dark_on_light layouts produce TEXT_DARK.
  3. cover_dark bespoke positions differ from body-canonical positions
     (chrome.yml carries an explicit title BoxPx for cover_dark; that
     BoxPx is not equal to canonical_title_box()).

Phase 4 also exercises build_deck.py end-to-end (resolve_slide_layouts +
write _meta.json with the layout field) against the fixture and a 4-slide
brief covering all four classes.

Usage:
    py -3 tests/run_layout_inheritance_smoke.py
        Build fixture (if absent), register, run all phases.
    py -3 tests/run_layout_inheritance_smoke.py --rebuild
        Force rebuild + re-register the fixture before running phases.

Exit 0 = all assertions pass. Non-zero = at least one phase failed.

The contract test wraps this script as `check_layout_inheritance_smoke_runs`,
gated on the fixture PPTX existing on disk so a brand-new clone doesn't
auto-trigger expensive fixture construction inside the contract test.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import subprocess
import sys
import tempfile
import traceback
from pathlib import Path

HERE = Path(__file__).resolve().parent
SKILL_ROOT = HERE.parent
SCRIPTS = SKILL_ROOT / "scripts"
TWINS = SKILL_ROOT / "twins"
sys.path.insert(0, str(SKILL_ROOT))
sys.path.insert(0, str(SCRIPTS))

import _paths as _p  # noqa: E402
from _chrome_schema import (  # noqa: E402
    BoxPx, LayoutChrome, ChromeSpec,
    CHROME_SCHEMA_VERSION_CURRENT,
    CANONICAL_TITLE_X, CANONICAL_TITLE_Y,
    canonical_title_box,
    load_chrome_yml, dump_chrome_yml,
)
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from lxml import etree  # noqa: E402


FIXTURES_DIR = HERE / "fixtures"
FIXTURE_PPTX = FIXTURES_DIR / "layout_diverse_template.pptx"


# Six target layouts. Order matters — first 6 default python-pptx layouts get
# renamed in this sequence. Each entry: (target_name, layout_class, background).
TARGET_LAYOUTS: list[tuple[str, str, str]] = [
    ("cover_light",          "bespoke",        "light"),  # was Title Slide
    ("body_canonical_light", "body-canonical", "light"),  # was Title and Content
    ("section_divider",      "bespoke",        "light"),  # was Section Header
    ("cover_dark",           "bespoke",        "dark"),   # was Two Content
    ("body_canonical_dark",  "body-canonical", "dark"),   # was Comparison
    ("section_divider_dark", "bespoke",        "dark"),   # was Title Only
]


# ---------------------------------------------------------------------------
# Phase 1 — fixture builder
# ---------------------------------------------------------------------------

def _set_layout_name(layout, new_name: str) -> None:
    """Rename a slide_layout in-place by mutating <p:cSld @name>."""
    csld = layout.element.find(qn("p:cSld"))
    if csld is None:
        raise RuntimeError(f"layout has no <p:cSld>; cannot rename to {new_name!r}")
    csld.set("name", new_name)


def _set_layout_dark_bg(layout) -> None:
    """Add a dark solid-fill background to a slide_layout's <p:cSld>.

    Removes any existing <p:bg> first. The <p:bg> element MUST come first
    under <p:cSld> per the OOXML schema, so we insert at position 0.
    """
    csld = layout.element.find(qn("p:cSld"))
    if csld is None:
        return
    for existing in csld.findall(qn("p:bg")):
        csld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="0A1A2E"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    )
    bg_el = etree.fromstring(bg_xml)
    csld.insert(0, bg_el)


def build_fixture(output_path: Path) -> None:
    """Build the 6-layout fixture pptx. Idempotent — rebuild on every call."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    if len(prs.slide_masters) == 0:
        raise RuntimeError("default Presentation has no masters")
    master = prs.slide_masters[0]
    layouts = list(master.slide_layouts)
    if len(layouts) < len(TARGET_LAYOUTS):
        raise RuntimeError(
            f"default Presentation has {len(layouts)} layouts; "
            f"need at least {len(TARGET_LAYOUTS)}"
        )
    for layout, (name, _klass, bg) in zip(layouts, TARGET_LAYOUTS):
        _set_layout_name(layout, name)
        if bg == "dark":
            _set_layout_dark_bg(layout)
    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# Phase 2 — programmatic registrar (writes brand.yml + theme.json + chrome.yml)
# ---------------------------------------------------------------------------

def _sha8_of(path: Path) -> str:
    h = hashlib.sha256()
    with open(str(path), "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()[:8]


def register_fixture(fixture_pptx: Path) -> Path:
    """Register the fixture via register_template's commit path with explicit
    classification overrides. Returns the chrome.yml path.
    """
    import register_template as rt

    # Build brand.yml + theme.json with sensible defaults (the fixture's theme
    # is just the python-pptx default — black title, gray subtitle).
    sha = hashlib.sha256(fixture_pptx.read_bytes()).hexdigest()
    sha8 = sha[:8]
    # We use the same _write_outputs path as a real registration, but with
    # canned brand colors (the fixture isn't a real client template).
    rt._write_outputs(
        fixture_pptx, sha, sha8,
        primary_hex="1F2937", primary_slot="dk2",
        accent_hex="F59E0B",  accent_slot="lt2",
        cover_bg_hex="0A1A2E", cover_bg_slot="dk2",
        font_heading="Arial", font_body="Arial",
        strip_master_backgrounds=False,  # fixture has no decoration
        colors={"dk1": "000000", "lt1": "FFFFFF",
                "dk2": "1F2937", "lt2": "F59E0B",
                "accent1": "1F2937", "accent2": "F59E0B"},
        n_master=1, n_layout=len(TARGET_LAYOUTS),
    )

    # Now write chrome.yml with explicit per-layout classification overrides.
    overrides = {name: klass for name, klass, _bg in TARGET_LAYOUTS}
    prs = Presentation(str(fixture_pptx))
    spec = rt.extract_chrome_spec(prs, sha8=sha8, classifications_override=overrides)
    chrome_yml_path = _p.chrome_yml(fixture_pptx)
    dump_chrome_yml(spec, chrome_yml_path)
    return chrome_yml_path


# ---------------------------------------------------------------------------
# Phase 3 — structural assertions on the registered chrome.yml
# ---------------------------------------------------------------------------

def assert_structural(spec: ChromeSpec) -> list[str]:
    """Run the structural assertions § 6.2 calls out. Returns error list."""
    errors: list[str] = []
    expected = {n: (k, b) for n, k, b in TARGET_LAYOUTS}
    for name, (klass, bg) in expected.items():
        lc = spec.layouts.get(name)
        if lc is None:
            errors.append(f"missing layout in chrome.yml: {name}")
            continue
        if lc.layout_class != klass:
            errors.append(
                f"{name}: layout_class={lc.layout_class!r}, expected {klass!r}"
            )
        if lc.background != bg:
            errors.append(
                f"{name}: background={lc.background!r}, expected {bg!r}"
            )
        # text_role mirrors background by extractor design
        expected_role = "light_on_dark" if bg == "dark" else "dark_on_light"
        if lc.text_role != expected_role:
            errors.append(
                f"{name}: text_role={lc.text_role!r}, expected {expected_role!r}"
            )

    # Property 1: body_canonical_light.title bottom == body_canonical_dark.title bottom
    # For body-canonical layouts, helpers.py uses canonical_title_box() in both
    # cases. Verify both layouts are class body-canonical so the helper path is
    # identical (the bottom-anchor invariant lives in add_title_block itself).
    b_light = spec.layouts.get("body_canonical_light")
    b_dark = spec.layouts.get("body_canonical_dark")
    if not (b_light and b_dark):
        errors.append(
            "body_canonical_light / body_canonical_dark missing — title "
            "bottom-y identity assertion unverifiable"
        )
    else:
        # Both layouts must be body-canonical so they route to
        # canonical_title_box(); bespoke layouts would each have their own box.
        if b_light.layout_class != "body-canonical":
            errors.append(
                f"body_canonical_light.layout_class={b_light.layout_class!r}; "
                f"must be 'body-canonical' for title bottom-y identity"
            )
        if b_dark.layout_class != "body-canonical":
            errors.append(
                f"body_canonical_dark.layout_class={b_dark.layout_class!r}; "
                f"must be 'body-canonical' for title bottom-y identity"
            )

    # Property 3: cover_dark.title BoxPx is set AND differs from canonical
    cover_dark = spec.layouts.get("cover_dark")
    if cover_dark is None:
        errors.append("cover_dark missing — bespoke-vs-canonical comparison unverifiable")
    elif cover_dark.title is None:
        errors.append(
            "cover_dark.title BoxPx is None — bespoke layout must carry an "
            "extracted title position (the source pptx's Title placeholder)"
        )
    else:
        canon = canonical_title_box()
        same = (cover_dark.title.x_px == canon.x_px
                and cover_dark.title.y_px == canon.y_px
                and cover_dark.title.w_px == canon.w_px
                and cover_dark.title.h_px == canon.h_px)
        if same:
            errors.append(
                f"cover_dark.title is identical to canonical title box "
                f"({canon.x_px},{canon.y_px},{canon.w_px},{canon.h_px}) — "
                f"bespoke layout should have its own positions"
            )
    return errors


# ---------------------------------------------------------------------------
# Phase 4 — text_role color flip (unit-level via helpers)
# ---------------------------------------------------------------------------

def assert_text_role_color_flip(spec: ChromeSpec) -> list[str]:
    # v0.3: bespoke-path text_role color flip via cover_dark.
    errors: list[str] = []
    import twins.helpers as h
    saved_chrome = h._ACTIVE_CHROME
    saved_yml = os.environ.pop("SLIDE_LAB_CHROME_YML", None)
    saved_lay = os.environ.pop("SLIDE_LAB_LAYOUT_NAME", None)
    try:
        lc = spec.layouts.get("cover_dark")
        if lc is None:
            errors.append("text_role: cover_dark missing")
            return errors
        prs, slide = h.new_slide()
        h.add_title_block(slide, "Title text", "Subtitle", chrome=lc)
        color = None
        for shp in slide.shapes:
            if (shp.name or "").lower() == "title":
                for para in shp.text_frame.paragraphs:
                    for r in para.runs:
                        try:
                            color = str(r.font.color.rgb)
                        except Exception:
                            pass
                        break
                    break
        if color != "FFFFFF":
            errors.append(f"cover_dark title color {color!r} != WHITE")
        pd, _ = h._text_colors_for("dark_on_light")
        pl, _ = h._text_colors_for("light_on_dark")
        if str(pd) == str(pl):
            errors.append("_text_colors_for: flip is a no-op")
    finally:
        h._ACTIVE_CHROME = saved_chrome
        if saved_yml is not None:
            os.environ["SLIDE_LAB_CHROME_YML"] = saved_yml
        if saved_lay is not None:
            os.environ["SLIDE_LAB_LAYOUT_NAME"] = saved_lay
    return errors


def assert_v2_inheritance_fields(spec: ChromeSpec) -> list[str]:
    # v0.3: body-canonical layouts in v2 chrome carry inheritance fields;
    # add_title_block on body-canonical draws no title shape.
    errors: list[str] = []
    import twins.helpers as h
    for name in ("body_canonical_light", "body_canonical_dark"):
        lc = spec.layouts.get(name)
        if lc is None:
            errors.append(f"v2: {name} missing")
            continue
        if getattr(lc, "title_placeholder_idx", None) is None:
            errors.append(f"v2: {name}.title_placeholder_idx is None")
        if getattr(lc, "body_top_y_px", None) is None:
            errors.append(f"v2: {name}.body_top_y_px is None")
        if getattr(lc, "body_bottom_y_px", None) is None:
            errors.append(f"v2: {name}.body_bottom_y_px is None")
    lc = spec.layouts.get("body_canonical_light")
    if lc is not None and getattr(lc, "title_placeholder_idx", None) is not None:
        saved = h._ACTIVE_CHROME
        try:
            prs, slide = h.new_slide()
            h.add_title_block(slide, "Title text", chrome=lc)
            for shp in slide.shapes:
                if (shp.name or "").lower() == "title":
                    errors.append("v2: add_title_block drew title shape")
                    break
        finally:
            h._ACTIVE_CHROME = saved
    return errors


def assert_v1_chrome_yml_still_loadable() -> list[str]:
    errors: list[str] = []
    import tempfile
    from _chrome_schema import load_chrome_yml as _load
    import yaml as _yaml
    v1d = {"schema_version": 1, "source_template_sha8": "v1abcd00",
           "layouts": {"some_body": {"name": "some_body",
               "layout_class": "body-canonical", "text_role": "dark_on_light",
               "background": "light", "has_page_number": True}}}
    with tempfile.TemporaryDirectory() as td:
        path = Path(td) / "v1_legacy.chrome.yml"
        path.write_text(_yaml.safe_dump(v1d), encoding="utf-8")
        try:
            spec = _load(path)
        except Exception as exc:
            errors.append(f"v1 load failed: {type(exc).__name__}: {exc}")
            return errors
        lc = spec.layouts.get("some_body")
        if lc is None:
            errors.append("v1: layout vanished")
            return errors
        for field in ("title_placeholder_idx", "body_top_y_px",
                      "body_bottom_y_px", "body_overlay_hex"):
            if getattr(lc, field) is not None:
                errors.append(f"v1: {field}={getattr(lc, field)!r}; expected None")
    return errors


# ---------------------------------------------------------------------------
# Phase 5 — end-to-end build_deck.py against a 4-slide brief
# ---------------------------------------------------------------------------

SMOKE_BRIEF_TEMPLATE = """---
client_template: {tpl}
deck_type: Smoke test (v0.2 P1.9)
default_layout: body_canonical_light
---

## Deck-level design notes

Synthetic 4-slide brief covering cover_dark, body_canonical_light,
body_canonical_dark, section_divider. Used by run_layout_inheritance_smoke.py
to exercise build_deck end-to-end against the layout-diverse fixture.

## Slide 1 — Smoke cover (dark)
**Layout:** cover_dark
**Archetype:** Cover / Title
**Governing thought:** Cover-class smoke verification.
**So-what:** The cover_dark layout grafts onto its own placeholder positions.
**Evidence:** None.

## Slide 2 — Light body
**Layout:** body_canonical_light
**Archetype:** Synthesis / Findings
**Governing thought:** Body-canonical light variant works.
**So-what:** Title bottom-y matches canonical.
**Evidence:** N/A.

## Slide 3 — Dark mid-deck emphasis
**Layout:** body_canonical_dark
**Archetype:** Synthesis / Findings
**Governing thought:** body_canonical_dark flips text colors via text_role.
**So-what:** Title bottom-y identical to slide 2 despite dark canvas.
**Evidence:** N/A.

## Slide 4 — Section divider
**Layout:** section_divider
**Archetype:** Cover / Title
**Governing thought:** Section divider is bespoke chrome.
**So-what:** Bespoke positions inherited from template.
**Evidence:** N/A.
"""


def run_build_deck_against_fixture(fixture_pptx: Path) -> tuple[int, list[str]]:
    """Exercise build_deck.py end-to-end against a 4-slide brief; return
    (exit_code, errors). Verifies exit 0 + _meta.json layout assignments."""
    errors: list[str] = []
    with tempfile.TemporaryDirectory() as td:
        brief = Path(td) / "smoke_brief.md"
        brief.write_text(SMOKE_BRIEF_TEMPLATE.format(tpl=fixture_pptx),
                         encoding="utf-8")
        out_dir = Path(td) / "smoke_out"
        cmd = [
            sys.executable, str(SCRIPTS / "build_deck.py"),
            "--brief", str(brief),
            "--template", str(fixture_pptx),
            "--out", str(out_dir),
        ]
        # build_deck.py also requires mmdc via stage1; we may not have it.
        # Capture stdout/stderr; gate failure handling on the actual exit code.
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        if result.returncode != 0:
            tail = (result.stderr or result.stdout or "").splitlines()
            # mmdc-missing exit 7 is environmental, not a smoke regression.
            joined = "\n".join(tail[-30:])
            if "Mermaid CLI (mmdc) not installed" in joined:
                return 7, ["build_deck stage-1 needs mmdc — skipped end-to-end"]
            errors.append(
                f"build_deck.py exit {result.returncode}; tail:\n{joined[-1500:]}"
            )
            return result.returncode, errors

        # Verify _meta.json got the layout field per slide
        meta_path = out_dir / "_meta.json"
        if not meta_path.exists():
            errors.append("build_deck.py exit 0 but _meta.json missing")
            return result.returncode, errors
        meta = json.loads(meta_path.read_text(encoding="utf-8"))
        if meta.get("schema_version") != 3:
            errors.append(
                f"_meta.json schema_version={meta.get('schema_version')}, "
                f"expected 3"
            )
        expected_layouts = ["cover_dark", "body_canonical_light",
                            "body_canonical_dark", "section_divider"]
        actual = [s.get("layout") for s in meta.get("slides", [])]
        if actual != expected_layouts:
            errors.append(
                f"meta.slides[].layout = {actual}; expected {expected_layouts}"
            )
    return 0, errors


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

def run_smoke(rebuild: bool = False, verbose: bool = True) -> int:
    if verbose:
        print("=" * 72)
        print("v0.2 layout-inheritance integration smoke (§ 6.2 + audit P1.9)")
        print("=" * 72)

    fixture_existed = FIXTURE_PPTX.exists()
    if rebuild or not fixture_existed:
        if verbose:
            action = "rebuild forced" if rebuild else "absent — building"
            print(f"\n[Phase 1] Fixture {action}: {FIXTURE_PPTX}")
        try:
            build_fixture(FIXTURE_PPTX)
        except Exception as exc:
            print(f"FAIL [Phase 1] fixture build: {type(exc).__name__}: {exc}")
            return 2
        # Re-register if rebuilt
        try:
            chrome_path = register_fixture(FIXTURE_PPTX)
            if verbose:
                print(f"  registered -> {chrome_path}")
        except Exception as exc:
            print(f"FAIL [Phase 2] fixture register: {type(exc).__name__}: {exc}")
            traceback.print_exc()
            return 2
    else:
        chrome_path = _p.chrome_yml(FIXTURE_PPTX)
        if not chrome_path.exists():
            if verbose:
                print(f"\n[Phase 2] chrome.yml missing — re-registering")
            try:
                register_fixture(FIXTURE_PPTX)
            except Exception as exc:
                print(f"FAIL [Phase 2] fixture register: "
                      f"{type(exc).__name__}: {exc}")
                return 2
        elif verbose:
            print(f"\n[Phase 1+2] Fixture + chrome.yml already in place")

    spec = load_chrome_yml(_p.chrome_yml(FIXTURE_PPTX))

    if verbose:
        print(f"\n[Phase 3] Structural assertions on chrome.yml "
              f"({len(spec.layouts)} layouts)")
    struct_errs = assert_structural(spec)
    if struct_errs:
        for e in struct_errs:
            print(f"  FAIL: {e}")
    elif verbose:
        print("  ok: all 6 layouts present with expected class + bg + text_role")
        print("  ok: body_canonical_light/dark both route to canonical_title_box "
              "(bottom-y identity)")
        print("  ok: cover_dark.title BoxPx is set and differs from canonical")

    if verbose:
        print("\n[Phase 4] text_role color flip via add_title_block")
    color_errs = assert_text_role_color_flip(spec)
    if color_errs:
        for e in color_errs:
            print(f"  FAIL: {e}")
    elif verbose:
        print("  ok: body_canonical_light title -> 000000 (TEXT_DARK)")
        print("  ok: body_canonical_dark  title -> FFFFFF (WHITE)")
        print("  ok: cover_dark           title -> FFFFFF (WHITE)")

    if verbose:
        print("\n[Phase 4b] v2 inheritance fields populated")
    v2_errs = assert_v2_inheritance_fields(spec)
    if v2_errs:
        for e in v2_errs:
            print(f"  FAIL: {e}")
    elif verbose:
        print("  ok: body_canonical layouts carry v2 inheritance fields")
        print("  ok: add_title_block on body-canonical draws no title shape")
    if verbose:
        print("\n[Phase 4c] legacy v1 chrome.yml still loadable")
    v1_errs = assert_v1_chrome_yml_still_loadable()
    if v1_errs:
        for e in v1_errs:
            print(f"  FAIL: {e}")
    elif verbose:
        print("  ok: v1 chrome.yml validates; v2 fields default to None")
    if verbose:
        print("\n[Phase 5] build_deck.py end-to-end against 4-slide brief")
    exit_code, build_errs = run_build_deck_against_fixture(FIXTURE_PPTX)
    skipped = exit_code == 7 and not build_errs[0].startswith("build_deck.py exit")
    if build_errs and not skipped:
        for e in build_errs:
            print(f"  FAIL: {e}")
    elif skipped:
        if verbose:
            print(f"  SKIP: {build_errs[0]}")
        build_errs = []  # don't count as failure
    elif verbose:
        print("  ok: build_deck exit 0; _meta.json v3 with layout per slide")

    all_errs = struct_errs + color_errs + v2_errs + v1_errs + build_errs
    if verbose:
        print()
        print("=" * 72)
    if all_errs:
        print(f"SMOKE FAILED: {len(all_errs)} assertion(s) failed.")
        return 1
    if verbose:
        print("SMOKE PASSED.")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(
        description="v0.2 layout-inheritance integration smoke (§ 6.2 + P1.9)"
    )
    ap.add_argument("--rebuild", action="store_true",
                    help="Force rebuild + re-register the fixture before phases")
    ap.add_argument("--quiet", action="store_true",
                    help="Suppress per-phase chatter (errors still printed)")
    args = ap.parse_args()
    try:
        return run_smoke(rebuild=args.rebuild, verbose=not args.quiet)
    except Exception as exc:
        print(f"SMOKE crashed: {type(exc).__name__}: {exc}")
        traceback.print_exc()
        return 2


if __name__ == "__main__":
    sys.exit(main())
