#!/usr/bin/env python3
"""Smoke test for compile_picks.py --splice-into (option 6b — external-deck redesign).

Builds a 3-slide "external" deck + a fake adopt out-dir holding one rebuilt slide
(slide 2) and an `adopted_source` _meta.json, then:
  - splices the rebuilt slide into position 2 and asserts the output has the SAME
    3 slides with slide 2 replaced and slides 1 & 3 untouched (the whole point:
    a splice preserves the deck, unlike the clear-and-rebuild compile);
  - asserts the original file is untouched;
  - asserts the adopted-deck GUARD fires: a plain compile (no --splice-into) on
    an adopted out-dir refuses with exit 2 instead of wiping the other slides.

Run:  py -3 slide-builder/tests/run_splice_smoke.py   (python3 on macOS/Linux)
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

from pptx import Presentation
from pptx.util import Inches


def _titled_slide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb.text_frame.text = text
    return slide


def _make_external(path: Path):
    prs = Presentation()
    for t in ["External slide 1", "External slide 2 (OLD)", "External slide 3"]:
        _titled_slide(prs, t)
    prs.save(str(path))


def _make_rebuilt(path: Path):
    prs = Presentation()
    _titled_slide(prs, "REBUILT slide 2 (NEW)")
    prs.save(str(path))


def _slide_text(slide) -> str:
    return " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def main() -> int:
    tmp = Path(tempfile.mkdtemp(prefix="splice_smoke_"))
    try:
        original = tmp / "client_pmo.pptx"
        _make_external(original)

        out_dir = tmp / "adopt_out"
        (out_dir / "slide_02").mkdir(parents=True)
        _make_rebuilt(out_dir / "slide_02" / "option_A.pptx")
        meta = {
            "template": str(original),
            "adopted_source": str(original),
            "slide_count": 3,
            "slides": [{"n": 1, "layout": ""}, {"n": 2, "layout": ""}, {"n": 3, "layout": ""}],
        }
        (out_dir / "_meta.json").write_text(json.dumps(meta), encoding="utf-8")

        print("[1] splice rebuilt slide 2 into the external deck")
        spliced = tmp / "client_pmo_slidelab.pptx"
        r = subprocess.run(
            [sys.executable, str(SCRIPTS / "compile_picks.py"), "--out", str(out_dir),
             "--picks", json.dumps({"slide_02": "A"}), "--splice-into", str(original),
             "--final", str(spliced)],
            capture_output=True, text=True)
        assert r.returncode == 0, f"splice failed:\n{r.stdout}\n{r.stderr}"
        assert spliced.exists(), "spliced deck not written"

        print("[2] verify count preserved + slide 2 replaced + 1 & 3 intact")
        prs = Presentation(str(spliced))
        slides = list(prs.slides)
        assert len(slides) == 3, f"expected 3 slides, got {len(slides)}"
        assert "External slide 1" in _slide_text(slides[0]), _slide_text(slides[0])
        assert "REBUILT slide 2" in _slide_text(slides[1]), f"slide 2 not replaced: {_slide_text(slides[1])!r}"
        assert "OLD" not in _slide_text(slides[1]), "old slide 2 still present"
        assert "External slide 3" in _slide_text(slides[2]), _slide_text(slides[2])
        print("    ok: 3 slides; slide 2 = REBUILT; slides 1 & 3 unchanged")

        print("[3] original file untouched")
        orig = Presentation(str(original))
        assert "OLD" in _slide_text(list(orig.slides)[1]), "original was modified!"
        print("    ok: original intact")

        print("[4] guard: plain compile on an adopted deck refuses (exit 2)")
        r2 = subprocess.run(
            [sys.executable, str(SCRIPTS / "compile_picks.py"), "--out", str(out_dir),
             "--picks", json.dumps({"slide_02": "A"})],
            capture_output=True, text=True)
        assert r2.returncode == 2, f"guard should exit 2, got {r2.returncode}\n{r2.stdout}"
        assert "adopted" in r2.stdout.lower(), r2.stdout
        assert "--splice-into" in r2.stdout, r2.stdout
        print("    ok: adopted-deck guard blocks the deck-wiping compile")

        print("\nSMOKE PASSED.")
        return 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
