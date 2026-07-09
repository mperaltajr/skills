#!/usr/bin/env python3
"""Smoke test for refresh_deck.py (option 6c — content-only PMO refresh).

Builds a tiny 2-slide deck with known text, runs `spec` to dump the fields,
edits one field's new_text, runs `apply`, and asserts:
  - the copy is a NEW file (original untouched),
  - the changed field got the new text (formatting preserved: still one run),
  - an untouched field is unchanged,
  - the drift guard skips a field whose recorded current_text no longer matches.

Run:  py -3 slide-builder/tests/run_refresh_smoke.py   (python3 on macOS/Linux)
Prints "SMOKE PASSED." on success; raises AssertionError otherwise.
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
SCRIPTS = HERE.parent / "scripts"
sys.path.insert(0, str(SCRIPTS))

from pptx import Presentation
from pptx.util import Inches, Pt


def _make_deck(path: Path) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title, body in [("Status: GREEN", "12 of 15 KPIs on track"),
                        ("Budget: $4.2M", "spend to date $1.1M")]:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        bb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6), Inches(1))
        bb.text_frame.text = body
    prs.save(str(path))


def main() -> int:
    tmp = Path(tempfile.mkdtemp(prefix="refresh_smoke_"))
    try:
        deck = tmp / "pmo_deck.pptx"
        _make_deck(deck)

        print("[1] spec — dump text fields")
        r = subprocess.run([sys.executable, str(SCRIPTS / "refresh_deck.py"), "spec", str(deck)],
                           capture_output=True, text=True)
        assert r.returncode == 0, r.stderr
        spec_path = deck.with_name("pmo_deck_refresh_spec.json")
        assert spec_path.exists(), "spec file not written"
        spec = json.loads(spec_path.read_text(encoding="utf-8"))
        # 2 slides x 2 text shapes = 4 fields
        assert len(spec["fields"]) == 4, spec["fields"]
        print(f"    ok: {len(spec['fields'])} fields dumped")

        print("[2] edit one field + one deliberately-drifted field, then apply")
        # change the real "Status: GREEN" title
        f_status = next(f for f in spec["fields"] if f["current_text"] == "Status: GREEN")
        f_status["new_text"] = "Status: AMBER"
        # a drift case: point at a real field but claim wrong current_text
        f_budget = next(f for f in spec["fields"] if f["current_text"] == "Budget: $4.2M")
        f_budget["current_text"] = "Budget: $9.9M"   # no longer matches the deck
        f_budget["new_text"] = "Budget: $5.0M"        # should be SKIPPED by drift guard
        spec_path.write_text(json.dumps(spec, indent=2), encoding="utf-8")

        r = subprocess.run([sys.executable, str(SCRIPTS / "refresh_deck.py"), "apply",
                            str(deck), str(spec_path)], capture_output=True, text=True)
        assert r.returncode == 0, r.stderr + r.stdout
        # find the dated copy
        copies = [p for p in tmp.glob("pmo_deck_refreshed_*.pptx")]
        assert len(copies) == 1, f"expected one dated copy, got {copies}"
        out = copies[0]
        assert out != deck, "must not overwrite original"
        print(f"    ok: wrote {out.name}")

        print("[3] verify change applied, format preserved, drift skipped, original untouched")
        prs_out = Presentation(str(out))
        slides = list(prs_out.slides)
        s0_title = list(slides[0].shapes)[0].text_frame
        assert s0_title.text == "Status: AMBER", s0_title.text
        assert len(s0_title.paragraphs[0].runs) == 1, "should collapse to one run"
        assert s0_title.paragraphs[0].runs[0].font.size == Pt(28), "font size must be preserved"
        s1_title = list(slides[1].shapes)[0].text_frame
        assert s1_title.text == "Budget: $4.2M", f"drifted field must be UNCHANGED, got {s1_title.text!r}"
        # original untouched
        prs_orig = Presentation(str(deck))
        assert list(list(prs_orig.slides)[0].shapes)[0].text_frame.text == "Status: GREEN", "original changed!"
        assert "1 skipped (text drifted)" in r.stdout, r.stdout
        print("    ok: change applied + font kept; drift skipped; original intact")

        print("\nSMOKE PASSED.")
        return 0
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    sys.exit(main())
