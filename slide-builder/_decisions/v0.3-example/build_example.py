"""Build a 3-slide OTC example with realistic content + light/dark/override variants.

Demonstrates v0.3 against Mario's freshly committed chrome.yml:
- Slide 1: 'Top 3 findings' — light body, default 'Use as default slide template'
- Slide 2: 'Where the cash gets stuck' — dark via per-slide override to 'Layout 01'
          (body_overlay_hex='0A1A2E', text flips to white via text_role=light_on_dark)
- Slide 3: '12-week fix path' — light body via per-slide override to '1_Layout 19'

Run from worktree root:
    py -3 slide-builder/_decisions/v0.3-example/build_example.py
"""
import sys
from pathlib import Path
from copy import deepcopy

HERE = Path(__file__).resolve().parent  # _decisions/v0.3-example
SKILL = HERE.parents[1]                 # slide-builder/
SCRIPTS = SKILL / "scripts"
sys.path.insert(0, str(SKILL))
sys.path.insert(0, str(SCRIPTS))

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from twins.composer import (
    _find_named_layout, _find_blank_layout,
    _clear_existing_slides, _populate_layout_placeholders,
    _insert_dark_overlay, _strip_layout_placeholders,
)

TEMPLATE = Path(r"C:/Users/m.a.peralta/OneDrive - Accenture/Library/FedEx/OTC/OTC Opportunity.pptx")
OUT_DIR = HERE / "build"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PPTX = OUT_DIR / "otc-v0.3-example-real-content.pptx"

# Read brand colors from the template's brand.yml (single source of truth).
# Falls back to hardcoded FedEx hexes if the sidecar isn't readable for some
# reason — example will still build but warn.
def _load_brand_hexes(template_path: Path) -> dict:
    try:
        from twins.client_theme import load_brand_sidecar
        return load_brand_sidecar(template_path)
    except Exception as e:
        print(f"  WARN: could not load brand.yml ({e}); using FedEx defaults")
        return {
            "primary_hex": "4D148C", "accent_hex": "FF6600",
            "cover_bg_hex": "4D148C", "dark_bg_hex": "4D148C",
        }

_BRAND = _load_brand_hexes(TEMPLATE)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    s = (hex_str or "").lstrip("#")
    if len(s) != 6:
        s = "000000"
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


PURPLE     = _hex_to_rgb(_BRAND["primary_hex"])       # brand primary
ORANGE     = _hex_to_rgb(_BRAND["accent_hex"])        # brand accent
DARK_BG    = _hex_to_rgb(_BRAND["dark_bg_hex"])       # dark variant ground
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
LIGHT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_DARK = RGBColor(0x88, 0x88, 0x88)
MUTED_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)

EMU_PER_PX = 9525  # at 96 DPI

def px_to_emu(px): return int(px * EMU_PER_PX)


def add_text_box(slide, *, x_px, y_px, w_px, h_px, text, font_size_pt=14,
                  color=DARK_TEXT, bold=False, anchor="top"):
    box = slide.shapes.add_textbox(
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(w_px), px_to_emu(h_px),
    )
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    if anchor == "middle":
        from pptx.enum.text import MSO_ANCHOR
        box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def add_rect(slide, *, x_px, y_px, w_px, h_px, fill_rgb, no_line=True):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu(x_px), px_to_emu(y_px),
        px_to_emu(w_px), px_to_emu(h_px),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_rgb
    if no_line:
        rect.line.fill.background()
    return rect


def build_slide_1(prs, layout):
    """Light body, top 3 findings as numbered blocks."""
    slide = prs.slides.add_slide(layout)
    # Title goes into the layout's inherited title placeholder
    _populate_layout_placeholders(
        slide,
        title="Top 3 findings — AR & collections gaps",
        page_num="1",
    )
    # Body content: 3 finding blocks in the body zone (~y=100 to y=660)
    findings = [
        ("01", "Cash application gap",
         "27% of incoming wires sit unmatched > 5 days. Hand-keyed remittance "
         "data is the bottleneck — no auto-match on customer + invoice."),
        ("02", "Disputed invoice volume",
         "Disputes run at 4.8% of invoiced revenue, vs. industry 1.5%. "
         "Disputes age in email inboxes; no shared queue or SLA."),
        ("03", "Credit-memo leakage",
         "$3.2M in credit memos issued in Q3 had no PO-level audit trail. "
         "Approval thresholds work; reconciliation back to AR doesn't."),
    ]
    block_top = 130
    block_h = 150
    block_gap = 12
    for i, (num, headline, body) in enumerate(findings):
        y = block_top + i * (block_h + block_gap)
        # Accent number bar on left
        add_rect(slide, x_px=58, y_px=y, w_px=60, h_px=block_h, fill_rgb=PURPLE)
        add_text_box(slide, x_px=64, y_px=y + 18, w_px=50, h_px=40,
                     text=num, font_size_pt=28, color=LIGHT_TEXT, bold=True)
        # Headline
        add_text_box(slide, x_px=140, y_px=y + 8, w_px=1080, h_px=36,
                     text=headline, font_size_pt=18, color=PURPLE, bold=True)
        # Body
        add_text_box(slide, x_px=140, y_px=y + 50, w_px=1080, h_px=block_h - 50,
                     text=body, font_size_pt=13, color=DARK_TEXT)
    return slide


def build_slide_2(prs, layout):
    """Dark variant — SAME default body layout as slide 1.
    Strip the layout's bright chrome (gradient bar, orange bar, placeholders),
    paint full-bleed dark overlay, draw title at the layout's title-placeholder
    position so it matches slide 1 vertically.
    Per Mario's decision: dark = same layout, strip chrome, full-bleed dark.
    """
    slide = prs.slides.add_slide(layout)
    # Strip everything inherited from the layout (decorative shapes + placeholders).
    # showMasterSp=0 is set inside; master decoration won't bleed either.
    _strip_layout_placeholders(slide)

    # Full-bleed dark overlay (covers the whole slide)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK_BG  # from brand.yml dark_bg_hex
    overlay.line.fill.background()
    overlay.name = "dark-variant-overlay"

    # Title at the same coords as slide 1's inherited title placeholder
    # (Use as default slide template: title at 0.42", 0.42", 12.50", 0.45")
    add_text_box(slide, x_px=40, y_px=40, w_px=1200, h_px=44,
                 text="Where the cash gets stuck — the funnel",
                 font_size_pt=24, color=LIGHT_TEXT, bold=False)

    # Funnel: 4 bars. Values placed OUTSIDE bars (right of bar end) so they
    # don't collide with the label on narrow bars.
    stages = [
        ("Invoiced this quarter",   "$248M", 1000),
        ("Collected on time",       "$181M", 730),
        ("Aged 30+ days",           "$44M",  450),
        ("Disputed / unresolved",   "$23M",  250),
    ]
    bar_top = 160
    bar_h = 56
    bar_gap = 18
    # On the FedEx purple background: top 2 bars in ORANGE (what came in),
    # bottom 2 bars in WHITE (what's stuck — the leakage). All values + label
    # text outside in WHITE so they're readable on the purple ground.
    for i, (label, value, w_px) in enumerate(stages):
        y = bar_top + i * (bar_h + bar_gap)
        if i < 2:
            bar_color = ORANGE
            label_color = LIGHT_TEXT
        else:
            bar_color = LIGHT_TEXT  # white bars for the leakage
            label_color = PURPLE     # purple text inside white bar
        add_rect(slide, x_px=80, y_px=y, w_px=w_px, h_px=bar_h, fill_rgb=bar_color)
        # Label inside bar (left)
        add_text_box(slide, x_px=96, y_px=y + 14, w_px=w_px - 32, h_px=bar_h - 8,
                     text=label, font_size_pt=15, color=label_color, bold=True)
        # Value OUTSIDE bar to the right (always white on the purple ground)
        add_text_box(slide, x_px=80 + w_px + 14, y_px=y + 14, w_px=200, h_px=bar_h - 8,
                     text=value, font_size_pt=18, color=LIGHT_TEXT, bold=True)

    add_text_box(slide, x_px=80, y_px=480, w_px=1120, h_px=80,
                 text="Bottom of funnel = 27% of invoiced revenue not converted to cash within "
                      "30 days. Each $1 of leakage flows from 3 process gaps: cash-app, dispute "
                      "ageing, credit-memo recon.",
                 font_size_pt=13, color=MUTED_LIGHT)
    return slide


def build_slide_3(prs, layout):
    """Light body via alternate layout — 4-phase roadmap."""
    slide = prs.slides.add_slide(layout)
    _populate_layout_placeholders(
        slide,
        title="Sequenced fix path — 12 weeks",
        page_num="3",
    )
    phases = [
        ("Discover", "Weeks 1-3", "Map current AR flows. Quantify leakage by customer & SKU."),
        ("Design",   "Weeks 4-6", "Cash-app auto-match rules. Dispute queue + SLA framework."),
        ("Pilot",    "Weeks 7-9", "Run new flow with 3 high-volume customers. Tune match rates."),
        ("Scale",    "Wks 10-12", "Roll to top 50 customers. Stand up monthly leakage dashboard."),
    ]
    card_top = 160
    card_h = 340
    card_w = 268
    card_x_start = 80
    card_gap = 18
    for i, (phase, weeks, body) in enumerate(phases):
        x = card_x_start + i * (card_w + card_gap)
        # Card background
        add_rect(slide, x_px=x, y_px=card_top, w_px=card_w, h_px=card_h,
                 fill_rgb=RGBColor(0xF5, 0xF0, 0xF8))
        # Header band (purple)
        add_rect(slide, x_px=x, y_px=card_top, w_px=card_w, h_px=60, fill_rgb=PURPLE)
        # Phase name in header
        add_text_box(slide, x_px=x + 16, y_px=card_top + 14, w_px=card_w - 32, h_px=36,
                     text=phase, font_size_pt=18, color=LIGHT_TEXT, bold=True)
        # Weeks subtitle
        add_text_box(slide, x_px=x + 16, y_px=card_top + 76, w_px=card_w - 32, h_px=28,
                     text=weeks, font_size_pt=12, color=ORANGE, bold=True)
        # Body
        add_text_box(slide, x_px=x + 16, y_px=card_top + 110, w_px=card_w - 32, h_px=card_h - 130,
                     text=body, font_size_pt=12, color=DARK_TEXT)
    return slide


def main():
    if not TEMPLATE.exists():
        print(f"ERROR: template not found: {TEMPLATE}")
        return 1

    prs = Presentation(str(TEMPLATE))
    _clear_existing_slides(prs)

    # Find layouts. All three slides use the SAME default body layout for
    # body slot consistency. Slide 3 demos per-slide override to a different
    # body-canonical layout (1_Layout 19).
    body_default = _find_named_layout(prs, "Use as default slide template")
    layout_19    = _find_named_layout(prs, "1_Layout 19")
    if not body_default:
        print("ERROR: 'Use as default slide template' not found"); return 1
    if not layout_19:
        print("ERROR: '1_Layout 19' not found"); return 1

    print(f"Building 3 slides against {TEMPLATE.name}")
    s1 = build_slide_1(prs, body_default)
    print(f"  Slide 1: {body_default.name!r} (light body, default)")
    s2 = build_slide_2(prs, body_default)
    print(f"  Slide 2: {body_default.name!r} (dark variant: strip chrome + full-bleed overlay)")
    s3 = build_slide_3(prs, layout_19)
    print(f"  Slide 3: {layout_19.name!r} (light body, per-slide override)")

    prs.save(str(OUT_PPTX))
    size_kb = OUT_PPTX.stat().st_size / 1024
    print(f"\nSaved: {OUT_PPTX}  ({size_kb:.1f} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
