"""Diagram 2 -- Swimlane (cross-functional, 3 lanes x 4 steps L->R)."""
from pathlib import Path
import sys
sys.path.insert(0, r"C:\Users\m.a.peralta\.claude\skills\slide-builder")
from twins.helpers import (
    new_slide, add_text, add_rect, add_title_block, add_footer,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    TEXT_DARK, TEXT_MID, TEXT_FAINT, CARD_BG, CARD_BORDER, WHITE,
    px_to_emu,
)
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt


def _box(slide, sid, x, y, w, h, fill, border=None, line_w=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h),
    )
    shape.name = sid
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _arrow(slide, sid, x, y, w, h, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h),
    )
    shape.name = sid
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def build():
    prs, slide = new_slide()

    add_title_block(
        slide,
        title="Order-to-cash <strong>cycles across three functions</strong> with two critical hand-offs",
        subtitle="Sales captures intent; Operations fulfils; Finance closes the loop on revenue",
    )

    LABEL_W = 130
    BODY_X = 64 + LABEL_W + 6
    BODY_W = 1216 - BODY_X
    LANE_TOP = 150
    LANE_H = 150
    LANE_GAP = 8
    N_COLS = 4
    COL_W = BODY_W // N_COLS

    lanes = [
        ("Sales",       BRAND_PRIMARY),
        ("Operations",  BRAND_PRIMARY_MID),
        ("Finance",     BRAND_ACCENT_SOFT),
    ]
    label_text_colors = [WHITE, WHITE, TEXT_DARK]

    for i, (name, fill) in enumerate(lanes):
        ly = LANE_TOP + i * (LANE_H + LANE_GAP)
        _box(slide, f"lane-label-{i}", 64, ly, LABEL_W, LANE_H, fill)
        add_text(
            slide, f"lane-label-txt-{i}", name,
            x_px=64, y_px=ly, w_px=LABEL_W, h_px=LANE_H,
            font_size_px=15, color=label_text_colors[i], bold=True,
            align="center", anchor="middle",
        )
        _box(slide, f"lane-body-{i}", BODY_X, ly, BODY_W, LANE_H, CARD_BG)

    steps = [
        (0, 0, "Capture order",       "Sales rep logs deal in CRM\nwith SKU mix and terms",         False),
        (1, 1, "Allocate & ship",     "Ops confirms stock, releases\nto warehouse, books carrier",  False),
        (2, 2, "Invoice customer",    "Finance generates invoice\nand applies tax / discounts",     False),
        (2, 3, "Collect & post cash", "AR clears payment, posts\nto GL, closes the order",          True),
    ]

    STEP_W = 175
    STEP_H = 100

    step_centers = {}
    for lane_idx, col, header, body, is_accent in steps:
        ly = LANE_TOP + lane_idx * (LANE_H + LANE_GAP)
        bx = BODY_X + col * COL_W + (COL_W - STEP_W) // 2
        by = ly + (LANE_H - STEP_H) // 2
        fill = BRAND_ACCENT if is_accent else WHITE
        border = BRAND_ACCENT if is_accent else CARD_BORDER
        title_color = WHITE if is_accent else TEXT_DARK
        body_color = WHITE if is_accent else TEXT_MID
        _box(slide, f"step-{lane_idx}-{col}", bx, by, STEP_W, STEP_H,
             fill=fill, border=border, line_w=1.25)
        add_text(
            slide, f"step-h-{lane_idx}-{col}", header,
            x_px=bx + 8, y_px=by + 8, w_px=STEP_W - 16, h_px=22,
            font_size_px=13, color=title_color, bold=True, align="center",
        )
        add_text(
            slide, f"step-b-{lane_idx}-{col}", body,
            x_px=bx + 8, y_px=by + 32, w_px=STEP_W - 16, h_px=STEP_H - 40,
            font_size_px=11, color=body_color, align="center",
        )
        step_centers[(lane_idx, col)] = (
            bx + STEP_W // 2, by + STEP_H // 2, by, by + STEP_H, bx, bx + STEP_W,
        )

    flow = [(0, 0), (1, 1), (2, 2), (2, 3)]
    ARROW_W = 36
    ARROW_H = 22
    for (la, ca), (lb, cb) in zip(flow, flow[1:]):
        cxa, cya, ta, ba, lxa, rxa = step_centers[(la, ca)]
        cxb, cyb, tb, bb, lxb, rxb = step_centers[(lb, cb)]
        if la != lb:
            v_x = (rxa + lxb) // 2
            _box(slide, f"vconn-{la}-{ca}",
                 v_x - 2, min(cya, cyb), 4, abs(cyb - cya),
                 fill=BRAND_PRIMARY_MID)
            _arrow(slide, f"arrow-{la}-{ca}-to-{lb}-{cb}",
                   v_x + 2, cyb - ARROW_H // 2, lxb - v_x - 4, ARROW_H,
                   BRAND_PRIMARY_MID)
            _box(slide, f"hstub-{la}-{ca}",
                 rxa, cya - 2, v_x - rxa, 4,
                 fill=BRAND_PRIMARY_MID)
        else:
            _arrow(slide, f"arrow-{la}-{ca}-to-{lb}-{cb}",
                   rxa + 4, cya - ARROW_H // 2, lxb - rxa - 8, ARROW_H,
                   BRAND_PRIMARY_MID)

    add_footer(slide, page_num=2,
               source="Process discovery workshops, FY25 Q1",
               footnote="Lane transitions between Sales -> Ops and Ops -> Finance are the two control hand-offs")
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "diagram2-swimlane.pptx"
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
