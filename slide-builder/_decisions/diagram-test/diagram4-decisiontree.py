"""Diagram 4 -- Decision tree (root -> 2 branches -> 4 leaves)."""
from pathlib import Path
import sys
sys.path.insert(0, r"C:\Users\m.a.peralta\.claude\skills\slide-builder")
from twins.helpers import (
    new_slide, add_text, add_rect, add_title_block, add_footer,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    TEXT_DARK, TEXT_MID, TEXT_FAINT, CARD_BG, CARD_BORDER, WHITE,
    px_to_emu,
)
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt


def _box(slide, sid, x, y, w, h, fill, border=None, line_w=1.0, radius=False):
    s = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        s, px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h),
    )
    shape.name = sid
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _line(slide, sid, x1, y1, x2, y2, color, width_pt=1.25):
    conn = slide.shapes.add_connector(
        1, px_to_emu(x1), px_to_emu(y1), px_to_emu(x2), px_to_emu(y2),
    )
    conn.name = sid
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn


def build():
    prs, slide = new_slide()

    add_title_block(
        slide,
        title="Market-entry decision <strong>resolves to four operating outcomes</strong> under current data",
        subtitle="Decision logic for entering the LATAM specialty foods segment, FY26 plan",
    )

    ROOT_W, ROOT_H = 340, 70
    root_cx = 640
    root_y = 140
    root_x = root_cx - ROOT_W // 2

    _box(slide, "root", root_x, root_y, ROOT_W, ROOT_H,
         fill=BRAND_PRIMARY, border=None)
    add_text(
        slide, "root-q", "Should we enter the LATAM market in FY26?",
        x_px=root_x + 12, y_px=root_y + 8, w_px=ROOT_W - 24, h_px=ROOT_H - 16,
        font_size_px=14, color=WHITE, bold=True, align="center", anchor="middle",
    )

    BR_W, BR_H = 240, 68
    br_y = 295
    yes_cx, no_cx = 360, 920

    _box(slide, "branch-yes", yes_cx - BR_W // 2, br_y, BR_W, BR_H,
         fill=BRAND_PRIMARY_MID)
    add_text(
        slide, "branch-yes-title", "YES -- Enter market",
        x_px=yes_cx - BR_W // 2 + 8, y_px=br_y + 6, w_px=BR_W - 16, h_px=22,
        font_size_px=13, color=WHITE, bold=True, align="center",
    )
    add_text(
        slide, "branch-yes-cond", "if local margin > 14%\nand currency hedge available",
        x_px=yes_cx - BR_W // 2 + 8, y_px=br_y + 28, w_px=BR_W - 16, h_px=BR_H - 32,
        font_size_px=11, color=BRAND_ACCENT_SOFT, align="center",
    )

    _box(slide, "branch-no", no_cx - BR_W // 2, br_y, BR_W, BR_H,
         fill=BRAND_PRIMARY_MID)
    add_text(
        slide, "branch-no-title", "NO -- Defer entry",
        x_px=no_cx - BR_W // 2 + 8, y_px=br_y + 6, w_px=BR_W - 16, h_px=22,
        font_size_px=13, color=WHITE, bold=True, align="center",
    )
    add_text(
        slide, "branch-no-cond", "if margin < 14%\nor FX volatility > 8%",
        x_px=no_cx - BR_W // 2 + 8, y_px=br_y + 28, w_px=BR_W - 16, h_px=BR_H - 32,
        font_size_px=11, color=BRAND_ACCENT_SOFT, align="center",
    )

    add_text(slide, "edge-yes", "Yes",
             x_px=root_cx - 130, y_px=root_y + ROOT_H + 16, w_px=60, h_px=20,
             font_size_px=11, color=TEXT_MID, italic=True, align="center")
    add_text(slide, "edge-no", "No",
             x_px=root_cx + 70, y_px=root_y + ROOT_H + 16, w_px=60, h_px=20,
             font_size_px=11, color=TEXT_MID, italic=True, align="center")

    root_bot_x, root_bot_y = root_cx, root_y + ROOT_H
    for bx, sid in ((yes_cx, "yes"), (no_cx, "no")):
        _line(slide, f"r2b-{sid}", root_bot_x, root_bot_y, bx, br_y,
              BRAND_PRIMARY_MID, width_pt=1.5)

    LF_W, LF_H = 215, 95
    leaf_y = 475

    leaves = [
        (170, yes_cx, "Full launch", "Open 2 plants, hire 180 FTE\nyear-1 revenue ~ $42M", False),
        (430, yes_cx, "Pilot region", "Single-country test in Mexico\nde-risk before full rollout", True),
        (740, no_cx,  "Wait 12 months", "Re-scan in FY27 once\ncurrency stabilises", False),
        (1080, no_cx, "Exit consideration", "Redirect capex to APAC\nor return to shareholders", False),
    ]

    for i, (cx, parent_cx, title, desc, is_accent) in enumerate(leaves):
        x = cx - LF_W // 2
        fill = BRAND_ACCENT if is_accent else WHITE
        border = BRAND_ACCENT if is_accent else CARD_BORDER
        title_color = WHITE if is_accent else BRAND_PRIMARY
        desc_color = WHITE if is_accent else TEXT_MID
        _box(slide, f"leaf-{i}", x, leaf_y, LF_W, LF_H,
             fill=fill, border=border, line_w=1.5)
        suffix = "  *" if is_accent else ""
        add_text(
            slide, f"leaf-title-{i}", title + suffix,
            x_px=x + 8, y_px=leaf_y + 8, w_px=LF_W - 16, h_px=22,
            font_size_px=13, color=title_color, bold=True, align="center",
        )
        add_text(
            slide, f"leaf-desc-{i}", desc,
            x_px=x + 8, y_px=leaf_y + 32, w_px=LF_W - 16, h_px=LF_H - 40,
            font_size_px=11, color=desc_color, align="center",
        )
        _line(slide, f"b2l-{i}", parent_cx, br_y + BR_H, cx, leaf_y,
              BRAND_PRIMARY_MID, width_pt=1.25)

    add_text(
        slide, "recommended-tag", "* Recommended path",
        x_px=64, y_px=620, w_px=240, h_px=20,
        font_size_px=11, color=BRAND_ACCENT, bold=True, italic=True,
    )

    add_footer(slide, page_num=4,
               source="Strategy team market-entry model v3.2",
               footnote="Decision thresholds calibrated to FY24 actuals; refresh quarterly")
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "diagram4-decisiontree.pptx"
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
