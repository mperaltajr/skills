"""Diagram 1 -- Org chart (3-level hierarchy)."""
from pathlib import Path
import sys
sys.path.insert(0, r"C:\Users\m.a.peralta\.claude\skills\slide-builder")
from twins.helpers import (
    new_slide, add_text, add_rect, add_title_block, add_footer,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    TEXT_DARK, TEXT_MID, TEXT_FAINT, CARD_BG, CARD_BORDER, WHITE,
    px_to_emu,
)
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt


def _box(slide, sid, x, y, w, h, fill, border=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu(x), px_to_emu(y), px_to_emu(w), px_to_emu(h),
    )
    shape.name = sid
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _line(slide, sid, x1, y1, x2, y2, color, width_pt=1.25):
    conn = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        px_to_emu(x1), px_to_emu(y1),
        px_to_emu(x2), px_to_emu(y2),
    )
    conn.name = sid
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn


def _node(slide, sid, x, y, w, h, title, role, *, fill, title_color, role_color, border=None):
    _box(slide, f"{sid}-bg", x, y, w, h, fill, border=border)
    add_text(
        slide, f"{sid}-title", title,
        x_px=x + 6, y_px=y + 8, w_px=w - 12, h_px=22,
        font_size_px=13, color=title_color, bold=True, align="center",
    )
    add_text(
        slide, f"{sid}-role", role,
        x_px=x + 6, y_px=y + 30, w_px=w - 12, h_px=h - 36,
        font_size_px=10, color=role_color, align="center",
    )


def build():
    prs, slide = new_slide()

    add_title_block(
        slide,
        title="Proposed <strong>leadership structure</strong> for the integrated operating model",
        subtitle="CEO, three division heads, and six functional managers under unified governance",
    )

    L1_W, L1_H = 220, 70
    L2_W, L2_H = 200, 70
    L3_W, L3_H = 175, 68

    canvas_cx = 640

    l1_y = 150
    l1_x = canvas_cx - L1_W // 2
    _node(
        slide, "ceo", l1_x, l1_y, L1_W, L1_H,
        title="J. Reeves",
        role="Chief Executive Officer\nAccountable for enterprise P&L",
        fill=BRAND_PRIMARY, title_color=WHITE, role_color=BRAND_ACCENT_SOFT,
    )

    l2_y = 310
    l2_centers = [280, 640, 1000]
    l2_specs = [
        ("ops",     "M. Carrillo", "VP, Operations\nManufacturing, supply chain, QA"),
        ("strat",   "A. Okafor",   "VP, Strategy\nCorporate dev and transformation"),
        ("fin",     "S. Yamamoto", "VP, Finance\nFP&A, treasury, controllership"),
    ]
    l2_positions = []
    for cx, (sid, name, role) in zip(l2_centers, l2_specs):
        x = cx - L2_W // 2
        _node(
            slide, sid, x, l2_y, L2_W, L2_H,
            title=name, role=role,
            fill=BRAND_PRIMARY_MID, title_color=WHITE, role_color=BRAND_ACCENT_SOFT,
        )
        l2_positions.append((cx, sid))

    l3_y = 480
    managers = [
        ("ops-mfg",  180,  "K. Patel",    "Mfg. Director\nPlants & throughput"),
        ("ops-sc",   385,  "R. Singh",    "Supply Chain Dir.\nLogistics & vendor mgmt"),
        ("str-cdev", 540,  "L. Bianchi",  "Corp Dev Lead\nM&A and partnerships"),
        ("str-tx",   745,  "D. Muller",   "Transformation Lead\nPMO and value capture"),
        ("fin-fpa",  900,  "T. Nguyen",   "FP&A Director\nForecast & planning"),
        ("fin-ctrl", 1105, "E. Brooks",   "Controller\nGL, close, statutory"),
    ]
    parent_map = {
        "ops-mfg": 280, "ops-sc": 280,
        "str-cdev": 640, "str-tx": 640,
        "fin-fpa": 1000, "fin-ctrl": 1000,
    }

    for sid, cx, name, role in managers:
        x = cx - L3_W // 2
        _node(
            slide, sid, x, l3_y, L3_W, L3_H,
            title=name, role=role,
            fill=WHITE, title_color=TEXT_DARK, role_color=TEXT_MID,
            border=CARD_BORDER,
        )

    ceo_bottom_y = l1_y + L1_H
    bus_y = l2_y - 22
    _line(slide, "ceo-stub", canvas_cx, ceo_bottom_y, canvas_cx, bus_y, BRAND_PRIMARY_MID)
    _line(slide, "l2-bus", l2_centers[0], bus_y, l2_centers[-1], bus_y, BRAND_PRIMARY_MID)
    for cx, _ in l2_positions:
        _line(slide, f"l2-stub-{cx}", cx, bus_y, cx, l2_y, BRAND_PRIMARY_MID)

    l3_bus_y = l3_y - 22
    by_parent = {}
    for sid, cx, *_ in managers:
        by_parent.setdefault(parent_map[sid], []).append((sid, cx))

    for parent_cx, children in by_parent.items():
        parent_bottom_y = l2_y + L2_H
        _line(slide, f"l3-pstub-{parent_cx}", parent_cx, parent_bottom_y, parent_cx, l3_bus_y, BRAND_PRIMARY_MID)
        xs = sorted(cx for _, cx in children)
        _line(slide, f"l3-bus-{parent_cx}", xs[0], l3_bus_y, xs[-1], l3_bus_y, BRAND_PRIMARY_MID)
        for sid, cx in children:
            _line(slide, f"l3-cstub-{sid}", cx, l3_bus_y, cx, l3_y, BRAND_PRIMARY_MID)

    add_footer(slide, page_num=1,
               source="Proposed Q3 operating model, OneCo executive workshop",
               footnote="Reporting lines reflect proposed Day-1 structure; final roles subject to board approval")
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "diagram1-orgchart.pptx"
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
