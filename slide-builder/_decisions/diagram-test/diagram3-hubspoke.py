"""Diagram 3 -- Hub-and-spoke (Porter 5-force variant, 4 compass satellites)."""
from pathlib import Path
import sys
import math
sys.path.insert(0, r"C:\Users\m.a.peralta\.claude\skills\slide-builder")
from twins.helpers import (
    new_slide, add_text, add_rect, add_title_block, add_footer,
    BRAND_PRIMARY, BRAND_PRIMARY_MID, BRAND_ACCENT, BRAND_ACCENT_SOFT,
    TEXT_DARK, TEXT_MID, TEXT_FAINT, CARD_BG, CARD_BORDER, WHITE,
    px_to_emu,
)
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt


def _circle(slide, sid, cx, cy, d, fill, border=None, border_w=1.5):
    x = cx - d // 2
    y = cy - d // 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px_to_emu(x), px_to_emu(y), px_to_emu(d), px_to_emu(d),
    )
    shape.name = sid
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_w)
    else:
        shape.line.fill.background()
    return shape


def _arrow_to(slide, sid, x1, y1, x2, y2, color, width_pt=2.0):
    conn = slide.shapes.add_connector(
        1,
        px_to_emu(x1), px_to_emu(y1),
        px_to_emu(x2), px_to_emu(y2),
    )
    conn.name = sid
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


def build():
    prs, slide = new_slide()

    add_title_block(
        slide,
        title="Four competitive forces <strong>squeeze industry margins</strong> simultaneously",
        subtitle="Porter framing applied to the specialty chemicals sector, FY25 outlook",
    )

    CX = 640
    CY = 410

    HUB_D = 170
    SAT_D = 145

    OFFSET_V = 200
    OFFSET_H = 290
    positions = {
        "top":    (CX, CY - OFFSET_V),
        "right":  (CX + OFFSET_H, CY),
        "bottom": (CX, CY + OFFSET_V),
        "left":   (CX - OFFSET_H, CY),
    }

    sats = [
        ("top",    "New Entrants", "Low capex barriers attract\nregional Asian producers"),
        ("right",  "Buyers",       "Top-5 customers concentrate\n62% of revenue and dictate price"),
        ("bottom", "Substitutes",  "Bio-based polymers gain share\nin 3 of 8 end markets"),
        ("left",   "Suppliers",    "Two feedstock vendors control\n80% of upstream supply"),
    ]

    for key, _, _ in sats:
        sx, sy = positions[key]
        angle = math.atan2(CY - sy, CX - sx)
        start_x = sx + (SAT_D / 2) * math.cos(angle)
        start_y = sy + (SAT_D / 2) * math.sin(angle)
        end_x = CX - (HUB_D / 2) * math.cos(angle)
        end_y = CY - (HUB_D / 2) * math.sin(angle)
        _arrow_to(slide, f"arrow-{key}", start_x, start_y, end_x, end_y,
                  BRAND_PRIMARY_MID, width_pt=2.25)

    for key, title, desc in sats:
        sx, sy = positions[key]
        _circle(slide, f"sat-{key}", sx, sy, SAT_D,
                fill=CARD_BG, border=CARD_BORDER, border_w=1.5)
        add_text(
            slide, f"sat-title-{key}", title,
            x_px=sx - SAT_D // 2 + 8, y_px=sy - SAT_D // 2 + 32,
            w_px=SAT_D - 16, h_px=22,
            font_size_px=13, color=BRAND_PRIMARY, bold=True, align="center",
        )
        add_text(
            slide, f"sat-desc-{key}", desc,
            x_px=sx - SAT_D // 2 + 8, y_px=sy - SAT_D // 2 + 58,
            w_px=SAT_D - 16, h_px=SAT_D - 74,
            font_size_px=10, color=TEXT_MID, align="center",
        )

    _circle(slide, "hub", CX, CY, HUB_D, fill=BRAND_PRIMARY)
    add_text(
        slide, "hub-label", "Our Firm",
        x_px=CX - HUB_D // 2, y_px=CY - 24,
        w_px=HUB_D, h_px=24,
        font_size_px=18, color=WHITE, bold=True, align="center",
    )
    add_text(
        slide, "hub-sub", "Specialty\nChemicals",
        x_px=CX - HUB_D // 2, y_px=CY + 2,
        w_px=HUB_D, h_px=44,
        font_size_px=11, color=BRAND_ACCENT_SOFT, align="center",
    )

    add_footer(slide, page_num=3,
               source="Industry analysis, FY25 strategic review",
               footnote="Force intensity scored qualitatively on management interviews and IBISWorld data")
    return prs


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "diagram3-hubspoke.pptx"
    prs = build()
    prs.save(str(out))
    print(f"Wrote: {out}")
