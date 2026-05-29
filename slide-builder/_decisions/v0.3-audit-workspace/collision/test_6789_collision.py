"""Tests 6, 7, 8, 9 — collision detector unit tests."""
import sys
from pathlib import Path

SKILL = Path(r"C:\Users\m.a.peralta\.claude\skills\.claude\worktrees\agent-abaae9e70ed8f6544\slide-builder")
sys.path.insert(0, str(SKILL))
sys.path.insert(0, str(SKILL / "scripts"))

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import finalize_deck
from finalize_deck import (
    _check_dark_variant_collisions,
    _rgb_distance,
    _DARK_VARIANT_COLLISION_THRESHOLD,
    DarkVariantCollisionError,
)

DARK_BG = "4D148C"

def _new_slide():
    prs = Presentation()
    prs.slide_width = Emu(1280 * 9525)
    prs.slide_height = Emu(720 * 9525)
    blank = None
    for L in prs.slide_layouts:
        try:
            if L.name == "Blank":
                blank = L
                break
        except Exception:
            pass
    if blank is None:
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)
    # Remove inherited placeholders
    sp_tree = slide.shapes._spTree
    for shp in list(slide.shapes):
        try:
            sp_tree.remove(shp.element)
        except Exception:
            pass
    return prs, slide

def _add_overlay(slide):
    """Add the dark overlay (should be skipped by the detector)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Emu(1280 * 9525), Emu(720 * 9525))
    s.name = "dark-variant-overlay"
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x4D, 0x14, 0x8C)
    s.line.fill.background()
    return s

def _add_text(slide, text, color_rgb, fill_rgb=None, name="tb"):
    tb = slide.shapes.add_textbox(Emu(100*9525), Emu(100*9525),
                                    Emu(400*9525), Emu(80*9525))
    tb.name = name
    if fill_rgb is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill_rgb
    else:
        tb.fill.background()
    tb.text_frame.text = text
    p = tb.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].font.color.rgb = color_rgb
    return tb

def _add_filled_rect(slide, fill_rgb, name="rect"):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(50*9525), Emu(200*9525),
                                 Emu(300*9525), Emu(80*9525))
    s.name = name
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    s.line.fill.background()
    return s

print("=" * 70)
print("TEST 9 — threshold sanity")
print("=" * 70)
d0 = _rgb_distance("4D148C", "4D148C")
d_far = _rgb_distance("4D148C", "FF6600")
d_near = _rgb_distance("4D148C", "4F1890")
print(f"  dist('4D148C','4D148C') = {d0:.3f} (expect 0)")
print(f"  dist('4D148C','FF6600') = {d_far:.3f} (expect >200)")
print(f"  dist('4D148C','4F1890') = {d_near:.3f} (expect <10)")
assert d0 == 0.0
assert d_far > 200
assert d_near < 10
print(f"  threshold = {_DARK_VARIANT_COLLISION_THRESHOLD}")
assert _DARK_VARIANT_COLLISION_THRESHOLD == 40.0
print("TEST 9: PASS")

print("\n" + "=" * 70)
print("TEST 6 - clean dark slide (build_example.py style) -> 0 issues")
print("=" * 70)
prs, slide = _new_slide()
_add_overlay(slide)
# White bars on purple bg (like build_example slide 2)
_add_filled_rect(slide, RGBColor(0xFF, 0xFF, 0xFF), name="bar-white-0")
_add_filled_rect(slide, RGBColor(0xFF, 0x66, 0x00), name="bar-orange-0")
# Purple text inside the white bar — should be SKIPPED (white has solid fill)
_add_text(slide, "label inside bar", RGBColor(0x4D, 0x14, 0x8C),
            fill_rgb=RGBColor(0xFF, 0xFF, 0xFF), name="label-purple-on-white")
# White text on no-fill text box — safe (white is far from purple)
_add_text(slide, "white text on dark", RGBColor(0xFF, 0xFF, 0xFF),
            fill_rgb=None, name="white-tb")
issues_clean = _check_dark_variant_collisions(slide, DARK_BG, 1)
print(f"  issues: {len(issues_clean)}")
for it in issues_clean:
    print(f"    - {it}")
assert len(issues_clean) == 0, f"clean slide should be clean, got {issues_clean}"
print("TEST 6: PASS")

print("\n" + "=" * 70)
print("TEST 7a - colliding slide (shape fill + text on no-fill box) -> 2 issues")
print("=" * 70)
prs, slide = _new_slide()
_add_overlay(slide)
# A purple-filled shape on purple bg — collision
_add_filled_rect(slide, RGBColor(0x4D, 0x14, 0x8C), name="purple-rect")
# Purple text on no-fill text box — collision
_add_text(slide, "invisible purple text",
            RGBColor(0x4D, 0x14, 0x8C), fill_rgb=None, name="purple-tb-nofill")
issues_collide = _check_dark_variant_collisions(slide, DARK_BG, 1)
print(f"  issues: {len(issues_collide)}")
for it in issues_collide:
    print(f"    - {it}")
assert len(issues_collide) == 2, f"expected exactly 2 issues, got {len(issues_collide)}"
print("TEST 7a: PASS (unit)")

print("\n" + "=" * 70)
print("TEST 8 - false-positive guard (purple text on white-filled box) -> 0 issues")
print("=" * 70)
prs, slide = _new_slide()
_add_overlay(slide)
# White-filled text box with purple text inside — text is on white, not on dark bg
_add_text(slide, "purple on white bar",
            RGBColor(0x4D, 0x14, 0x8C), fill_rgb=RGBColor(0xFF, 0xFF, 0xFF),
            name="white-bar-with-purple-text")
issues_fp = _check_dark_variant_collisions(slide, DARK_BG, 1)
print(f"  issues: {len(issues_fp)}")
for it in issues_fp:
    print(f"    - {it}")
assert len(issues_fp) == 0, f"false-positive: expected 0, got {issues_fp}"
print("TEST 8: PASS")

print("\nAll unit tests passed.")
