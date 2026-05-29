"""Show what's in the themed colliding pptx and what the detector sees."""
import sys
from pathlib import Path
from pptx import Presentation

SKILL = Path(r"C:\Users\m.a.peralta\.claude\skills\.claude\worktrees\agent-abaae9e70ed8f6544\slide-builder")
sys.path.insert(0, str(SKILL))
sys.path.insert(0, str(SKILL / "scripts"))

from finalize_deck import _check_dark_variant_collisions

THEMED = Path(r"C:\Users\m.a.peralta\.claude\skills\.claude\worktrees\agent-abaae9e70ed8f6544\slide-builder\_decisions\v0.3-audit-workspace\collision\build2\slide_01\option_A.pptx")
RAW = Path(r"C:\Users\m.a.peralta\.claude\skills\.claude\worktrees\agent-abaae9e70ed8f6544\slide-builder\_decisions\v0.3-audit-workspace\collision\build2\slide_01\_raw\option_A.pptx")

print("=== RAW colliding option_A (pre-theme) ===")
prs = Presentation(str(RAW))
slide = prs.slides[0]
print(f"shapes: {len(slide.shapes)}")
for i, s in enumerate(slide.shapes):
    nm = (s.name or "?").strip()
    print(f"  [{i}] {nm!r} type={s.shape_type}")

print("\n=== THEMED colliding option_A (post-finalize) ===")
prs = Presentation(str(THEMED))
slide = prs.slides[0]
print(f"shapes: {len(slide.shapes)}")
for i, s in enumerate(slide.shapes):
    nm = (s.name or "?").strip()
    txt = ""
    try:
        if s.has_text_frame:
            txt = (s.text_frame.text or "")[:50]
    except Exception:
        pass
    print(f"  [{i}] {nm!r} type={s.shape_type} text={txt!r}")

issues = _check_dark_variant_collisions(slide, "4D148C", 1)
print(f"\n_check_dark_variant_collisions issues: {len(issues)}")
for it in issues:
    print(f"  - {it}")
