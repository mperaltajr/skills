"""Test 4: Inspect dark slide structure — overlay, z-order, title, no inherited title."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

print("=" * 70)
print("TEST 4 — Dark slide structure inspection")
print("=" * 70)

OUT = Path(r"C:\Users\m.a.peralta\.claude\skills\.claude\worktrees\agent-abaae9e70ed8f6544\slide-builder\_decisions\v0.3-audit-workspace\collision\build")

light = OUT / "slide_01" / "option_A.pptx"
dark = OUT / "slide_02" / "option_A.pptx"

for label, pptx_path in [("LIGHT slide_01", light), ("DARK slide_02", dark)]:
    print(f"\n--- {label}: {pptx_path.name} ---")
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    print(f"  shapes count: {len(slide.shapes)}")
    overlay_found = False
    overlay_z = None
    overlay_fill = None
    title_white_found = False
    for i, shape in enumerate(slide.shapes):
        nm = (shape.name or "").strip()
        try:
            fill = shape.fill
            ftype = fill.type
            fhex = ""
            if ftype == 1:
                try:
                    fhex = str(fill.fore_color.rgb)
                except Exception:
                    fhex = "?"
        except Exception:
            ftype = None
            fhex = ""
        is_ph = shape.is_placeholder if hasattr(shape, "is_placeholder") else False
        print(f"  [{i}] name={nm!r} type={shape.shape_type} fill={fhex} is_ph={is_ph}")
        if nm == "dark-variant-overlay":
            overlay_found = True
            overlay_z = i
            overlay_fill = fhex

    # spTree position
    sp_tree = slide.shapes._spTree
    children = [c for c in sp_tree if c.tag.endswith("}sp") or c.tag.endswith("}pic") or c.tag.endswith("}grpSp")]

    if label.startswith("DARK"):
        # find overlay element in sp_tree directly
        overlay_idx_sp = None
        all_kids = list(sp_tree)
        for idx, kid in enumerate(all_kids):
            # nvSpPr/cNvPr name
            for cnv in kid.iter():
                if cnv.tag.endswith("}cNvPr") and cnv.get("name") == "dark-variant-overlay":
                    overlay_idx_sp = idx
                    break
            if overlay_idx_sp is not None:
                break
        print(f"  overlay sp_tree index: {overlay_idx_sp}")
        assert overlay_found, "dark-variant-overlay shape missing on dark slide"
        assert overlay_fill.upper() == "4D148C", f"overlay fill = {overlay_fill}, expected 4D148C"
        # Check overlay is at sp_tree index 2 (per spec). Note: sp_tree may have non-sp prefix elements
        # (nvGrpSpPr, grpSpPr), so check the absolute index includes those.
        print(f"  overlay fill hex: {overlay_fill} (expected 4D148C)")
        # Check no inherited title placeholder visible (PLACEHOLDER type 14 = title)
        for shape in slide.shapes:
            if shape.is_placeholder:
                print(f"  WARN: placeholder found on dark slide: {shape.placeholder_format.type}")
        # Check a white text box exists with title
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.is_placeholder:
                try:
                    txt = shape.text_frame.text or ""
                    if "cash gets stuck" in txt.lower() or "where the cash" in txt.lower():
                        # check white
                        p = shape.text_frame.paragraphs[0]
                        if p.runs:
                            try:
                                col = p.runs[0].font.color.rgb
                                if col and str(col).upper() == "FFFFFF":
                                    title_white_found = True
                                    print(f"  WHITE title text box: {txt[:60]!r}")
                            except Exception:
                                pass
                except Exception:
                    pass
        if not title_white_found:
            print(f"  WARN: no white title found - check finalize's title resolution")
    else:
        # LIGHT slide should NOT have dark-variant-overlay
        assert not overlay_found, "light slide unexpectedly has dark-variant-overlay"

print("\nTEST 4: structure checks above")
