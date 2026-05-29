"""Dump dark slide internals to understand what happened."""
import sys
from pathlib import Path
from pptx import Presentation

OUT = Path(r"C:\Users\m.a.peralta\.claude\skills\.claude\worktrees\agent-abaae9e70ed8f6544\slide-builder\_decisions\v0.3-audit-workspace\collision\build")
dark = OUT / "slide_02" / "option_A.pptx"

# Compare to raw (pre-theme) pptx
raw_dark = OUT / "slide_02" / "_raw" / "option_A.pptx"

print("=== RAW dark (pre-theme) ===")
if raw_dark.exists():
    prs = Presentation(str(raw_dark))
    slide = prs.slides[0]
    print(f"shapes: {len(slide.shapes)}")
    for i, s in enumerate(slide.shapes):
        print(f"  [{i}] {(s.name or '?').strip()!r} type={s.shape_type}")
else:
    print("raw not stashed yet")

print("\n=== THEMED dark (post-finalize) ===")
prs = Presentation(str(dark))
slide = prs.slides[0]
print(f"shapes: {len(slide.shapes)}")
for i, s in enumerate(slide.shapes):
    nm = (s.name or "?").strip()
    txt = ""
    try:
        if s.has_text_frame:
            txt = (s.text_frame.text or "")[:60]
    except Exception:
        pass
    print(f"  [{i}] {nm!r} type={s.shape_type} text={txt!r}")

# Print sp_tree raw XML for the dark slide
print("\n=== sp_tree element names ===")
sp_tree = slide.shapes._spTree
for i, kid in enumerate(sp_tree):
    name_attr = ""
    for cnv in kid.iter():
        if cnv.tag.endswith("}cNvPr") and cnv.get("name"):
            name_attr = cnv.get("name")
            break
    print(f"  [{i}] tag={kid.tag.split('}')[-1]} name={name_attr!r}")
